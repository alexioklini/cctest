"""
doc_convert.py — Convert binary documents (PDF, DOCX) to markdown so the
MemPalace miner (which only reads .md/.txt/code extensions) can pick them
up automatically.

Pre-mine pass for the project-sync daemon: walks each input folder, finds
unsupported binary files, writes companion .md files into a hidden
`.brain-extracted/` subdirectory. Idempotent via (mtime, size) hash so
re-runs are cheap and dedup naturally.

Public surface:
    SUPPORTED_EXTS               — extensions we convert (case-insensitive)
    convert_folder(root) -> dict — sync_status-style summary
    sweep_stale(root)    -> int  — drop converted .md files whose source vanished
    extract_md_path(root, source) -> str — resolve where a source file's md lives

Frontmatter convention written to every produced .md:
    <!-- brain-source: /abs/path/to/original.pdf -->
    <!-- brain-source-mtime: 1730000000 -->
    <!-- brain-source-size: 12345 -->

The frontmatter is HTML-comment so it survives the markdown miner (which
treats lines as text) without polluting drawer content noticeably; the
source-resolver in claude_cli's _build_system_prompt block can grep for
`brain-source:` to map drawer source_files back to original PDFs.

Failures in conversion are isolated per-file: one broken PDF doesn't break
the whole pass.
"""
from __future__ import annotations

import os
import time
from dataclasses import dataclass, field
from typing import Iterable

EXTRACT_SUBDIR = ".brain-extracted"

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".eml", ".msg"}

# Per-sheet row policy for xlsx extraction. Bank operations rely on
# medium-to-large lookup tables (retention schedules, role matrices,
# tariff lists) where the rows ARE the policy — a small preview throws
# away exactly the information we want indexed. So we extract every
# row, with two safety guards:
#   * XLSX_WARN_ROWS_PER_SHEET — soft warn in daemon log if exceeded
#     (helps spot accidental CSV-export-as-xlsx file drops).
#   * XLSX_MAX_ROWS_PER_SHEET  — hard cap at extraction time. Anything
#     this large is a database export, not a policy artifact, and the
#     drawer storage would be useless anyway (each row mostly numeric).
XLSX_WARN_ROWS_PER_SHEET = 5_000
XLSX_MAX_ROWS_PER_SHEET = 100_000
# Hard cap per cell so a single megablob doesn't blow up a row.
XLSX_CELL_MAX_CHARS = 200

# Files we should never try to convert even if their extension matches —
# guards against name collisions if a user has e.g. a docx named "file.docx"
# but it's actually corrupt.
SKIP_NAMES = {".DS_Store"}


@dataclass
class ConvertResult:
    converted: int = 0
    skipped_unchanged: int = 0
    skipped_unreadable: int = 0
    failed: int = 0
    stale_removed: int = 0
    seen_total: int = 0
    failures: list[str] = field(default_factory=list)
    elapsed_s: float = 0.0


def _stat_key(path: str) -> tuple[int, int]:
    """Cheap change detector: (mtime_ns, size). Re-conversion only happens
    when one of these changes."""
    try:
        st = os.stat(path)
        return int(st.st_mtime), int(st.st_size)
    except OSError:
        return 0, 0


def _read_md_source_stat(md_path: str) -> tuple[int, int]:
    """Pull the brain-source-mtime / brain-source-size from a converted
    file's frontmatter. Returns (0, 0) if absent."""
    try:
        with open(md_path, "r", encoding="utf-8", errors="replace") as f:
            head = f.read(2048)
    except OSError:
        return 0, 0
    mt = sz = 0
    for line in head.splitlines():
        line = line.strip()
        if line.startswith("<!-- brain-source-mtime:"):
            try:
                mt = int(line.split(":", 1)[1].rstrip(" -->").strip())
            except ValueError:
                pass
        elif line.startswith("<!-- brain-source-size:"):
            try:
                sz = int(line.split(":", 1)[1].rstrip(" -->").strip())
            except ValueError:
                pass
        if mt and sz:
            break
    return mt, sz


def _md_path_for(root: str, source_path: str) -> str:
    """Compute the converted-md path under <root>/.brain-extracted/.
    The relative layout under .brain-extracted/ mirrors the source's
    relative layout under root, with .md appended to the original name
    (so collisions across extensions like file.pdf vs file.docx don't
    overwrite each other).
    """
    rel = os.path.relpath(source_path, root)
    base = os.path.join(root, EXTRACT_SUBDIR, rel) + ".md"
    return base


def _source_for_md(root: str, md_path: str) -> str:
    """Inverse of _md_path_for — used only by sweep_stale to detect orphans."""
    rel = os.path.relpath(md_path, os.path.join(root, EXTRACT_SUBDIR))
    if rel.endswith(".md"):
        rel = rel[:-3]
    return os.path.join(root, rel)


def _frontmatter(source_path: str, mtime: int, size: int) -> str:
    return (
        f"<!-- brain-source: {source_path} -->\n"
        f"<!-- brain-source-mtime: {mtime} -->\n"
        f"<!-- brain-source-size: {size} -->\n\n"
    )


# ── Per-format extractors ────────────────────────────────────────────────────

def _extract_pdf(path: str) -> tuple[str, str | None]:
    """Returns (markdown_text, error_msg or None). Empty text + None means
    the PDF had no extractable text (likely scanned image). Empty text +
    error means a real failure."""
    try:
        import fitz  # type: ignore
    except ImportError:
        return "", "pymupdf not installed (pip install pymupdf)"
    try:
        doc = fitz.open(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    parts = []
    title = (doc.metadata or {}).get("title") or os.path.splitext(os.path.basename(path))[0]
    parts.append(f"# {title}\n")
    try:
        for i, page in enumerate(doc, 1):
            t = page.get_text("text") or ""
            t = t.strip()
            if t:
                parts.append(f"\n\n## Page {i}\n\n{t}\n")
    finally:
        doc.close()
    text = "\n".join(parts).strip()
    if len(text.splitlines()) <= 1:
        # Only the title line — likely scanned PDF or empty.
        return "", None
    return text + "\n", None


def _extract_docx(path: str) -> tuple[str, str | None]:
    try:
        import docx  # type: ignore
    except ImportError:
        return "", "python-docx not installed (pip install python-docx)"
    try:
        d = docx.Document(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    title = os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]
    for p in d.paragraphs:
        t = (p.text or "").rstrip()
        if not t:
            parts.append("")
            continue
        # Heading levels: bold + first run + style starting with "Heading"
        style_name = (getattr(p.style, "name", "") or "").lower()
        if style_name.startswith("heading"):
            try:
                level = int(style_name.replace("heading", "").strip() or "2")
            except ValueError:
                level = 2
            level = max(2, min(6, level))
            parts.append(f"\n{'#' * level} {t}\n")
        else:
            parts.append(t)
    text = "\n".join(parts).strip()
    if len(text.splitlines()) <= 1:
        return "", None
    return text + "\n", None


def _extract_pptx(path: str) -> tuple[str, str | None]:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return "", "python-pptx not installed (pip install python-pptx)"
    try:
        prs = Presentation(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    title = os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]
    slide_count = 0
    for i, slide in enumerate(prs.slides, 1):
        slide_count += 1
        # Try to surface the slide title from its title placeholder when
        # present; fall back to "Slide N".
        slide_title = ""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_title = (slide.shapes.title.text_frame.text or "").strip()
        except Exception:
            slide_title = ""
        heading = f"## Slide {i}" + (f" — {slide_title}" if slide_title else "")
        parts.append(f"\n{heading}\n")
        for shape in slide.shapes:
            # Skip the title shape we already used.
            if shape == slide.shapes.title:
                continue
            try:
                if not shape.has_text_frame:
                    continue
            except Exception:
                continue
            for para in shape.text_frame.paragraphs:
                t = (para.text or "").rstrip()
                if not t:
                    continue
                # Indent bullets a bit so list structure carries through.
                level = getattr(para, "level", 0) or 0
                bullet = "  " * level + "- " if level >= 0 else ""
                parts.append(f"{bullet}{t}")
        # Speaker notes — often hold the actual policy/explanation text.
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"\n_Speaker notes:_ {notes}")
        except Exception:
            pass
    text = "\n".join(parts).strip()
    # Title + 0 slides → empty.
    if slide_count == 0 or len(text.splitlines()) <= 1:
        return "", None
    return text + "\n", None


def _extract_xlsx(path: str) -> tuple[str, str | None]:
    """Render every sheet's header + first N rows as a markdown table.

    Spreadsheets are mostly numeric/tabular and don't fit the drawer model
    cleanly. We extract a structured preview per sheet — enough for the
    embedding to learn the topic and the KG extractor to see e.g. retention
    schedules, role matrices, threshold tables — and stop. Users who need
    full-fidelity row-by-row analysis should call read_document on the
    original .xlsx rather than relying on the converted markdown.
    """
    try:
        import openpyxl  # type: ignore
    except ImportError:
        return "", "openpyxl not installed (pip install openpyxl)"
    try:
        # data_only=True so cached formula values come through (the user
        # sees the displayed value, not the formula expression).
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    title = os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]

    def _cell(v):
        if v is None:
            return ""
        s = str(v)
        if len(s) > XLSX_CELL_MAX_CHARS:
            s = s[:XLSX_CELL_MAX_CHARS] + "…"
        # Markdown table separator + escape pipes inside cells.
        return s.replace("|", "\\|").replace("\n", " ").strip()

    sheet_count = 0
    for sheet in wb.worksheets:
        try:
            rows = sheet.iter_rows(values_only=True)
            header = next(rows, None)
            if header is None:
                continue
            sheet_count += 1
            parts.append(f"\n## Sheet: {sheet.title}\n")
            cells = [_cell(c) for c in header]
            # Drop trailing empty columns so the table doesn't have ragged
            # padding from sparsely-used right side of the sheet.
            while cells and cells[-1] == "":
                cells.pop()
            if not cells:
                # Empty header row, skip; sheet may be data-without-header.
                cells = ["col_1"]
            n_cols = len(cells)
            parts.append("| " + " | ".join(cells) + " |")
            parts.append("|" + "|".join(["---"] * n_cols) + "|")
            shown = 0
            warned = False
            for r in rows:
                if shown >= XLSX_MAX_ROWS_PER_SHEET:
                    parts.append(
                        f"\n_(truncated at {XLSX_MAX_ROWS_PER_SHEET:,} "
                        f"rows — use read_document on the original .xlsx "
                        f"for the rest)_\n")
                    break
                rc = [_cell(c) for c in r[:n_cols]]
                # Pad short rows to keep table valid.
                while len(rc) < n_cols:
                    rc.append("")
                if all(c == "" for c in rc):
                    continue  # skip blank rows
                parts.append("| " + " | ".join(rc) + " |")
                shown += 1
                if not warned and shown == XLSX_WARN_ROWS_PER_SHEET:
                    print(
                        f"[doc-convert] {path}: sheet `{sheet.title}` "
                        f"has >{XLSX_WARN_ROWS_PER_SHEET:,} rows — "
                        f"continuing extraction up to "
                        f"{XLSX_MAX_ROWS_PER_SHEET:,}",
                        flush=True)
                    warned = True
        except Exception as e:
            parts.append(f"\n_(failed to read sheet `{sheet.title}`: "
                         f"{type(e).__name__})_\n")
    try:
        wb.close()
    except Exception:
        pass
    if sheet_count == 0:
        return "", None
    return "\n".join(parts) + "\n", None


def _extract_eml(path: str) -> tuple[str, str | None]:
    """Parse RFC 822 / MIME messages. Stdlib `email` only — no extras."""
    import email
    from email import policy as _policy
    try:
        with open(path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=_policy.default)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    return _email_to_markdown(msg), None


def _extract_msg(path: str) -> tuple[str, str | None]:
    """Outlook .msg files. Requires `extract-msg`; if absent, returns a
    helpful error rather than silently skipping."""
    try:
        import extract_msg  # type: ignore
    except ImportError:
        return "", "extract-msg not installed (pip install extract-msg)"
    try:
        m = extract_msg.Message(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    title = m.subject or os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]
    if m.sender: parts.append(f"**From:** {m.sender}")
    if m.to: parts.append(f"**To:** {m.to}")
    if m.cc: parts.append(f"**Cc:** {m.cc}")
    if m.date: parts.append(f"**Date:** {m.date}")
    if m.subject: parts.append(f"**Subject:** {m.subject}")
    parts.append("")
    body = (m.body or "").strip()
    if body:
        parts.append(body)
    try:
        m.close()
    except Exception:
        pass
    text = "\n".join(parts).strip()
    if len(text.splitlines()) <= 1:
        return "", None
    return text + "\n", None


def _email_to_markdown(msg) -> str:
    """Shared rendering for both .eml (parsed via stdlib email) and any
    other future `email.message.Message` source."""
    title = msg.get("Subject") or "(no subject)"
    parts = [f"# {title}\n"]
    for hdr in ("From", "To", "Cc", "Date", "Subject"):
        v = msg.get(hdr)
        if v:
            parts.append(f"**{hdr}:** {v}")
    parts.append("")
    body = ""
    try:
        # Prefer the plain text body part.
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type() or ""
                disp = (part.get("Content-Disposition") or "").lower()
                if "attachment" in disp:
                    continue
                if ctype == "text/plain":
                    body = part.get_content()
                    break
            if not body:
                # Fall back to html → strip tags lightly.
                for part in msg.walk():
                    if part.get_content_type() == "text/html":
                        import re
                        html = part.get_content()
                        body = re.sub(r"<[^>]+>", "", html or "")
                        break
        else:
            body = msg.get_content() or ""
    except Exception:
        body = ""
    body = (body or "").strip()
    if body:
        parts.append(body)
    return "\n".join(parts) + "\n"


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".eml": _extract_eml,
    ".msg": _extract_msg,
}


def _iter_source_files(root: str) -> Iterable[str]:
    """Yield every supported source path in `root` (recursive). Skips the
    .brain-extracted subdir itself."""
    skip_prefix = os.path.join(root, EXTRACT_SUBDIR)
    for dirpath, _dirs, files in os.walk(root):
        if dirpath == skip_prefix or dirpath.startswith(skip_prefix + os.sep):
            continue
        for fn in files:
            if fn in SKIP_NAMES:
                continue
            ext = os.path.splitext(fn)[1].lower()
            if ext in SUPPORTED_EXTS:
                yield os.path.join(dirpath, fn)


# ── Public API ───────────────────────────────────────────────────────────────

def convert_folder(root: str, *, log_prefix: str = "[doc-convert]") -> ConvertResult:
    """Walk `root` for supported binary docs, write `.md` siblings under
    `<root>/.brain-extracted/<rel>.md`. Idempotent via (mtime, size) hash.
    Returns a result summary; does not raise on per-file failures.
    """
    res = ConvertResult()
    if not root or not os.path.isdir(root):
        return res
    t0 = time.time()
    extract_root = os.path.join(root, EXTRACT_SUBDIR)
    for src in _iter_source_files(root):
        res.seen_total += 1
        ext = os.path.splitext(src)[1].lower()
        extractor = _EXTRACTORS.get(ext)
        if extractor is None:
            continue
        cur_mt, cur_sz = _stat_key(src)
        if cur_mt == 0 and cur_sz == 0:
            res.skipped_unreadable += 1
            continue
        md_path = _md_path_for(root, src)
        if os.path.isfile(md_path):
            old_mt, old_sz = _read_md_source_stat(md_path)
            if old_mt == cur_mt and old_sz == cur_sz:
                res.skipped_unchanged += 1
                continue
        # (Re)convert.
        try:
            os.makedirs(os.path.dirname(md_path), exist_ok=True)
        except OSError as e:
            res.failed += 1
            res.failures.append(f"{src}: mkdir failed {e}")
            continue
        try:
            text, err = extractor(src)
        except Exception as e:
            res.failed += 1
            res.failures.append(f"{src}: {type(e).__name__}: {e}")
            continue
        if err:
            res.failed += 1
            res.failures.append(f"{src}: {err}")
            continue
        if not text:
            # Empty after extraction — e.g. scanned PDF. Write a marker file
            # so we don't retry every cycle, but keep it small enough that the
            # miner doesn't flood drawer storage.
            text = (
                "# " + os.path.splitext(os.path.basename(src))[0] + "\n\n"
                "_(no extractable text — possibly a scanned image; OCR not run)_\n"
            )
        body = _frontmatter(src, cur_mt, cur_sz) + text
        try:
            tmp_path = md_path + ".tmp"
            with open(tmp_path, "w", encoding="utf-8") as f:
                f.write(body)
            os.replace(tmp_path, md_path)
            res.converted += 1
        except OSError as e:
            res.failed += 1
            res.failures.append(f"{src}: write failed {e}")
            continue

    res.elapsed_s = time.time() - t0
    if res.converted or res.failed:
        print(
            f"{log_prefix} {root}: converted={res.converted} "
            f"unchanged={res.skipped_unchanged} failed={res.failed} "
            f"seen={res.seen_total} elapsed={res.elapsed_s:.1f}s",
            flush=True)
    return res


def sweep_stale(root: str, *, log_prefix: str = "[doc-convert]") -> int:
    """Walk `<root>/.brain-extracted/` and drop any .md whose source path
    no longer exists. Returns count removed.
    """
    extract_root = os.path.join(root, EXTRACT_SUBDIR)
    if not os.path.isdir(extract_root):
        return 0
    removed = 0
    for dirpath, _dirs, files in os.walk(extract_root):
        for fn in files:
            if not fn.endswith(".md"):
                continue
            md_path = os.path.join(dirpath, fn)
            src = _source_for_md(root, md_path)
            if not os.path.isfile(src):
                try:
                    os.remove(md_path)
                    removed += 1
                except OSError:
                    pass
    if removed:
        print(f"{log_prefix} {root}: stale_removed={removed}", flush=True)
    return removed


def resolve_original_source(md_path: str) -> str | None:
    """Given a converted .md path, read its frontmatter and return the
    original source path. None if no marker present."""
    try:
        with open(md_path, "r", encoding="utf-8", errors="replace") as f:
            head = f.read(1024)
    except OSError:
        return None
    for line in head.splitlines():
        line = line.strip()
        if line.startswith("<!-- brain-source:"):
            return line.split(":", 1)[1].rstrip(" -->").strip() or None
    return None
