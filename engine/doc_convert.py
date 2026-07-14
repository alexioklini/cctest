"""
doc_convert.py — Convert binary documents (PDF, DOCX) to markdown so the
MemPalace miner (which only reads .md/.txt/code extensions) can pick them
up automatically.

Pre-mine pass for the project-sync daemon: walks each input folder, finds
unsupported binary files, writes companion .md files into a hidden
`.brain-extracted/` subdirectory. Idempotent via (mtime, size) hash so
re-runs are cheap and dedup naturally.

Public surface:
    SUPPORTED_EXTS               — extensions we convert (case-insensitive)
    convert_folder(root) -> dict — sync_status-style summary
    sweep_stale(root)    -> int  — drop converted .md files whose source vanished
    extract_md_path(root, source) -> str — resolve where a source file's md lives

Frontmatter convention written to every produced .md:
    <!-- brain-source: /abs/path/to/original.pdf -->
    <!-- brain-source-mtime: 1730000000 -->
    <!-- brain-source-size: 12345 -->

The frontmatter is HTML-comment so it survives the markdown miner (which
treats lines as text) without polluting drawer content noticeably; the
source-resolver in claude_cli's _build_system_prompt block can grep for
`brain-source:` to map drawer source_files back to original PDFs.

Failures in conversion are isolated per-file: one broken PDF doesn't break
the whole pass.
"""
from __future__ import annotations

import base64
import json
import os
import threading
import time
from dataclasses import dataclass, field
from typing import Iterable

from engine.context import get_request_context

EXTRACT_SUBDIR = ".brain-extracted"

# Hard wall-clock cap for an in-process PDF text extractor (pymupdf4llm / fitz /
# pdfplumber). Some large, table-dense PDFs make pymupdf4llm's layout analysis
# run for MINUTES at 100% CPU (observed: a 37-page struck-off-companies list
# that plain fitz.get_text reads in 0.1s). Those extractors run in worker
# threads, so signal.alarm can't guard them (SIGALRM is main-thread only) — we
# run the call in a daemon thread and abandon it on timeout, falling through to
# the next/faster engine. The abandoned thread keeps burning CPU until it
# finishes, but it no longer blocks the turn.
_PDF_EXTRACT_TIMEOUT_SECS = 60
# PDFs with more than this many pages get a coarse progress ping before the
# (single, whole-doc) pymupdf4llm subprocess call, and per-page progress in the
# fitz fallback. Smaller docs extract sub-second so progress isn't needed.
_PDF_PERPAGE_THRESHOLD = 8


class _ExtractTimeout(Exception):
    pass


def _extract_progress(phase: str, *, pct: float | None = None,
                      current: int | None = None, total: int | None = None,
                      note: str = "") -> None:
    """Thin wrapper over the generic report_tool_progress for the extraction
    chain — phase = the active backend ('pymupdf4llm' / 'fitz' / 'OCR' / …)."""
    from engine.context import report_tool_progress
    report_tool_progress(phase=phase, pct=pct, current=current, total=total, note=note)


def _run_with_timeout(fn, timeout_secs: float):
    """Run `fn()` in a daemon thread; return its result, or raise
    _ExtractTimeout if it doesn't finish within `timeout_secs`. The worker
    thread is NOT killed (Python can't) — it's abandoned (daemon) so the
    process can still exit; the caller falls back to a faster path."""
    box: dict = {}

    def _target():
        try:
            box["result"] = fn()
        except BaseException as e:  # capture so the caller sees the real error
            box["error"] = e

    t = threading.Thread(target=_target, daemon=True)
    t.start()
    t.join(timeout_secs)
    if t.is_alive():
        raise _ExtractTimeout(f"extractor exceeded {timeout_secs:.0f}s")
    if "error" in box:
        raise box["error"]
    return box.get("result")

SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".xls", ".xlsb",
                  ".csv", ".tsv", ".eml", ".msg", ".epub", ".zip"}

# Ad-hoc cache for files outside any project input folder (chat attachments,
# arbitrary paths the agent reads via read_document). Each entry is a
# companion .md keyed by (abs_path, mtime, size) so re-uploads with the
# same temp filename get a fresh extraction automatically. 30-day LRU
# eviction keyed on atime — see evict_adhoc_cache().
ADHOC_CACHE_DIR = os.path.expanduser("~/.brain-agent/extracted-cache")
ADHOC_CACHE_TTL_SECS = 30 * 24 * 3600

# Per-sheet row policy for xlsx extraction. Bank operations rely on
# medium-to-large lookup tables (retention schedules, role matrices,
# tariff lists) where the rows ARE the policy — a small preview throws
# away exactly the information we want indexed. So we extract every
# row, with two safety guards:
#   * XLSX_WARN_ROWS_PER_SHEET — soft warn in daemon log if exceeded
#     (helps spot accidental CSV-export-as-xlsx file drops).
#   * XLSX_MAX_ROWS_PER_SHEET  — hard cap at extraction time. Anything
#     this large is a database export, not a policy artifact, and the
#     drawer storage would be useless anyway (each row mostly numeric).
XLSX_WARN_ROWS_PER_SHEET = 5_000
XLSX_MAX_ROWS_PER_SHEET = 100_000
# Hard cap per cell so a single megablob doesn't blow up a row.
XLSX_CELL_MAX_CHARS = 200

# Files we should never try to convert even if their extension matches —
# guards against name collisions if a user has e.g. a docx named "file.docx"
# but it's actually corrupt.
SKIP_NAMES = {".DS_Store"}


@dataclass
class ConvertResult:
    converted: int = 0
    skipped_unchanged: int = 0
    skipped_unreadable: int = 0
    failed: int = 0
    stale_removed: int = 0
    seen_total: int = 0
    failures: list[str] = field(default_factory=list)
    elapsed_s: float = 0.0


def _stat_key(path: str) -> tuple[int, int]:
    """Cheap change detector: (mtime_ns, size). Re-conversion only happens
    when one of these changes."""
    try:
        st = os.stat(path)
        return int(st.st_mtime), int(st.st_size)
    except OSError:
        return 0, 0


def _read_md_source_stat(md_path: str) -> tuple[int, int]:
    """Pull the brain-source-mtime / brain-source-size from a converted
    file's frontmatter. Returns (0, 0) if absent."""
    try:
        with open(md_path, "r", encoding="utf-8", errors="replace") as f:
            head = f.read(2048)
    except OSError:
        return 0, 0
    mt = sz = 0
    for line in head.splitlines():
        line = line.strip()
        if line.startswith("<!-- brain-source-mtime:"):
            try:
                mt = int(line.split(":", 1)[1].rstrip(" -->").strip())
            except ValueError:
                pass
        elif line.startswith("<!-- brain-source-size:"):
            try:
                sz = int(line.split(":", 1)[1].rstrip(" -->").strip())
            except ValueError:
                pass
        if mt and sz:
            break
    return mt, sz


def _md_path_for(root: str, source_path: str) -> str:
    """Compute the converted-md path under <root>/.brain-extracted/.
    The relative layout under .brain-extracted/ mirrors the source's
    relative layout under root, with .md appended to the original name
    (so collisions across extensions like file.pdf vs file.docx don't
    overwrite each other).
    """
    rel = os.path.relpath(source_path, root)
    base = os.path.join(root, EXTRACT_SUBDIR, rel) + ".md"
    return base


def _source_for_md(root: str, md_path: str) -> str:
    """Inverse of _md_path_for — used only by sweep_stale to detect orphans."""
    rel = os.path.relpath(md_path, os.path.join(root, EXTRACT_SUBDIR))
    if rel.endswith(".md"):
        rel = rel[:-3]
    return os.path.join(root, rel)


def _frontmatter(source_path: str, mtime: int, size: int,
                 backend: str = "") -> str:
    parts = [
        f"<!-- brain-source: {source_path} -->\n",
        f"<!-- brain-source-mtime: {mtime} -->\n",
        f"<!-- brain-source-size: {size} -->\n",
    ]
    if backend:
        parts.append(f"<!-- brain-converter: {backend} -->\n")
    parts.append("\n")
    return "".join(parts)


# ── Markitdown wrapper (preferred backend when available) ────────────────────
#
# Microsoft markitdown produces materially better markdown for LLM consumption
# than fitz/python-docx/python-pptx — preserves table structure, heading
# hierarchy, OCR fallback, and bullet semantics. Wired via subprocess (the
# CLI is at /opt/homebrew/bin/markitdown installed under its own python; we
# avoid importing the package because its python and ours may not match).
#
# Used as the preferred backend for .pdf/.docx/.pptx/.xlsx; falls through
# to per-format extractors on:
#   - subprocess error / non-zero exit
#   - empty stdout (e.g. binary too damaged for markitdown to handle)
#   - markitdown not on PATH (degraded but functional)
#
# Toggleable via config.json -> conversion.use_markitdown (default true when
# detected on PATH).
import shutil
import subprocess

_MARKITDOWN_BIN = shutil.which("markitdown")
_MARKITDOWN_TIMEOUT_SECS = 120
# Per-extension extractor choice: which formats try markitdown FIRST (vs. go
# straight to Brain's own _extract_* fallback). This is the DEFAULT; an admin
# can override it per type in config.json -> conversion.markitdown_exts (a list
# of extensions). Resolved via _markitdown_exts() so a Settings change takes
# effect without a code edit.
#
# .xlsx/.csv/.tsv/.eml are deliberately NOT defaults → Brain's own extractors
# win. Reason: markitdown flattens FOOTER-grouped xlsx reports and loses
# member→group membership (the e487a415 Kostenstellen bug; our _extract_xlsx
# emits an explicit grouping note instead), and leaks MIME headers on .eml. The
# remaining defaults (.pdf/.docx/.pptx/.msg/.epub/.zip) extract better via
# markitdown empirically. (.epub/.zip have no own extractor — _extract_markitdown_only
# just re-invokes markitdown, so removing them from the set would break them;
# they're filtered out of the editable matrix in the UI.)
_DEFAULT_MARKITDOWN_EXTS = {".pdf", ".docx", ".pptx", ".msg", ".epub", ".zip"}
# Formats that have a real own-code extractor, so flipping markitdown off for
# them is meaningful (the UI matrix is built from this set).
_MARKITDOWN_OPTIONAL_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm",
                            ".xls", ".xlsb", ".csv", ".tsv", ".eml", ".msg"}


def _markitdown_exts() -> set:
    """Effective markitdown-first extension set: config override if present,
    else the default. config.json -> conversion.markitdown_exts is a list of
    extensions (with or without leading dot); absent = use defaults."""
    try:
        import brain as _brain
        conv = (_brain._server_config() or {}).get("conversion") or {}
    except Exception:
        conv = {}
    raw = conv.get("markitdown_exts")
    if raw is None:
        return set(_DEFAULT_MARKITDOWN_EXTS)
    out = set()
    for e in raw:
        e = ("." + e.lstrip(".")).lower()
        out.add(e)
    # .epub/.zip have no own extractor — always keep them markitdown-first
    # regardless of config, or they'd produce nothing.
    out.update({".epub", ".zip"})
    return out


# Back-compat shim: some call sites still read the module constant directly.
# Kept as the DEFAULT; the live decision uses _markitdown_exts().
_MARKITDOWN_EXTS = _DEFAULT_MARKITDOWN_EXTS

# ── OCR fallback for scanned PDFs ────────────────────────────────────────────
#
# Markitdown + fitz both produce empty output on image-only PDFs (no text
# layer). Without OCR these become unreadable to the agent — the companion
# .md gets a 1-line "no extractable text" marker, MemPalace mines nothing,
# the agent has no content to cite.
#
# Wired to Mistral OCR (`mistral-ocr-latest`): purpose-built model, ~1 USD
# per 1000 pages, ~2-3s per page, returns structured markdown with tables
# preserved. Routed via the configured Mistral provider — same auth path
# as any other Mistral call, so PII routing applies normally.
#
# Triggers ONLY on PDFs that come back nearly empty from markitdown+fitz
# (per-page chars below OCR_TRIGGER_CHARS_PER_PAGE). Per-cycle page cap
# guards against runaway costs from someone pointing at a 10k-page archive.
_OCR_TRIGGER_CHARS_PER_PAGE = 50
_OCR_TIMEOUT_SECS = 300

# Per-cycle counter (bumped by _extract_with_ocr, reset by callers that
# want to rate-limit). Module-global so daemon cycles can read+reset.
_ocr_pages_this_cycle = 0


def _ocr_config() -> dict:
    """Read the `ocr` block from config.json.

    FAIL-LOUD: provider/model are NEVER fabricated — they come from config or
    are empty (callers error cleanly on empty, and the Doctor's config-model-ref
    check flags an OCR provider/model that doesn't resolve). The only "default"
    here is `engine='none'` when the section is missing entirely, which means
    "OCR off" — a fail-SAFE switch, not a guessed model. Non-model knobs
    (caps/dpi/cost) keep numeric defaults; those aren't model references."""
    try:
        cfg_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "config.json")
        with open(cfg_path, "r", encoding="utf-8") as f:
            cfg = json.load(f)
        ocr = cfg.get("ocr") or {}
    except (OSError, ValueError):
        ocr = {}
    return {
        # engine unset → OCR disabled (no hardcoded provider/model guess).
        # mlx_ocr | mistral_ocr | local_vision | auto | none
        "engine": ocr.get("engine", "none"),
        "provider": ocr.get("provider", ""),
        "model": ocr.get("model", ""),
        # mlx_ocr: a purpose-built OCR model run IN-PROCESS via mlx_vlm (no
        # server hop). Unlike local_vision it does not borrow a chat model —
        # see engine/mlx_ocr.py for why that is both faster and smaller.
        "mlx_ocr_model": ocr.get("mlx_ocr_model", ""),
        "mlx_ocr_max_tokens": int(ocr.get("mlx_ocr_max_tokens", 4096)),
        # Longest edge handed to the OCR model (0 = no cap). See
        # _cap_image_edge: a 200-DPI page render is ~5x slower than a 1600px
        # one for the SAME text. Applies to rendered PDF pages only — a plain
        # image file is used as-is.
        "mlx_ocr_max_edge_px": int(ocr.get("mlx_ocr_max_edge_px", 1600)),
        "max_pages_per_cycle": int(ocr.get("max_pages_per_cycle", 1000)),
        "trigger_chars_per_page": int(ocr.get("trigger_chars_per_page", _OCR_TRIGGER_CHARS_PER_PAGE)),
        # NB: OCR page billing moved to a PER-MODEL rate (models.<id>.cost_per_page_usd,
        # read via quotas._unit_rate) — the old global ocr.cost_per_page_usd is gone.
        # Local-vision fallback: render PDF page → image → vision LLM with
        # OCR prompt. Used when engine='local_vision' or when
        # engine='auto' + GDPR/PII gate forces local.
        "local_vision_model": ocr.get("local_vision_model", ""),
        "local_vision_render_dpi": int(ocr.get("local_vision_render_dpi", 200)),
        "local_vision_max_tokens": int(ocr.get("local_vision_max_tokens", 4096)),
    }


def reset_ocr_cycle_counter() -> int:
    """Return + reset pages-OCR'd-this-cycle. Daemon hooks call at cycle
    start to enforce per-cycle caps."""
    global _ocr_pages_this_cycle
    n = _ocr_pages_this_cycle
    _ocr_pages_this_cycle = 0
    return n


def _extract_with_mistral_ocr(path: str, *, mime: str = "application/pdf"
                              ) -> tuple[str, str | None, int]:
    """OCR a PDF (or an image — pass its `mime`) via Mistral's /v1/ocr
    endpoint. Returns (markdown, error, pages_processed). On any failure
    returns ("", "<reason>", 0) so the caller can fall through to the
    empty-marker path.

    Reads provider config from config.json -> providers[<provider>] (api_key
    + base_url). Per-cycle page cap honored via _ocr_pages_this_cycle.
    """
    global _ocr_pages_this_cycle
    cfg = _ocr_config()
    if cfg["engine"] != "mistral_ocr":
        return "", "ocr disabled", 0
    # Fail loud on missing config — no fabricated provider/model.
    if not cfg["provider"]:
        return "", "ocr.provider not configured", 0
    if not cfg["model"]:
        return "", "ocr.model not configured", 0
    if _ocr_pages_this_cycle >= cfg["max_pages_per_cycle"]:
        return "", f"per-cycle cap {cfg['max_pages_per_cycle']} reached", 0

    # Provider lookup.
    try:
        cfg_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "config.json")
        with open(cfg_path, "r", encoding="utf-8") as f:
            full = json.load(f)
        prov = (full.get("providers") or {}).get(cfg["provider"]) or {}
        api_key = prov.get("api_key", "")
        base_url = prov.get("base_url", "")
    except (OSError, ValueError) as e:
        return "", f"provider config read failed: {e}", 0
    if not base_url:
        return "", f"provider {cfg['provider']} has no base_url", 0
    if not api_key:
        return "", f"provider {cfg['provider']} has no api_key", 0

    import base64
    import urllib.request
    import urllib.error
    try:
        with open(path, "rb") as f:
            data_b64 = base64.b64encode(f.read()).decode("ascii")
    except OSError as e:
        return "", f"read failed: {e}", 0

    # The /ocr endpoint takes a PDF as `document_url` and an image as
    # `image_url` — same call, different document type.
    if mime == "application/pdf":
        document = {"type": "document_url",
                    "document_url": f"data:{mime};base64,{data_b64}"}
    else:
        document = {"type": "image_url",
                    "image_url": f"data:{mime};base64,{data_b64}"}
    payload = {
        "model": cfg["model"],
        "document": document,
        "include_image_base64": False,
    }
    req = urllib.request.Request(
        base_url.rstrip("/") + "/ocr",
        data=json.dumps(payload).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        },
    )
    try:
        with urllib.request.urlopen(req, timeout=_OCR_TIMEOUT_SECS) as resp:
            body = json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as e:
        msg = e.read().decode("utf-8", errors="replace")[:200]
        return "", f"ocr http {e.code}: {msg}", 0
    except (urllib.error.URLError, OSError, ValueError) as e:
        return "", f"ocr request failed: {type(e).__name__}: {e}", 0

    pages = body.get("pages") or []
    parts = []
    for p in pages:
        idx = p.get("index", 0)
        md = (p.get("markdown") or "").strip()
        if md:
            parts.append(f"## Page {idx + 1}\n\n{md}")
    text = "\n\n".join(parts)
    pages_processed = int((body.get("usage_info") or {}).get("pages_processed") or len(pages))
    _ocr_pages_this_cycle += pages_processed

    # Cost row to costs.db. Fail-soft: OCR shouldn't error out because
    # billing logging is unavailable.
    try:
        # Per-MODEL page rate (cost_per_page_usd on the model entry) — the OCR
        # price lives on the model now, not a global ocr.* knob.
        from engine.quotas import _unit_rate
        _log_ocr_cost(
            model=cfg["model"],
            provider=cfg["provider"],
            pages=pages_processed,
            cost_usd=pages_processed * _unit_rate(cfg["model"], "cost_per_page_usd"),
        )
    except Exception as e:
        print(f"[doc-convert] OCR cost-log failed: "
              f"{type(e).__name__}: {e}", flush=True)

    if not text:
        return "", "ocr empty output", pages_processed
    return text + "\n", None, pages_processed


_LOCAL_OCR_PROMPT = (
    "You are an OCR engine. Extract all text from this page image and return "
    "it as clean markdown. Preserve heading levels, list bullets, and table "
    "structure (use markdown tables with | separators). Do not add commentary, "
    "do not describe images. If the page contains a table, output it as a "
    "markdown table. If the page is empty, return an empty response."
)


def _vision_provider_cfg(model: str) -> tuple[str, str, str, str | None]:
    """Resolve (wire_model, base_url, api_key, error) for a vision model.

    Reads config.json directly — same shape (and same reason) as
    _extract_with_mistral_ocr: avoids depending on engine.provider's
    module-globals being initialized outside the daemon.
    """
    try:
        cfg_path = os.path.join(os.path.dirname(os.path.dirname(__file__)),
                                "config.json")
        with open(cfg_path, "r", encoding="utf-8") as f:
            full = json.load(f)
    except (OSError, ValueError) as e:
        return "", "", "", f"config read failed: {e}"
    model_cfg = (full.get("models") or {}).get(model) or {}
    provider_name = model_cfg.get("provider", "")
    if not provider_name:
        return "", "", "", f"model '{model}' has no provider in config"
    prov = (full.get("providers") or {}).get(provider_name) or {}
    base_url = prov.get("base_url", "")
    if not base_url:
        return "", "", "", f"no base_url for provider '{provider_name}'"
    wire_model = model_cfg.get("base_model_id") or model.split("/", 1)[-1]
    return wire_model, base_url, prov.get("api_key", ""), None


def _vision_ocr_round(img_b64: str, mime: str, *, wire_model: str,
                      base_url: str, api_key: str,
                      max_tokens: int) -> tuple[str, str | None]:
    """One image → text round against an OpenAI-compatible vision endpoint.

    The single wire shape used for BOTH a rendered PDF page and a standalone
    image file, so the two paths can't drift apart.
    """
    import urllib.request
    import urllib.error

    payload = {
        "model": wire_model,
        "messages": [{
            "role": "user",
            "content": [
                {"type": "text", "text": _LOCAL_OCR_PROMPT},
                {"type": "image_url",
                 "image_url": {"url": f"data:{mime};base64,{img_b64}"}},
            ],
        }],
        "max_tokens": max_tokens,
        "temperature": 0.0,
        "stream": False,
    }
    req = urllib.request.Request(
        base_url.rstrip("/") + "/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={"Authorization": f"Bearer {api_key}",
                 "Content-Type": "application/json"},
    )
    try:
        with urllib.request.urlopen(req, timeout=_OCR_TIMEOUT_SECS) as resp:
            body = json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as e:
        return "", f"http {e.code}: {e.read().decode('utf-8', errors='replace')[:200]}"
    except (urllib.error.URLError, OSError, ValueError) as e:
        return "", f"request failed: {type(e).__name__}: {e}"
    choices = body.get("choices") or []
    if not choices:
        return "", "empty response"
    return ((choices[0].get("message") or {}).get("content") or "").strip(), None


def _extract_with_local_vision(path: str) -> tuple[str, str | None, int]:
    """OCR a PDF locally by rendering each page to an image and sending it
    to a vision-capable LLM (e.g. Gemma-4 via oMLX). Returns (markdown,
    error, pages_processed). On any failure returns ("", "<reason>", 0).

    Used when ocr.engine='local_vision' or when ocr.engine='auto' and the
    GDPR gate forces a local model. Slower than Mistral OCR (~10-20s/page
    on Gemma-4 26B) and lower quality on tables/multi-column layouts —
    pick this only when data must not leave the host.
    """
    global _ocr_pages_this_cycle
    cfg = _ocr_config()
    model = cfg.get("local_vision_model") or ""
    if not model:
        return "", "local_vision_model not configured", 0
    if _ocr_pages_this_cycle >= cfg["max_pages_per_cycle"]:
        return "", f"per-cycle cap {cfg['max_pages_per_cycle']} reached", 0

    try:
        import fitz  # type: ignore
    except ImportError:
        return "", "pymupdf not installed (pip install pymupdf)", 0

    # Direct provider lookup from config.json — same shape as
    # _extract_with_mistral_ocr. Avoids depending on engine.provider's
    # module-globals being initialized (the daemon thread is fine, but
    # a bare import in a one-shot context isn't).
    try:
        cfg_path = os.path.join(os.path.dirname(os.path.dirname(__file__)), "config.json")
        with open(cfg_path, "r", encoding="utf-8") as f:
            full = json.load(f)
    except (OSError, ValueError) as e:
        return "", f"config read failed: {e}", 0
    model_cfg = (full.get("models") or {}).get(model) or {}
    provider_name = model_cfg.get("provider", "")
    if not provider_name:
        return "", f"model '{model}' has no provider in config", 0
    prov = (full.get("providers") or {}).get(provider_name) or {}
    api_key = prov.get("api_key", "")
    base_url = prov.get("base_url", "")
    if not base_url:
        return "", f"no base_url for provider '{provider_name}'", 0
    # Strip provider prefix if model id is "provider/model_id" — the
    # actual wire-level model name is base_model_id when present.
    wire_model = model_cfg.get("base_model_id") or model.split("/", 1)[-1]

    import base64
    import urllib.request
    import urllib.error

    try:
        doc = fitz.open(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}", 0

    parts: list[str] = []
    pages_processed = 0
    dpi = cfg["local_vision_render_dpi"]
    max_tokens = cfg["local_vision_max_tokens"]
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    try:
        for i, page in enumerate(doc, 1):
            if _ocr_pages_this_cycle >= cfg["max_pages_per_cycle"]:
                # Don't fail the whole call — return what we have.
                parts.append(f"## Page {i}\n\n_(skipped — OCR cycle cap reached)_")
                continue
            try:
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                png_bytes = pix.tobytes("png")
            except Exception as e:
                parts.append(f"## Page {i}\n\n_(render failed: {type(e).__name__}: {e})_")
                continue
            img_b64 = base64.b64encode(png_bytes).decode("ascii")
            data_uri = f"data:image/png;base64,{img_b64}"

            payload = {
                "model": wire_model,
                "messages": [{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": _LOCAL_OCR_PROMPT},
                        {"type": "image_url", "image_url": {"url": data_uri}},
                    ],
                }],
                "max_tokens": max_tokens,
                "temperature": 0.0,
                "stream": False,
            }
            req = urllib.request.Request(
                base_url.rstrip("/") + "/chat/completions",
                data=json.dumps(payload).encode("utf-8"),
                headers={
                    "Authorization": f"Bearer {api_key}",
                    "Content-Type": "application/json",
                },
            )
            try:
                with urllib.request.urlopen(req, timeout=_OCR_TIMEOUT_SECS) as resp:
                    body = json.loads(resp.read().decode("utf-8"))
            except urllib.error.HTTPError as e:
                err_msg = e.read().decode("utf-8", errors="replace")[:200]
                parts.append(f"## Page {i}\n\n_(OCR http {e.code}: {err_msg})_")
                continue
            except (urllib.error.URLError, OSError, ValueError) as e:
                parts.append(f"## Page {i}\n\n_(OCR request failed: {type(e).__name__}: {e})_")
                continue
            choices = body.get("choices") or []
            content = ""
            if choices:
                msg = choices[0].get("message") or {}
                content = (msg.get("content") or "").strip()
            if content:
                parts.append(f"## Page {i}\n\n{content}")
            else:
                parts.append(f"## Page {i}\n\n_(empty OCR output)_")
            pages_processed += 1
            _ocr_pages_this_cycle += 1
    finally:
        doc.close()

    text = "\n\n".join(parts).strip()
    if not text or pages_processed == 0:
        return "", "local-vision: no usable output", pages_processed

    # Local OCR is free at the cost level (own GPU), but we still log it so
    # ops can see throughput / per-cycle volume in the same dashboard.
    try:
        _log_ocr_cost(model=model, provider=provider_name,
                      pages=pages_processed, cost_usd=0.0)
    except Exception as e:
        print(f"[doc-convert] local-vision OCR cost-log failed: "
              f"{type(e).__name__}: {e}", flush=True)

    return text + "\n", None, pages_processed


_IMAGE_MIME = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
}


def _cap_image_edge(path: str, max_edge: int) -> None:
    """Downscale an image in place so its longest edge is <= max_edge.

    A page rendered at 200 DPI is ~2500px wide, and the OCR model's runtime
    scales with the pixel count: MEASURED on the same scanned passport —
        2500px  14.7s        1600px  2.9s        1200px  1.5s
    with byte-identical, fully correct text at every size. So the big render is
    pure waste. We cap rather than lower the DPI because fitz renders sharper
    than a downscale of a low-DPI render, and dense text pages still want the
    detail; 0 disables the cap.
    """
    if not max_edge or max_edge <= 0:
        return
    try:
        from PIL import Image
        with Image.open(path) as im:
            longest = max(im.size)
            if longest <= max_edge:
                return
            ratio = max_edge / longest
            im.resize((max(1, int(im.width * ratio)),
                       max(1, int(im.height * ratio))),
                      Image.LANCZOS).save(path)
    except Exception as e:
        # Cosmetic optimisation — a failure here just means a slower OCR.
        print(f"[doc-convert] image downscale failed for {path}: "
              f"{type(e).__name__}: {e}", flush=True)


def _tesseract_sees_text(path: str) -> bool | None:
    """Does a DETERMINISTIC OCR find any word at all? None = can't tell.

    The sanity gate on the LLM's output. Tesseract never invents text — where
    it reads nothing, there is nothing legible. The vision model does invent:
    on a 420x160 unreadable crop of a passport it produced a holder named
    "Pham Van Pham" and a "Type of Airport: New York City", neither of which is
    in the image. Measured over the 10 real scans, Tesseract returned ZERO words
    on exactly the worst one — so "Tesseract sees nothing" is a clean signal
    that anything the model says about this image is unfounded.
    Only that extreme is used: mid-range confidence does NOT predict a
    hallucination (the other bad scan scored a perfectly normal 53%), so we do
    not gate on a confidence threshold — that would throw away good reads.
    """
    try:
        import pytesseract
        from PIL import Image
        with Image.open(path) as im:
            d = pytesseract.image_to_data(
                im, output_type=pytesseract.Output.DICT)
        return any(t.strip() and int(c) > 0
                   for t, c in zip(d["text"], d["conf"]))
    except Exception:
        return None      # tesseract absent/broken → no opinion, keep the text


def _mlx_ocr_extract(path: str, cfg: dict) -> tuple[str, str | None]:
    """One image → text via the in-process MLX OCR model. (text, error)."""
    from engine import mlx_ocr as _mlx
    text, err = _mlx.extract(
        path, repo=cfg.get("mlx_ocr_model") or _mlx.DEFAULT_MODEL,
        max_tokens=cfg.get("mlx_ocr_max_tokens", 4096))
    if err:
        return "", err
    if text and _tesseract_sees_text(path) is False:
        print(f"[doc-convert] OCR verworfen (kein lesbarer Text im Bild, "
              f"Modellausgabe unbelegt): {os.path.basename(path)}", flush=True)
        return "", None
    if text:
        # Local model on our own GPU → $0, but still logged so OCR throughput
        # shows up in the same dashboard as the cloud engine.
        try:
            _log_ocr_cost(model=cfg.get("mlx_ocr_model") or _mlx.DEFAULT_MODEL,
                          provider="local-mlx-ocr", pages=1, cost_usd=0.0)
        except Exception:
            pass
    return text, None


def _pdf_mlx_ocr(src: str, cfg: dict) -> tuple[str, str | None]:
    """Scanned PDF → text via the in-process MLX OCR model.

    A PDF page is not an image, so it has to be rendered first — same fitz
    render the local_vision path does, reusing its DPI knob.
    """
    global _ocr_pages_this_cycle
    try:
        import fitz  # type: ignore
    except ImportError:
        return "", "pymupdf not installed (pip install pymupdf)"
    try:
        doc = fitz.open(src)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"

    zoom = cfg["local_vision_render_dpi"] / 72.0
    matrix = fitz.Matrix(zoom, zoom)
    parts: list[str] = []
    done = 0
    try:
        for i, page in enumerate(doc, 1):
            if _ocr_pages_this_cycle >= cfg["max_pages_per_cycle"]:
                parts.append(f"## Page {i}\n\n_(skipped — OCR cycle cap reached)_")
                continue
            _extract_progress("OCR", current=i, total=doc.page_count,
                              note=f"Seite {i} (MLX-OCR)")
            tmp = ""
            try:
                pix = page.get_pixmap(matrix=matrix, alpha=False)
                # mlx_vlm reads image FILES, so the rendered page spills to a
                # tempfile; unique name (concurrent mines) + always removed.
                import tempfile
                fd, tmp = tempfile.mkstemp(suffix=".png", prefix="brain-ocr-")
                os.close(fd)
                pix.save(tmp)
                _cap_image_edge(tmp, cfg["mlx_ocr_max_edge_px"])
                text, err = _mlx_ocr_extract(tmp, cfg)
            except Exception as e:
                parts.append(f"## Page {i}\n\n_(render failed: {type(e).__name__}: {e})_")
                continue
            finally:
                if tmp and os.path.isfile(tmp):
                    try:
                        os.unlink(tmp)
                    except OSError:
                        pass
            if err:
                parts.append(f"## Page {i}\n\n_(OCR failed: {err})_")
                continue
            parts.append(f"## Page {i}\n\n{text}" if text
                         else f"## Page {i}\n\n_(empty OCR output)_")
            done += 1
            _ocr_pages_this_cycle += 1
    finally:
        doc.close()

    out = "\n\n".join(parts).strip()
    if not out or done == 0:
        return "", "mlx-ocr: no usable output"
    return out + "\n", None


def _extract_image_ocr(path: str) -> tuple[str, str, str | None]:
    """Read the TEXT out of a standalone image (scan, photo of a document).

    Returns (text, backend, error); empty text = nothing readable, caller
    keeps its metadata-only output. Routed by the SAME config.json -> ocr
    block as the PDF path (engine: mistral_ocr | local_vision | auto | none),
    so an operator who pins OCR to a local model for privacy gets that for
    images too — which is the point, since the images that carry text here
    are passports and ID scans.

    Why this exists: a project image used to yield Pillow metadata ONLY
    (dimensions/format), so a scanned passport entered the corpus as
    "821 x 852 JPEG" with zero content and the KG extracted 0 triples from it.
    """
    ext = os.path.splitext(path)[1].lower()
    mime = _IMAGE_MIME.get(ext)
    if not mime:
        return "", "", None          # .svg et al — not a raster scan
    cfg = _ocr_config()
    engine = cfg["engine"]
    if engine == "none":
        return "", "", None

    global _ocr_pages_this_cycle
    if _ocr_pages_this_cycle >= cfg["max_pages_per_cycle"]:
        return "", "", f"per-cycle cap {cfg['max_pages_per_cycle']} reached"

    err = ""
    # In-process MLX OCR — the fast path: an image needs no rendering, so the
    # file goes straight to the OCR model.
    if engine == "mlx_ocr":
        text, m_err = _mlx_ocr_extract(path, cfg)
        if text:
            _ocr_pages_this_cycle += 1
            return text, "mlx-ocr (1p)", None
        return "", "", m_err

    # Cloud OCR first (when allowed): Mistral's /ocr takes an image data-URI
    # the same way it takes a PDF one — only the document type differs.
    if engine in ("mistral_ocr", "auto"):
        _extract_progress("OCR", note="Bild — Cloud-OCR")
        text, err, pages = _extract_with_mistral_ocr(path, mime=mime)
        if text:
            _ocr_pages_this_cycle += pages
            return text, f"mistral-ocr ({pages}p)", None

    if engine in ("local_vision", "auto"):
        model = cfg.get("local_vision_model") or ""
        if not model:
            return "", "", err or "local_vision_model not configured"
        _extract_progress("OCR", note="Bild — lokales Vision-Modell")
        wire_model, base_url, api_key, cfg_err = _vision_provider_cfg(model)
        if cfg_err:
            return "", "", err or cfg_err
        try:
            with open(path, "rb") as f:
                img_b64 = base64.b64encode(f.read()).decode("ascii")
        except OSError as e:
            return "", "", f"read failed: {e}"
        text, v_err = _vision_ocr_round(
            img_b64, mime, wire_model=wire_model, base_url=base_url,
            api_key=api_key, max_tokens=cfg["local_vision_max_tokens"])
        if text:
            _ocr_pages_this_cycle += 1
            try:
                _log_ocr_cost(model=model, provider="", pages=1, cost_usd=0.0)
            except Exception:
                pass
            return text + "\n", "local-vision (1p)", None
        err = err or v_err

    return "", "", err or None


def _log_ocr_cost(*, model: str, provider: str, pages: int, cost_usd: float) -> None:
    """Forward to brain's live CostTracker. Pulls agent/session/user from
    thread-locals when called from a chat thread; falls back to ('main', '', '')
    for daemon calls."""
    try:
        from brain import _cost_tracker, _current_agent
    except ImportError:
        return
    if _cost_tracker is None:
        return
    agent_id = "main"
    session_id = ""
    user_id = ""
    _ctx = get_request_context()
    if _ctx is not None:
        agent = _ctx.current_agent or _current_agent
        if agent is not None:
            agent_id = getattr(agent, "agent_id", "main")
        session_id = _ctx.current_session_id or ""
        user_id = _ctx.current_user_id or ""
    _cost_tracker.log_ocr(agent=agent_id, session_id=session_id, model=model,
                          provider=provider, pages=pages, cost_usd=cost_usd,
                          user_id=user_id)


def _extract_with_markitdown(path: str) -> tuple[str, str | None]:
    """Returns (markdown_text, error_msg or None). On non-zero exit or
    timeout returns ("", "<reason>") so the caller can fall through to the
    legacy per-format extractor."""
    if not _MARKITDOWN_BIN:
        return "", "markitdown not on PATH"
    try:
        proc = subprocess.run(
            [_MARKITDOWN_BIN, path],
            capture_output=True,
            timeout=_MARKITDOWN_TIMEOUT_SECS,
            check=False,
        )
    except subprocess.TimeoutExpired:
        return "", f"markitdown timeout after {_MARKITDOWN_TIMEOUT_SECS}s"
    except OSError as e:
        return "", f"markitdown spawn failed: {type(e).__name__}: {e}"
    if proc.returncode != 0:
        err = (proc.stderr or b"").decode("utf-8", errors="replace").strip()
        return "", f"markitdown exit {proc.returncode}: {err[:200]}"
    text = (proc.stdout or b"").decode("utf-8", errors="replace").strip()
    if not text:
        return "", "markitdown empty output"
    return text + "\n", None


# ── Per-format extractors ────────────────────────────────────────────────────

def _parse_index_selection(spec: str | None) -> set[int] | None:
    """Parse a 1-based selection string like "1,3,5-7" into a set of ints.
    None/empty → None (meaning "all", no filtering). Malformed parts are
    skipped silently rather than raising — a bad selection should degrade to
    fewer pages, never crash the read."""
    if not spec or not str(spec).strip():
        return None
    out: set[int] = set()
    for part in str(spec).split(","):
        part = part.strip()
        if not part:
            continue
        try:
            if "-" in part:
                a, b = part.split("-", 1)
                for i in range(int(a), int(b) + 1):
                    out.add(i)
            else:
                out.add(int(part))
        except ValueError:
            continue
    return out or None


def _render_pdf_tables(path: str, page_sel: set[int] | None) -> dict[int, list[str]]:
    """Per-page pdfplumber table reconstruction → {page_num: [markdown, ...]}.
    Empty dict when pdfplumber is unavailable or no tables found. Ported from
    read_document's inline include_tables path so the unified pipeline keeps
    table fidelity for table-heavy PDFs."""
    tables_by_page: dict[int, list[str]] = {}
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return tables_by_page
    try:
        with pdfplumber.open(path) as plumb:
            for i, page in enumerate(plumb.pages, 1):
                if page_sel is not None and i not in page_sel:
                    continue
                found = page.extract_tables()
                if not found:
                    continue
                rendered = []
                for t in found:
                    rows = [[(c or "").replace("|", "\\|").replace("\n", " ").strip()
                             for c in row] for row in t]
                    if not rows:
                        continue
                    width = max(len(r) for r in rows)
                    for r in rows:
                        r.extend([""] * (width - len(r)))
                    header = "| " + " | ".join(rows[0]) + " |"
                    sep = "| " + " | ".join("---" for _ in range(width)) + " |"
                    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
                    rendered.append(header + "\n" + sep + ("\n" + body if body else ""))
                if rendered:
                    tables_by_page[i] = rendered
    except Exception:
        # pdfplumber is best-effort enrichment — never let it fail the read.
        return {}
    return tables_by_page


def _extract_pdf(path: str, *, pages: str | None = None,
                 include_tables: bool = False, emit_meta: bool = False,
                 page_marker: str = "## Page") -> tuple[str, str | None]:
    """Returns (markdown_text, error_msg or None). Empty text + None means
    the PDF had no extractable text (likely scanned image). Empty text +
    error means a real failure.

    Mining calls with defaults → byte-stable legacy output (`# title` +
    `## Page N` headings, no metadata, no tables). read_document/
    read_attachment pass:
      * pages          — 1-based selection "1,3,5-7" (None = all)
      * include_tables — opt-in pdfplumber per-page table reconstruction
      * emit_meta      — prepend a **title/author/page_count** metadata block
      * page_marker    — heading prefix per page ("--- Page" for the
                         read_document `--- Page N ---` format the citation/
                         page logic parses)
    """
    try:
        import fitz  # type: ignore
    except ImportError:
        return "", "pymupdf not installed (pip install pymupdf)"
    try:
        doc = fitz.open(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    page_sel = _parse_index_selection(pages)
    parts = []
    meta = doc.metadata or {}
    title = meta.get("title") or os.path.splitext(os.path.basename(path))[0]
    page_count = doc.page_count
    if emit_meta:
        meta_lines = [f"**title:** {title}"]
        if meta.get("author"):
            meta_lines.append(f"**author:** {meta.get('author')}")
        meta_lines.append(f"**page_count:** {page_count}")
        parts.append("\n".join(meta_lines) + "\n")
    else:
        parts.append(f"# {title}\n")
    # `--- Page N ---` is a flat marker on its own line; `## Page N` is a
    # markdown heading. Match the two historical formats exactly.
    flat_marker = page_marker.strip().startswith("---")
    try:
        for i, page in enumerate(doc, 1):
            if page_sel is not None and i not in page_sel:
                continue
            t = page.get_text("text") or ""
            t = t.strip()
            if t:
                if flat_marker:
                    parts.append(f"\n\n{page_marker} {i} ---\n\n{t}\n")
                else:
                    parts.append(f"\n\n{page_marker} {i}\n\n{t}\n")
    finally:
        doc.close()

    if include_tables:
        tables_by_page = _render_pdf_tables(path, page_sel)
        if tables_by_page:
            for pnum in sorted(tables_by_page):
                for j, md in enumerate(tables_by_page[pnum], 1):
                    parts.append(f"### Table (page {pnum}, #{j})\n{md}")

    text = "\n".join(parts).strip()
    if len(text.splitlines()) <= 1:
        # Only the title/meta line — likely scanned PDF or empty.
        return "", None
    return text + "\n", None


def _pymupdf4llm_is_blank(md: str) -> bool:
    """True if pymupdf4llm output carries no real text — only its image
    placeholder lines (`**==> picture [W x H] intentionally omitted <==**`)
    and blank lines. An image-only/scanned PDF emits ONE such placeholder per
    embedded image fragment (a scanned page is often 100+ images), so the raw
    line count looks substantial even though there is zero extractable text.
    Without this, the `splitlines() <= 1` emptiness check below sees hundreds
    of placeholder lines, declares the doc non-empty, and the caller never
    reaches the OCR fallback. Stripping placeholders first lets a genuinely
    image-only PDF be detected as scanned → OCR runs."""
    for line in md.splitlines():
        s = line.strip()
        if not s:
            continue
        if s.startswith("**==>") and "intentionally omitted" in s and s.endswith("<==**"):
            continue
        # Empty markdown headers: pymupdf4llm sometimes emits a run of bare
        # "##" / "###" lines (heading markers with NO text) for a PDF whose
        # real text it failed to lift — strip them too, else the first such
        # line reads as "real content" and the fitz/OCR fallback never fires
        # (observed on a Wiener-Privatbank termination letter: pymupdf4llm gave
        # only empty "##" headers + one picture placeholder = 77 chars of junk,
        # while fitz lifted the full 1578-char text). A header WITH text
        # ("## Vertrag") is still real content → returns False as before.
        if s.lstrip("#").strip() == "":
            continue
        return False
    return True


# Worker script run in a child process to extract one PDF (or a page subset)
# via pymupdf4llm. Kept as a -c string (no separate file to ship/maintain).
# Reads a JSON job from stdin {path, idxs|null, out}, writes the markdown to
# the `out` temp file, prints a one-line JSON status to stdout. Runs in its OWN
# process so the parent can SIGKILL it on timeout and RECLAIM the CPU — a plain
# daemon-thread timeout only abandons the thread, which keeps grinding at 100%
# on pathological PDFs (the v9.156.x incident: one web-fetched PDF pegged a core
# for minutes after the 60s timeout, starving the chat turn).
_PYMUPDF4LLM_WORKER = r"""
import sys, json
job = json.load(sys.stdin)
try:
    import pymupdf4llm
    kwargs = {}
    if job.get("idxs") is not None:
        kwargs["pages"] = job["idxs"]
    md = pymupdf4llm.to_markdown(job["path"], **kwargs)
    with open(job["out"], "w", encoding="utf-8") as f:
        f.write(md or "")
    print(json.dumps({"ok": True}))
except Exception as e:
    print(json.dumps({"ok": False, "err": "%s: %s" % (type(e).__name__, e)}))
"""


def _pymupdf4llm_subprocess(path: str, idxs: list[int] | None,
                            timeout_secs: float) -> tuple[str, str | None]:
    """Run pymupdf4llm.to_markdown in a child process, hard-killed on timeout.

    `idxs` = 0-based page indices to extract, or None for the whole document.
    Returns (markdown, error). On timeout returns ("", "<timeout reason>") — the
    child is SIGKILLed by subprocess.run so its CPU is reclaimed immediately
    (unlike the thread-based _run_with_timeout, which can only abandon the
    still-grinding worker). Mirrors _extract_with_markitdown's contract so the
    caller's fallback chain (→ fitz → OCR) is unchanged.
    """
    import sys
    import tempfile
    out_fd, out_path = tempfile.mkstemp(suffix=".md", prefix="pmupdf4llm-")
    os.close(out_fd)
    try:
        job = json.dumps({"path": path, "idxs": idxs, "out": out_path})
        try:
            proc = subprocess.run(
                [sys.executable, "-c", _PYMUPDF4LLM_WORKER],
                input=job.encode("utf-8"),
                capture_output=True,
                timeout=timeout_secs,
                check=False,
            )
        except subprocess.TimeoutExpired:
            # subprocess.run has already killed the child here — CPU reclaimed.
            return "", f"pymupdf4llm exceeded {timeout_secs:.0f}s (killed)"
        except OSError as e:
            return "", f"pymupdf4llm spawn failed: {type(e).__name__}: {e}"
        if proc.returncode != 0:
            err = (proc.stderr or b"").decode("utf-8", errors="replace").strip()
            return "", f"pymupdf4llm worker exit {proc.returncode}: {err[:200]}"
        # Parse the worker's status line (last line of stdout).
        status = {}
        try:
            line = (proc.stdout or b"").decode("utf-8", errors="replace").strip().splitlines()
            if line:
                status = json.loads(line[-1])
        except (ValueError, IndexError):
            pass
        if not status.get("ok"):
            return "", f"pymupdf4llm: {status.get('err', 'worker produced no status')}"
        try:
            with open(out_path, "r", encoding="utf-8") as f:
                return f.read(), None
        except OSError as e:
            return "", f"pymupdf4llm read output failed: {type(e).__name__}: {e}"
    finally:
        try:
            os.remove(out_path)
        except OSError:
            pass


def _extract_pdf_pymupdf4llm(path: str, *, pages: str | None = None) -> tuple[str, str | None]:
    """PDF → markdown via pymupdf4llm (a fitz wrapper). Renders tables + layout
    to GitHub-flavoured markdown far better than markitdown on structured docs
    (financial reports etc. — verified on the WPB Konzernbilanz). Returns
    (markdown, error). Empty text + None = no extractable text (scanned →
    caller's OCR path). Missing dep / failure → error so the caller falls back.
    NOTE: pymupdf4llm/PyMuPDF is AGPL-3.0 (Artifex)."""
    try:
        import fitz  # type: ignore
    except ImportError:
        return "", "pymupdf4llm not installed (pip install pymupdf4llm)"
    page_sel = _parse_index_selection(pages)
    # Page list (0-based) to process; None → all pages. fitz.open is a fast,
    # GIL-releasing read just to count pages and pick the branch — the heavy
    # pymupdf4llm layout analysis runs in the subprocess below.
    try:
        doc = fitz.open(path)
        n_total = doc.page_count
        doc.close()
    except Exception as e:
        return "", f"pymupdf4llm open: {type(e).__name__}: {e}"
    if page_sel is not None:
        idxs = sorted(p - 1 for p in page_sel if 1 <= p <= n_total)
    else:
        idxs = list(range(n_total))

    # pymupdf4llm.to_markdown runs in ONE child process (hard-killable on
    # timeout) instead of a daemon thread — see _pymupdf4llm_subprocess. The
    # layout analysis is CPU-bound and uninterruptible from Python, so the old
    # thread-based _run_with_timeout could only ABANDON a runaway, leaving it to
    # peg a core indefinitely (the v9.156.x freeze). A subprocess timeout SIGKILLs
    # the child and reclaims the CPU.
    #
    # A single whole-doc call handles all sizes: measured ~1.8s for an 18-page
    # doc, vs ~8.6s if each page were its own subprocess (spawn overhead ×N).
    # The old per-page loop existed to bound a single pathological page and to
    # emit page-i/N progress; subprocess isolation already bounds the WHOLE call
    # (one timeout, one kill), so the per-page spawning was pure overhead. We
    # keep a coarse progress ping for big docs so the UI isn't silent.
    if len(idxs) > _PDF_PERPAGE_THRESHOLD:
        _extract_progress("pymupdf4llm", note=f"{len(idxs)} Seiten — Layout-Analyse")
    sub_idxs = idxs if page_sel is not None else None
    md, err = _pymupdf4llm_subprocess(path, sub_idxs, _PDF_EXTRACT_TIMEOUT_SECS)
    if err:
        # Timeout → propagate as _ExtractTimeout so _do_extract switches to fitz
        # (not markitdown); other (hard) errors return as a plain error for the
        # same fallback.
        if "exceeded" in err:
            raise _ExtractTimeout(err)
        return "", err
    md = (md or "").strip()
    if len(md.splitlines()) <= 1 or _pymupdf4llm_is_blank(md):
        return "", None   # scanned/empty → caller's OCR path
    if _pdf_has_no_text_layer(path, page_sel):
        # pymupdf4llm silently falls back to TESSERACT on a page with no text
        # layer ("Using Tesseract for OCR processing" on stderr) and returns its
        # output as if it were extracted text. That output is markedly worse
        # than our configured OCR engine — on a scanned passport Tesseract read
        # "05.02.1847" (year off by a century) and "S6068370F" for the passport
        # number, where GLM-OCR reads both correctly. Because a non-empty
        # result counts as success, the OCR chain below was NEVER reached and
        # every scan in the corpus silently got the worse text.
        # No text layer ⇒ whatever came back IS Tesseract ⇒ drop it and let the
        # caller run the real OCR.
        return "", None
    return md + "\n", None


def _pdf_has_no_text_layer(path: str, page_sel: set | None = None) -> bool:
    """True when NO selected page carries an actual text layer (i.e. the doc is
    a scan). Cheap: `page.get_text()` reads the existing text objects, it does
    not render or OCR."""
    try:
        import fitz  # type: ignore
        doc = fitz.open(path)
    except Exception:
        return False      # can't tell → don't discard pymupdf4llm's output
    try:
        for i, page in enumerate(doc, 1):
            if page_sel is not None and i not in page_sel:
                continue
            if (page.get_text() or "").strip():
                return False
        return True
    except Exception:
        return False
    finally:
        doc.close()


def _extract_pdf_fitz_fast(path: str, *, pages: str | None = None) -> tuple[str, str | None]:
    """Fastest possible PDF text read: bare `page.get_text()` per page, no table
    reconstruction. Used as the timeout fallback when pymupdf4llm/pdfplumber
    hang on a big PDF that nonetheless has a clean text layer (the 4aad5750
    case). Returns (text, error); empty text + None = no text layer (→ OCR)."""
    try:
        import fitz  # type: ignore
    except ImportError:
        return "", "PyMuPDF (fitz) not installed"
    page_sel = _parse_index_selection(pages)
    try:
        doc = fitz.open(path)
    except Exception as e:
        return "", f"fitz open: {type(e).__name__}: {e}"
    parts = []
    sel = [i for i in range(1, doc.page_count + 1)
           if page_sel is None or i in page_sel]
    try:
        for n, i in enumerate(sel, 1):
            if len(sel) > _PDF_PERPAGE_THRESHOLD:
                _extract_progress("fitz", current=n, total=len(sel),
                                  note=f"Seite {i}")
            parts.append(doc[i - 1].get_text() or "")
    finally:
        doc.close()
    text = "\n".join(p for p in parts if p).strip()
    return (text + "\n", None) if text else ("", None)


def _pdf_engine() -> str:
    """Which PDF→text backend to use: 'pymupdf4llm' (default — best tables),
    'markitdown', or 'fitz' (the plain page.get_text legacy path). Overridable
    via config.json -> conversion.pdf_engine."""
    try:
        import brain as _brain
        conv = (_brain._server_config() or {}).get("conversion") or {}
        eng = str(conv.get("pdf_engine") or "").strip().lower()
        if eng in ("pymupdf4llm", "markitdown", "fitz"):
            return eng
    except Exception:
        pass
    return "pymupdf4llm"


def _extract_docx(path: str) -> tuple[str, str | None]:
    try:
        import docx  # type: ignore
    except ImportError:
        return "", "python-docx not installed (pip install python-docx)"
    try:
        d = docx.Document(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    title = os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]
    for p in d.paragraphs:
        t = (p.text or "").rstrip()
        if not t:
            parts.append("")
            continue
        # Heading levels: bold + first run + style starting with "Heading"
        style_name = (getattr(p.style, "name", "") or "").lower()
        if style_name.startswith("heading"):
            try:
                level = int(style_name.replace("heading", "").strip() or "2")
            except ValueError:
                level = 2
            level = max(2, min(6, level))
            parts.append(f"\n{'#' * level} {t}\n")
        else:
            parts.append(t)
    text = "\n".join(parts).strip()
    if len(text.splitlines()) <= 1:
        return "", None
    return text + "\n", None


def _extract_pptx(path: str, *, slides: str | None = None
                  ) -> tuple[str, str | None]:
    """`slides` (default None = all slides, mining behavior) is a 1-based
    selection string like "1,3,5-7" — read_document/read_attachment `slides=`
    selection. Slides outside the set are skipped."""
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return "", "python-pptx not installed (pip install python-pptx)"
    try:
        prs = Presentation(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    slide_sel = _parse_index_selection(slides)
    title = os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]
    slide_count = 0
    for i, slide in enumerate(prs.slides, 1):
        if slide_sel is not None and i not in slide_sel:
            continue
        slide_count += 1
        # Try to surface the slide title from its title placeholder when
        # present; fall back to "Slide N".
        slide_title = ""
        try:
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                slide_title = (slide.shapes.title.text_frame.text or "").strip()
        except Exception:
            slide_title = ""
        heading = f"## Slide {i}" + (f" — {slide_title}" if slide_title else "")
        parts.append(f"\n{heading}\n")
        for shape in slide.shapes:
            # Skip the title shape we already used.
            if shape == slide.shapes.title:
                continue
            try:
                if not shape.has_text_frame:
                    continue
            except Exception:
                continue
            for para in shape.text_frame.paragraphs:
                t = (para.text or "").rstrip()
                if not t:
                    continue
                # Indent bullets a bit so list structure carries through.
                level = getattr(para, "level", 0) or 0
                bullet = "  " * level + "- " if level >= 0 else ""
                parts.append(f"{bullet}{t}")
        # Speaker notes — often hold the actual policy/explanation text.
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    parts.append(f"\n_Speaker notes:_ {notes}")
        except Exception:
            pass
    text = "\n".join(parts).strip()
    # Title + 0 slides → empty.
    if slide_count == 0 or len(text.splitlines()) <= 1:
        return "", None
    return text + "\n", None


def _detect_footer_groups(header, data_rows):
    """Detect a FOOTER-grouped report: a 'key' column filled ONLY on TOTAL rows
    (where a 'label' column is empty), each preceded by a block of MEMBER rows
    (key empty, label filled). Returns (key_name, label_name, [(key, [labels])])
    or None when the pattern isn't clearly present (→ normal table, no note).

    Conservative on purpose: requires ≥3 total rows, members ≥ totals, at most
    one structural violation, and every group non-empty — so it fires on real
    grouped reports (e.g. Kostenstellen-Teilnehmer) but never reshapes an
    ordinary flat table. Verified on the live e487a415 files (22 groups, no
    false positive on flat/key-value tables)."""
    if not header or len(header) < 2:
        return None
    n = len(header)

    def s(v):
        return str(v).strip() if v is not None else ""

    norm = []
    for r in data_rows:
        rr = [s(c) for c in (r or [])]
        rr += [""] * (n - len(rr))
        norm.append(rr[:n])

    best = None
    for kc in range(n):
        for lc in range(n):
            if kc == lc:
                continue
            members = []
            groups = []
            totals = 0
            member_rows = 0
            bad = 0
            for r in norm:
                if not any(r):              # blank separator row
                    if members:
                        bad += 1            # members with no total before a blank = orphan
                        members = []
                    continue
                kv, lv = r[kc], r[lc]
                if kv and not lv:           # TOTAL/footer row
                    totals += 1
                    groups.append((kv, list(members)))
                    members = []
                elif lv and not kv:         # MEMBER row
                    member_rows += 1
                    members.append(lv)
                else:                        # both filled / partial → not the pattern
                    bad += 1
            if members:
                bad += 1
            if (totals >= 3 and member_rows >= totals and bad <= 1
                    and groups and all(len(m) >= 1 for _, m in groups)):
                score = totals + member_rows
                if not best or score > best[0]:
                    best = (score, kc, lc, groups)
    if not best:
        return None
    _, kc, lc, groups = best
    kname = s(header[kc]) or f"col{kc + 1}"
    lname = s(header[lc]) or f"col{lc + 1}"
    return (kname, lname, groups)


def _trim_placeholder_columns(header, data_rows, cell_fn):
    """Trim dead trailing columns off (header, data_rows).

    Excel reports max_column up to the sheet limit (16384) when stray
    formatting or auto-named placeholder headers (e.g. "Spalte41" …
    "Spalte16347") extend past the real data — openpyxl then yields
    thousands of columns whose header is non-empty but whose every DATA
    cell is blank. A trailing-EMPTY-header trim can't catch these (the
    headers are named), so the flattened table exploded to ~1.5 MB and
    every downstream pass crawled over 16k columns (~60 s). Compute the
    last column any header OR data cell actually uses, then slice header +
    rows to it up front — so this is the ONLY full-width pass. Real columns
    kept, dead placeholders gone. We keep header columns only up to one past
    the last data column OR the last CONTIGUOUS real header — whichever is
    larger — so a legit trailing column that happens to be data-empty in the
    preview isn't lost, while the 16k-placeholder tail is.

    `cell_fn` is the caller's cell-normalizer (str + strip; caps don't affect
    emptiness). Shared by _extract_xlsx (mining/read — output must stay
    byte-stable, pinned in tests/test_xlsx_tools.py) and engine/tools/
    xlsx_tools.py.
    """
    _data_last = -1
    for _r in data_rows:
        for _ci in range(len(_r) - 1, _data_last, -1):
            if cell_fn(_r[_ci]) != "":
                _data_last = _ci
                break
    # Real headers are contiguous from col 0; the placeholder tail begins
    # at the first header whose name matches Excel's auto pattern
    # "Spalte<N>" AND which has no data. Find where the contiguous
    # non-auto header block ends.
    import re as _re
    _hdr_real_last = -1
    for _ci, _hv in enumerate(header):
        _hs = cell_fn(_hv)
        if _hs == "" or _re.fullmatch(r"Spalte\d+", _hs):
            break
        _hdr_real_last = _ci
    _last_used = max(_data_last, _hdr_real_last)
    if _last_used < 0:
        _last_used = 0
    if _last_used < len(header) - 1:
        header = header[:_last_used + 1]
        data_rows = [_r[:_last_used + 1] for _r in data_rows]
    return header, data_rows


def _extract_xlsx(path: str, *, caps: bool = True,
                  sheet: str | None = None) -> tuple[str, str | None]:
    """Render every sheet's header + first N rows as a markdown table.

    Spreadsheets are mostly numeric/tabular and don't fit the drawer model
    cleanly. We extract a structured preview per sheet — enough for the
    embedding to learn the topic and the KG extractor to see e.g. retention
    schedules, role matrices, threshold tables — and stop. Users who need
    full-fidelity row-by-row analysis should call read_document on the
    original .xlsx rather than relying on the converted markdown.

    `caps` (default True, mining behavior) applies the per-sheet row cap and
    per-cell char cap. read_document/read_attachment pass caps=False for
    full-fidelity row-by-row output. `sheet` restricts output to a single
    named sheet (read_document/read_attachment `sheet=` selection); unknown
    or None name = all sheets. Pipe-escaping, blank-row skipping, and
    per-sheet error isolation apply regardless of caps.
    """
    try:
        import openpyxl  # type: ignore
    except ImportError:
        return "", "openpyxl not installed (pip install openpyxl)"
    try:
        # data_only=True so cached formula values come through (the user
        # sees the displayed value, not the formula expression).
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    title = os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]
    cell_cap = XLSX_CELL_MAX_CHARS if caps else None
    row_cap = XLSX_MAX_ROWS_PER_SHEET if caps else None

    def _cell(v):
        if v is None:
            return ""
        s = str(v)
        if cell_cap is not None and len(s) > cell_cap:
            s = s[:cell_cap] + "…"
        # Markdown table separator + escape pipes inside cells.
        return s.replace("|", "\\|").replace("\n", " ").strip()

    if sheet and sheet in wb.sheetnames:
        target_sheets = [wb[sheet]]
    else:
        target_sheets = wb.worksheets

    sheet_count = 0
    for sheet in target_sheets:
        try:
            rows = sheet.iter_rows(values_only=True)
            header = next(rows, None)
            if header is None:
                continue
            sheet_count += 1
            parts.append(f"\n## Sheet: {sheet.title}\n")
            # Buffer the data rows so we can both (a) detect a footer-grouping
            # structure and (b) render the table. Memory note: read_only +
            # generator was for streaming, but sheets we extract are already
            # capped/previewed; buffering one sheet's rows is fine.
            data_rows = list(rows)
            # Trim dead trailing columns BEFORE any per-column work below —
            # the v9.261.0 wide-sheet fix, shared with xlsx_tools (see
            # _trim_placeholder_columns for the full story).
            header, data_rows = _trim_placeholder_columns(
                header, data_rows, _cell)
            # Footer-grouping note: some reports encode group membership by
            # POSITION — a "group key" column (e.g. Kostenstelle) is filled only
            # on a TOTAL/footer row BELOW its member rows (whose key cell is
            # blank). Flattening to a table drops which member belongs to which
            # group (the e487a415 bug: model couldn't map person→Kostenstelle).
            # When the pattern is CLEARLY present, prepend an explicit
            # "group: members" list so the model (and KG) get the membership
            # the flat table loses. No-op on normal tables (returns None).
            try:
                _fg = _detect_footer_groups(header, data_rows)
            except Exception:
                _fg = None
            if _fg:
                _kname, _lname, _groups = _fg
                parts.append(
                    f"_Hinweis zur Struktur: Diese Tabelle ist nach **{_kname}** "
                    f"gruppiert — der {_kname}-Wert steht jeweils in der "
                    f"Summenzeile UNTER seinen Mitgliedern (deren {_kname}-Zelle "
                    f"leer ist). Zuordnung der Mitglieder zur Gruppe:_\n")
                for _k, _mem in _groups:
                    _names = ", ".join(m for m in _mem)
                    parts.append(f"- **{_kname} {_cell(_k)}**: {_cell(_names)}")
                parts.append("")  # blank line before the raw table
            cells = [_cell(c) for c in header]
            # Drop trailing empty columns so the table doesn't have ragged
            # padding from sparsely-used right side of the sheet. (The wide-sheet
            # placeholder case is already handled by the up-front trim above.)
            while cells and cells[-1] == "":
                cells.pop()
            if not cells:
                # Empty header row, skip; sheet may be data-without-header.
                cells = ["col_1"]
            n_cols = len(cells)
            parts.append("| " + " | ".join(cells) + " |")
            parts.append("|" + "|".join(["---"] * n_cols) + "|")
            shown = 0
            warned = False
            for r in data_rows:
                if row_cap is not None and shown >= row_cap:
                    parts.append(
                        f"\n_(truncated at {row_cap:,} "
                        f"rows — use read_document on the original .xlsx "
                        f"for the rest)_\n")
                    break
                rc = [_cell(c) for c in r[:n_cols]]
                # Pad short rows to keep table valid.
                while len(rc) < n_cols:
                    rc.append("")
                if all(c == "" for c in rc):
                    continue  # skip blank rows
                parts.append("| " + " | ".join(rc) + " |")
                shown += 1
                if not warned and shown == XLSX_WARN_ROWS_PER_SHEET:
                    print(
                        f"[doc-convert] {path}: sheet `{sheet.title}` "
                        f"has >{XLSX_WARN_ROWS_PER_SHEET:,} rows — "
                        f"continuing extraction up to "
                        f"{XLSX_MAX_ROWS_PER_SHEET:,}",
                        flush=True)
                    warned = True
        except Exception as e:
            parts.append(f"\n_(failed to read sheet `{sheet.title}`: "
                         f"{type(e).__name__})_\n")
    try:
        wb.close()
    except Exception:
        pass
    # Append any VBA macro source (xlsm/xls with macros). No-op for plain xlsx.
    vba = _extract_vba(path)
    if vba:
        parts.append(vba)
    if sheet_count == 0 and not vba:
        return "", None
    return "\n".join(parts) + "\n", None


def _extract_xlsb(path: str, *, caps: bool = True,
                  sheet: str | None = None) -> tuple[str, str | None]:
    """Binary workbook (.xlsb) via pyxlsb (openpyxl can't read the binary
    format). Mirrors _extract_xlsx's output shape: a markdown table per sheet
    (header + capped rows), VBA source appended if present. pyxlsb is read-only
    + values-only, so there's no formula text and no footer-group detection
    pass (the value buffering would defeat its streaming design); the flat
    table + VBA is the contract for binary workbooks."""
    try:
        from pyxlsb import open_workbook  # type: ignore
    except ImportError:
        return "", "pyxlsb not installed (pip install pyxlsb)"
    cell_cap = XLSX_CELL_MAX_CHARS if caps else None
    row_cap = XLSX_MAX_ROWS_PER_SHEET if caps else None
    title = os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]

    def _cell(v):
        if v is None:
            return ""
        s = str(v)
        if cell_cap is not None and len(s) > cell_cap:
            s = s[:cell_cap] + "…"
        return s.replace("|", "\\|").replace("\n", " ").strip()

    sheet_count = 0
    try:
        with open_workbook(path) as wb:
            names = wb.sheets
            targets = [sheet] if (sheet and sheet in names) else names
            for sname in targets:
                try:
                    with wb.get_sheet(sname) as ws:
                        rows_iter = ws.rows()
                        header = next(rows_iter, None)
                        if header is None:
                            continue
                        sheet_count += 1
                        parts.append(f"\n## Sheet: {sname}\n")
                        hdr = [_cell(c.v) for c in header]
                        while hdr and hdr[-1] == "":
                            hdr.pop()
                        if not hdr:
                            continue
                        parts.append("| " + " | ".join(hdr) + " |")
                        parts.append("| " + " | ".join("---" for _ in hdr) + " |")
                        n = 0
                        for row in rows_iter:
                            if row_cap is not None and n >= row_cap:
                                parts.append(f"\n_(truncated at {row_cap:,} rows)_\n")
                                break
                            vals = [_cell(c.v) for c in row][:len(hdr)]
                            vals += [""] * (len(hdr) - len(vals))
                            if any(v for v in vals):
                                parts.append("| " + " | ".join(vals) + " |")
                                n += 1
                except Exception as e:
                    parts.append(f"\n_(failed to read sheet `{sname}`: "
                                 f"{type(e).__name__})_\n")
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    vba = _extract_vba(path)
    if vba:
        parts.append(vba)
    if sheet_count == 0 and not vba:
        return "", None
    return "\n".join(parts) + "\n", None


def list_vba_modules(path: str) -> list[dict]:
    """VBA module sources of a macro-enabled Office file as structured data
    [{name, code}] — the JSON twin of _extract_vba's markdown (which stays
    byte-stable for mining). Macros are NEVER executed; read of the stored
    source only. Empty list when there are no macros / oletools is absent.
    Feeds the bottom-panel VBA viewer (GET /v1/files/xlsm-vba, v9.265.0)."""
    try:
        from oletools.olevba import VBA_Parser  # type: ignore
    except ImportError:
        return []
    vp = None
    try:
        vp = VBA_Parser(path)
        if not vp.detect_vba_macros():
            return []
        modules = []
        for (_fname, _stream, vba_name, vba_code) in vp.extract_macros():
            code = (vba_code or "").strip()
            if code:
                modules.append({"name": vba_name, "code": code})
        return modules
    except Exception:
        return []
    finally:
        if vp is not None:
            try:
                vp.close()
            except Exception:
                pass


def _extract_vba(path: str) -> str:
    """Extract VBA macro source from a macro-enabled Office file as a markdown
    section (one fenced code block per module), or '' when there are no macros /
    oletools is absent. Macros are NEVER executed — this reads the stored source
    only. Used by the xlsm/xls/xlsb extractors so the agent can read + reason
    about automation logic in a workbook."""
    try:
        from oletools.olevba import VBA_Parser  # type: ignore
    except ImportError:
        return ""
    vp = None
    try:
        vp = VBA_Parser(path)
        if not vp.detect_vba_macros():
            return ""
        blocks = []
        for (_fname, _stream, vba_name, vba_code) in vp.extract_macros():
            code = (vba_code or "").strip()
            if not code:
                continue
            blocks.append(f"### Makro: {vba_name}\n```vba\n{code}\n```")
        if not blocks:
            return ""
        return "\n## VBA-Makros (Quellcode — wird NICHT ausgeführt)\n\n" + "\n\n".join(blocks)
    except Exception:
        return ""
    finally:
        if vp is not None:
            try:
                vp.close()
            except Exception:
                pass


def _extract_eml(path: str) -> tuple[str, str | None]:
    """Parse RFC 822 / MIME messages. Stdlib `email` only — no extras."""
    import email
    from email import policy as _policy
    try:
        with open(path, "rb") as f:
            msg = email.message_from_binary_file(f, policy=_policy.default)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    return _email_to_markdown(msg), None


def _extract_msg(path: str) -> tuple[str, str | None]:
    """Outlook .msg files. Requires `extract-msg`; if absent, returns a
    helpful error rather than silently skipping."""
    try:
        import extract_msg  # type: ignore
    except ImportError:
        return "", "extract-msg not installed (pip install extract-msg)"
    try:
        m = extract_msg.Message(path)
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    title = m.subject or os.path.splitext(os.path.basename(path))[0]
    parts = [f"# {title}\n"]
    if m.sender: parts.append(f"**From:** {m.sender}")
    if m.to: parts.append(f"**To:** {m.to}")
    if m.cc: parts.append(f"**Cc:** {m.cc}")
    if m.date: parts.append(f"**Date:** {m.date}")
    if m.subject: parts.append(f"**Subject:** {m.subject}")
    parts.append("")
    body = (m.body or "").strip()
    if body:
        parts.append(body)
    try:
        m.close()
    except Exception:
        pass
    text = "\n".join(parts).strip()
    if len(text.splitlines()) <= 1:
        return "", None
    return text + "\n", None


def _extract_csv(path: str, *, caps: bool = True) -> tuple[str, str | None]:
    """CSV/TSV → one markdown table. Pipe-escaping + newline-strip mirror
    _extract_xlsx so a cell containing `|` or a newline can't corrupt the
    table. `caps` (default True, mining behavior) bounds rows/cells; read
    paths pass caps=False for full fidelity."""
    import csv as _csv
    delimiter = "\t" if path.lower().endswith(".tsv") else ","
    cell_cap = XLSX_CELL_MAX_CHARS if caps else None
    row_cap = XLSX_MAX_ROWS_PER_SHEET if caps else None
    try:
        with open(path, "r", newline="", errors="replace") as f:
            rows = list(_csv.reader(f, delimiter=delimiter))
    except Exception as e:
        return "", f"open failed: {type(e).__name__}: {e}"
    if not rows:
        return "", None

    def _cell(v: str) -> str:
        s = "" if v is None else str(v)
        if cell_cap is not None and len(s) > cell_cap:
            s = s[:cell_cap] + "…"
        return s.replace("|", "\\|").replace("\n", " ").strip()

    n_cols = max(len(r) for r in rows)
    header = [_cell(c) for c in rows[0]]
    while len(header) < n_cols:
        header.append("")
    parts = ["| " + " | ".join(header) + " |",
             "|" + "|".join(["---"] * n_cols) + "|"]
    for r in rows[1:]:
        if row_cap is not None and len(parts) - 2 >= row_cap:
            parts.append(f"\n_(truncated at {row_cap:,} rows)_\n")
            break
        rc = [_cell(c) for c in r[:n_cols]]
        while len(rc) < n_cols:
            rc.append("")
        parts.append("| " + " | ".join(rc) + " |")
    return "\n".join(parts) + "\n", None


def _extract_markitdown_only(path: str) -> tuple[str, str | None]:
    """Fallback for formats that have no native extractor (.epub, .zip).
    By the time _do_extract calls this, markitdown has already been tried
    (these exts are in _MARKITDOWN_EXTS), so reaching here means markitdown
    was absent or failed — return a helpful error, not silent empty."""
    ext = os.path.splitext(path)[1].lower()
    if not _MARKITDOWN_BIN:
        return "", (f"Reading {ext} requires the markitdown package "
                    "(pip3 install 'markitdown[outlook]').")
    return "", f"markitdown could not extract {ext}"


def _email_to_markdown(msg) -> str:
    """Shared rendering for both .eml (parsed via stdlib email) and any
    other future `email.message.Message` source."""
    title = msg.get("Subject") or "(no subject)"
    parts = [f"# {title}\n"]
    for hdr in ("From", "To", "Cc", "Date", "Subject"):
        v = msg.get(hdr)
        if v:
            parts.append(f"**{hdr}:** {v}")
    parts.append("")
    body = ""
    try:
        # Prefer the plain text body part.
        if msg.is_multipart():
            for part in msg.walk():
                ctype = part.get_content_type() or ""
                disp = (part.get("Content-Disposition") or "").lower()
                if "attachment" in disp:
                    continue
                if ctype == "text/plain":
                    body = part.get_content()
                    break
            if not body:
                # Fall back to html → strip tags lightly.
                for part in msg.walk():
                    if part.get_content_type() == "text/html":
                        import re
                        html = part.get_content()
                        body = re.sub(r"<[^>]+>", "", html or "")
                        break
        else:
            body = msg.get_content() or ""
    except Exception:
        body = ""
    body = (body or "").strip()
    if body:
        parts.append(body)
    return "\n".join(parts) + "\n"


_EXTRACTORS = {
    ".pdf": _extract_pdf,
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xlsm": _extract_xlsx,  # macro-enabled workbook — openpyxl reads cells the
                            # same way; VBA source is appended by _extract_vba
    ".xls": _extract_xlsx,   # legacy alias; true .xls isn't openpyxl-readable
    ".xlsb": _extract_xlsb,  # binary workbook — pyxlsb (openpyxl can't read it)
    ".csv": _extract_csv,
    ".tsv": _extract_csv,
    ".eml": _extract_eml,
    ".msg": _extract_msg,
    ".epub": _extract_markitdown_only,
    ".zip": _extract_markitdown_only,
}


def _iter_source_files(root: str) -> Iterable[str]:
    """Yield every supported source path in `root` (recursive). Skips the
    .brain-extracted subdir itself."""
    skip_prefix = os.path.join(root, EXTRACT_SUBDIR)
    for dirpath, _dirs, files in os.walk(root):
        if dirpath == skip_prefix or dirpath.startswith(skip_prefix + os.sep):
            continue
        for fn in files:
            if fn in SKIP_NAMES:
                continue
            ext = os.path.splitext(fn)[1].lower()
            if ext in SUPPORTED_EXTS:
                yield os.path.join(dirpath, fn)


# ── Single-file conversion (used by read_document + convert_folder) ──────────

def _adhoc_cache_path(src: str, mtime: int, size: int) -> str:
    """Companion-md path under ADHOC_CACHE_DIR for a file outside any
    project input folder. Key includes (mtime, size) so re-uploads at the
    same temp path get a fresh extraction automatically."""
    import hashlib
    key = f"{os.path.abspath(src)}:{mtime}:{size}".encode("utf-8")
    return os.path.join(ADHOC_CACHE_DIR, hashlib.sha256(key).hexdigest() + ".md")


def _pdf_ocr_or_empty(src: str) -> tuple[str, str, str | None]:
    """OCR a PDF that produced NO text from every text extractor (genuine
    image-only scan). Returns (text, backend, error). Routing by
    config.json -> ocr.engine: mistral_ocr (cloud) | local_vision (local) |
    auto (cloud→local) | none (→ empty marker). Shared by both the normal
    end-of-chain path and the pymupdf4llm-timeout→fitz-empty path."""
    ocr_cfg = _ocr_config()
    engine = ocr_cfg["engine"]
    ocr_text = ""
    ocr_err = ""
    pages = 0
    backend = ""
    if engine == "mlx_ocr":
        _extract_progress("OCR", note="Gescanntes PDF — MLX-OCR (lokal)")
        ocr_text, ocr_err = _pdf_mlx_ocr(src, ocr_cfg)
        if ocr_text:
            return ocr_text, "mlx-ocr", None
        if ocr_err:
            print(f"[doc-convert] OCR fallback failed for {src}: {ocr_err}",
                  flush=True)
        return "", "fitz/legacy", None
    if engine in ("mistral_ocr", "auto"):
        _extract_progress("OCR", note="Gescanntes PDF — Cloud-OCR")
        ocr_text, ocr_err, pages = _extract_with_mistral_ocr(src)
        if ocr_text:
            backend = f"mistral-ocr ({pages}p)"
    if not ocr_text and engine in ("local_vision", "auto"):
        _extract_progress("OCR", note="Gescanntes PDF — lokales Vision-Modell")
        local_text, local_err, local_pages = _extract_with_local_vision(src)
        if local_text:
            ocr_text = local_text
            pages = local_pages
            backend = f"local-vision ({pages}p)"
        else:
            ocr_err = ocr_err or local_err or "no ocr engine produced output"
    if ocr_text:
        return ocr_text, backend, None
    if engine != "none" and ocr_err:
        print(f"[doc-convert] OCR fallback failed for {src}: {ocr_err}", flush=True)
    return "", "fitz/legacy", None


def _do_extract(src: str, *, use_markitdown: bool = True,
                caps: bool = True, sheet: str | None = None,
                slides: str | None = None, pages: str | None = None,
                include_tables: bool = False, emit_meta: bool = False,
                page_marker: str = "## Page"
                ) -> tuple[str, str, str | None]:
    """Run markitdown (when enabled+supported) then per-format fallback,
    then OCR if both came back empty for PDFs.
    Returns (text, backend, error). Empty text + None error = no extractable
    content even after OCR — caller writes a marker file.

    Mining calls with all defaults (caps=True, no selection, no meta) → the
    fallback extractors produce byte-stable legacy output. read_document /
    read_attachment pass caps=False + the relevant selection/meta knobs.
    These knobs ONLY affect the per-format fallback — when markitdown
    succeeds (the common case for clean office files) its output is returned
    verbatim and the knobs are inert."""
    ext = os.path.splitext(src)[1].lower()
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        return "", "", f"unsupported extension: {ext}"
    md_enabled = bool(use_markitdown and _MARKITDOWN_BIN)
    text = ""
    # PDF engine choice (config conversion.pdf_engine, default pymupdf4llm).
    # pymupdf4llm renders tables/layout far better than markitdown on structured
    # PDFs; 'markitdown' keeps the old default; 'fitz' forces the plain legacy
    # _extract_pdf path. Tried BEFORE markitdown so the choice actually wins;
    # falls through to markitdown/fitz on empty/error (e.g. scanned → OCR).
    if ext == ".pdf" and _pdf_engine() == "pymupdf4llm":
        # SELF-CONTAINED pdf chain — deterministic fallback order:
        #   pymupdf4llm (tables/layout)  → on TIMEOUT or EMPTY →
        #   fitz get_text (fast, text layer)  → on EMPTY (true scan) →
        #   OCR.
        # markitdown is deliberately NOT in this path: it's slower AND worse on
        # the same PDFs (it bottoms out on pdfminer just like pymupdf4llm), so
        # falling to it after a pymupdf4llm timeout reproduced the original
        # 60s+ stall (chat 4aad5750). fitz is the correct next step.
        _timed_out = False
        p_text, p_err = "", None
        # Call DIRECTLY — no _run_with_timeout wrapper. _extract_pdf_pymupdf4llm
        # now runs pymupdf4llm in a hard-killable subprocess that owns the
        # timeout and raises _ExtractTimeout itself; a second thread-based
        # timeout on top would double the per-page budget and abandon a thread
        # that's only blocked in subprocess.run. The subprocess is the single
        # timeout authority, and it SIGKILLs the runaway (reclaiming the CPU).
        try:
            p_text, p_err = _extract_pdf_pymupdf4llm(src, pages=pages)
        except _ExtractTimeout:
            _timed_out = True
        if not _timed_out and not p_err and p_text:
            return p_text, "pymupdf4llm", None
        if p_err and not _timed_out:
            # Hard dep/error from pymupdf4llm (e.g. not installed) — let the
            # generic chain (markitdown / fitz legacy) handle it below.
            pass
        else:
            # Timeout OR empty result from pymupdf4llm on a (likely) text PDF →
            # bare fitz get_text. Reads the text layer in well under a second on
            # the very PDFs that hang pymupdf4llm.
            if _timed_out:
                print(f"[doc-convert] pymupdf4llm timed out after "
                      f"{_PDF_EXTRACT_TIMEOUT_SECS}s on {src} — bare fitz read "
                      f"(subprocess killed)",
                      flush=True)
            _extract_progress("fitz", note="Wechsel zu fitz (Textebene)")
            # Call fitz DIRECTLY: bare get_text is inherently fast (sub-second
            # even for dozens of pages) and releases the GIL in its C core.
            f_text, f_err = _extract_pdf_fitz_fast(src, pages=pages)
            if not f_err and f_text:
                return f_text, "fitz/fast", None
            if f_err:
                return "", "", f_err
            # fitz found no text layer → genuine scan → OCR.
            return _pdf_ocr_or_empty(src)
    if md_enabled and ext in _markitdown_exts() and not (ext == ".pdf" and _pdf_engine() == "fitz"):
        _extract_progress("markitdown", note="Extrahiere mit markitdown")
        text, err = _extract_with_markitdown(src)
        if not err and text:
            return text, "markitdown", None
    # Per-format fallback. Pass only the knobs each extractor accepts so a
    # default mining call stays byte-stable.
    extractor_kwargs: dict = {}
    if ext in (".xlsx", ".xlsm", ".xls", ".xlsb"):
        extractor_kwargs = {"caps": caps, "sheet": sheet}
    elif ext in (".csv", ".tsv"):
        extractor_kwargs = {"caps": caps}
    elif ext == ".pptx":
        extractor_kwargs = {"slides": slides}
    elif ext == ".pdf":
        extractor_kwargs = {"pages": pages, "include_tables": include_tables,
                            "emit_meta": emit_meta, "page_marker": page_marker}
    try:
        if ext == ".pdf":
            # Guard the fitz/pdfplumber path too (pdfplumber table detection can
            # hang on dense tables). On timeout, drop to a bare fitz get_text —
            # the fastest possible read of the text layer — before considering
            # the PDF "empty" (which would wrongly route a text PDF to OCR).
            try:
                text, err = _run_with_timeout(
                    lambda: extractor(src, **extractor_kwargs),
                    _PDF_EXTRACT_TIMEOUT_SECS)
            except _ExtractTimeout:
                print(f"[doc-convert] PDF extractor timed out after "
                      f"{_PDF_EXTRACT_TIMEOUT_SECS}s on {src} — bare fitz read",
                      flush=True)
                text, err = _extract_pdf_fitz_fast(src, pages=pages)
        else:
            text, err = extractor(src, **extractor_kwargs)
    except Exception as e:
        return "", "", f"{type(e).__name__}: {e}"
    if err:
        # Hard error from extractor (file corrupt, missing dep). Don't try
        # OCR — OCR is for image-only PDFs, not broken files.
        return "", "", err
    if text:
        return text, "fitz/legacy", None

    # Empty output from every text extractor → genuine image-only PDF → OCR.
    if ext == ".pdf":
        return _pdf_ocr_or_empty(src)
    return "", "fitz/legacy", None


def convert_one(src: str, *, project_root: str | None = None,
                use_markitdown: bool = True) -> tuple[str | None, str | None]:
    """Convert a single binary document to a companion .md, idempotently.

    Returns (md_path, error). md_path is None on hard failure; error is None
    on success.

    Companion location:
      - If `project_root` is given AND `src` is inside it → standard
        `<project_root>/.brain-extracted/<rel>.md` (matches the project-sync
        daemon layout — same byte-for-byte output for the same source).
      - Otherwise → ad-hoc cache under ADHOC_CACHE_DIR keyed by
        (abs_path, mtime, size).

    Idempotent: if the companion exists and its frontmatter (mtime, size)
    matches the source, returns immediately without re-extracting.

    Mirrors the convert_folder() logic for one file so behavior is identical
    between bulk pre-mining and on-demand reads.
    """
    if not src or not os.path.isfile(src):
        return None, f"source not a file: {src}"
    src = os.path.abspath(src)
    cur_mt, cur_sz = _stat_key(src)
    if cur_mt == 0 and cur_sz == 0:
        return None, "stat failed"

    # Pick companion location. Project files reuse the daemon's layout so
    # one PDF never has two companions.
    if project_root and os.path.commonpath([project_root, src]) == os.path.abspath(project_root):
        md_path = _md_path_for(os.path.abspath(project_root), src)
    else:
        md_path = _adhoc_cache_path(src, cur_mt, cur_sz)

    # Fast path: companion exists with matching (mtime, size).
    if os.path.isfile(md_path):
        old_mt, old_sz = _read_md_source_stat(md_path)
        if old_mt == cur_mt and old_sz == cur_sz:
            # Touch atime so LRU eviction sees recent use.
            try:
                os.utime(md_path, None)
            except OSError:
                pass
            return md_path, None

    try:
        os.makedirs(os.path.dirname(md_path), exist_ok=True)
    except OSError as e:
        return None, f"mkdir failed: {e}"

    text, backend, err = _do_extract(src, use_markitdown=use_markitdown)
    if err:
        return None, err
    if not text:
        # Scanned PDF / empty doc — write a marker so we don't retry every
        # call. Same behavior as convert_folder.
        text = (
            "# " + os.path.splitext(os.path.basename(src))[0] + "\n\n"
            "_(no extractable text — possibly a scanned image; OCR not run)_\n"
        )
        backend = backend or "empty"

    body = _frontmatter(src, cur_mt, cur_sz, backend=backend) + text
    try:
        tmp_path = md_path + ".tmp"
        with open(tmp_path, "w", encoding="utf-8") as f:
            f.write(body)
        os.replace(tmp_path, md_path)
    except OSError as e:
        return None, f"write failed: {e}"
    return md_path, None


def evict_adhoc_cache(*, ttl_secs: int = ADHOC_CACHE_TTL_SECS,
                      log_prefix: str = "[doc-convert]") -> int:
    """Drop ad-hoc cache entries older than `ttl_secs` (atime-based).
    Project companions under `.brain-extracted/` are never touched here.
    Returns count removed. Safe to call from any periodic daemon."""
    if not os.path.isdir(ADHOC_CACHE_DIR):
        return 0
    cutoff = time.time() - ttl_secs
    removed = 0
    for fn in os.listdir(ADHOC_CACHE_DIR):
        if not fn.endswith(".md"):
            continue
        p = os.path.join(ADHOC_CACHE_DIR, fn)
        try:
            st = os.stat(p)
        except OSError:
            continue
        if st.st_atime < cutoff:
            try:
                os.remove(p)
                removed += 1
            except OSError:
                pass
    if removed:
        print(f"{log_prefix} adhoc-cache evicted={removed}", flush=True)
    return removed


# ── Public API ───────────────────────────────────────────────────────────────

def convert_folder(root: str, *, log_prefix: str = "[doc-convert]",
                   use_markitdown: bool = True) -> ConvertResult:
    """Walk `root` for supported binary docs, write `.md` siblings under
    `<root>/.brain-extracted/<rel>.md`. Idempotent via (mtime, size) hash.
    Returns a result summary; does not raise on per-file failures.

    When `use_markitdown` is true and the markitdown CLI is on PATH, prefer
    it for .pdf/.docx/.pptx/.xlsx (better table + heading fidelity for LLM
    retrieval). Falls through to the per-format extractor on any failure.
    """
    res = ConvertResult()
    if not root or not os.path.isdir(root):
        return res
    t0 = time.time()
    md_enabled = bool(use_markitdown and _MARKITDOWN_BIN)
    md_used = 0
    md_fallback = 0
    extract_root = os.path.join(root, EXTRACT_SUBDIR)
    for src in _iter_source_files(root):
        res.seen_total += 1
        ext = os.path.splitext(src)[1].lower()
        extractor = _EXTRACTORS.get(ext)
        if extractor is None:
            continue
        cur_mt, cur_sz = _stat_key(src)
        if cur_mt == 0 and cur_sz == 0:
            res.skipped_unreadable += 1
            continue
        md_path = _md_path_for(root, src)
        if os.path.isfile(md_path):
            old_mt, old_sz = _read_md_source_stat(md_path)
            if old_mt == cur_mt and old_sz == cur_sz:
                res.skipped_unchanged += 1
                continue
        # (Re)convert.
        try:
            os.makedirs(os.path.dirname(md_path), exist_ok=True)
        except OSError as e:
            res.failed += 1
            res.failures.append(f"{src}: mkdir failed {e}")
            continue
        text = ""
        err: str | None = None
        backend = ""
        if md_enabled and ext in _markitdown_exts():
            text, err = _extract_with_markitdown(src)
            if err or not text:
                md_fallback += 1
                text = ""
                err = None
            else:
                backend = "markitdown"
                md_used += 1
        if not text:
            try:
                text, err = extractor(src)
            except Exception as e:
                res.failed += 1
                res.failures.append(f"{src}: {type(e).__name__}: {e}")
                continue
            if err:
                res.failed += 1
                res.failures.append(f"{src}: {err}")
                continue
            if text:
                backend = backend or "fitz/legacy"
        if not text:
            # Empty after extraction — e.g. scanned PDF. Write a marker file
            # so we don't retry every cycle, but keep it small enough that the
            # miner doesn't flood drawer storage.
            text = (
                "# " + os.path.splitext(os.path.basename(src))[0] + "\n\n"
                "_(no extractable text — possibly a scanned image; OCR not run)_\n"
            )
            backend = backend or "empty"
        body = _frontmatter(src, cur_mt, cur_sz, backend=backend) + text
        try:
            tmp_path = md_path + ".tmp"
            with open(tmp_path, "w", encoding="utf-8") as f:
                f.write(body)
            os.replace(tmp_path, md_path)
            res.converted += 1
        except OSError as e:
            res.failed += 1
            res.failures.append(f"{src}: write failed {e}")
            continue

    res.elapsed_s = time.time() - t0
    if res.converted or res.failed:
        md_note = ""
        if md_enabled:
            md_note = f" markitdown={md_used} fallback={md_fallback}"
        elif use_markitdown and not _MARKITDOWN_BIN:
            md_note = " markitdown=unavailable"
        print(
            f"{log_prefix} {root}: converted={res.converted} "
            f"unchanged={res.skipped_unchanged} failed={res.failed} "
            f"seen={res.seen_total} elapsed={res.elapsed_s:.1f}s{md_note}",
            flush=True)
    return res


def sweep_stale(root: str, *, log_prefix: str = "[doc-convert]") -> int:
    """Walk `<root>/.brain-extracted/` and drop any .md whose source path
    no longer exists. Returns count removed.
    """
    extract_root = os.path.join(root, EXTRACT_SUBDIR)
    if not os.path.isdir(extract_root):
        return 0
    removed = 0
    for dirpath, _dirs, files in os.walk(extract_root):
        for fn in files:
            if not fn.endswith(".md"):
                continue
            md_path = os.path.join(dirpath, fn)
            src = _source_for_md(root, md_path)
            if not os.path.isfile(src):
                try:
                    os.remove(md_path)
                    removed += 1
                except OSError:
                    pass
    if removed:
        print(f"{log_prefix} {root}: stale_removed={removed}", flush=True)
    return removed


def resolve_original_source(md_path: str) -> str | None:
    """Given a converted .md path, read its frontmatter and return the
    original source path. None if no marker present."""
    try:
        with open(md_path, "r", encoding="utf-8", errors="replace") as f:
            head = f.read(1024)
    except OSError:
        return None
    for line in head.splitlines():
        line = line.strip()
        if line.startswith("<!-- brain-source:"):
            return line.split(":", 1)[1].rstrip(" -->").strip() or None
    return None
