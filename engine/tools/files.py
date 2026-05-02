# Extracted from claude_cli.py — file/document/shell/exec tools
# This module is imported by claude_cli.py which sets up all shared globals
# before these functions are called. Cross-module references resolved via
# the claude_cli module namespace at call time.

import os
import re
import sys
import time
import subprocess
import fnmatch
import glob as globmod


def tool_read_file(args: dict) -> str:
    node_result = _route_to_node("read_file", args)
    if node_result is not None:
        return node_result
    path = args.get("path", "")
    offset = args.get("offset", 1)
    limit = args.get("limit", 400)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        # Cache hit only when caller asked for the whole file (default offset=1
        # AND no explicit limit override — i.e. they accepted the default 400).
        # If the model is paginating, always re-read on disk so it gets the
        # window it asked for.
        _full_read = (int(offset or 1) == 1 and ("limit" not in args))
        if _full_read and os.path.exists(path):
            _stub = _read_doc_cache_lookup(path)
            if _stub is not None:
                return _stub
        if not os.path.exists(path):
            return _err(
                f"File not found: {path}. "
                "Do NOT retry with the same path — the file does not exist on disk. "
                "Hallucinated paths are a frequent cause: the path was inferred "
                "from the project's folder structure or filename schema rather "
                "than copied from a real source. "
                "USE ONLY paths returned in `read_path` (or `read_path_original` "
                "as fallback) of a prior `mempalace_query` drawer. "
                "If no drawer pointed at this content, the project does not "
                "contain it — refuse per REFUSAL DISCIPLINE instead of guessing "
                "another path."
            )
        with open(path, "r", errors="replace") as f:
            lines = f.readlines()
        total = len(lines)
        start = max(0, offset - 1)
        end = start + limit if limit else total
        selected = lines[start:end]
        # Number lines
        numbered = []
        for i, line in enumerate(selected, start=start + 1):
            numbered.append(f"{i:>6}\t{line.rstrip()}")
        content = "\n".join(numbered)
        if _full_read and end >= total:
            try:
                _round = int(getattr(_thread_local, "tool_round", 0) or 0)
            except Exception:
                _round = 0
            _read_doc_cache_store(path, content, tool_round=_round)
        return _ok({"path": path, "total_lines": total, "showing": f"{start+1}-{min(end, total)}", "content": content})
    except Exception as e:
        return _err(f"read_file: {e}")


def _maybe_qmd_reindex(path: str) -> None:
    """If path is a .md file inside an agent dir, trigger debounced QMD reindex."""
    if not path.endswith(".md"):
        return
    agents_dir = os.path.realpath(AGENTS_DIR)
    real_path = os.path.realpath(path)
    if not real_path.startswith(agents_dir + os.sep):
        return
    rel = real_path[len(agents_dir) + 1:]
    collection = rel.split(os.sep)[0]
    _qmd_debounced_embed(collection)


def _get_artifact_session_folder(session_id: str) -> str:
    """Return session folder name for artifacts: <date>_<session_prefix>"""
    cache_key = f"_artifact_folder_{session_id}"
    cached = getattr(_thread_local, cache_key, None)
    if cached:
        return cached
    from datetime import datetime as _dt
    folder = f"{_dt.now().strftime('%Y-%m-%d')}_{session_id[:8]}"
    setattr(_thread_local, cache_key, folder)
    return folder


def tool_write_file(args: dict) -> str:
    node_result = _route_to_node("write_file", args)
    if node_result is not None:
        return node_result
    path = args.get("path", "")
    content = args.get("content", "")
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            # Default relative paths to artifacts session folder during chat
            session_id = getattr(_thread_local, 'current_session_id', None)
            agent = getattr(_thread_local, 'current_agent', None) or _current_agent
            if session_id and agent:
                folder = _get_artifact_session_folder(session_id)
                artifact_dir = os.path.join(AGENTS_DIR, agent.agent_id, "artifacts", folder)
                os.makedirs(artifact_dir, exist_ok=True)
                path = os.path.join(artifact_dir, path)
            else:
                path = os.path.abspath(path)
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        with open(path, "w") as f:
            f.write(content)
        size = os.path.getsize(path)
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        _after_file_write(path, "created", agent.agent_id if agent else "main")
        return _ok({"path": path, "size": size, "status": "written"})
    except Exception as e:
        return _err(f"write_file: {e}")


def tool_edit_file(args: dict) -> str:
    path = args.get("path", "")
    old_string = args.get("old_string", "")
    new_string = args.get("new_string", "")
    replace_all = args.get("replace_all", False)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        with open(path, "r") as f:
            content = f.read()
        count = content.count(old_string)
        if count == 0:
            return _err(f"edit_file: old_string not found in {path}")
        if count > 1 and not replace_all:
            return _err(f"edit_file: old_string found {count} times — use replace_all=true or provide a more specific match")
        if replace_all:
            new_content = content.replace(old_string, new_string)
        else:
            new_content = content.replace(old_string, new_string, 1)
        with open(path, "w") as f:
            f.write(new_content)
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        _after_file_write(path, "modified", agent.agent_id if agent else "main")
        return _ok({"path": path, "replacements": count if replace_all else 1, "status": "edited"})
    except Exception as e:
        return _err(f"edit_file: {e}")


def _describe_image_with_vision(image_data_b64: str, media_type: str, filename: str) -> str:
    """Use a vision-capable model to describe an image attachment."""
    vision_model = getattr(_thread_local, 'attachment_image_model', '') or ''
    if not vision_model:
        return f"(Image: {filename} — no image model configured. Set attachments.image_model in config.)"

    # Build multimodal message with image
    content_blocks = [
        {
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": media_type,
                "data": image_data_b64,
            },
        },
        {"type": "text", "text": f"Describe this image ({filename}) in detail. Include any text, data, diagrams, or visual elements you see."},
    ]

    try:
        result = _run_delegate(
            messages=[{"role": "user", "content": content_blocks}],
            model=vision_model,
            system_prompt="You are a precise image description assistant. Describe the image content thoroughly and concisely.",
            inference_params={"max_tokens": 2048, "temperature": 0.1},
            tools=False,
        )
        if result:
            return f"**Image: {filename}**\n\n{result}"
        return f"(Image: {filename} — vision model returned no description)"
    except Exception as e:
        return f"(Image: {filename} — vision error: {e})"


def tool_read_attachment(args: dict) -> str:
    """Read a user-attached file from the session attachment store."""
    name = args.get("name", "")
    if not name:
        return _err("read_attachment: 'name' is required")
    attachments = getattr(_thread_local, 'attachments', None) or {}
    if name not in attachments:
        # Try case-insensitive match
        for k in attachments:
            if k.lower() == name.lower():
                name = k
                break
        else:
            available = list(attachments.keys())
            if available:
                return _err(f"Attachment '{name}' not found. Available: {', '.join(available)}")
            return _err("No attachments in this session.")
    att = attachments[name]
    content = att.get("content", "")
    encoding = att.get("encoding", "text")
    media_type = att.get("media_type", "text/plain")
    ext = os.path.splitext(name)[1].lower()

    # --- Binary files (base64 encoded) ---
    if encoding == "base64":
        import base64 as b64_mod
        import io as io_mod
        try:
            raw_bytes = b64_mod.b64decode(content)
        except Exception as e:
            return _err(f"Failed to decode attachment: {e}")

        # PDF
        if ext == ".pdf" or "pdf" in media_type.lower():
            try:
                import fitz  # pymupdf
                doc = fitz.open(stream=raw_bytes, filetype="pdf")
                pages = []
                for i, page in enumerate(doc):
                    text = page.get_text()
                    if text.strip():
                        pages.append(f"--- Page {i+1} ---\n{text}")
                doc.close()
                return "\n\n".join(pages) if pages else "(PDF has no extractable text)"
            except ImportError:
                return _err("Install pymupdf for PDF support: pip3 install pymupdf")

        # DOCX
        if ext == ".docx":
            try:
                import docx
                doc = docx.Document(io_mod.BytesIO(raw_bytes))
                return DocumentParser.parse_docx.__func__(None) if False else _parse_docx_from_bytes(raw_bytes)
            except ImportError:
                return _err("Install python-docx for DOCX support: pip3 install python-docx")

        # XLSX
        if ext == ".xlsx":
            try:
                return _parse_xlsx_from_bytes(raw_bytes, sheet=args.get("sheet"))
            except ImportError:
                return _err("Install openpyxl for XLSX support: pip3 install openpyxl")

        # PPTX
        if ext == ".pptx":
            try:
                return _parse_pptx_from_bytes(raw_bytes)
            except ImportError:
                return _err("Install python-pptx for PPTX support: pip3 install python-pptx")

        # Images — use vision model
        if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg") or media_type.startswith("image/"):
            return _describe_image_with_vision(content, media_type, name)

        # Unknown binary
        return f"(Binary file: {name}, {len(raw_bytes)} bytes, type: {media_type})"

    # --- Text files ---
    return content


def _parse_docx_from_bytes(raw_bytes: bytes) -> str:
    """Parse DOCX from in-memory bytes."""
    import io, docx
    doc = docx.Document(io.BytesIO(raw_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                level = 1
                try:
                    level = int(para.style.name.replace("Heading", "").strip()) or 1
                except ValueError:
                    pass
                text = "#" * level + " " + text
            paragraphs.append(text)
    # Tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            # Add header separator after first row
            header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
            rows.insert(1, header_sep)
            paragraphs.append("\n".join(rows))
    return "\n\n".join(paragraphs)


def _parse_xlsx_from_bytes(raw_bytes: bytes, sheet: str | None = None) -> str:
    """Parse XLSX from in-memory bytes."""
    import io, openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    sheets = [wb[sheet]] if sheet and sheet in wb.sheetnames else wb.worksheets
    parts = []
    for ws in sheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # Build markdown table
        header = rows[0]
        cols = [str(c) if c is not None else "" for c in header]
        lines = [f"**Sheet: {ws.title}**"]
        lines.append("| " + " | ".join(cols) + " |")
        lines.append("| " + " | ".join(["---"] * len(cols)) + " |")
        for row in rows[1:200]:  # Cap at 200 rows
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("| " + " | ".join(cells) + " |")
        if len(rows) > 201:
            lines.append(f"*... ({len(rows) - 201} more rows)*")
        parts.append("\n".join(lines))
    wb.close()
    return "\n\n".join(parts) if parts else "(Empty spreadsheet)"


def _parse_pptx_from_bytes(raw_bytes: bytes) -> str:
    """Parse PPTX from in-memory bytes."""
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append("| " + " | ".join(cells) + " |")
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    # Notes
    for i, slide in enumerate(prs.slides, 1):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slides.append(f"[Slide {i} notes] {notes}")
    return "\n\n".join(slides) if slides else "(Empty presentation)"


def tool_read_document(args: dict) -> str:
    """Format-aware document reader."""
    # Accept legacy / drawer-vocab synonyms — drawers from mempalace_query
    # carry the field as `source_file`, and the model frequently passes that
    # name verbatim. Treating an empty `path` as CWD silently expanded into
    # "Is a directory" errors that the model then ignored.
    path = args.get("path", "") or args.get("source_file", "") or args.get("file", "")
    if not path or not str(path).strip():
        return _err(
            "read_document: 'path' is required (got empty). When following up "
            "on a mempalace_query drawer, take its `source_file` value and "
            "pass it as the `path` argument — JOIN with the input-folder "
            "absolute path if `source_file` is relative.")
    path = str(path).strip()
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        if os.path.isdir(path):
            return _err(
                f"read_document: '{path}' is a directory, not a file. The "
                "model commonly passes a base path here when it meant to "
                "pass a specific document — re-issue the call with the full "
                "file path including extension.")
        if not os.path.exists(path):
            return _err(
                f"File not found: {path}. "
                "Do NOT retry with the same path — the file does not exist on disk. "
                "Hallucinated paths are a frequent cause: the path was inferred "
                "from the project's folder structure or filename schema rather "
                "than copied from a real source. "
                "USE ONLY paths returned in `read_path` (or `read_path_original` "
                "as fallback) of a prior `mempalace_query` drawer. "
                "If no drawer pointed at this content, the project does not "
                "contain it — refuse per REFUSAL DISCIPLINE instead of guessing "
                "another path."
            )

        # Per-session cache: a previous turn already read this exact file,
        # mtime+size unchanged → return a stub instead of re-streaming all
        # bytes back into the model's context. Skip the cache if the model
        # is paginating (offset / limit / pages / sheet / slides) — those
        # signal it wants a specific window, not the whole file.
        _is_paginated = any(
            args.get(k) for k in ("offset", "limit", "pages", "sheet", "slides")
        )
        if not _is_paginated:
            _stub = _read_doc_cache_lookup(path)
            if _stub is not None:
                return _stub

        def _ok_and_cache(payload: dict) -> str:
            # Cache only the full-file shape so subsequent paginated reads
            # always hit disk; that's what _is_paginated already gates above.
            if not _is_paginated:
                try:
                    _round = int(getattr(_thread_local, "tool_round", 0) or 0)
                except Exception:
                    _round = 0
                _read_doc_cache_store(path, str(payload.get("content", "") or ""), tool_round=_round)
            return _ok(payload)

        ext = os.path.splitext(path)[1].lower()

        if ext == ".pdf":
            try:
                import fitz
            except ImportError:
                return _err("Install pymupdf: pip3 install pymupdf")
            doc = fitz.open(path)
            meta = {
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "page_count": doc.page_count,
            }
            pages_param = args.get("pages", "")
            page_indices = None
            if pages_param:
                page_indices = set()
                for part in pages_param.split(","):
                    part = part.strip()
                    if "-" in part:
                        a, b = part.split("-", 1)
                        for i in range(int(a), int(b) + 1):
                            page_indices.add(i)
                    else:
                        page_indices.add(int(part))
            page_texts = []
            for i, page in enumerate(doc, 1):
                if page_indices and i not in page_indices:
                    continue
                page_texts.append(f"--- Page {i} ---\n{page.get_text()}")
            doc.close()

            include_tables = bool(args.get("include_tables"))
            tables_by_page: dict[int, list[str]] = {}
            tables_note = ""
            if include_tables:
                try:
                    import pdfplumber
                    with pdfplumber.open(path) as plumb:
                        for i, page in enumerate(plumb.pages, 1):
                            if page_indices and i not in page_indices:
                                continue
                            found = page.extract_tables()
                            if not found:
                                continue
                            rendered = []
                            for t in found:
                                rows = [[(c or "").replace("|", "\\|").replace("\n", " ").strip() for c in row] for row in t]
                                if not rows:
                                    continue
                                width = max(len(r) for r in rows)
                                for r in rows:
                                    r.extend([""] * (width - len(r)))
                                header = "| " + " | ".join(rows[0]) + " |"
                                sep = "| " + " | ".join("---" for _ in range(width)) + " |"
                                body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
                                rendered.append(header + "\n" + sep + ("\n" + body if body else ""))
                            if rendered:
                                tables_by_page[i] = rendered
                except ImportError:
                    tables_note = "\n\n*(pdfplumber not installed — install with: pip3 install pdfplumber)*"

            if tables_by_page:
                merged = []
                for block in page_texts:
                    first = block.split("\n", 1)[0]
                    page_num = None
                    if first.startswith("--- Page ") and first.endswith(" ---"):
                        try:
                            page_num = int(first[9:-4])
                        except ValueError:
                            pass
                    merged.append(block)
                    if page_num is not None and page_num in tables_by_page:
                        for j, md in enumerate(tables_by_page[page_num], 1):
                            merged.append(f"### Table (page {page_num}, #{j})\n{md}")
                content = "\n\n".join(merged)
            else:
                content = "\n\n".join(page_texts) + tables_note

            meta_str = "\n".join(f"**{k}:** {v}" for k, v in meta.items() if v)
            return _ok_and_cache({"path": path, "format": "pdf", "metadata": meta_str, "content": content})

        elif ext == ".docx":
            try:
                import docx
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            doc = docx.Document(path)
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    if para.style and para.style.name and para.style.name.startswith("Heading"):
                        level = 1
                        try:
                            level = int(para.style.name.replace("Heading", "").strip()) or 1
                        except ValueError:
                            pass
                        text = "#" * level + " " + text
                    paragraphs.append(text)
            # Extract tables
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    max_cols = max(len(r) for r in rows)
                    for r in rows:
                        while len(r) < max_cols:
                            r.append("")
                    header = "| " + " | ".join(rows[0]) + " |"
                    sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
                    table_lines = [header, sep]
                    for r in rows[1:]:
                        table_lines.append("| " + " | ".join(r) + " |")
                    paragraphs.append("\n".join(table_lines))
            content = "\n\n".join(paragraphs)
            return _ok_and_cache({"path": path, "format": "docx", "content": content})

        elif ext in (".xlsx", ".xls"):
            sheet = args.get("sheet")
            content = DocumentParser.parse_xlsx(path, sheet=sheet)
            return _ok_and_cache({"path": path, "format": "xlsx", "content": content})

        elif ext == ".pptx":
            slides = args.get("slides")
            content = DocumentParser.parse_pptx(path, slides=slides)
            return _ok_and_cache({"path": path, "format": "pptx", "content": content})

        elif ext in (".csv", ".tsv"):
            content = DocumentParser.parse_csv(path)
            return _ok_and_cache({"path": path, "format": "csv", "content": content})

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"):
            meta_text = DocumentParser.parse_image(path)
            vision_note = "\n\n*(For AI-powered image description, include this image directly in your chat message)*"
            return _ok({"path": path, "format": "image", "content": meta_text + vision_note})

        elif ext == ".svg":
            content = DocumentParser.parse_svg(path)
            return _ok_and_cache({"path": path, "format": "svg", "content": content})

        elif ext == ".eml":
            import email
            from email import policy
            with open(path, "rb") as f:
                msg = email.message_from_bytes(f.read(), policy=policy.default)
            headers = {
                "from": str(msg.get("From", "")),
                "to": str(msg.get("To", "")),
                "cc": str(msg.get("Cc", "")),
                "subject": str(msg.get("Subject", "")),
                "date": str(msg.get("Date", "")),
            }
            body_part = msg.get_body(preferencelist=("plain", "html"))
            body = body_part.get_content() if body_part else ""
            if body_part and body_part.get_content_type() == "text/html":
                body = re.sub(r"<[^>]+>", " ", body)
                body = re.sub(r"\s+", " ", body).strip()
            attachments = [p.get_filename() for p in msg.iter_attachments() if p.get_filename()]
            meta_str = "\n".join(f"**{k}:** {v}" for k, v in headers.items() if v)
            if attachments:
                meta_str += f"\n**attachments:** {', '.join(attachments)}"
            return _ok_and_cache({"path": path, "format": "eml", "content": f"{meta_str}\n\n{body}"})

        elif ext in (".msg", ".epub", ".zip"):
            try:
                from markitdown import MarkItDown
            except ImportError:
                return _err(
                    f"Reading {ext} files requires the markitdown package, which is not "
                    "included in airgapped installer builds. On dev machines with internet: "
                    "pip3 install 'markitdown[outlook]'."
                )
            md = MarkItDown()
            result = md.convert(path)
            fmt = ext.lstrip(".")
            return _ok_and_cache({"path": path, "format": fmt, "content": result.text_content})

        else:
            # Plain-text / markdown / unknown-extension read. Honor explicit
            # offset+limit when the model paginates; otherwise read the whole
            # file. The previous hard-cap at 500 lines truncated mid-document
            # silently — on the kg-real-policies ISMS Handbuch (1903 lines)
            # the section 2.13 percent-to-score table sits at line 1267,
            # WAY past line 500, so the model never saw the table even when
            # it correctly chose read_document on the .md companion. The
            # tool_result_char_limit middleware will compact downstream if
            # the content is too large for the round budget.
            offset = int(args.get("offset", 1) or 1)
            limit = args.get("limit")
            with open(path, "r", errors="replace") as f:
                lines = f.readlines()
            total = len(lines)
            start = max(0, offset - 1)
            end = start + int(limit) if limit else total
            selected = lines[start:end]
            numbered = []
            for i, line in enumerate(selected, start=start + 1):
                numbered.append(f"{i:>6}\t{line.rstrip()}")
            content = "\n".join(numbered)
            shown = f"{start+1}-{min(end, total)}"
            return _ok_and_cache({"path": path, "format": "text",
                        "total_lines": total, "showing": shown,
                        "content": content})
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"read_document: {e}")


def tool_write_document(args: dict) -> str:
    """Create documents from markdown content."""
    path = args.get("path", "")
    content = args.get("content", "")
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        ext = os.path.splitext(path)[1].lower()

        if ext == ".docx":
            try:
                import docx
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            doc = docx.Document()
            lines = content.split("\n")
            i = 0
            while i < len(lines):
                line = lines[i]
                # Headings
                heading_match = re.match(r'^(#{1,6})\s+(.*)', line)
                if heading_match:
                    level = len(heading_match.group(1))
                    doc.add_heading(heading_match.group(2), level=level)
                    i += 1
                    continue
                # Table detection
                if "|" in line and i + 1 < len(lines) and re.match(r'^\|[\s\-:|]+\|', lines[i + 1]):
                    table_rows = []
                    while i < len(lines) and "|" in lines[i]:
                        cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                        if not re.match(r'^[\s\-:|]+$', lines[i].strip().strip("|")):
                            table_rows.append(cells)
                        i += 1
                    if table_rows:
                        max_cols = max(len(r) for r in table_rows)
                        table = doc.add_table(rows=len(table_rows), cols=max_cols)
                        table.style = "Table Grid"
                        for ri, row_data in enumerate(table_rows):
                            for ci, cell_val in enumerate(row_data):
                                if ci < max_cols:
                                    table.rows[ri].cells[ci].text = cell_val
                    continue
                # Regular paragraph with inline formatting
                stripped = line.strip()
                if stripped:
                    para = doc.add_paragraph()
                    parts = re.split(r'(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*.*?\*)', stripped)
                    for part in parts:
                        if part.startswith("***") and part.endswith("***"):
                            run = para.add_run(part[3:-3])
                            run.bold = True
                            run.italic = True
                        elif part.startswith("**") and part.endswith("**"):
                            run = para.add_run(part[2:-2])
                            run.bold = True
                        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
                            run = para.add_run(part[1:-1])
                            run.italic = True
                        else:
                            para.add_run(part)
                i += 1
            doc.save(path)

        elif ext == ".xlsx":
            try:
                import openpyxl
            except ImportError:
                return _err("Install openpyxl: pip3 install openpyxl")
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            sections = re.split(r'^##\s+(.+)$', content, flags=re.MULTILINE)
            if len(sections) < 3:
                ws = wb.create_sheet("Sheet1")
                _write_md_table_to_sheet(ws, content)
            else:
                for si in range(1, len(sections), 2):
                    sheet_name = sections[si].strip()
                    if sheet_name.lower().startswith("sheet:"):
                        sheet_name = sheet_name[6:].strip()
                    sheet_content = sections[si + 1] if si + 1 < len(sections) else ""
                    ws = wb.create_sheet(sheet_name[:31])
                    _write_md_table_to_sheet(ws, sheet_content)
            if not wb.sheetnames:
                wb.create_sheet("Sheet1")
            wb.save(path)

        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return _err("Install python-pptx: pip3 install python-pptx")
            prs = Presentation()
            slides_content = re.split(r'^#\s+(.+)$', content, flags=re.MULTILINE)
            if len(slides_content) < 3:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = "Slide 1"
                slide.placeholders[1].text = content.strip()
            else:
                for si in range(1, len(slides_content), 2):
                    title = slides_content[si].strip()
                    body = slides_content[si + 1].strip() if si + 1 < len(slides_content) else ""
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = title
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    body_lines = [l for l in body.split("\n") if l.strip()]
                    for li, bline in enumerate(body_lines):
                        bline = bline.strip()
                        bline = re.sub(r'^[-*]\s+', '', bline)
                        if li == 0:
                            tf.text = bline
                        else:
                            p = tf.add_paragraph()
                            p.text = bline
            prs.save(path)

        elif ext == ".pdf":
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
            except ImportError:
                return _err("Install reportlab: pip3 install reportlab")
            doc_pdf = SimpleDocTemplate(path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 12))
                    continue
                heading_match = re.match(r'^(#{1,6})\s+(.*)', line)
                if heading_match:
                    level = len(heading_match.group(1))
                    style_name = f"Heading{min(level, 6)}"
                    if style_name not in styles:
                        style_name = "Heading1"
                    story.append(Paragraph(heading_match.group(2), styles[style_name]))
                else:
                    line = re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', line)
                    line = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', line)
                    line = re.sub(r'\*(.+?)\*', r'<i>\1</i>', line)
                    story.append(Paragraph(line, styles["Normal"]))
            doc_pdf.build(story)

        else:
            return _err(f"write_document: unsupported format '{ext}'. Supported: .docx, .xlsx, .pptx, .pdf")

        size = os.path.getsize(path)
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        _after_file_write(path, "created", agent.agent_id if agent else "main")
        return _ok({"path": path, "size": size, "format": ext.lstrip("."), "status": "written"})
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"write_document: {e}")


def _write_md_table_to_sheet(ws, md_text: str) -> None:
    """Helper: parse markdown table text and write rows to an openpyxl worksheet."""
    row_idx = 1
    for line in md_text.split("\n"):
        line = line.strip()
        if not line or not line.startswith("|"):
            continue
        if re.match(r'^\|[\s\-:|]+\|$', line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        for ci, val in enumerate(cells, 1):
            try:
                ws.cell(row=row_idx, column=ci, value=float(val) if "." in val else int(val))
            except (ValueError, TypeError):
                ws.cell(row=row_idx, column=ci, value=val)
        row_idx += 1


def tool_edit_document(args: dict) -> str:
    """Targeted edits to existing documents."""
    path = args.get("path", "")
    action = args.get("action", "")
    params = args.get("params", {})
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        if not os.path.exists(path):
            return _err(f"File not found: {path}")
        ext = os.path.splitext(path)[1].lower()

        if ext == ".docx":
            if action != "replace_text":
                return _err(f"edit_document: unsupported action '{action}' for DOCX. Use: replace_text")
            try:
                import docx
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            old_text = params.get("old_text", "")
            new_text = params.get("new_text", "")
            if not old_text:
                return _err("edit_document: 'old_text' required in params")
            doc = docx.Document(path)
            count = 0
            for para in doc.paragraphs:
                if old_text in para.text:
                    full_text = para.text
                    new_full = full_text.replace(old_text, new_text)
                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = new_full
                    else:
                        para.add_run(new_full)
                    count += 1
            doc.save(path)
            agent = getattr(_thread_local, 'current_agent', None) or _current_agent
            _after_file_write(path, "modified", agent.agent_id if agent else "main")
            return _ok({"path": path, "action": action, "replacements": count, "status": "edited"})

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
            except ImportError:
                return _err("Install openpyxl: pip3 install openpyxl")
            wb = openpyxl.load_workbook(path)

            if action == "update_cell":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                cell_ref = params.get("cell", "")
                value = params.get("value", "")
                if not cell_ref:
                    return _err("edit_document: 'cell' required (e.g. 'A1')")
                if sheet_name not in wb.sheetnames:
                    return _err(f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}")
                ws = wb[sheet_name]
                try:
                    ws[cell_ref] = float(value) if isinstance(value, str) and value.replace(".", "", 1).replace("-", "", 1).isdigit() else value
                except Exception:
                    ws[cell_ref] = value
                wb.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "cell": cell_ref, "sheet": sheet_name, "status": "edited"})

            elif action == "add_row":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                values = params.get("values", [])
                if not values:
                    return _err("edit_document: 'values' array required")
                if sheet_name not in wb.sheetnames:
                    return _err(f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}")
                ws = wb[sheet_name]
                ws.append(values)
                wb.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "sheet": sheet_name, "row_added": len(values), "status": "edited"})

            else:
                return _err(f"edit_document: unsupported action '{action}' for XLSX. Use: update_cell, add_row")

        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return _err("Install python-pptx: pip3 install python-pptx")
            prs = Presentation(path)

            if action == "update_slide":
                slide_index = int(params.get("slide_index", 1))
                if slide_index < 1 or slide_index > len(prs.slides):
                    return _err(f"Slide index {slide_index} out of range (1-{len(prs.slides)})")
                slide = prs.slides[slide_index - 1]
                title = params.get("title")
                body = params.get("body")
                if title and slide.shapes.title:
                    slide.shapes.title.text = title
                if body:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            shape.text_frame.text = body
                            break
                prs.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "slide_index": slide_index, "status": "edited"})

            elif action == "add_slide":
                title = params.get("title", "New Slide")
                body = params.get("body", "")
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                if body:
                    slide.placeholders[1].text = body
                prs.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "slide_count": len(prs.slides), "status": "edited"})

            else:
                return _err(f"edit_document: unsupported action '{action}' for PPTX. Use: update_slide, add_slide")

        else:
            return _err(f"edit_document: unsupported format '{ext}'. Supported: .docx, .xlsx, .xls, .pptx")
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"edit_document: {e}")


def tool_list_directory(args: dict) -> str:
    node_result = _route_to_node("list_directory", args)
    if node_result is not None:
        return node_result
    path = args.get("path", ".")
    pattern = args.get("pattern")
    recursive = args.get("recursive", False)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)

        if pattern:
            if recursive or "**" in pattern:
                full_pattern = os.path.join(path, pattern)
                entries = globmod.glob(full_pattern, recursive=True)
            else:
                full_pattern = os.path.join(path, pattern)
                entries = globmod.glob(full_pattern)
        elif recursive:
            entries = []
            for root, dirs, files in os.walk(path):
                # Skip hidden dirs
                dirs[:] = [d for d in dirs if not d.startswith(".")]
                for f in files:
                    if not f.startswith("."):
                        entries.append(os.path.join(root, f))
        else:
            entries = [os.path.join(path, e) for e in os.listdir(path)]

        results = []
        for entry in sorted(entries)[:500]:
            try:
                st = os.stat(entry)
                is_dir = os.path.isdir(entry)
                results.append({
                    "name": os.path.relpath(entry, path),
                    "type": "directory" if is_dir else "file",
                    "size": st.st_size if not is_dir else None,
                })
            except OSError:
                results.append({"name": os.path.relpath(entry, path), "type": "unknown"})

        return _ok({"path": path, "count": len(results), "entries": results})
    except Exception as e:
        return _err(f"list_directory: {e}")


def tool_search_files(args: dict) -> str:
    pattern = args.get("pattern", "")
    path = args.get("path", ".")
    file_glob = args.get("glob")
    case_insensitive = args.get("case_insensitive", False)
    max_results = args.get("max_results", 50)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)

        flags = re.IGNORECASE if case_insensitive else 0
        regex = re.compile(pattern, flags)

        matches = []
        files_searched = 0

        if os.path.isfile(path):
            file_list = [path]
        else:
            file_list = []
            for root, dirs, files in os.walk(path):
                dirs[:] = [d for d in dirs if not d.startswith(".") and d not in ("node_modules", "__pycache__", ".git")]
                for f in files:
                    if f.startswith("."):
                        continue
                    fp = os.path.join(root, f)
                    if file_glob and not fnmatch.fnmatch(f, file_glob):
                        continue
                    file_list.append(fp)

        for fp in file_list:
            if len(matches) >= max_results:
                break
            files_searched += 1
            try:
                with open(fp, "r", errors="replace") as fh:
                    for lineno, line in enumerate(fh, 1):
                        if regex.search(line):
                            matches.append({
                                "file": os.path.relpath(fp, path) if not os.path.isfile(path) else fp,
                                "line": lineno,
                                "text": line.rstrip()[:500],
                            })
                            if len(matches) >= max_results:
                                break
            except (OSError, UnicodeDecodeError):
                continue

        return _ok({"pattern": pattern, "path": path, "files_searched": files_searched,
                     "match_count": len(matches), "matches": matches})
    except re.error as e:
        return _err(f"search_files: invalid regex: {e}")
    except Exception as e:
        return _err(f"search_files: {e}")


def _strip_ansi(text: str) -> str:
    """Remove all ANSI escape sequences, control chars, and normalize whitespace."""
    text = re.sub(r"\033\[[0-9;]*[a-zA-Z]", "", text)
    text = re.sub(r"\033\[\?[0-9;]*[a-zA-Z]", "", text)
    text = re.sub(r"\033\([A-Z]", "", text)
    text = re.sub(r"\033][^\a]*\a", "", text)  # OSC sequences
    text = re.sub(r"\r", "", text)  # Carriage returns
    text = text.replace("\t", "    ")  # Expand tabs
    return text


def _build_shell_command(command: str) -> tuple[list | str, bool]:
    """Build the shell invocation based on execute_command config.

    Returns (cmd, shell_flag) for subprocess.Popen.
    If login_shell is True, wraps the command in a login shell invocation
    so that ~/.zprofile, ~/.zshrc etc. are sourced (giving full PATH).
    """
    _exec_cfg = get_tool_config().get("execute_command", {})
    use_login_shell = _exec_cfg.get("login_shell", True)
    if use_login_shell:
        shell_path = _exec_cfg.get("shell_path", "") or os.environ.get("SHELL", "/bin/zsh")
        return [shell_path, "-l", "-c", command], False
    return command, True


def _streaming_execute_command(command: str, timeout: int, cwd: str | None,
                               event_callback, tool_use_id: str) -> str:
    """Execute command with streaming output via event_callback."""
    env = os.environ.copy()
    env["TERM"] = "dumb"
    env["NO_COLOR"] = "1"
    env["PAGER"] = "cat"
    env["COLUMNS"] = "200"
    env["LINES"] = "50"

    shell_cmd, shell_flag = _build_shell_command(command)
    proc = subprocess.Popen(
        shell_cmd, shell=shell_flag, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        stdin=subprocess.DEVNULL, cwd=cwd, env=env,
        start_new_session=True,
    )
    output_lines = []
    import io
    deadline = time.time() + timeout
    try:
        # Read stdout line by line, emitting events
        for raw_line in iter(proc.stdout.readline, b''):
            if time.time() > deadline:
                import signal as sig
                try:
                    os.killpg(proc.pid, sig.SIGKILL)
                except OSError:
                    proc.kill()
                proc.wait(timeout=5)
                line_text = _strip_ansi(raw_line.decode("utf-8", errors="replace"))
                output_lines.append(line_text)
                output = "".join(output_lines)
                if len(output) > 50000:
                    output = output[:50000] + "\n... (truncated)"
                return _err(f"execute_command: timed out after {timeout}s\n{output}")
            line_text = _strip_ansi(raw_line.decode("utf-8", errors="replace"))
            output_lines.append(line_text)
            if event_callback:
                event_callback("tool_output", {
                    "tool_use_id": tool_use_id,
                    "line": line_text.rstrip("\n"),
                })
        proc.stdout.close()
        proc.wait(timeout=5)
    except subprocess.TimeoutExpired:
        import signal as sig
        try:
            os.killpg(proc.pid, sig.SIGKILL)
        except OSError:
            proc.kill()
        proc.wait(timeout=5)

    output = "".join(output_lines)
    stderr_data = proc.stderr.read() if proc.stderr else b""
    proc.stderr.close()
    if stderr_data:
        err_text = _strip_ansi(stderr_data.decode("utf-8", errors="replace"))
        output += ("\n--- stderr ---\n" + err_text) if output else err_text
    if len(output) > 50000:
        output = output[:50000] + "\n... (truncated)"
    return _ok({"command": command, "exit_code": proc.returncode, "output": output})


def tool_execute_command(args: dict) -> str:
    node_result = _route_to_node("execute_command", args)
    if node_result is not None:
        return node_result
    command = args.get("command", "")
    cwd = args.get("cwd")
    # Read default timeout from tools_config (per-call timeout still overrides)
    _exec_cfg = get_tool_config().get("execute_command", {})
    default_timeout = _exec_cfg.get("timeout", 15)
    timeout = args.get("timeout", default_timeout)
    # Check banned commands
    banned = _exec_cfg.get("banned_commands", [])
    for b in banned:
        if b and b in command:
            return _err(f"execute_command: command contains banned pattern '{b}'")
    try:
        if cwd:
            cwd = os.path.expanduser(cwd)

        # Use streaming version if event_callback is available
        ecb = getattr(_thread_local, 'event_callback', None)
        tuid = getattr(_thread_local, 'tool_use_id', None)
        if ecb and tuid:
            return _streaming_execute_command(command, timeout, cwd, ecb, tuid)

        # Force non-interactive environment
        env = os.environ.copy()
        env["TERM"] = "dumb"
        env["NO_COLOR"] = "1"
        env["PAGER"] = "cat"
        env["COLUMNS"] = "200"
        env["LINES"] = "50"

        shell_cmd, shell_flag = _build_shell_command(command)
        proc = subprocess.Popen(
            shell_cmd, shell=shell_flag, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            stdin=subprocess.DEVNULL, cwd=cwd, env=env,
            start_new_session=True,  # own process group so we can kill the tree
        )
        try:
            stdout, stderr = proc.communicate(timeout=timeout)
        except subprocess.TimeoutExpired:
            # Kill the entire process group
            import signal as sig
            try:
                os.killpg(proc.pid, sig.SIGKILL)
            except OSError:
                proc.kill()
            stdout, stderr = proc.communicate(timeout=5)
            output = _strip_ansi(stdout.decode("utf-8", errors="replace"))
            if stderr:
                output += "\n--- stderr ---\n" + _strip_ansi(stderr.decode("utf-8", errors="replace"))
            if len(output) > 50000:
                output = output[:50000] + "\n... (truncated)"
            return _err(f"execute_command: timed out after {timeout}s (partial output below). Use non-interactive commands, e.g. 'top -l 1' not 'top'.\n{output}")

        output = _strip_ansi(stdout.decode("utf-8", errors="replace"))
        if stderr:
            err_text = _strip_ansi(stderr.decode("utf-8", errors="replace"))
            output += ("\n--- stderr ---\n" + err_text) if output else err_text
        if len(output) > 50000:
            output = output[:50000] + "\n... (truncated)"
        return _ok({"command": command, "exit_code": proc.returncode, "output": output})
    except Exception as e:
        return _err(f"execute_command: {e}")


def tool_python_exec(args: dict) -> str:
    """Execute Python code in an isolated subprocess with artifact folder as cwd."""
    code = args.get("code", "")
    if not code.strip():
        return _err("python_exec: no code provided")

    _cfg = get_tool_config().get("python_exec", {})
    timeout = args.get("timeout", _cfg.get("timeout", 30))
    max_output = _cfg.get("max_output_chars", 50000)
    venv_path = _cfg.get("venv_path", "")

    # Working dir = session artifact folder so files written by code become artifacts
    session_id = getattr(_thread_local, 'current_session_id', None)
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    if session_id and agent:
        folder = _get_artifact_session_folder(session_id)
        work_dir = os.path.join(AGENTS_DIR, agent.agent_id, "artifacts", folder)
    else:
        import tempfile
        work_dir = os.path.join(tempfile.gettempdir(), "brain-pyexec")
    os.makedirs(work_dir, exist_ok=True)

    # Snapshot existing files before execution
    pre_files = set(os.listdir(work_dir))

    # Save script as a numbered artifact (persisted for reuse)
    counter = 1
    while os.path.exists(os.path.join(work_dir, f"script_{counter}.py")):
        counter += 1
    script_name = f"script_{counter}.py"
    script_path = os.path.join(work_dir, script_name)
    with open(script_path, "w") as f:
        f.write(code)
    agent_id = (agent.agent_id if agent else "main")
    _after_file_write(script_path, "created", agent_id)

    env = os.environ.copy()
    env["PYTHONDONTWRITEBYTECODE"] = "1"
    env["PYTHONUNBUFFERED"] = "1"
    if venv_path and os.path.isdir(venv_path):
        env["PYTHONPATH"] = venv_path + ((":" + env.get("PYTHONPATH", "")) if env.get("PYTHONPATH") else "")

    try:
        proc = subprocess.Popen(
            [sys.executable, script_path],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            stdin=subprocess.DEVNULL, cwd=work_dir, env=env,
            start_new_session=True,
        )
        try:
            stdout, stderr = proc.communicate(timeout=timeout)
        except subprocess.TimeoutExpired:
            import signal as sig
            try:
                os.killpg(proc.pid, sig.SIGKILL)
            except OSError:
                proc.kill()
            stdout, stderr = proc.communicate(timeout=5)
            output = stdout.decode("utf-8", errors="replace")
            if stderr:
                output += "\n--- stderr ---\n" + stderr.decode("utf-8", errors="replace")
            if len(output) > max_output:
                output = output[:max_output] + "\n... (truncated)"
            return _err(f"python_exec: timed out after {timeout}s\n{output}")

        output = stdout.decode("utf-8", errors="replace")
        if stderr:
            err_text = stderr.decode("utf-8", errors="replace")
            # Clean tracebacks: replace script path, remove echoed source lines
            err_text = err_text.replace(script_path, "<python_exec>")
            err_lines = err_text.splitlines()
            cleaned = []
            skip_next = False
            for line in err_lines:
                if '<python_exec>' in line:
                    cleaned.append(line)
                    skip_next = True
                elif skip_next and line.startswith('    ') and not line.strip().startswith('File '):
                    skip_next = False
                    continue
                else:
                    skip_next = False
                    cleaned.append(line)
            err_text = "\n".join(cleaned)
            output += ("\n--- stderr ---\n" + err_text) if output else err_text
        if len(output) > max_output:
            output = output[:max_output] + "\n... (truncated)"

        result = {"exit_code": proc.returncode, "output": output, "script": script_name}

        # Register any new files created by the script as artifacts
        post_files = set(os.listdir(work_dir))
        new_files = sorted(post_files - pre_files - {script_name})
        if new_files and agent:
            created = []
            for fname in new_files:
                fpath = os.path.join(work_dir, fname)
                if os.path.isfile(fpath):
                    _after_file_write(fpath, "created", agent_id)
                    created.append(fname)
            if created:
                result["artifacts"] = created

        # Always save stdout as an artifact when the script didn't write any files
        if output and proc.returncode == 0 and not new_files and agent:
            try:
                artifact_path = os.path.join(work_dir, "output.txt")
                counter = 1
                while os.path.exists(artifact_path):
                    artifact_path = os.path.join(work_dir, f"output_{counter}.txt")
                    counter += 1
                with open(artifact_path, "w") as af:
                    af.write(output)
                _after_file_write(artifact_path, "created", agent_id)
                result["artifacts"] = [os.path.basename(artifact_path)]
                # For large outputs, replace inline data with a reference so the
                # summariser doesn't ingest a megabyte of stdout. Small outputs
                # stay inline so the summariser can describe them meaningfully.
                if len(output) > 1000:
                    lines = output.splitlines()
                    result["output"] = (
                        f"Output saved as artifact {os.path.basename(artifact_path)} "
                        f"({len(lines)} lines, {len(output):,} chars). "
                        f"The user can view it directly. Summarize what was computed, do NOT repeat the data."
                    )
            except Exception:
                pass

        return _ok(result)
    except Exception as e:
        return _err(f"python_exec: {e}")
