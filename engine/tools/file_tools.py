# File / shell / python / document tool bodies (extracted from brain.py, E1).
#
# The 10 file-system tools (read/write/edit_file, list_directory, search_files,
# execute_command, python_exec, read/write/edit_document) plus the helpers
# private to this cluster. This is a PURE relocation — JSON envelopes, error
# strings, security gates, and side-effects (artifact registration, GDPR
# anonymisation) are byte-identical to the pre-E1 brain.py.
#
# Seams:
#   - `_ok` / `_err` / `_get_artifact_session_folder` / `_record_session_read_path`
#     come from engine.tool_exec (the C2 extraction; re-exported on brain too).
#   - `get_request_context` comes from engine.context (low-level base, no cycle).
#   - brain-runtime symbols (`_after_file_write`, `_gdpr_anon_tool_text`,
#     `_route_to_node`, `get_tool_config`, `AGENTS_DIR`, `_current_agent`,
#     `DocumentParser`) are reached lazily via `import brain as _brain` inside
#     function bodies — a top-level `import brain` would be a cycle (brain.py
#     imports this module for TOOL_DISPATCH).
#   - `engine.doc_convert._do_extract` (read_document's binary path) imports
#     directly from engine/.
#
# brain.py re-exports all 10 tool functions via
# `from engine.tools.file_tools import (...)` so `brain.tool_read_file` etc.,
# the TOOL_DISPATCH entries, and any in-brain bare-name calls resolve unchanged.

from __future__ import annotations

import fnmatch
import glob as globmod
import json
import os
import re
import subprocess
import sys
import time

from engine.context import get_request_context
from engine.tool_exec import (
    _ok,
    _err,
    _get_artifact_session_folder,
    _record_session_read_path,
    register_tool_process,
    unregister_tool_process,
)


def _cwd_base() -> str | None:
    """The base directory a RELATIVE path should resolve against. In Code Mode
    the request context carries `working_dir` → that's the base (so list/read/
    search/glob operate inside the project's working directory, matching where
    the file-writing tools land). Otherwise None → callers keep their existing
    `os.path.abspath` (process cwd) behavior unchanged."""
    try:
        _wd = get_request_context().working_dir
    except Exception:
        _wd = None
    return _wd if (_wd and os.path.isdir(_wd)) else None


def _resolve_under_cwd(path: str) -> str:
    """Expand `path` and, if relative, resolve it under the code-mode working
    directory when one is active; else fall back to process-cwd abspath
    (unchanged behavior for non-code-mode)."""
    path = os.path.expanduser(path)
    if os.path.isabs(path):
        return path
    base = _cwd_base()
    return os.path.join(base, path) if base else os.path.abspath(path)


def _read_matched_regions(md_path: str, matched: set, radius: int = 2):
    """Return the union of chunk windows around each matched chunk_index of
    md_path, joined in order with a marker at skipped gaps. None on any failure
    (caller then falls back to a full read — never silently truncates).
    """
    import brain as _brain
    try:
        cfg = _brain._load_mempalace_config()
        palace = cfg.get("palace_path", "")
        if not palace or not os.path.isdir(palace):
            return None
        ok, _ = _brain._ensure_mempalace_importable()
        if not ok:
            return None
        from mempalace.palace import get_collection
        col = get_collection(palace, create=False)
        if col is None:
            return None
        got = col.get(where={"source_file": md_path},
                      include=["documents", "metadatas"])
        docs = got.get("documents") or []
        metas = got.get("metadatas") or []
        rows = []
        for d, m in zip(docs, metas):
            if not isinstance(d, str):
                continue
            try:
                ci = int((m or {}).get("chunk_index"))
            except (TypeError, ValueError):
                continue
            rows.append((ci, d))
        if not rows:
            return None
        rows.sort(key=lambda t: t[0])
        # Small-file shortcut: if the file is small, trimming buys nothing —
        # the gap-markers + lost context cost more than the few KB saved, and
        # the model is better off seeing the whole thing. Return None (full read)
        # when the doc has few chunks OR is small in total characters.
        _total_chars = sum(len(d) for _ci, d in rows)
        if len(rows) <= 8 or _total_chars <= 6000:
            return None
        # Union of [ci-radius .. ci+radius] over all matched chunks.
        keep = set()
        for ci in matched:
            for j in range(ci - radius, ci + radius + 1):
                keep.add(j)
        out, prev = [], None
        for ci, d in rows:
            if ci not in keep:
                continue
            if prev is not None and ci > prev + 1:
                out.append(f"\n[... {ci - prev - 1} chunk(s) omitted — not matched ...]\n")
            out.append(d)
            prev = ci
        if not out:
            return None
        stitched = "\n\n".join(out)
        # Worth-it gate: many matched chunks (or wide radius) make the union of
        # regions add up to ~the whole file — once you count overlap + the gap
        # markers, trimming saves little or nothing. Only trim when the result
        # is meaningfully smaller than the full file; otherwise return None so
        # the caller reads the whole thing (full context, no fragmentation).
        _full_chars = sum(len(d) for _ci, d in rows)
        if len(stitched) >= 0.75 * _full_chars:
            return None
        return stitched
    except Exception:
        return None


def tool_read_file(args: dict) -> str:
    import brain as _brain
    node_result = _brain._route_to_node("read_file", args)
    if node_result is not None:
        return node_result
    path = args.get("path", "")
    offset = args.get("offset", 1)
    limit = args.get("limit", 400)
    try:
        path = _resolve_under_cwd(path)
        if not os.path.exists(path):
            return _err(
                f"File not found: {path}. "
                "Do NOT retry with the same path — the file does not exist on disk. "
                "Hallucinated paths are a frequent cause: the path was inferred "
                "from the project's folder structure or filename schema rather "
                "than copied from a real source. "
                "USE ONLY paths returned in `read_path` (or `read_path_original` "
                "as fallback) of a prior `mempalace_query` drawer. "
                "If no drawer pointed at this content, the project does not "
                "contain it — refuse per REFUSAL DISCIPLINE instead of guessing "
                "another path."
            )
        with open(path, "r", errors="replace") as f:
            lines = f.readlines()
        total = len(lines)
        start = max(0, offset - 1)
        end = start + limit if limit else total
        selected = lines[start:end]
        # Number lines
        numbered = []
        for i, line in enumerate(selected, start=start + 1):
            numbered.append(f"{i:>6}\t{line.rstrip()}")
        content = "\n".join(numbered)
        _record_session_read_path(path)
        content_anon = _brain._gdpr_anon_tool_text(content, f"file:{os.path.basename(path)}")
        return _ok({"path": path, "total_lines": total, "showing": f"{start+1}-{min(end, total)}", "content": content_anon})
    except Exception as e:
        return _err(f"read_file: {e}")


def tool_write_file(args: dict) -> str:
    import brain as _brain
    node_result = _brain._route_to_node("write_file", args)
    if node_result is not None:
        return node_result
    path = args.get("path", "")
    content = args.get("content", "")
    try:
        # Hard guard: the resolved write path MUST be inside the session artifact
        # folder (relative names default into it; absolute / .. escapes refused).
        path, err = _enforce_artifact_path(path, "write_file")
        if err:
            return err
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        with open(path, "w") as f:
            f.write(content)
        size = os.path.getsize(path)
        agent = get_request_context().current_agent or _brain._current_agent
        _brain._after_file_write(path, "created", agent.agent_id if agent else "main")
        return _ok({"path": path, "size": size, "status": "written"})
    except Exception as e:
        return _err(f"write_file: {e}")


def tool_edit_file(args: dict) -> str:
    import brain as _brain
    path = args.get("path", "")
    old_string = args.get("old_string", "")
    new_string = args.get("new_string", "")
    replace_all = args.get("replace_all", False)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        with open(path, "r") as f:
            content = f.read()
        count = content.count(old_string)
        if count == 0:
            return _err(f"edit_file: old_string not found in {path}")
        if count > 1 and not replace_all:
            return _err(f"edit_file: old_string found {count} times — use replace_all=true or provide a more specific match")
        if replace_all:
            new_content = content.replace(old_string, new_string)
        else:
            new_content = content.replace(old_string, new_string, 1)
        with open(path, "w") as f:
            f.write(new_content)
        agent = get_request_context().current_agent or _brain._current_agent
        _brain._after_file_write(path, "modified", agent.agent_id if agent else "main")
        return _ok({"path": path, "replacements": count if replace_all else 1, "status": "edited"})
    except Exception as e:
        return _err(f"edit_file: {e}")


# Max chars a SINGLE whole-document read returns inline before it's spilled to
# disk + previewed. ~120k chars ≈ 30k tokens — generous for one doc, small
# enough that several reads + the running conversation still fit a 256k window.
# (A 3.3 MB PDF extracts to ~768k chars / ~190k tokens — one such read alone
# overflows the window, which is the bug this guards.) Overridable via
# config.json conversion.read_document_budget_chars.
_DOC_READ_BUDGET_CHARS_DEFAULT = 120_000
_DOC_READ_PREVIEW_CHARS = 8_000


def _doc_read_budget_chars() -> int:
    try:
        import brain as _b
        v = int(((_b._server_config() or {}).get("conversion") or {})
                .get("read_document_budget_chars", 0) or 0)
        return v if v > 0 else _DOC_READ_BUDGET_CHARS_DEFAULT
    except Exception:
        return _DOC_READ_BUDGET_CHARS_DEFAULT


_DOC_READ_BUDGET_CHARS = _DOC_READ_BUDGET_CHARS_DEFAULT  # re-resolved per call below


def _spill_oversized_read(path: str, payload: dict, content: str) -> str:
    """A whole-doc read exceeded the single-read budget: write the full content
    to the session's tool-results dir and return a preview + instructions to
    read it in ranges. Keeps the file fully available (re-read / grep / citation
    validator) without forcing ~190k tokens into one prompt."""
    import brain as _brain
    fmt = payload.get("format", "")
    is_pdf = fmt == "pdf" or path.lower().endswith(".pdf")
    size_kb = len(content) // 1024
    sid = get_request_context().current_session_id or ""
    agent = get_request_context().current_agent or getattr(_brain, "_current_agent", None)
    agent_id = agent.agent_id if agent else "main"
    spill_path = ""
    try:
        results_dir = os.path.join(
            _brain.AGENTS_DIR, agent_id, "artifacts",
            _get_artifact_session_folder(sid), "tool-results")
        os.makedirs(results_dir, exist_ok=True)
        base = os.path.basename(path)
        spill_path = os.path.join(results_dir, base + ".fulltext.md")
        with open(spill_path, "w", encoding="utf-8") as fh:
            fh.write(content)
    except OSError:
        spill_path = ""
    preview = content[:_DOC_READ_PREVIEW_CHARS]
    if is_pdf:
        how = ("Read it in parts: call read_document again with `pages` "
               "(e.g. pages=\"1-10\", then \"11-20\", …) to page through it")
    else:
        how = ("Read it in parts: call read_document again with `offset`+`limit` "
               "(line numbers) to page through it")
    spill_line = (f"Full text saved to: {spill_path} — "
                  f"read_document/read_file that path with offset+limit, or grep "
                  f"it via execute_command.\n" if spill_path else "")
    note = (
        f"[Document too large to return in one read ({size_kb}KB, "
        f"~{len(content)//4} tokens) — it would overflow the model context. "
        f"{how}; do NOT re-read the same range.]\n"
        f"{spill_line}"
        f"Preview (first {_DOC_READ_PREVIEW_CHARS} chars):\n{preview}\n…")
    out = {k: v for k, v in payload.items() if k != "content"}
    out["content"] = note
    out["truncated"] = True
    out["full_chars"] = len(content)
    if spill_path:
        out["full_text_path"] = spill_path
    return _ok(out)


def tool_read_document(args: dict) -> str:
    """Format-aware document reader."""
    global _DOC_READ_BUDGET_CHARS
    _DOC_READ_BUDGET_CHARS = _doc_read_budget_chars()
    import brain as _brain
    # Accept legacy / drawer-vocab synonyms — drawers from mempalace_query
    # carry the field as `source_file`, and the model frequently passes that
    # name verbatim. Treating an empty `path` as CWD silently expanded into
    # "Is a directory" errors that the model then ignored.
    path = args.get("path", "") or args.get("source_file", "") or args.get("file", "")
    if not path or not str(path).strip():
        return _err(
            "read_document: 'path' is required (got empty). When following up "
            "on a mempalace_query drawer, take its `source_file` value and "
            "pass it as the `path` argument — JOIN with the input-folder "
            "absolute path if `source_file` is relative.")
    path = str(path).strip()
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        if os.path.isdir(path):
            return _err(
                f"read_document: '{path}' is a directory, not a file. The "
                "model commonly passes a base path here when it meant to "
                "pass a specific document — re-issue the call with the full "
                "file path including extension.")
        if not os.path.exists(path):
            return _err(
                f"File not found: {path}. "
                "Do NOT retry with the same path — the file does not exist on disk. "
                "Hallucinated paths are a frequent cause: the path was inferred "
                "from the project's folder structure or filename schema rather "
                "than copied from a real source. "
                "USE ONLY paths returned in `read_path` (or `read_path_original` "
                "as fallback) of a prior `mempalace_query` drawer. "
                "If no drawer pointed at this content, the project does not "
                "contain it — refuse per REFUSAL DISCIPLINE instead of guessing "
                "another path."
            )

        def _ok_and_cache(payload: dict) -> str:
            # Record path so the citation validator can grep this file for
            # verbatim quotes. Anonymisation applies on return so the model
            # sees pseudonymised text in cross-session reads.
            _record_session_read_path(path)
            _src = f"attachment:{os.path.basename(path)}"
            raw = str(payload.get("content", "") or "")
            anon = _brain._gdpr_anon_tool_text(raw, _src)
            if anon is not raw:
                payload = dict(payload)
                payload["content"] = anon
            # Single-read size guard. A whole-document read of a large file
            # (e.g. a 3.3 MB PDF → ~190k tokens) returned VERBATIM overflows the
            # model's context window in ONE call — the turn dies at round 0 with
            # a provider "prompt too large" 400, before the model can do
            # anything. The message-level budget (_apply_tool_result_budget) is
            # a no-op on the interactive path (tool results are ephemeral in the
            # sidecar, never in session.messages), so the clamp has to happen
            # HERE at the tool-result point the sidecar actually sends. When the
            # content exceeds the budget we spill the FULL (already-anonymised)
            # text to the session's tool-results dir and return a preview plus a
            # format-appropriate instruction to read it in ranges — the model
            # can then narrow with pages= (PDF) or offset/limit (text), or grep
            # the spilled file. Nothing is lost; it's just not all forced into
            # one prompt. An explicit pages/offset/limit selection is left
            # alone (the model is already narrowing).
            content = str(payload.get("content", "") or "")
            narrowed = bool(args.get("pages") or args.get("offset")
                            or args.get("limit") or args.get("sheet")
                            or args.get("slides"))
            if (not narrowed) and len(content) > _DOC_READ_BUDGET_CHARS:
                return _spill_oversized_read(path, payload, content)
            return _ok(payload)

        ext = os.path.splitext(path)[1].lower()

        # Office/PDF formats route through the unified doc_convert pipeline
        # (markitdown-first → per-format fallback), shared with project mining.
        # read_document passes caps=False (full fidelity) + the selection/meta
        # knobs each format supports.
        # All doc_convert-supported binary/tabular formats route through the
        # one unified pipeline (markitdown-first → per-format fallback), shared
        # with project mining + classification scan. read_document passes
        # caps=False (full fidelity) + the selection/meta knobs each format
        # supports. .xls aliases to the xlsx extractor; .msg/.epub/.zip are
        # markitdown-first (no inline MarkItDown() call here anymore).
        from engine.doc_convert import SUPPORTED_EXTS as _SUPPORTED_EXTS
        if ext in _SUPPORTED_EXTS:
            from engine.doc_convert import _do_extract
            kwargs: dict = {"caps": False}
            if ext in (".xlsx", ".xlsm", ".xls", ".xlsb"):
                kwargs["sheet"] = args.get("sheet")
            elif ext == ".pptx":
                kwargs["slides"] = args.get("slides")
            elif ext == ".pdf":
                kwargs["pages"] = args.get("pages", "") or None
                kwargs["include_tables"] = bool(args.get("include_tables"))
                kwargs["emit_meta"] = True
                kwargs["page_marker"] = "--- Page"
            text, backend, err = _do_extract(path, **kwargs)
            if err:
                return _err(f"read_document: {err}")
            # Normalise the reported format for the aliases.
            fmt = {".xls": "xlsx", ".xlsm": "xlsx", ".xlsb": "xlsx",
                   ".tsv": "csv"}.get(ext, ext.lstrip("."))
            # `backend` ("markitdown" / "fitz/legacy" / "mistral-ocr (Np)" / …)
            # records which of the two extraction surfaces produced this text.
            # Surfaced in the chat tool-block + persisted via metadata.tools.
            return _ok_and_cache({"path": path, "format": fmt,
                                  "backend": backend, "content": text})

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"):
            meta_text = _brain.DocumentParser.parse_image(path)
            vision_note = "\n\n*(For AI-powered image description, include this image directly in your chat message)*"
            _content = _brain._gdpr_anon_tool_text(
                meta_text, f"attachment:{os.path.basename(path)}") + vision_note
            return _ok({"path": path, "format": "image", "content": _content})

        elif ext == ".svg":
            content = _brain.DocumentParser.parse_svg(path)
            return _ok_and_cache({"path": path, "format": "svg", "content": content})

        else:
            # Matched-regions auto-read: when this exact .md file came from a
            # mempalace_query this session, the relevant content is the few
            # chunks that matched — typically SCATTERED across the doc (a
            # Löschkonzept matched chunks 2/18/20/48). Return the union of small
            # windows around those matched chunks instead of the whole file —
            # the model gets every relevant region at a fraction of the bytes,
            # automatically (no new tool / no model decision). Falls back to a
            # full read when (a) the model paginated explicitly (offset/limit),
            # or (b) this file wasn't a query hit (regions unknown) — so an
            # ad-hoc read or a "give me everything" never silently truncates.
            _regions = (set() if (args.get("offset") or args.get("limit"))
                        else _brain._get_match_regions(path))
            if _regions:
                _rr = _read_matched_regions(path, _regions, radius=2)
                if _rr is not None:
                    return _ok_and_cache({"path": path, "format": "text-regions",
                                          "matched_chunks": sorted(_regions),
                                          "content": _rr})
            # Plain-text / markdown / unknown-extension read. Honor explicit
            # offset+limit when the model paginates; otherwise read the whole
            # file. The previous hard-cap at 500 lines truncated mid-document
            # silently — on the kg-real-policies ISMS Handbuch (1903 lines)
            # the section 2.13 percent-to-score table sits at line 1267,
            # WAY past line 500, so the model never saw the table even when
            # it correctly chose read_document on the .md companion. Oversized
            # results are spilled to disk + previewed by _apply_tool_result_budget
            # on subsequent turns if they exceed the round budget.
            offset = int(args.get("offset", 1) or 1)
            limit = args.get("limit")
            with open(path, "r", errors="replace") as f:
                lines = f.readlines()
            total = len(lines)
            start = max(0, offset - 1)
            end = start + int(limit) if limit else total
            selected = lines[start:end]
            numbered = []
            for i, line in enumerate(selected, start=start + 1):
                numbered.append(f"{i:>6}\t{line.rstrip()}")
            content = "\n".join(numbered)
            shown = f"{start+1}-{min(end, total)}"
            result = {"path": path, "format": "text",
                      "total_lines": total, "showing": shown,
                      "content": content}
            # Fail-LOUD on a misused `pages` arg: `pages` only applies to real
            # PDFs (page-indexed). On a .md/.txt it was silently ignored, so a
            # model trying to narrow to "page 41" got the whole file back every
            # time and kept retrying with different page ranges (the fc3fa95b
            # 16-round / 500k-token loop). Tell it plainly so it stops guessing.
            if args.get("pages"):
                result["note"] = (f"`pages` has no effect on this {ext or 'text'} "
                                  f"file — it has lines, not pages ({total} lines total). "
                                  f"To read a slice, use offset + limit (line numbers).")
            return _ok_and_cache(result)
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"read_document: {e}")


_MD_IMAGE_RE = re.compile(r'^\s*!\[([^\]]*)\]\(([^)\s]+)\)\s*$')

# Built-in fallback style — used when no preset is named / found. Mirrors the
# corporate.yaml keys so the apply-code can assume the full shape.
_DEFAULT_DOC_STYLE = {
    "fonts": {"body": "Calibri", "heading": "Calibri", "mono": "Consolas"},
    "sizes": {"body": 11, "h1": 20, "h2": 16, "h3": 13},
    "colors": {"heading": "#1F3864", "body": "#222222", "accent": "#2E74B5",
               "table_header_bg": "#1F3864", "table_header_text": "#FFFFFF"},
    "docx": {"table_style": "Light Grid Accent 1", "heading_bold": True,
             # Opus-near polish, all deterministic (CLAUDE.md rule 5). The MODEL
             # writes plain markdown; code applies these. zebra_fill = alternating
             # body-row shading; rule_color = ---/heading-underline colour;
             # strip_emoji keeps regulatory docs clean; risk_badges colours a
             # Bewertung/Risiko/Rating column by cell value (gering/mittel/
             # erhöht/hoch); cover renders a title page from the FIRST # H1 +
             # frontmatter; toc inserts a Word TOC field after the cover.
             "zebra_fill": "#EDF1F8", "rule_color": "#B4C6E7",
             "strip_emoji": True, "risk_badges": True, "cover": True, "toc": True},
    "pdf": {"page_size": "letter", "margin_inch": 1.0},
    "pptx": {"title_color": "#1F3864", "body_color": "#222222",
             "accent": "#2E74B5", "background": "#FFFFFF"},
    "mermaid": {"theme": "default", "background": "white"},
    # Running header/footer + logo, applied to paginated formats (docx/pdf) and,
    # for the logo + footer text, to pptx slides. Empty text/file = not rendered.
    # `text` supports {page} and {date} tokens (page-number / current date).
    "header": {"text": "", "align": "left", "font_size": 9, "color": "#666666"},
    "footer": {"text": "", "align": "center", "font_size": 9,
               "color": "#666666", "page_numbers": False,
               # `classification: true` adds an automatic, content-derived
               # sensitivity line ("Klassifizierung: Vertraulich") on its own
               # footer line. Heuristic-only (engine/classification.py), no LLM.
               "classification": True},
    "logo": {"file": "", "width_inch": 1.2, "align": "right", "position": "header"},
}


def _doc_style_dir(agent_id: str) -> str:
    import brain as _brain
    return os.path.join(_brain.AGENTS_DIR, agent_id, "skills", "doc-styles")


def _preset_exists(name: str, agent_id: str) -> bool:
    if not name:
        return False
    base = _doc_style_dir(agent_id)
    return any(os.path.isfile(os.path.join(base, f"{name}.{e}")) for e in ("yaml", "yml"))


def _resolve_default_style(explicit: str) -> str:
    """Pick the style preset NAME to apply to a document. Resolution order:
      1. explicit `style` arg from the tool call (the model named one) — wins.
      2. the current project's `doc_style` (project.json), if in a project.
      3. config.json → doc_styles.default (global default preset).
      4. 'corporate' if that preset file exists (sensible built-in default).
      5. '' → built-in look (no preset).
    This is why a report is styled even when the model OMITS style= (the common
    case): the tool applies a default rather than the bare built-in look."""
    import brain as _brain
    explicit = (explicit or "").strip()
    if explicit:
        return explicit
    agent = get_request_context().current_agent or _brain._current_agent
    agent_id = agent.agent_id if agent else "main"
    # 2. project preset
    try:
        proj_name = get_request_context().project
        if proj_name:
            proj = _brain.ProjectManager.get_project(agent_id, proj_name)
            cand = ((proj or {}).get("doc_style") or "").strip()
            if cand and _preset_exists(cand, agent_id):
                return cand
    except Exception:
        pass
    # 3. global config default
    try:
        cfg_default = ((_brain._server_config().get("doc_styles") or {}).get("default") or "").strip()
        if cfg_default and _preset_exists(cfg_default, agent_id):
            return cfg_default
    except Exception:
        pass
    # 4. corporate, if present
    if _preset_exists("corporate", agent_id):
        return "corporate"
    # 5. built-in
    return ""


def _load_doc_style(name: str):
    """Load a document style preset (agents/<agent>/skills/doc-styles/<name>.yaml),
    deep-merged over the built-in defaults so callers always get the full shape.
    name '' / None / not-found → built-in defaults. Deterministic — the MODEL
    never interprets this; the doc tools apply it in code (CLAUDE.md rule 5)."""
    import copy
    import brain as _brain
    style = copy.deepcopy(_DEFAULT_DOC_STYLE)
    name = (name or "").strip()
    if not name:
        return style
    try:
        import yaml
        agent = get_request_context().current_agent or _brain._current_agent
        agent_id = agent.agent_id if agent else "main"
        base = os.path.join(_brain.AGENTS_DIR, agent_id, "skills", "doc-styles")
        for fn in (f"{name}.yaml", f"{name}.yml"):
            p = os.path.join(base, fn)
            if os.path.isfile(p):
                with open(p, encoding="utf-8") as f:
                    loaded = yaml.safe_load(f) or {}
                for k, v in loaded.items():
                    if isinstance(v, dict) and isinstance(style.get(k), dict):
                        style[k].update(v)
                    else:
                        style[k] = v
                break
    except Exception:
        pass
    return style


def _resolve_reference_docx(filename: str) -> str | None:
    """Resolve a reference .docx to lift styling FROM. Looks in the current
    project's instruction-files/ dir (the owner-uploaded template/reference).
    `filename` empty → auto-pick the project's sole/first instruction-file .docx.
    Returns an existing absolute path or None (no project / no match / not .docx)."""
    import brain as _brain
    try:
        proj_name = get_request_context().project
        if not proj_name:
            return None
        agent = get_request_context().current_agent or _brain._current_agent
        agent_id = agent.agent_id if agent else "main"
        idir = _brain.ProjectManager._instruction_files_dir(agent_id, proj_name)
        if not os.path.isdir(idir):
            return None
        filename = (filename or "").strip()
        if filename:
            cand = os.path.join(idir, os.path.basename(filename))
            return cand if (os.path.isfile(cand) and cand.lower().endswith(".docx")) else None
        # auto-pick: first .docx among the project's instruction files
        for fn in sorted(os.listdir(idir)):
            if fn.lower().endswith(".docx"):
                return os.path.join(idir, fn)
        return None
    except Exception:
        return None


def _load_doc_style_from_reference(ref_path: str):
    """Build a doc-style dict by reading the named-style DEFINITIONS out of a
    reference .docx (the look the user wants to match) instead of a brand preset.
    Lifts: body font/size/color (Normal style), heading font/size/color/bold
    (Heading 1-3), and page margins. Deep-merged over the built-in defaults so
    the returned shape always matches _DEFAULT_DOC_STYLE — the apply-loop in
    tool_write_document consumes it unchanged.

    Deterministic (CLAUDE.md rule 5): the model never interprets fonts/colors;
    code reads them. Returns (style_dict, note) — `note` is a short human string
    for the tool result (what matched / why it fell back). On any failure returns
    (built-in default style, reason) so the caller can fall back loudly, never
    silently producing a clobbered doc.

    Scope limit: python-docx exposes named style + section-margin definitions,
    NOT the full visual template (themes, cover pages, complex section layouts).
    'Match' therefore means body/heading typography + margins, not a pixel clone."""
    import copy
    style = copy.deepcopy(_DEFAULT_DOC_STYLE)
    try:
        import docx
    except ImportError:
        return style, "python-docx not installed — applied built-in default"
    try:
        ref = docx.Document(ref_path)
    except Exception as e:
        return style, f"could not open reference ({e}) — applied built-in default"

    def _style_font(name):
        try:
            return ref.styles[name].font
        except Exception:
            return None

    def _doc_default_font_size():
        """The document's docDefaults (rPrDefault/rPr) — the font/size every style
        inherits when it sets none explicitly. python-docx's styles['Normal'].font
        does NOT surface these, so a doc whose body font lives only in docDefaults
        (the common Word case) would otherwise fall back to Calibri. Returns
        (font_name|None, size_pt|None)."""
        try:
            from docx.oxml.ns import qn
            dd = ref.styles.element.find(qn("w:docDefaults"))
            rpr = dd.find(qn("w:rPrDefault")) if dd is not None else None
            r = rpr.find(qn("w:rPr")) if rpr is not None else None
            if r is None:
                return None, None
            rf = r.find(qn("w:rFonts"))
            fname = rf.get(qn("w:ascii")) if rf is not None else None
            sz = r.find(qn("w:sz"))
            spt = None
            if sz is not None and sz.get(qn("w:val")):
                spt = round(int(sz.get(qn("w:val"))) / 2)  # half-points → pt
            return fname, spt
        except Exception:
            return None, None

    def _hexkey(rgb):
        # docx RGBColor → '#RRGGBB'; None when the style inherits (no explicit color)
        try:
            return "#" + str(rgb) if rgb is not None else None
        except Exception:
            return None

    matched = []
    _dd_font, _dd_size = _doc_default_font_size()
    nf = _style_font("Normal")
    _body_font = (nf.name if nf is not None else None) or _dd_font
    if _body_font:
        style["fonts"]["body"] = _body_font; matched.append("body font")
    _body_size = None
    if nf is not None and nf.size is not None:
        try:
            _body_size = round(nf.size.pt)
        except Exception:
            _body_size = None
    if _body_size is None:
        _body_size = _dd_size
    if _body_size:
        style["sizes"]["body"] = _body_size; matched.append("body size")
    if nf is not None:
        bc = _hexkey(getattr(nf.color, "rgb", None))
        if bc:
            style["colors"]["body"] = bc; matched.append("body color")
    # Headings: use Heading 1 as the canonical heading look; size each level from
    # its own style when present so the hierarchy is preserved. Heading font falls
    # back to the doc default (then body) when the Heading style names none.
    h1f = _style_font("Heading 1")
    _heading_font = (h1f.name if h1f is not None else None) or _dd_font
    if _heading_font:
        style["fonts"]["heading"] = _heading_font; matched.append("heading font")
    if h1f is not None:
        hc = _hexkey(getattr(h1f.color, "rgb", None))
        if hc:
            style["colors"]["heading"] = hc; matched.append("heading color")
        if h1f.bold is not None:
            style["docx"]["heading_bold"] = bool(h1f.bold)
    for _lvl, _szkey in ((1, "h1"), (2, "h2"), (3, "h3")):
        hf = _style_font(f"Heading {_lvl}")
        if hf is not None and hf.size is not None:
            try:
                style["sizes"][_szkey] = round(hf.size.pt)
            except Exception:
                pass
    # Page margins (first section) → pdf.margin_inch is the only margin knob in
    # the shape; docx margins are applied via header/footer section, so capture
    # the top/left as a representative inch value for any paginated reuse.
    try:
        sec = ref.sections[0]
        if sec.left_margin is not None:
            style["pdf"]["margin_inch"] = round(sec.left_margin.inches, 2)
            matched.append("margins")
    except Exception:
        pass

    note = ("matched " + ", ".join(matched)) if matched else \
        "reference had no explicit named-style overrides — applied built-in default"
    return style, note


def _hex_rgb(h: str):
    """'#RRGGBB' → (r,g,b) ints; tolerant of missing '#'. None on bad input."""
    try:
        h = (h or "").lstrip("#")
        if len(h) != 6:
            return None
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    except Exception:
        return None


def _resolve_doc_image(src: str, doc_dir: str) -> str | None:
    """Resolve an image path from a markdown ![alt](src) inside write_document.
    `src` is usually a RELATIVE filename produced by render_diagram in the same
    artifact folder, so try doc_dir first; then expanduser/abs. Returns an
    existing path or None. (Remote URLs are not fetched here — the model should
    render_diagram or attach the file.)"""
    if not src or src.startswith(("http://", "https://", "data:")):
        return None
    for cand in (os.path.join(doc_dir, src), os.path.expanduser(src), os.path.abspath(src)):
        if os.path.isfile(cand):
            return cand
    return None


def _resolve_style_logo(style: dict) -> str | None:
    """Resolve a doc-style preset's logo image to an existing path.
    `logo.file` is a bare filename stored alongside the preset YAML
    (agents/<agent>/skills/doc-styles/<file>), uploaded via the GUI editor.
    Returns the absolute path or None (missing / empty / not found)."""
    import brain as _brain
    fn = ((style.get("logo") or {}).get("file") or "").strip()
    if not fn:
        return None
    try:
        agent = get_request_context().current_agent or _brain._current_agent
        agent_id = agent.agent_id if agent else "main"
        p = os.path.join(_brain.AGENTS_DIR, agent_id, "skills", "doc-styles",
                         os.path.basename(fn))
        return p if os.path.isfile(p) else None
    except Exception:
        return None


def _hdrftr_text(raw: str, *, page: bool = False) -> str:
    """Substitute {date} (and strip {page} for renderers that inject the page
    number as a field/canvas draw separately). `page=True` leaves {page} for the
    caller to handle; otherwise removes it. Returns '' for falsy input."""
    import datetime
    s = str(raw or "")
    if not s:
        return ""
    s = s.replace("{date}", datetime.date.today().isoformat())
    if not page:
        s = s.replace("{page}", "")
    return s


def _apply_pptx_logo_footer(prs, style: dict):
    """Place the preset's logo (any non-'none' position → a corner image) and
    footer text band on every slide of a presentation. Slides have no native
    page header/footer, so this draws plain shapes; {page} → slide number."""
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as _PRGB
    from pptx.enum.text import PP_ALIGN
    logo = style.get("logo") or {}
    ftr = style.get("footer") or {}
    logo_path = _resolve_style_logo(style)
    show_logo = bool(logo_path) and (logo.get("position") or "header").lower() != "none"
    ftr_text = str(ftr.get("text") or "").strip()
    page_nums = bool(ftr.get("page_numbers"))
    if not (show_logo or ftr_text or page_nums):
        return
    sw, sh = prs.slide_width, prs.slide_height
    margin = Inches(0.3)
    f_align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
               "right": PP_ALIGN.RIGHT}.get((ftr.get("align") or "center").lower(),
                                            PP_ALIGN.CENTER)
    f_col = _hex_rgb(ftr.get("color"))
    f_size = float(ftr.get("font_size") or 9)
    for idx, slide in enumerate(prs.slides, start=1):
        if show_logo:
            try:
                w = Inches(float(logo.get("width_inch") or 1.2))
                a = (logo.get("align") or "right").lower()
                top = margin if (logo.get("position") or "header").lower() != "footer" \
                    else sh - margin - w
                x = margin if a == "left" else (sw - margin - w if a == "right"
                                                else int((sw - w) / 2))
                slide.shapes.add_picture(logo_path, x, top, width=w)
            except Exception:
                pass
        f_raw = ftr_text
        if page_nums and "{page}" not in f_raw:
            f_raw = (f_raw + "  Folie {page}").strip() if f_raw else "Folie {page}"
        if f_raw:
            try:
                txt = _hdrftr_text(f_raw, page=True).replace("{page}", str(idx))
                box = slide.shapes.add_textbox(margin, sh - Inches(0.45),
                                               sw - 2 * margin, Inches(0.35))
                tf = box.text_frame
                tf.word_wrap = False
                p = tf.paragraphs[0]
                p.alignment = f_align
                run = p.add_run(); run.text = txt
                run.font.size = Pt(f_size)
                if f_col:
                    run.font.color.rgb = _PRGB(*f_col)
            except Exception:
                pass


# EU AI Act transparency disclosure — always present on generated documents so a
# reader can tell the content was AI-assisted. Fixed text; preset may override via
# footer.ai_disclosure_text, or disable with footer.ai_disclosure: false.
_AI_DISCLOSURE_DEFAULT = (
    "Dieses Dokument wurde mit Unterstützung künstlicher Intelligenz erstellt."
)


def _classification_footer_line(content: str, style: dict) -> str:
    """Derive an automatic, CONTENT-based sensitivity classification for a
    generated document and return a footer line like `Klassifizierung: Vertraulich`
    — or "" when disabled.

    Uses the existing ARL content heuristic (engine/classification.py, regex +
    keyword + PII, NO LLM). The generated doc carries no explicit marker, so we
    take the content `heuristic_level`. A corporate report floors at 'internal' —
    freshly generated business content is never labelled "Öffentlich".
    """
    ftr = style.get("footer") or {}
    if not ftr.get("classification", False):
        return ""
    try:
        from engine import classification as _cls
        import brain as _brain
        try:
            cfg = _brain._server_config()
        except Exception:
            cfg = None
        res = _cls.detect_classification(content or "", cfg=cfg)
        level = (res.get("content_signals") or {}).get("heuristic_level") or "internal"
        if level == "public":
            level = "internal"
        label = _cls.LEVEL_LABEL_DE.get(level, "Intern")
        return f"Klassifizierung: {label}"
    except Exception:
        return ""


def _ai_disclosure_footer_line(style: dict) -> str:
    """The EU AI Act transparency line, automatic on every generated document
    unless the preset disables it (footer.ai_disclosure: false). Preset may set a
    custom footer.ai_disclosure_text."""
    ftr = style.get("footer") or {}
    if not ftr.get("ai_disclosure", True):
        return ""
    return str(ftr.get("ai_disclosure_text") or _AI_DISCLOSURE_DEFAULT).strip()


def _auto_footer_lines(content: str, style: dict) -> list[str]:
    """Ordered automatic footer lines (each rendered on its OWN line, in the
    footer's font/size): classification first, then the AI-disclosure. The page
    number is handled separately (it must be the LAST line)."""
    lines = []
    cls = _classification_footer_line(content, style)
    if cls:
        lines.append(cls)
    ai = _ai_disclosure_footer_line(style)
    if ai:
        lines.append(ai)
    return lines


def _make_pdf_hdrftr_cb(style: dict, pagesize, margin, inch, hexcolor, content: str = ""):
    """Build a reportlab onPage(canvas, doc) callback drawing the preset's running
    header/footer text (with {page}/{date} tokens) + logo on every page. Returns
    None if the preset has no header/footer/logo to render."""
    hdr = style.get("header") or {}
    ftr = style.get("footer") or {}
    logo = style.get("logo") or {}
    logo_path = _resolve_style_logo(style)
    logo_pos = (logo.get("position") or "header").lower()
    hdr_text = str(hdr.get("text") or "").strip()
    ftr_text = str(ftr.get("text") or "").strip()
    page_nums = bool(ftr.get("page_numbers"))
    auto_lines = _auto_footer_lines(content, style)  # classification + AI-disclosure
    has_logo = bool(logo_path) and logo_pos in ("header", "footer")
    if not (hdr_text or ftr_text or page_nums or auto_lines or has_logo):
        return None
    pw, ph = pagesize
    left, right = margin, pw - margin

    def _xpos(align, text_w):
        a = (align or "left").lower()
        if a == "center":
            return (left + right) / 2 - text_w / 2
        if a == "right":
            return right - text_w
        return left

    def _draw_text(canvas, spec, raw, y, page_num):
        txt = _hdrftr_text(raw, page=True).replace("{page}", str(page_num))
        if not txt:
            return
        size = float(spec.get("font_size") or 9)
        canvas.setFont("Helvetica", size)
        col = spec.get("color")
        if col:
            try:
                canvas.setFillColor(hexcolor(col))
            except Exception:
                pass
        tw = canvas.stringWidth(txt, "Helvetica", size)
        canvas.drawString(_xpos(spec.get("align"), tw), y, txt)

    def _draw_logo(canvas, y_baseline, in_header):
        if not has_logo:
            return
        try:
            from reportlab.lib.utils import ImageReader
            img = ImageReader(logo_path)
            iw, ih = img.getSize()
            w = float(logo.get("width_inch") or 1.2) * inch
            h = w * (ih / iw) if iw else w
            a = (logo.get("align") or "right").lower()
            x = left if a == "left" else (right - w if a == "right"
                                          else (left + right) / 2 - w / 2)
            # Header logo grows upward from the header baseline; footer downward.
            y = y_baseline if in_header else (y_baseline - h)
            canvas.drawImage(img, x, y, width=w, height=h,
                             preserveAspectRatio=True, mask="auto")
        except Exception:
            pass

    def _cb(canvas, doc):
        canvas.saveState()
        pn = canvas.getPageNumber()
        hy = ph - margin + 0.25 * inch   # header baseline, inside top margin
        fy = margin - 0.35 * inch        # footer baseline, inside bottom margin
        if hdr_text:
            _draw_text(canvas, hdr, hdr_text, hy, pn)
        # Footer is a STACK of lines, all in the footer's font/size, bottom-
        # anchored and growing upward (so the page number sits lowest):
        #   [footer text] / [classification] / [AI-disclosure] / [Seite - N]
        # Build bottom→top; the page line is last (drawn at the lowest baseline).
        _explicit_page = "{page}" in ftr_text
        _want_page_line = page_nums and not _explicit_page
        _fsize = float(ftr.get("font_size") or 9)
        _line_h = _fsize + 2
        _stack = []  # top→bottom order
        if ftr_text:
            _stack.append(ftr_text)
        _stack.extend(auto_lines)
        if _want_page_line:
            _stack.append(_PAGE_LINE_PREFIX + "{page}")
        # Lowest line at fy; each earlier line one line-height higher.
        n = len(_stack)
        for i, line in enumerate(_stack):
            _y = fy + (n - 1 - i) * _line_h
            _draw_text(canvas, ftr, line, _y, pn)
        if logo_pos == "header":
            _draw_logo(canvas, hy, in_header=True)
        elif logo_pos == "footer":
            _draw_logo(canvas, fy, in_header=False)
        canvas.restoreState()
    return _cb


def _html_inline_md(text: str) -> str:
    """Convert the inline-markdown subset (bold/italic, `code`, links) to HTML,
    escaping the rest. Mirrors what the docx/pdf branches handle."""
    import html as _html
    out = _html.escape(text)
    out = re.sub(r'\*\*\*(.+?)\*\*\*', r'<strong><em>\1</em></strong>', out)
    out = re.sub(r'\*\*(.+?)\*\*', r'<strong>\1</strong>', out)
    out = re.sub(r'\*(.+?)\*', r'<em>\1</em>', out)
    out = re.sub(r'`([^`]+?)`', r'<code>\1</code>', out)
    out = re.sub(r'\[([^\]]+)\]\((https?://[^)\s]+)\)',
                 r'<a href="\2">\1</a>', out)
    return out


def _html_logo_data_uri(logo_path: str) -> str:
    """Read a logo image into a base64 data URI (so the HTML is self-contained)."""
    import base64
    import mimetypes
    try:
        ctype = mimetypes.guess_type(logo_path)[0] or "image/png"
        with open(logo_path, "rb") as f:
            b64 = base64.b64encode(f.read()).decode("ascii")
        return f"data:{ctype};base64,{b64}"
    except Exception:
        return ""


def _looks_like_html(content: str) -> bool:
    """Heuristic: is `content` already an HTML document (vs markdown)? True when
    it opens with a doctype/<html>, or carries block-level HTML structure (head/
    body/several tags). Avoids markdown-escaping raw HTML the model handed us."""
    head = (content or "").lstrip()[:2000].lower()
    if head.startswith("<!doctype html") or head.startswith("<html"):
        return True
    if "<html" in head or "<body" in head or "<head" in head:
        return True
    # Fallback: many block tags → it's HTML, not markdown that happens to have one.
    return len(re.findall(r'</?(?:div|p|table|section|header|footer|h[1-6]|ul|ol|li|span|article|main|style|script)\b', head)) >= 4


def _html_inline_local_imgs(content: str, doc_dir: str) -> str:
    """Rewrite <img src="localfile"> / <img src='localfile'> to base64 data URIs
    (so the HTML is self-contained + the diagram travels with it). Leaves http(s)/
    data: srcs untouched. Best-effort — an unresolvable src stays as-is."""
    def repl(m):
        quote = m.group(1)
        src = m.group(2)
        if src.startswith(("http://", "https://", "data:")):
            return m.group(0)
        resolved = _resolve_doc_image(src, doc_dir)
        if not resolved:
            return m.group(0)
        uri = _html_logo_data_uri(resolved)  # generic image→data-uri reader
        return f'src={quote}{uri}{quote}' if uri else m.group(0)
    return re.sub(r'''src=(["'])([^"']+)\1''', repl, content)


def _finalize_raw_html(content: str, style: dict, doc_dir: str) -> str:
    """Write raw HTML through (NOT markdown), but make it portable + on-brand:
    inline local <img> as base64, and inject the preset's header/footer/logo
    bands just inside <body> (and the supporting CSS into <head>) when present.
    If there's no <body>, the content is returned with images inlined only."""
    out = _html_inline_local_imgs(content, doc_dir)
    # Build the same header/footer/logo bands the markdown path produces. Reuse
    # _render_markdown_html on an empty body to harvest the band markup + CSS.
    hdr = style.get("header") or {}
    ftr = style.get("footer") or {}
    logo = style.get("logo") or {}
    has_chrome = (str(hdr.get("text") or "").strip() or str(ftr.get("text") or "").strip()
                  or ftr.get("page_numbers") or _resolve_style_logo(style))
    if not has_chrome or not re.search(r'<body[^>]*>', out, re.IGNORECASE):
        return out
    scaffold = _render_markdown_html("", style, doc_dir)
    m_css = re.search(r'<style>(.*?)</style>', scaffold, re.DOTALL)
    m_hdr = re.search(r'(<div class="doc-header".*?</div>)', scaffold, re.DOTALL)
    m_ftr = re.search(r'(<div class="doc-footer".*?</div>)', scaffold, re.DOTALL)
    # Inject the band CSS into <head> (so .doc-header/.doc-footer/.pagenum work),
    # scoped to those classes so it doesn't fight the document's own styling.
    if m_css and re.search(r'</head>', out, re.IGNORECASE):
        band_css = ("\n.doc-header{border-bottom:1px solid #eee;padding-bottom:8px;margin-bottom:20px;overflow:hidden}"
                    ".doc-footer{border-top:1px solid #eee;padding-top:8px;margin-top:28px;overflow:hidden}"
                    ".doc-header img,.doc-footer img{vertical-align:middle}"
                    ".pagenum::after{content:\"\"}@media print{.pagenum::after{content:counter(page)}}\n")
        out = re.sub(r'</head>', f'<style>{band_css}</style></head>', out, count=1, flags=re.IGNORECASE)
    if m_hdr:
        band = m_hdr.group(1)
        out = re.sub(r'(<body[^>]*>)', lambda mm: mm.group(1) + "\n" + band,
                     out, count=1, flags=re.IGNORECASE)
    if m_ftr:
        band = m_ftr.group(1)
        out = re.sub(r'</body>', lambda mm: band + "\n</body>",
                     out, count=1, flags=re.IGNORECASE)
    return out


def _render_markdown_html(content: str, style: dict, doc_dir: str) -> str:
    """Render markdown → a self-contained, styled HTML document applying the
    preset (fonts/colors/sizes/table colors + header/footer/logo). Deterministic
    — same subset as the docx/pdf branches; the model just writes plain markdown.
    Embedded `![alt](file)` images become base64 <img> so the file travels alone.
    {date} resolves; {page} works only when PRINTED (CSS @page counter)."""
    import html as _html
    import base64
    import mimetypes
    f = style.get("fonts") or {}
    sz = style.get("sizes") or {}
    c = style.get("colors") or {}
    hdr = style.get("header") or {}
    ftr = style.get("footer") or {}
    logo = style.get("logo") or {}
    logo_path = _resolve_style_logo(style)
    logo_pos = (logo.get("position") or "header").lower()
    logo_uri = _html_logo_data_uri(logo_path) if logo_path else ""

    def _logo_img(where):
        if not (logo_uri and logo_pos == where and logo_pos != "none"):
            return ""
        a = (logo.get("align") or "right").lower()
        w = float(logo.get("width_inch") or 1.2)
        css = f"max-height:{w * 24:.0f}px;width:auto;"
        if a == "center":
            css += "display:block;margin:0 auto;"
        else:
            css += f"float:{a};"
        return f'<img class="doc-logo" src="{logo_uri}" alt="" style="{css}">'

    def _band(spec, where):
        txt = _hdrftr_text(spec.get("text"), page=True)
        # Page number renders on its OWN line as `Seite - N` (print-only CSS
        # counter), in the same font/size as the footer text — unless the preset
        # text already embeds an explicit {page} token (then honour it inline).
        page_line_html = ""
        if where == "footer" and spec.get("page_numbers") and "{page}" not in txt:
            page_line_html = (f'<div class="doc-pageline">'
                              f'{_html.escape(_PAGE_LINE_PREFIX)}'
                              f'<span class="pagenum"></span></div>')
        # Automatic footer lines (classification + AI-disclosure), each own line.
        auto_html = ""
        if where == "footer":
            for _auto in _auto_footer_lines(content, style):
                auto_html += f'<div class="doc-autoline">{_html.escape(_auto)}</div>'
        # {page} → a print-only CSS counter span; on screen it shows nothing.
        parts = txt.split("{page}")
        html_txt = _html.escape(parts[0])
        for seg in parts[1:]:
            html_txt += '<span class="pagenum"></span>' + _html.escape(seg)
        logo_html = _logo_img(where)
        if not (html_txt.strip() or logo_html or page_line_html or auto_html):
            return ""
        align = spec.get("align") or ("center" if where == "footer" else "left")
        size = spec.get("font_size") or 9
        col = spec.get("color") or "#666666"
        return (f'<div class="doc-{where}" style="text-align:{_html.escape(str(align))};'
                f'font-size:{float(size):.0f}pt;color:{_html.escape(str(col))}">'
                f'{logo_html}{html_txt}{auto_html}{page_line_html}</div>')

    # ── body: reuse the same line-walk as the other branches ──
    body = []
    lines = content.split("\n")
    # Strip a leading YAML frontmatter block (--- … ---) — non-standard framing
    # the model sometimes prefixes; without this it renders as literal <p> text.
    if lines and lines[0].strip() == "---":
        for j in range(1, min(len(lines), 20)):
            if lines[j].strip() == "---":
                lines = lines[j + 1:]
                break
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped_full = line.strip()
        # Horizontal rule: a line of only ---/***/___ (3+) → <hr> (was <p>---</p>).
        if re.match(r'^(-{3,}|\*{3,}|_{3,})$', stripped_full):
            body.append('<hr>')
            i += 1
            continue
        # Blockquote: one or more consecutive `> ` lines → a single <blockquote>.
        if re.match(r'^>\s?', line):
            quote_lines = []
            while i < len(lines) and re.match(r'^>\s?', lines[i]):
                quote_lines.append(re.sub(r'^>\s?', '', lines[i]))
                i += 1
            inner = "".join(f'<p>{_html_inline_md(q.strip())}</p>'
                             for q in quote_lines if q.strip())
            body.append(f'<blockquote>{inner}</blockquote>')
            continue
        # Unordered / ordered list: consecutive `- `/`* `/`+ ` or `1. ` lines.
        m_li = re.match(r'^(\s*)([-*+]|\d+\.)\s+(.*)', line)
        if m_li:
            ordered = bool(re.match(r'^\d+\.$', m_li.group(2)))
            items = []
            while i < len(lines):
                mm = re.match(r'^(\s*)([-*+]|\d+\.)\s+(.*)', lines[i])
                if not mm:
                    break
                items.append(mm.group(3).strip())
                i += 1
            tag = 'ol' if ordered else 'ul'
            lis = "".join(f'<li>{_html_inline_md(it)}</li>' for it in items)
            body.append(f'<{tag}>{lis}</{tag}>')
            continue
        m_img = _MD_IMAGE_RE.match(line)
        if m_img:
            src = _resolve_doc_image(m_img.group(2), doc_dir)
            if src:
                try:
                    ct = mimetypes.guess_type(src)[0] or "image/png"
                    with open(src, "rb") as fh:
                        uri = f"data:{ct};base64,{base64.b64encode(fh.read()).decode('ascii')}"
                    body.append(f'<p class="figure"><img src="{uri}" '
                                f'alt="{_html.escape(m_img.group(1))}"></p>')
                except Exception:
                    body.append(f'<p>[Bild: {_html.escape(m_img.group(1) or m_img.group(2))}]</p>')
            else:
                body.append(f'<p>[Bild nicht gefunden: {_html.escape(m_img.group(2))}]</p>')
            i += 1
            continue
        # Fenced code block: gather to the closing fence. A Mermaid block →
        # render to PNG + inline it; any other code block → <pre><code>.
        m_fence = re.match(r'^(```+|~~~+)\s*([A-Za-z0-9_-]*)\s*$', line)
        if m_fence:
            fence, flang = m_fence.group(1)[0], m_fence.group(2).lower()
            i += 1
            code_lines = []
            while i < len(lines) and not lines[i].strip().startswith(fence * 3):
                code_lines.append(lines[i])
                i += 1
            i += 1  # skip closing fence
            code_text = "\n".join(code_lines)
            _blk = {"type": "code", "text": code_text, "lang": flang}
            _emb = None
            if _is_mermaid_block(_blk, style):
                _png = _render_mermaid_block_to_file(code_text, style, doc_dir)
                if _png:
                    _pp = _resolve_doc_image(_png, doc_dir)
                    try:
                        ct = mimetypes.guess_type(_pp)[0] or "image/png"
                        with open(_pp, "rb") as fh:
                            uri = f"data:{ct};base64,{base64.b64encode(fh.read()).decode('ascii')}"
                        _emb = f'<p class="figure"><img src="{uri}" alt="Diagramm"></p>'
                    except Exception:
                        _emb = None
            if _emb:
                body.append(_emb)
            else:
                body.append(f'<pre><code>{_html.escape(code_text)}</code></pre>')
            continue
        m_h = re.match(r'^(#{1,6})\s+(.*)', line)
        if m_h:
            lvl = len(m_h.group(1))
            body.append(f'<h{lvl}>{_html_inline_md(m_h.group(2))}</h{lvl}>')
            i += 1
            continue
        if "|" in line and i + 1 < len(lines) and re.match(r'^\|[\s\-:|]+\|', lines[i + 1]):
            rows = []
            while i < len(lines) and "|" in lines[i]:
                stripped = lines[i].strip().strip("|")
                if not re.match(r'^[\s\-:|]+$', stripped):
                    rows.append([cell.strip() for cell in stripped.split("|")])
                i += 1
            if rows:
                thead = "".join(f"<th>{_html_inline_md(cell)}</th>" for cell in rows[0])
                tbody = "".join(
                    "<tr>" + "".join(f"<td>{_html_inline_md(cell)}</td>" for cell in r) + "</tr>"
                    for r in rows[1:])
                body.append(f'<table><thead><tr>{thead}</tr></thead><tbody>{tbody}</tbody></table>')
            continue
        stripped = line.strip()
        if stripped:
            body.append(f'<p>{_html_inline_md(stripped)}</p>')
        i += 1

    fb = _html.escape(f.get("body") or "Calibri")
    fh = _html.escape(f.get("heading") or f.get("body") or "Calibri")
    fm = _html.escape(f.get("mono") or "Consolas")
    css = f"""
    @page {{ margin: 2cm; }}
    body {{ font-family: '{fb}', Arial, sans-serif; color: {c.get('body', '#222')};
      font-size: {sz.get('body', 11)}pt; line-height: 1.5; max-width: 820px;
      margin: 0 auto; padding: 24px; background: #fff; }}
    h1,h2,h3,h4,h5,h6 {{ font-family: '{fh}', Arial, sans-serif;
      color: {c.get('heading', '#1F3864')}; line-height: 1.2; margin: 1.1em 0 .4em; }}
    h1 {{ font-size: {sz.get('h1', 20)}pt; }} h2 {{ font-size: {sz.get('h2', 16)}pt; }}
    h3 {{ font-size: {sz.get('h3', 13)}pt; }}
    a {{ color: {c.get('accent', '#2E74B5')}; }}
    code {{ font-family: '{fm}', monospace; background: #f4f4f4; padding: 1px 4px; border-radius: 3px; }}
    p.figure {{ text-align: center; }} p.figure img {{ max-width: 100%; height: auto; }}
    hr {{ border: none; border-top: 1px solid #ddd; margin: 1.4em 0; }}
    blockquote {{ margin: 12px 0; padding: 4px 16px; border-left: 3px solid {c.get('accent', '#2E74B5')};
      color: #555; background: #fafafa; }}
    blockquote p {{ margin: .3em 0; }}
    ul, ol {{ margin: .6em 0 .6em 1.4em; padding: 0; }}
    li {{ margin: .2em 0; }}
    table {{ border-collapse: collapse; width: 100%; margin: 12px 0; }}
    th, td {{ border: 1px solid #ddd; padding: 5px 9px; text-align: left; }}
    th {{ background: {c.get('table_header_bg', '#1F3864')};
      color: {c.get('table_header_text', '#FFFFFF')}; }}
    tbody tr:nth-child(even) {{ background: #fafafa; }}
    .doc-header {{ border-bottom: 1px solid #eee; padding-bottom: 8px; margin-bottom: 20px; overflow: hidden; }}
    .doc-footer {{ border-top: 1px solid #eee; padding-top: 8px; margin-top: 28px; overflow: hidden; }}
    .pagenum::after {{ content: ""; }}
    @media print {{ .pagenum::after {{ content: counter(page); }} }}
    """
    header_band = _band(hdr, "header")
    footer_band = _band(ftr, "footer")
    title = ""
    m_title = re.search(r'^#\s+(.+)$', content, flags=re.MULTILINE)
    if m_title:
        # Strip inline-markdown markers (**bold**/*italic*/`code`) — the <title>
        # is plain text, so they'd otherwise show up literally in the browser tab.
        raw_title = re.sub(r'[*`]', '', m_title.group(1)).strip()
        title = _html.escape(raw_title)
    return (f'<!DOCTYPE html>\n<html lang="de">\n<head>\n<meta charset="UTF-8">\n'
            f'<meta name="viewport" content="width=device-width, initial-scale=1.0">\n'
            f'<title>{title or "Dokument"}</title>\n<style>{css}</style>\n</head>\n<body>\n'
            f'{header_band}\n' + "\n".join(body) + f'\n{footer_band}\n</body>\n</html>\n')


def _docx_align(name: str):
    """Map 'left|center|right' → docx WD_ALIGN_PARAGRAPH (default LEFT)."""
    from docx.enum.text import WD_ALIGN_PARAGRAPH as _A
    return {"left": _A.LEFT, "center": _A.CENTER, "right": _A.RIGHT}.get(
        (name or "left").lower(), _A.LEFT)


# Footer page-number line: rendered as `Seite - N` on its own line, the number a
# live Word PAGE field in the same font/size as the surrounding footer text.
_PAGE_LINE_PREFIX = "Seite - "


def _docx_add_page_field(paragraph, *, size=None, color=None):
    """Append a Word PAGE field (live page number) to a paragraph run, optionally
    styled to a point size + RGB colour so it matches the footer text.
    python-docx has no field API, so emit the field XML directly."""
    from docx.shared import Pt, RGBColor
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    run = paragraph.add_run()
    if size:
        run.font.size = Pt(float(size))
    if color:
        run.font.color.rgb = RGBColor(*color)
    fld_begin = OxmlElement("w:fldChar"); fld_begin.set(qn("w:fldCharType"), "begin")
    instr = OxmlElement("w:instrText"); instr.set(qn("xml:space"), "preserve")
    instr.text = "PAGE"
    fld_end = OxmlElement("w:fldChar"); fld_end.set(qn("w:fldCharType"), "end")
    run._r.append(fld_begin); run._r.append(instr); run._r.append(fld_end)


def _docx_fill_hdrftr(container, spec: dict, *, logo_path, logo_spec, page_token: bool,
                      page_line: bool = False):
    """Populate a docx header/footer container (section.header/footer) with the
    preset text (+{page}/{date} tokens), alignment, size, color, and optional
    logo picture. `container` is a _Header/_Footer; reuses its first paragraph.

    `page_line=True` renders the live page number on ITS OWN paragraph as
    `Seite - N` (same font/size/colour as the footer text), instead of trailing
    the footer text on the same line."""
    from docx.shared import Pt, RGBColor, Inches
    para = container.paragraphs[0] if container.paragraphs else container.add_paragraph()
    para.alignment = _docx_align(spec.get("align"))
    size = spec.get("font_size"); col = _hex_rgb(spec.get("color"))
    text = str(spec.get("text") or "")
    # Logo first (so a left/center logo sits before the text run).
    if logo_path:
        try:
            w = float((logo_spec or {}).get("width_inch") or 1.2)
            para.alignment = _docx_align((logo_spec or {}).get("align"))
            para.add_run().add_picture(logo_path, width=Inches(w))
            if text:
                para.add_run("  ")
        except Exception:
            pass

    def _style_run(r):
        if size:
            r.font.size = Pt(float(size))
        if col:
            r.font.color.rgb = RGBColor(*col)

    def _tighten(p):
        # Multi-line footers stack as separate paragraphs; kill the default
        # space-before/after + loosen-to-1.0 line spacing so the lines sit close.
        pf = p.paragraph_format
        pf.space_before = Pt(0)
        pf.space_after = Pt(0)
        from docx.enum.text import WD_LINE_SPACING
        pf.line_spacing_rule = WD_LINE_SPACING.SINGLE

    _tighten(para)

    # Text, splitting on {page} so the page field renders live (when not on its
    # own line). With page_line, {page} tokens are stripped from the main text and
    # the number goes on a dedicated second paragraph below.
    segments = text.split("{page}") if (page_token and not page_line) else [_hdrftr_text(text)]
    for idx, seg in enumerate(segments):
        seg = _hdrftr_text(seg) if page_token else seg
        if seg:
            _style_run(para.add_run(seg))
        if page_token and not page_line and idx < len(segments) - 1:
            _docx_add_page_field(para)

    # Automatic footer lines (classification, AI-disclosure), each on its own
    # line in the same font/size/colour as the footer text, tightly spaced.
    for _auto in (spec.get("_auto_lines") or []):
        ap = container.add_paragraph()
        ap.alignment = _docx_align(spec.get("align"))
        _tighten(ap)
        _style_run(ap.add_run(str(_auto)))

    if page_line:
        # Dedicated page-number line: `Seite - N`, same font/size/colour.
        pl = container.add_paragraph()
        pl.alignment = _docx_align(spec.get("align"))
        _tighten(pl)
        _style_run(pl.add_run(_PAGE_LINE_PREFIX))
        _docx_add_page_field(pl, size=size, color=col)


def _apply_docx_header_footer(doc, style: dict, doc_dir: str, content: str = ""):
    """Apply preset header/footer/logo to every section of a docx."""
    hdr = style.get("header") or {}
    ftr = style.get("footer") or {}
    logo = style.get("logo") or {}
    logo_path = _resolve_style_logo(style)
    logo_pos = (logo.get("position") or "header").lower()
    has_hdr = bool(str(hdr.get("text") or "").strip()) or (logo_path and logo_pos == "header")
    # Automatic footer lines (classification + AI-disclosure), each own line.
    auto_lines = _auto_footer_lines(content, style)
    # Footer renders if it has text, page numbers, automatic lines, or the logo.
    ftr_text = str(ftr.get("text") or "").strip()
    has_ftr = (bool(ftr_text) or bool(ftr.get("page_numbers")) or bool(auto_lines)
               or (logo_path and logo_pos == "footer"))
    if not (has_hdr or has_ftr):
        return
    # A tall header logo overflows Word's default 0.5" header distance and spills
    # into the body. Measure the rendered logo height and, if the header carries
    # it, push the header distance + top margin down to clear it (+ a little air).
    from docx.shared import Inches
    _logo_h_in = 0.0
    if logo_path and logo_pos in ("header", "footer"):
        try:
            from PIL import Image
            with Image.open(logo_path) as _im:
                _iw, _ih = _im.size
            _w_in = float(logo.get("width_inch") or 1.2)
            _logo_h_in = _w_in * (_ih / _iw) if _iw else _w_in
        except Exception:
            _logo_h_in = float(logo.get("width_inch") or 1.2)  # square fallback
    for sec in doc.sections:
        if logo_pos == "header" and _logo_h_in > 0:
            try:
                _hdist = Inches(0.4)
                sec.header_distance = _hdist
                # logo bottom = header_distance + logo height; body must start a
                # comfortable gap below that so a heading never crowds the logo.
                _need_top = _hdist + Inches(_logo_h_in) + Inches(0.35)
                if int(sec.top_margin) < int(_need_top):
                    sec.top_margin = _need_top
            except Exception:
                pass
        elif logo_pos == "footer" and _logo_h_in > 0:
            try:
                _fdist = Inches(0.35)
                sec.footer_distance = _fdist
                _need_bot = _fdist + Inches(_logo_h_in) + Inches(0.15)
                if int(sec.bottom_margin) < int(_need_bot):
                    sec.bottom_margin = _need_bot
            except Exception:
                pass
        if has_hdr:
            _docx_fill_hdrftr(sec.header, hdr,
                              logo_path=logo_path if logo_pos == "header" else None,
                              logo_spec=logo, page_token=True)
        if has_ftr:
            # The page number renders on its OWN line as `Seite - N` (page_line),
            # not trailing the footer text. If the preset text itself contains an
            # explicit {page} token we honour that inline instead.
            f_spec = dict(ftr)
            f_spec["_auto_lines"] = auto_lines
            _explicit_page = "{page}" in str(f_spec.get("text") or "")
            _want_page_line = bool(ftr.get("page_numbers")) and not _explicit_page
            _docx_fill_hdrftr(sec.footer, f_spec,
                              logo_path=logo_path if logo_pos == "footer" else None,
                              logo_spec=logo, page_token=True, page_line=_want_page_line)


# ── Markdown → block model (markdown-it-py) ────────────────────────────────
# markdown-it-py is already installed (markitdown dependency). It ONLY tokenises
# (text → token list) — it renders nothing itself. We translate its token stream
# into a small, format-agnostic block model below; the docx + pdf renderers each
# consume that SAME model, so the two paths stay in lock-step and we keep full
# control of the look (badges/zebra/cover/KPI all live in OUR renderer, not the
# lib). This replaced a hand-rolled line parser that silently dropped lists,
# blockquotes, code fences and links. Our non-standard ::kpi + cover frontmatter
# are stripped in a pre-pass BEFORE markdown-it ever sees the text.

# An inline run carries text + bold/italic/mono/link — built from markdown-it's
# inline children so headings/cells/paragraphs render identically everywhere.
class _Run:
    __slots__ = ("text", "bold", "italic", "mono", "href")
    def __init__(self, text, bold=False, italic=False, mono=False, href=None):
        self.text, self.bold, self.italic, self.mono, self.href = text, bold, italic, mono, href


def _inline_tokens_to_runs(tok):
    """Flatten a markdown-it 'inline' token's children into a list of _Run.
    Handles strong/em/code_inline/link/softbreak. Nested emphasis composes."""
    runs = []
    if tok is None or not getattr(tok, "children", None):
        # plain inline content with no markup
        if tok is not None and tok.content:
            runs.append(_Run(tok.content))
        return runs
    bold = ital = 0
    href = [None]
    for ch in tok.children:
        t = ch.type
        if t == "strong_open": bold += 1
        elif t == "strong_close": bold = max(0, bold - 1)
        elif t == "em_open": ital += 1
        elif t == "em_close": ital = max(0, ital - 1)
        elif t == "link_open":
            href[0] = dict(ch.attrs).get("href") if ch.attrs else None
        elif t == "link_close":
            href[0] = None
        elif t == "code_inline":
            runs.append(_Run(ch.content, mono=True, bold=bool(bold), italic=bool(ital), href=href[0]))
        elif t == "softbreak" or t == "hardbreak":
            runs.append(_Run(" "))
        elif t == "text":
            if ch.content:
                runs.append(_Run(ch.content, bold=bool(bold), italic=bool(ital), href=href[0]))
        elif t == "image":
            # inline images are rare in our reports; keep alt text as a fallback
            alt = ch.content or ""
            if alt:
                runs.append(_Run(alt, italic=True))
    return runs or ([_Run(tok.content)] if tok.content else [])


def _runs_plain(runs):
    return "".join(r.text for r in runs)


def _md_parser():
    from markdown_it import MarkdownIt
    return (MarkdownIt("commonmark", {"html": False})
            .enable("table").enable("strikethrough"))


def _markdown_to_blocks(md_text):
    """Parse markdown into a flat list of block dicts the renderers consume:
      {"type":"heading","level":n,"runs":[...]}
      {"type":"paragraph","runs":[...]}
      {"type":"hr"}
      {"type":"code","text":str,"lang":str}
      {"type":"quote","blocks":[...]}            # nested blocks
      {"type":"list","ordered":bool,"items":[[blocks...], ...],"level":n}
      {"type":"table","rows":[[runs,...], ...]}  # row 0 = header
    Lists nest via item-blocks that may themselves contain a 'list'."""
    toks = _md_parser().parse(md_text)
    # Walk the flat token stream with an explicit container stack.
    root = []
    stack = [root]            # current block-list to append to
    list_stack = []           # active list dicts
    item_stack = []           # active list-item block-lists
    pending_table = None
    cur_row = None
    in_header = False

    def top(): return stack[-1]

    i = 0
    while i < len(toks):
        t = toks[i]
        ty = t.type
        if ty == "heading_open":
            lvl = int(t.tag[1])
            inline = toks[i + 1] if i + 1 < len(toks) else None
            top().append({"type": "heading", "level": lvl,
                          "runs": _inline_tokens_to_runs(inline)})
            i += 3; continue
        if ty == "paragraph_open":
            inline = toks[i + 1] if i + 1 < len(toks) else None
            top().append({"type": "paragraph", "runs": _inline_tokens_to_runs(inline)})
            i += 3; continue
        if ty == "hr":
            top().append({"type": "hr"}); i += 1; continue
        if ty == "fence" or ty == "code_block":
            top().append({"type": "code", "text": t.content.rstrip("\n"),
                          "lang": (t.info or "").strip()})
            i += 1; continue
        if ty == "bullet_list_open" or ty == "ordered_list_open":
            lst = {"type": "list", "ordered": ty.startswith("ordered"),
                   "items": [], "level": len(list_stack)}
            top().append(lst)
            list_stack.append(lst)
            i += 1; continue
        if ty == "list_item_open":
            item_blocks = []
            list_stack[-1]["items"].append(item_blocks)
            stack.append(item_blocks)   # children append into this item
            i += 1; continue
        if ty == "list_item_close":
            stack.pop()
            i += 1; continue
        if ty == "bullet_list_close" or ty == "ordered_list_close":
            list_stack.pop()
            i += 1; continue
        if ty == "blockquote_open":
            qblocks = []
            top().append({"type": "quote", "blocks": qblocks})
            stack.append(qblocks)
            i += 1; continue
        if ty == "blockquote_close":
            stack.pop()
            i += 1; continue
        if ty == "table_open":
            pending_table = {"type": "table", "rows": []}
            top().append(pending_table)
            i += 1; continue
        if ty == "table_close":
            pending_table = None
            i += 1; continue
        if ty in ("thead_open", "thead_close", "tbody_open", "tbody_close"):
            in_header = ty == "thead_open" or (ty == "tbody_close" and in_header)
            i += 1; continue
        if ty == "tr_open":
            cur_row = []; i += 1; continue
        if ty == "tr_close":
            if pending_table is not None and cur_row is not None:
                pending_table["rows"].append(cur_row)
            cur_row = None; i += 1; continue
        if ty in ("th_open", "td_open"):
            inline = toks[i + 1] if i + 1 < len(toks) else None
            if cur_row is not None:
                cur_row.append(_inline_tokens_to_runs(inline))
            i += 3; continue
        i += 1
    return root


# ── Opus-near docx polish (all deterministic) ──────────────────────────────
# Leading emoji + variation selectors a model tends to prefix onto headings
# (📌📊📜🏢🔍📈🛡️📉🎯📎 …). Stripped for regulatory docs when docx.strip_emoji on.
_EMOJI_RE = re.compile(
    "[\U0001F000-\U0001FAFF\U00002600-\U000027BF\U0001F1E6-\U0001F1FF"
    "\U00002190-\U000021FF\U00002B00-\U00002BFF\U0000FE00-\U0000FE0F\U0000200D]+")
_INLINE_MD_SPLIT = re.compile(r'(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*[^*]+?\*|`[^`]+?`)')


def _clean_heading_text(text: str, strip_emoji: bool) -> str:
    """Strip a leading emoji (+ trailing whitespace) from a heading when enabled.
    Only removes emoji at the START — emoji inside running prose is left alone."""
    t = text
    if strip_emoji:
        t = _EMOJI_RE.sub("", t).strip()
        t = re.sub(r"^[\s•\-–—]+", "", t)  # leftover bullet/dash glue
    return t.strip()


def _add_inline_md_runs(paragraph, text: str, *, base_bold=False, base_italic=False,
                        color=None, mono_font="Consolas"):
    """Parse inline **bold**/*italic*/***both***/`code` and append styled runs to
    `paragraph`. The SINGLE place inline markdown is rendered — used by paragraphs,
    HEADINGS and TABLE CELLS alike (the old code only did this for paragraphs, so
    ** / * leaked verbatim into headings + table headers — the visible bug)."""
    from docx.shared import RGBColor
    parts = _INLINE_MD_SPLIT.split(text)
    for part in parts:
        if not part:
            continue
        b, it, mono, txt = base_bold, base_italic, False, part
        if part.startswith("***") and part.endswith("***") and len(part) > 6:
            b, it, txt = True, True, part[3:-3]
        elif part.startswith("**") and part.endswith("**") and len(part) > 4:
            b, txt = True, part[2:-2]
        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
            it, txt = True, part[1:-1]
        elif part.startswith("`") and part.endswith("`") and len(part) > 2:
            mono, txt = True, part[1:-1]
        run = paragraph.add_run(txt)
        run.bold = b
        run.italic = it
        if mono:
            run.font.name = mono_font
        if color:
            rgb = _hex_rgb(color)
            if rgb:
                run.font.color.rgb = RGBColor(*rgb)


def _docx_add_runs(paragraph, runs, *, base_bold=False, color=None, mono_font="Consolas"):
    """Render a list of _Run (from the markdown-it block model) into a docx
    paragraph — bold/italic/mono + clickable hyperlinks. The token-based sibling
    of _add_inline_md_runs (which parses raw **md** text); both coexist so the
    cover/KPI helpers can keep using the text variant."""
    from docx.shared import RGBColor
    for r in runs:
        if r.href:
            _docx_add_hyperlink(paragraph, r.text, r.href, color=color)
            continue
        run = paragraph.add_run(r.text)
        run.bold = bool(r.bold or base_bold)
        run.italic = bool(r.italic)
        if r.mono:
            run.font.name = mono_font
        if color:
            rgb = _hex_rgb(color)
            if rgb:
                run.font.color.rgb = RGBColor(*rgb)


def _docx_add_hyperlink(paragraph, text, url, *, color=None):
    """Append a real clickable hyperlink run (python-docx has no API for it).
    Falls back to a plain blue-underlined run if relationship wiring fails."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import RGBColor
    try:
        part = paragraph.part
        r_id = part.relate_to(
            url, "http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            is_external=True)
        hyperlink = OxmlElement("w:hyperlink")
        hyperlink.set(qn("r:id"), r_id)
        new_run = OxmlElement("w:r")
        rPr = OxmlElement("w:rPr")
        c = OxmlElement("w:color"); c.set(qn("w:val"), "2E74B5"); rPr.append(c)
        u = OxmlElement("w:u"); u.set(qn("w:val"), "single"); rPr.append(u)
        new_run.append(rPr)
        t = OxmlElement("w:t"); t.text = text; new_run.append(t)
        hyperlink.append(new_run)
        paragraph._p.append(hyperlink)
    except Exception:
        run = paragraph.add_run(text)
        run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)
        run.underline = True


def _docx_render_list(doc, lst, style, *, level=0):
    """Render a (possibly nested) list block as real Word list paragraphs.
    Bullet → 'List Bullet[ N]', ordered → 'List Number[ N]'. Nested lists recurse
    with a deeper Word list style so Word shows the indent + sub-bullets."""
    base = "List Number" if lst.get("ordered") else "List Bullet"
    lvl = min(level, 2)
    style_name = base if lvl == 0 else f"{base} {lvl + 1}"
    mono = style["fonts"].get("mono", "Consolas")
    for item in lst["items"]:
        first_para_done = False
        for blk in item:
            bt = blk["type"]
            if bt == "paragraph" and not first_para_done:
                try:
                    p = doc.add_paragraph(style=style_name)
                except KeyError:
                    p = doc.add_paragraph(style=base)
                _docx_add_runs(p, blk["runs"], mono_font=mono)
                first_para_done = True
            elif bt == "list":
                _docx_render_list(doc, blk, style, level=level + 1)
            elif bt == "paragraph":
                # continuation paragraph inside the same item
                p = doc.add_paragraph()
                p.paragraph_format.left_indent = __import__("docx").shared.Inches(0.25 * (lvl + 1))
                _docx_add_runs(p, blk["runs"], mono_font=mono)
            elif bt == "code":
                _docx_render_code(doc, blk, style)
            elif bt == "quote":
                for qb in blk["blocks"]:
                    if qb.get("type") == "paragraph":
                        _docx_render_quote_para(doc, qb["runs"], style)


def _docx_render_quote_para(doc, runs, style):
    """A blockquote paragraph: left indent + a left accent border + italic grey."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import Inches
    p = doc.add_paragraph()
    p.paragraph_format.left_indent = Inches(0.3)
    pPr = p._p.get_or_add_pPr()
    pbdr = OxmlElement("w:pBdr")
    left = OxmlElement("w:left")
    left.set(qn("w:val"), "single"); left.set(qn("w:sz"), "18")
    left.set(qn("w:space"), "8")
    left.set(qn("w:color"), (style["docx"].get("rule_color") or "#B4C6E7").lstrip("#"))
    pbdr.append(left); pPr.append(pbdr)
    _docx_add_runs(p, runs, mono_font=style["fonts"].get("mono", "Consolas"))
    for r in p.runs:
        r.italic = True


def _is_mermaid_block(blk: dict, style: dict) -> bool:
    """A code block that holds a Mermaid diagram — either fenced as ```mermaid or
    a bare ```gantt/flowchart/sequenceDiagram… block the model wrote without the
    `mermaid` language tag (the bdb6a1e7 case). Gated on `mermaid.embed` (default
    on) so a preset can opt out of auto-rendering."""
    if (style.get("mermaid") or {}).get("embed", True) is False:
        return False
    if blk.get("type") != "code":
        return False
    lang = (blk.get("lang") or "").strip().lower()
    if lang in ("mermaid", "mmd"):
        return True
    if lang and lang not in ("", "text", "plain"):
        return False  # an explicit non-mermaid language → leave as code
    try:
        from engine.tools.image_gen import looks_like_mermaid
        return looks_like_mermaid(blk.get("text") or "")
    except Exception:
        return False


_MERMAID_DOC_SEQ = {}  # doc_dir → counter, for unique auto-diagram filenames


def _render_mermaid_block_to_file(code: str, style: dict, doc_dir: str, fmt: str = "png"):
    """Render a Mermaid code block to an image FILE inside the document folder and
    return its basename (so the existing image-embed path can place it), or None on
    any failure (caller then falls back to rendering the raw code block)."""
    try:
        from engine.tools.image_gen import render_mermaid_file
        _n = _MERMAID_DOC_SEQ.get(doc_dir, 0) + 1
        _MERMAID_DOC_SEQ[doc_dir] = _n
        out_name = f"_auto_diagram_{_n}.{fmt}"
        out_path = os.path.join(doc_dir, out_name)
        merm = style.get("mermaid") or {}
        bg = (merm.get("background") or "white").lower()
        res = render_mermaid_file(code, out_path=out_path, full_style=style,
                                  fmt=fmt, theme=merm.get("theme"),
                                  background=bg, explicit_theme=False)
        if res and os.path.exists(out_path):
            try:
                import brain as _brain
                _ag = _brain.get_request_context().current_agent or _brain._current_agent
                _brain._after_file_write(out_path, "created", _ag.agent_id if _ag else "main")
            except Exception:
                pass
            return out_name
    except Exception:
        pass
    return None


def _docx_add_fitted_picture(doc, img_path):
    """Insert an image scaled to fit the section's usable width AND height (so a
    wide/tall auto-rendered diagram doesn't overflow the page). Centred."""
    from docx.shared import Emu
    sec = doc.sections[0]
    usable_w = int(sec.page_width) - int(sec.left_margin) - int(sec.right_margin)
    usable_h = int(sec.page_height) - int(sec.top_margin) - int(sec.bottom_margin)
    usable_h = int(usable_h * 0.92)  # leave a little vertical air
    try:
        from PIL import Image
        with Image.open(img_path) as im:
            iw, ih = im.size  # pixels
        # python-docx maps px→EMU at 96 dpi by default (9525 EMU/px).
        nat_w = iw * 9525
        nat_h = ih * 9525
        ratio = min(usable_w / nat_w, usable_h / nat_h, 1.0)
        p = doc.add_paragraph()
        p.alignment = 1  # center
        p.add_run().add_picture(img_path, width=Emu(int(nat_w * ratio)))
    except Exception:
        from docx.shared import Inches
        doc.add_picture(img_path, width=Inches(6.0))


def _docx_render_code(doc, blk, style):
    """A fenced code block: monospace runs on a light-shaded paragraph, one Word
    line per source line (no \\n in a run — Word needs separate breaks)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import Pt
    mono = style["fonts"].get("mono", "Consolas")
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear"); shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), "F2F2F2")
    pPr.append(shd)
    lines = blk["text"].split("\n")
    for li, ln in enumerate(lines):
        run = p.add_run(ln)
        run.font.name = mono
        run.font.size = Pt(max(8, int(style["sizes"].get("body", 11)) - 1))
        if li < len(lines) - 1:
            run.add_break()


def _shade_cell(cell, hex_color):
    """Set a table cell's background fill (w:shd) — python-docx has no API for it."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    rgb = _hex_rgb(hex_color)
    if not rgb:
        return
    tcPr = cell._tc.get_or_add_tcPr()
    shd = OxmlElement("w:shd")
    shd.set(qn("w:val"), "clear")
    shd.set(qn("w:color"), "auto")
    shd.set(qn("w:fill"), "%02X%02X%02X" % rgb)
    tcPr.append(shd)


def _section_usable_width_emu(doc):
    """Page width minus left/right margins (EMU) for the first section — the room
    a full-width table actually has. Falls back to ~6.5in (Letter, 1in margins)."""
    try:
        sec = doc.sections[0]
        w = int(sec.page_width) - int(sec.left_margin) - int(sec.right_margin)
        return w if w > 914400 else int(6.5 * 914400)
    except Exception:
        return int(6.5 * 914400)


def _docx_keep_with_next(paragraph):
    """Mark a paragraph 'keep with next' so a heading never strands at the bottom
    of a page, split from the content it titles (Schusterjunge avoidance)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    pPr = paragraph._p.get_or_add_pPr()
    if pPr.find(qn("w:keepNext")) is None:
        pPr.append(OxmlElement("w:keepNext"))
    # widowControl: Word default is on, but a heading style may have cleared it.
    if pPr.find(qn("w:keepLines")) is None:
        pPr.append(OxmlElement("w:keepLines"))


def _docx_row_cant_split(row):
    """Forbid a single table row from breaking across a page (a row's cells stay
    together — kills the 'one line on the old page, rest on the next' artefact)."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    trPr = row._tr.get_or_add_trPr()
    if trPr.find(qn("w:cantSplit")) is None:
        trPr.append(OxmlElement("w:cantSplit"))


def _docx_row_repeat_header(row):
    """Mark a row as a repeating header (re-drawn at the top of every page the
    table spans) — so a multi-page table keeps its column titles."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    trPr = row._tr.get_or_add_trPr()
    if trPr.find(qn("w:tblHeader")) is None:
        trPr.append(OxmlElement("w:tblHeader"))


def _content_col_widths(rows, max_cols, total_emu, *, min_emu):
    """Distribute table width proportionally to each column's longest cell, with a
    floor per column so a short-header column never collapses to a 1-char-wide
    sliver that wraps every word. Returns a list of EMU widths summing to total."""
    longest = [1] * max_cols
    for r in rows:
        for c in range(max_cols):
            cell = r[c] if c < len(r) else ""
            # cell may be runs (list) or already-plain text
            txt = _runs_plain(cell) if isinstance(cell, list) else str(cell)
            # longest single line drives width (a wrapped paragraph is fine)
            longest[c] = max(longest[c], max((len(seg) for seg in txt.split("\n")), default=1))
    # Cap any single column's influence so one giant cell doesn't starve others.
    capped = [min(w, 60) for w in longest]
    floor = min(min_emu, total_emu // max_cols)  # never demand more than equal share
    free = total_emu - floor * max_cols
    if free < 0:  # too many columns for the floor — fall back to equal split
        return [total_emu // max_cols] * max_cols
    denom = sum(capped) or 1
    return [floor + int(free * (w / denom)) for w in capped]


def _docx_apply_col_widths(table, widths_emu):
    """Pin column widths on a docx table. Word reads a FIXED-layout table's widths
    from the <w:tblGrid> (in twips) AND the per-cell <w:tcW>; we set BOTH so every
    renderer (Word, LibreOffice) honours them and a narrow column can't collapse."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import Emu
    table.autofit = False
    tbl = table._tbl
    tblPr = tbl.tblPr
    layout = tblPr.find(qn("w:tblLayout"))
    if layout is None:
        layout = OxmlElement("w:tblLayout")
        tblPr.append(layout)
    layout.set(qn("w:type"), "fixed")
    # Rewrite the table grid (EMU → twips: 1 twip = 635 EMU).
    twips = [max(1, int(round(w / 635))) for w in widths_emu]
    grid = tbl.find(qn("w:tblGrid"))
    if grid is not None:
        tbl.remove(grid)
    grid = OxmlElement("w:tblGrid")
    for tw in twips:
        gc = OxmlElement("w:gridCol")
        gc.set(qn("w:w"), str(tw))
        grid.append(gc)
    # tblGrid must directly follow tblPr.
    tblPr.addnext(grid)
    for ci, w in enumerate(widths_emu):
        for row in table.rows:
            if ci < len(row.cells):
                row.cells[ci].width = Emu(w)


# Risk-badge palette: (text-colour, fill). Order matters — most specific first
# ('sehr gut' before 'gut', 'hoch' substring caught after 'erhöht').
_RISK_BADGES = [
    (("sehr gut", "gering", "niedrig", "low"), ("548235", "E2EFDA")),
    (("erhöht", "erhoht", "elevated"), ("C55A11", "FCE4D6")),
    (("hoch", "high", "hohes"), ("C00000", "F8D7DA")),
    (("mittel", "angemessen", "medium", "moderat"), ("BF8F00", "FFF2CC")),
]


def _risk_badge(value: str):
    """Map a Bewertung/Risiko cell value → (fg_hex, bg_hex) or None if no match."""
    v = (value or "").strip().lower()
    if not v:
        return None
    for keys, (fg, bg) in _RISK_BADGES:
        if any(k in v for k in keys):
            return fg, bg
    return None


_BADGE_COL_HINTS = ("bewertung", "risiko", "rating", "einstufung", "risk", "stufe")

# A heading that introduces the document's change/version history. When the model
# writes one (per the doc-styles convention), the renderer starts it on a fresh
# page so the history table isn't split across a page boundary.
_VERSION_HISTORY_RE = re.compile(
    r'^\s*(versions?historie|änderungs?historie|aenderungs?historie|'
    r'dokumenten?historie|revisionshistorie|change\s*history|version\s*history|'
    r'revision\s*history|document\s*history)\s*$',
    re.IGNORECASE)


def _is_version_history_heading(text: str) -> bool:
    return bool(_VERSION_HISTORY_RE.match((text or "").strip()))


def _add_hrule(doc, color):
    """Add a horizontal divider paragraph (--- in markdown) as a bottom border."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    p = doc.add_paragraph()
    pPr = p._p.get_or_add_pPr()
    pbdr = OxmlElement("w:pBdr")
    bottom = OxmlElement("w:bottom")
    bottom.set(qn("w:val"), "single")
    bottom.set(qn("w:sz"), "6")
    bottom.set(qn("w:space"), "1")
    bottom.set(qn("w:color"), (color or "#B4C6E7").lstrip("#"))
    pbdr.append(bottom)
    pPr.append(pbdr)
    return p


def _docx_force_update_fields(doc):
    """Set <w:updateFields val="true"/> in settings.xml so Word/LibreOffice
    recalculates ALL fields (the TOC, page numbers) on open — without this the
    TOC stays an empty placeholder forever (the user never sees a TOC). Idempotent."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    settings = doc.settings.element
    # Word ignores updateFields unless it's the first child of <w:settings>.
    existing = settings.find(qn("w:updateFields"))
    if existing is not None:
        existing.set(qn("w:val"), "true")
        return
    uf = OxmlElement("w:updateFields")
    uf.set(qn("w:val"), "true")
    settings.insert(0, uf)


def _docx_bookmark_heading(paragraph, name):
    """Wrap a heading paragraph's runs in a bookmark so a TOC PAGEREF can target
    it. The bookmarkStart MUST go AFTER <w:pPr> (pPr must stay the first child of
    <w:p> per the OOXML schema) — inserting it at index 0 puts it before pPr, which
    makes Word bind the bookmark wrong and resolve every PAGEREF to page 1.
    Returns the bookmark name."""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    bid = str(abs(hash(name)) % 1_000_000)
    start = OxmlElement("w:bookmarkStart")
    start.set(qn("w:id"), bid)
    start.set(qn("w:name"), name)
    end = OxmlElement("w:bookmarkEnd")
    end.set(qn("w:id"), bid)
    p = paragraph._p
    pPr = p.find(qn("w:pPr"))
    if pPr is not None:
        pPr.addnext(start)   # bookmarkStart immediately AFTER pPr
    else:
        p.insert(0, start)
    p.append(end)
    return name


def _add_toc_field(doc, headings=None):
    """Insert a NATIVE Word Table-of-Contents field (heading levels 1-3). Word
    generates the entries + correct page numbers itself from the Heading 1-3
    styles when the document is opened (the field is marked dirty + settings.xml
    requests a recalc), so the TOC fills in automatically and F9 only re-flows.
    A short placeholder shows until that first open. Deterministic, no model input.

    The earlier approach pre-built manual PAGEREF entries, but those resolved
    every entry to page 1 (bookmarks placed before pagination) — the native field
    is what makes the page numbers correct in Word. `headings` is kept for
    signature compatibility but unused now.
    """
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    from docx.shared import Pt
    title = doc.add_paragraph()
    _add_inline_md_runs(title, "Inhaltsverzeichnis")
    for r in title.runs:
        r.bold = True
        r.font.size = Pt(16)
    p = doc.add_paragraph()
    run = p.add_run()
    fb = OxmlElement("w:fldChar"); fb.set(qn("w:fldCharType"), "begin")
    fb.set(qn("w:dirty"), "true")
    instr = OxmlElement("w:instrText"); instr.set(qn("xml:space"), "preserve")
    instr.text = 'TOC \\o "1-3" \\h \\z \\u'
    sep = OxmlElement("w:fldChar"); sep.set(qn("w:fldCharType"), "separate")
    ph = OxmlElement("w:t"); ph.set(qn("xml:space"), "preserve")
    ph.text = "Inhaltsverzeichnis \u2014 wird beim \u00d6ffnen in Word automatisch erzeugt (sonst F9)."
    fe = OxmlElement("w:fldChar"); fe.set(qn("w:fldCharType"), "end")
    run._r.append(fb); run._r.append(instr); run._r.append(sep)
    run._r.append(ph); run._r.append(fe)
    try:
        _docx_force_update_fields(doc)
    except Exception:
        pass


def _extract_cover_and_body(content, dx):
    """Pre-pass over RAW markdown (before markdown-it): if the doc is substantial
    and starts with an H1 (+ optional 'Key: value' frontmatter), peel those off
    for the cover page and return (title, [(k,v)...], body_markdown). The cover
    lines are NON-standard framing we render ourselves; the rest goes to
    markdown-it. Returns (None, [], content) when no cover applies."""
    if not dx.get("cover", True):
        return None, [], content
    lines = content.split("\n")
    n_headings = sum(1 for l in lines if re.match(r'^#{1,6}\s', l))
    n_nonblank = sum(1 for l in lines if l.strip())
    title, tidx = None, None
    for li, ln in enumerate(lines[:40]):
        m = re.match(r'^#\s+(.*)', ln)
        if m:
            title, tidx = m.group(1).strip(), li
            break
    if title is None:
        return None, [], content
    # frontmatter directly under the H1 (until blank / --- / non-kv line)
    front, consumed_to = [], tidx
    for li in range(tidx + 1, min(tidx + 12, len(lines))):
        raw = lines[li].strip()
        if not raw or raw == "---":
            consumed_to = li if raw == "---" else li - 1
            break
        fm = re.match(r'^\*{0,2}([A-Za-zÄÖÜäöü][^:]{1,40}?)\*{0,2}:\s*(.+)$', raw)
        if fm and "|" not in raw and not raw.startswith("::"):
            front.append((fm.group(1).strip(), re.sub(r'\*+', '', fm.group(2)).strip()))
            consumed_to = li
        else:
            consumed_to = li - 1
            break
    substantial = (n_headings >= 4) or bool(front) or (n_nonblank >= 30)
    if not substantial:
        return None, [], content
    body = "\n".join(lines[consumed_to + 1:])
    return title, front, body


def _kpi_match(line: str):
    """A '::kpi VALUE | LABEL | risk' convention line → (value, label, badge) or
    None. Lets the model flag a headline metric as a coloured stat box WITHOUT
    free-form layout — deterministic trigger, any model can emit it."""
    m = re.match(r'^\s*::kpi\s+(.*)$', line)
    if not m:
        return None
    parts = [p.strip() for p in m.group(1).split("|")]
    value = parts[0] if parts else ""
    label = parts[1] if len(parts) > 1 else ""
    badge = parts[2] if len(parts) > 2 else (label or value)
    return value, label, badge


def _emit_kpi_strip(doc, kpis, style):
    """Render a row of coloured KPI boxes from collected ::kpi lines."""
    from docx.shared import Pt, RGBColor
    if not kpis:
        return
    table = doc.add_table(rows=1, cols=len(kpis))
    table.autofit = True
    for ci, (value, label, badge) in enumerate(kpis):
        cell = table.rows[0].cells[ci]
        pal = _risk_badge(badge) or ("44546A", "EDF1F8")
        _shade_cell(cell, "#" + pal[1])
        cell.paragraphs[0].alignment = 1  # center
        rv = cell.paragraphs[0].add_run(value)
        rv.bold = True
        rv.font.size = Pt(20)
        crgb = _hex_rgb("#" + pal[0])
        if crgb:
            rv.font.color.rgb = RGBColor(*crgb)
        if label:
            pl = cell.add_paragraph()
            pl.alignment = 1
            rl = pl.add_run(label.upper())
            rl.font.size = Pt(8.5)
            lrgb = _hex_rgb(style["colors"].get("heading"))
            if lrgb:
                rl.font.color.rgb = RGBColor(*lrgb)


def _render_cover_page(doc, style, title, frontmatter):
    """Render a deterministic title page from the FIRST # H1 + leading key:value
    'frontmatter' lines (Stichtag: …, Verantwortlich: …). Followed by a page
    break. No model layout — code composes it from content already in the md."""
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    accent = style["colors"].get("accent", "#2E74B5")
    navy = style["colors"].get("heading", "#1F3864")
    for _ in range(6):
        doc.add_paragraph()
    # confidential kicker
    kick = doc.add_paragraph()
    rk = kick.add_run("VERTRAULICH · INTERNE RISIKOANALYSE")
    rk.bold = True
    rk.font.size = Pt(10)
    krgb = _hex_rgb(style["colors"].get("accent"))
    if krgb:
        rk.font.color.rgb = RGBColor(*krgb)
    _add_hrule(doc, navy)
    # title — the cover title is always bold, so drop inline **/* markers outright
    # (don't leave them verbatim like the old heading path did).
    _ttl = _clean_heading_text(title, style["docx"].get("strip_emoji", True))
    _ttl = re.sub(r'[*`]+', '', _ttl).strip()
    tp = doc.add_paragraph()
    rt = tp.add_run(_ttl)
    rt.bold = True
    rt.font.size = Pt(30)
    rt.font.name = style["fonts"].get("heading", "Calibri")
    trgb = _hex_rgb(navy)
    if trgb:
        rt.font.color.rgb = RGBColor(*trgb)
    # frontmatter meta lines
    if frontmatter:
        doc.add_paragraph()
        for key, val in frontmatter:
            mp = doc.add_paragraph()
            rkk = mp.add_run(key.upper() + "   ")
            rkk.bold = True
            rkk.font.size = Pt(9.5)
            srgb = _hex_rgb(style["colors"].get("body"))
            mp.add_run(val).font.size = Pt(10.5)
    doc.add_page_break()


def tool_write_document(args: dict) -> str:
    """Create documents from markdown content. Markdown ![alt](file) image
    references are EMBEDDED (docx/pptx/pdf) — pair with render_diagram to put
    data-accurate diagrams into a report/presentation."""
    import brain as _brain
    path = args.get("path", "")
    content = args.get("content", "")
    # Editorial report layout: style='report' (alias 'editorial') routes HTML
    # output through the SAME renderer Deep Research + Studio use
    # (engine.report_html.render_report_html) — the warm editorial look (drop-cap,
    # gradient headings, sticky TOC, aurora background) instead of the Word/PDF-
    # like doc-styles preset. Detected from the RAW style arg (before the
    # default-preset fill) so it only fires when explicitly asked for. Only
    # meaningful for .html; other formats fall through to the normal preset path.
    _raw_style = (args.get("style", "") or "").strip().lower()
    _editorial = _raw_style in ("report", "editorial")
    # Resolve the preset to apply — explicit style= wins, else a sensible default
    # (project/global/'corporate') so output is styled even when the model omits
    # style= (the common case). _load_doc_style('') would give the bare built-in.
    _style_name = "" if _editorial else _resolve_default_style(args.get("style", ""))
    # style='reference' / 'reference:<filename>' → lift the look FROM a reference
    # .docx (a project instruction-file template) instead of applying a brand
    # preset. Only meaningful for .docx output; for other formats it degrades to
    # the built-in default (no Word styles to lift), noted below.
    _ref_note = ""
    if _style_name == "reference" or _style_name.startswith("reference:"):
        _ref_file = _style_name.split(":", 1)[1] if ":" in _style_name else ""
        _ref_path = _resolve_reference_docx(_ref_file)
        if _ref_path:
            _style, _ref_note = _load_doc_style_from_reference(_ref_path)
            _ref_note = f"Referenz-Stil aus {os.path.basename(_ref_path)}: {_ref_note}"
        else:
            _style = _load_doc_style("")
            _ref_note = ("Keine Referenz-.docx in den Projekt-Begleitdateien gefunden "
                         "— Standard-Stil angewandt.")
    else:
        _style = _load_doc_style(_style_name)
    try:
        # Hard guard: resolved path MUST be inside the session artifact folder
        # (relative names default into it; absolute / .. escapes refused). This
        # also fixes the old relative→repo-root default.
        path, err = _enforce_artifact_path(path, "write_document")
        if err:
            return err
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        ext = os.path.splitext(path)[1].lower()
        _doc_dir = os.path.dirname(path) or "."

        if ext == ".docx":
            try:
                import docx
                from docx.shared import Pt, RGBColor
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            doc = docx.Document()
            # Apply style: default body font + size on the Normal style; heading
            # font/color/size on Heading 1-3. Deterministic — model writes plain md.
            try:
                _ds = _style
                _normal = doc.styles["Normal"]
                _normal.font.name = _ds["fonts"]["body"]
                _normal.font.size = Pt(_ds["sizes"]["body"])
                _bodyc = _hex_rgb(_ds["colors"]["body"])
                if _bodyc:
                    _normal.font.color.rgb = RGBColor(*_bodyc)
                _hc = _hex_rgb(_ds["colors"]["heading"])
                for _lvl, _szkey in ((1, "h1"), (2, "h2"), (3, "h3")):
                    try:
                        _hs = doc.styles[f"Heading {_lvl}"]
                        _hs.font.name = _ds["fonts"]["heading"]
                        _hs.font.size = Pt(_ds["sizes"][_szkey])
                        _hs.font.bold = bool(_ds["docx"].get("heading_bold", True))
                        if _hc:
                            _hs.font.color.rgb = RGBColor(*_hc)
                    except KeyError:
                        pass
            except Exception:
                pass
            # Running header / footer + logo (deterministic). python-docx exposes
            # section.header/footer paragraphs; the page-number is a Word field.
            try:
                _apply_docx_header_footer(doc, _style, _doc_dir, content)
            except Exception:
                pass
            _dx = _style["docx"]
            _strip_emoji = bool(_dx.get("strip_emoji", True))
            _zebra = _dx.get("zebra_fill")
            _rule_col = _dx.get("rule_color", "#B4C6E7")
            _hdr_bg = _style["colors"].get("table_header_bg", "#1F3864")
            _hdr_fg = _style["colors"].get("table_header_text", "#FFFFFF")
            _do_badges = bool(_dx.get("risk_badges", True))
            _mono = _style["fonts"].get("mono", "Consolas")

            # Cover page + TOC (substantial reports only — see _extract_cover_and_body).
            _cover_title, _cover_front, _body_md = _extract_cover_and_body(content, _dx)

            # Parse the body with markdown-it (real lists/quotes/code/links/tables)
            # then render the block model. ::kpi lines arrive as plain paragraphs —
            # we recognise + group them here into coloured stat-strips.
            _blocks = _markdown_to_blocks(_body_md)

            # Pre-derive the (level, text, bookmark) list from headings ≤ level 3 so
            # the TOC can be PRE-POPULATED (visible + page-numbered) before the body
            # is rendered; the same bookmark names are pinned on the headings below.
            _toc_headings = []
            _bm_for_block = {}
            for _bi, _blk in enumerate(_blocks):
                if _blk.get("type") == "heading" and _blk.get("level", 9) <= 3:
                    _htxt = _clean_heading_text(_runs_plain(_blk["runs"]), _strip_emoji)
                    _htxt = re.sub(r'[*`]+', '', _htxt).strip()
                    if not _htxt or _is_version_history_heading(_runs_plain(_blk["runs"])):
                        continue  # skip the version-history appendix in the TOC
                    _bm = f"_Toc_h{_bi}"
                    _bm_for_block[_bi] = _bm
                    _toc_headings.append((_blk["level"], _htxt, _bm))

            if _cover_title is not None:
                _render_cover_page(doc, _style, _cover_title, _cover_front)
                if _dx.get("toc", True):
                    _add_toc_field(doc, _toc_headings)
                    doc.add_page_break()

            def _docx_render_table_block(blk):
                rows = blk.get("rows") or []
                if not rows:
                    return
                max_cols = max(len(r) for r in rows)
                table = doc.add_table(rows=len(rows), cols=max_cols)
                try:
                    table.style = _dx.get("table_style") or "Table Grid"
                except (KeyError, Exception):
                    table.style = "Table Grid"
                # Content-proportional column widths with a per-column floor, so a
                # narrow-header column doesn't collapse and wrap to one char/line.
                try:
                    from docx.shared import Inches as _In
                    _usable = _section_usable_width_emu(doc)
                    _widths = _content_col_widths(rows, max_cols, _usable,
                                                  min_emu=int(0.7 * 914400))
                    _docx_apply_col_widths(table, _widths)
                except Exception:
                    pass
                # Badge column by evidence (same rule as before, now over runs).
                _badge_col = None
                if _do_badges and len(rows) > 1:
                    _hdr_l = [_runs_plain(rows[0][c]).strip().lower() if c < len(rows[0]) else ""
                              for c in range(max_cols)]
                    _best, _best_hits = None, 0
                    for _ci in range(max_cols):
                        _hits = sum(1 for _r in rows[1:]
                                    if _ci < len(_r) and _risk_badge(_runs_plain(_r[_ci])))
                        _hinted = any(_hdr_l[_ci] == hint or _hdr_l[_ci].endswith(hint)
                                      for hint in _BADGE_COL_HINTS)
                        _score = _hits + (0.5 if _hinted else 0)
                        if _hits >= max(1, (len(rows) - 1) // 2) and _score > _best_hits:
                            _best, _best_hits = _ci, _score
                    _badge_col = _best
                for ri, row in enumerate(rows):
                    is_header = ri == 0
                    # Keep each row intact across a page break (no orphaned single
                    # line); repeat the header row atop every page the table spans.
                    try:
                        _docx_row_cant_split(table.rows[ri])
                        if is_header:
                            _docx_row_repeat_header(table.rows[ri])
                    except Exception:
                        pass
                    for ci in range(max_cols):
                        cell = table.rows[ri].cells[ci]
                        cell.text = ""
                        para = cell.paragraphs[0]
                        runs = row[ci] if ci < len(row) else []
                        cell_plain = _runs_plain(runs)
                        badge = (not is_header and ci == _badge_col and _risk_badge(cell_plain)) or None
                        if is_header:
                            _shade_cell(cell, _hdr_bg)
                            _docx_add_runs(para, runs, base_bold=True, color=_hdr_fg, mono_font=_mono)
                        elif badge:
                            _shade_cell(cell, "#" + badge[1])
                            para.alignment = 1
                            _docx_add_runs(para, runs, base_bold=True, color="#" + badge[0], mono_font=_mono)
                        else:
                            if _zebra and ri % 2 == 0:
                                _shade_cell(cell, _zebra)
                            _docx_add_runs(para, runs, mono_font=_mono)

            _kpis = []
            def _flush_kpis():
                if _kpis:
                    _emit_kpi_strip(doc, _kpis, _style)
                    _kpis.clear()

            for _blk_idx, blk in enumerate(_blocks):
                bt = blk["type"]
                # A paragraph that is only '::kpi …' lines → collect into a strip.
                if bt == "paragraph":
                    txt = _runs_plain(blk["runs"])
                    kpi_lines = [_kpi_match(l) for l in txt.split("\n")]
                    if kpi_lines and all(kpi_lines):
                        _kpis.extend(kpi_lines)
                        continue
                    _flush_kpis()
                    # An image-only paragraph (![alt](file)) → embed the picture.
                    _im = _MD_IMAGE_RE.match(txt.strip())
                    if _im:
                        _ip = _resolve_doc_image(_im.group(2), _doc_dir)
                        if _ip and _ip.lower().endswith(".svg"):
                            doc.add_paragraph(f"[Bild {_im.group(2)} — für DOCX bitte als PNG rendern (render_diagram format=png)]")
                        elif _ip:
                            try:
                                from docx.shared import Inches
                                doc.add_picture(_ip, width=Inches(6.0))
                            except Exception:
                                doc.add_paragraph(f"[Bild: {_im.group(1) or _im.group(2)}]")
                        else:
                            doc.add_paragraph(f"[Bild nicht gefunden: {_im.group(2)}]")
                        continue
                    para = doc.add_paragraph()
                    _docx_add_runs(para, blk["runs"], mono_font=_mono)
                    continue
                _flush_kpis()
                if bt == "heading":
                    htext = _clean_heading_text(_runs_plain(blk["runs"]), _strip_emoji)
                    # Start the version-history section on a fresh page so its
                    # table stays whole (and reads as a clear document appendix).
                    if _is_version_history_heading(_runs_plain(blk["runs"])):
                        doc.add_page_break()
                    h = doc.add_heading("", level=min(blk["level"], 9))
                    # keep inline emphasis inside the heading, but feed the
                    # emoji-cleaned plain text (headings rarely carry links)
                    _add_inline_md_runs(h, htext)
                    # Pin the bookmark the TOC PAGEREF entry targets (same name).
                    _bm = _bm_for_block.get(_blk_idx)
                    if _bm:
                        try:
                            _docx_bookmark_heading(h, _bm)
                        except Exception:
                            pass
                    # Don't let a heading strand alone at a page bottom.
                    try:
                        _docx_keep_with_next(h)
                    except Exception:
                        pass
                elif bt == "hr":
                    _add_hrule(doc, _rule_col)
                elif bt == "list":
                    _docx_render_list(doc, blk, _style)
                elif bt == "quote":
                    for qb in blk["blocks"]:
                        if qb.get("type") == "paragraph":
                            _docx_render_quote_para(doc, qb["runs"], _style)
                        elif qb.get("type") == "list":
                            _docx_render_list(doc, qb, _style)
                elif bt == "code":
                    # A Mermaid diagram block → render to PNG + embed as a picture
                    # (don't dump raw `gantt …` source as a code block). Falls back
                    # to the code block if rendering is unavailable / fails.
                    _merm_png = (_render_mermaid_block_to_file(blk["text"], _style, _doc_dir)
                                 if _is_mermaid_block(blk, _style) else None)
                    if _merm_png:
                        _mp = _resolve_doc_image(_merm_png, _doc_dir)
                        try:
                            _docx_add_fitted_picture(doc, _mp)
                        except Exception:
                            _docx_render_code(doc, blk, _style)
                    else:
                        _docx_render_code(doc, blk, _style)
                elif bt == "table":
                    _docx_render_table_block(blk)
            _flush_kpis()
            doc.save(path)

        elif ext == ".xlsx":
            try:
                import openpyxl
            except ImportError:
                return _err("Install openpyxl: pip3 install openpyxl")
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            sections = re.split(r'^##\s+(.+)$', content, flags=re.MULTILINE)
            if len(sections) < 3:
                ws = wb.create_sheet("Sheet1")
                _write_md_table_to_sheet(ws, content)
            else:
                for si in range(1, len(sections), 2):
                    sheet_name = sections[si].strip()
                    if sheet_name.lower().startswith("sheet:"):
                        sheet_name = sheet_name[6:].strip()
                    sheet_content = sections[si + 1] if si + 1 < len(sections) else ""
                    ws = wb.create_sheet(sheet_name[:31])
                    _write_md_table_to_sheet(ws, sheet_content)
            if not wb.sheetnames:
                wb.create_sheet("Sheet1")
            wb.save(path)

        elif ext == ".pptx":
            try:
                from pptx import Presentation
                from pptx.dml.color import RGBColor as _PRGB
            except ImportError:
                return _err("Install python-pptx: pip3 install python-pptx")
            prs = Presentation()
            _title_rgb = _hex_rgb(_style["pptx"].get("title_color"))
            _body_rgb = _hex_rgb(_style["pptx"].get("body_color"))

            def _style_slide(slide):
                """Apply preset title/body colors to a slide's placeholders."""
                try:
                    if slide.shapes.title and _title_rgb:
                        for _p in slide.shapes.title.text_frame.paragraphs:
                            for _r in _p.runs:
                                _r.font.color.rgb = _PRGB(*_title_rgb)
                    for _ph in slide.placeholders:
                        if _ph == slide.shapes.title:
                            continue
                        if _body_rgb and _ph.has_text_frame:
                            for _p in _ph.text_frame.paragraphs:
                                for _r in _p.runs:
                                    _r.font.color.rgb = _PRGB(*_body_rgb)
                except Exception:
                    pass
            from pptx.util import Inches as _PInches, Pt as _PPt
            _pptx_strip_emoji = bool(_style["docx"].get("strip_emoji", True))
            _pptx_do_badges = bool(_style["docx"].get("risk_badges", True))
            _pptx_hdr_bg = _hex_rgb(_style["colors"].get("table_header_bg", "#1F3864"))
            _pptx_hdr_fg = _hex_rgb(_style["colors"].get("table_header_text", "#FFFFFF"))
            _pptx_zebra = _hex_rgb(_style["docx"].get("zebra_fill", "#EDF1F8"))

            # Inline runs → pptx paragraph runs (bold/italic/mono/link), the pptx
            # sibling of _docx_add_runs. Same markdown-it block model as docx/pdf.
            def _pptx_add_runs(paragraph, runs):
                for r in runs:
                    run = paragraph.add_run()
                    run.text = r.text
                    run.font.bold = bool(r.bold)
                    run.font.italic = bool(r.italic)
                    if r.mono:
                        run.font.name = _style["fonts"].get("mono", "Consolas")
                    if r.href:
                        try:
                            run.hyperlink.address = r.href
                        except Exception:
                            pass

            def _pptx_emit_blocks(tf, blocks, level=0):
                """Render block list into a slide text-frame. Lists/quotes/code
                become indented bullets (slide-appropriate); tables are handled by
                the caller (they become real pptx tables, not text)."""
                for blk in blocks:
                    bt = blk["type"]
                    if bt == "paragraph":
                        p = tf.add_paragraph()
                        p.level = min(level, 4)
                        _pptx_add_runs(p, blk["runs"])
                    elif bt == "list":
                        for item in blk["items"]:
                            wrote_first = False
                            for sub in item:
                                if sub["type"] == "paragraph" and not wrote_first:
                                    p = tf.add_paragraph()
                                    p.level = min(level, 4)
                                    _pptx_add_runs(p, sub["runs"])
                                    wrote_first = True
                                elif sub["type"] == "list":
                                    _pptx_emit_blocks(tf, [sub], level + 1)
                    elif bt == "quote":
                        for qb in blk["blocks"]:
                            if qb.get("type") == "paragraph":
                                p = tf.add_paragraph()
                                p.level = min(level, 4)
                                _pptx_add_runs(p, qb["runs"])
                                for _r in p.runs:
                                    _r.font.italic = True
                    elif bt == "code":
                        for cl in blk["text"].split("\n"):
                            p = tf.add_paragraph()
                            p.level = min(level + 1, 4)
                            run = p.add_run(); run.text = cl or " "
                            run.font.name = _style["fonts"].get("mono", "Consolas")
                            run.font.size = _PPt(14)
                    elif bt == "hr":
                        tf.add_paragraph()  # blank separator line

            def _pptx_add_table(slide, blk, top_in):
                """Render a markdown table as a REAL pptx table with header fill,
                zebra rows and risk-badge cell colouring (parity with docx/pdf)."""
                rows = blk.get("rows") or []
                if not rows:
                    return top_in
                nrows, ncols = len(rows), max(len(r) for r in rows)
                width = _PInches(9.0); height = _PInches(0.4 * nrows)
                gtbl = slide.shapes.add_table(nrows, ncols, _PInches(0.5),
                                              _PInches(top_in), width, height).table
                # badge column by evidence
                _bcol = None
                if _pptx_do_badges and nrows > 1:
                    _hl = [_runs_plain(rows[0][c]).strip().lower() if c < len(rows[0]) else ""
                           for c in range(ncols)]
                    _best, _bsc = None, 0
                    for _ci in range(ncols):
                        _hits = sum(1 for _r in rows[1:]
                                    if _ci < len(_r) and _risk_badge(_runs_plain(_r[_ci])))
                        _hint = any(_hl[_ci] == h or _hl[_ci].endswith(h) for h in _BADGE_COL_HINTS)
                        _sc = _hits + (0.5 if _hint else 0)
                        if _hits >= max(1, (nrows - 1) // 2) and _sc > _bsc:
                            _best, _bsc = _ci, _sc
                    _bcol = _best
                for ri, row in enumerate(rows):
                    for ci in range(ncols):
                        cell = gtbl.cell(ri, ci)
                        runs = row[ci] if ci < len(row) else []
                        cell.text = ""
                        p = cell.text_frame.paragraphs[0]
                        is_header = ri == 0
                        badge = (not is_header and ci == _bcol and _risk_badge(_runs_plain(runs))) or None
                        if is_header:
                            if _pptx_hdr_bg:
                                cell.fill.solid(); cell.fill.fore_color.rgb = _PRGB(*_pptx_hdr_bg)
                            _pptx_add_runs(p, runs)
                            for _r in p.runs:
                                _r.font.bold = True
                                if _pptx_hdr_fg:
                                    _r.font.color.rgb = _PRGB(*_pptx_hdr_fg)
                        elif badge:
                            cell.fill.solid(); cell.fill.fore_color.rgb = _PRGB(*_hex_rgb("#" + badge[1]))
                            _pptx_add_runs(p, runs)
                            for _r in p.runs:
                                _r.font.bold = True
                                _r.font.color.rgb = _PRGB(*_hex_rgb("#" + badge[0]))
                        else:
                            if _pptx_zebra and ri % 2 == 0:
                                cell.fill.solid(); cell.fill.fore_color.rgb = _PRGB(*_pptx_zebra)
                            _pptx_add_runs(p, runs)
                        for _r in p.runs:
                            _r.font.size = _PPt(12)
                return top_in + 0.45 * nrows + 0.3

            # Slide boundary = a level-1 OR level-2 heading (`# ` / `## `). Decks
            # are commonly written `# Deck Title` + `## Slide 1` + `## Slide 2`, so
            # both start a new slide; deeper headings (### …) stay in the body.
            slides_content = re.split(r'^#{1,2}\s+(.+)$', content, flags=re.MULTILINE)
            if len(slides_content) < 3:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = "Slide 1"
                tf = slide.placeholders[1].text_frame
                tf.clear()
                _blocks = _markdown_to_blocks(content.strip())
                _first = tf.paragraphs[0]
                _para_blocks = [b for b in _blocks if b["type"] != "table"]
                if _para_blocks and _para_blocks[0]["type"] == "paragraph":
                    _pptx_add_runs(_first, _para_blocks[0]["runs"])
                    _para_blocks = _para_blocks[1:]
                _pptx_emit_blocks(tf, _para_blocks)
            else:
                for si in range(1, len(slides_content), 2):
                    title = _clean_heading_text(slides_content[si].strip(), _pptx_strip_emoji)
                    body = slides_content[si + 1].strip() if si + 1 < len(slides_content) else ""
                    _blocks = _markdown_to_blocks(body)
                    # Separate images (placed as pictures) from flow blocks.
                    _img_paths = []
                    _flow = []
                    for b in _blocks:
                        if b["type"] == "paragraph":
                            _m = _MD_IMAGE_RE.match(_runs_plain(b["runs"]).strip())
                            if _m:
                                _ip = _resolve_doc_image(_m.group(2), _doc_dir)
                                if _ip and not _ip.lower().endswith(".svg"):
                                    _img_paths.append(_ip)
                                    continue
                        _flow.append(b)
                    _tables = [b for b in _flow if b["type"] == "table"]
                    _text_blocks = [b for b in _flow if b["type"] != "table"]
                    _has_text = any(
                        (b["type"] in ("paragraph", "list", "quote", "code")) and
                        (_runs_plain(b["runs"]).strip() if b["type"] == "paragraph" else True)
                        for b in _text_blocks)
                    if _img_paths and not _has_text and not _tables:
                        # Picture-focused slide.
                        slide = prs.slides.add_slide(prs.slide_layouts[5])
                        slide.shapes.title.text = title
                        slide.shapes.add_picture(_img_paths[0], _PInches(1.0), _PInches(1.6),
                                                 height=_PInches(5.0))
                    else:
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = title
                        tf = slide.placeholders[1].text_frame
                        tf.clear()
                        # First text block fills paragraphs[0]; rest append.
                        if _text_blocks:
                            first = _text_blocks[0]
                            if first["type"] == "paragraph":
                                _pptx_add_runs(tf.paragraphs[0], first["runs"])
                                _pptx_emit_blocks(tf, _text_blocks[1:])
                            else:
                                _pptx_emit_blocks(tf, _text_blocks)
                        # Tables become real pptx tables, stacked below the text.
                        _top = 4.2 if _has_text else 1.8
                        for _tb in _tables:
                            _top = _pptx_add_table(slide, _tb, _top)
                        if _img_paths and not _tables:
                            try:
                                slide.shapes.add_picture(_img_paths[0], _PInches(5.2),
                                                         _PInches(1.8), height=_PInches(4.0))
                            except Exception:
                                pass
            for _sl in prs.slides:
                _style_slide(_sl)
            try:
                _apply_pptx_logo_footer(prs, _style)
            except Exception:
                pass
            prs.save(path)

        elif ext == ".pdf":
            try:
                from reportlab.lib.pagesizes import letter, A4
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as _RLImage, Table as _RLTable, TableStyle as _RLTableStyle, PageBreak as _RLPageBreak
                from reportlab.platypus.flowables import HRFlowable as _RLHRFlowable
                from reportlab.platypus.tableofcontents import TableOfContents as _RLToC
                from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle as _RLParaStyle
                from reportlab.lib.units import inch as _rl_inch
                from reportlab.lib.colors import HexColor as _RLHex
                from reportlab.lib import colors as _rl_colors
                from reportlab.lib.enums import TA_LEFT as _RL_TA_LEFT, TA_CENTER as _RL_TA_CENTER
            except ImportError:
                return _err("Install reportlab: pip3 install reportlab")
            _pdf_dx = _style["docx"]  # polish keys are shared (strip_emoji/badges/cover/…)
            _pdf_strip_emoji = bool(_pdf_dx.get("strip_emoji", True))
            _pdf_rule_hex = _pdf_dx.get("rule_color", "#B4C6E7")
            _pdf_do_badges = bool(_pdf_dx.get("risk_badges", True))
            _psize = A4 if str(_style["pdf"].get("page_size", "letter")).lower() == "a4" else letter
            _marg = float(_style["pdf"].get("margin_inch", 1.0)) * _rl_inch
            _pdf_want_toc = bool(_pdf_dx.get("toc", True))
            _pdf_toc_active = False  # True once a ToC flowable is actually placed
            # Reserve bottom space for the full footer stack (text + classification
            # + AI-disclosure + page line) so the running footer never collides
            # with the body. Footer baseline is margin-0.35in; each extra line adds
            # one line-height above it.
            _ftr = _style.get("footer") or {}
            _ftr_lines = (1 if str(_ftr.get("text") or "").strip() else 0)
            _ftr_lines += len(_auto_footer_lines(content, _style))
            if _ftr.get("page_numbers") and "{page}" not in str(_ftr.get("text") or ""):
                _ftr_lines += 1
            _ftr_fsize = float(_ftr.get("font_size") or 9)
            _ftr_reserve = max(0, _ftr_lines - 1) * (_ftr_fsize + 2) + 0.45 * _rl_inch
            _bot_marg = max(_marg, _ftr_reserve)
            doc_pdf = SimpleDocTemplate(path, pagesize=_psize, topMargin=_marg,
                                        bottomMargin=_bot_marg, leftMargin=_marg, rightMargin=_marg)
            # PDF table-of-contents: reportlab fills a ToC flowable over a 2-pass
            # build (multiBuild) driven by 'TOCEntry' notifications. We tag each
            # heading flowable with a bookmark key + emit the notify in an
            # afterFlowable hook bound onto this template instance.
            _pdf_toc_seq = [0]
            def _pdf_after_flowable(flowable):
                _txt = getattr(flowable, "_toc_text", None)
                if _txt is not None:
                    doc_pdf.notify("TOCEntry", (getattr(flowable, "_toc_level", 0),
                                                _txt, doc_pdf.page,
                                                getattr(flowable, "_toc_key", None)))
            doc_pdf.afterFlowable = _pdf_after_flowable
            styles = getSampleStyleSheet()
            # Apply style fonts/colors to the reportlab paragraph styles.
            try:
                styles["Normal"].fontName = _style["fonts"]["body"] if False else styles["Normal"].fontName  # keep core font (reportlab needs registered TTFs for custom)
                styles["Normal"].fontSize = _style["sizes"]["body"]
                _bc = _style["colors"].get("body")
                if _bc:
                    styles["Normal"].textColor = _RLHex(_bc)
                _hc = _style["colors"].get("heading")
                for _hn in ("Heading1", "Heading2", "Heading3"):
                    if _hn in styles:
                        if _hc:
                            styles[_hn].textColor = _RLHex(_hc)
                        # Keep a heading on the same page as the content it titles
                        # (no heading stranded at a page bottom).
                        styles[_hn].keepWithNext = 1
            except Exception:
                pass
            # Inline markdown (**bold**/*italic*) → reportlab mini-HTML markup.
            # Escape raw &/< first (reportlab Paragraph parses XML — a bare '&'
            # in 'M&P AM' would swallow following text as a bogus entity), then
            # apply the bold/italic tags so they survive as real markup.
            def _pdf_inline(_s):
                _s = _s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")
                _s = re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', _s)
                _s = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', _s)
                _s = re.sub(r'\*(.+?)\*', r'<i>\1</i>', _s)
                return _s

            def _esc(_s):
                return _s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

            # Runs (markdown-it block model) → reportlab mini-HTML markup, with
            # real <a href> links + monospace via <font face>. The token sibling
            # of _pdf_inline (which parses raw **md** text).
            def _pdf_runs(runs):
                out = []
                for r in runs:
                    t = _esc(r.text)
                    if r.mono:
                        t = f'<font face="Courier">{t}</font>'
                    if r.bold:
                        t = f"<b>{t}</b>"
                    if r.italic:
                        t = f"<i>{t}</i>"
                    if r.href:
                        href = _esc(r.href)
                        t = f'<a href="{href}" color="#2E74B5"><u>{t}</u></a>'
                    out.append(t)
                return "".join(out)

            # A cell paragraph style so long cell text wraps inside the column
            # instead of overflowing (reportlab Table needs flowables to wrap).
            _cell_style = _RLParaStyle(
                "TableCell", parent=styles["Normal"],
                fontSize=max(7, int(_style["sizes"]["body"]) - 2), leading=max(9, int(_style["sizes"]["body"])),
                alignment=_RL_TA_LEFT)
            _hdr_text_hex = _style["colors"].get("table_header_text", "#FFFFFF")
            _hdr_cell_style = _RLParaStyle(
                "TableHeaderCell", parent=_cell_style, textColor=_RLHex(_hdr_text_hex))
            _hdr_bg_hex = _style["colors"].get("table_header_bg") or _style["colors"].get("heading", "#1F3864")
            _grid_hex = _style["colors"].get("accent", "#999999")

            story = []
            lines = content.split("\n")

            # KPI strip helper: build a coloured 1-row table from ::kpi lines.
            def _pdf_emit_kpis(kpis):
                if not kpis:
                    return
                cells, cmds = [], []
                for ci, (value, label, badge) in enumerate(kpis):
                    pal = _risk_badge(badge) or ("44546A", "EDF1F8")
                    vsty = _RLParaStyle(f"Kpi{ci}", parent=styles["Normal"], fontSize=22,
                                        leading=24, alignment=_RL_TA_CENTER,
                                        textColor=_RLHex("#" + pal[0]), fontName="Helvetica-Bold")
                    lsty = _RLParaStyle(f"KpiL{ci}", parent=styles["Normal"], fontSize=8,
                                        leading=10, alignment=_RL_TA_CENTER,
                                        textColor=_RLHex(_style["colors"].get("heading", "#1F3864")))
                    inner = _RLTable([[Paragraph(value, vsty)], [Paragraph((label or "").upper(), lsty)]])
                    inner.setStyle(_RLTableStyle([("TOPPADDING", (0, 0), (-1, -1), 6),
                                                  ("BOTTOMPADDING", (0, 0), (-1, -1), 6)]))
                    cells.append(inner)
                    cmds.append(("BACKGROUND", (ci, 0), (ci, 0), _RLHex("#" + pal[1])))
                avail_w = _psize[0] - 2 * _marg
                strip = _RLTable([cells], colWidths=[avail_w / len(kpis)] * len(kpis))
                strip.setStyle(_RLTableStyle([
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                    ("LEFTPADDING", (0, 0), (-1, -1), 6), ("RIGHTPADDING", (0, 0), (-1, -1), 6),
                    ("INNERGRID", (0, 0), (-1, -1), 3, _rl_colors.white),
                ] + cmds))
                story.append(strip)
                story.append(Spacer(1, 14))

            # ── Cover page + TOC (shared substance gate with docx) ──────────────
            _pdf_title, _pdf_front, _pdf_body = _extract_cover_and_body(content, _pdf_dx)
            if _pdf_title is not None:
                _navy = _style["colors"].get("heading", "#1F3864")
                _acc = _style["colors"].get("accent", "#2E74B5")
                story.append(Spacer(1, 2.0 * _rl_inch))
                story.append(Paragraph("VERTRAULICH · INTERNE RISIKOANALYSE", _RLParaStyle(
                    "CoverKick", parent=styles["Normal"], fontSize=10, textColor=_RLHex(_acc), fontName="Helvetica-Bold")))
                story.append(_RLHRFlowable(width="100%", thickness=2, color=_RLHex(_navy), spaceBefore=6, spaceAfter=18))
                _ttl_txt = re.sub(r'[*`]+', '', _clean_heading_text(_pdf_title, _pdf_strip_emoji))
                story.append(Paragraph(_ttl_txt.replace("&", "&amp;").replace("<", "&lt;"), _RLParaStyle(
                    "CoverTitle", parent=styles["Normal"], fontSize=28, leading=32, textColor=_RLHex(_navy), fontName="Helvetica-Bold")))
                if _pdf_front:
                    story.append(Spacer(1, 0.4 * _rl_inch))
                    for _k, _v in _pdf_front:
                        _ke = _k.upper().replace("&", "&amp;").replace("<", "&lt;")
                        _ve = _v.replace("&", "&amp;").replace("<", "&lt;")
                        story.append(Paragraph(f"<b>{_ke}</b>&nbsp;&nbsp;&nbsp;{_ve}", _RLParaStyle(
                            "CoverMeta", parent=styles["Normal"], fontSize=10, leading=18)))
                story.append(_RLPageBreak())
                if _pdf_want_toc:
                    story.append(Paragraph("Inhaltsverzeichnis", _RLParaStyle(
                        "TocTitle", parent=styles["Normal"], fontSize=16, leading=20,
                        textColor=_RLHex(_navy), fontName="Helvetica-Bold", spaceAfter=12)))
                    _toc = _RLToC()
                    _toc.levelStyles = [
                        _RLParaStyle("TOC0", parent=styles["Normal"], fontSize=11, leading=16,
                                     leftIndent=0, firstLineIndent=0, spaceBefore=4,
                                     textColor=_RLHex(_style["colors"].get("heading", "#1F3864"))),
                        _RLParaStyle("TOC1", parent=styles["Normal"], fontSize=10, leading=14, leftIndent=16),
                        _RLParaStyle("TOC2", parent=styles["Normal"], fontSize=9.5, leading=13, leftIndent=32,
                                     textColor=_RLHex("#666666")),
                    ]
                    story.append(_toc)
                    story.append(_RLPageBreak())
                    _pdf_toc_active = True

            # ── Render the body via the markdown-it block model ─────────────────
            _pdf_hkey = [0]
            _pdf_kpis = []

            def _pdf_flush_kpis():
                if _pdf_kpis:
                    _pdf_emit_kpis(list(_pdf_kpis))
                    _pdf_kpis.clear()

            def _pdf_table_flowable(blk):
                rows = blk.get("rows") or []
                if not rows:
                    return None
                max_cols = max(len(r) for r in rows)
                _bcol = None
                if _pdf_do_badges and len(rows) > 1:
                    _hl = [_runs_plain(rows[0][c]).strip().lower() if c < len(rows[0]) else ""
                           for c in range(max_cols)]
                    _best, _bsc = None, 0
                    for _ci in range(max_cols):
                        _hits = sum(1 for _r in rows[1:]
                                    if _ci < len(_r) and _risk_badge(_runs_plain(_r[_ci])))
                        _hint = any(_hl[_ci] == h or _hl[_ci].endswith(h) for h in _BADGE_COL_HINTS)
                        _sc = _hits + (0.5 if _hint else 0)
                        if _hits >= max(1, (len(rows) - 1) // 2) and _sc > _bsc:
                            _best, _bsc = _ci, _sc
                    _bcol = _best
                data, _badge_cmds = [], []
                for ri, row in enumerate(rows):
                    cells = []
                    cstyle = _hdr_cell_style if ri == 0 else _cell_style
                    for ci in range(max_cols):
                        runs = row[ci] if ci < len(row) else []
                        plain = _runs_plain(runs)
                        _bdg = (ri > 0 and ci == _bcol and _risk_badge(plain)) or None
                        _cs = cstyle
                        if _bdg:
                            _cs = _RLParaStyle(f"Badge{ri}{ci}", parent=_cell_style,
                                               textColor=_RLHex("#" + _bdg[0]),
                                               alignment=_RL_TA_CENTER, fontName="Helvetica-Bold")
                            _badge_cmds.append(("BACKGROUND", (ci, ri), (ci, ri), _RLHex("#" + _bdg[1])))
                        cells.append(Paragraph(_pdf_runs(runs) or "&nbsp;", _cs))
                    data.append(cells)
                avail_w = _psize[0] - 2 * _marg
                # Content-proportional widths with a per-column floor so a narrow
                # header doesn't wrap to one char/line; sum is pinned to avail_w.
                _cw = _content_col_widths(rows, max_cols, int(avail_w),
                                          min_emu=int(0.6 * _rl_inch))
                _wsum = sum(_cw) or 1
                _cw = [w * avail_w / _wsum for w in _cw]
                tbl = _RLTable(data, colWidths=_cw, repeatRows=1)
                tbl.setStyle(_RLTableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), _RLHex(_hdr_bg_hex)),
                    ("GRID", (0, 0), (-1, -1), 0.5, _RLHex(_grid_hex)),
                    ("VALIGN", (0, 0), (-1, -1), "TOP"),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1), [_rl_colors.white, _RLHex(_pdf_dx.get("zebra_fill", "#F2F4F8"))]),
                    ("LEFTPADDING", (0, 0), (-1, -1), 5), ("RIGHTPADDING", (0, 0), (-1, -1), 5),
                    ("TOPPADDING", (0, 0), (-1, -1), 3), ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
                ] + _badge_cmds))
                return tbl

            def _pdf_render_list(blk, depth=0):
                """Flat-render a (nested) list with an indent + bullet/number prefix.
                reportlab ListFlowable is finicky with mixed content; a prefixed,
                indented paragraph per item is robust and reads correctly."""
                bullet = "•"
                for n, item in enumerate(blk["items"], 1):
                    prefix = f"{n}." if blk.get("ordered") else bullet
                    first = True
                    for sub in item:
                        st = sub["type"]
                        if st == "paragraph" and first:
                            sty = _RLParaStyle(f"LI{depth}", parent=styles["Normal"],
                                               leftIndent=14 + depth * 16, firstLineIndent=-12,
                                               spaceBefore=1, spaceAfter=1)
                            story.append(Paragraph(f"{prefix}&nbsp;&nbsp;" + _pdf_runs(sub["runs"]), sty))
                            first = False
                        elif st == "list":
                            _pdf_render_list(sub, depth + 1)
                        elif st == "paragraph":
                            sty = _RLParaStyle(f"LIc{depth}", parent=styles["Normal"],
                                               leftIndent=14 + depth * 16 + 12)
                            story.append(Paragraph(_pdf_runs(sub["runs"]), sty))

            def _pdf_render_quote(blk):
                qsty = _RLParaStyle("Quote", parent=styles["Normal"], leftIndent=18,
                                    textColor=_RLHex("#555555"), borderColor=_RLHex(_pdf_rule_hex),
                                    borderWidth=0, leading=15)
                for qb in blk["blocks"]:
                    if qb.get("type") == "paragraph":
                        story.append(Paragraph("<i>" + _pdf_runs(qb["runs"]) + "</i>", qsty))
                    elif qb.get("type") == "list":
                        _pdf_render_list(qb)

            def _pdf_render_code(blk):
                csty = _RLParaStyle("Code", parent=styles["Normal"], fontName="Courier",
                                    fontSize=max(8, int(_style["sizes"].get("body", 11)) - 1),
                                    leading=12, leftIndent=8, backColor=_RLHex("#F2F2F2"),
                                    borderPadding=4)
                # Preserve leading indentation per line (reportlab collapses
                # runs of spaces in its mini-HTML, so swap them for &nbsp;).
                _clines = []
                for _cl in _esc(blk["text"]).split("\n"):
                    _lead = len(_cl) - len(_cl.lstrip(" "))
                    _clines.append("&nbsp;" * _lead + _cl.lstrip(" "))
                safe = "<br/>".join(_clines)
                story.append(Paragraph(safe or "&nbsp;", csty))
                story.append(Spacer(1, 8))

            for blk in _markdown_to_blocks(_pdf_body):
                bt = blk["type"]
                if bt == "paragraph":
                    txt = _runs_plain(blk["runs"])
                    kpi_lines = [_kpi_match(l) for l in txt.split("\n")]
                    if kpi_lines and all(kpi_lines):
                        _pdf_kpis.extend(kpi_lines)
                        continue
                    _pdf_flush_kpis()
                    _im = _MD_IMAGE_RE.match(txt.strip())
                    if _im:
                        _ip = _resolve_doc_image(_im.group(2), _doc_dir)
                        if _ip and not _ip.lower().endswith(".svg"):
                            try:
                                img = _RLImage(_ip)
                                maxw = 6.5 * _rl_inch
                                if img.drawWidth > maxw:
                                    ratio = maxw / img.drawWidth
                                    img.drawWidth = maxw; img.drawHeight *= ratio
                                story.append(img); story.append(Spacer(1, 12))
                            except Exception:
                                story.append(Paragraph(f"[Bild: {_im.group(1) or _im.group(2)}]", styles["Normal"]))
                        else:
                            story.append(Paragraph(
                                f"[Bild {_im.group(2)} — für PDF bitte als PNG rendern (render_diagram format=png)]",
                                styles["Normal"]))
                        continue
                    story.append(Paragraph(_pdf_runs(blk["runs"]) or "&nbsp;", styles["Normal"]))
                    continue
                _pdf_flush_kpis()
                if bt == "heading":
                    level = blk["level"]
                    style_name = f"Heading{min(level, 6)}"
                    if style_name not in styles:
                        style_name = "Heading1"
                    # Version-history section starts on a fresh page (table stays
                    # whole) and is kept OUT of the TOC (it's a document appendix).
                    _is_vhist = _is_version_history_heading(_runs_plain(blk["runs"]))
                    if _is_vhist:
                        story.append(_RLPageBreak())
                    htext = _clean_heading_text(_runs_plain(blk["runs"]), _pdf_strip_emoji)
                    _plain = re.sub(r'[*`]+', '', htext).strip()
                    if _pdf_toc_active and level <= 3 and _plain and not _is_vhist:
                        _pdf_hkey[0] += 1
                        _key = f"h{_pdf_hkey[0]}"
                        _para = Paragraph(f'<a name="{_key}"/>' + _esc(htext), styles[style_name])
                        _para._toc_text = _plain
                        _para._toc_level = level - 1
                        _para._toc_key = _key
                        story.append(_para)
                    else:
                        story.append(Paragraph(_esc(htext), styles[style_name]))
                elif bt == "hr":
                    story.append(_RLHRFlowable(width="100%", thickness=0.75,
                                               color=_RLHex(_pdf_rule_hex), spaceBefore=6, spaceAfter=10))
                elif bt == "list":
                    _pdf_render_list(blk)
                    story.append(Spacer(1, 6))
                elif bt == "quote":
                    _pdf_render_quote(blk)
                elif bt == "code":
                    # Mermaid diagram → render to PNG + embed; else a code block.
                    _merm_png = (_render_mermaid_block_to_file(blk["text"], _style, _doc_dir)
                                 if _is_mermaid_block(blk, _style) else None)
                    _mp = _resolve_doc_image(_merm_png, _doc_dir) if _merm_png else None
                    if _mp and not _mp.lower().endswith(".svg"):
                        try:
                            img = _RLImage(_mp)
                            # Fit within the printable area by BOTH width and height
                            # (a wide gantt/flowchart otherwise overflows the frame
                            # → reportlab "too large on page" error).
                            maxw = 6.5 * _rl_inch
                            maxh = (_psize[1] - 2 * _marg) - 0.6 * _rl_inch
                            _ratio = min(maxw / img.drawWidth, maxh / img.drawHeight, 1.0)
                            img.drawWidth *= _ratio; img.drawHeight *= _ratio
                            story.append(img); story.append(Spacer(1, 12))
                        except Exception:
                            _pdf_render_code(blk)
                    else:
                        _pdf_render_code(blk)
                elif bt == "table":
                    _t = _pdf_table_flowable(blk)
                    if _t is not None:
                        story.append(_t); story.append(Spacer(1, 12))
            _pdf_flush_kpis()
            _pdf_cb = _make_pdf_hdrftr_cb(_style, _psize, _marg, _rl_inch, _RLHex, content)
            # multiBuild (2 passes) resolves the ToC page numbers; plain build
            # otherwise. The header/footer canvas callback runs on every page.
            _build = doc_pdf.multiBuild if _pdf_toc_active else doc_pdf.build
            if _pdf_cb:
                _build(story, onFirstPage=_pdf_cb, onLaterPages=_pdf_cb)
            else:
                _build(story)

        elif ext in (".html", ".htm"):
            # Two content modes, auto-detected:
            #  • RAW HTML (the model passed a full <html> document) → write it
            #    through, only inlining local <img src> as base64 + injecting the
            #    preset's header/footer/logo bands so the file is portable + on-
            #    brand. (Treating raw HTML as markdown would escape every tag and
            #    show the source as text — the v9.152.0 bug.)
            #  • MARKDOWN → render to a self-contained styled HTML document
            #    applying the full preset (fonts/colors/sizes/tables + chrome).
            if _editorial and not _looks_like_html(content):
                # Editorial report layout (style='report'/'editorial') — same
                # renderer as Deep Research + Studio. Needs MARKDOWN (raw HTML
                # can't be re-flowed into this layout → falls through to the
                # preset path below). Title = first '# H1' in the content, else
                # the filename stem.
                from engine import report_html
                _m = re.search(r"^\s*#\s+(.+)$", content, re.MULTILINE)
                _title = (_m.group(1).strip() if _m
                          else os.path.splitext(os.path.basename(path))[0].replace("_", " "))
                html_doc = report_html.render_report_html(content, _title, category="report", doc_dir=_doc_dir)
            elif _looks_like_html(content):
                html_doc = _finalize_raw_html(content, _style, _doc_dir)
            else:
                html_doc = _render_markdown_html(content, _style, _doc_dir)
            with open(path, "w", encoding="utf-8") as f:
                f.write(html_doc)

        else:
            return _err(f"write_document: unsupported format '{ext}'. Supported: .docx, .xlsx, .pptx, .pdf, .html")

        size = os.path.getsize(path)
        agent = get_request_context().current_agent or _brain._current_agent
        _brain._after_file_write(path, "created", agent.agent_id if agent else "main")
        _res = {"path": path, "size": size, "format": ext.lstrip("."), "status": "written"}
        if _ref_note:
            _res["style"] = _ref_note
        return _ok(_res)
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"write_document: {e}")


def _write_md_table_to_sheet(ws, md_text: str) -> None:
    """Helper: parse markdown table text and write rows to an openpyxl worksheet."""
    row_idx = 1
    for line in md_text.split("\n"):
        line = line.strip()
        if not line or not line.startswith("|"):
            continue
        if re.match(r'^\|[\s\-:|]+\|$', line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        for ci, val in enumerate(cells, 1):
            try:
                ws.cell(row=row_idx, column=ci, value=float(val) if "." in val else int(val))
            except (ValueError, TypeError):
                ws.cell(row=row_idx, column=ci, value=val)
        row_idx += 1


def tool_edit_document(args: dict) -> str:
    """Targeted edits to existing documents."""
    import brain as _brain
    path = args.get("path", "")
    action = args.get("action", "")
    params = args.get("params", {})
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        if not os.path.exists(path):
            return _err(f"File not found: {path}")
        ext = os.path.splitext(path)[1].lower()

        if ext == ".docx":
            if action != "replace_text":
                return _err(f"edit_document: unsupported action '{action}' for DOCX. Use: replace_text")
            try:
                import docx
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            old_text = params.get("old_text", "")
            new_text = params.get("new_text", "")
            if not old_text:
                return _err("edit_document: 'old_text' required in params")
            doc = docx.Document(path)
            count = 0
            for para in doc.paragraphs:
                if old_text in para.text:
                    full_text = para.text
                    new_full = full_text.replace(old_text, new_text)
                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = new_full
                    else:
                        para.add_run(new_full)
                    count += 1
            doc.save(path)
            agent = get_request_context().current_agent or _brain._current_agent
            _brain._after_file_write(path, "modified", agent.agent_id if agent else "main")
            return _ok({"path": path, "action": action, "replacements": count, "status": "edited"})

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
            except ImportError:
                return _err("Install openpyxl: pip3 install openpyxl")
            wb = openpyxl.load_workbook(path)

            if action == "update_cell":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                cell_ref = params.get("cell", "")
                value = params.get("value", "")
                if not cell_ref:
                    return _err("edit_document: 'cell' required (e.g. 'A1')")
                if sheet_name not in wb.sheetnames:
                    return _err(f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}")
                ws = wb[sheet_name]
                try:
                    ws[cell_ref] = float(value) if isinstance(value, str) and value.replace(".", "", 1).replace("-", "", 1).isdigit() else value
                except Exception:
                    ws[cell_ref] = value
                wb.save(path)
                agent = get_request_context().current_agent or _brain._current_agent
                _brain._after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "cell": cell_ref, "sheet": sheet_name, "status": "edited"})

            elif action == "add_row":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                values = params.get("values", [])
                if not values:
                    return _err("edit_document: 'values' array required")
                if sheet_name not in wb.sheetnames:
                    return _err(f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}")
                ws = wb[sheet_name]
                ws.append(values)
                wb.save(path)
                agent = get_request_context().current_agent or _brain._current_agent
                _brain._after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "sheet": sheet_name, "row_added": len(values), "status": "edited"})

            else:
                return _err(f"edit_document: unsupported action '{action}' for XLSX. Use: update_cell, add_row")

        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return _err("Install python-pptx: pip3 install python-pptx")
            prs = Presentation(path)

            if action == "update_slide":
                slide_index = int(params.get("slide_index", 1))
                if slide_index < 1 or slide_index > len(prs.slides):
                    return _err(f"Slide index {slide_index} out of range (1-{len(prs.slides)})")
                slide = prs.slides[slide_index - 1]
                title = params.get("title")
                body = params.get("body")
                if title and slide.shapes.title:
                    slide.shapes.title.text = title
                if body:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            shape.text_frame.text = body
                            break
                prs.save(path)
                agent = get_request_context().current_agent or _brain._current_agent
                _brain._after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "slide_index": slide_index, "status": "edited"})

            elif action == "add_slide":
                title = params.get("title", "New Slide")
                body = params.get("body", "")
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                if body:
                    slide.placeholders[1].text = body
                prs.save(path)
                agent = get_request_context().current_agent or _brain._current_agent
                _brain._after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "slide_count": len(prs.slides), "status": "edited"})

            else:
                return _err(f"edit_document: unsupported action '{action}' for PPTX. Use: update_slide, add_slide")

        else:
            return _err(f"edit_document: unsupported format '{ext}'. Supported: .docx, .xlsx, .xls, .pptx")
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"edit_document: {e}")


def tool_list_directory(args: dict) -> str:
    import brain as _brain
    node_result = _brain._route_to_node("list_directory", args)
    if node_result is not None:
        return node_result
    path = args.get("path", ".")
    pattern = args.get("pattern")
    recursive = args.get("recursive", False)
    try:
        path = _resolve_under_cwd(path)

        if pattern:
            if recursive or "**" in pattern:
                full_pattern = os.path.join(path, pattern)
                entries = globmod.glob(full_pattern, recursive=True)
            else:
                full_pattern = os.path.join(path, pattern)
                entries = globmod.glob(full_pattern)
        elif recursive:
            entries = []
            for root, dirs, files in os.walk(path):
                # Skip hidden dirs
                dirs[:] = [d for d in dirs if not d.startswith(".")]
                for f in files:
                    if not f.startswith("."):
                        entries.append(os.path.join(root, f))
        else:
            entries = [os.path.join(path, e) for e in os.listdir(path)]

        results = []
        for entry in sorted(entries)[:500]:
            try:
                st = os.stat(entry)
                is_dir = os.path.isdir(entry)
                results.append({
                    "name": os.path.relpath(entry, path),
                    "type": "directory" if is_dir else "file",
                    "size": st.st_size if not is_dir else None,
                })
            except OSError:
                results.append({"name": os.path.relpath(entry, path), "type": "unknown"})

        return _ok({"path": path, "count": len(results), "entries": results})
    except Exception as e:
        return _err(f"list_directory: {e}")


def tool_search_files(args: dict) -> str:
    pattern = args.get("pattern", "")
    path = args.get("path", ".")
    file_glob = args.get("glob")
    case_insensitive = args.get("case_insensitive", False)
    max_results = args.get("max_results", 50)
    try:
        path = _resolve_under_cwd(path)

        flags = re.IGNORECASE if case_insensitive else 0
        regex = re.compile(pattern, flags)

        matches = []
        files_searched = 0

        if os.path.isfile(path):
            file_list = [path]
        else:
            file_list = []
            for root, dirs, files in os.walk(path):
                dirs[:] = [d for d in dirs if not d.startswith(".") and d not in ("node_modules", "__pycache__", ".git")]
                for f in files:
                    if f.startswith("."):
                        continue
                    fp = os.path.join(root, f)
                    if file_glob and not fnmatch.fnmatch(f, file_glob):
                        continue
                    file_list.append(fp)

        for fp in file_list:
            if len(matches) >= max_results:
                break
            files_searched += 1
            try:
                with open(fp, "r", errors="replace") as fh:
                    for lineno, line in enumerate(fh, 1):
                        if regex.search(line):
                            matches.append({
                                "file": os.path.relpath(fp, path) if not os.path.isfile(path) else fp,
                                "line": lineno,
                                "text": line.rstrip()[:500],
                            })
                            if len(matches) >= max_results:
                                break
            except (OSError, UnicodeDecodeError):
                continue

        return _ok({"pattern": pattern, "path": path, "files_searched": files_searched,
                     "match_count": len(matches), "matches": matches})
    except re.error as e:
        return _err(f"search_files: invalid regex: {e}")
    except Exception as e:
        return _err(f"search_files: {e}")


def _strip_ansi(text: str) -> str:
    """Remove all ANSI escape sequences, control chars, and normalize whitespace."""
    text = re.sub(r"\033\[[0-9;]*[a-zA-Z]", "", text)
    text = re.sub(r"\033\[\?[0-9;]*[a-zA-Z]", "", text)
    text = re.sub(r"\033\([A-Z]", "", text)
    text = re.sub(r"\033][^\a]*\a", "", text)  # OSC sequences
    text = re.sub(r"\r", "", text)  # Carriage returns
    text = text.replace("\t", "    ")  # Expand tabs
    return text


def _build_shell_command(command: str) -> tuple[list | str, bool]:
    """Build the shell invocation based on execute_command config.

    Returns (cmd, shell_flag) for subprocess.Popen.
    If login_shell is True, wraps the command in a login shell invocation
    so that ~/.zprofile, ~/.zshrc etc. are sourced (giving full PATH).
    """
    import brain as _brain
    _exec_cfg = _brain.get_tool_config().get("execute_command", {})
    use_login_shell = _exec_cfg.get("login_shell", True)
    if use_login_shell:
        shell_path = _exec_cfg.get("shell_path", "") or os.environ.get("SHELL", "/bin/zsh")
        return [shell_path, "-l", "-c", command], False
    return command, True


def _streaming_execute_command(command: str, timeout: int, cwd: str | None,
                               event_callback, tool_use_id: str) -> str:
    """Execute command with streaming output via event_callback."""
    import brain as _brain
    env = os.environ.copy()
    env["TERM"] = "dumb"
    env["NO_COLOR"] = "1"
    env["PAGER"] = "cat"
    env["COLUMNS"] = "200"
    env["LINES"] = "50"

    shell_cmd, shell_flag = _build_shell_command(command)
    proc = subprocess.Popen(
        shell_cmd, shell=shell_flag, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        stdin=subprocess.DEVNULL, cwd=cwd, env=env,
        start_new_session=True,
    )
    _proc_key = register_tool_process(proc)  # per-tool kill registration
    output_lines = []
    import io
    deadline = time.time() + timeout
    try:
        # Read stdout line by line, emitting events
        for raw_line in iter(proc.stdout.readline, b''):
            if time.time() > deadline:
                import signal as sig
                try:
                    os.killpg(proc.pid, sig.SIGKILL)
                except OSError:
                    proc.kill()
                proc.wait(timeout=5)
                line_text = _strip_ansi(raw_line.decode("utf-8", errors="replace"))
                output_lines.append(line_text)
                output = "".join(output_lines)
                if len(output) > 50000:
                    output = output[:50000] + "\n... (truncated)"
                return _err(f"execute_command: timed out after {timeout}s\n{output}")
            line_text = _strip_ansi(raw_line.decode("utf-8", errors="replace"))
            output_lines.append(line_text)
            if event_callback:
                event_callback("tool_output", {
                    "tool_use_id": tool_use_id,
                    "line": line_text.rstrip("\n"),
                })
        proc.stdout.close()
        proc.wait(timeout=5)
    except subprocess.TimeoutExpired:
        import signal as sig
        try:
            os.killpg(proc.pid, sig.SIGKILL)
        except OSError:
            proc.kill()
        proc.wait(timeout=5)
    finally:
        unregister_tool_process(_proc_key)

    # External per-tool kill mid-stream → process group SIGKILLed.
    if proc.returncode is not None and proc.returncode < 0:
        import signal as sig
        if proc.returncode == -sig.SIGKILL:
            return _err("execute_command: cancelled by user (process killed).")

    output = "".join(output_lines)
    stderr_data = proc.stderr.read() if proc.stderr else b""
    proc.stderr.close()
    if stderr_data:
        err_text = _strip_ansi(stderr_data.decode("utf-8", errors="replace"))
        output += ("\n--- stderr ---\n" + err_text) if output else err_text
    if len(output) > 50000:
        output = output[:50000] + "\n... (truncated)"
    output = _brain._gdpr_anon_tool_text(output, "execute_command:stdout")
    return _ok({"command": command, "exit_code": proc.returncode, "output": output})


def tool_execute_command(args: dict) -> str:
    import brain as _brain
    command = args.get("command", "")
    cwd = args.get("cwd")
    # Read default timeout from tools_config (per-call timeout still overrides)
    _exec_cfg = _brain.get_tool_config().get("execute_command", {})
    default_timeout = _exec_cfg.get("timeout", 15)
    timeout = args.get("timeout", default_timeout)
    # Check banned commands
    banned = _exec_cfg.get("banned_commands", [])
    for b in banned:
        if b and b in command:
            return _err(f"execute_command: command contains banned pattern '{b}'")
    try:
        if cwd:
            cwd = os.path.expanduser(cwd)
        else:
            # Default cwd = the session write root: the code-mode working_dir if
            # set, else the session artifact folder (matches python_exec) so
            # outputs land in the Artifacts panel without the model having to
            # guess the path. Falls back to Brain's cwd when there's no session.
            _wr, _ = _resolve_artifact_dir()
            if _wr:
                cwd = _wr
                try:
                    os.makedirs(cwd, exist_ok=True)
                except OSError:
                    pass

        # Snapshot artifact folder if cwd lands inside it — files appearing
        # post-exec auto-register as artifacts (mirrors tool_python_exec). We
        # only snapshot when cwd IS the artifact folder; commands run outside
        # it shouldn't pollute the artifact panel.
        _artifact_cwd = None
        _pre_files: dict[str, tuple[float, int]] = {}
        _session_id = get_request_context().current_session_id
        _agent = get_request_context().current_agent or _brain._current_agent
        if cwd and _session_id and _agent:
            _expected = os.path.join(_brain.AGENTS_DIR, _agent.agent_id, "artifacts",
                                     _get_artifact_session_folder(_session_id))
            if os.path.realpath(cwd) == os.path.realpath(_expected):
                _artifact_cwd = _expected
                _pre_files = _snapshot_dir(_artifact_cwd)

        # Use streaming version if event_callback is available
        ecb = get_request_context().event_callback
        tuid = get_request_context().tool_use_id
        if ecb and tuid:
            _res = _streaming_execute_command(command, timeout, cwd, ecb, tuid)
            _register_new_artifacts(_artifact_cwd, _pre_files, _agent)
            _stray = _stray_write_warning(command, _artifact_cwd)
            if _stray:
                _res = _append_to_tool_result(_res, _stray)
            return _res

        # Force non-interactive environment
        env = os.environ.copy()
        env["TERM"] = "dumb"
        env["NO_COLOR"] = "1"
        env["PAGER"] = "cat"
        env["COLUMNS"] = "200"
        env["LINES"] = "50"

        shell_cmd, shell_flag = _build_shell_command(command)
        try:
            from engine.context import report_tool_progress
            report_tool_progress(phase="Läuft", note=(command or "").split()[0] if command else "Befehl")
        except Exception:
            pass
        proc = subprocess.Popen(
            shell_cmd, shell=shell_flag, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            stdin=subprocess.DEVNULL, cwd=cwd, env=env,
            start_new_session=True,  # own process group so we can kill the tree
        )
        _proc_key = register_tool_process(proc)  # per-tool kill registration
        try:
            stdout, stderr = proc.communicate(timeout=timeout)
        except subprocess.TimeoutExpired:
            # Kill the entire process group
            import signal as sig
            try:
                os.killpg(proc.pid, sig.SIGKILL)
            except OSError:
                proc.kill()
            stdout, stderr = proc.communicate(timeout=5)
            output = _strip_ansi(stdout.decode("utf-8", errors="replace"))
            if stderr:
                output += "\n--- stderr ---\n" + _strip_ansi(stderr.decode("utf-8", errors="replace"))
            if len(output) > 50000:
                output = output[:50000] + "\n... (truncated)"
            return _err(f"execute_command: timed out after {timeout}s (partial output below). Use non-interactive commands, e.g. 'top -l 1' not 'top'.\n{output}")
        finally:
            unregister_tool_process(_proc_key)
        # External per-tool kill → process group SIGKILLed mid-run.
        if proc.returncode is not None and proc.returncode < 0:
            import signal as sig
            if proc.returncode == -sig.SIGKILL:
                return _err("execute_command: cancelled by user (process killed).")

        output = _strip_ansi(stdout.decode("utf-8", errors="replace"))
        if stderr:
            err_text = _strip_ansi(stderr.decode("utf-8", errors="replace"))
            output += ("\n--- stderr ---\n" + err_text) if output else err_text
        if len(output) > 50000:
            output = output[:50000] + "\n... (truncated)"
        # Transparent-anonymisation seam (same rationale as tool_python_exec):
        # `cat /tmp/brain-attachments/.../report.docx` or `grep -r foo /tmp`
        # could surface raw PII bytes to the LLM. Wrap stdout through the
        # same mapping that read_document uses.
        output = _brain._gdpr_anon_tool_text(output, "execute_command:stdout")
        _new_artifacts = _register_new_artifacts(_artifact_cwd, _pre_files, _agent)
        result = {"command": command, "exit_code": proc.returncode, "output": output}
        if _new_artifacts:
            result["artifacts"] = _new_artifacts
        _stray = _stray_write_warning(command, _artifact_cwd)
        if _stray:
            result["output"] = (result.get("output") or "") + _stray
        return _ok(result)
    except Exception as e:
        return _err(f"execute_command: {e}")


def _snapshot_dir(dir_path: str) -> dict[str, tuple[float, int]]:
    """Map each file in dir_path to (mtime, size). Used to detect files that a
    python_exec script / execute_command run created OR overwrote in place — a
    plain name-set diff would miss in-place rewrites of pre-existing artifacts."""
    snap: dict[str, tuple[float, int]] = {}
    try:
        with os.scandir(dir_path) as it:
            for e in it:
                try:
                    if e.is_file():
                        st = e.stat()
                        snap[e.name] = (st.st_mtime, st.st_size)
                except OSError:
                    continue
    except OSError:
        pass
    return snap


def _changed_files(dir_path: str, pre: dict[str, tuple[float, int]],
                   exclude: set | None = None) -> list[tuple[str, bool]]:
    """Return [(filename, was_new), ...] for files that are new since `pre` or
    whose (mtime, size) changed (in-place overwrite). `pre` is a _snapshot_dir
    mapping. was_new=False marks an overwrite (caller emits 'modified')."""
    exclude = exclude or set()
    out: list[tuple[str, bool]] = []
    for fname, sig in sorted(_snapshot_dir(dir_path).items()):
        if fname in exclude:
            continue
        prev = pre.get(fname)
        if prev is None:
            out.append((fname, True))
        elif prev != sig:
            out.append((fname, False))
    return out


def _register_new_artifacts(artifact_cwd: str | None,
                            pre_files: dict[str, tuple[float, int]],
                            agent) -> list[str]:
    """Diff artifact_cwd against the pre-exec snapshot; register each created
    OR overwritten file via `_after_file_write` so it shows in the Artifacts
    panel. Returns the registered basenames (for inclusion in the tool result)."""
    import brain as _brain
    if not artifact_cwd or not agent:
        return []
    changed = _changed_files(artifact_cwd, pre_files)
    if not changed:
        return []
    agent_id = agent.agent_id
    created = []
    for fname, was_new in changed:
        fpath = os.path.join(artifact_cwd, fname)
        if os.path.isfile(fpath):
            try:
                _brain._after_file_write(
                    fpath, "created" if was_new else "modified", agent_id)
                created.append(fname)
            except Exception:
                pass
    return created


# Absolute path literals in python_exec code / execute_command commands.
# We don't redirect (the user may legitimately want an absolute write) — we
# only WARN, so the model self-corrects on the next turn. See _stray_write_warning.
_ABS_PATH_RE = re.compile(r"""['"`(\s=](/(?:[\w.\-+ ]+/)*[\w.\-+]+\.[\w]{1,8})\b""")


def _stray_write_warning(text: str, artifact_dir: str | None) -> str:
    """Return a warning string when `text` (python_exec code or a shell command)
    references absolute file paths that point OUTSIDE the session artifact
    folder, OR "" when there's nothing to warn about.

    Rationale: python_exec / execute_command run with cwd = artifact folder and
    only diff THAT folder for new artifacts. If the model's script does
    `doc.save("/Users/.../repo/report.docx")` (an absolute path it invented from
    the old, misleading "Current working directory" line), the file lands outside
    the artifact folder, never appears in the post-exec diff, and is invisible to
    the Artifacts panel — while the model believes it succeeded. We surface that
    back into the tool result so the model re-saves with a relative filename.

    Non-invasive by design (user chose "warn, don't redirect"): an absolute path
    the user explicitly asked for is left alone; the warning just informs.
    """
    if not text or not artifact_dir:
        return ""
    try:
        art_real = os.path.realpath(artifact_dir)
    except OSError:
        return ""
    stray: list[str] = []
    seen: set[str] = set()
    for m in _ABS_PATH_RE.finditer(text):
        p = m.group(1)
        if p in seen:
            continue
        seen.add(p)
        # Inside the artifact folder? fine. /tmp, /dev, /proc etc. are not
        # user-facing artifact writes — skip them (avoid noise on temp files).
        if p.startswith(("/tmp/", "/private/tmp/", "/var/", "/dev/", "/proc/",
                         "/usr/", "/etc/", "/System/", "/Library/")):
            continue
        try:
            if os.path.realpath(p).startswith(art_real):
                continue
        except OSError:
            pass
        stray.append(p)
    if not stray:
        return ""
    paths = ", ".join(stray[:5])
    return (
        f"\n\n⚠️ WARNING: this wrote to an absolute path OUTSIDE your session "
        f"artifact folder ({paths}). Files written there do NOT appear in the "
        f"Artifacts panel and the user cannot see or download them. Re-save the "
        f"output using a RELATIVE filename (just the filename, e.g. "
        f"`report.docx`) so it lands in your artifact folder — unless the user "
        f"explicitly asked for that absolute path."
    )


def _resolve_artifact_dir():
    """Return (write_root, agent_id) for the current chat/scheduled session, or
    (None, agent_id) when there's no session context (CLI one-shots / warmup).
    The write_root is the ONLY folder a file-writing tool may write into.

    Code Mode: when the request context carries a `working_dir` (set by
    apply_domain_context for a code_mode project), THAT is the write root — the
    agent reads/edits/creates files directly in the user's working directory
    instead of the session artifact folder."""
    import brain as _brain
    _wd = get_request_context().working_dir
    if _wd and os.path.isdir(_wd):
        agent = get_request_context().current_agent or _brain._current_agent
        return _wd, (agent.agent_id if agent else "main")
    session_id = get_request_context().current_session_id
    agent = get_request_context().current_agent or _brain._current_agent
    agent_id = agent.agent_id if agent else "main"
    if session_id and agent:
        folder = _get_artifact_session_folder(session_id)
        return os.path.join(_brain.AGENTS_DIR, agent_id, "artifacts", folder), agent_id
    return None, agent_id


def _path_within(child: str, parent: str) -> bool:
    """True if `child` is `parent` or lives underneath it (boundary-safe — uses
    realpath + a separator guard so /x/artifacts-evil ≠ inside /x/artifacts)."""
    try:
        c = os.path.realpath(child)
        p = os.path.realpath(parent)
    except OSError:
        return False
    return c == p or c.startswith(p + os.sep)


def _enforce_artifact_path(raw_path: str, tool: str):
    """Resolve a model-supplied write path to its FINAL location and require it to
    be inside the session artifact folder. Returns (final_path, error_or_None):

      • relative input → joined into the artifact dir (the intended default).
      • absolute (or relative that escapes via ..) input → must already resolve
        inside the artifact dir, else an error is returned and the write is
        REFUSED (no file is written). Covers both absolute AND relative paths —
        the check is on the RESOLVED path, not the literal.
      • no session context (CLI/warmup) → no artifact dir to enforce against;
        fall back to abspath, no restriction.

    This is the hard block: writing anywhere but the artifact folder is not
    allowed, so the user always sees/downloads what the agent produced.
    """
    p = os.path.expanduser(raw_path or "")
    artifact_dir, _ = _resolve_artifact_dir()
    if not artifact_dir:
        # No chat/scheduled session — nothing to scope to; keep prior behavior.
        return (p if os.path.isabs(p) else os.path.abspath(p)), None
    os.makedirs(artifact_dir, exist_ok=True)
    final = p if os.path.isabs(p) else os.path.join(artifact_dir, p)
    if not _path_within(final, artifact_dir):
        return None, _err(
            f"{tool}: writing outside your session artifact folder is not allowed. "
            f"The path '{raw_path}' resolves outside it. Use a RELATIVE filename "
            f"(just `name.ext`, e.g. `report.html`) so the file lands in your "
            f"artifact folder — do NOT pass an absolute path or one with '..'.")
    return final, None


def _append_to_tool_result(res_json: str, suffix: str) -> str:
    """Append `suffix` to the `output` field of a JSON tool result string
    (as produced by `_ok`). Falls back to a raw string append if the result
    isn't the expected JSON-object-with-output shape."""
    if not suffix:
        return res_json
    try:
        obj = json.loads(res_json)
        if isinstance(obj, dict):
            obj["output"] = (obj.get("output") or "") + suffix
            return json.dumps(obj, ensure_ascii=False)
    except (ValueError, TypeError):
        pass
    return res_json + suffix


def tool_python_exec(args: dict) -> str:
    """Execute Python code in an isolated subprocess with artifact folder as cwd."""
    import brain as _brain
    code = args.get("code", "")
    if not code.strip():
        return _err("python_exec: no code provided")

    _cfg = _brain.get_tool_config().get("python_exec", {})
    timeout = args.get("timeout", _cfg.get("timeout", 30))
    max_output = _cfg.get("max_output_chars", 50000)
    venv_path = _cfg.get("venv_path", "")

    # Working dir = the session write root: the code-mode working_dir if set,
    # else the session artifact folder (so files written by code become
    # artifacts). Falls back to a temp dir when there's no session context.
    _wr, _ = _resolve_artifact_dir()
    if _wr:
        work_dir = _wr
    else:
        import tempfile
        work_dir = os.path.join(tempfile.gettempdir(), "brain-pyexec")
    os.makedirs(work_dir, exist_ok=True)

    # Snapshot existing files before execution. Capture (mtime, size) per file
    # — not just names — so the post-exec diff catches files the script
    # OVERWROTE in place, not only freshly-created ones. (A re-run that rewrites
    # an already-existing artifact, e.g. swapping image URLs in a prior .html,
    # must register as a new version, not vanish from the Artifacts panel.)
    pre_files = _snapshot_dir(work_dir)

    # Save script as a numbered artifact (persisted for reuse)
    counter = 1
    while os.path.exists(os.path.join(work_dir, f"script_{counter}.py")):
        counter += 1
    script_name = f"script_{counter}.py"
    script_path = os.path.join(work_dir, script_name)
    with open(script_path, "w") as f:
        f.write(code)
    # `agent` and `session_id` were both referenced later in this function but
    # never defined — a NameError on EVERY python_exec call (caught by the
    # characterization test; `agent` at the _after_file_write below, `session_id`
    # at the stray-write gate further down). Resolve both from the request
    # context the same way every other tool here does.
    agent = get_request_context().current_agent or getattr(_brain, "_current_agent", None)
    agent_id = (agent.agent_id if agent else "main")
    session_id = get_request_context().current_session_id or ""
    _brain._after_file_write(script_path, "created", agent_id)

    env = os.environ.copy()
    env["PYTHONDONTWRITEBYTECODE"] = "1"
    env["PYTHONUNBUFFERED"] = "1"
    if venv_path and os.path.isdir(venv_path):
        env["PYTHONPATH"] = venv_path + ((":" + env.get("PYTHONPATH", "")) if env.get("PYTHONPATH") else "")

    _proc_key = None
    try:
        try:
            from engine.context import report_tool_progress
            report_tool_progress(phase="Läuft", note="Python-Skript")
        except Exception:
            pass
        proc = subprocess.Popen(
            [sys.executable, script_path],
            stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            stdin=subprocess.DEVNULL, cwd=work_dir, env=env,
            start_new_session=True,
        )
        # Register for per-tool kill: a cancel SIGKILLs this process group and
        # communicate() returns promptly with a negative returncode (the signal).
        _proc_key = register_tool_process(proc)
        try:
            stdout, stderr = proc.communicate(timeout=timeout)
        except subprocess.TimeoutExpired:
            import signal as sig
            try:
                os.killpg(proc.pid, sig.SIGKILL)
            except OSError:
                proc.kill()
            stdout, stderr = proc.communicate(timeout=5)
            output = stdout.decode("utf-8", errors="replace")
            if stderr:
                output += "\n--- stderr ---\n" + stderr.decode("utf-8", errors="replace")
            if len(output) > max_output:
                output = output[:max_output] + "\n... (truncated)"
            return _err(f"python_exec: timed out after {timeout}s\n{output}")
        finally:
            unregister_tool_process(_proc_key)
        # External per-tool kill: communicate() returned because the process
        # group was SIGKILLed (negative returncode = -signal). Surface a clear
        # cancelled result so the model knows this tool was aborted on purpose.
        if proc.returncode is not None and proc.returncode < 0:
            import signal as sig
            if proc.returncode == -sig.SIGKILL:
                return _err("python_exec: cancelled by user (process killed).")

        output = stdout.decode("utf-8", errors="replace")
        if stderr:
            err_text = stderr.decode("utf-8", errors="replace")
            # Clean tracebacks: replace script path, remove echoed source lines
            err_text = err_text.replace(script_path, "<python_exec>")
            err_lines = err_text.splitlines()
            cleaned = []
            skip_next = False
            for line in err_lines:
                if '<python_exec>' in line:
                    cleaned.append(line)
                    skip_next = True
                elif skip_next and line.startswith('    ') and not line.strip().startswith('File '):
                    skip_next = False
                    continue
                else:
                    skip_next = False
                    cleaned.append(line)
            err_text = "\n".join(cleaned)
            output += ("\n--- stderr ---\n" + err_text) if output else err_text
        if len(output) > max_output:
            output = output[:max_output] + "\n... (truncated)"

        # Transparent-anonymisation seam: any PII the script printed via
        # stdout reaches the LLM through this `output` field. Wrap it
        # through the same `_gdpr_anon_tool_text` helper that read_document
        # / read_file use — closes the python_exec end-around where the
        # model bypasses read_document by `open(path, 'rb')` and printing
        # the bytes back. Note: artifacts written to disk get their own
        # post-write deanonymisation pass via `_after_file_write` →
        # `make_gdpr_after_file_write_cb`, so a script that writes a file
        # with restored values still gets surfaced to the user correctly.
        output = _brain._gdpr_anon_tool_text(output, "python_exec:stdout")

        result = {"exit_code": proc.returncode, "output": output, "script": script_name}

        # Warn when the script saved to an absolute path outside the artifact
        # folder — such files never show up in the post-exec diff below and are
        # invisible to the Artifacts panel (the recurring "model wrote to repo
        # root" failure). Non-invasive: we warn, we don't move the file.
        _stray = _stray_write_warning(code, work_dir if (session_id and agent) else None)
        if _stray:
            result["output"] = (result.get("output") or "") + _stray

        # Register files the script created OR overwrote in place as artifacts.
        # _changed_files compares the post-exec snapshot against pre_files by
        # (mtime, size), so a re-run that rewrites an existing artifact bumps
        # its version instead of being skipped (the script itself is excluded).
        changed = _changed_files(work_dir, pre_files, exclude={script_name})
        if changed and agent:
            created = []
            for fname, was_new in changed:
                fpath = os.path.join(work_dir, fname)
                if os.path.isfile(fpath):
                    _brain._after_file_write(
                        fpath, "created" if was_new else "modified", agent_id)
                    created.append(fname)
            if created:
                result["artifacts"] = created

        # Always save stdout as an artifact when the script didn't write any files
        if output and proc.returncode == 0 and not changed and agent:
            try:
                artifact_path = os.path.join(work_dir, "output.txt")
                counter = 1
                while os.path.exists(artifact_path):
                    artifact_path = os.path.join(work_dir, f"output_{counter}.txt")
                    counter += 1
                with open(artifact_path, "w") as af:
                    af.write(output)
                _brain._after_file_write(artifact_path, "created", agent_id)
                result["artifacts"] = [os.path.basename(artifact_path)]
                # For large outputs, replace inline data with a reference so the
                # summariser doesn't ingest a megabyte of stdout. Small outputs
                # stay inline so the summariser can describe them meaningfully.
                if len(output) > 1000:
                    lines = output.splitlines()
                    result["output"] = (
                        f"Output saved as artifact {os.path.basename(artifact_path)} "
                        f"({len(lines)} lines, {len(output):,} chars). "
                        f"The user can view it directly. Summarize what was computed, do NOT repeat the data."
                    )
            except Exception:
                pass

        return _ok(result)
    except Exception as e:
        return _err(f"python_exec: {e}")
