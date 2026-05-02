# Extracted from claude_cli.py — MemoryStore, project/note/ingest management
# Cross-module deps (e.g. _err, _ok, _thread_local, AGENTS_DIR, _run_delegate,
# _delegate_api_key, _models_config, AgentConfig, _current_agent,
# _get_agent_team_info, _scheduler, _sched_conn, list_agents,
# resolve_provider_for_model, send_message_with_fallback, gdpr_pick_model_for_background,
# GDPRBlockedError, _get_reranker_model) are resolved via claude_cli namespace.

import os
import re
import sys
import json
import time
import shutil
import fnmatch
import hashlib
import sqlite3
import datetime
import threading
import subprocess
import urllib.request
import urllib.error


def _extract_json_from_llm(text: str, expect_array: bool = False):
    """Robustly extract JSON object or array from LLM response text.

    Handles markdown code fences, nested objects, surrounding text.
    Returns parsed dict/list or None on failure.
    """
    if not text:
        return None
    # Strip markdown code fences first
    stripped = re.sub(r'```(?:json)?\s*', '', text)
    stripped = stripped.replace('```', '')

    # Try parsing the entire stripped text first
    try:
        parsed = json.loads(stripped.strip())
        if expect_array and isinstance(parsed, list):
            return parsed
        if not expect_array and isinstance(parsed, dict):
            return parsed
    except (json.JSONDecodeError, ValueError):
        pass

    # Use json.JSONDecoder.raw_decode to find the first valid JSON structure
    decoder = json.JSONDecoder()
    target_char = '[' if expect_array else '{'
    for i, ch in enumerate(stripped):
        if ch == target_char:
            try:
                obj, _ = decoder.raw_decode(stripped[i:])
                if expect_array and isinstance(obj, list):
                    return obj
                if not expect_array and isinstance(obj, dict):
                    return obj
            except (json.JSONDecodeError, ValueError):
                continue
    return None

# QMD HTTP MCP daemon endpoint
_QMD_URL = "http://localhost:8181/mcp"
_QMD_HEADERS = {
    "Content-Type": "application/json",
    "Accept": "application/json, text/event-stream",
}

# Shared MCP session ID (set on first successful init)
_qmd_session_id: str | None = None
_qmd_session_lock = threading.RLock()
# Per-collection debounce timers for embedding after writes
_qmd_embed_timers: dict[str, threading.Timer] = {}
_qmd_embed_lock = threading.Lock()

# Files to skip when indexing (not memory files)
_QMD_IGNORE_FILES = {"soul.md", "tools.md"}


def _qmd_rpc(method: str, params: dict | None = None) -> dict | None:
    """Send a JSON-RPC request to the QMD MCP HTTP daemon. Returns result or None on failure."""
    global _qmd_session_id
    payload = {"jsonrpc": "2.0", "id": 1, "method": method, "params": params or {}}
    data = json.dumps(payload).encode()
    headers = dict(_QMD_HEADERS)
    with _qmd_session_lock:
        if _qmd_session_id:
            headers["Mcp-Session-Id"] = _qmd_session_id
    req = urllib.request.Request(_QMD_URL, data=data, headers=headers)
    try:
        with urllib.request.urlopen(req, timeout=10) as resp:
            # Capture session ID from response (under lock)
            sid = resp.headers.get("Mcp-Session-Id")
            if sid:
                with _qmd_session_lock:
                    _qmd_session_id = sid
            body = json.loads(resp.read().decode())
            if "error" in body:
                import logging
                logging.warning("QMD RPC error: %s", body["error"])
                return None
            return body.get("result")
    except Exception as e:
        import logging
        logging.debug("QMD unreachable: %s", e)
        return None


def _qmd_init_session() -> bool:
    """Initialize an MCP session with QMD. Returns True if successful.
    Thread-safe: only one session is created even under concurrent access."""
    with _qmd_session_lock:
        if _qmd_session_id:
            return True  # Another thread already initialized
        result = _qmd_rpc("initialize", {
            "protocolVersion": "2024-11-05",
            "capabilities": {},
            "clientInfo": {"name": "brain-agent", "version": "1.0"},
        })
        return result is not None


def _qmd_ensure_collection(name: str, directory: str):
    """Register a collection with QMD if it doesn't exist (via CLI)."""
    try:
        result = subprocess.run(
            ["qmd", "collection", "show", name],
            capture_output=True, text=True, timeout=5,
        )
        if result.returncode == 0:
            return  # Already exists
        subprocess.run(
            ["qmd", "collection", "add", directory,
             "--name", name, "--pattern", "*.md",
             "--ignore", ",".join(_QMD_IGNORE_FILES)],
            capture_output=True, text=True, timeout=15,
        )
    except Exception as e:
        import logging
        logging.debug("QMD collection setup failed for %s: %s", name, e)


def _qmd_debounced_embed(collection: str):
    """Schedule a debounced qmd update+embed for a collection (2s delay).
    Each collection gets its own timer so concurrent writes to different
    collections don't cancel each other's embed."""
    def _do_embed():
        try:
            subprocess.run(
                ["qmd", "update"], capture_output=True, text=True, timeout=30,
            )
            subprocess.run(
                ["qmd", "embed", "-c", collection],
                capture_output=True, text=True, timeout=60,
            )
        except Exception as e:
            import logging
            logging.debug("QMD embed failed for %s: %s", collection, e)
        finally:
            with _qmd_embed_lock:
                _qmd_embed_timers.pop(collection, None)
    with _qmd_embed_lock:
        old = _qmd_embed_timers.get(collection)
        if old:
            old.cancel()
        timer = threading.Timer(2.0, _do_embed)
        timer.daemon = True
        _qmd_embed_timers[collection] = timer
        timer.start()


def _parse_frontmatter(raw: str) -> tuple[dict, str]:
    """Parse YAML frontmatter from a markdown file. Returns (metadata, body)."""
    # Normalize line endings for cross-platform compatibility
    raw = raw.replace("\r\n", "\n")
    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n(.*)$', raw, re.DOTALL)
    if fm_match:
        fm_text, body = fm_match.groups()
        fm = {}
        for line in fm_text.split("\n"):
            # Only parse top-level keys (skip indented/nested YAML lines)
            if ":" in line and not line.startswith((" ", "\t", "-")):
                k, v = line.split(":", 1)
                fm[k.strip()] = v.strip()
        return fm, body.strip()
    return {}, raw.strip()


def _yaml_escape(value: str) -> str:
    """Escape a string for safe inclusion in YAML frontmatter."""
    if not value:
        return '""'
    # Quote if contains special YAML chars
    if any(c in value for c in (':', '#', '{', '}', '[', ']', ',', '&', '*', '?', '|', '-', '<', '>', '=', '!', '%', '@', '`', '\n')):
        escaped = value.replace('\\', '\\\\').replace('"', '\\"').replace('\n', ' ')
        return f'"{escaped}"'
    return value


class MemoryStore:
    """Per-agent memory store backed by QMD hybrid search and markdown files."""

    _ensured_collections: set[str] = set()
    _ensured_lock = threading.Lock()

    def __init__(self, agent_id: str = "main", base_dir: str | None = None,
                 user_id: str | None = None, team_ids: list[str] | None = None):
        self.agent_id = agent_id
        if base_dir:
            self.dir = base_dir
        else:
            self.dir = os.path.join(AGENTS_DIR, agent_id)
        os.makedirs(self.dir, exist_ok=True)
        self._collection = agent_id

        # Multi-user scoping
        self.user_id = user_id
        self.team_ids = team_ids or []
        self._user_dir = None
        self._team_dirs: list[tuple[str, str]] = []  # (dir_path, team_id)
        if user_id:
            self._user_dir = os.path.join(self.dir, "users", user_id)
            os.makedirs(self._user_dir, exist_ok=True)
        for tid in self.team_ids:
            td = os.path.join(self.dir, "teams", tid)
            os.makedirs(td, exist_ok=True)
            self._team_dirs.append((td, tid))

        # Ensure QMD knows about this collection (once per collection, background)
        with MemoryStore._ensured_lock:
            already_ensured = agent_id in MemoryStore._ensured_collections
            if not already_ensured:
                MemoryStore._ensured_collections.add(agent_id)
        if not already_ensured:
            threading.Thread(
                target=_qmd_ensure_collection,
                args=(self._collection, self.dir),
                daemon=True,
            ).start()

    def _make_id(self, name: str) -> str:
        """Generate a stable ID from name."""
        return hashlib.sha256(name.encode()).hexdigest()[:12]

    def _name_to_filename(self, name: str) -> str:
        """Convert a memory name to a safe filename with hash suffix to avoid collisions."""
        safe = re.sub(r'[^\w\s-]', '', name).strip().lower()
        safe = re.sub(r'[\s]+', '_', safe)
        # Add short hash to prevent collisions between similar names
        h = hashlib.sha256(name.encode()).hexdigest()[:6]
        base = safe[:50]
        return f"{base}_{h}.md" if base else f"{h}.md"

    def _find_file_for_name(self, name: str) -> str | None:
        """Find existing file for a memory name by checking frontmatter."""
        for fname in os.listdir(self.dir):
            if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
                continue
            fpath = os.path.join(self.dir, fname)
            try:
                with open(fpath, "r") as f:
                    raw = f.read(500)  # only need frontmatter
                fm, _ = _parse_frontmatter(raw)
                if fm.get("name") == name:
                    return fpath
            except Exception:
                continue
        return None

    def _resolve_store_dir(self, scope: str = "global") -> str:
        """Resolve target directory based on scope: global, user, team:<id>."""
        if scope == "user" and self._user_dir:
            return self._user_dir
        if scope.startswith("team:"):
            tid = scope[5:]
            for td, t_id in self._team_dirs:
                if t_id == tid:
                    return td
        return self.dir  # global

    def _all_scan_dirs(self) -> list[str]:
        """Get all directories to scan for this user's memories (global + user + team)."""
        dirs = [self.dir]
        if self._user_dir and os.path.isdir(self._user_dir):
            dirs.append(self._user_dir)
        for td, _ in self._team_dirs:
            if os.path.isdir(td):
                dirs.append(td)
        return dirs

    def store(self, name: str, content: str, description: str = "",
              mem_type: str = "general", scope: str = "global") -> dict:
        """Store or update a memory. Writes .md file and triggers QMD reindex.
        scope: 'global', 'user', or 'team:<team_id>'"""
        mem_id = self._make_id(name)
        filename = self._name_to_filename(name)
        target_dir = self._resolve_store_dir(scope)
        file_path = os.path.join(target_dir, filename)

        # Check if memory already exists under a different filename (migration)
        existing = self._find_file_for_name(name)
        if existing and existing != file_path:
            file_path = existing  # update in place
            filename = os.path.basename(existing)

        # Write markdown file with properly escaped frontmatter
        md_content = f"""---
name: {_yaml_escape(name)}
description: {_yaml_escape(description)}
type: {mem_type}
agent: {self.agent_id}
---

{content}
"""
        with open(file_path, "w") as f:
            f.write(md_content)

        # Trigger debounced QMD update+embed
        _qmd_debounced_embed(self._collection)

        # --- Entity extraction auto-linking (Mechanism 2) ---
        try:
            entities = _extract_entities(content)
            if entities:
                # Find other files sharing entities
                matches = _find_entity_matches(self.agent_id, filename, entities)
                for other_fname in matches[:10]:  # limit to avoid excessive linking
                    other_path = os.path.join(self.dir, other_fname)
                    if os.path.exists(other_path):
                        _add_related_to_file(file_path, other_fname, "same_topic")
                        _add_related_to_file(other_path, filename, "same_topic")
                # Update entity index with new file
                _update_entity_index(self.agent_id, filename, entities)
                if matches:
                    _qmd_debounced_embed(self._collection)
        except Exception as e:
            import logging
            logging.warning(f"Entity linking failed for {filename}: {e}")  # best-effort, never block store

        return {"id": mem_id, "name": name, "file": filename, "status": "stored"}

    def recall(self, query: str, limit: int = 10, mem_type: str | None = None) -> list[dict]:
        """Search memories using QMD hybrid search (BM25 + vector + reranking).
        Falls back to file-scan substring matching if QMD is unreachable."""
        # Try QMD first
        results = self._qmd_query(query, limit, mem_type)
        if results is not None:
            self._stamp_last_recalled_bg(results)
            return results
        # Fallback: scan files
        results = self._fallback_search(query, limit, mem_type)
        self._stamp_last_recalled_bg(results)
        return results

    def _stamp_last_recalled_bg(self, results: list[dict]):
        """Stamp last_recalled date on recalled files in a background thread."""
        paths = [r.get("file_path") for r in results if r.get("file_path")]
        if paths:
            threading.Thread(target=self._stamp_last_recalled, args=(paths,),
                             daemon=True, name="stamp_recalled").start()

    def _stamp_last_recalled(self, file_paths: list[str]):
        """Update last_recalled frontmatter field on recalled memory files. Best-effort."""
        now = datetime.datetime.now().strftime("%Y-%m-%d")
        for fpath in file_paths:
            try:
                with open(fpath, "r") as f:
                    raw = f.read()
                fm_match = re.match(r'^(---\s*\n)(.*?)(\n---\s*\n)(.*)$', raw, re.DOTALL)
                if not fm_match:
                    continue
                opener, fm_text, closer, body = fm_match.groups()
                if "last_recalled:" in fm_text:
                    fm_text = re.sub(r'last_recalled:.*', f'last_recalled: {now}', fm_text)
                else:
                    fm_text = fm_text.rstrip() + f"\nlast_recalled: {now}"
                with open(fpath, "w") as f:
                    f.write(opener + fm_text + closer + body)
            except Exception as e:
                import logging
                logging.debug(f"Failed to stamp last_recalled on {fpath}: {e}")
                continue

    def _qmd_query(self, query: str, limit: int, mem_type: str | None) -> list[dict] | None:
        """Query QMD via MCP HTTP. Returns list of results or None if unavailable."""
        global _qmd_session_id
        # Ensure session
        if not _qmd_session_id:
            if not _qmd_init_session():
                return None

        # Sanitize query: strip newlines, quotes, markdown — QMD silently returns empty on these
        clean_q = query.replace('\n', ' ').replace('\r', ' ').replace('"', '').replace("'", "")
        clean_q = re.sub(r'[#*`~\[\]{}()]', '', clean_q).strip()
        if not clean_q:
            return []

        searches = [
            {"type": "lex", "query": clean_q},
            {"type": "vec", "query": clean_q},
        ]
        result = _qmd_rpc("tools/call", {
            "name": "query",
            "arguments": {
                "searches": searches,
                "collections": [self._collection],
                "limit": limit * 2 if mem_type else limit,  # over-fetch if filtering
            },
        })
        if not result:
            # Session may have expired, retry once with lock to prevent stampede
            with _qmd_session_lock:
                _qmd_session_id = None
            if not _qmd_init_session():
                return None
            result = _qmd_rpc("tools/call", {
                "name": "query",
                "arguments": {
                    "searches": searches,
                    "collections": [self._collection],
                    "limit": limit * 2 if mem_type else limit,
                },
            })
            if not result:
                return None

        # Parse structured results
        structured = result.get("structuredContent", {})
        qmd_results = structured.get("results", [])
        memories = []
        for r in qmd_results:
            file_rel = r.get("file", "")
            # Strip collection prefix (e.g. "main/foo.md" -> "foo.md")
            if "/" in file_rel:
                fname = file_rel.split("/", 1)[1]
            else:
                fname = file_rel
            fpath = os.path.join(self.dir, fname)
            # Read the actual file for full content + frontmatter
            try:
                with open(fpath, "r") as f:
                    raw = f.read()
                fm, body = _parse_frontmatter(raw)
            except FileNotFoundError:
                fm = {"name": r.get("title", fname), "type": "general"}
                body = r.get("snippet", "")

            mem = {
                "id": self._make_id(fm.get("name", fname)),
                "name": fm.get("name", fname.replace(".md", "")),
                "description": fm.get("description", ""),
                "type": fm.get("type", "general"),
                "content": body,
                "file_path": fpath,
                "score": r.get("score", 0),
            }
            # Post-filter by type
            if mem_type and mem["type"] != mem_type:
                continue
            memories.append(mem)
            if len(memories) >= limit:
                break
        return memories

    def _fallback_search(self, query: str, limit: int, mem_type: str | None) -> list[dict]:
        """Fallback: scan .md files and do substring matching."""
        terms = query.lower().split()
        if not terms:
            return self.list_all(mem_type)[:limit]
        results = []
        # Directories to scan: all user-visible dirs + chats-indexed subdir
        scan_dirs = list(self._all_scan_dirs())
        chats_dir = os.path.join(self.dir, "chats-indexed")
        if os.path.isdir(chats_dir):
            scan_dirs.append(chats_dir)
        for scan_dir in scan_dirs:
            for fname in os.listdir(scan_dir):
                if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
                    continue
                fpath = os.path.join(scan_dir, fname)
                try:
                    # Cap file read to 32KB to prevent OOM on large files
                    with open(fpath, "r") as f:
                        raw = f.read(32768)
                    fm, body = _parse_frontmatter(raw)
                    mtype = fm.get("type", "general")
                    if mem_type and mtype != mem_type:
                        continue
                    searchable = (fm.get("name", "") + " " + fm.get("description", "") + " " + body).lower()
                    hits = sum(1 for t in terms if t in searchable)
                    if hits > 0:
                        results.append({
                            "id": self._make_id(fm.get("name", fname)),
                            "name": fm.get("name", fname.replace(".md", "")),
                            "description": fm.get("description", ""),
                            "type": mtype,
                            "content": body,
                            "file_path": fpath,
                            "score": hits / len(terms),
                        })
                except (UnicodeDecodeError, OSError) as e:
                    import logging
                    logging.debug(f"Fallback search skipping {fname}: {e}")
                    continue
                except Exception:
                    continue
        results.sort(key=lambda x: x["score"], reverse=True)
        return results[:limit]

    def delete(self, name: str) -> dict:
        """Delete a memory by name."""
        filename = self._name_to_filename(name)
        file_path = os.path.join(self.dir, filename)
        if not os.path.exists(file_path):
            # Try scanning for a file with matching frontmatter name
            found = self._find_file_for_name(name)
            if found:
                file_path = found
            else:
                return {"error": f"Memory '{name}' not found"}
        os.remove(file_path)
        # Trigger QMD reindex
        _qmd_debounced_embed(self._collection)
        return {"name": name, "status": "deleted"}

    def list_all(self, mem_type: str | None = None) -> list[dict]:
        """List all memories by scanning .md files (no QMD needed).
        Scans global + user + team directories."""
        results = []
        for scan_dir in self._all_scan_dirs():
            # Determine scope label
            scope = "global"
            if scan_dir == self._user_dir:
                scope = "user"
            else:
                for td, tid in self._team_dirs:
                    if scan_dir == td:
                        scope = f"team:{tid}"
                        break
            try:
                entries = os.listdir(scan_dir)
            except OSError:
                continue
            for fname in entries:
                if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
                    continue
                fpath = os.path.join(scan_dir, fname)
                try:
                    with open(fpath, "r") as f:
                        raw = f.read()
                    fm, body = _parse_frontmatter(raw)
                    mtype = fm.get("type", "general")
                    if mem_type and mtype != mem_type:
                        continue
                    mtime = os.path.getmtime(fpath)
                    results.append({
                        "id": self._make_id(fm.get("name", fname)),
                        "name": fm.get("name", fname.replace(".md", "")),
                        "description": fm.get("description", ""),
                        "type": mtype,
                        "content": body,
                        "updated_at": datetime.datetime.fromtimestamp(mtime).isoformat(),
                        "scope": scope,
                    })
                except Exception:
                    continue
        results.sort(key=lambda x: x.get("updated_at", ""), reverse=True)
        return results

    def reindex(self) -> dict:
        """Trigger QMD update+embed for this collection."""
        try:
            r1 = subprocess.run(
                ["qmd", "update"], capture_output=True, text=True, timeout=30,
            )
            r2 = subprocess.run(
                ["qmd", "embed", "-c", self._collection],
                capture_output=True, text=True, timeout=60,
            )
            return {"agent": self.agent_id, "status": "reindexed",
                    "update": r1.returncode == 0, "embed": r2.returncode == 0}
        except Exception as e:
            return {"agent": self.agent_id, "status": "error", "error": str(e)}


# Global memory store instance (set in _run_interactive)
_memory_store: MemoryStore | None = None


# ─── Projects System ──────────────────────────────────────────────────

class ProjectManager:
    """CRUD operations for per-agent projects."""

    @staticmethod
    def _project_dir(agent_id: str, name: str) -> str:
        return os.path.join(AGENTS_DIR, agent_id, "projects", name)

    @staticmethod
    def _projects_base(agent_id: str) -> str:
        return os.path.join(AGENTS_DIR, agent_id, "projects")

    @staticmethod
    def list_projects(agent_id: str, user_id: str | None = None,
                      user_team_ids: list[str] | None = None) -> list[dict]:
        """List projects for an agent, optionally filtered by user access.
        user_id=None means admin (sees all). user_team_ids are team IDs the user belongs to."""
        base = ProjectManager._projects_base(agent_id)
        if not os.path.isdir(base):
            return []
        projects = []
        for name in sorted(os.listdir(base)):
            pdir = os.path.join(base, name)
            if not os.path.isdir(pdir) or name.startswith("."):
                continue
            cfg_path = os.path.join(pdir, "project.json")
            cfg = {}
            if os.path.exists(cfg_path):
                try:
                    with open(cfg_path, "r") as f:
                        cfg = json.load(f)
                except (OSError, json.JSONDecodeError):
                    pass
            # Count ingested docs
            ingested_dir = os.path.join(pdir, "ingested")
            chunk_count = 0
            if os.path.isdir(ingested_dir):
                chunk_count = sum(1 for fn in os.listdir(ingested_dir) if fn.startswith("ingest-") and fn.endswith(".md"))
            # Count memory files
            mem_count = sum(1 for fn in os.listdir(pdir) if fn.endswith(".md") and not fn.startswith("ingest-") and fn not in _QMD_IGNORE_FILES)
            # Count ingested source files (unique sources, not chunks)
            doc_count = 0
            if os.path.isdir(ingested_dir):
                seen_sources = set()
                for fn in os.listdir(ingested_dir):
                    if fn.startswith("ingest-") and fn.endswith(".md"):
                        # Source hash is between first and second dash
                        parts_fn = fn.split("-", 2)
                        if len(parts_fn) >= 2:
                            seen_sources.add(parts_fn[1])
                doc_count = len(seen_sources) if seen_sources else chunk_count
            visibility = cfg.get("visibility", "global")
            owner_uid = cfg.get("owner_user_id", "")
            owner_tid = cfg.get("owner_team_id", "")
            projects.append({
                "name": name,
                "description": cfg.get("description", ""),
                "instructions": cfg.get("instructions", ""),
                "icon": cfg.get("icon", "folder"),
                "created_at": cfg.get("created_at", ""),
                "tags": cfg.get("tags", []),
                "watch_folders": cfg.get("watch_folders", []),
                "input_folders": cfg.get("input_folders", []) or [],
                "input_folders_last_scan": cfg.get("input_folders_last_scan", ""),
                "sync_status": cfg.get("sync_status", {}) or {},
                "status": cfg.get("status", "active"),
                "chunks": chunk_count,
                "doc_count": doc_count,
                "memories": mem_count,
                "visibility": visibility,
                "owner_user_id": owner_uid,
                "owner_team_id": owner_tid,
                "extra_member_user_ids": cfg.get("extra_member_user_ids", []) or [],
                "excluded_user_ids": cfg.get("excluded_user_ids", []) or [],
            })
        # Filter by user access if user_id provided
        if user_id is not None:
            team_set = set(user_team_ids or [])
            def _accessible(p):
                if p["owner_user_id"] == user_id:
                    return True
                extras = p.get("extra_member_user_ids") or []
                if user_id in extras:
                    return True
                if p["visibility"] == "global":
                    return user_id not in (p.get("excluded_user_ids") or [])
                if p["visibility"] == "team":
                    return p["owner_team_id"] in team_set
                # visibility == "user"
                return False
            projects = [p for p in projects if _accessible(p)]
        return projects

    @staticmethod
    def create_project(agent_id: str, name: str, description: str = "",
                       config: dict | None = None,
                       visibility: str = "global",
                       owner_user_id: str = "",
                       owner_team_id: str = "") -> dict:
        """Create a new project directory with project.json.
        visibility: 'global', 'user', or 'team'"""
        import uuid
        # Validate name
        safe_name = re.sub(r'[^\w\s-]', '', name).strip().lower().replace(' ', '-')
        if not safe_name:
            return {"error": "Invalid project name"}
        pdir = ProjectManager._project_dir(agent_id, safe_name)
        if os.path.exists(pdir):
            return {"error": f"Project '{safe_name}' already exists"}
        os.makedirs(pdir, exist_ok=True)
        os.makedirs(os.path.join(pdir, "ingested"), exist_ok=True)
        os.makedirs(os.path.join(pdir, "notes"), exist_ok=True)
        cfg = {
            "id": uuid.uuid4().hex[:12],
            "name": config.get("name", name) if config else name,
            "description": description,
            "icon": (config or {}).get("icon", "📁"),
            "created_at": datetime.datetime.now(datetime.timezone.utc).isoformat(),
            "watch_folders": (config or {}).get("watch_folders", []),
            "tags": (config or {}).get("tags", []),
            "model": (config or {}).get("model"),
            "visibility": visibility,
            "owner_user_id": owner_user_id,
            "owner_team_id": owner_team_id,
            "extra_member_user_ids": (config or {}).get("extra_member_user_ids", []) or [],
            "excluded_user_ids": (config or {}).get("excluded_user_ids", []) or [],
        }
        cfg_path = os.path.join(pdir, "project.json")
        with open(cfg_path, "w") as f:
            json.dump(cfg, f, indent=2)
        # Register QMD collection for this project
        collection_name = f"{agent_id}/{safe_name}"
        threading.Thread(
            target=_qmd_ensure_collection,
            args=(collection_name, pdir),
            daemon=True,
        ).start()
        return {"name": safe_name, "status": "created", "path": pdir}

    @staticmethod
    def get_project(agent_id: str, name: str) -> dict | None:
        """Read project.json for a project."""
        import uuid
        pdir = ProjectManager._project_dir(agent_id, name)
        if not os.path.isdir(pdir):
            return None
        cfg_path = os.path.join(pdir, "project.json")
        if not os.path.exists(cfg_path):
            # Auto-create minimal project.json for dirs that exist without one
            cfg = {"name": name, "description": "", "created_at": datetime.datetime.now(datetime.timezone.utc).isoformat()}
            try:
                with open(cfg_path, "w") as f:
                    json.dump(cfg, f, indent=2)
            except OSError:
                pass
        try:
            with open(cfg_path, "r") as f:
                cfg = json.load(f)
            # Backfill a stable globally-unique id on first read. Used as the
            # MemPalace wing key so renaming the project doesn't strand its
            # drawers, and so two same-named projects under different agents
            # never collide. Persisted lazily — if the disk file is read-only
            # we still hand back the in-memory id, but old chats may produce
            # new ids on each read until the file becomes writable.
            if not cfg.get("id"):
                cfg["id"] = uuid.uuid4().hex[:12]
                try:
                    with open(cfg_path, "w") as f:
                        json.dump(cfg, f, indent=2)
                except OSError:
                    pass
            # Add computed stats
            ingested_dir = os.path.join(pdir, "ingested")
            chunk_count = 0
            if os.path.isdir(ingested_dir):
                chunk_count = sum(1 for fn in os.listdir(ingested_dir) if fn.startswith("ingest-") and fn.endswith(".md"))
            cfg["chunks"] = chunk_count
            cfg["dir"] = pdir
            return cfg
        except (OSError, json.JSONDecodeError):
            return None

    @staticmethod
    def update_project(agent_id: str, name: str, updates: dict) -> dict:
        """Update project.json fields."""
        pdir = ProjectManager._project_dir(agent_id, name)
        cfg_path = os.path.join(pdir, "project.json")
        if not os.path.exists(cfg_path):
            return {"error": f"Project '{name}' not found"}
        try:
            with open(cfg_path, "r") as f:
                cfg = json.load(f)
            for k in ("description", "watch_folders", "tags", "model", "name", "icon",
                       "status", "instructions",
                       "input_folders", "input_folders_last_scan", "sync_status",
                       "visibility", "owner_user_id", "owner_team_id",
                       "extra_member_user_ids", "excluded_user_ids"):
                if k in updates:
                    cfg[k] = updates[k]
            with open(cfg_path, "w") as f:
                json.dump(cfg, f, indent=2)
            return {"name": name, "status": "updated"}
        except (OSError, json.JSONDecodeError) as e:
            return {"error": str(e)}

    @staticmethod
    def delete_project(agent_id: str, name: str) -> dict:
        """Soft-delete a project (move to .trash)."""
        pdir = ProjectManager._project_dir(agent_id, name)
        if not os.path.isdir(pdir):
            return {"error": f"Project '{name}' not found"}
        trash_dir = os.path.join(AGENTS_DIR, ".trash")
        os.makedirs(trash_dir, exist_ok=True)
        dest = os.path.join(trash_dir, f"{agent_id}_project_{name}_{int(time.time())}")
        shutil.move(pdir, dest)
        # Remove QMD collection
        collection_name = f"{agent_id}/{name}"
        def _find_qmd():
            p = shutil.which("qmd")
            if p:
                return p
            # Common locations when running under launchd
            for candidate in [
                os.path.expanduser("~/.nvm/versions/node/v22.20.0/bin/qmd"),
                "/opt/homebrew/bin/qmd",
                "/usr/local/bin/qmd",
            ]:
                if os.path.isfile(candidate):
                    return candidate
            return "qmd"
        try:
            subprocess.run(
                [_find_qmd(), "collection", "remove", collection_name],
                capture_output=True, text=True, timeout=5,
            )
        except Exception:
            pass
        return {"name": name, "status": "deleted", "moved_to": dest}


class NoteManager:
    """CRUD operations for project notes (markdown files in notes/ subdirectory)."""

    @staticmethod
    def _notes_dir(agent_id: str, project_name: str) -> str:
        return os.path.join(AGENTS_DIR, agent_id, "projects", project_name, "notes")

    @staticmethod
    def list_notes(agent_id: str, project_name: str) -> list[dict]:
        """Walk notes/ tree recursively, return metadata for each .md file."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        if not os.path.isdir(notes_dir):
            return []
        results = []
        for dirpath, _, filenames in os.walk(notes_dir):
            for fname in sorted(filenames):
                if not fname.endswith(".md"):
                    continue
                fpath = os.path.join(dirpath, fname)
                rel_path = os.path.relpath(fpath, notes_dir)
                try:
                    stat = os.stat(fpath)
                    with open(fpath, "r", errors="replace") as f:
                        raw = f.read(2000)
                    fm, _ = _parse_frontmatter(raw)
                    results.append({
                        "path": rel_path,
                        "name": fm.get("name", fname.replace(".md", "")),
                        "type": fm.get("type", "note"),
                        "size": stat.st_size,
                        "created_at": fm.get("created_at", ""),
                        "updated_at": fm.get("updated_at", ""),
                    })
                except Exception:
                    continue
        return results

    @staticmethod
    def get_note(agent_id: str, project_name: str, path: str) -> dict | None:
        """Read a note file and return its content with metadata."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        if not os.path.isfile(fpath):
            return None
        try:
            stat = os.stat(fpath)
            with open(fpath, "r", errors="replace") as f:
                raw = f.read()
            fm, body = _parse_frontmatter(raw)
            return {
                "path": path,
                "name": fm.get("name", os.path.basename(path).replace(".md", "")),
                "content": body,
                "frontmatter": fm,
                "size": stat.st_size,
                "created_at": fm.get("created_at", ""),
                "updated_at": fm.get("updated_at", ""),
            }
        except Exception:
            return None

    @staticmethod
    def create_note(agent_id: str, project_name: str, path: str, content: str = "") -> dict:
        """Create a note with YAML frontmatter, entity extraction, and QMD reindex."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        # Create parent dirs if needed
        os.makedirs(os.path.dirname(fpath), exist_ok=True)
        if os.path.exists(fpath):
            return {"error": f"Note '{path}' already exists"}

        now = datetime.datetime.now(datetime.timezone.utc).isoformat()
        name = os.path.basename(path).replace(".md", "")
        md_content = f"""---
name: {_yaml_escape(name)}
type: note
created_at: {now}
updated_at: {now}
agent: {agent_id}
project: {project_name}
---

{content}
"""
        with open(fpath, "w") as f:
            f.write(md_content)

        # Entity extraction + auto-link
        try:
            entities = _extract_entities(content)
            if entities:
                filename = os.path.basename(fpath)
                matches = _find_entity_matches(agent_id, filename, entities)
                for other_fname in matches[:10]:
                    other_path = os.path.join(os.path.dirname(fpath), other_fname)
                    if not os.path.exists(other_path):
                        # Try agent dir
                        other_path = os.path.join(AGENTS_DIR, agent_id, other_fname)
                    if os.path.exists(other_path):
                        _add_related_to_file(fpath, other_fname, "same_topic")
                        _add_related_to_file(other_path, filename, "same_topic")
                _update_entity_index(agent_id, filename, entities)
        except Exception:
            pass

        # Add same_folder relationships with sibling notes
        try:
            folder = os.path.dirname(fpath)
            fname = os.path.basename(fpath)
            for sibling in os.listdir(folder):
                if sibling.endswith(".md") and sibling != fname:
                    sibling_path = os.path.join(folder, sibling)
                    _add_related_to_file(fpath, sibling, "same_folder")
                    _add_related_to_file(sibling_path, fname, "same_folder")
        except Exception:
            pass

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"path": path, "status": "created"}

    @staticmethod
    def update_note(agent_id: str, project_name: str, path: str, content: str) -> dict:
        """Update a note's content, preserving frontmatter and updating timestamp."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        if not os.path.isfile(fpath):
            return {"error": f"Note '{path}' not found"}

        try:
            with open(fpath, "r") as f:
                raw = f.read()
            fm, _ = _parse_frontmatter(raw)
        except Exception:
            fm = {}

        now = datetime.datetime.now(datetime.timezone.utc).isoformat()
        fm["updated_at"] = now

        # Rebuild frontmatter
        fm_lines = []
        for k, v in fm.items():
            if k == "related":
                continue  # related is multi-line, handle separately
            fm_lines.append(f"{k}: {v}")

        # Preserve related section if it exists
        related_section = ""
        fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', raw, re.DOTALL)
        if fm_match:
            fm_text = fm_match.group(1)
            rel_match = re.search(r'(related:.*)', fm_text, re.DOTALL)
            if rel_match:
                related_section = "\n" + rel_match.group(1)

        md_content = f"---\n" + "\n".join(fm_lines) + related_section + f"\n---\n\n{content}\n"
        with open(fpath, "w") as f:
            f.write(md_content)

        # Re-run entity extraction
        try:
            entities = _extract_entities(content)
            if entities:
                filename = os.path.basename(fpath)
                _update_entity_index(agent_id, filename, entities)
        except Exception:
            pass

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"path": path, "status": "updated"}

    @staticmethod
    def delete_note(agent_id: str, project_name: str, path: str) -> dict:
        """Remove a note file and trigger QMD reindex."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        if not os.path.isfile(fpath):
            return {"error": f"Note '{path}' not found"}

        os.remove(fpath)

        # Clean up empty parent directories
        parent = os.path.dirname(fpath)
        while parent != notes_dir:
            try:
                if not os.listdir(parent):
                    os.rmdir(parent)
                    parent = os.path.dirname(parent)
                else:
                    break
            except Exception:
                break

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"path": path, "status": "deleted"}

    @staticmethod
    def rename_note(agent_id: str, project_name: str, old_path: str, new_path: str) -> dict:
        """Rename/move a note file and update its frontmatter name."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        old_fpath = os.path.join(notes_dir, old_path)
        new_fpath = os.path.join(notes_dir, new_path)
        if not os.path.isfile(old_fpath):
            return {"error": f"Note '{old_path}' not found"}
        if os.path.exists(new_fpath):
            return {"error": f"Note '{new_path}' already exists"}

        # Create parent dirs for new path if needed
        os.makedirs(os.path.dirname(new_fpath), exist_ok=True)
        os.rename(old_fpath, new_fpath)

        # Update frontmatter name field
        try:
            with open(new_fpath, "r") as f:
                raw = f.read()
            new_name = os.path.basename(new_path).replace(".md", "")
            raw = re.sub(r'^(name:\s*).*$', rf'\g<1>{_yaml_escape(new_name)}', raw, count=1, flags=re.MULTILINE)
            now = datetime.datetime.now(datetime.timezone.utc).isoformat()
            raw = re.sub(r'^(updated_at:\s*).*$', rf'\g<1>{now}', raw, count=1, flags=re.MULTILINE)
            with open(new_fpath, "w") as f:
                f.write(raw)
        except Exception:
            pass

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"old_path": old_path, "new_path": new_path, "status": "renamed"}

    @staticmethod
    def create_folder(agent_id: str, project_name: str, folder_path: str) -> dict:
        """Create a folder within the notes directory."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, folder_path)
        os.makedirs(fpath, exist_ok=True)
        return {"path": folder_path, "status": "created"}


# ─── Document Ingestion Engine ────────────────────────────────────────

class DocumentParser:
    """Parse various document formats to plain text."""

    @staticmethod
    def parse_pdf(path: str) -> str:
        """Parse PDF to text using pymupdf."""
        try:
            import fitz  # pymupdf
        except ImportError:
            raise ImportError("Install pymupdf for PDF support: pip3 install pymupdf")
        doc = fitz.open(path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n\n".join(pages)

    @staticmethod
    def parse_docx(path: str) -> str:
        """Parse DOCX to text using python-docx."""
        try:
            import docx
        except ImportError:
            raise ImportError("Install python-docx for DOCX support: pip3 install python-docx")
        doc = docx.Document(path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    level = 1
                    try:
                        level = int(para.style.name.replace("Heading", "").strip()) or 1
                    except ValueError:
                        pass
                    text = "#" * level + " " + text
                paragraphs.append(text)
        return "\n\n".join(paragraphs)

    @staticmethod
    def parse_txt(path: str) -> str:
        """Parse plain text file."""
        with open(path, "r", errors="replace") as f:
            return f.read()

    @staticmethod
    def parse_md(path: str) -> str:
        """Parse markdown file (keep as-is)."""
        with open(path, "r", errors="replace") as f:
            return f.read()

    @staticmethod
    def parse_html(content: str) -> str:
        """Strip HTML tags and extract text content."""
        from html.parser import HTMLParser

        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self._parts: list[str] = []
                self._skip = False
                self._skip_tags = {"script", "style", "nav", "header", "footer", "noscript"}

            def handle_starttag(self, tag, attrs):
                if tag in self._skip_tags:
                    self._skip = True
                elif tag in ("p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr"):
                    self._parts.append("\n")
                    if tag.startswith("h"):
                        level = int(tag[1])
                        self._parts.append("#" * level + " ")

            def handle_endtag(self, tag):
                if tag in self._skip_tags:
                    self._skip = False
                elif tag in ("p", "div", "h1", "h2", "h3", "h4", "h5", "h6"):
                    self._parts.append("\n")

            def handle_data(self, data):
                if not self._skip:
                    self._parts.append(data)

        extractor = _TextExtractor()
        extractor.feed(content)
        text = "".join(extractor._parts)
        # Clean up whitespace
        lines = [line.strip() for line in text.split("\n")]
        cleaned = "\n".join(lines)
        # Collapse multiple blank lines
        cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
        return cleaned.strip()

    @staticmethod
    def parse_xlsx(path: str, sheet: str | None = None) -> str:
        """Parse XLSX/XLS to markdown tables using openpyxl."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("Install openpyxl for XLSX support: pip3 install openpyxl")
        wb = openpyxl.load_workbook(path, read_only=True, data_only=False)
        parts = []
        sheet_names = wb.sheetnames
        parts.append(f"**Sheets:** {', '.join(sheet_names)}\n")
        target_sheets = [sheet] if sheet and sheet in sheet_names else sheet_names
        for sname in target_sheets:
            ws = wb[sname]
            parts.append(f"## Sheet: {sname}\n")
            rows = []
            for row in ws.iter_rows(values_only=False):
                cells = []
                for cell in row:
                    val = cell.value
                    if val is None:
                        cells.append("")
                    elif isinstance(val, str) and val.startswith("="):
                        # Show formula
                        cells.append(f"{val}")
                    else:
                        cells.append(str(val))
                rows.append(cells)
            if not rows:
                parts.append("*(empty sheet)*\n")
                continue
            # Build markdown table
            max_cols = max(len(r) for r in rows) if rows else 0
            # Pad rows to same length
            for r in rows:
                while len(r) < max_cols:
                    r.append("")
            # Header row
            header = "| " + " | ".join(rows[0]) + " |"
            sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
            table_lines = [header, sep]
            for r in rows[1:]:
                table_lines.append("| " + " | ".join(r) + " |")
            parts.append("\n".join(table_lines) + "\n")
        wb.close()
        return "\n".join(parts)

    @staticmethod
    def parse_pptx(path: str, slides: str | None = None) -> str:
        """Parse PPTX to text using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Install python-pptx for PPTX support: pip3 install python-pptx")
        prs = Presentation(path)
        total_slides = len(prs.slides)
        parts = [f"**Slides:** {total_slides}\n"]
        # Parse slide range
        slide_indices = None
        if slides:
            slide_indices = set()
            for part in slides.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    for i in range(int(a), int(b) + 1):
                        slide_indices.add(i)
                else:
                    slide_indices.add(int(part))
        for idx, slide in enumerate(prs.slides, 1):
            if slide_indices and idx not in slide_indices:
                continue
            parts.append(f"## Slide {idx}")
            # Title
            if slide.shapes.title:
                parts.append(f"**Title:** {slide.shapes.title.text}")
            # Body text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        max_cols = max(len(r) for r in rows)
                        for r in rows:
                            while len(r) < max_cols:
                                r.append("")
                        header = "| " + " | ".join(rows[0]) + " |"
                        sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
                        table_lines = [header, sep]
                        for r in rows[1:]:
                            table_lines.append("| " + " | ".join(r) + " |")
                        parts.append("\n".join(table_lines))
            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"*Speaker Notes:* {notes}")
            parts.append("")
        return "\n".join(parts)

    @staticmethod
    def parse_csv(path: str) -> str:
        """Parse CSV/TSV to markdown table."""
        import csv
        delimiter = "\t" if path.lower().endswith(".tsv") else ","
        with open(path, "r", newline="", errors="replace") as f:
            reader = csv.reader(f, delimiter=delimiter)
            rows = [row for row in reader]
        if not rows:
            return "*(empty file)*"
        max_cols = max(len(r) for r in rows)
        for r in rows:
            while len(r) < max_cols:
                r.append("")
        header = "| " + " | ".join(rows[0]) + " |"
        sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
        lines = [header, sep]
        for r in rows[1:]:
            lines.append("| " + " | ".join(r) + " |")
        return "\n".join(lines)

    @staticmethod
    def parse_image(path: str) -> str:
        """Parse image metadata using Pillow. Returns metadata text."""
        try:
            from PIL import Image
        except ImportError:
            raise ImportError("Install Pillow for image support: pip3 install Pillow")
        img = Image.open(path)
        width, height = img.size
        fmt = img.format or os.path.splitext(path)[1].lstrip(".")
        mode = img.mode
        info_parts = [
            f"**Image:** {os.path.basename(path)}",
            f"**Dimensions:** {width} x {height}",
            f"**Format:** {fmt}",
            f"**Mode:** {mode}",
        ]
        # Extract EXIF if available
        try:
            exif = img.getexif()
            if exif:
                for tag_id, value in list(exif.items())[:10]:
                    try:
                        from PIL.ExifTags import TAGS
                        tag_name = TAGS.get(tag_id, str(tag_id))
                        info_parts.append(f"**{tag_name}:** {value}")
                    except Exception:
                        pass
        except Exception:
            pass
        img.close()
        return "\n".join(info_parts)

    @staticmethod
    def parse_svg(path: str) -> str:
        """Parse SVG to extract text elements and metadata."""
        from xml.etree import ElementTree
        tree = ElementTree.parse(path)
        root = tree.getroot()
        ns = {"svg": "http://www.w3.org/2000/svg"}
        parts = [f"**SVG:** {os.path.basename(path)}"]
        # Get dimensions
        w = root.get("width", "")
        h = root.get("height", "")
        vb = root.get("viewBox", "")
        if w and h:
            parts.append(f"**Dimensions:** {w} x {h}")
        if vb:
            parts.append(f"**ViewBox:** {vb}")
        # Extract title and desc
        for tag in ("title", "desc"):
            el = root.find(f"svg:{tag}", ns) or root.find(tag)
            if el is not None and el.text:
                parts.append(f"**{tag.capitalize()}:** {el.text.strip()}")
        # Extract all text elements
        texts = []
        for text_el in list(root.iter(f"{{{ns['svg']}}}text")) + list(root.iter("text")):
            t = "".join(text_el.itertext()).strip()
            if t:
                texts.append(t)
        if texts:
            parts.append(f"\n**Text content:**")
            for t in texts:
                parts.append(f"- {t}")
        return "\n".join(parts)

    @staticmethod
    def parse_url(url: str) -> str:
        """Fetch URL and parse HTML to text."""
        req_headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        }
        req = urllib.request.Request(url, headers=req_headers)
        with urllib.request.urlopen(req, timeout=30) as resp:
            raw = resp.read()
            encoding = resp.headers.get("Content-Encoding", "")
            if encoding == "gzip":
                import gzip
                raw = gzip.decompress(raw)
            charset = resp.headers.get_content_charset() or "utf-8"
            html = raw.decode(charset, errors="replace")
        return DocumentParser.parse_html(html)

    @staticmethod
    def parse(path_or_url: str) -> tuple[str, str]:
        """Auto-detect format and parse. Returns (text, source_type)."""
        if path_or_url.startswith(("http://", "https://")):
            return DocumentParser.parse_url(path_or_url), "url"
        ext = os.path.splitext(path_or_url)[1].lower()
        parsers = {
            ".pdf": ("pdf", DocumentParser.parse_pdf),
            ".docx": ("docx", DocumentParser.parse_docx),
            ".txt": ("txt", DocumentParser.parse_txt),
            ".md": ("md", DocumentParser.parse_md),
            ".html": ("html", lambda p: DocumentParser.parse_html(open(p, "r", errors="replace").read())),
            ".htm": ("html", lambda p: DocumentParser.parse_html(open(p, "r", errors="replace").read())),
            ".xlsx": ("xlsx", DocumentParser.parse_xlsx),
            ".xls": ("xlsx", DocumentParser.parse_xlsx),
            ".pptx": ("pptx", DocumentParser.parse_pptx),
            ".csv": ("csv", DocumentParser.parse_csv),
            ".tsv": ("csv", DocumentParser.parse_csv),
            ".png": ("image", DocumentParser.parse_image),
            ".jpg": ("image", DocumentParser.parse_image),
            ".jpeg": ("image", DocumentParser.parse_image),
            ".gif": ("image", DocumentParser.parse_image),
            ".webp": ("image", DocumentParser.parse_image),
            ".bmp": ("image", DocumentParser.parse_image),
            ".svg": ("svg", DocumentParser.parse_svg),
        }
        if ext not in parsers:
            raise ValueError(f"Unsupported format: {ext}. Supported: {', '.join(parsers.keys())}")
        source_type, parser_fn = parsers[ext]
        return parser_fn(path_or_url), source_type


class DocumentChunker:
    """Split text into overlapping chunks with section header preservation."""

    @staticmethod
    def chunk(text: str, chunk_size: int = 1500, chunk_overlap: int = 200,
              min_chunk_size: int = 100) -> list[dict]:
        """Split text into chunks. chunk_size/overlap are in ~tokens (chars/4 approximation).
        Returns list of {text, index, total, header}."""
        # Convert token counts to char approximation
        max_chars = chunk_size * 4
        overlap_chars = chunk_overlap * 4
        min_chars = min_chunk_size * 4

        # Split into paragraphs
        paragraphs = re.split(r'\n\s*\n', text)
        paragraphs = [p.strip() for p in paragraphs if p.strip()]

        chunks: list[dict] = []
        current_parts: list[str] = []
        current_len = 0
        last_header = ""

        def _flush():
            nonlocal current_parts, current_len
            if not current_parts:
                return
            chunk_text = "\n\n".join(current_parts)
            if len(chunk_text) < min_chars:
                return
            # Prepend section header if available
            if last_header and not chunk_text.startswith("#"):
                chunk_text = last_header + "\n\n" + chunk_text
            chunks.append({
                "text": chunk_text,
                "index": len(chunks),
                "total": 0,  # filled in later
                "header": last_header,
            })
            # Keep overlap: take text from end of current chunk
            if overlap_chars > 0:
                overlap_text = chunk_text[-overlap_chars:]
                current_parts = [overlap_text]
                current_len = len(overlap_text)
            else:
                current_parts = []
                current_len = 0

        for para in paragraphs:
            # Track section headers
            header_match = re.match(r'^(#{1,6}\s+.+)', para.split('\n')[0])
            if header_match:
                last_header = header_match.group(1)

            # If a single paragraph exceeds max_chars, split it
            if len(para) > max_chars:
                # Flush current buffer first
                if current_parts:
                    _flush()
                # Split on sentences
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sent in sentences:
                    if len(sent) > max_chars:
                        # Split on words as last resort
                        words = sent.split()
                        for word in words:
                            if current_len + len(word) + 1 > max_chars:
                                _flush()
                            current_parts.append(word)
                            current_len += len(word) + 1
                    else:
                        if current_len + len(sent) + 1 > max_chars:
                            _flush()
                        current_parts.append(sent)
                        current_len += len(sent) + 1
            else:
                if current_len + len(para) + 2 > max_chars:
                    _flush()
                current_parts.append(para)
                current_len += len(para) + 2

        # Flush remaining
        if current_parts:
            chunk_text = "\n\n".join(current_parts)
            if len(chunk_text) >= min_chars:
                if last_header and not chunk_text.startswith("#"):
                    chunk_text = last_header + "\n\n" + chunk_text
                chunks.append({
                    "text": chunk_text,
                    "index": len(chunks),
                    "total": 0,
                    "header": last_header,
                })

        # Fill in total count
        for c in chunks:
            c["total"] = len(chunks)

        return chunks


class IngestManager:
    """Ingest files and URLs into agent or project memory as chunked markdown."""

    @staticmethod
    def _source_hash(source: str) -> str:
        """6-char hash of source name/URL."""
        return hashlib.sha256(source.encode()).hexdigest()[:6]

    @staticmethod
    def _ingest_dir(agent_id: str, project_name: str | None = None) -> str:
        """Get the directory where ingested chunks are stored."""
        if project_name:
            d = os.path.join(AGENTS_DIR, agent_id, "projects", project_name, "ingested")
        else:
            d = os.path.join(AGENTS_DIR, agent_id, "ingested")
        os.makedirs(d, exist_ok=True)
        return d

    @staticmethod
    def _collection_name(agent_id: str, project_name: str | None = None) -> str:
        """QMD collection name for embedding."""
        if project_name:
            return f"{agent_id}/{project_name}"
        return agent_id

    @staticmethod
    def ingest_file(agent_id: str, file_path: str,
                    project_name: str | None = None,
                    tags: list[str] | None = None,
                    chunk_size: int = 1500, chunk_overlap: int = 200) -> dict:
        """Parse, chunk, and store a file as ingested memory chunks."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        source_name = os.path.basename(file_path)
        try:
            text, source_type = DocumentParser.parse(file_path)
        except (ImportError, ValueError) as e:
            return {"error": str(e)}
        return IngestManager._store_chunks(
            agent_id, project_name, source_name, source_type, text,
            tags=tags, chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        )

    @staticmethod
    def ingest_url(agent_id: str, url: str,
                   project_name: str | None = None,
                   tags: list[str] | None = None,
                   chunk_size: int = 1500, chunk_overlap: int = 200) -> dict:
        """Fetch URL, parse HTML, chunk, and store."""
        try:
            text = DocumentParser.parse_url(url)
        except Exception as e:
            return {"error": f"Failed to fetch URL: {e}"}
        return IngestManager._store_chunks(
            agent_id, project_name, url, "url", text,
            tags=tags, chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        )

    @staticmethod
    def _store_chunks(agent_id: str, project_name: str | None,
                      source: str, source_type: str, text: str,
                      tags: list[str] | None = None,
                      chunk_size: int = 1500, chunk_overlap: int = 200) -> dict:
        """Chunk text and write as ingest-*.md files with frontmatter."""
        src_hash = IngestManager._source_hash(source)
        ingest_dir = IngestManager._ingest_dir(agent_id, project_name)
        collection = IngestManager._collection_name(agent_id, project_name)

        # Delete existing chunks for this source (re-ingest)
        existing = [f for f in os.listdir(ingest_dir) if f.startswith(f"ingest-{src_hash}-") and f.endswith(".md")]
        for f in existing:
            os.remove(os.path.join(ingest_dir, f))

        # Chunk
        chunks = DocumentChunker.chunk(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        if not chunks:
            return {"error": "No content extracted from document"}

        all_tags = ["ingested"]
        if tags:
            all_tags.extend(tags)
        # Add source name as tag (sanitized)
        safe_source_tag = re.sub(r'[^\w-]', '', source.split("/")[-1].split(".")[0].lower())
        if safe_source_tag:
            all_tags.append(safe_source_tag)
        tags_yaml = "\n".join(f"  - {t}" for t in all_tags)

        now = datetime.datetime.now(datetime.timezone.utc).isoformat()
        files_written = []
        for chunk in chunks:
            idx = chunk["index"]
            total = chunk["total"]
            title = chunk["header"] or f"{source} - Chunk {idx + 1}"

            # Build related links
            related_lines = []
            if idx > 0:
                prev_file = f"ingest-{src_hash}-{idx - 1:03d}.md"
                related_lines.append(f"  - file: {prev_file}\n    type: prev_chunk")
            if idx < total - 1:
                next_file = f"ingest-{src_hash}-{idx + 1:03d}.md"
                related_lines.append(f"  - file: {next_file}\n    type: next_chunk")
            if idx != 0:
                first_file = f"ingest-{src_hash}-000.md"
                related_lines.append(f"  - file: {first_file}\n    type: same_source")
            related_yaml = ""
            if related_lines:
                related_yaml = "related:\n" + "\n".join(related_lines) + "\n"

            filename = f"ingest-{src_hash}-{idx:03d}.md"
            md_content = f"""---
title: {_yaml_escape(title)}
source: {_yaml_escape(source)}
source_type: {source_type}
ingested_at: "{now}"
chunk_index: {idx}
total_chunks: {total}
agent: {agent_id}
tags:
{tags_yaml}
{related_yaml}---

{chunk['text']}
"""
            fpath = os.path.join(ingest_dir, filename)
            with open(fpath, "w") as f:
                f.write(md_content)
            files_written.append(filename)

        # Trigger QMD indexing
        _qmd_debounced_embed(collection)

        word_count = len(text.split())
        return {
            "source": source,
            "source_type": source_type,
            "source_hash": src_hash,
            "chunks": len(chunks),
            "words": word_count,
            "files": files_written,
            "agent": agent_id,
            "project": project_name,
            "status": "ingested",
        }

    @staticmethod
    def list_ingested(agent_id: str, project_name: str | None = None) -> list[dict]:
        """List ingested documents grouped by source."""
        ingest_dir = IngestManager._ingest_dir(agent_id, project_name)
        if not os.path.isdir(ingest_dir):
            return []
        # Group by source hash
        groups: dict[str, dict] = {}
        for fname in os.listdir(ingest_dir):
            if not fname.startswith("ingest-") or not fname.endswith(".md"):
                continue
            fpath = os.path.join(ingest_dir, fname)
            try:
                with open(fpath, "r") as f:
                    raw = f.read(800)
                fm, _ = _parse_frontmatter(raw)
            except Exception:
                continue
            source = fm.get("source", "unknown")
            src_hash = fname.split("-")[1] if "-" in fname else "?"
            if src_hash not in groups:
                groups[src_hash] = {
                    "source": source,
                    "source_type": fm.get("source_type", "unknown"),
                    "source_hash": src_hash,
                    "chunks": 0,
                    "ingested_at": fm.get("ingested_at", ""),
                    "tags": [],
                }
            groups[src_hash]["chunks"] += 1
            # Parse tags from frontmatter
            tags_str = fm.get("tags", "")
            if isinstance(tags_str, str) and tags_str:
                for t in tags_str.split(","):
                    t = t.strip().strip("-").strip()
                    if t and t not in groups[src_hash]["tags"]:
                        groups[src_hash]["tags"].append(t)
        return sorted(groups.values(), key=lambda x: x.get("ingested_at", ""), reverse=True)

    @staticmethod
    def delete_ingested(agent_id: str, source_hash: str,
                        project_name: str | None = None) -> dict:
        """Delete all chunks for a source hash."""
        ingest_dir = IngestManager._ingest_dir(agent_id, project_name)
        if not os.path.isdir(ingest_dir):
            return {"error": "No ingested documents found"}
        deleted = 0
        source_name = ""
        for fname in os.listdir(ingest_dir):
            if fname.startswith(f"ingest-{source_hash}-") and fname.endswith(".md"):
                if not source_name:
                    fpath = os.path.join(ingest_dir, fname)
                    try:
                        with open(fpath, "r") as f:
                            fm, _ = _parse_frontmatter(f.read(500))
                        source_name = fm.get("source", "unknown")
                    except Exception:
                        pass
                os.remove(os.path.join(ingest_dir, fname))
                deleted += 1
        collection = IngestManager._collection_name(agent_id, project_name)
        _qmd_debounced_embed(collection)
        return {"source": source_name, "source_hash": source_hash, "deleted": deleted}


# ─── Watched Folders (Auto-Ingestion) ────────────────────────────────

class IngestWatcher:
    """Background thread that polls watched folders and auto-ingests new/modified files."""

    POLL_INTERVAL = 30  # seconds

    def __init__(self):
        self._stop = threading.Event()
        self._thread = None

    def start(self):
        """Start the background watcher thread."""
        self._stop.clear()
        self._thread = threading.Thread(target=self._run_loop, daemon=True, name="ingest_watcher")
        self._thread.start()

    def stop(self):
        self._stop.set()
        if self._thread:
            self._thread.join(timeout=2)

    def _run_loop(self):
        """Poll all watched folders across all agents and projects."""
        while not self._stop.is_set():
            try:
                self._scan_all()
            except Exception as e:
                import logging
                logging.debug("IngestWatcher error: %s", e)
            self._stop.wait(self.POLL_INTERVAL)

    def _scan_all(self):
        """Scan watched folders for all agents and projects."""
        if not os.path.isdir(AGENTS_DIR):
            return
        for agent_name in os.listdir(AGENTS_DIR):
            if agent_name.startswith("."):
                continue
            agent_dir = os.path.join(AGENTS_DIR, agent_name)
            if not os.path.isdir(agent_dir):
                continue
            # Check agent-level watches (from agent.json)
            agent_json_path = os.path.join(agent_dir, "agent.json")
            if os.path.exists(agent_json_path):
                try:
                    with open(agent_json_path, "r") as f:
                        agent_cfg = json.load(f)
                    watches = agent_cfg.get("ingest_watch", [])
                    if watches:
                        self._process_watches(agent_name, None, watches, agent_dir)
                except (OSError, json.JSONDecodeError):
                    pass
            # Check project-level watches
            projects_dir = os.path.join(agent_dir, "projects")
            if os.path.isdir(projects_dir):
                for proj_name in os.listdir(projects_dir):
                    proj_dir = os.path.join(projects_dir, proj_name)
                    proj_json = os.path.join(proj_dir, "project.json")
                    if not os.path.exists(proj_json):
                        continue
                    try:
                        with open(proj_json, "r") as f:
                            proj_cfg = json.load(f)
                        watches = proj_cfg.get("watch_folders", [])
                        if watches:
                            self._process_watches(agent_name, proj_name, watches, proj_dir)
                    except (OSError, json.JSONDecodeError):
                        pass

    def _process_watches(self, agent_id: str, project_name: str | None,
                         watches: list[dict], base_dir: str):
        """Process watched folders, detect changes, ingest as needed."""
        registry_path = os.path.join(base_dir, "ingest_registry.json")
        registry = {}
        if os.path.exists(registry_path):
            try:
                with open(registry_path, "r") as f:
                    registry = json.load(f)
            except (OSError, json.JSONDecodeError):
                pass
        watches_reg = registry.get("watches", {})
        changed = False

        for watch in watches:
            watch_path = watch.get("path", "")
            if not watch_path or not os.path.isdir(watch_path):
                continue
            pattern = watch.get("pattern", "*")
            recursive = watch.get("recursive", False)
            tags = watch.get("tags", [])
            chunk_size = watch.get("chunk_size", 1500)

            # Get or create registry entry for this watch
            wreg = watches_reg.get(watch_path, {"files": {}, "last_scan": ""})

            # Scan for matching files
            if recursive:
                matched_files = []
                for root, _dirs, files in os.walk(watch_path):
                    for fn in files:
                        if fnmatch.fnmatch(fn, pattern):
                            matched_files.append(os.path.join(root, fn))
            else:
                matched_files = [
                    os.path.join(watch_path, fn) for fn in os.listdir(watch_path)
                    if fnmatch.fnmatch(fn, pattern) and os.path.isfile(os.path.join(watch_path, fn))
                ]

            current_files = set()
            for fpath in matched_files:
                fname = os.path.basename(fpath)
                current_files.add(fname)
                try:
                    stat = os.stat(fpath)
                except OSError:
                    continue
                prev = wreg["files"].get(fname, {})
                if prev.get("mtime") == stat.st_mtime and prev.get("size") == stat.st_size:
                    continue  # unchanged

                # New or modified file — ingest
                try:
                    result = IngestManager.ingest_file(
                        agent_id, fpath, project_name=project_name,
                        tags=tags, chunk_size=chunk_size,
                    )
                    if "error" not in result:
                        wreg["files"][fname] = {
                            "mtime": stat.st_mtime,
                            "size": stat.st_size,
                            "hash": result.get("source_hash", ""),
                            "chunks": result.get("chunks", 0),
                            "ingested_at": datetime.datetime.now(datetime.timezone.utc).isoformat(),
                        }
                        changed = True
                except Exception as e:
                    import logging
                    logging.debug("IngestWatcher: failed to ingest %s: %s", fpath, e)

            # Detect deleted files
            for fname in list(wreg["files"].keys()):
                if fname not in current_files:
                    src_hash = wreg["files"][fname].get("hash", "")
                    if src_hash:
                        IngestManager.delete_ingested(agent_id, src_hash, project_name)
                    del wreg["files"][fname]
                    changed = True

            wreg["last_scan"] = datetime.datetime.now(datetime.timezone.utc).isoformat()
            watches_reg[watch_path] = wreg

        if changed:
            registry["watches"] = watches_reg
            try:
                with open(registry_path, "w") as f:
                    json.dump(registry, f, indent=2)
            except OSError:
                pass


# Global IngestWatcher instance (started alongside scheduler in server.py)
_ingest_watcher: IngestWatcher | None = None


def _get_memory_store() -> MemoryStore | None:
    """Get the active memory store: thread-local (delegation/scheduler) or global (main thread)."""
    return getattr(_thread_local, 'memory_store', None) or _memory_store


def tool_memory_store(args: dict) -> str:
    """Store a memory. When a project is active, writes to project directory."""
    ms = _get_memory_store()
    if not ms:
        return _err("Memory store not initialized")
    name = args.get("name", "")
    content = args.get("content", "")
    description = args.get("description", "")
    mem_type = args.get("type", "general")
    if not name or not content:
        return _err("memory_store: name and content are required")
    # If project is active, store in project directory
    project = getattr(_thread_local, 'project', None)
    if project:
        agent_id = ms.agent_id
        proj_dir = os.path.join(AGENTS_DIR, agent_id, "projects", project)
        if os.path.isdir(proj_dir):
            proj_store = MemoryStore(agent_id=f"{agent_id}/{project}", base_dir=proj_dir)
            result = proj_store.store(name, content, description, mem_type)
            result["project"] = project
            return _ok(result)
    result = ms.store(name, content, description, mem_type)
    # Trigger near-term memory summary refresh when user-facing memories are stored
    # (skip if this IS the memory summary being written)
    if name != "Memory Summary" and mem_type in ("user", "feedback", "project"):
        try:
            agent_id = ms.agent_id if hasattr(ms, 'agent_id') else "main"
            trigger_memory_summary_refresh(agent_id)
        except Exception:
            pass
    return _ok(result)


def _graph_expand_results(results: list[dict], base_dir: str, ingest_dir: str,
                          max_hops: int = 1) -> list[dict]:
    """Follow 'related' frontmatter links from matched results for context expansion."""
    seen_files = {r.get("file_path", "") for r in results}
    expanded = list(results)
    frontier = list(results)
    for _hop in range(max_hops):
        next_frontier = []
        for r in frontier:
            fpath = r.get("file_path", "")
            if not fpath or not os.path.exists(fpath):
                continue
            try:
                with open(fpath, "r") as f:
                    raw = f.read(2000)
                fm, _ = _parse_frontmatter(raw)
            except Exception:
                continue
            # Parse related field (simple YAML list parsing)
            related_raw = fm.get("related", "")
            if not related_raw:
                continue
            # related is stored as multi-line YAML in frontmatter, parse linked files
            related_files = re.findall(r'file:\s*(\S+\.md)', raw)
            for rel_file in related_files:
                # Try ingest_dir first, then base_dir
                for search_dir in (ingest_dir, base_dir):
                    rel_path = os.path.join(search_dir, rel_file)
                    if rel_path in seen_files or not os.path.exists(rel_path):
                        continue
                    seen_files.add(rel_path)
                    try:
                        with open(rel_path, "r") as f:
                            rel_raw = f.read()
                        rel_fm, rel_body = _parse_frontmatter(rel_raw)
                        mem = {
                            "id": hashlib.sha256(rel_fm.get("name", rel_file).encode()).hexdigest()[:12],
                            "name": rel_fm.get("name", rel_fm.get("title", rel_file.replace(".md", ""))),
                            "description": rel_fm.get("description", ""),
                            "type": rel_fm.get("type", "general"),
                            "content": rel_body,
                            "file_path": rel_path,
                            "score": max(0, (r.get("score", 0.5) - 0.2)),
                            "source_scope": "related",
                        }
                        expanded.append(mem)
                        next_frontier.append(mem)
                    except Exception:
                        continue
                    break  # found in one dir, skip the other
        frontier = next_frontier
    return expanded


def tool_memory_recall(args: dict) -> str:
    """Recall memories by searching. When a project is active, searches project first."""
    ms = _get_memory_store()
    if not ms:
        return _err("Memory store not initialized")
    query = args.get("query", "")
    limit = args.get("limit", 10)
    mem_type = args.get("type")
    mode = args.get("mode", "")

    # Project-scoped search: search project collection first, then agent
    project = getattr(_thread_local, 'project', None)
    if project and query:
        agent_id = ms.agent_id
        proj_dir = os.path.join(AGENTS_DIR, agent_id, "projects", project)
        if os.path.isdir(proj_dir):
            proj_store = MemoryStore(agent_id=f"{agent_id}/{project}", base_dir=proj_dir)
            # Also search ingested subdir
            ingest_dir = os.path.join(proj_dir, "ingested")
            proj_results = proj_store.recall(query, limit, mem_type)
            # Tag project results
            for r in proj_results:
                r["source_scope"] = "project"
            # Then agent-level results
            agent_results = ms.recall(query, max(2, limit - len(proj_results)), mem_type)
            for r in agent_results:
                r["source_scope"] = "agent"
            results = proj_results + agent_results
            # Always expand via graph relationships (follow related links 1 hop)
            if results:
                results = _graph_expand_results(results, proj_dir, ingest_dir,
                                                max_hops=2 if mode == "graph" else 1)
            for r in results:
                if r.get("content") and len(r["content"]) > 4000:
                    r["content"] = r["content"][:4000] + "..."
            return _ok({"query": query, "project": project, "results": results[:limit], "count": len(results[:limit])})

    if not query:
        results = ms.list_all(mem_type)
        return _ok({"query": "", "results": results, "count": len(results)})
    results = ms.recall(query, limit, mem_type)

    # Always expand via graph relationships (1 hop default, 2 hops for explicit graph mode)
    if results:
        agent_id = ms.agent_id
        agent_dir = os.path.join(AGENTS_DIR, agent_id)
        ingest_dir = os.path.join(agent_dir, "ingested")
        results = _graph_expand_results(results, agent_dir, ingest_dir,
                                        max_hops=2 if mode == "graph" else 1)

    for r in results:
        if r.get("content") and len(r["content"]) > 4000:
            r["content"] = r["content"][:4000] + "..."

    # --- Co-recall tracking (Mechanism 3) ---
    if query and len(results) >= 2:
        try:
            result_files = [os.path.basename(r.get("file_path", "")) for r in results if r.get("file_path")]
            agent_id = ms.agent_id
            agent_dir = os.path.join(AGENTS_DIR, agent_id)
            threading.Thread(
                target=_record_recall_cooccurrence,
                args=(result_files, agent_id, agent_dir),
                daemon=True,
            ).start()
        except Exception:
            pass  # Co-recall tracking is best-effort

    return _ok({"query": query, "results": results, "count": len(results)})


def tool_memory_delete(args: dict) -> str:
    """Delete a memory."""
    ms = _get_memory_store()
    if not ms:
        return _err("Memory store not initialized")
    name = args.get("name", "")
    if not name:
        return _err("memory_delete: name is required")
    result = ms.delete(name)
    return _ok(result)


def tool_memory_shared(args: dict) -> str:
    """Access shared memory — global (main) or team (team head) scope."""
    action = args.get("action", "recall")
    scope = args.get("scope", "global")

    # Determine which agent's memory to use
    if scope == "team":
        # Find the team head for the calling agent
        caller_id = getattr(_thread_local, "delegate_agent_id", None)
        if not caller_id:
            agent = getattr(_thread_local, 'current_agent', None) or _current_agent
            caller_id = agent.agent_id if agent else "main"
        team_info = _get_agent_team_info(caller_id)
        if not team_info:
            return _err("memory_shared: agent is not in any team — use scope='global' instead")
        team_head_id = team_info["head"]
        target_agent = AgentConfig(team_head_id)
        source_label = f"{team_info['name']} (team)"
    else:
        target_agent = AgentConfig("main")
        source_label = "main (shared)"

    shared_store = MemoryStore(agent_id=target_agent.agent_id, base_dir=target_agent.memory_dir)

    if action == "store":
        name = args.get("name", "")
        content = args.get("content", "")
        description = args.get("description", "")
        mem_type = args.get("type", "general")
        if not name or not content:
            return _err("memory_shared store: name and content are required")
        result = shared_store.store(name, content, description, mem_type)
        result["source"] = source_label
        return _ok(result)
    else:  # recall
        query = args.get("query", "")
        limit = args.get("limit", 10)
        mem_type = args.get("type")
        if not query:
            results = shared_store.list_all(mem_type)
        else:
            results = shared_store.recall(query, limit, mem_type)
            # Graph expansion on shared memory too
            if results:
                shared_dir = os.path.join(AGENTS_DIR, target_agent.agent_id)
                shared_ingest = os.path.join(shared_dir, "ingested")
                results = _graph_expand_results(results, shared_dir, shared_ingest, max_hops=1)
            for r in results:
                if r.get("content") and len(r["content"]) > 4000:
                    r["content"] = r["content"][:4000] + "..."
        return _ok({"query": query, "source": source_label, "results": results[:limit], "count": len(results[:limit])})


def tool_use_skill(args: dict) -> str:
    """Load a skill's instructions into context."""
    skill_name = args.get("skill", "")
    if not skill_name:
        return _err("use_skill: skill name is required")
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    if not agent:
        return _err("use_skill: no active agent")

    body = agent.load_skill(skill_name)
    if body is None:
        available = [s.get("slug", s["name"]) for s in agent.list_skills()]
        return _err(f"use_skill: skill '{skill_name}' not found. Available: {', '.join(available) or 'none'}")

    return _ok({"skill": skill_name, "instructions": body})
