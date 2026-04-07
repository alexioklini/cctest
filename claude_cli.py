#!/usr/bin/env python3
"""Brain Agent — Agentic CLI for interacting with LLM APIs."""

VERSION = "5.9.0"
VERSION_DATE = "2026-04-07"
CHANGELOG = [
    ("5.9.0", "2026-04-07", "Chat file attachments — attach files directly in the chat composer. Files are saved to a session-scoped temp directory on the server; the agent reads them on demand via read_document (PDF, DOCX, XLSX, PPTX, CSV) or read_file (text/code). Web UI: file input accepts 30+ extensions, binary formats (PDF, DOCX, XLSX, PPTX) sent as base64, text files as UTF-8. File preview chips in composer with remove buttons. Attached files shown on sent user messages with extension badges. Configurable vision model for image attachments (Settings → Server → Attachments). Documents tool group added to default agent config."),
    ("5.8.0", "2026-04-06", "Model management overhaul — models now have a configurable display_name (default = cleaned model ID). All UI surfaces show models as 'displayName (provider)' format — selectors, dropdowns, status bar, spinners, agent config. Models tab redesigned: grouped by provider with collapsible sections, sorted by display name, inline editable display names, per-model remove button. Manual model add form for providers without /models endpoint (model ID + provider + display name with datalist autocomplete). Code mode dropdown uses modelsConfig instead of flat model list. Provider tab badges use compact names. Backend: display_name field added to _match_known_model()."),
    ("5.7.0", "2026-04-03", "Token optimization suite — comprehensive per-agent token usage controls via new Tokens tab in agent config. Tool group filtering sends only relevant tools to the LLM (13 groups: core, memory, context, web, email, documents, delegation, code_graph, git, scheduler, mcp, skills, nodes). System prompt trimmed: tools.md reduced from 1500 to 400 tokens, memory summary cap configurable per agent. Anthropic prompt caching via cache_control on system prompt blocks. System prompt cached per-session (60s TTL) to avoid disk I/O on tool loops. Memory summary scheduled tasks restricted to memory-only tools (4 instead of 39). Compact threshold configurable per agent. SDK duplicate tools cleaned up. Kilo API provider added (OpenAI-compatible gateway). Context fill bar and manual compact button in chat footer. Background pipeline model selectors with fallback in Memory tab GUI. Fixed SDK token count inflating context display (was reporting API tokens_in instead of conversation estimate)."),
    ("5.6.0", "2026-03-31", "Code Mode overhaul — full-featured coding assistant experience. Folder browser GUI for project selection (breadcrumb navigation, lazy-loaded directory listing via /v1/files/tree). Fixed SSE streaming (two-line event:/data: format matching server output). Tool calls display identically to main chat (gear/check icons, expandable args/results, proper .open toggle). Streaming indicator with wave animation, model name, tool labels, elapsed timer, and stop button. Folder-based project system — sessions tagged with folder path, 'All Projects' expands to show discovered projects with session counts, selecting a project filters sessions. Session management: archive and delete buttons on hover in sidebar. Code mode sessions properly scoped by folder path via project field."),
    ("5.5.0", "2026-03-31", "Claude.ai-style Projects — full project workspace system modeled after Claude.ai. Projects list view with card grid, search, Your Projects/Archived tabs, sort by activity/name/created. Project detail view with back navigation, description (show more/less), chat composer, conversation list, and right panel with Instructions and Files sections. Custom instructions per project — editable via modal, stored in project.json, injected into system prompt for all project conversations. File upload via multipart form (replaced deprecated cgi module with manual boundary parser for Python 3.13+). Files displayed with document icons, deletable per file. Project-scoped conversations — sessions filtered by project field, new chats auto-scoped. Project CRUD — create modal with name/description/agent, archive (preserves data + QMD), delete (soft-delete to .trash, removes QMD collection). Context menus on project cards and chat items. API: GET/POST /v1/sessions?project=X for project-filtered session listing."),
    ("5.4.0", "2026-03-31", "Artifact system — Claude.ai-style artifact management. Files written with relative paths auto-land in agents/<name>/artifacts/<session_folder>/. Session-scoped SQLite registry with content snapshots (versioned blobs, up to 5MB per version). Resizable right panel with type-aware rendering: syntax-highlighted code (highlight.js), sandboxed HTML iframe, inline SVG, image display, rendered markdown. Version selector dropdown, copy/download/source-toggle actions. Artifact cards in chat messages (coral border, monitor icon) open panel on click. Artifacts excluded from QMD indexing and knowledge graph (not memory). Artifacts browse view in sidebar: full-page grid with content preview cards, type filter tabs (All/Code/HTML/Documents/Images/Markdown), agent filter chips, time-ago timestamps. Click-through from browse opens chat + artifact panel. API: GET /v1/artifacts, /v1/artifacts/browse, /v1/artifacts/<id>/content, /v1/artifacts/<id>/download."),
    ("5.3.0", "2026-03-31", "Claude.ai-style web UI + interactive agents. Complete UI rewrite: sidebar + multi-view layout (Welcome, Chat, Chats, Projects, Knowledge Graph, Customize) with Anthropic Sans/Serif/Mono fonts, warm light/dark themes. Tool call blocks show with full args during streaming and persist across page reloads (reconstructed from assistant message metadata). Tool display toggle works. Interactive mode: agents can ask clarifying questions via AskUserQuestion — sidecar intercepts with PreToolUse hook, emits user_input_needed SSE event, blocks until answer arrives via POST /answer/{query_id}. TUI renders questions with selectable options. New endpoints: memory CRUD, soul.md AI editing, MCP registry. Agent creation accepts model and display_name. Sidecar captures tool input_json_delta for full args on content_block_stop."),
    ("5.2.0", "2026-03-29", "Mission Control cockpit — the web UI is now a dashboard-first design inspired by mission control interfaces. Agent cards show live status, model, schedules, past actions (scrollable), projects, and cost. Chat and project views are full-screen modals with maximum screen space. Token Cost Feed table with per-agent breakdown. Consistent color palette (dark navy header, light cards, green/orange/purple accents) across cockpit, chat modal, config dialogs, and settings. Session cache for instant cockpit loads. Hover actions for archive/delete on sessions and projects. Team badges on agent cards. Agent ordering matches team hierarchy (main → teams → standalone). All chat input controls (attach, think, plan, tools toggle, refine) available in modal."),
    ("5.1.0", "2026-03-28", "Real-time streaming + Claude Code skills. Sidecar rewritten as REST API (POST /query, GET /events/{id}) — decouples event production from consumption for true token-by-token streaming. MCP tools served via /mcp JSON-RPC endpoint on the server. Hooks moved server-side into /mcp tools/call handler (SDK hook registration was the root cause of streaming buffering). Claude Code plugin integration: scan, browse, install, and toggle 121 CC plugins per agent via GUI. SDK integration audited against official docs: @tool decorator, allowed_tools wildcards, correct hook signatures."),
    ("5.0.0", "2026-03-28", "Full SDK migration complete — closed all remaining gaps from the Agent SDK transition. HTTP MCP server for 24 custom tools, chat summary + transcript indexing for SDK path, file change watcher, rate limiting + model fallback, trace spans + audit logging, all background tasks route through SDK, TUI + CLI one-shot + scheduled tasks route through sidecar with graceful direct-API fallback."),
    ("4.5.0", "2026-03-27", "Agent SDK integration — all agents now use the Anthropic Agent SDK (Claude Code) as the agentic loop backend. Multi-provider support: Claude via CLIProxyAPI (Max subscription), MiniMax, oMLX local models, and Gemini via CLIProxyAPI. Real-time token streaming via a lean sidecar process. SDK badge in status bar and message footers. Provider-aware env var routing. System prompt extracted into reusable _build_system_prompt()."),
    ("4.4.6", "2026-03-26", "Token consumption guardrails — base64 image data is stripped from tool results after processing, individual tool results are capped at 30K chars, accumulated results are compressed when they exceed 50K tokens, mid-turn compaction runs every 3 tool rounds, and CLIProxyAPI gets a tighter 8-round tool limit to protect the OAuth quota. Telegram dedup guard prevents duplicate sessions from 409 polling conflicts."),
    ("4.4.5", "2026-03-26", "File previews now support images (JPEG, PNG, GIF, WebP, SVG) and office documents (DOCX, XLSX, PPTX, PDF, CSV). Sidebar chat attachments show preview and download buttons, file counts are accurate from first load, and the list is stable across re-renders."),
    ("4.4.4", "2026-03-26", "Fixed sidebar file attachments disappearing when the accordion was opened — files created in a chat session are now fetched from a dedicated endpoint that includes the full message history, so compacted messages no longer hide previously created files."),
    ("4.4.3", "2026-03-26", "Fixed cost tracking showing $0 for all models — auto-discovery was writing cost_input=0 to config, which silently overrode the built-in Anthropic rate table. Config zeros are now treated as unset and fall through to correct rates. Added missing model IDs (Haiku 4.5, Sonnet 4.0, Opus 4.0) to the built-in table, plus prefix patterns as a catch-all for future versions. Historical costs for the past week were retroactively corrected."),
    ("4.4.2", "2026-03-26", "Smarter memory I/O — the nightly background pipeline now reads only what it needs from memory files instead of loading everything into RAM. The autodream consolidation run also reuses a single memory scan across all its passes, reducing redundant filesystem work."),
    ("4.4.1", "2026-03-26", "Significant token savings — background tasks that only produce structured output (memory deduplication, conflict detection, relationship discovery, context summarization) no longer receive the full tool schema on every call, saving ~8,000 tokens per invocation. Smart fallback chains route time-sensitive tasks to fast cloud models and nightly tasks to free local models."),
    ("4.4.0", "2026-03-26", "Reduced Anthropic token consumption across the board — shorter context windows, a capped memory summary injected only at the start of each turn, earlier context compaction, and configurable model selection for all background pipelines. GUI model selectors added for memory summary, relationship discovery, and autodream."),
    ("4.3.1", "2026-03-25", "Stability fixes — resolved a session corruption issue when model fallback was triggered mid-tool-loop. Rate limit and overload errors from the API now retry gracefully instead of switching models immediately. Shell commands now run in a login shell so your PATH and environment are always available."),
    ("4.3.0", "2026-03-24", "Autodream — a nightly memory consolidation pipeline that automatically deduplicates overlapping memories, flags stale ones, detects contradictions, and identifies reusable procedures worth turning into skills. Results are summarised in a health report with scoring. A new Memory Health dashboard in Settings shows per-agent stats and recall frequency."),
    ("4.2.0", "2026-03-23", "Smarter code graph — the knowledge graph now generates plain-English summaries for every function and class, classifies code into architecture layers (API, service, data, UI, util, test), and produces a guided reading order. New context fill indicator and manual compaction controls added to the chat footer."),
    ("4.1.0", "2026-03-23", "Chat reliability improvements — conversations no longer get corrupted when a tool loop fails mid-way. Partial responses are preserved on cancel or error. Message metadata (model, tokens, cost, thinking) is now persisted and restored when reopening a chat. Thinking depth can be controlled (off / low / medium / high)."),
    ("4.0.0", "2026-03-23", "Universal file support — agents can now read, write, and edit Excel, PowerPoint, CSV, images, and SVG files in addition to PDF and Word. A full code structure graph powered by Tree-sitter covers 14 languages, with tools to query relationships, trace call chains, and analyse the blast radius of any change."),
    ("3.7.0", "2026-03-23", "Extensible hooks — attach shell scripts to any tool call or file write event. Hooks can inspect, block, or react to agent actions without modifying core code. A hooks UI in agent config makes wiring them up straightforward."),
    ("3.6.0", "2026-03-22", "Lossless context — long conversations are no longer truncated into a flat summary. A DAG-based hierarchy preserves the full conversation tree, letting you search and drill back into any compacted segment with context_search and context_recall."),
    ("3.5.0", "2026-03-22", "Full-text chat search across all sessions, with semantic search via QMD and SQLite fallback. Knowledge graph relationship discovery upgraded to a two-stage pipeline — QMD finds candidates, an LLM classifies the actual relationship type."),
    ("3.4.0", "2026-03-22", "Remote nodes — Brain Agent can now distribute tasks to other machines running node.py. Nodes are managed from Settings with configurable tokens, allowed tools, and concurrency limits. macOS launchd install/uninstall built in."),
    ("3.3.0", "2026-03-22", "Richer projects — AI can directly edit notes using the same file tools it uses everywhere else. Chat transcripts are indexed for semantic search. Sessions now show LLM-generated summaries in the sidebar, and files created during a conversation appear as downloadable attachments."),
    ("3.2.0", "2026-03-22", "Project Notes — a full note-taking system inside Brain Agent, with a rich text editor, AI chat sidebar per note, and automatic knowledge graph integration. Notes live alongside agent memory and are searchable."),
    ("3.1.0", "2026-03-21", "Memories now form automatically from conversations — corrections, decisions, and references are detected and stored without prompting. The knowledge graph gains a visual canvas view. Sidebar redesigned around Projects and Chats. Continuous session summarisation added."),
    ("3.0.0", "2026-03-20", "Major platform expansion — provider fallback with retry, backup and restore, webhook and email notifications, full observability tracing, dynamic MCP client connections, image upload and vision support, and a multi-messaging adapter framework."),
    ("2.1.0", "2026-03-20", "Agent workflows — define multi-stage pipelines in YAML with optional human approval gates between stages. Web UI sidebar moved to the left with a consolidated status bar. Mobile layout improved."),
    ("2.0.0", "2026-03-20", "Projects — organise work into scoped contexts. Ingest PDFs, Word docs, web pages, and watched folders. The knowledge graph tracks relationships between documents and memories, and chat history is scoped per project."),
    ("1.7.0", "2026-03-20", "Plan mode for read-only review before executing. Web search results are cached to avoid redundant calls. Cost tracking and rate limiting added per agent. Custom slash commands and LLM-powered input refinement."),
    ("1.6.0", "2026-03-20", "TUI reaches feature parity with the Web UI — 30+ slash commands, autocomplete popup menus in both interfaces."),
    ("1.5.3", "2026-03-20", "Concurrency hardening — thread-safe agent context, per-collection QMD debouncing, YAML-safe memory frontmatter, hash suffixes on memory filenames to prevent collisions, concurrent scheduler execution."),
    ("1.5.2", "2026-03-20", "Fixed memory summary refresh and QMD index path normalisation. QMD collection health stats surfaced in the settings UI."),
    ("1.5.1", "2026-03-18", "MiniMax provider support. Add Model UI. Fixed a QMD session leak and an issue where shared memory returned only metadata instead of full content. Telegram bot runs in-process."),
    ("1.5.0", "2026-03-18", "Settings dashboard — a full admin panel for Server, QMD, Models, Telegram, and Providers. Agent activity indicators, QMD document browser with per-file index health, smart model routing, and a self-healing QMD index keeper."),
    ("1.4.0", "2026-03-17", "QMD hybrid memory search (BM25 + vector + LLM reranking) replaces SQLite FTS5. Improved SSE error handling and server resilience."),
    ("1.2.0", "2026-03-16", "Multi-provider routing, Gmail integration, scheduler dashboard, SQLite resilience improvements, Cloudflare Zero Trust deployment."),
    ("1.1.0", "2026-03-14", "MCP support — connect any MCP server via stdio or SSE transport, scoped per agent or globally."),
    ("1.0.0", "2026-03-14", "Async agent delegation — agents run in background threads with task status tracking and cancellation."),
    ("0.9.0", "2026-03-14", "Skills system — SKILL.md files loaded on demand, available per agent or globally."),
    ("0.8.0", "2026-03-14", "Multi-agent system — define agents with soul.md personalities, delegate tasks between them, switch with /agent."),
    ("0.7.0", "2026-03-14", "Persistent memory — store and recall information across sessions with per-agent isolation."),
    ("0.6.0", "2026-03-14", "Context window management with automatic compaction when the window fills."),
    ("0.5.0", "2026-03-14", "Full tool suite — file read/write/edit, shell execution, web search, and fetch."),
    ("0.4.0", "2026-03-14", "Escape to cancel in-flight requests. Dynamic terminal rendering and startup greeting."),
    ("0.3.0", "2026-03-13", "Exa web search with agentic tool-use loop."),
    ("0.2.0", "2026-03-12", "Interactive TUI with spinner, markdown rendering, and model switching."),
    ("0.1.0", "2026-03-10", "Initial release — streaming chat, model fallback, SSE parsing."),
]

import argparse
import collections
import datetime
import fnmatch
import logging
import glob as globmod
import json
import os
import random
import re
import select
import signal
import shutil
import socket
import subprocess
import sys
import termios
import threading
import time
import tty
import urllib.request
import urllib.error


# --- Web Result Cache ---

class WebCache:
    """Thread-safe LRU cache with TTL for web results."""

    def __init__(self, max_entries: int = 200, ttl: int = 900):
        self._cache = collections.OrderedDict()
        self._lock = threading.Lock()
        self.max_entries = max_entries
        self.ttl = ttl
        self.hits = 0
        self.misses = 0

    def get(self, key: str) -> dict | None:
        with self._lock:
            entry = self._cache.get(key)
            if entry is None:
                self.misses += 1
                return None
            ts, value = entry
            if time.time() - ts > self.ttl:
                del self._cache[key]
                self.misses += 1
                return None
            self._cache.move_to_end(key)
            self.hits += 1
            return value

    def put(self, key: str, value: dict):
        with self._lock:
            if key in self._cache:
                del self._cache[key]
            self._cache[key] = (time.time(), value)
            while len(self._cache) > self.max_entries:
                self._cache.popitem(last=False)

    def clear(self):
        with self._lock:
            self._cache.clear()
            self.hits = 0
            self.misses = 0

    def stats(self) -> dict:
        with self._lock:
            return {
                "entries": len(self._cache),
                "max_entries": self.max_entries,
                "ttl": self.ttl,
                "hits": self.hits,
                "misses": self.misses,
                "hit_rate": round(self.hits / max(1, self.hits + self.misses) * 100, 1),
            }


_web_cache = WebCache()


# --- Plan Mode ---

READONLY_TOOLS = frozenset({
    "read_file", "list_directory", "search_files", "web_fetch", "exa_search",
    "memory_recall", "memory_shared", "task_status", "list_nodes",
    "context_search", "context_detail", "context_recall", "schedule_list",
    "schedule_history", "use_skill", "gmail_inbox", "gmail_read", "gmail_search",
    "read_document",
    "code_graph_build", "code_graph_query",
})

PLAN_MODE_PROMPT = (
    "\n\nPLAN MODE ACTIVE: You are in read-only planning mode. "
    "You may ONLY use read-only tools (read_file, list_directory, search_files, "
    "web_fetch, exa_search, memory_recall, memory_shared, task_status, etc.). "
    "Do NOT attempt to write files, execute commands, store memory, send emails, "
    "or delegate tasks. Instead, describe a detailed plan of what you WOULD do, "
    "including specific file paths, commands, and steps.\n"
)


# --- Tool Definitions ---

TOOL_DEFINITIONS = [
    {
        "name": "read_file",
        "description": (
            "Read the contents of a file. Returns the full text content. "
            "Use offset and limit to read a specific range of lines from large files."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute or relative file path to read"},
                "offset": {"type": "integer", "description": "Line number to start reading from (1-based, default: 1)"},
                "limit": {"type": "integer", "description": "Maximum number of lines to read (default: all)"},
                "node": {"type": "string", "description": "Remote node name or 'tag:NAME' to execute on a remote node instead of locally"},
            },
            "required": ["path"],
        },
    },
    {
        "name": "write_file",
        "description": (
            "Create a new file or overwrite an existing file with the given content. "
            "Creates parent directories automatically if they don't exist."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "File path to write to"},
                "content": {"type": "string", "description": "The full content to write to the file"},
                "node": {"type": "string", "description": "Remote node name or 'tag:NAME' to execute on a remote node instead of locally"},
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "edit_file",
        "description": (
            "Edit an existing file by replacing an exact string match with new content. "
            "The old_string must match exactly (including whitespace/indentation). "
            "Use replace_all=true to replace every occurrence."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "File path to edit"},
                "old_string": {"type": "string", "description": "Exact string to find and replace"},
                "new_string": {"type": "string", "description": "Replacement string"},
                "replace_all": {"type": "boolean", "description": "Replace all occurrences (default: false)"},
            },
            "required": ["path", "old_string", "new_string"],
        },
    },
    {
        "name": "list_directory",
        "description": (
            "List files and directories at a given path. "
            "Supports glob patterns (e.g. '*.py', '**/*.js'). "
            "Returns file names, sizes, and types."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Directory path to list (default: current directory)"},
                "pattern": {"type": "string", "description": "Glob pattern to filter results (e.g. '*.py', '**/*.ts')"},
                "recursive": {"type": "boolean", "description": "List recursively (default: false)"},
                "node": {"type": "string", "description": "Remote node name or 'tag:NAME' to execute on a remote node instead of locally"},
            },
            "required": [],
        },
    },
    {
        "name": "search_files",
        "description": (
            "Search for a regex pattern across files. Returns matching lines with file paths and line numbers. "
            "Similar to grep/ripgrep. Use glob to filter which files to search."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "Regex pattern to search for"},
                "path": {"type": "string", "description": "Directory or file to search in (default: current directory)"},
                "glob": {"type": "string", "description": "Glob pattern to filter files (e.g. '*.py')"},
                "case_insensitive": {"type": "boolean", "description": "Case-insensitive search (default: false)"},
                "max_results": {"type": "integer", "description": "Maximum number of matches to return (default: 50)"},
            },
            "required": ["pattern"],
        },
    },
    {
        "name": "execute_command",
        "description": (
            "Execute a shell command and return its output (stdout + stderr). "
            "Commands run in the current working directory with no TTY (non-interactive). "
            "IMPORTANT: Only use non-interactive commands. For example use 'top -l 1' (not 'top'), "
            "'ps aux' (not 'htop'), 'cat' (not 'less'). "
            "Use this for: running scripts, git commands, package managers, compiling, testing, "
            "system administration, or any shell operation."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "command": {"type": "string", "description": "The shell command to execute"},
                "cwd": {"type": "string", "description": "Working directory for the command (default: current directory)"},
                "timeout": {"type": "integer", "description": "Timeout in seconds (default: 120)"},
                "node": {"type": "string", "description": "Remote node name or 'tag:NAME' to execute on a remote node instead of locally"},
            },
            "required": ["command"],
        },
    },
    {
        "name": "web_fetch",
        "description": (
            "Fetch content from a URL. Returns the response body as text. "
            "Works with web pages, APIs, raw files, etc."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "The URL to fetch"},
                "method": {"type": "string", "description": "HTTP method (default: GET)", "enum": ["GET", "POST", "PUT", "DELETE", "PATCH"]},
                "headers": {"type": "object", "description": "Additional HTTP headers as key-value pairs"},
                "body": {"type": "string", "description": "Request body (for POST/PUT/PATCH)"},
                "max_length": {"type": "integer", "description": "Max response length in characters (default: 50000)"},
            },
            "required": ["url"],
        },
    },
    {
        "name": "gmail_inbox",
        "description": "List recent emails from Gmail inbox. Returns subject, from, date for each email.",
        "input_schema": {
            "type": "object",
            "properties": {
                "limit": {"type": "integer", "description": "Number of emails to return (default: 10)"},
                "folder": {"type": "string", "description": "Mailbox folder (default: INBOX)"},
            },
            "required": [],
        },
    },
    {
        "name": "gmail_read",
        "description": "Read a specific email by its ID. Returns full body, attachments list, headers.",
        "input_schema": {
            "type": "object",
            "properties": {
                "id": {"type": "string", "description": "Email ID from gmail_inbox or gmail_search"},
                "folder": {"type": "string", "description": "Mailbox folder (default: INBOX)"},
            },
            "required": ["id"],
        },
    },
    {
        "name": "gmail_search",
        "description": "Search emails using Gmail search syntax (from:, subject:, is:unread, after:, has:attachment, etc).",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Gmail search query"},
                "limit": {"type": "integer", "description": "Max results (default: 10)"},
            },
            "required": ["query"],
        },
    },
    {
        "name": "gmail_send",
        "description": "Send an email via Gmail.",
        "input_schema": {
            "type": "object",
            "properties": {
                "to": {"type": "string", "description": "Recipient email address"},
                "subject": {"type": "string", "description": "Email subject"},
                "body": {"type": "string", "description": "Email body (plain text)"},
                "cc": {"type": "string", "description": "CC email address (optional)"},
            },
            "required": ["to", "subject", "body"],
        },
    },
    {
        "name": "gmail_reply",
        "description": "Reply to an existing email by its ID. Preserves threading.",
        "input_schema": {
            "type": "object",
            "properties": {
                "id": {"type": "string", "description": "Email ID to reply to"},
                "body": {"type": "string", "description": "Reply body (plain text)"},
            },
            "required": ["id", "body"],
        },
    },
    {
        "name": "exa_search",
        "description": (
            "Search the web using Exa AI for current, relevant information. "
            "Use this tool whenever the user asks to search the web, look something up, "
            "find recent news, or get current information about any topic."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {
                    "type": "string",
                    "description": "The search query or topic to look up",
                },
                "num_results": {
                    "type": "integer",
                    "description": "Number of search results to return (default: 5)",
                    "minimum": 1,
                    "maximum": 20,
                },
                "category": {
                    "type": "string",
                    "description": "Optional category: news, research paper, tweet, company, people",
                    "enum": ["news", "research paper", "tweet", "company", "people"],
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "memory_store",
        "description": (
            "Store a memory for later recall. Use this to remember important information, "
            "user preferences, decisions, project context, or anything that should persist "
            "across conversations. Memories are searchable by keyword."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "name": {"type": "string", "description": "Short unique name for this memory (used as identifier)"},
                "content": {"type": "string", "description": "The memory content to store"},
                "description": {"type": "string", "description": "One-line description for search indexing"},
                "type": {"type": "string", "description": "Memory type", "enum": ["user", "project", "feedback", "reference", "general"]},
            },
            "required": ["name", "content"],
        },
    },
    {
        "name": "memory_recall",
        "description": (
            "Search and recall stored memories. Use this when you need context from previous "
            "conversations, user preferences, project decisions, or any previously stored information. "
            "Returns matching memories ranked by relevance, plus related memories discovered via "
            "the knowledge graph (automatically follows relationship links for richer context)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query (keywords). Leave empty to list all memories."},
                "limit": {"type": "integer", "description": "Max results (default: 10)"},
                "type": {"type": "string", "description": "Filter by memory type", "enum": ["user", "project", "feedback", "reference", "general"]},
            },
            "required": [],
        },
    },
    {
        "name": "memory_delete",
        "description": "Delete a stored memory by name.",
        "input_schema": {
            "type": "object",
            "properties": {
                "name": {"type": "string", "description": "Name of the memory to delete"},
            },
            "required": ["name"],
        },
    },
    {
        "name": "memory_shared",
        "description": (
            "Access shared memory at global or team scope. Global scope (default) accesses the main agent's "
            "memory containing infrastructure, user preferences, and project-wide decisions. "
            "Team scope accesses the team head's memory for team-level knowledge. "
            "Use this when you need context that isn't in your own memory."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "description": "Action to perform", "enum": ["recall", "store"]},
                "scope": {"type": "string", "description": "Memory scope: 'global' for main agent (default), 'team' for team head's memory", "enum": ["global", "team"], "default": "global"},
                "query": {"type": "string", "description": "Search query for recall, or empty to list all"},
                "name": {"type": "string", "description": "Memory name (required for store)"},
                "content": {"type": "string", "description": "Content to store (required for store)"},
                "description": {"type": "string", "description": "One-line description (for store)"},
                "type": {"type": "string", "description": "Memory type", "enum": ["user", "project", "feedback", "reference", "general"]},
                "limit": {"type": "integer", "description": "Max results for recall (default: 10)"},
            },
            "required": ["action"],
        },
    },
    {
        "name": "context_search",
        "description": "Search through compacted conversation history by keyword. Returns matching message excerpts from earlier in the conversation that have been summarized away.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search keyword or phrase"},
                "limit": {"type": "integer", "description": "Max results (default: 10)"},
            },
            "required": ["query"],
        },
    },
    {
        "name": "context_detail",
        "description": "Expand a specific context summary to see the original messages it was created from. Use summary IDs from the conversation context header.",
        "input_schema": {
            "type": "object",
            "properties": {
                "summary_id": {"type": "string", "description": "The summary ID to expand"},
            },
            "required": ["summary_id"],
        },
    },
    {
        "name": "context_recall",
        "description": "Deep recall: search compacted conversation history and get a focused answer about a specific topic from earlier in the conversation. Uses a sub-LLM call to analyze original messages.",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "What to recall from earlier conversation"},
            },
            "required": ["query"],
        },
    },
    {
        "name": "read_document",
        "description": (
            "Format-aware document reader for PDF, DOCX, XLSX, PPTX, CSV/TSV, and image files. "
            "Returns structured content: PDF pages, DOCX paragraphs/tables, XLSX sheets as markdown tables, "
            "PPTX slides with notes, CSV as markdown table, image metadata + vision description. "
            "For unknown extensions, falls back to plain text read."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "File path to the document"},
                "sheet": {"type": "string", "description": "Sheet name for XLSX (default: all sheets)"},
                "pages": {"type": "string", "description": "Page range for PDF, e.g. '1-5' or '1,3,7'"},
                "slides": {"type": "string", "description": "Slide range for PPTX, e.g. '1-10' or '2,5'"},
            },
            "required": ["path"],
        },
    },
    {
        "name": "write_document",
        "description": (
            "Create a new document from markdown content. Dispatches by file extension: "
            ".docx (headings, tables, bold/italic), .xlsx (markdown tables to sheets), "
            ".pptx (# sections to slides), .pdf (basic formatted PDF via reportlab)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Output file path (extension determines format)"},
                "content": {"type": "string", "description": "Markdown content to convert into the document"},
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "edit_document",
        "description": (
            "Targeted edits to existing documents. Actions by format: "
            "DOCX: replace_text (find/replace in paragraphs). "
            "XLSX: update_cell (sheet, cell, value), add_row (sheet, values). "
            "PPTX: update_slide (slide_index, title, body), add_slide (title, body)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the document to edit"},
                "action": {
                    "type": "string",
                    "description": "Edit action to perform",
                    "enum": ["replace_text", "update_cell", "add_row", "update_slide", "add_slide"],
                },
                "params": {
                    "type": "object",
                    "description": (
                        "Action-specific parameters. "
                        "replace_text: {old_text, new_text}. "
                        "update_cell: {sheet, cell, value}. "
                        "add_row: {sheet, values (array)}. "
                        "update_slide: {slide_index (1-based), title, body}. "
                        "add_slide: {title, body}."
                    ),
                },
            },
            "required": ["path", "action", "params"],
        },
    },
    {
        "name": "delegate_task",
        "description": (
            "Delegate a task to another agent. Runs in a background thread with its own context. "
            "By default waits for result (wait=true). Set wait=false for async execution, "
            "then use task_status to poll for completion."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "agent": {"type": "string", "description": "Target agent ID (e.g. 'research', 'health')"},
                "task": {"type": "string", "description": "Task description for the target agent"},
                "wait": {"type": "boolean", "description": "Wait for result (default: true). Set false for async."},
                "model": {"type": "string", "description": "Override model for this task (optional)"},
            },
            "required": ["agent", "task"],
        },
    },
    {
        "name": "task_status",
        "description": (
            "Check status of background tasks. Call with task_id to check a specific task, "
            "or without to list all tasks. Returns status (running/completed/cancelled/error) and result."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "task_id": {"type": "string", "description": "Task ID to check (optional, lists all if empty)"},
            },
            "required": [],
        },
    },
    {
        "name": "task_cancel",
        "description": "Cancel a running background task.",
        "input_schema": {
            "type": "object",
            "properties": {
                "task_id": {"type": "string", "description": "Task ID to cancel"},
            },
            "required": ["task_id"],
        },
    },
    {
        "name": "use_skill",
        "description": (
            "Load a skill's instructions into context. Skills provide specialized knowledge "
            "for specific tasks (e.g. github, docker, swift). Call this BEFORE performing a task "
            "that matches a skill. The skill's instructions will be returned as text — follow them."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "skill": {"type": "string", "description": "Name of the skill to load"},
            },
            "required": ["skill"],
        },
    },
    {
        "name": "schedule_list",
        "description": "List all scheduled tasks with their status, next run time, and configuration.",
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "list_nodes",
        "description": (
            "List all registered remote nodes with their status, hostname, OS, tags, "
            "allowed tools, and resource usage. Use this to check what remote nodes are "
            "available before routing commands to them."
        ),
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "schedule_history",
        "description": "Get execution history for scheduled tasks. Shows status, results, and timestamps.",
        "input_schema": {
            "type": "object",
            "properties": {
                "name": {"type": "string", "description": "Filter by schedule name (optional)"},
                "limit": {"type": "integer", "description": "Max results (default: 20)"},
            },
            "required": [],
        },
    },
    {
        "name": "mcp_connect",
        "description": (
            "Connect to an MCP server at runtime. Discovers tools from the server and makes them "
            "available as mcp_<name>_<tool> tools. Use transport='sse' for HTTP servers, 'stdio' for local commands."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "MCP server URL (for SSE) or command (for stdio)"},
                "name": {"type": "string", "description": "Friendly name for this connection"},
                "transport": {"type": "string", "description": "Transport type: 'sse' (default) or 'stdio'", "enum": ["sse", "stdio"]},
                "persist": {"type": "boolean", "description": "Save to mcp.json for reconnect on restart (default: false)"},
            },
            "required": ["url", "name"],
        },
    },
    {
        "name": "mcp_disconnect",
        "description": "Disconnect from a runtime MCP server. Its tools will no longer be available.",
        "input_schema": {
            "type": "object",
            "properties": {
                "name": {"type": "string", "description": "Name of the MCP server to disconnect"},
            },
            "required": ["name"],
        },
    },
    {
        "name": "mcp_servers",
        "description": "List all connected MCP servers with their tools and status.",
        "input_schema": {
            "type": "object",
            "properties": {},
            "required": [],
        },
    },
    {
        "name": "code_graph_build",
        "description": (
            "Build or rebuild the code structure graph for a directory. Parses source files "
            "using Tree-sitter AST parsing to extract functions, classes, imports, and call "
            "relationships. Supports Python, JavaScript, TypeScript, Go, Rust, Java, and more. "
            "Use incremental=true (default) to only re-parse changed files."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Directory to parse (absolute path)"},
                "incremental": {"type": "boolean", "description": "Only re-parse changed files (default: true)"},
            },
            "required": ["path"],
        },
    },
    {
        "name": "code_graph_query",
        "description": (
            "Query the code structure graph for structural relationships. Find callers/callees "
            "of a function, imports, inheritance, test coverage, and more. Build the graph first "
            "with code_graph_build."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query_type": {
                    "type": "string",
                    "enum": ["callers_of", "callees_of", "imports_of", "importers_of",
                             "tests_for", "inheritors_of", "children_of", "file_summary"],
                    "description": "Type of structural query",
                },
                "target": {"type": "string", "description": "Qualified name or function/class name to query"},
                "limit": {"type": "integer", "description": "Max results (default: 20)"},
            },
            "required": ["query_type", "target"],
        },
    },
    {
        "name": "code_graph_impact",
        "description": (
            "Blast-radius analysis: given a list of changed files, find all functions, classes, "
            "and files that could be affected. Uses BFS traversal of the code graph. "
            "Build the graph first with code_graph_build."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "files": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": "List of changed file paths",
                },
                "depth": {"type": "integer", "description": "Max traversal depth (default: 2)"},
            },
            "required": ["files"],
        },
    },
    {
        "name": "code_graph_enhance",
        "description": (
            "Enhance the code graph with LLM-generated summaries, architecture layer classification, "
            "and a guided tour. Actions: 'all' (default), 'summaries' (LLM descriptions per function/class), "
            "'layers' (classify as api/service/data/ui/util/test), 'tour' (dependency-ordered walkthrough)."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["all", "summaries", "layers", "tour"],
                    "description": "What to generate (default: all)",
                },
                "batch_size": {"type": "integer", "description": "Max files to summarize per run (default: 20)"},
                "root_dir": {"type": "string", "description": "Root directory for tour (default: last build dir)"},
            },
            "required": [],
        },
    },
    {
        "name": "git_command",
        "description": (
            "Execute git operations with structured output. Actions:\n"
            "- status: working tree status (modified, staged, untracked files)\n"
            "- diff: show changes (optional file path, staged=true for staged only)\n"
            "- log: commit history (limit, author, since, path filters)\n"
            "- branch: list/create/switch branches (name, create=true, switch=true)\n"
            "- commit: create commit (message required, files=[] to stage specific files, all=true for -a)\n"
            "- stash: stash/pop/list (sub_action: save/pop/list/drop)\n"
            "- blame: annotate file lines (path, line_start, line_end)\n"
            "- show: show commit details (ref)\n"
            "- tag: list/create tags (name, message)\n"
            "- remote: list remotes or show remote info"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["status", "diff", "log", "branch", "commit", "stash", "blame", "show", "tag", "remote"],
                    "description": "Git operation to perform",
                },
                "message": {"type": "string", "description": "Commit/tag message"},
                "files": {"type": "array", "items": {"type": "string"}, "description": "Specific files to stage/diff/blame"},
                "path": {"type": "string", "description": "File path for diff/blame/log"},
                "name": {"type": "string", "description": "Branch/tag name"},
                "ref": {"type": "string", "description": "Commit ref for show/diff (default: HEAD)"},
                "limit": {"type": "integer", "description": "Max entries for log (default: 20)"},
                "author": {"type": "string", "description": "Filter log by author"},
                "since": {"type": "string", "description": "Filter log since date (e.g., '1 week ago')"},
                "staged": {"type": "boolean", "description": "Show only staged changes for diff"},
                "create": {"type": "boolean", "description": "Create new branch/tag"},
                "switch": {"type": "boolean", "description": "Switch to branch"},
                "all": {"type": "boolean", "description": "Stage all changes for commit (-a)"},
                "sub_action": {"type": "string", "description": "Sub-action for stash (save/pop/list/drop)"},
                "line_start": {"type": "integer", "description": "Start line for blame"},
                "line_end": {"type": "integer", "description": "End line for blame"},
            },
            "required": ["action"],
        },
    },
    {
        "name": "github_command",
        "description": (
            "Interact with GitHub via the gh CLI. Requires gh to be installed and authenticated. Actions:\n"
            "- pr_list: list open PRs (limit, state, author)\n"
            "- pr_create: create PR (title, body, base, head, draft)\n"
            "- pr_view: view PR details (number)\n"
            "- pr_merge: merge a PR (number, method=merge|squash|rebase)\n"
            "- pr_review: list PR reviews/comments (number)\n"
            "- issue_list: list issues (limit, state, labels)\n"
            "- issue_create: create issue (title, body, labels)\n"
            "- issue_view: view issue details (number)\n"
            "- repo_view: show repo info\n"
            "- release_list: list releases\n"
            "- workflow_list: list GitHub Actions workflows\n"
            "- workflow_run: view workflow run status (run_id)\n"
            "- api: raw GitHub API call (endpoint, method)"
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {
                    "type": "string",
                    "enum": ["pr_list", "pr_create", "pr_view", "pr_merge", "pr_review",
                             "issue_list", "issue_create", "issue_view",
                             "repo_view", "release_list", "workflow_list", "workflow_run", "api"],
                    "description": "GitHub operation to perform",
                },
                "number": {"type": "integer", "description": "PR or issue number"},
                "title": {"type": "string", "description": "PR/issue title"},
                "body": {"type": "string", "description": "PR/issue body"},
                "base": {"type": "string", "description": "Base branch for PR (default: main)"},
                "head": {"type": "string", "description": "Head branch for PR"},
                "draft": {"type": "boolean", "description": "Create PR as draft"},
                "method": {"type": "string", "description": "Merge method (merge/squash/rebase)"},
                "state": {"type": "string", "description": "Filter by state (open/closed/all)"},
                "labels": {"type": "string", "description": "Comma-separated labels"},
                "author": {"type": "string", "description": "Filter by author"},
                "limit": {"type": "integer", "description": "Max results (default: 20)"},
                "run_id": {"type": "string", "description": "Workflow run ID"},
                "endpoint": {"type": "string", "description": "API endpoint for raw call (e.g., repos/{owner}/{repo}/issues)"},
                "api_method": {"type": "string", "description": "HTTP method for API call (GET/POST/PATCH)"},
            },
            "required": ["action"],
        },
    },
    {
        "name": "tool_search",
        "description": (
            "Search for available tools by name or description. Use this when you need a "
            "tool that isn't in your current tool list. Returns matching tool schemas that "
            "will be available on subsequent turns. Useful when MCP tools are deferred."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query to match against tool names and descriptions"},
                "max_results": {"type": "integer", "description": "Maximum results to return (default: 5)"},
            },
            "required": ["query"],
        },
    },
]

# Build OpenAI-compatible format automatically
TOOL_DEFINITIONS_OPENAI = []
for _td in TOOL_DEFINITIONS:
    TOOL_DEFINITIONS_OPENAI.append({
        "type": "function",
        "function": {
            "name": _td["name"],
            "description": _td["description"],
            "parameters": {
                "type": _td["input_schema"]["type"],
                "properties": _td["input_schema"]["properties"],
                "required": _td["input_schema"].get("required", []),
            },
        },
    })

# Tool name → definition index for fast lookup
_TOOL_DEF_INDEX = {td["name"]: td for td in TOOL_DEFINITIONS}
_TOOL_DEF_OPENAI_INDEX = {td["function"]["name"]: td for td in TOOL_DEFINITIONS_OPENAI}

# Tool groups for per-agent filtering (agents can specify groups or individual tool names)
TOOL_GROUPS = {
    "core": {"read_file", "write_file", "edit_file", "list_directory", "search_files",
             "execute_command"},
    "memory": {"memory_store", "memory_recall", "memory_delete", "memory_shared"},
    "context": {"context_search", "context_detail", "context_recall"},
    "web": {"web_fetch", "exa_search"},
    "email": {"gmail_inbox", "gmail_read", "gmail_search", "gmail_send", "gmail_reply"},
    "documents": {"read_document", "write_document", "edit_document"},
    "delegation": {"delegate_task", "task_status", "task_cancel"},
    "code_graph": {"code_graph_build", "code_graph_query", "code_graph_impact",
                   "code_graph_enhance"},
    "git": {"git_command", "github_command"},
    "scheduler": {"schedule_list", "schedule_history"},
    "mcp": {"mcp_connect", "mcp_disconnect", "mcp_servers"},
    "skills": {"use_skill"},
    "nodes": {"list_nodes"},
}

# Default tool groups included for all agents (if no explicit config)
DEFAULT_TOOL_GROUPS = {"core", "memory", "context", "web", "delegation", "git", "skills",
                       "nodes", "scheduler", "mcp"}


TOKEN_CONFIG_DEFAULTS = {
    "tool_groups": None,           # None = all tools, list = specific groups from TOOL_GROUPS
    "extra_tools": None,           # Additional individual tool names beyond groups
    "include_tools_guide": True,   # Inject tools.md into system prompt
    "include_memory_summary": True, # Inject memory summary into system prompt
    "memory_summary_cap": 3000,    # Max chars for memory summary injection
    "compact_threshold": None,     # None = use default (0.60), float = override
    "prompt_caching": True,        # Use Anthropic cache_control on system prompt
    "scheduled_task_tools": True,  # Include full tool schema in scheduled tasks
}


def _get_token_config(agent_id: str | None = None) -> dict:
    """Get token optimization config for an agent, merged with defaults."""
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    if not agent:
        return dict(TOKEN_CONFIG_DEFAULTS)
    cfg = agent.config.get("token_config", {})
    result = dict(TOKEN_CONFIG_DEFAULTS)
    if isinstance(cfg, dict):
        for k in TOKEN_CONFIG_DEFAULTS:
            if k in cfg:
                result[k] = cfg[k]
    return result


def _get_agent_tool_names(agent_id: str | None = None) -> set[str] | None:
    """Get the set of allowed tool names for an agent based on its config.
    Returns None if no filtering is configured (all tools allowed)."""
    tcfg = _get_token_config(agent_id)
    tool_groups = tcfg.get("tool_groups")
    extra_tools = tcfg.get("extra_tools")
    if not tool_groups and not extra_tools:
        return None  # No filtering configured — all tools
    names = set()
    if tool_groups:
        for g in tool_groups:
            names.update(TOOL_GROUPS.get(g, set()))
    if extra_tools:
        names.update(extra_tools)
    return names


def _should_defer_mcp(defer_setting, mcp_tools: list[dict], model: str,
                      is_openai: bool = False) -> bool:
    """Decide whether to defer MCP tool schemas.

    defer_setting: True (always defer), False (never), "auto" (defer when MCP tokens > 10% of context)
    """
    if defer_setting is True:
        return bool(mcp_tools)
    if defer_setting is False:
        return False
    # "auto" mode: defer when MCP tool schemas would exceed 10% of context window
    if not mcp_tools:
        return False
    mcp_schema_chars = sum(len(json.dumps(t)) for t in mcp_tools)
    mcp_schema_tokens = mcp_schema_chars // 4
    max_ctx = get_model_max_context(model)
    threshold = max_ctx * 0.10  # 10% of context window
    return mcp_schema_tokens > threshold


def _filter_tools(tool_list: list[dict], allowed: set[str] | None,
                  is_openai: bool = False) -> list[dict]:
    """Filter a tool definition list to only include allowed tools.
    Returns tools sorted by name for prompt cache stability."""
    if allowed is None:
        filtered = list(tool_list)
    elif is_openai:
        filtered = [t for t in tool_list if t["function"]["name"] in allowed]
    else:
        filtered = [t for t in tool_list if t["name"] in allowed]
    # Sort deterministically by tool name for Anthropic prompt cache stability.
    # Consistent ordering prevents cache misses when tool list is assembled in different order.
    if is_openai:
        filtered.sort(key=lambda t: t.get("function", {}).get("name", ""))
    else:
        filtered.sort(key=lambda t: t.get("name", ""))
    return filtered


# --- Tool Execution ---

def _ok(result: dict) -> str:
    return json.dumps(result, ensure_ascii=False)


def _err(msg: str) -> str:
    return json.dumps({"error": msg}, ensure_ascii=False)


def _route_to_node(tool_name: str, args: dict) -> str | None:
    """If args contain 'node', route to remote node via server API. Returns result string or None for local."""
    node = args.pop("node", None)
    if not node:
        return None
    try:
        import urllib.request
        body = json.dumps({"node": node, "tool": tool_name, "params": args}).encode("utf-8")
        req = urllib.request.Request(
            "http://127.0.0.1:8420/v1/nodes/execute",
            data=body,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=300) as resp:
            result = json.loads(resp.read().decode("utf-8"))
        if "error" in result:
            return _err(f"Node '{node}': {result['error']}")
        return _ok(result)
    except Exception as e:
        return _err(f"Node routing error: {e}")


def tool_read_file(args: dict) -> str:
    node_result = _route_to_node("read_file", args)
    if node_result is not None:
        return node_result
    path = args.get("path", "")
    offset = args.get("offset", 1)
    limit = args.get("limit", 400)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        with open(path, "r", errors="replace") as f:
            lines = f.readlines()
        total = len(lines)
        start = max(0, offset - 1)
        end = start + limit if limit else total
        selected = lines[start:end]
        # Number lines
        numbered = []
        for i, line in enumerate(selected, start=start + 1):
            numbered.append(f"{i:>6}\t{line.rstrip()}")
        content = "\n".join(numbered)
        return _ok({"path": path, "total_lines": total, "showing": f"{start+1}-{min(end, total)}", "content": content})
    except Exception as e:
        return _err(f"read_file: {e}")


def _maybe_qmd_reindex(path: str) -> None:
    """If path is a .md file inside an agent dir, trigger debounced QMD reindex."""
    if not path.endswith(".md"):
        return
    agents_dir = os.path.realpath(AGENTS_DIR)
    real_path = os.path.realpath(path)
    if not real_path.startswith(agents_dir + os.sep):
        return
    rel = real_path[len(agents_dir) + 1:]
    collection = rel.split(os.sep)[0]
    _qmd_debounced_embed(collection)


def _get_artifact_session_folder(session_id: str) -> str:
    """Return session folder name for artifacts: <date>_<session_prefix>"""
    cache_key = f"_artifact_folder_{session_id}"
    cached = getattr(_thread_local, cache_key, None)
    if cached:
        return cached
    from datetime import datetime as _dt
    folder = f"{_dt.now().strftime('%Y-%m-%d')}_{session_id[:8]}"
    setattr(_thread_local, cache_key, folder)
    return folder


def tool_write_file(args: dict) -> str:
    node_result = _route_to_node("write_file", args)
    if node_result is not None:
        return node_result
    path = args.get("path", "")
    content = args.get("content", "")
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            # Default relative paths to artifacts session folder during chat
            session_id = getattr(_thread_local, 'current_session_id', None)
            agent = getattr(_thread_local, 'current_agent', None) or _current_agent
            if session_id and agent:
                folder = _get_artifact_session_folder(session_id)
                artifact_dir = os.path.join(AGENTS_DIR, agent.agent_id, "artifacts", folder)
                os.makedirs(artifact_dir, exist_ok=True)
                path = os.path.join(artifact_dir, path)
            else:
                path = os.path.abspath(path)
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        with open(path, "w") as f:
            f.write(content)
        size = os.path.getsize(path)
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        _after_file_write(path, "created", agent.agent_id if agent else "main")
        return _ok({"path": path, "size": size, "status": "written"})
    except Exception as e:
        return _err(f"write_file: {e}")


def tool_edit_file(args: dict) -> str:
    path = args.get("path", "")
    old_string = args.get("old_string", "")
    new_string = args.get("new_string", "")
    replace_all = args.get("replace_all", False)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        with open(path, "r") as f:
            content = f.read()
        count = content.count(old_string)
        if count == 0:
            return _err(f"edit_file: old_string not found in {path}")
        if count > 1 and not replace_all:
            return _err(f"edit_file: old_string found {count} times — use replace_all=true or provide a more specific match")
        if replace_all:
            new_content = content.replace(old_string, new_string)
        else:
            new_content = content.replace(old_string, new_string, 1)
        with open(path, "w") as f:
            f.write(new_content)
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        _after_file_write(path, "modified", agent.agent_id if agent else "main")
        return _ok({"path": path, "replacements": count if replace_all else 1, "status": "edited"})
    except Exception as e:
        return _err(f"edit_file: {e}")


def _describe_image_with_vision(image_data_b64: str, media_type: str, filename: str) -> str:
    """Use a vision-capable model to describe an image attachment."""
    vision_model = getattr(_thread_local, 'attachment_image_model', '') or ''
    if not vision_model:
        return f"(Image: {filename} — no image model configured. Set attachments.image_model in config.)"

    # Build multimodal message with image
    content_blocks = [
        {
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": media_type,
                "data": image_data_b64,
            },
        },
        {"type": "text", "text": f"Describe this image ({filename}) in detail. Include any text, data, diagrams, or visual elements you see."},
    ]

    try:
        result = _run_delegate(
            messages=[{"role": "user", "content": content_blocks}],
            model=vision_model,
            system_prompt="You are a precise image description assistant. Describe the image content thoroughly and concisely.",
            inference_params={"max_tokens": 2048, "temperature": 0.1},
            tools=False,
        )
        if result:
            return f"**Image: {filename}**\n\n{result}"
        return f"(Image: {filename} — vision model returned no description)"
    except Exception as e:
        return f"(Image: {filename} — vision error: {e})"


def tool_read_attachment(args: dict) -> str:
    """Read a user-attached file from the session attachment store."""
    name = args.get("name", "")
    if not name:
        return _err("read_attachment: 'name' is required")
    attachments = getattr(_thread_local, 'attachments', None) or {}
    if name not in attachments:
        # Try case-insensitive match
        for k in attachments:
            if k.lower() == name.lower():
                name = k
                break
        else:
            available = list(attachments.keys())
            if available:
                return _err(f"Attachment '{name}' not found. Available: {', '.join(available)}")
            return _err("No attachments in this session.")
    att = attachments[name]
    content = att.get("content", "")
    encoding = att.get("encoding", "text")
    media_type = att.get("media_type", "text/plain")
    ext = os.path.splitext(name)[1].lower()

    # --- Binary files (base64 encoded) ---
    if encoding == "base64":
        import base64 as b64_mod
        import io as io_mod
        try:
            raw_bytes = b64_mod.b64decode(content)
        except Exception as e:
            return _err(f"Failed to decode attachment: {e}")

        # PDF
        if ext == ".pdf" or "pdf" in media_type.lower():
            try:
                import fitz  # pymupdf
                doc = fitz.open(stream=raw_bytes, filetype="pdf")
                pages = []
                for i, page in enumerate(doc):
                    text = page.get_text()
                    if text.strip():
                        pages.append(f"--- Page {i+1} ---\n{text}")
                doc.close()
                return "\n\n".join(pages) if pages else "(PDF has no extractable text)"
            except ImportError:
                return _err("Install pymupdf for PDF support: pip3 install pymupdf")

        # DOCX
        if ext == ".docx":
            try:
                import docx
                doc = docx.Document(io_mod.BytesIO(raw_bytes))
                return DocumentParser.parse_docx.__func__(None) if False else _parse_docx_from_bytes(raw_bytes)
            except ImportError:
                return _err("Install python-docx for DOCX support: pip3 install python-docx")

        # XLSX
        if ext == ".xlsx":
            try:
                return _parse_xlsx_from_bytes(raw_bytes, sheet=args.get("sheet"))
            except ImportError:
                return _err("Install openpyxl for XLSX support: pip3 install openpyxl")

        # PPTX
        if ext == ".pptx":
            try:
                return _parse_pptx_from_bytes(raw_bytes)
            except ImportError:
                return _err("Install python-pptx for PPTX support: pip3 install python-pptx")

        # Images — use vision model
        if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".svg") or media_type.startswith("image/"):
            return _describe_image_with_vision(content, media_type, name)

        # Unknown binary
        return f"(Binary file: {name}, {len(raw_bytes)} bytes, type: {media_type})"

    # --- Text files ---
    return content


def _parse_docx_from_bytes(raw_bytes: bytes) -> str:
    """Parse DOCX from in-memory bytes."""
    import io, docx
    doc = docx.Document(io.BytesIO(raw_bytes))
    paragraphs = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                level = 1
                try:
                    level = int(para.style.name.replace("Heading", "").strip()) or 1
                except ValueError:
                    pass
                text = "#" * level + " " + text
            paragraphs.append(text)
    # Tables
    for table in doc.tables:
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append("| " + " | ".join(cells) + " |")
        if rows:
            # Add header separator after first row
            header_sep = "| " + " | ".join(["---"] * len(table.rows[0].cells)) + " |"
            rows.insert(1, header_sep)
            paragraphs.append("\n".join(rows))
    return "\n\n".join(paragraphs)


def _parse_xlsx_from_bytes(raw_bytes: bytes, sheet: str | None = None) -> str:
    """Parse XLSX from in-memory bytes."""
    import io, openpyxl
    wb = openpyxl.load_workbook(io.BytesIO(raw_bytes), read_only=True, data_only=True)
    sheets = [wb[sheet]] if sheet and sheet in wb.sheetnames else wb.worksheets
    parts = []
    for ws in sheets:
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        # Build markdown table
        header = rows[0]
        cols = [str(c) if c is not None else "" for c in header]
        lines = [f"**Sheet: {ws.title}**"]
        lines.append("| " + " | ".join(cols) + " |")
        lines.append("| " + " | ".join(["---"] * len(cols)) + " |")
        for row in rows[1:200]:  # Cap at 200 rows
            cells = [str(c) if c is not None else "" for c in row]
            lines.append("| " + " | ".join(cells) + " |")
        if len(rows) > 201:
            lines.append(f"*... ({len(rows) - 201} more rows)*")
        parts.append("\n".join(lines))
    wb.close()
    return "\n\n".join(parts) if parts else "(Empty spreadsheet)"


def _parse_pptx_from_bytes(raw_bytes: bytes) -> str:
    """Parse PPTX from in-memory bytes."""
    import io
    from pptx import Presentation
    prs = Presentation(io.BytesIO(raw_bytes))
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        texts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append("| " + " | ".join(cells) + " |")
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    # Notes
    for i, slide in enumerate(prs.slides, 1):
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slides.append(f"[Slide {i} notes] {notes}")
    return "\n\n".join(slides) if slides else "(Empty presentation)"


def tool_read_document(args: dict) -> str:
    """Format-aware document reader."""
    path = args.get("path", "")
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        if not os.path.exists(path):
            return _err(f"File not found: {path}")
        ext = os.path.splitext(path)[1].lower()

        if ext == ".pdf":
            try:
                import fitz
            except ImportError:
                return _err("Install pymupdf: pip3 install pymupdf")
            doc = fitz.open(path)
            meta = {
                "title": doc.metadata.get("title", ""),
                "author": doc.metadata.get("author", ""),
                "page_count": doc.page_count,
            }
            pages_param = args.get("pages", "")
            page_indices = None
            if pages_param:
                page_indices = set()
                for part in pages_param.split(","):
                    part = part.strip()
                    if "-" in part:
                        a, b = part.split("-", 1)
                        for i in range(int(a), int(b) + 1):
                            page_indices.add(i)
                    else:
                        page_indices.add(int(part))
            page_texts = []
            for i, page in enumerate(doc, 1):
                if page_indices and i not in page_indices:
                    continue
                page_texts.append(f"--- Page {i} ---\n{page.get_text()}")
            doc.close()
            content = "\n\n".join(page_texts)
            meta_str = "\n".join(f"**{k}:** {v}" for k, v in meta.items() if v)
            return _ok({"path": path, "format": "pdf", "metadata": meta_str, "content": content})

        elif ext == ".docx":
            try:
                import docx
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            doc = docx.Document(path)
            paragraphs = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    if para.style and para.style.name and para.style.name.startswith("Heading"):
                        level = 1
                        try:
                            level = int(para.style.name.replace("Heading", "").strip()) or 1
                        except ValueError:
                            pass
                        text = "#" * level + " " + text
                    paragraphs.append(text)
            # Extract tables
            for table in doc.tables:
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    max_cols = max(len(r) for r in rows)
                    for r in rows:
                        while len(r) < max_cols:
                            r.append("")
                    header = "| " + " | ".join(rows[0]) + " |"
                    sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
                    table_lines = [header, sep]
                    for r in rows[1:]:
                        table_lines.append("| " + " | ".join(r) + " |")
                    paragraphs.append("\n".join(table_lines))
            content = "\n\n".join(paragraphs)
            return _ok({"path": path, "format": "docx", "content": content})

        elif ext in (".xlsx", ".xls"):
            sheet = args.get("sheet")
            content = DocumentParser.parse_xlsx(path, sheet=sheet)
            return _ok({"path": path, "format": "xlsx", "content": content})

        elif ext == ".pptx":
            slides = args.get("slides")
            content = DocumentParser.parse_pptx(path, slides=slides)
            return _ok({"path": path, "format": "pptx", "content": content})

        elif ext in (".csv", ".tsv"):
            content = DocumentParser.parse_csv(path)
            return _ok({"path": path, "format": "csv", "content": content})

        elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"):
            meta_text = DocumentParser.parse_image(path)
            vision_note = "\n\n*(For AI-powered image description, include this image directly in your chat message)*"
            return _ok({"path": path, "format": "image", "content": meta_text + vision_note})

        elif ext == ".svg":
            content = DocumentParser.parse_svg(path)
            return _ok({"path": path, "format": "svg", "content": content})

        else:
            # Fallback to plain text read (same as read_file)
            with open(path, "r", errors="replace") as f:
                lines = f.readlines()
            total = len(lines)
            limit = 500
            selected = lines[:limit]
            numbered = []
            for i, line in enumerate(selected, start=1):
                numbered.append(f"{i:>6}\t{line.rstrip()}")
            content = "\n".join(numbered)
            return _ok({"path": path, "format": "text", "total_lines": total, "content": content})
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"read_document: {e}")


def tool_write_document(args: dict) -> str:
    """Create documents from markdown content."""
    path = args.get("path", "")
    content = args.get("content", "")
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        ext = os.path.splitext(path)[1].lower()

        if ext == ".docx":
            try:
                import docx
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            doc = docx.Document()
            lines = content.split("\n")
            i = 0
            while i < len(lines):
                line = lines[i]
                # Headings
                heading_match = re.match(r'^(#{1,6})\s+(.*)', line)
                if heading_match:
                    level = len(heading_match.group(1))
                    doc.add_heading(heading_match.group(2), level=level)
                    i += 1
                    continue
                # Table detection
                if "|" in line and i + 1 < len(lines) and re.match(r'^\|[\s\-:|]+\|', lines[i + 1]):
                    table_rows = []
                    while i < len(lines) and "|" in lines[i]:
                        cells = [c.strip() for c in lines[i].strip().strip("|").split("|")]
                        if not re.match(r'^[\s\-:|]+$', lines[i].strip().strip("|")):
                            table_rows.append(cells)
                        i += 1
                    if table_rows:
                        max_cols = max(len(r) for r in table_rows)
                        table = doc.add_table(rows=len(table_rows), cols=max_cols)
                        table.style = "Table Grid"
                        for ri, row_data in enumerate(table_rows):
                            for ci, cell_val in enumerate(row_data):
                                if ci < max_cols:
                                    table.rows[ri].cells[ci].text = cell_val
                    continue
                # Regular paragraph with inline formatting
                stripped = line.strip()
                if stripped:
                    para = doc.add_paragraph()
                    parts = re.split(r'(\*\*\*.*?\*\*\*|\*\*.*?\*\*|\*.*?\*)', stripped)
                    for part in parts:
                        if part.startswith("***") and part.endswith("***"):
                            run = para.add_run(part[3:-3])
                            run.bold = True
                            run.italic = True
                        elif part.startswith("**") and part.endswith("**"):
                            run = para.add_run(part[2:-2])
                            run.bold = True
                        elif part.startswith("*") and part.endswith("*") and len(part) > 2:
                            run = para.add_run(part[1:-1])
                            run.italic = True
                        else:
                            para.add_run(part)
                i += 1
            doc.save(path)

        elif ext == ".xlsx":
            try:
                import openpyxl
            except ImportError:
                return _err("Install openpyxl: pip3 install openpyxl")
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            sections = re.split(r'^##\s+(.+)$', content, flags=re.MULTILINE)
            if len(sections) < 3:
                ws = wb.create_sheet("Sheet1")
                _write_md_table_to_sheet(ws, content)
            else:
                for si in range(1, len(sections), 2):
                    sheet_name = sections[si].strip()
                    if sheet_name.lower().startswith("sheet:"):
                        sheet_name = sheet_name[6:].strip()
                    sheet_content = sections[si + 1] if si + 1 < len(sections) else ""
                    ws = wb.create_sheet(sheet_name[:31])
                    _write_md_table_to_sheet(ws, sheet_content)
            if not wb.sheetnames:
                wb.create_sheet("Sheet1")
            wb.save(path)

        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return _err("Install python-pptx: pip3 install python-pptx")
            prs = Presentation()
            slides_content = re.split(r'^#\s+(.+)$', content, flags=re.MULTILINE)
            if len(slides_content) < 3:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = "Slide 1"
                slide.placeholders[1].text = content.strip()
            else:
                for si in range(1, len(slides_content), 2):
                    title = slides_content[si].strip()
                    body = slides_content[si + 1].strip() if si + 1 < len(slides_content) else ""
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = title
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    body_lines = [l for l in body.split("\n") if l.strip()]
                    for li, bline in enumerate(body_lines):
                        bline = bline.strip()
                        bline = re.sub(r'^[-*]\s+', '', bline)
                        if li == 0:
                            tf.text = bline
                        else:
                            p = tf.add_paragraph()
                            p.text = bline
            prs.save(path)

        elif ext == ".pdf":
            try:
                from reportlab.lib.pagesizes import letter
                from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
                from reportlab.lib.styles import getSampleStyleSheet
            except ImportError:
                return _err("Install reportlab: pip3 install reportlab")
            doc_pdf = SimpleDocTemplate(path, pagesize=letter)
            styles = getSampleStyleSheet()
            story = []
            for line in content.split("\n"):
                line = line.strip()
                if not line:
                    story.append(Spacer(1, 12))
                    continue
                heading_match = re.match(r'^(#{1,6})\s+(.*)', line)
                if heading_match:
                    level = len(heading_match.group(1))
                    style_name = f"Heading{min(level, 6)}"
                    if style_name not in styles:
                        style_name = "Heading1"
                    story.append(Paragraph(heading_match.group(2), styles[style_name]))
                else:
                    line = re.sub(r'\*\*\*(.+?)\*\*\*', r'<b><i>\1</i></b>', line)
                    line = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', line)
                    line = re.sub(r'\*(.+?)\*', r'<i>\1</i>', line)
                    story.append(Paragraph(line, styles["Normal"]))
            doc_pdf.build(story)

        else:
            return _err(f"write_document: unsupported format '{ext}'. Supported: .docx, .xlsx, .pptx, .pdf")

        size = os.path.getsize(path)
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        _after_file_write(path, "created", agent.agent_id if agent else "main")
        return _ok({"path": path, "size": size, "format": ext.lstrip("."), "status": "written"})
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"write_document: {e}")


def _write_md_table_to_sheet(ws, md_text: str) -> None:
    """Helper: parse markdown table text and write rows to an openpyxl worksheet."""
    row_idx = 1
    for line in md_text.split("\n"):
        line = line.strip()
        if not line or not line.startswith("|"):
            continue
        if re.match(r'^\|[\s\-:|]+\|$', line):
            continue
        cells = [c.strip() for c in line.strip("|").split("|")]
        for ci, val in enumerate(cells, 1):
            try:
                ws.cell(row=row_idx, column=ci, value=float(val) if "." in val else int(val))
            except (ValueError, TypeError):
                ws.cell(row=row_idx, column=ci, value=val)
        row_idx += 1


def tool_edit_document(args: dict) -> str:
    """Targeted edits to existing documents."""
    path = args.get("path", "")
    action = args.get("action", "")
    params = args.get("params", {})
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)
        if not os.path.exists(path):
            return _err(f"File not found: {path}")
        ext = os.path.splitext(path)[1].lower()

        if ext == ".docx":
            if action != "replace_text":
                return _err(f"edit_document: unsupported action '{action}' for DOCX. Use: replace_text")
            try:
                import docx
            except ImportError:
                return _err("Install python-docx: pip3 install python-docx")
            old_text = params.get("old_text", "")
            new_text = params.get("new_text", "")
            if not old_text:
                return _err("edit_document: 'old_text' required in params")
            doc = docx.Document(path)
            count = 0
            for para in doc.paragraphs:
                if old_text in para.text:
                    full_text = para.text
                    new_full = full_text.replace(old_text, new_text)
                    for run in para.runs:
                        run.text = ""
                    if para.runs:
                        para.runs[0].text = new_full
                    else:
                        para.add_run(new_full)
                    count += 1
            doc.save(path)
            agent = getattr(_thread_local, 'current_agent', None) or _current_agent
            _after_file_write(path, "modified", agent.agent_id if agent else "main")
            return _ok({"path": path, "action": action, "replacements": count, "status": "edited"})

        elif ext in (".xlsx", ".xls"):
            try:
                import openpyxl
            except ImportError:
                return _err("Install openpyxl: pip3 install openpyxl")
            wb = openpyxl.load_workbook(path)

            if action == "update_cell":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                cell_ref = params.get("cell", "")
                value = params.get("value", "")
                if not cell_ref:
                    return _err("edit_document: 'cell' required (e.g. 'A1')")
                if sheet_name not in wb.sheetnames:
                    return _err(f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}")
                ws = wb[sheet_name]
                try:
                    ws[cell_ref] = float(value) if isinstance(value, str) and value.replace(".", "", 1).replace("-", "", 1).isdigit() else value
                except Exception:
                    ws[cell_ref] = value
                wb.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "cell": cell_ref, "sheet": sheet_name, "status": "edited"})

            elif action == "add_row":
                sheet_name = params.get("sheet", wb.sheetnames[0])
                values = params.get("values", [])
                if not values:
                    return _err("edit_document: 'values' array required")
                if sheet_name not in wb.sheetnames:
                    return _err(f"Sheet '{sheet_name}' not found. Available: {wb.sheetnames}")
                ws = wb[sheet_name]
                ws.append(values)
                wb.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "sheet": sheet_name, "row_added": len(values), "status": "edited"})

            else:
                return _err(f"edit_document: unsupported action '{action}' for XLSX. Use: update_cell, add_row")

        elif ext == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return _err("Install python-pptx: pip3 install python-pptx")
            prs = Presentation(path)

            if action == "update_slide":
                slide_index = int(params.get("slide_index", 1))
                if slide_index < 1 or slide_index > len(prs.slides):
                    return _err(f"Slide index {slide_index} out of range (1-{len(prs.slides)})")
                slide = prs.slides[slide_index - 1]
                title = params.get("title")
                body = params.get("body")
                if title and slide.shapes.title:
                    slide.shapes.title.text = title
                if body:
                    for shape in slide.shapes:
                        if shape.has_text_frame and shape != slide.shapes.title:
                            shape.text_frame.text = body
                            break
                prs.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "slide_index": slide_index, "status": "edited"})

            elif action == "add_slide":
                title = params.get("title", "New Slide")
                body = params.get("body", "")
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = title
                if body:
                    slide.placeholders[1].text = body
                prs.save(path)
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                _after_file_write(path, "modified", agent.agent_id if agent else "main")
                return _ok({"path": path, "action": action, "slide_count": len(prs.slides), "status": "edited"})

            else:
                return _err(f"edit_document: unsupported action '{action}' for PPTX. Use: update_slide, add_slide")

        else:
            return _err(f"edit_document: unsupported format '{ext}'. Supported: .docx, .xlsx, .xls, .pptx")
    except ImportError as e:
        return _err(str(e))
    except Exception as e:
        return _err(f"edit_document: {e}")


def tool_list_directory(args: dict) -> str:
    node_result = _route_to_node("list_directory", args)
    if node_result is not None:
        return node_result
    path = args.get("path", ".")
    pattern = args.get("pattern")
    recursive = args.get("recursive", False)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)

        if pattern:
            if recursive or "**" in pattern:
                full_pattern = os.path.join(path, pattern)
                entries = globmod.glob(full_pattern, recursive=True)
            else:
                full_pattern = os.path.join(path, pattern)
                entries = globmod.glob(full_pattern)
        elif recursive:
            entries = []
            for root, dirs, files in os.walk(path):
                # Skip hidden dirs
                dirs[:] = [d for d in dirs if not d.startswith(".")]
                for f in files:
                    if not f.startswith("."):
                        entries.append(os.path.join(root, f))
        else:
            entries = [os.path.join(path, e) for e in os.listdir(path)]

        results = []
        for entry in sorted(entries)[:500]:
            try:
                st = os.stat(entry)
                is_dir = os.path.isdir(entry)
                results.append({
                    "name": os.path.relpath(entry, path),
                    "type": "directory" if is_dir else "file",
                    "size": st.st_size if not is_dir else None,
                })
            except OSError:
                results.append({"name": os.path.relpath(entry, path), "type": "unknown"})

        return _ok({"path": path, "count": len(results), "entries": results})
    except Exception as e:
        return _err(f"list_directory: {e}")


def tool_search_files(args: dict) -> str:
    pattern = args.get("pattern", "")
    path = args.get("path", ".")
    file_glob = args.get("glob")
    case_insensitive = args.get("case_insensitive", False)
    max_results = args.get("max_results", 50)
    try:
        path = os.path.expanduser(path)
        if not os.path.isabs(path):
            path = os.path.abspath(path)

        flags = re.IGNORECASE if case_insensitive else 0
        regex = re.compile(pattern, flags)

        matches = []
        files_searched = 0

        if os.path.isfile(path):
            file_list = [path]
        else:
            file_list = []
            for root, dirs, files in os.walk(path):
                dirs[:] = [d for d in dirs if not d.startswith(".") and d not in ("node_modules", "__pycache__", ".git")]
                for f in files:
                    if f.startswith("."):
                        continue
                    fp = os.path.join(root, f)
                    if file_glob and not fnmatch.fnmatch(f, file_glob):
                        continue
                    file_list.append(fp)

        for fp in file_list:
            if len(matches) >= max_results:
                break
            files_searched += 1
            try:
                with open(fp, "r", errors="replace") as fh:
                    for lineno, line in enumerate(fh, 1):
                        if regex.search(line):
                            matches.append({
                                "file": os.path.relpath(fp, path) if not os.path.isfile(path) else fp,
                                "line": lineno,
                                "text": line.rstrip()[:500],
                            })
                            if len(matches) >= max_results:
                                break
            except (OSError, UnicodeDecodeError):
                continue

        return _ok({"pattern": pattern, "path": path, "files_searched": files_searched,
                     "match_count": len(matches), "matches": matches})
    except re.error as e:
        return _err(f"search_files: invalid regex: {e}")
    except Exception as e:
        return _err(f"search_files: {e}")


def _strip_ansi(text: str) -> str:
    """Remove all ANSI escape sequences, control chars, and normalize whitespace."""
    text = re.sub(r"\033\[[0-9;]*[a-zA-Z]", "", text)
    text = re.sub(r"\033\[\?[0-9;]*[a-zA-Z]", "", text)
    text = re.sub(r"\033\([A-Z]", "", text)
    text = re.sub(r"\033][^\a]*\a", "", text)  # OSC sequences
    text = re.sub(r"\r", "", text)  # Carriage returns
    text = text.replace("\t", "    ")  # Expand tabs
    return text


def _build_shell_command(command: str) -> tuple[list | str, bool]:
    """Build the shell invocation based on execute_command config.

    Returns (cmd, shell_flag) for subprocess.Popen.
    If login_shell is True, wraps the command in a login shell invocation
    so that ~/.zprofile, ~/.zshrc etc. are sourced (giving full PATH).
    """
    _exec_cfg = get_tool_config().get("execute_command", {})
    use_login_shell = _exec_cfg.get("login_shell", True)
    if use_login_shell:
        shell_path = _exec_cfg.get("shell_path", "") or os.environ.get("SHELL", "/bin/zsh")
        return [shell_path, "-l", "-c", command], False
    return command, True


def _streaming_execute_command(command: str, timeout: int, cwd: str | None,
                               event_callback, tool_use_id: str) -> str:
    """Execute command with streaming output via event_callback."""
    env = os.environ.copy()
    env["TERM"] = "dumb"
    env["NO_COLOR"] = "1"
    env["PAGER"] = "cat"
    env["COLUMNS"] = "200"
    env["LINES"] = "50"

    shell_cmd, shell_flag = _build_shell_command(command)
    proc = subprocess.Popen(
        shell_cmd, shell=shell_flag, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        stdin=subprocess.DEVNULL, cwd=cwd, env=env,
        start_new_session=True,
    )
    output_lines = []
    import io
    deadline = time.time() + timeout
    try:
        # Read stdout line by line, emitting events
        for raw_line in iter(proc.stdout.readline, b''):
            if time.time() > deadline:
                import signal as sig
                try:
                    os.killpg(proc.pid, sig.SIGKILL)
                except OSError:
                    proc.kill()
                proc.wait(timeout=5)
                line_text = _strip_ansi(raw_line.decode("utf-8", errors="replace"))
                output_lines.append(line_text)
                output = "".join(output_lines)
                if len(output) > 50000:
                    output = output[:50000] + "\n... (truncated)"
                return _err(f"execute_command: timed out after {timeout}s\n{output}")
            line_text = _strip_ansi(raw_line.decode("utf-8", errors="replace"))
            output_lines.append(line_text)
            if event_callback:
                event_callback("tool_output", {
                    "tool_use_id": tool_use_id,
                    "line": line_text.rstrip("\n"),
                })
        proc.stdout.close()
        proc.wait(timeout=5)
    except subprocess.TimeoutExpired:
        import signal as sig
        try:
            os.killpg(proc.pid, sig.SIGKILL)
        except OSError:
            proc.kill()
        proc.wait(timeout=5)

    output = "".join(output_lines)
    stderr_data = proc.stderr.read() if proc.stderr else b""
    proc.stderr.close()
    if stderr_data:
        err_text = _strip_ansi(stderr_data.decode("utf-8", errors="replace"))
        output += ("\n--- stderr ---\n" + err_text) if output else err_text
    if len(output) > 50000:
        output = output[:50000] + "\n... (truncated)"
    return _ok({"command": command, "exit_code": proc.returncode, "output": output})


def tool_execute_command(args: dict) -> str:
    node_result = _route_to_node("execute_command", args)
    if node_result is not None:
        return node_result
    command = args.get("command", "")
    cwd = args.get("cwd")
    # Read default timeout from tools_config (per-call timeout still overrides)
    _exec_cfg = get_tool_config().get("execute_command", {})
    default_timeout = _exec_cfg.get("timeout", 15)
    timeout = args.get("timeout", default_timeout)
    # Check banned commands
    banned = _exec_cfg.get("banned_commands", [])
    for b in banned:
        if b and b in command:
            return _err(f"execute_command: command contains banned pattern '{b}'")
    try:
        if cwd:
            cwd = os.path.expanduser(cwd)

        # Use streaming version if event_callback is available
        ecb = getattr(_thread_local, 'event_callback', None)
        tuid = getattr(_thread_local, 'tool_use_id', None)
        if ecb and tuid:
            return _streaming_execute_command(command, timeout, cwd, ecb, tuid)

        # Force non-interactive environment
        env = os.environ.copy()
        env["TERM"] = "dumb"
        env["NO_COLOR"] = "1"
        env["PAGER"] = "cat"
        env["COLUMNS"] = "200"
        env["LINES"] = "50"

        shell_cmd, shell_flag = _build_shell_command(command)
        proc = subprocess.Popen(
            shell_cmd, shell=shell_flag, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
            stdin=subprocess.DEVNULL, cwd=cwd, env=env,
            start_new_session=True,  # own process group so we can kill the tree
        )
        try:
            stdout, stderr = proc.communicate(timeout=timeout)
        except subprocess.TimeoutExpired:
            # Kill the entire process group
            import signal as sig
            try:
                os.killpg(proc.pid, sig.SIGKILL)
            except OSError:
                proc.kill()
            stdout, stderr = proc.communicate(timeout=5)
            output = _strip_ansi(stdout.decode("utf-8", errors="replace"))
            if stderr:
                output += "\n--- stderr ---\n" + _strip_ansi(stderr.decode("utf-8", errors="replace"))
            if len(output) > 50000:
                output = output[:50000] + "\n... (truncated)"
            return _err(f"execute_command: timed out after {timeout}s (partial output below). Use non-interactive commands, e.g. 'top -l 1' not 'top'.\n{output}")

        output = _strip_ansi(stdout.decode("utf-8", errors="replace"))
        if stderr:
            err_text = _strip_ansi(stderr.decode("utf-8", errors="replace"))
            output += ("\n--- stderr ---\n" + err_text) if output else err_text
        if len(output) > 50000:
            output = output[:50000] + "\n... (truncated)"
        return _ok({"command": command, "exit_code": proc.returncode, "output": output})
    except Exception as e:
        return _err(f"execute_command: {e}")


# --- Gmail Tools ---

import imaplib
import smtplib
import email
import email.mime.text
import email.mime.multipart
from email.header import decode_header as _decode_header


def _gmail_config():
    """Load Gmail credentials from tools_config, falling back to gmail.json."""
    # Check tools_config first
    tcfg = get_tool_config().get("gmail", {})
    if tcfg.get("email") and tcfg.get("app_password"):
        return {"email": tcfg["email"], "app_password": tcfg["app_password"]}
    # Fall back to gmail.json
    config_path = os.path.join(AGENTS_DIR, "main", "gmail.json")
    if not os.path.exists(config_path):
        return None
    try:
        with open(config_path) as f:
            return json.load(f)
    except (json.JSONDecodeError, OSError):
        return None


def _decode_mime_header(raw):
    """Decode a MIME-encoded header value."""
    if not raw:
        return ""
    parts = _decode_header(raw)
    decoded = []
    for data, charset in parts:
        if isinstance(data, bytes):
            decoded.append(data.decode(charset or "utf-8", errors="replace"))
        else:
            decoded.append(str(data))
    return " ".join(decoded)


def _get_email_body(msg):
    """Extract plain text body from email message."""
    if msg.is_multipart():
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/plain" and "attachment" not in str(part.get("Content-Disposition", "")):
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    return payload.decode(charset, errors="replace")
        # Fallback to HTML
        for part in msg.walk():
            ct = part.get_content_type()
            if ct == "text/html":
                payload = part.get_payload(decode=True)
                if payload:
                    charset = part.get_content_charset() or "utf-8"
                    text = payload.decode(charset, errors="replace")
                    # Strip HTML tags roughly
                    text = re.sub(r'<[^>]+>', ' ', text)
                    text = re.sub(r'\s+', ' ', text).strip()
                    return text
    else:
        payload = msg.get_payload(decode=True)
        if payload:
            charset = msg.get_content_charset() or "utf-8"
            return payload.decode(charset, errors="replace")
    return ""


def tool_gmail_inbox(args: dict) -> str:
    """List recent emails from Gmail inbox."""
    cfg = _gmail_config()
    if not cfg:
        return _err("Gmail not configured. Create agents/main/gmail.json with email and app_password.")
    limit = args.get("limit", 10)
    folder = args.get("folder", "INBOX")
    try:
        imap = imaplib.IMAP4_SSL("imap.gmail.com")
        imap.login(cfg["email"], cfg["app_password"])
        imap.select(folder, readonly=True)
        _, data = imap.search(None, "ALL")
        ids = data[0].split()
        ids = ids[-limit:]  # most recent
        ids.reverse()

        emails = []
        for eid in ids:
            _, msg_data = imap.fetch(eid, "(RFC822.HEADER)")
            raw = msg_data[0][1]
            msg = email.message_from_bytes(raw)
            emails.append({
                "id": eid.decode(),
                "from": _decode_mime_header(msg.get("From", "")),
                "subject": _decode_mime_header(msg.get("Subject", "")),
                "date": msg.get("Date", ""),
            })
        imap.logout()
        return _ok({"folder": folder, "count": len(emails), "emails": emails})
    except Exception as e:
        return _err(f"gmail_inbox: {e}")


def tool_gmail_read(args: dict) -> str:
    """Read a specific email by ID."""
    cfg = _gmail_config()
    if not cfg:
        return _err("Gmail not configured.")
    email_id = args.get("id", "")
    folder = args.get("folder", "INBOX")
    if not email_id:
        return _err("gmail_read: email id is required")
    try:
        imap = imaplib.IMAP4_SSL("imap.gmail.com")
        imap.login(cfg["email"], cfg["app_password"])
        imap.select(folder, readonly=True)
        _, msg_data = imap.fetch(email_id.encode(), "(RFC822)")
        raw = msg_data[0][1]
        msg = email.message_from_bytes(raw)
        body = _get_email_body(msg)
        # Truncate long bodies
        if len(body) > 10000:
            body = body[:10000] + "\n...(truncated)"
        attachments = []
        if msg.is_multipart():
            for part in msg.walk():
                fn = part.get_filename()
                if fn:
                    attachments.append(_decode_mime_header(fn))
        imap.logout()
        return _ok({
            "id": email_id,
            "from": _decode_mime_header(msg.get("From", "")),
            "to": _decode_mime_header(msg.get("To", "")),
            "cc": _decode_mime_header(msg.get("Cc", "")),
            "subject": _decode_mime_header(msg.get("Subject", "")),
            "date": msg.get("Date", ""),
            "body": body,
            "attachments": attachments,
            "message_id": msg.get("Message-ID", ""),
        })
    except Exception as e:
        return _err(f"gmail_read: {e}")


def tool_gmail_search(args: dict) -> str:
    """Search emails using Gmail search syntax."""
    cfg = _gmail_config()
    if not cfg:
        return _err("Gmail not configured.")
    query = args.get("query", "")
    limit = args.get("limit", 10)
    if not query:
        return _err("gmail_search: query is required")
    try:
        imap = imaplib.IMAP4_SSL("imap.gmail.com")
        imap.login(cfg["email"], cfg["app_password"])
        imap.select("INBOX", readonly=True)
        # Gmail supports X-GM-RAW for full Gmail search syntax
        _, data = imap.search(None, f'X-GM-RAW "{query}"')
        ids = data[0].split()
        ids = ids[-limit:]
        ids.reverse()

        emails = []
        for eid in ids:
            _, msg_data = imap.fetch(eid, "(RFC822.HEADER)")
            raw = msg_data[0][1]
            msg = email.message_from_bytes(raw)
            emails.append({
                "id": eid.decode(),
                "from": _decode_mime_header(msg.get("From", "")),
                "subject": _decode_mime_header(msg.get("Subject", "")),
                "date": msg.get("Date", ""),
            })
        imap.logout()
        return _ok({"query": query, "count": len(emails), "emails": emails})
    except Exception as e:
        return _err(f"gmail_search: {e}")


def tool_gmail_send(args: dict) -> str:
    """Send an email via Gmail SMTP."""
    cfg = _gmail_config()
    if not cfg:
        return _err("Gmail not configured.")
    to = args.get("to", "")
    subject = args.get("subject", "")
    body = args.get("body", "")
    cc = args.get("cc", "")
    if not to or not subject:
        return _err("gmail_send: to and subject are required")
    try:
        msg = email.mime.multipart.MIMEMultipart()
        msg["From"] = cfg["email"]
        msg["To"] = to
        msg["Subject"] = subject
        if cc:
            msg["Cc"] = cc
        msg.attach(email.mime.text.MIMEText(body, "plain"))

        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
            smtp.login(cfg["email"], cfg["app_password"])
            recipients = [to] + ([cc] if cc else [])
            smtp.send_message(msg, to_addrs=recipients)

        return _ok({"status": "sent", "to": to, "subject": subject})
    except Exception as e:
        return _err(f"gmail_send: {e}")


def tool_gmail_reply(args: dict) -> str:
    """Reply to an email."""
    cfg = _gmail_config()
    if not cfg:
        return _err("Gmail not configured.")
    email_id = args.get("id", "")
    body = args.get("body", "")
    if not email_id or not body:
        return _err("gmail_reply: id and body are required")
    try:
        # Fetch original email to get headers
        imap = imaplib.IMAP4_SSL("imap.gmail.com")
        imap.login(cfg["email"], cfg["app_password"])
        imap.select("INBOX", readonly=True)
        _, msg_data = imap.fetch(email_id.encode(), "(RFC822)")
        raw = msg_data[0][1]
        original = email.message_from_bytes(raw)
        imap.logout()

        # Build reply
        reply_to = original.get("Reply-To") or original.get("From", "")
        subject = original.get("Subject", "")
        if not subject.lower().startswith("re:"):
            subject = f"Re: {subject}"

        msg = email.mime.multipart.MIMEMultipart()
        msg["From"] = cfg["email"]
        msg["To"] = reply_to
        msg["Subject"] = subject
        msg["In-Reply-To"] = original.get("Message-ID", "")
        msg["References"] = original.get("Message-ID", "")
        msg.attach(email.mime.text.MIMEText(body, "plain"))

        with smtplib.SMTP_SSL("smtp.gmail.com", 465) as smtp:
            smtp.login(cfg["email"], cfg["app_password"])
            smtp.send_message(msg)

        return _ok({"status": "replied", "to": reply_to, "subject": subject})
    except Exception as e:
        return _err(f"gmail_reply: {e}")


def tool_web_fetch(args: dict) -> str:
    url = args.get("url", "")
    method = args.get("method", "GET")
    headers = args.get("headers", {})
    body = args.get("body")
    max_length = args.get("max_length", 50000)
    force_fresh = args.get("force_fresh", False)
    # Read timeout and max_size from tools_config
    _wf_cfg = get_tool_config().get("web_fetch", {})
    _wf_timeout = _wf_cfg.get("timeout", 30)
    _wf_max_size_mb = _wf_cfg.get("max_size_mb", 10)

    # Check cache for GET requests without body
    cache_key = url if method == "GET" and not body else None
    if cache_key and not force_fresh:
        cached = _web_cache.get(cache_key)
        if cached is not None:
            cached["cached"] = True
            return _ok(cached)

    try:
        req_headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        }
        req_headers.update(headers)
        data = body.encode("utf-8") if body else None
        req = urllib.request.Request(url, data=data, headers=req_headers, method=method)
        with urllib.request.urlopen(req, timeout=_wf_timeout) as resp:
            raw = resp.read(_wf_max_size_mb * 1024 * 1024)
            encoding = resp.headers.get("Content-Encoding", "")
            if encoding == "gzip":
                import gzip
                raw = gzip.decompress(raw)
            charset = resp.headers.get_content_charset() or "utf-8"
            text = raw.decode(charset, errors="replace")
        if len(text) > max_length:
            text = text[:max_length] + "\n... (truncated)"
        result = {"url": url, "status": resp.status, "length": len(text), "content": text}
        if cache_key:
            _web_cache.put(cache_key, dict(result))
        return _ok(result)
    except urllib.error.HTTPError as e:
        body_text = ""
        try:
            body_text = e.read().decode("utf-8", errors="replace")[:5000]
        except Exception:
            pass
        return _err(f"web_fetch: HTTP {e.code} {e.reason}\n{body_text}")
    except Exception as e:
        return _err(f"web_fetch: {e}")


def exa_search(query: str, num_results: int = 5, category: str | None = None,
               force_fresh: bool = False) -> str:
    """Execute an Exa web search and return JSON results. Uses stdlib only."""
    # Check cache
    cache_key = f"exa:{query}:{num_results}:{category or ''}"
    if not force_fresh:
        cached = _web_cache.get(cache_key)
        if cached is not None:
            cached["cached"] = True
            return json.dumps(cached, indent=1)

    # Read API key from tools_config, fall back to env var, then hardcoded default
    _tcfg = get_tool_config().get("exa_search", {})
    api_key = _tcfg.get("api_key") or os.environ.get("EXA_API_KEY", "97dbd594-f7b4-4866-9a8e-6a297e3df576")

    body = {
        "query": query,
        "type": "auto",
        "num_results": num_results,
        "contents": {"highlights": {"max_characters": 4000}},
    }
    if category:
        body["category"] = category

    headers = {
        "x-api-key": api_key,
        "Content-Type": "application/json",
        "Accept": "application/json",
        "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/131.0.0.0 Safari/537.36",
        "Accept-Language": "en-US,en;q=0.9",
    }

    try:
        req = urllib.request.Request(
            "https://api.exa.ai/search",
            data=json.dumps(body).encode("utf-8"),
            headers=headers,
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=30) as resp:
            raw = resp.read()
            # Handle gzip encoding if server sends it anyway
            encoding = resp.headers.get("Content-Encoding", "")
            if encoding == "gzip":
                import gzip
                raw = gzip.decompress(raw)
            response_data = json.loads(raw.decode("utf-8"))

        results = []
        for r in response_data.get("results", []):
            highlights = r.get("highlights", [])
            snippet = " ".join(highlights) if highlights else ""
            results.append({
                "title": r.get("title", ""),
                "link": r.get("url", ""),
                "snippet": snippet,
            })

        search_info = {"query": query, "results": results, "result_count": len(results)}
        if category:
            search_info["category"] = category
        if not results:
            search_info["message"] = "No search results found. Try a different query."
        if results:
            _web_cache.put(cache_key, dict(search_info))
        return json.dumps(search_info, indent=1)

    except urllib.error.HTTPError as e:
        error_body = ""
        try:
            error_body = e.read().decode("utf-8")
        except Exception:
            pass
        return json.dumps({"query": query, "results": [], "error": f"HTTP {e.code}: {error_body}"})
    except Exception as e:
        return json.dumps({"query": query, "results": [], "error": str(e)})


# --- MCP Client ---

class MCPStdioClient:
    """MCP client over stdio — launches a subprocess and communicates via JSON-RPC."""

    def __init__(self, name: str, command: str, args: list[str] | None = None,
                 env: dict | None = None):
        self.name = name
        self.command = command
        self.args = args or []
        self.env = env
        self.process = None
        self._request_id = 0
        self._lock = threading.Lock()
        self.tools: list[dict] = []

    def start(self) -> bool:
        """Start the MCP server subprocess."""
        try:
            run_env = os.environ.copy()
            if self.env:
                run_env.update(self.env)
            # Use login shell to resolve commands installed via nvm/brew/etc.
            # Same approach as _build_shell_command() for execute_command.
            _exec_cfg = get_tool_config().get("execute_command", {})
            use_login_shell = _exec_cfg.get("login_shell", True)
            if use_login_shell:
                shell_path = _exec_cfg.get("shell_path", "") or os.environ.get("SHELL", "/bin/zsh")
                full_cmd = " ".join([self.command] + self.args)
                cmd = [shell_path, "-l", "-c", full_cmd]
            else:
                cmd = [self.command] + self.args
            self.process = subprocess.Popen(
                cmd,
                stdin=subprocess.PIPE, stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                env=run_env, start_new_session=True,
            )
            # Initialize
            resp = self._send_request("initialize", {
                "protocolVersion": "2024-11-05",
                "capabilities": {},
                "clientInfo": {"name": "brain-agent", "version": VERSION},
            })
            if resp and not resp.get("error"):
                # Send initialized notification
                self._send_notification("notifications/initialized", {})
                # List tools
                tools_resp = self._send_request("tools/list", {})
                if tools_resp and tools_resp.get("result"):
                    self.tools = tools_resp["result"].get("tools", [])
                return True
            # Capture stderr for diagnostics
            self._last_error = ""
            if self.process and self.process.stderr:
                try:
                    self._last_error = self.process.stderr.read(2000).decode("utf-8", errors="replace")
                except Exception:
                    pass
            if self.process:
                self.process.terminate()
                self.process = None
            return False
        except Exception as e:
            self._last_error = str(e)
            if self.process:
                try:
                    self.process.terminate()
                except Exception:
                    pass
                self.process = None
            return False

    def stop(self):
        """Stop the subprocess."""
        if self.process:
            try:
                self.process.terminate()
                self.process.wait(timeout=5)
            except Exception:
                try:
                    self.process.kill()
                except Exception:
                    pass
            self.process = None

    def call_tool(self, tool_name: str, arguments: dict) -> str:
        """Call a tool on the MCP server. Returns JSON string."""
        resp = self._send_request("tools/call", {
            "name": tool_name,
            "arguments": arguments,
        })
        if resp and resp.get("result"):
            content = resp["result"].get("content", [])
            # Extract text from content blocks
            texts = []
            for block in content:
                if isinstance(block, dict):
                    if block.get("type") == "text":
                        texts.append(block.get("text", ""))
                    else:
                        texts.append(json.dumps(block))
                elif isinstance(block, str):
                    texts.append(block)
            return json.dumps({"result": "\n".join(texts)})
        elif resp and resp.get("error"):
            return json.dumps({"error": resp["error"].get("message", str(resp["error"]))})
        return json.dumps({"error": "No response from MCP server"})

    def _send_request(self, method: str, params: dict) -> dict | None:
        with self._lock:
            self._request_id += 1
            msg = {
                "jsonrpc": "2.0",
                "id": self._request_id,
                "method": method,
                "params": params,
            }
            return self._send_and_receive(msg)

    def _send_notification(self, method: str, params: dict):
        with self._lock:
            msg = {
                "jsonrpc": "2.0",
                "method": method,
                "params": params,
            }
            self._write(msg)

    def _send_and_receive(self, msg: dict) -> dict | None:
        try:
            self._write(msg)
            return self._read()
        except Exception:
            return None

    def _write(self, msg: dict):
        if not self.process or not self.process.stdin:
            return
        data = json.dumps(msg)
        self.process.stdin.write(f"{data}\n".encode("utf-8"))
        self.process.stdin.flush()

    def _read(self, timeout: int = 30) -> dict | None:
        if not self.process or not self.process.stdout:
            return None
        import select as _select
        # Read lines until we get a JSON-RPC response (skip notifications)
        deadline = time.time() + timeout
        while time.time() < deadline:
            remaining = max(0.1, deadline - time.time())
            ready, _, _ = _select.select([self.process.stdout], [], [], remaining)
            if not ready:
                return None  # Timeout
            line = self.process.stdout.readline()
            if not line:
                return None
            line = line.decode("utf-8").strip()
            if not line:
                continue
            try:
                msg = json.loads(line)
                if "id" in msg:  # It's a response
                    return msg
                # Skip notifications
            except json.JSONDecodeError:
                continue
        return None  # Timeout


class MCPSSEClient:
    """MCP client over SSE/HTTP — connects to a running server."""

    def __init__(self, name: str, url: str, headers: dict | None = None):
        self.name = name
        self.url = url.rstrip("/")
        self.headers = headers or {}
        self._request_id = 0
        self.tools: list[dict] = []

    def start(self) -> bool:
        """Initialize connection and list tools."""
        try:
            resp = self._post("initialize", {
                "protocolVersion": "2024-11-05",
                "capabilities": {},
                "clientInfo": {"name": "brain-agent", "version": VERSION},
            })
            if resp and not resp.get("error"):
                self._post("notifications/initialized", {}, is_notification=True)
                tools_resp = self._post("tools/list", {})
                if tools_resp and tools_resp.get("result"):
                    self.tools = tools_resp["result"].get("tools", [])
                return True
            return False
        except Exception:
            return False

    def stop(self):
        pass  # No cleanup needed for HTTP

    def call_tool(self, tool_name: str, arguments: dict) -> str:
        """Call a tool on the MCP server."""
        resp = self._post("tools/call", {
            "name": tool_name,
            "arguments": arguments,
        })
        if resp and resp.get("result"):
            content = resp["result"].get("content", [])
            texts = []
            for block in content:
                if isinstance(block, dict):
                    if block.get("type") == "text":
                        texts.append(block.get("text", ""))
                    else:
                        texts.append(json.dumps(block))
                elif isinstance(block, str):
                    texts.append(block)
            return json.dumps({"result": "\n".join(texts)})
        elif resp and resp.get("error"):
            return json.dumps({"error": resp["error"].get("message", str(resp["error"]))})
        return json.dumps({"error": "No response from MCP server"})

    def _post(self, method: str, params: dict, is_notification: bool = False) -> dict | None:
        self._request_id += 1
        msg = {
            "jsonrpc": "2.0",
            "method": method,
            "params": params,
        }
        if not is_notification:
            msg["id"] = self._request_id

        headers = {
            "Content-Type": "application/json",
            "Accept": "application/json",
        }
        headers.update(self.headers)

        try:
            data = json.dumps(msg).encode("utf-8")
            req = urllib.request.Request(
                f"{self.url}/message" if "/message" not in self.url else self.url,
                data=data, headers=headers, method="POST",
            )
            with urllib.request.urlopen(req, timeout=30) as resp:
                body = resp.read().decode("utf-8")
                if body.strip():
                    return json.loads(body)
                return {} if is_notification else None
        except Exception:
            return None


class MCPManager:
    """Manages MCP server connections for an agent. Thread-safe."""

    def __init__(self):
        self._lock = threading.Lock()
        self.clients: dict[str, MCPStdioClient | MCPSSEClient] = {}
        self._tool_to_server: dict[str, str] = {}  # tool_name -> server_name

    def load_config(self, config_path: str) -> int:
        """Load MCP servers from a mcp.json config file. Returns count of servers started."""
        if not os.path.exists(config_path):
            return 0
        try:
            with open(config_path, "r") as f:
                config = json.load(f)
        except (json.JSONDecodeError, OSError):
            return 0

        count = 0
        for name, cfg in config.items():
            transport = cfg.get("transport", "stdio")
            if transport == "stdio":
                client = MCPStdioClient(
                    name=name,
                    command=cfg.get("command", ""),
                    args=cfg.get("args", []),
                    env=cfg.get("env"),
                )
            elif transport in ("sse", "http"):
                client = MCPSSEClient(
                    name=name,
                    url=cfg.get("url", ""),
                    headers=cfg.get("headers"),
                )
            else:
                continue

            if client.start():
                with self._lock:
                    self.clients[name] = client
                    for tool in client.tools:
                        tool_name = f"mcp_{name}_{tool['name']}"
                        self._tool_to_server[tool_name] = name
                count += 1
        return count

    def get_tool_definitions(self) -> list[dict]:
        """Get all MCP tool definitions in Anthropic format."""
        defs = []
        with self._lock:
            clients_snapshot = list(self.clients.items())
        for server_name, client in clients_snapshot:
            for tool in client.tools:
                prefixed_name = f"mcp_{server_name}_{tool['name']}"
                defs.append({
                    "name": prefixed_name,
                    "description": f"[MCP:{server_name}] {tool.get('description', '')}",
                    "input_schema": tool.get("inputSchema", {
                        "type": "object", "properties": {}, "required": [],
                    }),
                })
        return defs

    def get_tool_definitions_openai(self) -> list[dict]:
        """Get all MCP tool definitions in OpenAI format."""
        defs = []
        for td in self.get_tool_definitions():
            defs.append({
                "type": "function",
                "function": {
                    "name": td["name"],
                    "description": td["description"],
                    "parameters": {
                        "type": td["input_schema"].get("type", "object"),
                        "properties": td["input_schema"].get("properties", {}),
                        "required": td["input_schema"].get("required", []),
                    },
                },
            })
        return defs

    def call_tool(self, prefixed_name: str, arguments: dict) -> str:
        """Call an MCP tool by its prefixed name."""
        with self._lock:
            server_name = self._tool_to_server.get(prefixed_name)
            client = self.clients.get(server_name) if server_name else None
        if not server_name or not client:
            return json.dumps({"error": f"MCP tool '{prefixed_name}' not found"})
        prefix = f"mcp_{server_name}_"
        original_name = prefixed_name[len(prefix):]
        return client.call_tool(original_name, arguments)

    def is_mcp_tool(self, name: str) -> bool:
        with self._lock:
            return name in self._tool_to_server

    def list_servers(self) -> list[dict]:
        """List all connected MCP servers and their tools."""
        result = []
        with self._lock:
            clients_snapshot = list(self.clients.items())
        for name, client in clients_snapshot:
            result.append({
                "name": name,
                "transport": "stdio" if isinstance(client, MCPStdioClient) else "sse",
                "tools": [t["name"] for t in client.tools],
                "tool_count": len(client.tools),
            })
        return result

    def connect_runtime(self, url: str, name: str, transport: str = "sse") -> dict:
        """Connect to an MCP server at runtime. Returns status dict with discovered tools."""
        with self._lock:
            if name in self.clients:
                return {"error": f"Server '{name}' is already connected"}
        if transport == "stdio":
            parts = url.split()
            client = MCPStdioClient(name=name, command=parts[0], args=parts[1:] if len(parts) > 1 else [])
        else:
            client = MCPSSEClient(name=name, url=url)

        if client.start():
            with self._lock:
                self.clients[name] = client
                for tool in client.tools:
                    tool_name = f"mcp_{name}_{tool['name']}"
                    self._tool_to_server[tool_name] = name
            return {
                "status": "connected",
                "name": name,
                "transport": transport,
                "url": url,
                "tools": [{"name": t["name"], "description": t.get("description", "")} for t in client.tools],
                "tool_count": len(client.tools),
            }
        detail = getattr(client, '_last_error', '') or ''
        msg = f"Failed to connect to MCP server '{name}' at {url}"
        if detail:
            msg += f"\nDetail: {detail[:500]}"
        return {"error": msg}

    def disconnect_runtime(self, name: str) -> dict:
        """Disconnect a runtime MCP server."""
        with self._lock:
            if name not in self.clients:
                return {"error": f"Server '{name}' is not connected"}
            client = self.clients.pop(name)
            to_remove = [k for k, v in self._tool_to_server.items() if v == name]
            for k in to_remove:
                del self._tool_to_server[k]
        client.stop()
        return {"status": "disconnected", "name": name}

    def stop_all(self):
        """Stop all MCP server connections."""
        with self._lock:
            clients_to_stop = list(self.clients.values())
            self.clients.clear()
            self._tool_to_server.clear()
        for client in clients_to_stop:
            client.stop()


# Global MCP manager
_mcp_manager: MCPManager | None = None


# --- Agent System ---

AGENTS_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "agents")

# --- Tool Configuration ---

_TOOLS_CONFIG_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "tools_config.json")

_TOOLS_CONFIG_DEFAULTS = {
    "exa_search": {
        "enabled": True,
        "api_key": "",
        "default_num_results": 5,
    },
    "gmail": {
        "enabled": True,
        "email": "",
        "app_password": "",
    },
    "execute_command": {
        "enabled": True,
        "timeout": 120,
        "banned_commands": ["rm -rf /", "mkfs", "dd if="],
        "login_shell": True,       # Use login shell (sources ~/.zprofile, ~/.zshrc) for full PATH
        "shell_path": "",           # Shell binary path, empty = auto-detect from $SHELL (default: /bin/zsh)
    },
    "web_fetch": {
        "enabled": True,
        "timeout": 30,
        "max_size_mb": 10,
    },
    "refinement": {
        "enabled": True,
        "model": "",  # empty = auto-select (Haiku > Sonnet > cheapest)
    },
    "read_document": {
        "enabled": True,
        "max_file_size_mb": 50,
        "vision_model": "",  # for image description; empty = auto
    },
    "write_document": {
        "enabled": True,
    },
    "edit_document": {
        "enabled": True,
    },
    "code_graph": {
        "enabled": True,
        "exclude_dirs": "node_modules,.git,__pycache__,venv,.venv,dist,build",
        "max_file_size_kb": 500,
    },
}


def get_tool_config() -> dict:
    """Load tool configuration from tools_config.json, falling back to defaults."""
    cfg = {}
    if os.path.exists(_TOOLS_CONFIG_PATH):
        try:
            with open(_TOOLS_CONFIG_PATH) as f:
                cfg = json.load(f)
        except (json.JSONDecodeError, OSError):
            cfg = {}
    # Merge with defaults (defaults provide missing keys)
    merged = {}
    for tool_name, defaults in _TOOLS_CONFIG_DEFAULTS.items():
        tool_cfg = cfg.get(tool_name, {})
        merged[tool_name] = {**defaults, **tool_cfg}
    return merged


def save_tool_config(cfg: dict) -> dict:
    """Save tool configuration to tools_config.json. Returns saved config."""
    # Merge with existing to preserve fields not in the incoming payload
    existing = get_tool_config()
    for tool_name, tool_cfg in cfg.items():
        if tool_name in existing:
            existing[tool_name].update(tool_cfg)
        else:
            existing[tool_name] = tool_cfg
    try:
        with open(_TOOLS_CONFIG_PATH, "w") as f:
            json.dump(existing, f, indent=2)
    except OSError as e:
        return {"error": str(e)}
    return existing


def get_tool_status() -> dict:
    """Return status of each configurable tool, checking all fallback sources."""
    cfg = get_tool_config()
    status = {}
    for tool_name, tool_cfg in cfg.items():
        tool_cfg = dict(tool_cfg)  # copy to avoid mutating defaults
        enabled = tool_cfg.get("enabled", True)
        if not enabled:
            s = "disabled"
        elif tool_name == "exa_search":
            exa_key = tool_cfg.get("api_key") or os.environ.get("EXA_API_KEY", "")
            # Check hardcoded fallback in tool function
            if not exa_key:
                exa_key = "97dbd594-f7b4-4866-9a8e-6a297e3df576"  # built-in default
                tool_cfg["_source"] = "built-in default"
            elif not tool_cfg.get("api_key") and os.environ.get("EXA_API_KEY"):
                tool_cfg["_source"] = "environment variable"
            s = "configured" if exa_key else "not configured"
        elif tool_name == "gmail":
            has_gmail = bool(tool_cfg.get("email") and tool_cfg.get("app_password"))
            if not has_gmail:
                gmail_fb = _gmail_config()
                if gmail_fb and gmail_fb.get("email") and gmail_fb.get("app_password"):
                    has_gmail = True
                    tool_cfg["email"] = gmail_fb["email"]
                    tool_cfg["_source"] = "gmail.json"
            s = "configured" if has_gmail else "not configured"
        elif tool_name == "refinement":
            s = "configured" if tool_cfg.get("model") else "auto (Haiku > cheapest)"
        else:
            s = "configured"
        status[tool_name] = {"enabled": enabled, "status": s, "config": tool_cfg}
    return status


class AgentConfig:
    """Configuration and file management for a single agent."""

    def __init__(self, agent_id: str):
        self.agent_id = agent_id
        self.dir = os.path.join(AGENTS_DIR, agent_id)
        os.makedirs(self.dir, exist_ok=True)
        self._ensure_defaults()

    def _ensure_defaults(self):
        """Create default files if they don't exist."""
        soul_path = os.path.join(self.dir, "soul.md")
        if not os.path.exists(soul_path):
            with open(soul_path, "w") as f:
                f.write(f"""# {self.agent_id}

You are the **{self.agent_id}** agent.
Adapt your behavior to the tasks you are given.
""")
        config_path = os.path.join(self.dir, "agent.json")
        if not os.path.exists(config_path):
            with open(config_path, "w") as f:
                json.dump({
                    "description": f"{self.agent_id} agent",
                    "model": None,
                }, f, indent=2)

    @property
    def soul(self) -> str:
        """Load soul.md content."""
        path = os.path.join(self.dir, "soul.md")
        try:
            with open(path, "r") as f:
                return f.read()
        except OSError:
            return ""

    @property
    def tools_guide(self) -> str:
        """Load per-agent tools.md, falling back to global tools.md."""
        agent_tools = os.path.join(self.dir, "tools.md")
        if os.path.exists(agent_tools):
            try:
                with open(agent_tools, "r") as f:
                    return f.read()
            except OSError:
                pass
        # Fall back to global tools.md
        global_tools = os.path.join(os.path.dirname(os.path.abspath(__file__)), "tools.md")
        try:
            with open(global_tools, "r") as f:
                return f.read()
        except OSError:
            return ""

    @property
    def config(self) -> dict:
        """Load agent.json config."""
        path = os.path.join(self.dir, "agent.json")
        try:
            with open(path, "r") as f:
                return json.load(f)
        except (OSError, json.JSONDecodeError):
            return {}

    @property
    def description(self) -> str:
        return self.config.get("description", self.agent_id)

    @property
    def preferred_model(self) -> str | None:
        raw = self.config.get("model")
        if not raw:
            return None
        purpose = self.config.get("model_purpose")
        resolved = resolve_model(raw, purpose) if _models_config else raw
        return resolved or raw

    @property
    def max_context(self) -> int | None:
        """Return model's max_context if a preferred model is set."""
        model = self.preferred_model
        if model:
            return get_model_max_context(model)
        return None

    @property
    def memory_dir(self) -> str:
        return self.dir  # memory.db lives alongside soul.md

    @property
    def skills_dir(self) -> str:
        return os.path.join(self.dir, "skills")

    @property
    def mcp_config_path(self) -> str:
        return os.path.join(self.dir, "mcp.json")

    def load_commands(self) -> list[dict]:
        """Load custom slash commands from commands.json + .claude/commands/*.md.

        Supports both Brain Agent format (JSON) and Claude Code format (markdown
        with YAML frontmatter, $ARGUMENTS, !`command` interpolation).
        """
        commands = []

        # 1. Brain Agent format: commands.json
        path = os.path.join(self.dir, "commands.json")
        try:
            with open(path, "r") as f:
                data = json.load(f)
                if isinstance(data, list):
                    for cmd in data:
                        cmd["_format"] = "brain"
                    commands.extend(data)
        except (OSError, json.JSONDecodeError):
            pass

        # 2. Claude Code format: .claude/commands/*.md and agent commands/ dir
        for cmd_dir in [
            os.path.join(os.path.dirname(os.path.abspath(__file__)), ".claude", "commands"),
            os.path.join(self.dir, "commands"),
        ]:
            if not os.path.isdir(cmd_dir):
                continue
            for fname in sorted(os.listdir(cmd_dir)):
                if not fname.endswith(".md"):
                    continue
                fpath = os.path.join(cmd_dir, fname)
                if os.path.islink(fpath):
                    fpath = os.path.realpath(fpath)
                if not os.path.isfile(fpath):
                    continue
                try:
                    with open(fpath, "r") as f:
                        raw = f.read()
                    # Parse YAML frontmatter
                    fm = {}
                    body = raw
                    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n(.*)$', raw, re.DOTALL)
                    if fm_match:
                        for line in fm_match.group(1).split("\n"):
                            if ":" in line:
                                k, v = line.split(":", 1)
                                fm[k.strip()] = v.strip().strip('"').strip("'")
                        body = fm_match.group(2).strip()
                    cmd_name = fname[:-3]  # strip .md
                    # Skip if already defined in commands.json (brain format takes priority)
                    if any(c.get("name") == cmd_name for c in commands):
                        continue
                    commands.append({
                        "name": cmd_name,
                        "description": fm.get("description", ""),
                        "prompt": body,
                        "allowed_tools": fm.get("allowed-tools", ""),
                        "_format": "claude-code",
                        "_path": fpath,
                    })
                except OSError:
                    continue

        return commands

    def save_commands(self, commands: list[dict]):
        """Save custom slash commands to commands.json (Brain Agent format only)."""
        path = os.path.join(self.dir, "commands.json")
        # Only save brain-format commands
        brain_cmds = [c for c in commands if c.get("_format") != "claude-code"]
        # Strip internal fields
        clean = [{k: v for k, v in c.items() if not k.startswith("_")} for c in brain_cmds]
        with open(path, "w") as f:
            json.dump(clean, f, indent=2)

    @staticmethod
    def expand_command(cmd: dict, args: str = "") -> str:
        """Expand a command template with arguments and dynamic content.

        Supports both formats:
        - Brain Agent: {{variable}} substitution
        - Claude Code: $ARGUMENTS substitution + !`command` interpolation
        """
        template = cmd.get("prompt", cmd.get("template", ""))
        fmt = cmd.get("_format", "brain")

        if fmt == "claude-code":
            # Replace $ARGUMENTS with user args
            result = template.replace("$ARGUMENTS", args)

            # Interpolate !`command` — runs shell command and injects output
            import subprocess
            def _run_interpolation(match):
                shell_cmd = match.group(1)
                try:
                    proc = subprocess.run(
                        shell_cmd, shell=True, capture_output=True, text=True,
                        timeout=10, cwd=os.getcwd(),
                        env={**os.environ, "TERM": "dumb"},
                    )
                    return proc.stdout.strip()
                except Exception as e:
                    return f"(error: {e})"

            result = re.sub(r'!`([^`]+)`', _run_interpolation, result)
            return result

        else:
            # Brain Agent: {{variable}} substitution
            if "{{" in template and args:
                var_match = re.search(r'\{\{(\w+)\}\}', template)
                if var_match:
                    return template.replace("{{" + var_match.group(1) + "}}", args)
            return template + (" " + args if args else "")

    def list_skills(self) -> list[dict]:
        """List all skills for this agent (own + main's global skills)."""
        skills = {}
        # Load main's skills first (global)
        if self.agent_id != "main":
            main_skills_dir = os.path.join(AGENTS_DIR, "main", "skills")
            skills.update(self._scan_skills(main_skills_dir, source="main"))
        # Load own skills (override globals if same name)
        skills.update(self._scan_skills(self.skills_dir, source=self.agent_id))
        return list(skills.values())

    def _scan_skills(self, skills_dir: str, source: str) -> dict[str, dict]:
        """Scan a skills directory and return {name: skill_info}."""
        result = {}
        if not os.path.isdir(skills_dir):
            return result
        for name in sorted(os.listdir(skills_dir)):
            skill_dir = os.path.join(skills_dir, name)
            skill_file = os.path.join(skill_dir, "SKILL.md")
            if not os.path.isfile(skill_file):
                continue
            try:
                with open(skill_file, "r") as f:
                    raw = f.read()
                # Parse YAML frontmatter
                fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n(.*)$', raw, re.DOTALL)
                if fm_match:
                    fm_text, body = fm_match.groups()
                    fm = {}
                    for line in fm_text.split("\n"):
                        if ":" in line:
                            k, v = line.split(":", 1)
                            fm[k.strip()] = v.strip().strip('"').strip("'")
                    result[name] = {
                        "name": fm.get("name", name),
                        "slug": name,
                        "description": fm.get("description", ""),
                        "source": source,
                        "path": skill_file,
                    }
                else:
                    result[name] = {
                        "name": name,
                        "slug": name,
                        "description": "",
                        "source": source,
                        "path": skill_file,
                    }
            except OSError:
                continue
        return result

    def load_skill(self, skill_name: str) -> str | None:
        """Load the full SKILL.md body for a specific skill.
        Accepts either the directory name (slug) or the display name."""
        # Try direct match first (slug = directory name)
        own_path = os.path.join(self.skills_dir, skill_name, "SKILL.md")
        if os.path.isfile(own_path):
            return self._read_skill_body(own_path)
        # Fall back to main's skills
        if self.agent_id != "main":
            main_path = os.path.join(AGENTS_DIR, "main", "skills", skill_name, "SKILL.md")
            if os.path.isfile(main_path):
                return self._read_skill_body(main_path)
        # Try matching by display name → slug lookup
        for s in self.list_skills():
            if s.get("name", "").lower() == skill_name.lower() or s.get("slug", "").lower() == skill_name.lower():
                slug = s.get("slug", "")
                if slug:
                    own_path = os.path.join(self.skills_dir, slug, "SKILL.md")
                    if os.path.isfile(own_path):
                        return self._read_skill_body(own_path)
                    if self.agent_id != "main":
                        main_path = os.path.join(AGENTS_DIR, "main", "skills", slug, "SKILL.md")
                        if os.path.isfile(main_path):
                            return self._read_skill_body(main_path)
        return None

    @staticmethod
    def _read_skill_body(path: str) -> str:
        """Read a SKILL.md and return just the body (after frontmatter)."""
        with open(path, "r") as f:
            raw = f.read()
        fm_match = re.match(r'^---\s*\n.*?\n---\s*\n(.*)$', raw, re.DOTALL)
        if fm_match:
            return fm_match.group(1).strip()
        return raw.strip()


def scan_claude_code_skills() -> list[dict]:
    """Discover Claude Code skills/plugins from ~/.claude.

    Scans three sources:
    1. Plugin skills — cached plugins with SKILL.md files
    2. Plugin commands — marketplace plugin commands (slash commands)
    3. User commands — ~/.claude/commands/ markdown files

    Returns list of dicts with: name, slug, description, source, type, path, plugin, marketplace, enabled
    """
    home = os.path.expanduser("~")
    claude_dir = os.path.join(home, ".claude")
    results = []

    # Read installed_plugins.json to get install paths
    installed = {}
    try:
        with open(os.path.join(claude_dir, "plugins", "installed_plugins.json")) as f:
            data = json.load(f)
            installed = data.get("plugins", {})
    except (OSError, json.JSONDecodeError):
        pass

    # Read settings.json for enabled state
    enabled_plugins = {}
    try:
        with open(os.path.join(claude_dir, "settings.json")) as f:
            settings = json.load(f)
            enabled_plugins = settings.get("enabledPlugins", {})
    except (OSError, json.JSONDecodeError):
        pass

    # 1. Plugin skills — scan installed plugin cache dirs for skills/*/SKILL.md
    for plugin_key, installs in installed.items():
        if not installs:
            continue
        install = installs[0]  # Use first (latest) install
        install_path = install.get("installPath", "")
        if not os.path.isdir(install_path):
            continue

        plugin_name = plugin_key.split("@")[0] if "@" in plugin_key else plugin_key
        marketplace = plugin_key.split("@")[1] if "@" in plugin_key else ""
        is_enabled = enabled_plugins.get(plugin_key, False)

        skills_dir = os.path.join(install_path, "skills")
        if os.path.isdir(skills_dir):
            for skill_name in sorted(os.listdir(skills_dir)):
                skill_file = os.path.join(skills_dir, skill_name, "SKILL.md")
                if not os.path.isfile(skill_file):
                    continue
                # Parse frontmatter for name/description
                try:
                    with open(skill_file) as f:
                        raw = f.read(4096)
                    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', raw, re.DOTALL)
                    fm = {}
                    if fm_match:
                        for line in fm_match.group(1).split("\n"):
                            if ":" in line:
                                k, v = line.split(":", 1)
                                fm[k.strip()] = v.strip().strip('"').strip("'")
                except OSError:
                    fm = {}

                results.append({
                    "name": fm.get("name", skill_name),
                    "slug": f"{plugin_name}:{skill_name}",
                    "description": fm.get("description", ""),
                    "source": "claude-code",
                    "type": "skill",
                    "path": skill_file,
                    "plugin": plugin_name,
                    "marketplace": marketplace,
                    "enabled": is_enabled,
                })

        # 2. Plugin commands — commands/*.md in the marketplace plugin dir
        # Find commands in the marketplace source (not cache)
        for mp_dir in ["claude-plugins-official", "anthropic-agent-skills"]:
            cmd_dir = os.path.join(claude_dir, "plugins", "marketplaces", mp_dir, "plugins", plugin_name, "commands")
            if os.path.isdir(cmd_dir):
                for cmd_file in sorted(os.listdir(cmd_dir)):
                    if not cmd_file.endswith(".md"):
                        continue
                    cmd_path = os.path.join(cmd_dir, cmd_file)
                    cmd_name = cmd_file[:-3]  # strip .md
                    # Parse frontmatter
                    try:
                        with open(cmd_path) as f:
                            raw = f.read(4096)
                        fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', raw, re.DOTALL)
                        fm = {}
                        if fm_match:
                            for line in fm_match.group(1).split("\n"):
                                if ":" in line:
                                    k, v = line.split(":", 1)
                                    fm[k.strip()] = v.strip().strip('"').strip("'")
                    except OSError:
                        fm = {}

                    results.append({
                        "name": fm.get("name", cmd_name),
                        "slug": f"{plugin_name}:{cmd_name}",
                        "description": fm.get("description", ""),
                        "source": "claude-code",
                        "type": "command",
                        "path": cmd_path,
                        "plugin": plugin_name,
                        "marketplace": mp_dir,
                        "enabled": is_enabled,
                    })

    # 3. User commands — ~/.claude/commands/*.md
    user_cmd_dir = os.path.join(claude_dir, "commands")
    if os.path.isdir(user_cmd_dir):
        for entry in sorted(os.listdir(user_cmd_dir)):
            if entry.startswith("."):
                continue
            entry_path = os.path.join(user_cmd_dir, entry)
            if os.path.islink(entry_path):
                # Resolve symlink
                entry_path = os.path.realpath(entry_path)
            if os.path.isfile(entry_path) and entry_path.endswith(".md"):
                cmd_name = entry[:-3]
                try:
                    with open(entry_path) as f:
                        raw = f.read(4096)
                    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', raw, re.DOTALL)
                    fm = {}
                    if fm_match:
                        for line in fm_match.group(1).split("\n"):
                            if ":" in line:
                                k, v = line.split(":", 1)
                                fm[k.strip()] = v.strip().strip('"').strip("'")
                except OSError:
                    fm = {}
                results.append({
                    "name": fm.get("name", cmd_name),
                    "slug": cmd_name,
                    "description": fm.get("description", ""),
                    "source": "claude-code",
                    "type": "user-command",
                    "path": entry_path,
                    "plugin": "",
                    "marketplace": "",
                    "enabled": True,  # User commands are always enabled
                })
            elif os.path.isdir(entry_path):
                # Directory with SKILL.md
                skill_file = os.path.join(entry_path, "SKILL.md")
                if os.path.isfile(skill_file):
                    try:
                        with open(skill_file) as f:
                            raw = f.read(4096)
                        fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', raw, re.DOTALL)
                        fm = {}
                        if fm_match:
                            for line in fm_match.group(1).split("\n"):
                                if ":" in line:
                                    k, v = line.split(":", 1)
                                    fm[k.strip()] = v.strip().strip('"').strip("'")
                    except OSError:
                        fm = {}
                    results.append({
                        "name": fm.get("name", entry),
                        "slug": entry,
                        "description": fm.get("description", ""),
                        "source": "claude-code",
                        "type": "user-skill",
                        "path": skill_file,
                        "plugin": "",
                        "marketplace": "",
                        "enabled": True,
                    })

    # Deduplicate: same skill name across plugins → keep first occurrence
    seen_names = {}
    deduped = []
    for s in results:
        key = (s["name"], s["type"])
        if key in seen_names:
            continue
        seen_names[key] = True
        deduped.append(s)

    return deduped


def browse_claude_code_plugins(query: str = "") -> list[dict]:
    """Browse available Claude Code plugins from local marketplace manifests.

    Returns list of dicts with: name, description, category, marketplace, source, homepage, installed
    """
    home = os.path.expanduser("~")
    claude_dir = os.path.join(home, ".claude")
    results = []

    # Read installed_plugins.json to check install state
    installed_keys = set()
    try:
        with open(os.path.join(claude_dir, "plugins", "installed_plugins.json")) as f:
            data = json.load(f)
            installed_keys = set(data.get("plugins", {}).keys())
    except (OSError, json.JSONDecodeError):
        pass

    # Scan all marketplace manifests
    mp_dir = os.path.join(claude_dir, "plugins", "marketplaces")
    if not os.path.isdir(mp_dir):
        return results

    for mp_name in sorted(os.listdir(mp_dir)):
        manifest_path = os.path.join(mp_dir, mp_name, ".claude-plugin", "marketplace.json")
        if not os.path.isfile(manifest_path):
            continue
        try:
            with open(manifest_path) as f:
                manifest = json.load(f)
        except (OSError, json.JSONDecodeError):
            continue

        for plugin in manifest.get("plugins", []):
            name = plugin.get("name", "")
            description = plugin.get("description", "")
            # Filter by query
            if query:
                q = query.lower()
                if q not in name.lower() and q not in description.lower():
                    continue

            plugin_key = f"{name}@{mp_name}"
            results.append({
                "name": name,
                "description": description,
                "category": plugin.get("category", ""),
                "marketplace": mp_name,
                "homepage": plugin.get("homepage", ""),
                "source": plugin.get("source", {}),
                "installed": plugin_key in installed_keys,
            })

    return results


def install_claude_code_plugin(plugin_name: str, marketplace: str = "claude-plugins-official") -> dict:
    """Install a Claude Code plugin from a marketplace.

    Uses `claude plugins add` CLI if available, otherwise clones from git source.
    Returns dict with status/error.
    """
    import subprocess
    import shutil

    home = os.path.expanduser("~")
    claude_dir = os.path.join(home, ".claude")

    # Find the plugin in the marketplace manifest
    manifest_path = os.path.join(claude_dir, "plugins", "marketplaces", marketplace,
                                  ".claude-plugin", "marketplace.json")
    if not os.path.isfile(manifest_path):
        return {"error": f"Marketplace '{marketplace}' not found"}

    try:
        with open(manifest_path) as f:
            manifest = json.load(f)
    except (OSError, json.JSONDecodeError) as e:
        return {"error": f"Failed to read manifest: {e}"}

    plugin_info = None
    for p in manifest.get("plugins", []):
        if p.get("name") == plugin_name:
            plugin_info = p
            break
    if not plugin_info:
        return {"error": f"Plugin '{plugin_name}' not found in {marketplace}"}

    # Try using claude CLI first
    claude_bin = shutil.which("claude")
    if claude_bin:
        try:
            result = subprocess.run(
                [claude_bin, "plugins", "add", f"{plugin_name}@{marketplace}"],
                capture_output=True, text=True, timeout=60,
                env={**os.environ, "TERM": "dumb"},
            )
            if result.returncode == 0:
                return {"status": "installed", "plugin": plugin_name, "marketplace": marketplace,
                        "method": "claude-cli"}
        except (subprocess.TimeoutExpired, OSError):
            pass  # Fall through to manual install

    # Manual install: clone from git source
    source = plugin_info.get("source", {})
    if isinstance(source, str):
        # Local source (relative path in marketplace)
        src_dir = os.path.join(claude_dir, "plugins", "marketplaces", marketplace, source)
        if os.path.isdir(src_dir):
            # Copy to cache
            cache_dir = os.path.join(claude_dir, "plugins", "cache", marketplace,
                                      plugin_name, "local")
            os.makedirs(cache_dir, exist_ok=True)
            shutil.copytree(src_dir, cache_dir, dirs_exist_ok=True)
            # Register in installed_plugins.json
            _register_cc_plugin(plugin_name, marketplace, cache_dir)
            return {"status": "installed", "plugin": plugin_name, "marketplace": marketplace,
                    "method": "copy", "path": cache_dir}
        return {"error": f"Local source '{source}' not found"}

    elif isinstance(source, dict):
        git_url = source.get("url", "")
        if not git_url:
            return {"error": "No git URL in plugin source"}

        # Clone to cache
        cache_dir = os.path.join(claude_dir, "plugins", "cache", marketplace,
                                  plugin_name, "latest")
        if os.path.isdir(cache_dir):
            shutil.rmtree(cache_dir)
        os.makedirs(os.path.dirname(cache_dir), exist_ok=True)

        try:
            # For git-subdir sources, clone then extract subdir
            subdir = source.get("path", "")
            ref = source.get("ref", "main")

            result = subprocess.run(
                ["git", "clone", "--depth", "1", "--branch", ref,
                 git_url if git_url.endswith(".git") else f"https://github.com/{git_url}.git",
                 cache_dir],
                capture_output=True, text=True, timeout=60,
            )
            if result.returncode != 0:
                return {"error": f"git clone failed: {result.stderr[:200]}"}

            # If subdir specified, move it up
            if subdir:
                subdir_path = os.path.join(cache_dir, subdir)
                if os.path.isdir(subdir_path):
                    import tempfile
                    tmp = tempfile.mkdtemp()
                    shutil.copytree(subdir_path, os.path.join(tmp, plugin_name), dirs_exist_ok=True)
                    shutil.rmtree(cache_dir)
                    shutil.copytree(os.path.join(tmp, plugin_name), cache_dir, dirs_exist_ok=True)
                    shutil.rmtree(tmp)

            _register_cc_plugin(plugin_name, marketplace, cache_dir)
            return {"status": "installed", "plugin": plugin_name, "marketplace": marketplace,
                    "method": "git", "path": cache_dir}

        except (subprocess.TimeoutExpired, OSError) as e:
            return {"error": f"Install failed: {e}"}

    return {"error": "Unknown source format"}


def _register_cc_plugin(plugin_name: str, marketplace: str, install_path: str):
    """Register a plugin in ~/.claude/plugins/installed_plugins.json."""
    import datetime
    home = os.path.expanduser("~")
    ip_path = os.path.join(home, ".claude", "plugins", "installed_plugins.json")

    try:
        with open(ip_path) as f:
            data = json.load(f)
    except (OSError, json.JSONDecodeError):
        data = {"version": 2, "plugins": {}}

    plugin_key = f"{plugin_name}@{marketplace}"
    now = datetime.datetime.utcnow().isoformat() + "Z"
    data["plugins"][plugin_key] = [{
        "scope": "user",
        "installPath": install_path,
        "version": "latest",
        "installedAt": now,
        "lastUpdated": now,
        "isLocal": True,
    }]

    with open(ip_path, "w") as f:
        json.dump(data, f, indent=4)

    # Also enable in settings.json
    settings_path = os.path.join(home, ".claude", "settings.json")
    try:
        with open(settings_path) as f:
            settings = json.load(f)
    except (OSError, json.JSONDecodeError):
        settings = {}
    if "enabledPlugins" not in settings:
        settings["enabledPlugins"] = {}
    settings["enabledPlugins"][plugin_key] = True
    with open(settings_path, "w") as f:
        json.dump(settings, f, indent=4)


def list_agents() -> list[str]:
    """List all available agent IDs."""
    if not os.path.isdir(AGENTS_DIR):
        return ["main"]
    agents = []
    for name in sorted(os.listdir(AGENTS_DIR)):
        if os.path.isdir(os.path.join(AGENTS_DIR, name)) and not name.startswith("."):
            agents.append(name)
    if not agents:
        agents = ["main"]
    return agents


def get_agent_summaries() -> list[dict]:
    """Get agent ID + description + soul summary for all agents, with team metadata."""
    # First pass: collect raw summaries and scan for teams
    raw = []
    # team_id (the agent whose config holds the team) -> team config
    teams_cfg: dict[str, dict] = {}
    for agent_id in list_agents():
        cfg = AgentConfig(agent_id)
        config = cfg.config
        soul = cfg.soul.strip()
        summary = ""
        for line in soul.split("\n"):
            line = line.strip()
            if line and not line.startswith("#") and not line.startswith("---"):
                summary = line
                break
        entry = {
            "id": agent_id,
            "display_name": config.get("display_name", ""),
            "description": cfg.description,
            "soul_summary": summary,
            "model": cfg.preferred_model,
            "avatar": config.get("avatar"),
            "paused": config.get("paused", False),
            "agent_sdk": config.get("agent_sdk", {}).get("enabled", True) if isinstance(config.get("agent_sdk"), dict) else bool(config.get("agent_sdk", True)),
        }
        team_cfg = config.get("team")
        if isinstance(team_cfg, dict) and team_cfg.get("members"):
            teams_cfg[agent_id] = team_cfg
        raw.append(entry)

    # Second pass: compute team metadata per agent
    for entry in raw:
        aid = entry["id"]
        # Which teams is this agent a member of?
        member_of = []
        is_head_of = None
        for cfg_holder, tcfg in teams_cfg.items():
            members = tcfg.get("members", [])
            head_id = tcfg.get("head", members[0] if members else cfg_holder)
            if aid in members:
                team_name = tcfg.get("name", cfg_holder)
                member_of.append({"team_id": cfg_holder, "team_name": team_name})
                if aid == head_id:
                    is_head_of = cfg_holder
        entry["teams"] = member_of
        entry["is_team_head"] = is_head_of is not None
        if is_head_of:
            tcfg = teams_cfg[is_head_of]
            entry["team_config_holder"] = is_head_of
            entry["team_members"] = list(tcfg.get("members", []))
            entry["team_head"] = tcfg.get("head", entry["team_members"][0] if entry["team_members"] else is_head_of)
            entry["team_name"] = tcfg.get("name", "")
            entry["team_description"] = tcfg.get("description", "")
            entry["team_avatar"] = tcfg.get("avatar", "")
    return raw


def get_team_structure() -> dict:
    """Return hierarchical team structure for UI and API consumption."""
    agents = get_agent_summaries()
    agent_map = {a["id"]: a for a in agents}
    teams = {}
    in_team = set()  # agents that appear in at least one team

    # Find all agents that hold team configs
    for agent_id in list_agents():
        cfg = AgentConfig(agent_id).config
        team_cfg = cfg.get("team")
        if not isinstance(team_cfg, dict) or not team_cfg.get("members"):
            continue
        members_ids = team_cfg["members"]
        head_id = team_cfg.get("head", members_ids[0] if members_ids else agent_id)
        members = []
        for mid in members_ids:
            if mid in agent_map:
                members.append(agent_map[mid])
                in_team.add(mid)
        teams[agent_id] = {
            "head": head_id,
            "head_agent": agent_map.get(head_id),
            "members": members,
            "name": team_cfg.get("name") or agent_id,
            "description": team_cfg.get("description", ""),
            "avatar": team_cfg.get("avatar", ""),
            "config_holder": agent_id,
        }

    standalone = [a for a in agents if a["id"] not in in_team and a["id"] != "main"]
    main_agent = agent_map.get("main")
    return {"teams": teams, "standalone": standalone, "main": main_agent}


def _get_delegation_scope(caller_agent_id: str) -> list[str]:
    """Return list of agent IDs the caller is allowed to delegate to."""
    all_agents = list_agents()

    # Collect all team configs
    team_cfgs = {}  # config_holder_id -> team_cfg
    for aid in all_agents:
        cfg = AgentConfig(aid).config
        team_cfg = cfg.get("team")
        if isinstance(team_cfg, dict) and team_cfg.get("members"):
            team_cfgs[aid] = team_cfg

    in_team = set()
    head_ids = set()
    for holder, tcfg in team_cfgs.items():
        head_id = tcfg.get("head", tcfg["members"][0] if tcfg["members"] else holder)
        head_ids.add(head_id)
        for m in tcfg["members"]:
            in_team.add(m)

    if caller_agent_id == "main":
        # Main can delegate to team heads and standalone agents (not regular members)
        return [a for a in all_agents if a != "main" and (a in head_ids or a not in in_team)]

    # Check if caller is a team head
    for holder, tcfg in team_cfgs.items():
        head_id = tcfg.get("head", tcfg["members"][0] if tcfg["members"] else holder)
        if caller_agent_id == head_id:
            # Team head can delegate to its members (excluding self)
            return [m for m in tcfg["members"] if m != caller_agent_id and m in all_agents]

    # Regular member: can delegate to peers in same team + team head
    reachable = set()
    for holder, tcfg in team_cfgs.items():
        if caller_agent_id in tcfg["members"]:
            for m in tcfg["members"]:
                if m != caller_agent_id and m in all_agents:
                    reachable.add(m)
    return list(reachable) if reachable else [a for a in all_agents if a != caller_agent_id]


def _find_team_head(agent_id: str) -> str | None:
    """Find the team head for a given agent. Returns None if not in any team."""
    for aid in list_agents():
        cfg = AgentConfig(aid).config
        team_cfg = cfg.get("team")
        if isinstance(team_cfg, dict) and team_cfg.get("members"):
            if agent_id in team_cfg["members"]:
                return team_cfg.get("head", team_cfg["members"][0] if team_cfg["members"] else aid)
    return None


def _get_agent_team_info(agent_id: str) -> dict | None:
    """Get team info for an agent. Returns dict with name, head, members, is_head, or None."""
    for aid in list_agents():
        cfg = AgentConfig(aid).config
        team_cfg = cfg.get("team")
        if isinstance(team_cfg, dict) and team_cfg.get("members"):
            if agent_id in team_cfg["members"]:
                head = team_cfg.get("head", team_cfg["members"][0])
                return {
                    "name": team_cfg.get("name", aid),
                    "head": head,
                    "members": team_cfg["members"],
                    "is_head": agent_id == head,
                    "config_holder": aid,
                }
    return None


def build_agent_registry(for_agent_id: str | None = None) -> str:
    """Build a text block describing available agents for injection into system prompts.
    Respects team hierarchy: team heads see members, main sees heads + standalone."""
    agents = get_agent_summaries()
    if len(agents) <= 1:
        return ""

    caller = for_agent_id or "main"
    scope = _get_delegation_scope(caller)
    agent_map = {a["id"]: a for a in agents}

    # Check if caller is a team head
    caller_team_info = _get_agent_team_info(caller) if caller != "main" else None
    caller_is_head = caller_team_info and caller_team_info["is_head"]

    lines = ["AGENT REGISTRY — use delegate_task to send tasks to these agents:"]

    if caller == "main":
        # Group by teams + standalone
        struct = get_team_structure()
        if struct["teams"]:
            lines.append("  TEAMS:")
            for tid, team in struct["teams"].items():
                team_name = team.get("name", tid)
                head_id = team["head"]
                head_agent = team.get("head_agent", {})
                detail = head_agent.get("soul_summary") or head_agent.get("description", "")
                model_note = f" (model: {head_agent.get('model')})" if head_agent.get("model") else ""
                member_names = ", ".join(m["id"] for m in team["members"])
                lines.append(f"    - {head_id} (head of '{team_name}'): {detail}{model_note}")
                lines.append(f"      Members: {member_names}")
        if struct["standalone"]:
            lines.append("  STANDALONE AGENTS:")
            for a in struct["standalone"]:
                detail = a.get("soul_summary") or a.get("description", "")
                model_note = f" (model: {a['model']})" if a.get("model") else ""
                lines.append(f"    - {a['id']}: {detail}{model_note}")
    elif caller_is_head:
        lines.append("  YOUR TEAM MEMBERS:")
        for mid in scope:
            a = agent_map.get(mid, {})
            detail = a.get("soul_summary") or a.get("description", mid)
            model_note = f" (model: {a.get('model')})" if a.get("model") else ""
            lines.append(f"    - {mid}: {detail}{model_note}")
    else:
        for aid in scope:
            a = agent_map.get(aid, {})
            detail = a.get("soul_summary") or a.get("description", aid)
            model_note = f" (model: {a.get('model')})" if a.get("model") else ""
            lines.append(f"  - {aid}: {detail}{model_note}")

    lines.append("")
    lines.append(
        "Before performing a task, consider if another agent is better suited. "
        "Delegate when the task clearly matches another agent's specialty. "
        "Do NOT delegate simple tasks you can handle yourself."
    )
    return "\n".join(lines)


# Current active agent (set in _run_interactive)
_current_agent: AgentConfig | None = None


# --- Memory System (QMD hybrid search) ---

import hashlib
import urllib.request
import urllib.error


def _extract_json_from_llm(text: str, expect_array: bool = False):
    """Robustly extract JSON object or array from LLM response text.

    Handles markdown code fences, nested objects, surrounding text.
    Returns parsed dict/list or None on failure.
    """
    if not text:
        return None
    # Strip markdown code fences first
    stripped = re.sub(r'```(?:json)?\s*', '', text)
    stripped = stripped.replace('```', '')

    # Try parsing the entire stripped text first
    try:
        parsed = json.loads(stripped.strip())
        if expect_array and isinstance(parsed, list):
            return parsed
        if not expect_array and isinstance(parsed, dict):
            return parsed
    except (json.JSONDecodeError, ValueError):
        pass

    # Use json.JSONDecoder.raw_decode to find the first valid JSON structure
    decoder = json.JSONDecoder()
    target_char = '[' if expect_array else '{'
    for i, ch in enumerate(stripped):
        if ch == target_char:
            try:
                obj, _ = decoder.raw_decode(stripped[i:])
                if expect_array and isinstance(obj, list):
                    return obj
                if not expect_array and isinstance(obj, dict):
                    return obj
            except (json.JSONDecodeError, ValueError):
                continue
    return None

# QMD HTTP MCP daemon endpoint
_QMD_URL = "http://localhost:8181/mcp"
_QMD_HEADERS = {
    "Content-Type": "application/json",
    "Accept": "application/json, text/event-stream",
}

# Shared MCP session ID (set on first successful init)
_qmd_session_id: str | None = None
_qmd_session_lock = threading.RLock()
# Per-collection debounce timers for embedding after writes
_qmd_embed_timers: dict[str, threading.Timer] = {}
_qmd_embed_lock = threading.Lock()

# Files to skip when indexing (not memory files)
_QMD_IGNORE_FILES = {"soul.md", "tools.md"}


def _qmd_rpc(method: str, params: dict | None = None) -> dict | None:
    """Send a JSON-RPC request to the QMD MCP HTTP daemon. Returns result or None on failure."""
    global _qmd_session_id
    payload = {"jsonrpc": "2.0", "id": 1, "method": method, "params": params or {}}
    data = json.dumps(payload).encode()
    headers = dict(_QMD_HEADERS)
    with _qmd_session_lock:
        if _qmd_session_id:
            headers["Mcp-Session-Id"] = _qmd_session_id
    req = urllib.request.Request(_QMD_URL, data=data, headers=headers)
    try:
        with urllib.request.urlopen(req, timeout=10) as resp:
            # Capture session ID from response (under lock)
            sid = resp.headers.get("Mcp-Session-Id")
            if sid:
                with _qmd_session_lock:
                    _qmd_session_id = sid
            body = json.loads(resp.read().decode())
            if "error" in body:
                logging.warning("QMD RPC error: %s", body["error"])
                return None
            return body.get("result")
    except Exception as e:
        logging.debug("QMD unreachable: %s", e)
        return None


def _qmd_init_session() -> bool:
    """Initialize an MCP session with QMD. Returns True if successful.
    Thread-safe: only one session is created even under concurrent access."""
    with _qmd_session_lock:
        if _qmd_session_id:
            return True  # Another thread already initialized
        result = _qmd_rpc("initialize", {
            "protocolVersion": "2024-11-05",
            "capabilities": {},
            "clientInfo": {"name": "brain-agent", "version": "1.0"},
        })
        return result is not None


def _qmd_ensure_collection(name: str, directory: str):
    """Register a collection with QMD if it doesn't exist (via CLI)."""
    try:
        result = subprocess.run(
            ["qmd", "collection", "show", name],
            capture_output=True, text=True, timeout=5,
        )
        if result.returncode == 0:
            return  # Already exists
        subprocess.run(
            ["qmd", "collection", "add", directory,
             "--name", name, "--pattern", "*.md",
             "--ignore", ",".join(_QMD_IGNORE_FILES)],
            capture_output=True, text=True, timeout=15,
        )
    except Exception as e:
        logging.debug("QMD collection setup failed for %s: %s", name, e)


def _qmd_debounced_embed(collection: str):
    """Schedule a debounced qmd update+embed for a collection (2s delay).
    Each collection gets its own timer so concurrent writes to different
    collections don't cancel each other's embed."""
    def _do_embed():
        try:
            subprocess.run(
                ["qmd", "update"], capture_output=True, text=True, timeout=30,
            )
            subprocess.run(
                ["qmd", "embed", "-c", collection],
                capture_output=True, text=True, timeout=60,
            )
        except Exception as e:
            logging.debug("QMD embed failed for %s: %s", collection, e)
        finally:
            with _qmd_embed_lock:
                _qmd_embed_timers.pop(collection, None)
    with _qmd_embed_lock:
        old = _qmd_embed_timers.get(collection)
        if old:
            old.cancel()
        timer = threading.Timer(2.0, _do_embed)
        timer.daemon = True
        _qmd_embed_timers[collection] = timer
        timer.start()


def _parse_frontmatter(raw: str) -> tuple[dict, str]:
    """Parse YAML frontmatter from a markdown file. Returns (metadata, body)."""
    # Normalize line endings for cross-platform compatibility
    raw = raw.replace("\r\n", "\n")
    fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n(.*)$', raw, re.DOTALL)
    if fm_match:
        fm_text, body = fm_match.groups()
        fm = {}
        for line in fm_text.split("\n"):
            # Only parse top-level keys (skip indented/nested YAML lines)
            if ":" in line and not line.startswith((" ", "\t", "-")):
                k, v = line.split(":", 1)
                fm[k.strip()] = v.strip()
        return fm, body.strip()
    return {}, raw.strip()


def _yaml_escape(value: str) -> str:
    """Escape a string for safe inclusion in YAML frontmatter."""
    if not value:
        return '""'
    # Quote if contains special YAML chars
    if any(c in value for c in (':', '#', '{', '}', '[', ']', ',', '&', '*', '?', '|', '-', '<', '>', '=', '!', '%', '@', '`', '\n')):
        escaped = value.replace('\\', '\\\\').replace('"', '\\"').replace('\n', ' ')
        return f'"{escaped}"'
    return value


class MemoryStore:
    """Per-agent memory store backed by QMD hybrid search and markdown files."""

    _ensured_collections: set[str] = set()
    _ensured_lock = threading.Lock()

    def __init__(self, agent_id: str = "main", base_dir: str | None = None,
                 user_id: str | None = None, team_ids: list[str] | None = None):
        self.agent_id = agent_id
        if base_dir:
            self.dir = base_dir
        else:
            self.dir = os.path.join(AGENTS_DIR, agent_id)
        os.makedirs(self.dir, exist_ok=True)
        self._collection = agent_id

        # Multi-user scoping
        self.user_id = user_id
        self.team_ids = team_ids or []
        self._user_dir = None
        self._team_dirs: list[tuple[str, str]] = []  # (dir_path, team_id)
        if user_id:
            self._user_dir = os.path.join(self.dir, "users", user_id)
            os.makedirs(self._user_dir, exist_ok=True)
        for tid in self.team_ids:
            td = os.path.join(self.dir, "teams", tid)
            os.makedirs(td, exist_ok=True)
            self._team_dirs.append((td, tid))

        # Ensure QMD knows about this collection (once per collection, background)
        with MemoryStore._ensured_lock:
            already_ensured = agent_id in MemoryStore._ensured_collections
            if not already_ensured:
                MemoryStore._ensured_collections.add(agent_id)
        if not already_ensured:
            threading.Thread(
                target=_qmd_ensure_collection,
                args=(self._collection, self.dir),
                daemon=True,
            ).start()

    def _make_id(self, name: str) -> str:
        """Generate a stable ID from name."""
        return hashlib.sha256(name.encode()).hexdigest()[:12]

    def _name_to_filename(self, name: str) -> str:
        """Convert a memory name to a safe filename with hash suffix to avoid collisions."""
        safe = re.sub(r'[^\w\s-]', '', name).strip().lower()
        safe = re.sub(r'[\s]+', '_', safe)
        # Add short hash to prevent collisions between similar names
        h = hashlib.sha256(name.encode()).hexdigest()[:6]
        base = safe[:50]
        return f"{base}_{h}.md" if base else f"{h}.md"

    def _find_file_for_name(self, name: str) -> str | None:
        """Find existing file for a memory name by checking frontmatter."""
        for fname in os.listdir(self.dir):
            if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
                continue
            fpath = os.path.join(self.dir, fname)
            try:
                with open(fpath, "r") as f:
                    raw = f.read(500)  # only need frontmatter
                fm, _ = _parse_frontmatter(raw)
                if fm.get("name") == name:
                    return fpath
            except Exception:
                continue
        return None

    def _resolve_store_dir(self, scope: str = "global") -> str:
        """Resolve target directory based on scope: global, user, team:<id>."""
        if scope == "user" and self._user_dir:
            return self._user_dir
        if scope.startswith("team:"):
            tid = scope[5:]
            for td, t_id in self._team_dirs:
                if t_id == tid:
                    return td
        return self.dir  # global

    def _all_scan_dirs(self) -> list[str]:
        """Get all directories to scan for this user's memories (global + user + team)."""
        dirs = [self.dir]
        if self._user_dir and os.path.isdir(self._user_dir):
            dirs.append(self._user_dir)
        for td, _ in self._team_dirs:
            if os.path.isdir(td):
                dirs.append(td)
        return dirs

    def store(self, name: str, content: str, description: str = "",
              mem_type: str = "general", scope: str = "global") -> dict:
        """Store or update a memory. Writes .md file and triggers QMD reindex.
        scope: 'global', 'user', or 'team:<team_id>'"""
        mem_id = self._make_id(name)
        filename = self._name_to_filename(name)
        target_dir = self._resolve_store_dir(scope)
        file_path = os.path.join(target_dir, filename)

        # Check if memory already exists under a different filename (migration)
        existing = self._find_file_for_name(name)
        if existing and existing != file_path:
            file_path = existing  # update in place
            filename = os.path.basename(existing)

        # Write markdown file with properly escaped frontmatter
        md_content = f"""---
name: {_yaml_escape(name)}
description: {_yaml_escape(description)}
type: {mem_type}
agent: {self.agent_id}
---

{content}
"""
        with open(file_path, "w") as f:
            f.write(md_content)

        # Trigger debounced QMD update+embed
        _qmd_debounced_embed(self._collection)

        # --- Entity extraction auto-linking (Mechanism 2) ---
        try:
            entities = _extract_entities(content)
            if entities:
                # Find other files sharing entities
                matches = _find_entity_matches(self.agent_id, filename, entities)
                for other_fname in matches[:10]:  # limit to avoid excessive linking
                    other_path = os.path.join(self.dir, other_fname)
                    if os.path.exists(other_path):
                        _add_related_to_file(file_path, other_fname, "same_topic")
                        _add_related_to_file(other_path, filename, "same_topic")
                # Update entity index with new file
                _update_entity_index(self.agent_id, filename, entities)
                if matches:
                    _qmd_debounced_embed(self._collection)
        except Exception as e:
            logging.warning(f"Entity linking failed for {filename}: {e}")  # best-effort, never block store

        return {"id": mem_id, "name": name, "file": filename, "status": "stored"}

    def recall(self, query: str, limit: int = 10, mem_type: str | None = None) -> list[dict]:
        """Search memories using QMD hybrid search (BM25 + vector + reranking).
        Falls back to file-scan substring matching if QMD is unreachable."""
        # Try QMD first
        results = self._qmd_query(query, limit, mem_type)
        if results is not None:
            self._stamp_last_recalled_bg(results)
            return results
        # Fallback: scan files
        results = self._fallback_search(query, limit, mem_type)
        self._stamp_last_recalled_bg(results)
        return results

    def _stamp_last_recalled_bg(self, results: list[dict]):
        """Stamp last_recalled date on recalled files in a background thread."""
        paths = [r.get("file_path") for r in results if r.get("file_path")]
        if paths:
            threading.Thread(target=self._stamp_last_recalled, args=(paths,),
                             daemon=True, name="stamp_recalled").start()

    def _stamp_last_recalled(self, file_paths: list[str]):
        """Update last_recalled frontmatter field on recalled memory files. Best-effort."""
        now = datetime.datetime.now().strftime("%Y-%m-%d")
        for fpath in file_paths:
            try:
                with open(fpath, "r") as f:
                    raw = f.read()
                fm_match = re.match(r'^(---\s*\n)(.*?)(\n---\s*\n)(.*)$', raw, re.DOTALL)
                if not fm_match:
                    continue
                opener, fm_text, closer, body = fm_match.groups()
                if "last_recalled:" in fm_text:
                    fm_text = re.sub(r'last_recalled:.*', f'last_recalled: {now}', fm_text)
                else:
                    fm_text = fm_text.rstrip() + f"\nlast_recalled: {now}"
                with open(fpath, "w") as f:
                    f.write(opener + fm_text + closer + body)
            except Exception as e:
                logging.debug(f"Failed to stamp last_recalled on {fpath}: {e}")
                continue

    def _qmd_query(self, query: str, limit: int, mem_type: str | None) -> list[dict] | None:
        """Query QMD via MCP HTTP. Returns list of results or None if unavailable."""
        global _qmd_session_id
        # Ensure session
        if not _qmd_session_id:
            if not _qmd_init_session():
                return None

        # Sanitize query: strip newlines, quotes, markdown — QMD silently returns empty on these
        clean_q = query.replace('\n', ' ').replace('\r', ' ').replace('"', '').replace("'", "")
        clean_q = re.sub(r'[#*`~\[\]{}()]', '', clean_q).strip()
        if not clean_q:
            return []

        searches = [
            {"type": "lex", "query": clean_q},
            {"type": "vec", "query": clean_q},
        ]
        result = _qmd_rpc("tools/call", {
            "name": "query",
            "arguments": {
                "searches": searches,
                "collections": [self._collection],
                "limit": limit * 2 if mem_type else limit,  # over-fetch if filtering
            },
        })
        if not result:
            # Session may have expired, retry once with lock to prevent stampede
            with _qmd_session_lock:
                _qmd_session_id = None
            if not _qmd_init_session():
                return None
            result = _qmd_rpc("tools/call", {
                "name": "query",
                "arguments": {
                    "searches": searches,
                    "collections": [self._collection],
                    "limit": limit * 2 if mem_type else limit,
                },
            })
            if not result:
                return None

        # Parse structured results
        structured = result.get("structuredContent", {})
        qmd_results = structured.get("results", [])
        memories = []
        for r in qmd_results:
            file_rel = r.get("file", "")
            # Strip collection prefix (e.g. "main/foo.md" -> "foo.md")
            if "/" in file_rel:
                fname = file_rel.split("/", 1)[1]
            else:
                fname = file_rel
            fpath = os.path.join(self.dir, fname)
            # Read the actual file for full content + frontmatter
            try:
                with open(fpath, "r") as f:
                    raw = f.read()
                fm, body = _parse_frontmatter(raw)
            except FileNotFoundError:
                fm = {"name": r.get("title", fname), "type": "general"}
                body = r.get("snippet", "")

            mem = {
                "id": self._make_id(fm.get("name", fname)),
                "name": fm.get("name", fname.replace(".md", "")),
                "description": fm.get("description", ""),
                "type": fm.get("type", "general"),
                "content": body,
                "file_path": fpath,
                "score": r.get("score", 0),
            }
            # Post-filter by type
            if mem_type and mem["type"] != mem_type:
                continue
            memories.append(mem)
            if len(memories) >= limit:
                break
        return memories

    def _fallback_search(self, query: str, limit: int, mem_type: str | None) -> list[dict]:
        """Fallback: scan .md files and do substring matching."""
        terms = query.lower().split()
        if not terms:
            return self.list_all(mem_type)[:limit]
        results = []
        # Directories to scan: all user-visible dirs + chats-indexed subdir
        scan_dirs = list(self._all_scan_dirs())
        chats_dir = os.path.join(self.dir, "chats-indexed")
        if os.path.isdir(chats_dir):
            scan_dirs.append(chats_dir)
        for scan_dir in scan_dirs:
            for fname in os.listdir(scan_dir):
                if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
                    continue
                fpath = os.path.join(scan_dir, fname)
                try:
                    # Cap file read to 32KB to prevent OOM on large files
                    with open(fpath, "r") as f:
                        raw = f.read(32768)
                    fm, body = _parse_frontmatter(raw)
                    mtype = fm.get("type", "general")
                    if mem_type and mtype != mem_type:
                        continue
                    searchable = (fm.get("name", "") + " " + fm.get("description", "") + " " + body).lower()
                    hits = sum(1 for t in terms if t in searchable)
                    if hits > 0:
                        results.append({
                            "id": self._make_id(fm.get("name", fname)),
                            "name": fm.get("name", fname.replace(".md", "")),
                            "description": fm.get("description", ""),
                            "type": mtype,
                            "content": body,
                            "file_path": fpath,
                            "score": hits / len(terms),
                        })
                except (UnicodeDecodeError, OSError) as e:
                    logging.debug(f"Fallback search skipping {fname}: {e}")
                    continue
                except Exception:
                    continue
        results.sort(key=lambda x: x["score"], reverse=True)
        return results[:limit]

    def delete(self, name: str) -> dict:
        """Delete a memory by name."""
        filename = self._name_to_filename(name)
        file_path = os.path.join(self.dir, filename)
        if not os.path.exists(file_path):
            # Try scanning for a file with matching frontmatter name
            found = self._find_file_for_name(name)
            if found:
                file_path = found
            else:
                return {"error": f"Memory '{name}' not found"}
        os.remove(file_path)
        # Trigger QMD reindex
        _qmd_debounced_embed(self._collection)
        return {"name": name, "status": "deleted"}

    def list_all(self, mem_type: str | None = None) -> list[dict]:
        """List all memories by scanning .md files (no QMD needed).
        Scans global + user + team directories."""
        results = []
        for scan_dir in self._all_scan_dirs():
            # Determine scope label
            scope = "global"
            if scan_dir == self._user_dir:
                scope = "user"
            else:
                for td, tid in self._team_dirs:
                    if scan_dir == td:
                        scope = f"team:{tid}"
                        break
            try:
                entries = os.listdir(scan_dir)
            except OSError:
                continue
            for fname in entries:
                if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
                    continue
                fpath = os.path.join(scan_dir, fname)
                try:
                    with open(fpath, "r") as f:
                        raw = f.read()
                    fm, body = _parse_frontmatter(raw)
                    mtype = fm.get("type", "general")
                    if mem_type and mtype != mem_type:
                        continue
                    mtime = os.path.getmtime(fpath)
                    results.append({
                        "id": self._make_id(fm.get("name", fname)),
                        "name": fm.get("name", fname.replace(".md", "")),
                        "description": fm.get("description", ""),
                        "type": mtype,
                        "content": body,
                        "updated_at": datetime.datetime.fromtimestamp(mtime).isoformat(),
                        "scope": scope,
                    })
                except Exception:
                    continue
        results.sort(key=lambda x: x.get("updated_at", ""), reverse=True)
        return results

    def reindex(self) -> dict:
        """Trigger QMD update+embed for this collection."""
        try:
            r1 = subprocess.run(
                ["qmd", "update"], capture_output=True, text=True, timeout=30,
            )
            r2 = subprocess.run(
                ["qmd", "embed", "-c", self._collection],
                capture_output=True, text=True, timeout=60,
            )
            return {"agent": self.agent_id, "status": "reindexed",
                    "update": r1.returncode == 0, "embed": r2.returncode == 0}
        except Exception as e:
            return {"agent": self.agent_id, "status": "error", "error": str(e)}


# Global memory store instance (set in _run_interactive)
_memory_store: MemoryStore | None = None


# ─── Projects System ──────────────────────────────────────────────────

class ProjectManager:
    """CRUD operations for per-agent projects."""

    @staticmethod
    def _project_dir(agent_id: str, name: str) -> str:
        return os.path.join(AGENTS_DIR, agent_id, "projects", name)

    @staticmethod
    def _projects_base(agent_id: str) -> str:
        return os.path.join(AGENTS_DIR, agent_id, "projects")

    @staticmethod
    def list_projects(agent_id: str, user_id: str | None = None,
                      user_team_ids: list[str] | None = None) -> list[dict]:
        """List projects for an agent, optionally filtered by user access.
        user_id=None means admin (sees all). user_team_ids are team IDs the user belongs to."""
        base = ProjectManager._projects_base(agent_id)
        if not os.path.isdir(base):
            return []
        projects = []
        for name in sorted(os.listdir(base)):
            pdir = os.path.join(base, name)
            if not os.path.isdir(pdir) or name.startswith("."):
                continue
            cfg_path = os.path.join(pdir, "project.json")
            cfg = {}
            if os.path.exists(cfg_path):
                try:
                    with open(cfg_path, "r") as f:
                        cfg = json.load(f)
                except (OSError, json.JSONDecodeError):
                    pass
            # Count ingested docs
            ingested_dir = os.path.join(pdir, "ingested")
            chunk_count = 0
            if os.path.isdir(ingested_dir):
                chunk_count = sum(1 for fn in os.listdir(ingested_dir) if fn.startswith("ingest-") and fn.endswith(".md"))
            # Count memory files
            mem_count = sum(1 for fn in os.listdir(pdir) if fn.endswith(".md") and not fn.startswith("ingest-") and fn not in _QMD_IGNORE_FILES)
            # Count ingested source files (unique sources, not chunks)
            doc_count = 0
            if os.path.isdir(ingested_dir):
                seen_sources = set()
                for fn in os.listdir(ingested_dir):
                    if fn.startswith("ingest-") and fn.endswith(".md"):
                        # Source hash is between first and second dash
                        parts_fn = fn.split("-", 2)
                        if len(parts_fn) >= 2:
                            seen_sources.add(parts_fn[1])
                doc_count = len(seen_sources) if seen_sources else chunk_count
            visibility = cfg.get("visibility", "global")
            owner_uid = cfg.get("owner_user_id", "")
            owner_tid = cfg.get("owner_team_id", "")
            projects.append({
                "name": name,
                "description": cfg.get("description", ""),
                "instructions": cfg.get("instructions", ""),
                "icon": cfg.get("icon", "folder"),
                "created_at": cfg.get("created_at", ""),
                "tags": cfg.get("tags", []),
                "watch_folders": cfg.get("watch_folders", []),
                "status": cfg.get("status", "active"),
                "chunks": chunk_count,
                "doc_count": doc_count,
                "memories": mem_count,
                "visibility": visibility,
                "owner_user_id": owner_uid,
                "owner_team_id": owner_tid,
            })
        # Filter by user access if user_id provided
        if user_id is not None:
            team_set = set(user_team_ids or [])
            projects = [
                p for p in projects
                if p["visibility"] == "global"
                or (p["visibility"] == "user" and p["owner_user_id"] == user_id)
                or (p["visibility"] == "team" and p["owner_team_id"] in team_set)
            ]
        return projects

    @staticmethod
    def create_project(agent_id: str, name: str, description: str = "",
                       config: dict | None = None,
                       visibility: str = "global",
                       owner_user_id: str = "",
                       owner_team_id: str = "") -> dict:
        """Create a new project directory with project.json.
        visibility: 'global', 'user', or 'team'"""
        # Validate name
        safe_name = re.sub(r'[^\w\s-]', '', name).strip().lower().replace(' ', '-')
        if not safe_name:
            return {"error": "Invalid project name"}
        pdir = ProjectManager._project_dir(agent_id, safe_name)
        if os.path.exists(pdir):
            return {"error": f"Project '{safe_name}' already exists"}
        os.makedirs(pdir, exist_ok=True)
        os.makedirs(os.path.join(pdir, "ingested"), exist_ok=True)
        os.makedirs(os.path.join(pdir, "notes"), exist_ok=True)
        cfg = {
            "name": config.get("name", name) if config else name,
            "description": description,
            "icon": (config or {}).get("icon", "📁"),
            "created_at": datetime.datetime.now(datetime.timezone.utc).isoformat(),
            "watch_folders": (config or {}).get("watch_folders", []),
            "tags": (config or {}).get("tags", []),
            "model": (config or {}).get("model"),
            "visibility": visibility,
            "owner_user_id": owner_user_id,
            "owner_team_id": owner_team_id,
        }
        cfg_path = os.path.join(pdir, "project.json")
        with open(cfg_path, "w") as f:
            json.dump(cfg, f, indent=2)
        # Register QMD collection for this project
        collection_name = f"{agent_id}/{safe_name}"
        threading.Thread(
            target=_qmd_ensure_collection,
            args=(collection_name, pdir),
            daemon=True,
        ).start()
        return {"name": safe_name, "status": "created", "path": pdir}

    @staticmethod
    def get_project(agent_id: str, name: str) -> dict | None:
        """Read project.json for a project."""
        pdir = ProjectManager._project_dir(agent_id, name)
        if not os.path.isdir(pdir):
            return None
        cfg_path = os.path.join(pdir, "project.json")
        if not os.path.exists(cfg_path):
            # Auto-create minimal project.json for dirs that exist without one
            cfg = {"name": name, "description": "", "created_at": datetime.datetime.now(datetime.timezone.utc).isoformat()}
            try:
                with open(cfg_path, "w") as f:
                    json.dump(cfg, f, indent=2)
            except OSError:
                pass
        try:
            with open(cfg_path, "r") as f:
                cfg = json.load(f)
            # Add computed stats
            ingested_dir = os.path.join(pdir, "ingested")
            chunk_count = 0
            if os.path.isdir(ingested_dir):
                chunk_count = sum(1 for fn in os.listdir(ingested_dir) if fn.startswith("ingest-") and fn.endswith(".md"))
            cfg["chunks"] = chunk_count
            cfg["dir"] = pdir
            return cfg
        except (OSError, json.JSONDecodeError):
            return None

    @staticmethod
    def update_project(agent_id: str, name: str, updates: dict) -> dict:
        """Update project.json fields."""
        pdir = ProjectManager._project_dir(agent_id, name)
        cfg_path = os.path.join(pdir, "project.json")
        if not os.path.exists(cfg_path):
            return {"error": f"Project '{name}' not found"}
        try:
            with open(cfg_path, "r") as f:
                cfg = json.load(f)
            for k in ("description", "watch_folders", "tags", "model", "name", "icon", "status", "instructions"):
                if k in updates:
                    cfg[k] = updates[k]
            with open(cfg_path, "w") as f:
                json.dump(cfg, f, indent=2)
            return {"name": name, "status": "updated"}
        except (OSError, json.JSONDecodeError) as e:
            return {"error": str(e)}

    @staticmethod
    def delete_project(agent_id: str, name: str) -> dict:
        """Soft-delete a project (move to .trash)."""
        pdir = ProjectManager._project_dir(agent_id, name)
        if not os.path.isdir(pdir):
            return {"error": f"Project '{name}' not found"}
        trash_dir = os.path.join(AGENTS_DIR, ".trash")
        os.makedirs(trash_dir, exist_ok=True)
        dest = os.path.join(trash_dir, f"{agent_id}_project_{name}_{int(time.time())}")
        shutil.move(pdir, dest)
        # Remove QMD collection
        collection_name = f"{agent_id}/{name}"
        def _find_qmd():
            p = shutil.which("qmd")
            if p:
                return p
            # Common locations when running under launchd
            for candidate in [
                os.path.expanduser("~/.nvm/versions/node/v22.20.0/bin/qmd"),
                "/opt/homebrew/bin/qmd",
                "/usr/local/bin/qmd",
            ]:
                if os.path.isfile(candidate):
                    return candidate
            return "qmd"
        try:
            subprocess.run(
                [_find_qmd(), "collection", "remove", collection_name],
                capture_output=True, text=True, timeout=5,
            )
        except Exception:
            pass
        return {"name": name, "status": "deleted", "moved_to": dest}


class NoteManager:
    """CRUD operations for project notes (markdown files in notes/ subdirectory)."""

    @staticmethod
    def _notes_dir(agent_id: str, project_name: str) -> str:
        return os.path.join(AGENTS_DIR, agent_id, "projects", project_name, "notes")

    @staticmethod
    def list_notes(agent_id: str, project_name: str) -> list[dict]:
        """Walk notes/ tree recursively, return metadata for each .md file."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        if not os.path.isdir(notes_dir):
            return []
        results = []
        for dirpath, _, filenames in os.walk(notes_dir):
            for fname in sorted(filenames):
                if not fname.endswith(".md"):
                    continue
                fpath = os.path.join(dirpath, fname)
                rel_path = os.path.relpath(fpath, notes_dir)
                try:
                    stat = os.stat(fpath)
                    with open(fpath, "r", errors="replace") as f:
                        raw = f.read(2000)
                    fm, _ = _parse_frontmatter(raw)
                    results.append({
                        "path": rel_path,
                        "name": fm.get("name", fname.replace(".md", "")),
                        "type": fm.get("type", "note"),
                        "size": stat.st_size,
                        "created_at": fm.get("created_at", ""),
                        "updated_at": fm.get("updated_at", ""),
                    })
                except Exception:
                    continue
        return results

    @staticmethod
    def get_note(agent_id: str, project_name: str, path: str) -> dict | None:
        """Read a note file and return its content with metadata."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        if not os.path.isfile(fpath):
            return None
        try:
            stat = os.stat(fpath)
            with open(fpath, "r", errors="replace") as f:
                raw = f.read()
            fm, body = _parse_frontmatter(raw)
            return {
                "path": path,
                "name": fm.get("name", os.path.basename(path).replace(".md", "")),
                "content": body,
                "frontmatter": fm,
                "size": stat.st_size,
                "created_at": fm.get("created_at", ""),
                "updated_at": fm.get("updated_at", ""),
            }
        except Exception:
            return None

    @staticmethod
    def create_note(agent_id: str, project_name: str, path: str, content: str = "") -> dict:
        """Create a note with YAML frontmatter, entity extraction, and QMD reindex."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        # Create parent dirs if needed
        os.makedirs(os.path.dirname(fpath), exist_ok=True)
        if os.path.exists(fpath):
            return {"error": f"Note '{path}' already exists"}

        now = datetime.datetime.now(datetime.timezone.utc).isoformat()
        name = os.path.basename(path).replace(".md", "")
        md_content = f"""---
name: {_yaml_escape(name)}
type: note
created_at: {now}
updated_at: {now}
agent: {agent_id}
project: {project_name}
---

{content}
"""
        with open(fpath, "w") as f:
            f.write(md_content)

        # Entity extraction + auto-link
        try:
            entities = _extract_entities(content)
            if entities:
                filename = os.path.basename(fpath)
                matches = _find_entity_matches(agent_id, filename, entities)
                for other_fname in matches[:10]:
                    other_path = os.path.join(os.path.dirname(fpath), other_fname)
                    if not os.path.exists(other_path):
                        # Try agent dir
                        other_path = os.path.join(AGENTS_DIR, agent_id, other_fname)
                    if os.path.exists(other_path):
                        _add_related_to_file(fpath, other_fname, "same_topic")
                        _add_related_to_file(other_path, filename, "same_topic")
                _update_entity_index(agent_id, filename, entities)
        except Exception:
            pass

        # Add same_folder relationships with sibling notes
        try:
            folder = os.path.dirname(fpath)
            fname = os.path.basename(fpath)
            for sibling in os.listdir(folder):
                if sibling.endswith(".md") and sibling != fname:
                    sibling_path = os.path.join(folder, sibling)
                    _add_related_to_file(fpath, sibling, "same_folder")
                    _add_related_to_file(sibling_path, fname, "same_folder")
        except Exception:
            pass

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"path": path, "status": "created"}

    @staticmethod
    def update_note(agent_id: str, project_name: str, path: str, content: str) -> dict:
        """Update a note's content, preserving frontmatter and updating timestamp."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        if not os.path.isfile(fpath):
            return {"error": f"Note '{path}' not found"}

        try:
            with open(fpath, "r") as f:
                raw = f.read()
            fm, _ = _parse_frontmatter(raw)
        except Exception:
            fm = {}

        now = datetime.datetime.now(datetime.timezone.utc).isoformat()
        fm["updated_at"] = now

        # Rebuild frontmatter
        fm_lines = []
        for k, v in fm.items():
            if k == "related":
                continue  # related is multi-line, handle separately
            fm_lines.append(f"{k}: {v}")

        # Preserve related section if it exists
        related_section = ""
        fm_match = re.match(r'^---\s*\n(.*?)\n---\s*\n', raw, re.DOTALL)
        if fm_match:
            fm_text = fm_match.group(1)
            rel_match = re.search(r'(related:.*)', fm_text, re.DOTALL)
            if rel_match:
                related_section = "\n" + rel_match.group(1)

        md_content = f"---\n" + "\n".join(fm_lines) + related_section + f"\n---\n\n{content}\n"
        with open(fpath, "w") as f:
            f.write(md_content)

        # Re-run entity extraction
        try:
            entities = _extract_entities(content)
            if entities:
                filename = os.path.basename(fpath)
                _update_entity_index(agent_id, filename, entities)
        except Exception:
            pass

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"path": path, "status": "updated"}

    @staticmethod
    def delete_note(agent_id: str, project_name: str, path: str) -> dict:
        """Remove a note file and trigger QMD reindex."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, path)
        if not os.path.isfile(fpath):
            return {"error": f"Note '{path}' not found"}

        os.remove(fpath)

        # Clean up empty parent directories
        parent = os.path.dirname(fpath)
        while parent != notes_dir:
            try:
                if not os.listdir(parent):
                    os.rmdir(parent)
                    parent = os.path.dirname(parent)
                else:
                    break
            except Exception:
                break

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"path": path, "status": "deleted"}

    @staticmethod
    def rename_note(agent_id: str, project_name: str, old_path: str, new_path: str) -> dict:
        """Rename/move a note file and update its frontmatter name."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        old_fpath = os.path.join(notes_dir, old_path)
        new_fpath = os.path.join(notes_dir, new_path)
        if not os.path.isfile(old_fpath):
            return {"error": f"Note '{old_path}' not found"}
        if os.path.exists(new_fpath):
            return {"error": f"Note '{new_path}' already exists"}

        # Create parent dirs for new path if needed
        os.makedirs(os.path.dirname(new_fpath), exist_ok=True)
        os.rename(old_fpath, new_fpath)

        # Update frontmatter name field
        try:
            with open(new_fpath, "r") as f:
                raw = f.read()
            new_name = os.path.basename(new_path).replace(".md", "")
            raw = re.sub(r'^(name:\s*).*$', rf'\g<1>{_yaml_escape(new_name)}', raw, count=1, flags=re.MULTILINE)
            now = datetime.datetime.now(datetime.timezone.utc).isoformat()
            raw = re.sub(r'^(updated_at:\s*).*$', rf'\g<1>{now}', raw, count=1, flags=re.MULTILINE)
            with open(new_fpath, "w") as f:
                f.write(raw)
        except Exception:
            pass

        # Trigger QMD reindex
        collection = f"{agent_id}/{project_name}"
        _qmd_debounced_embed(collection)

        return {"old_path": old_path, "new_path": new_path, "status": "renamed"}

    @staticmethod
    def create_folder(agent_id: str, project_name: str, folder_path: str) -> dict:
        """Create a folder within the notes directory."""
        notes_dir = NoteManager._notes_dir(agent_id, project_name)
        fpath = os.path.join(notes_dir, folder_path)
        os.makedirs(fpath, exist_ok=True)
        return {"path": folder_path, "status": "created"}


# ─── Document Ingestion Engine ────────────────────────────────────────

class DocumentParser:
    """Parse various document formats to plain text."""

    @staticmethod
    def parse_pdf(path: str) -> str:
        """Parse PDF to text using pymupdf."""
        try:
            import fitz  # pymupdf
        except ImportError:
            raise ImportError("Install pymupdf for PDF support: pip3 install pymupdf")
        doc = fitz.open(path)
        pages = []
        for page in doc:
            pages.append(page.get_text())
        doc.close()
        return "\n\n".join(pages)

    @staticmethod
    def parse_docx(path: str) -> str:
        """Parse DOCX to text using python-docx."""
        try:
            import docx
        except ImportError:
            raise ImportError("Install python-docx for DOCX support: pip3 install python-docx")
        doc = docx.Document(path)
        paragraphs = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # Preserve heading structure
                if para.style and para.style.name and para.style.name.startswith("Heading"):
                    level = 1
                    try:
                        level = int(para.style.name.replace("Heading", "").strip()) or 1
                    except ValueError:
                        pass
                    text = "#" * level + " " + text
                paragraphs.append(text)
        return "\n\n".join(paragraphs)

    @staticmethod
    def parse_txt(path: str) -> str:
        """Parse plain text file."""
        with open(path, "r", errors="replace") as f:
            return f.read()

    @staticmethod
    def parse_md(path: str) -> str:
        """Parse markdown file (keep as-is)."""
        with open(path, "r", errors="replace") as f:
            return f.read()

    @staticmethod
    def parse_html(content: str) -> str:
        """Strip HTML tags and extract text content."""
        from html.parser import HTMLParser

        class _TextExtractor(HTMLParser):
            def __init__(self):
                super().__init__()
                self._parts: list[str] = []
                self._skip = False
                self._skip_tags = {"script", "style", "nav", "header", "footer", "noscript"}

            def handle_starttag(self, tag, attrs):
                if tag in self._skip_tags:
                    self._skip = True
                elif tag in ("p", "div", "br", "h1", "h2", "h3", "h4", "h5", "h6", "li", "tr"):
                    self._parts.append("\n")
                    if tag.startswith("h"):
                        level = int(tag[1])
                        self._parts.append("#" * level + " ")

            def handle_endtag(self, tag):
                if tag in self._skip_tags:
                    self._skip = False
                elif tag in ("p", "div", "h1", "h2", "h3", "h4", "h5", "h6"):
                    self._parts.append("\n")

            def handle_data(self, data):
                if not self._skip:
                    self._parts.append(data)

        extractor = _TextExtractor()
        extractor.feed(content)
        text = "".join(extractor._parts)
        # Clean up whitespace
        lines = [line.strip() for line in text.split("\n")]
        cleaned = "\n".join(lines)
        # Collapse multiple blank lines
        cleaned = re.sub(r'\n{3,}', '\n\n', cleaned)
        return cleaned.strip()

    @staticmethod
    def parse_xlsx(path: str, sheet: str | None = None) -> str:
        """Parse XLSX/XLS to markdown tables using openpyxl."""
        try:
            import openpyxl
        except ImportError:
            raise ImportError("Install openpyxl for XLSX support: pip3 install openpyxl")
        wb = openpyxl.load_workbook(path, read_only=True, data_only=False)
        parts = []
        sheet_names = wb.sheetnames
        parts.append(f"**Sheets:** {', '.join(sheet_names)}\n")
        target_sheets = [sheet] if sheet and sheet in sheet_names else sheet_names
        for sname in target_sheets:
            ws = wb[sname]
            parts.append(f"## Sheet: {sname}\n")
            rows = []
            for row in ws.iter_rows(values_only=False):
                cells = []
                for cell in row:
                    val = cell.value
                    if val is None:
                        cells.append("")
                    elif isinstance(val, str) and val.startswith("="):
                        # Show formula
                        cells.append(f"{val}")
                    else:
                        cells.append(str(val))
                rows.append(cells)
            if not rows:
                parts.append("*(empty sheet)*\n")
                continue
            # Build markdown table
            max_cols = max(len(r) for r in rows) if rows else 0
            # Pad rows to same length
            for r in rows:
                while len(r) < max_cols:
                    r.append("")
            # Header row
            header = "| " + " | ".join(rows[0]) + " |"
            sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
            table_lines = [header, sep]
            for r in rows[1:]:
                table_lines.append("| " + " | ".join(r) + " |")
            parts.append("\n".join(table_lines) + "\n")
        wb.close()
        return "\n".join(parts)

    @staticmethod
    def parse_pptx(path: str, slides: str | None = None) -> str:
        """Parse PPTX to text using python-pptx."""
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("Install python-pptx for PPTX support: pip3 install python-pptx")
        prs = Presentation(path)
        total_slides = len(prs.slides)
        parts = [f"**Slides:** {total_slides}\n"]
        # Parse slide range
        slide_indices = None
        if slides:
            slide_indices = set()
            for part in slides.split(","):
                part = part.strip()
                if "-" in part:
                    a, b = part.split("-", 1)
                    for i in range(int(a), int(b) + 1):
                        slide_indices.add(i)
                else:
                    slide_indices.add(int(part))
        for idx, slide in enumerate(prs.slides, 1):
            if slide_indices and idx not in slide_indices:
                continue
            parts.append(f"## Slide {idx}")
            # Title
            if slide.shapes.title:
                parts.append(f"**Title:** {slide.shapes.title.text}")
            # Body text
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
                if shape.has_table:
                    table = shape.table
                    rows = []
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    if rows:
                        max_cols = max(len(r) for r in rows)
                        for r in rows:
                            while len(r) < max_cols:
                                r.append("")
                        header = "| " + " | ".join(rows[0]) + " |"
                        sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
                        table_lines = [header, sep]
                        for r in rows[1:]:
                            table_lines.append("| " + " | ".join(r) + " |")
                        parts.append("\n".join(table_lines))
            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"*Speaker Notes:* {notes}")
            parts.append("")
        return "\n".join(parts)

    @staticmethod
    def parse_csv(path: str) -> str:
        """Parse CSV/TSV to markdown table."""
        import csv
        delimiter = "\t" if path.lower().endswith(".tsv") else ","
        with open(path, "r", newline="", errors="replace") as f:
            reader = csv.reader(f, delimiter=delimiter)
            rows = [row for row in reader]
        if not rows:
            return "*(empty file)*"
        max_cols = max(len(r) for r in rows)
        for r in rows:
            while len(r) < max_cols:
                r.append("")
        header = "| " + " | ".join(rows[0]) + " |"
        sep = "| " + " | ".join("---" for _ in range(max_cols)) + " |"
        lines = [header, sep]
        for r in rows[1:]:
            lines.append("| " + " | ".join(r) + " |")
        return "\n".join(lines)

    @staticmethod
    def parse_image(path: str) -> str:
        """Parse image metadata using Pillow. Returns metadata text."""
        try:
            from PIL import Image
        except ImportError:
            raise ImportError("Install Pillow for image support: pip3 install Pillow")
        img = Image.open(path)
        width, height = img.size
        fmt = img.format or os.path.splitext(path)[1].lstrip(".")
        mode = img.mode
        info_parts = [
            f"**Image:** {os.path.basename(path)}",
            f"**Dimensions:** {width} x {height}",
            f"**Format:** {fmt}",
            f"**Mode:** {mode}",
        ]
        # Extract EXIF if available
        try:
            exif = img.getexif()
            if exif:
                for tag_id, value in list(exif.items())[:10]:
                    try:
                        from PIL.ExifTags import TAGS
                        tag_name = TAGS.get(tag_id, str(tag_id))
                        info_parts.append(f"**{tag_name}:** {value}")
                    except Exception:
                        pass
        except Exception:
            pass
        img.close()
        return "\n".join(info_parts)

    @staticmethod
    def parse_svg(path: str) -> str:
        """Parse SVG to extract text elements and metadata."""
        from xml.etree import ElementTree
        tree = ElementTree.parse(path)
        root = tree.getroot()
        ns = {"svg": "http://www.w3.org/2000/svg"}
        parts = [f"**SVG:** {os.path.basename(path)}"]
        # Get dimensions
        w = root.get("width", "")
        h = root.get("height", "")
        vb = root.get("viewBox", "")
        if w and h:
            parts.append(f"**Dimensions:** {w} x {h}")
        if vb:
            parts.append(f"**ViewBox:** {vb}")
        # Extract title and desc
        for tag in ("title", "desc"):
            el = root.find(f"svg:{tag}", ns) or root.find(tag)
            if el is not None and el.text:
                parts.append(f"**{tag.capitalize()}:** {el.text.strip()}")
        # Extract all text elements
        texts = []
        for text_el in list(root.iter(f"{{{ns['svg']}}}text")) + list(root.iter("text")):
            t = "".join(text_el.itertext()).strip()
            if t:
                texts.append(t)
        if texts:
            parts.append(f"\n**Text content:**")
            for t in texts:
                parts.append(f"- {t}")
        return "\n".join(parts)

    @staticmethod
    def parse_url(url: str) -> str:
        """Fetch URL and parse HTML to text."""
        req_headers = {
            "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) AppleWebKit/537.36",
            "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        }
        req = urllib.request.Request(url, headers=req_headers)
        with urllib.request.urlopen(req, timeout=30) as resp:
            raw = resp.read()
            encoding = resp.headers.get("Content-Encoding", "")
            if encoding == "gzip":
                import gzip
                raw = gzip.decompress(raw)
            charset = resp.headers.get_content_charset() or "utf-8"
            html = raw.decode(charset, errors="replace")
        return DocumentParser.parse_html(html)

    @staticmethod
    def parse(path_or_url: str) -> tuple[str, str]:
        """Auto-detect format and parse. Returns (text, source_type)."""
        if path_or_url.startswith(("http://", "https://")):
            return DocumentParser.parse_url(path_or_url), "url"
        ext = os.path.splitext(path_or_url)[1].lower()
        parsers = {
            ".pdf": ("pdf", DocumentParser.parse_pdf),
            ".docx": ("docx", DocumentParser.parse_docx),
            ".txt": ("txt", DocumentParser.parse_txt),
            ".md": ("md", DocumentParser.parse_md),
            ".html": ("html", lambda p: DocumentParser.parse_html(open(p, "r", errors="replace").read())),
            ".htm": ("html", lambda p: DocumentParser.parse_html(open(p, "r", errors="replace").read())),
            ".xlsx": ("xlsx", DocumentParser.parse_xlsx),
            ".xls": ("xlsx", DocumentParser.parse_xlsx),
            ".pptx": ("pptx", DocumentParser.parse_pptx),
            ".csv": ("csv", DocumentParser.parse_csv),
            ".tsv": ("csv", DocumentParser.parse_csv),
            ".png": ("image", DocumentParser.parse_image),
            ".jpg": ("image", DocumentParser.parse_image),
            ".jpeg": ("image", DocumentParser.parse_image),
            ".gif": ("image", DocumentParser.parse_image),
            ".webp": ("image", DocumentParser.parse_image),
            ".bmp": ("image", DocumentParser.parse_image),
            ".svg": ("svg", DocumentParser.parse_svg),
        }
        if ext not in parsers:
            raise ValueError(f"Unsupported format: {ext}. Supported: {', '.join(parsers.keys())}")
        source_type, parser_fn = parsers[ext]
        return parser_fn(path_or_url), source_type


class DocumentChunker:
    """Split text into overlapping chunks with section header preservation."""

    @staticmethod
    def chunk(text: str, chunk_size: int = 1500, chunk_overlap: int = 200,
              min_chunk_size: int = 100) -> list[dict]:
        """Split text into chunks. chunk_size/overlap are in ~tokens (chars/4 approximation).
        Returns list of {text, index, total, header}."""
        # Convert token counts to char approximation
        max_chars = chunk_size * 4
        overlap_chars = chunk_overlap * 4
        min_chars = min_chunk_size * 4

        # Split into paragraphs
        paragraphs = re.split(r'\n\s*\n', text)
        paragraphs = [p.strip() for p in paragraphs if p.strip()]

        chunks: list[dict] = []
        current_parts: list[str] = []
        current_len = 0
        last_header = ""

        def _flush():
            nonlocal current_parts, current_len
            if not current_parts:
                return
            chunk_text = "\n\n".join(current_parts)
            if len(chunk_text) < min_chars:
                return
            # Prepend section header if available
            if last_header and not chunk_text.startswith("#"):
                chunk_text = last_header + "\n\n" + chunk_text
            chunks.append({
                "text": chunk_text,
                "index": len(chunks),
                "total": 0,  # filled in later
                "header": last_header,
            })
            # Keep overlap: take text from end of current chunk
            if overlap_chars > 0:
                overlap_text = chunk_text[-overlap_chars:]
                current_parts = [overlap_text]
                current_len = len(overlap_text)
            else:
                current_parts = []
                current_len = 0

        for para in paragraphs:
            # Track section headers
            header_match = re.match(r'^(#{1,6}\s+.+)', para.split('\n')[0])
            if header_match:
                last_header = header_match.group(1)

            # If a single paragraph exceeds max_chars, split it
            if len(para) > max_chars:
                # Flush current buffer first
                if current_parts:
                    _flush()
                # Split on sentences
                sentences = re.split(r'(?<=[.!?])\s+', para)
                for sent in sentences:
                    if len(sent) > max_chars:
                        # Split on words as last resort
                        words = sent.split()
                        for word in words:
                            if current_len + len(word) + 1 > max_chars:
                                _flush()
                            current_parts.append(word)
                            current_len += len(word) + 1
                    else:
                        if current_len + len(sent) + 1 > max_chars:
                            _flush()
                        current_parts.append(sent)
                        current_len += len(sent) + 1
            else:
                if current_len + len(para) + 2 > max_chars:
                    _flush()
                current_parts.append(para)
                current_len += len(para) + 2

        # Flush remaining
        if current_parts:
            chunk_text = "\n\n".join(current_parts)
            if len(chunk_text) >= min_chars:
                if last_header and not chunk_text.startswith("#"):
                    chunk_text = last_header + "\n\n" + chunk_text
                chunks.append({
                    "text": chunk_text,
                    "index": len(chunks),
                    "total": 0,
                    "header": last_header,
                })

        # Fill in total count
        for c in chunks:
            c["total"] = len(chunks)

        return chunks


class IngestManager:
    """Ingest files and URLs into agent or project memory as chunked markdown."""

    @staticmethod
    def _source_hash(source: str) -> str:
        """6-char hash of source name/URL."""
        return hashlib.sha256(source.encode()).hexdigest()[:6]

    @staticmethod
    def _ingest_dir(agent_id: str, project_name: str | None = None) -> str:
        """Get the directory where ingested chunks are stored."""
        if project_name:
            d = os.path.join(AGENTS_DIR, agent_id, "projects", project_name, "ingested")
        else:
            d = os.path.join(AGENTS_DIR, agent_id, "ingested")
        os.makedirs(d, exist_ok=True)
        return d

    @staticmethod
    def _collection_name(agent_id: str, project_name: str | None = None) -> str:
        """QMD collection name for embedding."""
        if project_name:
            return f"{agent_id}/{project_name}"
        return agent_id

    @staticmethod
    def ingest_file(agent_id: str, file_path: str,
                    project_name: str | None = None,
                    tags: list[str] | None = None,
                    chunk_size: int = 1500, chunk_overlap: int = 200) -> dict:
        """Parse, chunk, and store a file as ingested memory chunks."""
        if not os.path.exists(file_path):
            return {"error": f"File not found: {file_path}"}
        source_name = os.path.basename(file_path)
        try:
            text, source_type = DocumentParser.parse(file_path)
        except (ImportError, ValueError) as e:
            return {"error": str(e)}
        return IngestManager._store_chunks(
            agent_id, project_name, source_name, source_type, text,
            tags=tags, chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        )

    @staticmethod
    def ingest_url(agent_id: str, url: str,
                   project_name: str | None = None,
                   tags: list[str] | None = None,
                   chunk_size: int = 1500, chunk_overlap: int = 200) -> dict:
        """Fetch URL, parse HTML, chunk, and store."""
        try:
            text = DocumentParser.parse_url(url)
        except Exception as e:
            return {"error": f"Failed to fetch URL: {e}"}
        return IngestManager._store_chunks(
            agent_id, project_name, url, "url", text,
            tags=tags, chunk_size=chunk_size, chunk_overlap=chunk_overlap,
        )

    @staticmethod
    def _store_chunks(agent_id: str, project_name: str | None,
                      source: str, source_type: str, text: str,
                      tags: list[str] | None = None,
                      chunk_size: int = 1500, chunk_overlap: int = 200) -> dict:
        """Chunk text and write as ingest-*.md files with frontmatter."""
        src_hash = IngestManager._source_hash(source)
        ingest_dir = IngestManager._ingest_dir(agent_id, project_name)
        collection = IngestManager._collection_name(agent_id, project_name)

        # Delete existing chunks for this source (re-ingest)
        existing = [f for f in os.listdir(ingest_dir) if f.startswith(f"ingest-{src_hash}-") and f.endswith(".md")]
        for f in existing:
            os.remove(os.path.join(ingest_dir, f))

        # Chunk
        chunks = DocumentChunker.chunk(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        if not chunks:
            return {"error": "No content extracted from document"}

        all_tags = ["ingested"]
        if tags:
            all_tags.extend(tags)
        # Add source name as tag (sanitized)
        safe_source_tag = re.sub(r'[^\w-]', '', source.split("/")[-1].split(".")[0].lower())
        if safe_source_tag:
            all_tags.append(safe_source_tag)
        tags_yaml = "\n".join(f"  - {t}" for t in all_tags)

        now = datetime.datetime.now(datetime.timezone.utc).isoformat()
        files_written = []
        for chunk in chunks:
            idx = chunk["index"]
            total = chunk["total"]
            title = chunk["header"] or f"{source} - Chunk {idx + 1}"

            # Build related links
            related_lines = []
            if idx > 0:
                prev_file = f"ingest-{src_hash}-{idx - 1:03d}.md"
                related_lines.append(f"  - file: {prev_file}\n    type: prev_chunk")
            if idx < total - 1:
                next_file = f"ingest-{src_hash}-{idx + 1:03d}.md"
                related_lines.append(f"  - file: {next_file}\n    type: next_chunk")
            if idx != 0:
                first_file = f"ingest-{src_hash}-000.md"
                related_lines.append(f"  - file: {first_file}\n    type: same_source")
            related_yaml = ""
            if related_lines:
                related_yaml = "related:\n" + "\n".join(related_lines) + "\n"

            filename = f"ingest-{src_hash}-{idx:03d}.md"
            md_content = f"""---
title: {_yaml_escape(title)}
source: {_yaml_escape(source)}
source_type: {source_type}
ingested_at: "{now}"
chunk_index: {idx}
total_chunks: {total}
agent: {agent_id}
tags:
{tags_yaml}
{related_yaml}---

{chunk['text']}
"""
            fpath = os.path.join(ingest_dir, filename)
            with open(fpath, "w") as f:
                f.write(md_content)
            files_written.append(filename)

        # Trigger QMD indexing
        _qmd_debounced_embed(collection)

        word_count = len(text.split())
        return {
            "source": source,
            "source_type": source_type,
            "source_hash": src_hash,
            "chunks": len(chunks),
            "words": word_count,
            "files": files_written,
            "agent": agent_id,
            "project": project_name,
            "status": "ingested",
        }

    @staticmethod
    def list_ingested(agent_id: str, project_name: str | None = None) -> list[dict]:
        """List ingested documents grouped by source."""
        ingest_dir = IngestManager._ingest_dir(agent_id, project_name)
        if not os.path.isdir(ingest_dir):
            return []
        # Group by source hash
        groups: dict[str, dict] = {}
        for fname in os.listdir(ingest_dir):
            if not fname.startswith("ingest-") or not fname.endswith(".md"):
                continue
            fpath = os.path.join(ingest_dir, fname)
            try:
                with open(fpath, "r") as f:
                    raw = f.read(800)
                fm, _ = _parse_frontmatter(raw)
            except Exception:
                continue
            source = fm.get("source", "unknown")
            src_hash = fname.split("-")[1] if "-" in fname else "?"
            if src_hash not in groups:
                groups[src_hash] = {
                    "source": source,
                    "source_type": fm.get("source_type", "unknown"),
                    "source_hash": src_hash,
                    "chunks": 0,
                    "ingested_at": fm.get("ingested_at", ""),
                    "tags": [],
                }
            groups[src_hash]["chunks"] += 1
            # Parse tags from frontmatter
            tags_str = fm.get("tags", "")
            if isinstance(tags_str, str) and tags_str:
                for t in tags_str.split(","):
                    t = t.strip().strip("-").strip()
                    if t and t not in groups[src_hash]["tags"]:
                        groups[src_hash]["tags"].append(t)
        return sorted(groups.values(), key=lambda x: x.get("ingested_at", ""), reverse=True)

    @staticmethod
    def delete_ingested(agent_id: str, source_hash: str,
                        project_name: str | None = None) -> dict:
        """Delete all chunks for a source hash."""
        ingest_dir = IngestManager._ingest_dir(agent_id, project_name)
        if not os.path.isdir(ingest_dir):
            return {"error": "No ingested documents found"}
        deleted = 0
        source_name = ""
        for fname in os.listdir(ingest_dir):
            if fname.startswith(f"ingest-{source_hash}-") and fname.endswith(".md"):
                if not source_name:
                    fpath = os.path.join(ingest_dir, fname)
                    try:
                        with open(fpath, "r") as f:
                            fm, _ = _parse_frontmatter(f.read(500))
                        source_name = fm.get("source", "unknown")
                    except Exception:
                        pass
                os.remove(os.path.join(ingest_dir, fname))
                deleted += 1
        collection = IngestManager._collection_name(agent_id, project_name)
        _qmd_debounced_embed(collection)
        return {"source": source_name, "source_hash": source_hash, "deleted": deleted}


# ─── Watched Folders (Auto-Ingestion) ────────────────────────────────

class IngestWatcher:
    """Background thread that polls watched folders and auto-ingests new/modified files."""

    POLL_INTERVAL = 30  # seconds

    def __init__(self):
        self._stop = threading.Event()
        self._thread = None

    def start(self):
        """Start the background watcher thread."""
        self._stop.clear()
        self._thread = threading.Thread(target=self._run_loop, daemon=True, name="ingest_watcher")
        self._thread.start()

    def stop(self):
        self._stop.set()
        if self._thread:
            self._thread.join(timeout=2)

    def _run_loop(self):
        """Poll all watched folders across all agents and projects."""
        while not self._stop.is_set():
            try:
                self._scan_all()
            except Exception as e:
                logging.debug("IngestWatcher error: %s", e)
            self._stop.wait(self.POLL_INTERVAL)

    def _scan_all(self):
        """Scan watched folders for all agents and projects."""
        if not os.path.isdir(AGENTS_DIR):
            return
        for agent_name in os.listdir(AGENTS_DIR):
            if agent_name.startswith("."):
                continue
            agent_dir = os.path.join(AGENTS_DIR, agent_name)
            if not os.path.isdir(agent_dir):
                continue
            # Check agent-level watches (from agent.json)
            agent_json_path = os.path.join(agent_dir, "agent.json")
            if os.path.exists(agent_json_path):
                try:
                    with open(agent_json_path, "r") as f:
                        agent_cfg = json.load(f)
                    watches = agent_cfg.get("ingest_watch", [])
                    if watches:
                        self._process_watches(agent_name, None, watches, agent_dir)
                except (OSError, json.JSONDecodeError):
                    pass
            # Check project-level watches
            projects_dir = os.path.join(agent_dir, "projects")
            if os.path.isdir(projects_dir):
                for proj_name in os.listdir(projects_dir):
                    proj_dir = os.path.join(projects_dir, proj_name)
                    proj_json = os.path.join(proj_dir, "project.json")
                    if not os.path.exists(proj_json):
                        continue
                    try:
                        with open(proj_json, "r") as f:
                            proj_cfg = json.load(f)
                        watches = proj_cfg.get("watch_folders", [])
                        if watches:
                            self._process_watches(agent_name, proj_name, watches, proj_dir)
                    except (OSError, json.JSONDecodeError):
                        pass

    def _process_watches(self, agent_id: str, project_name: str | None,
                         watches: list[dict], base_dir: str):
        """Process watched folders, detect changes, ingest as needed."""
        registry_path = os.path.join(base_dir, "ingest_registry.json")
        registry = {}
        if os.path.exists(registry_path):
            try:
                with open(registry_path, "r") as f:
                    registry = json.load(f)
            except (OSError, json.JSONDecodeError):
                pass
        watches_reg = registry.get("watches", {})
        changed = False

        for watch in watches:
            watch_path = watch.get("path", "")
            if not watch_path or not os.path.isdir(watch_path):
                continue
            pattern = watch.get("pattern", "*")
            recursive = watch.get("recursive", False)
            tags = watch.get("tags", [])
            chunk_size = watch.get("chunk_size", 1500)

            # Get or create registry entry for this watch
            wreg = watches_reg.get(watch_path, {"files": {}, "last_scan": ""})

            # Scan for matching files
            if recursive:
                matched_files = []
                for root, _dirs, files in os.walk(watch_path):
                    for fn in files:
                        if fnmatch.fnmatch(fn, pattern):
                            matched_files.append(os.path.join(root, fn))
            else:
                matched_files = [
                    os.path.join(watch_path, fn) for fn in os.listdir(watch_path)
                    if fnmatch.fnmatch(fn, pattern) and os.path.isfile(os.path.join(watch_path, fn))
                ]

            current_files = set()
            for fpath in matched_files:
                fname = os.path.basename(fpath)
                current_files.add(fname)
                try:
                    stat = os.stat(fpath)
                except OSError:
                    continue
                prev = wreg["files"].get(fname, {})
                if prev.get("mtime") == stat.st_mtime and prev.get("size") == stat.st_size:
                    continue  # unchanged

                # New or modified file — ingest
                try:
                    result = IngestManager.ingest_file(
                        agent_id, fpath, project_name=project_name,
                        tags=tags, chunk_size=chunk_size,
                    )
                    if "error" not in result:
                        wreg["files"][fname] = {
                            "mtime": stat.st_mtime,
                            "size": stat.st_size,
                            "hash": result.get("source_hash", ""),
                            "chunks": result.get("chunks", 0),
                            "ingested_at": datetime.datetime.now(datetime.timezone.utc).isoformat(),
                        }
                        changed = True
                except Exception as e:
                    logging.debug("IngestWatcher: failed to ingest %s: %s", fpath, e)

            # Detect deleted files
            for fname in list(wreg["files"].keys()):
                if fname not in current_files:
                    src_hash = wreg["files"][fname].get("hash", "")
                    if src_hash:
                        IngestManager.delete_ingested(agent_id, src_hash, project_name)
                    del wreg["files"][fname]
                    changed = True

            wreg["last_scan"] = datetime.datetime.now(datetime.timezone.utc).isoformat()
            watches_reg[watch_path] = wreg

        if changed:
            registry["watches"] = watches_reg
            try:
                with open(registry_path, "w") as f:
                    json.dump(registry, f, indent=2)
            except OSError:
                pass


# Global IngestWatcher instance (started alongside scheduler in server.py)
_ingest_watcher: IngestWatcher | None = None


def _get_memory_store() -> MemoryStore | None:
    """Get the active memory store: thread-local (delegation/scheduler) or global (main thread)."""
    return getattr(_thread_local, 'memory_store', None) or _memory_store


def tool_memory_store(args: dict) -> str:
    """Store a memory. When a project is active, writes to project directory."""
    ms = _get_memory_store()
    if not ms:
        return _err("Memory store not initialized")
    name = args.get("name", "")
    content = args.get("content", "")
    description = args.get("description", "")
    mem_type = args.get("type", "general")
    if not name or not content:
        return _err("memory_store: name and content are required")
    # If project is active, store in project directory
    project = getattr(_thread_local, 'project', None)
    if project:
        agent_id = ms.agent_id
        proj_dir = os.path.join(AGENTS_DIR, agent_id, "projects", project)
        if os.path.isdir(proj_dir):
            proj_store = MemoryStore(agent_id=f"{agent_id}/{project}", base_dir=proj_dir)
            result = proj_store.store(name, content, description, mem_type)
            result["project"] = project
            return _ok(result)
    result = ms.store(name, content, description, mem_type)
    # Trigger near-term memory summary refresh when user-facing memories are stored
    # (skip if this IS the memory summary being written)
    if name != "Memory Summary" and mem_type in ("user", "feedback", "project"):
        try:
            agent_id = ms.agent_id if hasattr(ms, 'agent_id') else "main"
            trigger_memory_summary_refresh(agent_id)
        except Exception:
            pass
    return _ok(result)


def _graph_expand_results(results: list[dict], base_dir: str, ingest_dir: str,
                          max_hops: int = 1) -> list[dict]:
    """Follow 'related' frontmatter links from matched results for context expansion."""
    seen_files = {r.get("file_path", "") for r in results}
    expanded = list(results)
    frontier = list(results)
    for _hop in range(max_hops):
        next_frontier = []
        for r in frontier:
            fpath = r.get("file_path", "")
            if not fpath or not os.path.exists(fpath):
                continue
            try:
                with open(fpath, "r") as f:
                    raw = f.read(2000)
                fm, _ = _parse_frontmatter(raw)
            except Exception:
                continue
            # Parse related field (simple YAML list parsing)
            related_raw = fm.get("related", "")
            if not related_raw:
                continue
            # related is stored as multi-line YAML in frontmatter, parse linked files
            related_files = re.findall(r'file:\s*(\S+\.md)', raw)
            for rel_file in related_files:
                # Try ingest_dir first, then base_dir
                for search_dir in (ingest_dir, base_dir):
                    rel_path = os.path.join(search_dir, rel_file)
                    if rel_path in seen_files or not os.path.exists(rel_path):
                        continue
                    seen_files.add(rel_path)
                    try:
                        with open(rel_path, "r") as f:
                            rel_raw = f.read()
                        rel_fm, rel_body = _parse_frontmatter(rel_raw)
                        mem = {
                            "id": hashlib.sha256(rel_fm.get("name", rel_file).encode()).hexdigest()[:12],
                            "name": rel_fm.get("name", rel_fm.get("title", rel_file.replace(".md", ""))),
                            "description": rel_fm.get("description", ""),
                            "type": rel_fm.get("type", "general"),
                            "content": rel_body,
                            "file_path": rel_path,
                            "score": max(0, (r.get("score", 0.5) - 0.2)),
                            "source_scope": "related",
                        }
                        expanded.append(mem)
                        next_frontier.append(mem)
                    except Exception:
                        continue
                    break  # found in one dir, skip the other
        frontier = next_frontier
    return expanded


def tool_memory_recall(args: dict) -> str:
    """Recall memories by searching. When a project is active, searches project first."""
    ms = _get_memory_store()
    if not ms:
        return _err("Memory store not initialized")
    query = args.get("query", "")
    limit = args.get("limit", 10)
    mem_type = args.get("type")
    mode = args.get("mode", "")

    # Project-scoped search: search project collection first, then agent
    project = getattr(_thread_local, 'project', None)
    if project and query:
        agent_id = ms.agent_id
        proj_dir = os.path.join(AGENTS_DIR, agent_id, "projects", project)
        if os.path.isdir(proj_dir):
            proj_store = MemoryStore(agent_id=f"{agent_id}/{project}", base_dir=proj_dir)
            # Also search ingested subdir
            ingest_dir = os.path.join(proj_dir, "ingested")
            proj_results = proj_store.recall(query, limit, mem_type)
            # Tag project results
            for r in proj_results:
                r["source_scope"] = "project"
            # Then agent-level results
            agent_results = ms.recall(query, max(2, limit - len(proj_results)), mem_type)
            for r in agent_results:
                r["source_scope"] = "agent"
            results = proj_results + agent_results
            # Always expand via graph relationships (follow related links 1 hop)
            if results:
                results = _graph_expand_results(results, proj_dir, ingest_dir,
                                                max_hops=2 if mode == "graph" else 1)
            for r in results:
                if r.get("content") and len(r["content"]) > 4000:
                    r["content"] = r["content"][:4000] + "..."
            return _ok({"query": query, "project": project, "results": results[:limit], "count": len(results[:limit])})

    if not query:
        results = ms.list_all(mem_type)
        return _ok({"query": "", "results": results, "count": len(results)})
    results = ms.recall(query, limit, mem_type)

    # Always expand via graph relationships (1 hop default, 2 hops for explicit graph mode)
    if results:
        agent_id = ms.agent_id
        agent_dir = os.path.join(AGENTS_DIR, agent_id)
        ingest_dir = os.path.join(agent_dir, "ingested")
        results = _graph_expand_results(results, agent_dir, ingest_dir,
                                        max_hops=2 if mode == "graph" else 1)

    for r in results:
        if r.get("content") and len(r["content"]) > 4000:
            r["content"] = r["content"][:4000] + "..."

    # --- Co-recall tracking (Mechanism 3) ---
    if query and len(results) >= 2:
        try:
            result_files = [os.path.basename(r.get("file_path", "")) for r in results if r.get("file_path")]
            agent_id = ms.agent_id
            agent_dir = os.path.join(AGENTS_DIR, agent_id)
            threading.Thread(
                target=_record_recall_cooccurrence,
                args=(result_files, agent_id, agent_dir),
                daemon=True,
            ).start()
        except Exception:
            pass  # Co-recall tracking is best-effort

    return _ok({"query": query, "results": results, "count": len(results)})


def tool_memory_delete(args: dict) -> str:
    """Delete a memory."""
    ms = _get_memory_store()
    if not ms:
        return _err("Memory store not initialized")
    name = args.get("name", "")
    if not name:
        return _err("memory_delete: name is required")
    result = ms.delete(name)
    return _ok(result)


def tool_memory_shared(args: dict) -> str:
    """Access shared memory — global (main) or team (team head) scope."""
    action = args.get("action", "recall")
    scope = args.get("scope", "global")

    # Determine which agent's memory to use
    if scope == "team":
        # Find the team head for the calling agent
        caller_id = getattr(_thread_local, "delegate_agent_id", None)
        if not caller_id:
            agent = getattr(_thread_local, 'current_agent', None) or _current_agent
            caller_id = agent.agent_id if agent else "main"
        team_info = _get_agent_team_info(caller_id)
        if not team_info:
            return _err("memory_shared: agent is not in any team — use scope='global' instead")
        team_head_id = team_info["head"]
        target_agent = AgentConfig(team_head_id)
        source_label = f"{team_info['name']} (team)"
    else:
        target_agent = AgentConfig("main")
        source_label = "main (shared)"

    shared_store = MemoryStore(agent_id=target_agent.agent_id, base_dir=target_agent.memory_dir)

    if action == "store":
        name = args.get("name", "")
        content = args.get("content", "")
        description = args.get("description", "")
        mem_type = args.get("type", "general")
        if not name or not content:
            return _err("memory_shared store: name and content are required")
        result = shared_store.store(name, content, description, mem_type)
        result["source"] = source_label
        return _ok(result)
    else:  # recall
        query = args.get("query", "")
        limit = args.get("limit", 10)
        mem_type = args.get("type")
        if not query:
            results = shared_store.list_all(mem_type)
        else:
            results = shared_store.recall(query, limit, mem_type)
            # Graph expansion on shared memory too
            if results:
                shared_dir = os.path.join(AGENTS_DIR, target_agent.agent_id)
                shared_ingest = os.path.join(shared_dir, "ingested")
                results = _graph_expand_results(results, shared_dir, shared_ingest, max_hops=1)
            for r in results:
                if r.get("content") and len(r["content"]) > 4000:
                    r["content"] = r["content"][:4000] + "..."
        return _ok({"query": query, "source": source_label, "results": results[:limit], "count": len(results[:limit])})


def tool_use_skill(args: dict) -> str:
    """Load a skill's instructions into context."""
    skill_name = args.get("skill", "")
    if not skill_name:
        return _err("use_skill: skill name is required")
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    if not agent:
        return _err("use_skill: no active agent")

    body = agent.load_skill(skill_name)
    if body is None:
        available = [s.get("slug", s["name"]) for s in agent.list_skills()]
        return _err(f"use_skill: skill '{skill_name}' not found. Available: {', '.join(available) or 'none'}")

    return _ok({"skill": skill_name, "instructions": body})



# ─── Memory Summary ────────────────────────────────────────────────
# Automatic periodic synthesis of chat history and scheduled task results.
# Configured per-agent in agent.json: { "memory_summary": { "enabled": true, "frequency": "every 24h", "start_time": "03:00" } }

MEMORY_SUMMARY_DEFAULTS = {
    "enabled": True,
    "paused": False,
    "frequency": "every 24h",
    "start_time": "03:00",
}


def _get_memory_summary_config(agent_id: str) -> dict:
    """Read memory_summary settings from agent.json, merging with defaults."""
    cfg = AgentConfig(agent_id).config
    ms_cfg = cfg.get("memory_summary", {})
    result = dict(MEMORY_SUMMARY_DEFAULTS)
    if isinstance(ms_cfg, dict):
        for k in ("enabled", "paused", "frequency", "start_time", "model", "model_fallback"):
            if k in ms_cfg:
                result[k] = ms_cfg[k]
    elif isinstance(ms_cfg, bool):
        result["enabled"] = ms_cfg
    return result


def _cleanup_orphaned_chat_index_files(agent_id: str):
    """Remove chats-indexed files for sessions that no longer exist in the DB."""
    chats_dir = os.path.join(AGENTS_DIR, agent_id, "chats-indexed")
    if not os.path.isdir(chats_dir):
        return
    chat_db_path = os.path.join(AGENTS_DIR, "main", "chats.db")
    if not os.path.exists(chat_db_path):
        return
    try:
        conn = sqlite3.connect(chat_db_path, timeout=5)
        existing_sids = {row[0] for row in conn.execute("SELECT id FROM sessions")}
        conn.close()
    except Exception:
        return
    removed = 0
    for fname in os.listdir(chats_dir):
        if not fname.startswith("chat-") or not fname.endswith(".md"):
            continue
        sid = fname.split("-", 1)[1].rsplit("-", 1)[0]
        if sid not in existing_sids:
            try:
                os.remove(os.path.join(chats_dir, fname))
                removed += 1
            except OSError:
                pass
    if removed:
        logging.info("Chat index cleanup for %s: removed %d orphaned files", agent_id, removed)
        _qmd_debounced_embed(agent_id)


def _memory_summary_schedule_name(agent_id: str) -> str:
    return f"_memory_summary_{agent_id}"


def _memory_summary_schedule_str(cfg: dict) -> str:
    """Convert memory_summary config to a scheduler schedule string.
    e.g. 'daily 03:00' or 'every 24h' depending on whether start_time is set."""
    freq = cfg.get("frequency", "every 24h").strip().lower()
    start_time = cfg.get("start_time", "").strip()
    # If frequency is a simple 'every Xh/Xd' and start_time is provided, use 'daily HH:MM'
    if start_time and re.match(r'^\d{1,2}:\d{2}$', start_time):
        # For 24h frequency, use daily at start_time
        m = re.match(r'every\s+(\d+)\s*(h|hour|d|day)', freq)
        if m:
            val, unit = int(m.group(1)), m.group(2)[0]
            if (unit == 'h' and val == 24) or (unit == 'd' and val == 1):
                return f"daily {start_time}"
            elif unit == 'h' and val < 24:
                return freq  # sub-daily interval, ignore start_time
            elif unit == 'd' and val > 1:
                return freq  # multi-day, keep as interval
    return freq


def _gather_recent_chat_history(agent_id: str, hours: int = 48, max_sessions: int = 10,
                                 max_messages_per_session: int = 20) -> str:
    """Read recent chat sessions and messages for an agent directly from chats.db.
    Returns a formatted text digest suitable for inclusion in a prompt."""
    # Chat DB is always in main agent dir (shared across all agents)
    chat_db_path = os.path.join(AGENTS_DIR, "main", "chats.db")
    if not os.path.exists(chat_db_path):
        return ""
    conn = None
    try:
        conn = sqlite3.connect(chat_db_path, timeout=5)
        conn.row_factory = sqlite3.Row
        cutoff = time.time() - (hours * 3600)
        sessions = conn.execute(
            "SELECT id, title, model, status, created_at, last_active FROM sessions "
            "WHERE agent_id = ? AND last_active > ? "
            "AND status IN ('active', 'archived') "
            "ORDER BY last_active DESC LIMIT ?",
            (agent_id, cutoff, max_sessions)
        ).fetchall()
        # Note: 'incognito' status sessions are excluded by the IN clause above
        if not sessions:
            conn.close()
            return ""
        lines = []
        for sess in sessions:
            ts = datetime.datetime.fromtimestamp(sess["last_active"]).strftime("%Y-%m-%d %H:%M") if sess["last_active"] else "?"
            title = sess["title"] or "(untitled)"
            status_tag = " [archived]" if sess["status"] == "archived" else ""
            lines.append(f"\n### Session: {title} ({ts}){status_tag}")
            msgs = conn.execute(
                "SELECT role, content FROM messages WHERE session_id = ? ORDER BY id",
                (sess["id"],)
            ).fetchall()
            # Take first + last messages to stay within budget
            msg_list = list(msgs)
            if len(msg_list) > max_messages_per_session:
                selected = msg_list[:max_messages_per_session // 2] + msg_list[-(max_messages_per_session // 2):]
                lines.append(f"  ({len(msg_list)} messages total, showing first and last {max_messages_per_session // 2})")
            else:
                selected = msg_list
            for msg in selected:
                role = msg["role"].upper()
                content = msg["content"]
                # Parse JSON content (tool calls etc)
                try:
                    parsed = json.loads(content)
                    if isinstance(parsed, list):
                        # Multi-part message (text + tool_use blocks)
                        text_parts = [p.get("text", "") for p in parsed if isinstance(p, dict) and p.get("type") == "text"]
                        content = " ".join(text_parts).strip() or "(tool calls)"
                    elif isinstance(parsed, dict):
                        content = parsed.get("text", str(parsed))
                    elif isinstance(parsed, str):
                        content = parsed
                except (json.JSONDecodeError, TypeError):
                    pass
                # Truncate long messages
                if isinstance(content, str) and len(content) > 400:
                    content = content[:400] + "..."
                if content and content != "(tool calls)":
                    lines.append(f"  [{role}] {content}")
        return "\n".join(lines)
    except Exception as e:
        return f"(Error reading chat history: {e})"
    finally:
        if conn:
            try:
                conn.close()
            except Exception:
                pass


def _gather_recent_schedule_history(agent_id: str, hours: int = 48, limit: int = 15) -> str:
    """Read recent scheduled task execution results for an agent from scheduler.db.
    Returns a formatted text digest."""
    if not _scheduler:
        return ""
    try:
        with _sched_conn() as conn:
            conn.row_factory = sqlite3.Row
            cutoff = (datetime.datetime.now() - datetime.timedelta(hours=hours)).isoformat()
            rows = conn.execute(
                "SELECT schedule_name, task, status, result, started_at, finished_at "
                "FROM schedule_history WHERE agent = ? AND finished_at > ? "
                "AND schedule_name NOT LIKE '_memory_summary_%' "
                "ORDER BY finished_at DESC LIMIT ?",
                (agent_id, cutoff, limit)
            ).fetchall()
        if not rows:
            return ""
        lines = []
        for r in rows:
            ts = r["finished_at"][:16] if r["finished_at"] else "?"
            status = r["status"] or "unknown"
            name = r["schedule_name"] or "(unnamed)"
            lines.append(f"\n### Task: {name} [{status}] ({ts})")
            lines.append(f"  Prompt: {(r['task'] or '')[:200]}")
            result = r["result"] or ""
            if len(result) > 1000:
                result = result[:1000] + "..."
            if result:
                lines.append(f"  Result: {result}")
        return "\n".join(lines)
    except Exception as e:
        return f"(Error reading schedule history: {e})"


def _build_memory_summary_prompt(agent_id: str) -> str:
    """Build the task prompt for the memory summary scheduled task.
    Gathers real chat and schedule history and injects it into the prompt."""
    now = datetime.datetime.now()
    # Determine lookback based on frequency config
    ms_cfg = _get_memory_summary_config(agent_id)
    freq = ms_cfg.get("frequency", "every 24h")
    # Parse hours from frequency for lookback window (double it for overlap)
    lookback_hours = 48  # default
    m = re.match(r'every\s+(\d+)\s*(h|hour|d|day)', freq.lower())
    if m:
        val, unit = int(m.group(1)), m.group(2)[0]
        if unit == 'h':
            lookback_hours = val * 2
        elif unit == 'd':
            lookback_hours = val * 48

    chat_digest = _gather_recent_chat_history(agent_id, hours=lookback_hours)
    schedule_digest = _gather_recent_schedule_history(agent_id, hours=lookback_hours)

    # Build the data section
    data_section = ""
    if chat_digest:
        data_section += f"\n## Recent Conversations (last {lookback_hours}h)\n{chat_digest}\n"
    else:
        data_section += "\n## Recent Conversations\nNo conversations found in this period.\n"

    if schedule_digest:
        data_section += f"\n## Recent Scheduled Task Results (last {lookback_hours}h)\n{schedule_digest}\n"
    else:
        data_section += "\n## Recent Scheduled Task Results\nNo scheduled task executions found in this period.\n"

    return f"""Perform a Memory Summary update for agent '{agent_id}'.

Your job is to create or update a structured synthesis of recent activity. Below is the actual data from your recent conversations and scheduled task executions.

{data_section}

---

Now do the following:

1. Use memory_recall with query "Memory Summary" to find any existing summary.

2. Analyze the conversation and task data above, then write or update the synthesis using the following structured sections. Keep each section concise — omit a section if there is nothing relevant for it.

   **## User Profile & Context**
   Who the user is, their role, responsibilities, and domain knowledge. What kind of collaboration they prefer.

   **## Communication & Working Style**
   Communication preferences, response style, level of detail they expect, what to avoid.

   **## Technical Preferences**
   Coding style, preferred languages/tools/frameworks, conventions, architecture decisions.

   **## Active Projects & Ongoing Work**
   Current projects, their status, key decisions made, open questions, deadlines.

   **## Task Execution Insights**
   Patterns from scheduled task outcomes — what works, what fails, recurring issues.

   **## Key Decisions & Context**
   Important decisions, user feedback, and context that should inform future interactions.

3. Store the updated synthesis using memory_store with:
   - name: "Memory Summary"
   - type: "general"
   - description: "Auto-generated synthesis of recent conversations and task executions, updated periodically"
   - content: The structured synthesis (300-800 words)

Focus on actionable insights, not a chronological log. If an existing summary exists, integrate new information — preserve important older context while adding recent developments. Remove information from conversations that no longer appear in the data above (they may have been deleted by the user). Drop stale details that are no longer relevant.

Current date and time: {now.strftime('%Y-%m-%d %H:%M')}"""


def ensure_memory_summary_schedules():
    """Ensure scheduler entries exist for all agents with memory_summary enabled.
    Called once on server/CLI startup after scheduler is initialized."""
    if not _scheduler:
        return
    agents = list_agents()
    existing = {s["name"]: s for s in _scheduler.list_all()}

    # Remove orphaned schedules for agents that no longer exist
    for sched_name, sched in existing.items():
        if sched_name.startswith("_memory_summary_"):
            orphan_agent = sched_name[len("_memory_summary_"):]
            if orphan_agent not in agents:
                _scheduler.remove(sched_name)

    for agent_id in agents:
        sched_name = _memory_summary_schedule_name(agent_id)
        ms_cfg = _get_memory_summary_config(agent_id)

        if ms_cfg["enabled"] and not ms_cfg.get("paused"):
            schedule_str = _memory_summary_schedule_str(ms_cfg)
            task_prompt = _build_memory_summary_prompt(agent_id)
            # Use Crow-9B for memory summary (local, free); fallback to Sonnet
            primary = ms_cfg.get("model") or "Crow-9B-HERETIC-4.6-MLX-8bit"
            fallback = ms_cfg.get("model_fallback") or "claude-sonnet-4-6"
            summary_model, _ = _resolve_model_with_fallback(primary, fallback, "claude-sonnet-4-6")

            if sched_name in existing:
                old = existing[sched_name]
                if old["schedule"] != schedule_str or old.get("enabled") != 1 or old.get("model") != summary_model:
                    _scheduler.remove(sched_name)
                    _scheduler.add(sched_name, task_prompt, schedule_str,
                                   agent=agent_id, model=summary_model, timeout=600)
            else:
                _scheduler.add(sched_name, task_prompt, schedule_str,
                               agent=agent_id, model=summary_model, timeout=600)
        else:
            # Disabled or paused — remove schedule if exists
            if sched_name in existing:
                _scheduler.remove(sched_name)


def get_memory_summary(agent_id: str) -> str | None:
    """Load the latest memory summary for an agent, if it exists.
    Returns None if memory summary is paused."""
    # Check if paused
    ms_cfg = _get_memory_summary_config(agent_id)
    if ms_cfg.get("paused"):
        return None
    ms = MemoryStore(agent_id)
    # Try to find the memory summary file
    for fname in os.listdir(ms.dir):
        if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
            continue
        fpath = os.path.join(ms.dir, fname)
        try:
            with open(fpath, "r") as f:
                raw = f.read()
            fm, body = _parse_frontmatter(raw)
            if fm.get("name") == "Memory Summary":
                return body.strip()
        except (UnicodeDecodeError, OSError) as e:
            logging.debug(f"Error reading memory file {fname}: {e}")
            continue
        except Exception:
            continue
    return None


def reset_memory_summary(agent_id: str) -> dict:
    """Delete the memory summary file for an agent (Gap 5: reset)."""
    ms = MemoryStore(agent_id)
    deleted = []
    for fname in os.listdir(ms.dir):
        if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
            continue
        fpath = os.path.join(ms.dir, fname)
        try:
            with open(fpath, "r") as f:
                raw = f.read()
            fm, _ = _parse_frontmatter(raw)
            if fm.get("name") == "Memory Summary":
                os.remove(fpath)
                deleted.append(fname)
        except Exception:
            continue
    if deleted:
        _qmd_debounced_embed(agent_id)
    return {"agent": agent_id, "deleted": deleted, "status": "reset"}


def trigger_memory_summary_refresh(agent_id: str):
    """Trigger an immediate memory summary refresh for an agent.
    Runs directly in a background thread instead of waiting for the scheduler loop."""
    if not _scheduler:
        return
    ms_cfg = _get_memory_summary_config(agent_id)
    if not ms_cfg.get("enabled") or ms_cfg.get("paused"):
        return
    sched_name = _memory_summary_schedule_name(agent_id)
    # Check if the schedule exists
    existing = {s["name"]: s for s in _scheduler.list_all()}
    if sched_name not in existing:
        return
    # Execute immediately in a background thread instead of relying on scheduler loop
    task_row = existing[sched_name]
    def _run():
        try:
            _scheduler._execute_scheduled(task_row)
        except Exception as e:
            logging.warning(f"Memory summary refresh failed for {agent_id}: {e}")
    threading.Thread(target=_run, daemon=True, name=f"memory_summary_refresh_{agent_id}").start()


# ─── Relationship Discovery System ────────────────────────────────────

# --- Mechanism 2: Entity Extraction + Auto-linking ---

# Entity extraction regexes — tuned to reduce false positives
# Require at least 2 capitalized words with 3+ chars each (filters "The Quick", "In This")
_RE_CAPITALIZED = re.compile(r'\b[A-Z][a-z]{2,}(?:\s+[A-Z][a-z]{2,})+\b')
_RE_MENTIONS = re.compile(r'@\w+')
_RE_URLS = re.compile(r'https?://[^\s<>"\']+')
_RE_FILEPATHS = re.compile(r'(?:/[\w.-]+){2,}\.\w+')  # require at least 2 path segments
_RE_HASHTAGS = re.compile(r'#\w+')
# Common false positive phrases to filter out
_ENTITY_STOP_PHRASES = frozenset({
    "The Following", "In This", "For Example", "As Well", "In Order",
    "This Case", "Each Time", "At Least", "Make Sure", "Right Now",
    "Based Upon", "Other Than",
})


def _extract_entities(text: str) -> set[str]:
    """Extract entities from text using fast regex heuristics (no LLM).
    Returns set of entity strings."""
    entities = set()
    for m in _RE_CAPITALIZED.finditer(text):
        phrase = m.group()
        if phrase not in _ENTITY_STOP_PHRASES:
            entities.add(phrase)
    for m in _RE_MENTIONS.finditer(text):
        entities.add(m.group())
    for m in _RE_URLS.finditer(text):
        entities.add(m.group().rstrip('.,;)'))
    for m in _RE_FILEPATHS.finditer(text):
        entities.add(m.group())
    for m in _RE_HASHTAGS.finditer(text):
        entities.add(m.group())
    return entities


# Entity index: maps entity string -> set of filenames mentioning it
# Keyed per agent: _entity_indices[agent_id] = {entity: {filename, ...}}
_entity_indices: dict[str, dict[str, set[str]]] = {}
_entity_index_lock = threading.Lock()
_entity_index_initialized: set[str] = set()


def _rebuild_entity_index(agent_id: str):
    """Full rescan of memory files to build entity index for an agent."""
    agent_dir = os.path.join(AGENTS_DIR, agent_id)
    if not os.path.isdir(agent_dir):
        return
    index: dict[str, set[str]] = {}
    # Scan top-level agent dir
    for fname in os.listdir(agent_dir):
        if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
            continue
        # Skip ingested chunks — only agent-created memories
        if fname.startswith("ingest-"):
            continue
        fpath = os.path.join(agent_dir, fname)
        try:
            with open(fpath, "r") as f:
                raw = f.read()
            _, body = _parse_frontmatter(raw)
            entities = _extract_entities(body)
            for ent in entities:
                index.setdefault(ent, set()).add(fname)
        except Exception:
            continue
    # Also scan chats-indexed/ subdirectory for chat transcript entities
    chats_dir = os.path.join(agent_dir, "chats-indexed")
    if os.path.isdir(chats_dir):
        for fname in os.listdir(chats_dir):
            if not fname.endswith(".md"):
                continue
            fpath = os.path.join(chats_dir, fname)
            try:
                with open(fpath, "r") as f:
                    raw = f.read()
                _, body = _parse_frontmatter(raw)
                entities = _extract_entities(body)
                for ent in entities:
                    index.setdefault(ent, set()).add(fname)
            except Exception:
                continue
    with _entity_index_lock:
        _entity_indices[agent_id] = index
        _entity_index_initialized.add(agent_id)


def _ensure_entity_index(agent_id: str):
    """Lazy-initialize entity index on first access. Thread-safe with double-checked locking."""
    # Fast path: already initialized
    if agent_id in _entity_index_initialized:
        return
    # Slow path: acquire lock, check again, rebuild if needed
    with _entity_index_lock:
        if agent_id in _entity_index_initialized:
            return
        # Mark as initialized BEFORE rebuild (under lock) to prevent duplicate work.
        # _rebuild_entity_index stores the result under this same lock.
        _entity_index_initialized.add(agent_id)
    # Rebuild outside lock (I/O heavy) — only one thread gets here per agent_id
    _rebuild_entity_index(agent_id)


def _update_entity_index(agent_id: str, filename: str, entities: set[str]):
    """Incrementally update entity index when a memory is stored."""
    _ensure_entity_index(agent_id)
    with _entity_index_lock:
        idx = _entity_indices.get(agent_id, {})
        # Remove old entries for this filename
        for ent in list(idx.keys()):
            idx[ent].discard(filename)
            if not idx[ent]:
                del idx[ent]
        # Add new entries
        for ent in entities:
            idx.setdefault(ent, set()).add(filename)
        _entity_indices[agent_id] = idx


def _find_entity_matches(agent_id: str, filename: str, entities: set[str]) -> list[str]:
    """Find other files sharing entities with the given file. Returns list of filenames."""
    _ensure_entity_index(agent_id)
    matches = set()
    with _entity_index_lock:
        idx = _entity_indices.get(agent_id, {})
        for ent in entities:
            for other_file in idx.get(ent, set()):
                if other_file != filename:
                    matches.add(other_file)
    return list(matches)


def _add_related_to_file(file_path: str, rel_file: str, rel_type: str) -> bool:
    """Add a related entry to a memory file's frontmatter if not already present.
    Returns True if file was modified."""
    try:
        with open(file_path, "r") as f:
            raw = f.read()
    except Exception:
        return False

    # Check if relationship already exists
    existing_rels = re.findall(r'file:\s*(\S+\.md)', raw)
    if rel_file in existing_rels:
        return False

    # Find the end of frontmatter to insert before the closing ---
    fm_match = re.match(r'^(---\s*\n)(.*?)(\n---\s*\n)(.*)$', raw, re.DOTALL)
    if not fm_match:
        return False

    opener, fm_text, closer, body = fm_match.groups()

    # Check if related: section already exists
    new_entry = f"  - file: {rel_file}\n    type: {rel_type}"
    if "related:" in fm_text:
        # Append to existing related section
        fm_text = fm_text.rstrip() + "\n" + new_entry
    else:
        # Add new related section
        fm_text = fm_text.rstrip() + "\nrelated:\n" + new_entry

    new_raw = opener + fm_text + closer + body
    try:
        with open(file_path, "w") as f:
            f.write(new_raw)
        return True
    except Exception:
        return False


# --- Mechanism 3: Co-recall Linking ---

# Tracks co-occurrence of files recalled together
# Key: frozenset({file1, file2}), Value: count
_recall_cooccurrence: dict[frozenset, int] = {}
_recall_cooccurrence_lock = threading.Lock()
_CO_RECALL_THRESHOLD = 3  # Auto-link after this many co-recalls


def _record_recall_cooccurrence(result_files: list[str], agent_id: str, base_dir: str):
    """Record which files appeared together in a recall result.
    When threshold is reached, auto-add co_recalled relationship."""
    if len(result_files) < 2:
        return
    # Only consider first 10 results to limit combinatorial explosion
    files = result_files[:10]
    pairs_to_link = []
    with _recall_cooccurrence_lock:
        # Cap dict size to prevent unbounded memory growth
        if len(_recall_cooccurrence) > 50_000:
            _recall_cooccurrence.clear()
        for i in range(len(files)):
            for j in range(i + 1, len(files)):
                pair = frozenset({files[i], files[j]})
                _recall_cooccurrence[pair] = _recall_cooccurrence.get(pair, 0) + 1
                if _recall_cooccurrence[pair] == _CO_RECALL_THRESHOLD:
                    pairs_to_link.append((files[i], files[j]))

    # Auto-link pairs that reached threshold (outside lock)
    for f1, f2 in pairs_to_link:
        f1_path = os.path.join(base_dir, f1) if not os.path.isabs(f1) else f1
        f2_path = os.path.join(base_dir, f2) if not os.path.isabs(f2) else f2
        modified = False
        if os.path.exists(f1_path) and os.path.exists(f2_path):
            f1_name = os.path.basename(f1)
            f2_name = os.path.basename(f2)
            if _add_related_to_file(f1_path, f2_name, "co_recalled"):
                modified = True
            if _add_related_to_file(f2_path, f1_name, "co_recalled"):
                modified = True
        if modified:
            _qmd_debounced_embed(agent_id)


# --- Mechanism 1: LLM-based Relationship Discovery ---

RELATIONSHIP_DISCOVERY_DEFAULTS = {
    "enabled": True,
    "frequency": "every 24h",
    "start_time": "04:15",  # 15 min after memory summary default (03:00)
}


def _get_relationship_discovery_config(agent_id: str) -> dict:
    """Read relationship_discovery settings from agent.json, merging with defaults."""
    cfg = AgentConfig(agent_id).config
    rd_cfg = cfg.get("relationship_discovery", {})
    result = dict(RELATIONSHIP_DISCOVERY_DEFAULTS)
    if isinstance(rd_cfg, dict):
        for k in ("enabled", "frequency", "start_time", "model", "model_fallback"):
            if k in rd_cfg:
                result[k] = rd_cfg[k]
    elif isinstance(rd_cfg, bool):
        result["enabled"] = rd_cfg
    return result


AUTODREAM_DEFAULTS = {
    "enabled": True,
    "model": None,  # None = auto (cheapest/Haiku), or explicit model ID
    "stale_threshold_days": 30,
    "dedup_similarity_threshold": 0.85,
    "max_dedup_merges": 10,
    "max_conflict_checks": 30,
    "report_retention": 3,
}


def _get_autodream_config(agent_id: str) -> dict:
    """Read autodream settings from agent.json, merging with defaults."""
    cfg = AgentConfig(agent_id).config
    ad_cfg = cfg.get("autodream", {})
    result = dict(AUTODREAM_DEFAULTS)
    if isinstance(ad_cfg, dict):
        for k in list(AUTODREAM_DEFAULTS) + ["model_fallback"]:
            if k in ad_cfg:
                result[k] = ad_cfg[k]
    elif isinstance(ad_cfg, bool):
        result["enabled"] = ad_cfg
    return result


def _get_auto_memory_config(agent_id: str) -> dict:
    """Read auto_memory settings from agent.json, merging with defaults."""
    cfg = AgentConfig(agent_id).config
    am = cfg.get("auto_memory", {})
    return {
        "enabled": am.get("enabled", True),
        "min_message_length": am.get("min_message_length", 20),
        "model": am.get("model", ""),
        "model_fallback": am.get("model_fallback", ""),
    }


def _auto_memory_extract(agent_id: str, user_message: str, assistant_response: str):
    """Background: check if the conversation exchange contains info worth auto-storing.
    Uses lightweight heuristics first, then optional LLM for borderline cases."""
    # Set thread-local context for this background thread
    _thread_local.current_agent = AgentConfig(agent_id)
    _thread_local.memory_store = MemoryStore(agent_id)
    try:
        _auto_memory_extract_inner(agent_id, user_message, assistant_response)
    finally:
        _thread_local.current_agent = None
        _thread_local.memory_store = None


def _auto_memory_extract_inner(agent_id: str, user_message: str, assistant_response: str):
    """Inner implementation of auto-memory extraction."""

    # Check config
    am_cfg = _get_auto_memory_config(agent_id)
    if not am_cfg.get("enabled", True):
        return

    min_len = am_cfg.get("min_message_length", 20)

    # Skip if messages are too short (small talk, acknowledgments)
    if len(user_message) < min_len or len(assistant_response) < 50:
        return

    # Skip if this looks like a tool-heavy response (already captured by tool actions)
    if assistant_response.count('tool_use') > 3:
        return

    # Heuristic detection of memorable content
    memorable_patterns = []

    # 1. User corrections / feedback
    correction_words = ['don\'t', 'stop', 'no,', 'actually', 'instead', 'wrong', 'not like that',
                        'prefer', 'always', 'never', 'remember that', 'keep in mind']
    user_lower = user_message.lower()
    for word in correction_words:
        if word in user_lower:
            memorable_patterns.append(('feedback', f'User correction/preference detected: "{word}"'))
            break

    # 2. User shares personal/role info
    identity_patterns = ['i am ', 'i\'m a', 'my role', 'i work ', 'my name', 'my team',
                         'my company', 'we use', 'our team', 'our project']
    for pat in identity_patterns:
        if pat in user_lower:
            memorable_patterns.append(('user', f'User identity/role info detected: "{pat}"'))
            break

    # 3. Decisions / commitments
    decision_patterns = ['let\'s go with', 'we decided', 'the plan is', 'going forward',
                         'from now on', 'the approach', 'we\'ll use', 'agreed']
    for pat in decision_patterns:
        if pat in user_lower:
            memorable_patterns.append(('project', f'Decision detected: "{pat}"'))
            break

    # 4. References to external resources
    if any(p in user_lower for p in ['http://', 'https://', 'the repo', 'the doc', 'the wiki',
                                      'slack channel', 'linear', 'jira', 'confluence']):
        memorable_patterns.append(('reference', 'External resource reference detected'))

    if not memorable_patterns:
        return  # Nothing worth storing

    # Use a quick LLM call to extract and format the memory
    mem_type, trigger = memorable_patterns[0]

    # Build a focused extraction prompt
    prompt = (
        f"Extract a concise, actionable memory from this conversation exchange.\n\n"
        f"USER: {user_message[:500]}\n\n"
        f"A: {assistant_response[:500]}\n\n"
        f"Trigger: {trigger}\n"
        f"Memory type: {mem_type}\n\n"
        f"Rules:\n"
        f'- Output ONLY a JSON object: {{"name": "short-title", "content": "the key fact/preference/decision", "type": "{mem_type}", "description": "one-line summary"}}\n'
        f"- Content should be 1-3 sentences, factual, no fluff\n"
        f'- If the exchange doesn\'t actually contain memorable info, output: {{"skip": true}}\n'
        f"- Do NOT output anything except the JSON object"
    )

    try:
        # Use configured model, or find cheapest
        model = None
        am_model = am_cfg.get("model", "")
        am_fallback = am_cfg.get("model_fallback", "")
        if am_model:
            model, _ = _resolve_model_with_fallback(am_model, am_fallback, "claude-haiku-4-5-20251001")
        if not model and _models_config:
            for mid, cfg in sorted(_models_config.items(), key=lambda x: x[1].get('cost_input', 999)):
                if cfg.get('enabled', True):
                    ml = mid.lower()
                    if 'haiku' in ml:
                        model = mid
                        break
            if not model:
                for mid, cfg in sorted(_models_config.items(), key=lambda x: x[1].get('priority', 0)):
                    if cfg.get('enabled', True):
                        model = mid
                        break

        if not model or not _delegate_api_key:
            return

        ms = MemoryStore(agent_id)
        result = _run_delegate(
            messages=[{"role": "user", "content": prompt}],
            model=model,
            system_prompt="You are a memory extraction assistant. Output only valid JSON.",
            memory_store=ms,
            inference_params={"max_tokens": 256, "temperature": 0.1},
            tools=False,
        )

        if not result or 'skip' in result.lower():
            return

        # Parse JSON
        data = _extract_json_from_llm(result)
        if not data or data.get('skip'):
            return

        name = data.get('name', '').strip()
        content = data.get('content', '').strip()
        mem_type_extracted = data.get('type', mem_type)
        description = data.get('description', '').strip()

        if not name or not content:
            return

        # Check if similar memory already exists — skip if content is substantially the same
        existing = ms.recall(name, limit=3)
        for e in existing:
            if e.get('name', '').lower() == name.lower():
                # Same name exists — update only if new content is meaningfully different
                old_content = e.get('content', '')
                if old_content and content.strip().lower() == old_content.strip().lower():
                    return  # Identical content, skip
                # Content differs — allow the store to overwrite
                break

        # Store it
        ms.store(name, content, description, mem_type_extracted)
        logging.info(f"Auto-memory stored for {agent_id}: {name} ({mem_type_extracted})")

    except (json.JSONDecodeError, ValueError, KeyError) as e:
        logging.debug(f"Auto-memory extraction parse error for {agent_id}: {e}")
    except Exception as e:
        logging.warning(f"Auto-memory extraction failed for {agent_id}: {e}")


def _collect_agent_memories(agent_id: str, max_content: int = 2000) -> list[dict]:
    """Collect all memory files for an agent, walking subdirectories.
    Content is truncated at max_content chars — callers never need more than 2000."""
    agent_dir = os.path.join(AGENTS_DIR, agent_id)
    if not os.path.isdir(agent_dir):
        return []
    memories = []
    ignore_dirs = {"skills", ".trash", "ingested"}
    for dirpath, dirnames, filenames in os.walk(agent_dir):
        dirnames[:] = [d for d in dirnames if d not in ignore_dirs and not d.startswith(".")]
        for fname in sorted(filenames):
            if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
                continue
            if fname.startswith("ingest-"):
                continue
            fpath = os.path.join(dirpath, fname)
            rel = os.path.relpath(fpath, agent_dir)
            try:
                with open(fpath, "r") as f:
                    raw = f.read(max_content * 4)  # read only what we need (frontmatter + content)
                fm, body = _parse_frontmatter(raw)
                content = body.strip()
                memories.append({
                    "file": rel,
                    "name": fm.get("name", fname.replace(".md", "")),
                    "type": fm.get("type", "general"),
                    "description": fm.get("description", "").strip('"').strip("'"),
                    "content": content[:max_content],
                })
            except Exception:
                continue
    return memories


def _find_candidate_pairs(agent_id: str, memories: list[dict], max_candidates: int = 40) -> list[tuple[dict, dict, float]]:
    """Use QMD semantic search to find candidate pairs that might be related.
    For each memory, query QMD with its content and find similar files.
    Returns list of (mem_a, mem_b, score) tuples, deduplicated."""
    ms = MemoryStore(agent_id)
    seen_pairs = set()
    candidates = []

    for mem in memories:
        # Build a clean plaintext query for QMD
        def _clean(s):
            s = re.sub(r'["\'\(\)\[\]#*_|]', '', s)
            s = re.sub(r'\s*part\s+\d+/\d+\s*', '', s)
            s = s.replace('\n', ' ').replace('\r', ' ')
            return re.sub(r'\s+', ' ', s).strip()
        query = _clean(f"{mem['name']} {mem['description']} {mem['content'][:200]}")[:300]
        if not query:
            continue
        try:
            results = ms.recall(query, limit=5)
        except Exception:
            continue
        for r in results:
            other_file = os.path.relpath(r.get("file_path", ""), os.path.join(AGENTS_DIR, agent_id))
            if other_file == mem["file"]:
                continue
            # Find the matching memory dict
            other_mem = next((m for m in memories if m["file"] == other_file), None)
            if not other_mem:
                continue
            pair_key = tuple(sorted([mem["file"], other_mem["file"]]))
            if pair_key in seen_pairs:
                continue
            seen_pairs.add(pair_key)
            candidates.append((mem, other_mem, r.get("score", 0)))

    # Sort by score descending, take top candidates
    candidates.sort(key=lambda x: x[2], reverse=True)
    return candidates[:max_candidates]


def _build_relationship_discovery_prompt(agent_id: str) -> str:
    """Build prompt for LLM-based relationship discovery.
    Two-stage approach:
    1. QMD semantic search finds candidate pairs (scales to any number of files)
    2. LLM reads full content of candidates and classifies relationships
    """
    memories = _collect_agent_memories(agent_id)
    if len(memories) < 2:
        return ""

    # Stage 1: find candidate pairs via QMD embeddings
    candidates = _find_candidate_pairs(agent_id, memories)

    if not candidates:
        # Fallback: if QMD unavailable, send all files with truncated content
        listing = "\n\n".join(
            f"### {m['file']} (type: {m['type']})\n{m['content'][:500]}"
            for m in memories[:50]
        )
        return f"""Analyze these memories for agent '{agent_id}' and identify relationships.

{listing}

Output ONLY a JSON array. Each object: from_file, to_file, type (references|same_topic|depends_on|contradicts|extends), reason (under 20 words).
If no relationships, output []. Be selective — meaningful relationships only."""

    # Stage 2: build prompt with truncated content of candidate pairs
    _PAIR_CONTENT_LIMIT = 2000
    pair_sections = []
    for i, (a, b, score) in enumerate(candidates, 1):
        a_content = a['content'][:_PAIR_CONTENT_LIMIT] + ("…" if len(a['content']) > _PAIR_CONTENT_LIMIT else "")
        b_content = b['content'][:_PAIR_CONTENT_LIMIT] + ("…" if len(b['content']) > _PAIR_CONTENT_LIMIT else "")
        pair_sections.append(
            f"## Candidate Pair {i} (similarity: {score:.2f})\n\n"
            f"### File A: {a['file']} (type: {a['type']})\n{a_content}\n\n"
            f"### File B: {b['file']} (type: {b['type']})\n{b_content}"
        )

    pairs_text = "\n\n---\n\n".join(pair_sections)

    return f"""You are analyzing candidate pairs of memories for agent '{agent_id}'.
Each pair was identified as potentially related by semantic similarity.
Read the FULL content of each file and determine if a meaningful relationship exists.

{pairs_text}

## Task

For each pair where a meaningful relationship exists, output:
- from_file: filename of file A
- to_file: filename of file B
- type: one of (references, same_topic, depends_on, contradicts, extends)
- reason: brief explanation (under 20 words)

Skip pairs that are only superficially similar. Output ONLY a JSON array.
If no meaningful relationships, output [].

Example:
[{{"from_file": "notes/Meeting.md", "to_file": "chats-indexed/chat-abc-000.md", "type": "same_topic", "reason": "Both discuss macOS architecture on Apple Silicon"}}]"""


def _apply_discovered_relationships(agent_id: str, relationships: list[dict]):
    """Apply discovered relationships to memory files.
    Updates frontmatter of both files for each relationship."""
    if not relationships:
        return
    agent_dir = os.path.join(AGENTS_DIR, agent_id)
    modified = False
    for rel in relationships:
        from_file = rel.get("from_file", "")
        to_file = rel.get("to_file", "")
        rel_type = rel.get("type", "references")
        # Validate type
        if rel_type not in ("references", "same_topic", "depends_on", "contradicts", "extends"):
            rel_type = "references"
        if not from_file or not to_file:
            continue
        from_path = os.path.join(agent_dir, from_file)
        to_path = os.path.join(agent_dir, to_file)
        if not os.path.exists(from_path) or not os.path.exists(to_path):
            continue
        if _add_related_to_file(from_path, to_file, rel_type):
            modified = True
        # Add reverse link (bidirectional)
        reverse_type = rel_type
        if rel_type == "depends_on":
            reverse_type = "extends"  # reverse of depends_on
        elif rel_type == "extends":
            reverse_type = "depends_on"
        if _add_related_to_file(to_path, from_file, reverse_type):
            modified = True
    if modified:
        _qmd_debounced_embed(agent_id)


def trigger_relationship_discovery(agent_id: str):
    """Run LLM-based relationship discovery immediately in a background thread."""
    if not _delegate_api_key:
        logging.warning("Relationship discovery skipped: delegate API not configured")
        return
    prompt = _build_relationship_discovery_prompt(agent_id)
    if not prompt:
        return

    def _run():
        try:
            # Use Crow-4B for relationship discovery — JSON classification; fallback to Haiku
            target = AgentConfig(agent_id)
            rd_cfg = _get_relationship_discovery_config(agent_id)
            primary = rd_cfg.get("model") or "mistral-small"
            fallback = rd_cfg.get("model_fallback") or "claude-haiku-4-5-20251001"
            model, fallback_model = _resolve_model_with_fallback(primary, fallback, "claude-haiku-4-5-20251001")
            ms = MemoryStore(agent_id)
            logging.info(f"Relationship discovery starting for {agent_id} using {model} (fallback: {fallback_model})")
            result_text = _run_delegate_with_fallback(
                messages=[{"role": "user", "content": prompt}],
                primary_model=model, fallback_model=fallback_model,
                system_prompt="You are a relationship analysis assistant. Analyze the given memories and output ONLY a valid JSON array of relationships. No explanations, no markdown, just the JSON array.",
                memory_store=ms,
                inference_params={"max_tokens": 16384, "temperature": 0.2},
                tools=False,
            )
            if not result_text:
                logging.warning(f"Relationship discovery for {agent_id}: empty response")
                return
            if result_text.startswith("Delegation error:"):
                logging.warning(f"Relationship discovery for {agent_id}: {result_text}")
                return
            # Extract JSON from response (may have markdown code fences)
            relationships = _extract_json_from_llm(result_text, expect_array=True)
            if isinstance(relationships, list) and relationships:
                _apply_discovered_relationships(agent_id, relationships)
                logging.info(f"Relationship discovery for {agent_id}: found {len(relationships)} relationships")
            elif relationships is not None:
                logging.info(f"Relationship discovery for {agent_id}: no relationships found")
            else:
                logging.warning(f"Relationship discovery for {agent_id}: no JSON in response: {result_text[:200]}")
        except Exception as e:
            logging.exception(f"Relationship discovery failed for {agent_id}: {e}")

    threading.Thread(target=_run, daemon=True, name=f"rel_discovery_{agent_id}").start()


def _relationship_discovery_schedule_name(agent_id: str) -> str:
    return f"_relationship_discovery_{agent_id}"


def _relationship_discovery_schedule_str(cfg: dict) -> str:
    """Convert relationship_discovery config to a scheduler schedule string."""
    freq = cfg.get("frequency", "every 24h").strip().lower()
    start_time = cfg.get("start_time", "").strip()
    if start_time and re.match(r'^\d{1,2}:\d{2}$', start_time):
        m = re.match(r'every\s+(\d+)\s*(h|hour|d|day)', freq)
        if m:
            val, unit = int(m.group(1)), m.group(2)[0]
            if (unit == 'h' and val == 24) or (unit == 'd' and val == 1):
                return f"daily {start_time}"
            elif unit == 'h' and val < 24:
                return freq
            elif unit == 'd' and val > 1:
                return freq
    return freq


def _build_relationship_discovery_task_prompt(agent_id: str) -> str:
    """Build the scheduled task prompt for relationship discovery."""
    return f"""Perform relationship discovery for agent '{agent_id}'.

Use memory_recall with an empty query to list all memories, then analyze them for relationships.

For each pair of related memories, use memory_store to update the related frontmatter.
Focus on meaningful relationships: references, same_topic, depends_on, contradicts, extends.

Be selective — only report meaningful relationships, not superficial ones."""


def ensure_relationship_discovery_schedules():
    """Ensure scheduler entries exist for all agents with relationship_discovery enabled."""
    if not _scheduler:
        return
    agents = list_agents()
    existing = {s["name"]: s for s in _scheduler.list_all()}

    # Remove orphaned schedules for agents that no longer exist
    for sched_name in list(existing):
        if sched_name.startswith("_relationship_discovery_"):
            orphan_agent = sched_name[len("_relationship_discovery_"):]
            if orphan_agent not in agents:
                _scheduler.remove(sched_name)

    for agent_id in agents:
        sched_name = _relationship_discovery_schedule_name(agent_id)
        rd_cfg = _get_relationship_discovery_config(agent_id)

        if rd_cfg["enabled"]:
            schedule_str = _relationship_discovery_schedule_str(rd_cfg)
            task_prompt = _build_relationship_discovery_task_prompt(agent_id)
            # Use Crow-4B for RD scheduled task; fallback Haiku
            primary = rd_cfg.get("model") or "mistral-small"
            fallback = rd_cfg.get("model_fallback") or "claude-haiku-4-5-20251001"
            rd_model, _ = _resolve_model_with_fallback(primary, fallback, "claude-haiku-4-5-20251001")

            if sched_name in existing:
                old = existing[sched_name]
                if old["schedule"] != schedule_str or old.get("enabled") != 1 or old.get("model") != rd_model:
                    _scheduler.remove(sched_name)
                    _scheduler.add(sched_name, task_prompt, schedule_str,
                                   agent=agent_id, model=rd_model, timeout=600)
            else:
                _scheduler.add(sched_name, task_prompt, schedule_str,
                               agent=agent_id, model=rd_model, timeout=600)
        else:
            if sched_name in existing:
                _scheduler.remove(sched_name)


# ─── Autodream: Memory Consolidation System ─────────────────────────────


def _autodream_dedup(agent_id: str, ms: MemoryStore, config: dict, memories: list[dict] | None = None) -> dict:
    """Scan memories for semantic duplicates via QMD, merge with LLM."""
    if memories is None:
        memories = _collect_agent_memories(agent_id)
    if len(memories) < 2:
        return {"duplicates_found": 0, "merged": 0, "skipped": 0, "merge_log": []}

    threshold = config.get("dedup_similarity_threshold", 0.85)
    max_merges = config.get("max_dedup_merges", 10)

    # Find candidate pairs via QMD semantic search
    try:
        pairs = _find_candidate_pairs(agent_id, memories, max_candidates=60)
    except Exception:
        return {"duplicates_found": 0, "merged": 0, "skipped": 0, "merge_log": [],
                "error": "Failed to find candidate pairs"}

    # Filter to high-similarity pairs
    high_sim = [(a, b, score) for a, b, score in pairs if score >= threshold]

    duplicates_found = 0
    merged = 0
    skipped = 0
    merge_log = []
    deleted_files = set()

    # Resolve model (configured primary + fallback)
    model, fallback_model = _resolve_autodream_model(config)
    if not model or not _delegate_api_key:
        return {"duplicates_found": len(high_sim), "merged": 0, "skipped": len(high_sim),
                "merge_log": ["No model available for merge confirmation"]}

    for mem_a, mem_b, score in high_sim:
        if mem_a["file"] in deleted_files or mem_b["file"] in deleted_files:
            continue

        duplicates_found += 1

        if merged >= max_merges:
            skipped += 1
            continue
        prompt = (
            f"Are these two memories duplicates or near-duplicates that should be merged?\n\n"
            f"MEMORY A ({mem_a['name']}):\n{mem_a['content'][:800]}\n\n"
            f"MEMORY B ({mem_b['name']}):\n{mem_b['content'][:800]}\n\n"
            f"If they are duplicates, produce a merged version that preserves all unique information.\n"
            f'Output ONLY JSON: {{"is_duplicate": true/false, "merged_name": "...", "merged_content": "...", '
            f'"merged_description": "..."}}\n'
            f'If not duplicates: {{"is_duplicate": false}}'
        )
        try:
            result = _run_delegate_with_fallback(
                messages=[{"role": "user", "content": prompt}],
                primary_model=model, fallback_model=fallback_model,
                system_prompt="You are a memory deduplication assistant. Output only valid JSON.",
                inference_params={"max_tokens": 1024, "temperature": 0.1},
                tools=False,
            )
            if not result:
                skipped += 1
                continue
            data = _extract_json_from_llm(result)
            if not data or not data.get("is_duplicate"):
                skipped += 1
                continue

            # Determine which file is older (keep newer filename)
            agent_dir = os.path.join(AGENTS_DIR, agent_id)
            path_a = os.path.join(agent_dir, mem_a["file"])
            path_b = os.path.join(agent_dir, mem_b["file"])
            mtime_a = os.path.getmtime(path_a) if os.path.exists(path_a) else 0
            mtime_b = os.path.getmtime(path_b) if os.path.exists(path_b) else 0

            # Store merged content under the newer file's name
            merged_name = data.get("merged_name", mem_a["name"] if mtime_a >= mtime_b else mem_b["name"])
            ms.store(merged_name, data.get("merged_content", ""),
                     data.get("merged_description", ""), mem_a.get("type", "general"))

            # Delete the older file
            older = mem_b if mtime_a >= mtime_b else mem_a
            ms.delete(older["name"])
            deleted_files.add(older["file"])
            merged += 1
            merge_log.append(f"Merged '{mem_a['name']}' + '{mem_b['name']}' → '{merged_name}'")
        except Exception as e:
            skipped += 1
            merge_log.append(f"Error merging '{mem_a['name']}' + '{mem_b['name']}': {str(e)[:100]}")

    return {"duplicates_found": duplicates_found, "merged": merged,
            "skipped": skipped, "merge_log": merge_log}


def _autodream_staleness(agent_id: str, ms: MemoryStore, config: dict) -> dict:
    """Flag memories not recalled in N days as stale. No LLM calls."""
    threshold_days = config.get("stale_threshold_days", 30)
    now = datetime.datetime.now()
    agent_dir = os.path.join(AGENTS_DIR, agent_id)

    total = 0
    stale_count = 0
    newly_stale = 0
    stale_names = []
    # Skip system files (Memory Summary, Memory Health Report)
    skip_prefixes = ("Memory Summary", "Memory Health Report")

    for fname in os.listdir(agent_dir):
        if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
            continue
        if fname.startswith("ingest-"):
            continue
        fpath = os.path.join(agent_dir, fname)
        try:
            with open(fpath, "r") as f:
                raw = f.read()
            fm, body = _parse_frontmatter(raw)
        except Exception:
            continue

        name = fm.get("name", fname.replace(".md", ""))
        if any(name.startswith(p) for p in skip_prefixes):
            continue
        total += 1

        # Determine last activity date
        last_recalled_str = fm.get("last_recalled", "")
        if last_recalled_str:
            try:
                last_date = datetime.datetime.strptime(last_recalled_str.strip(), "%Y-%m-%d")
            except ValueError:
                last_date = datetime.datetime.fromtimestamp(os.path.getmtime(fpath))
        else:
            last_date = datetime.datetime.fromtimestamp(os.path.getmtime(fpath))

        days_since = (now - last_date).days
        already_stale = fm.get("stale", "").lower() == "true"

        if days_since > threshold_days:
            stale_count += 1
            stale_names.append(name)
            if not already_stale:
                # Stamp stale: true in frontmatter
                try:
                    fm_match = re.match(r'^(---\s*\n)(.*?)(\n---\s*\n)(.*)$', raw, re.DOTALL)
                    if fm_match:
                        opener, fm_text, closer, body_text = fm_match.groups()
                        if "stale:" not in fm_text:
                            fm_text = fm_text.rstrip() + "\nstale: true\n"
                        with open(fpath, "w") as f:
                            f.write(opener + fm_text + closer + body_text)
                        newly_stale += 1
                except Exception as e:
                    logging.debug(f"Failed to mark stale: {fpath}: {e}")
        elif already_stale:
            # Was stale but now recalled within threshold — remove stale flag
            try:
                fm_match = re.match(r'^(---\s*\n)(.*?)(\n---\s*\n)(.*)$', raw, re.DOTALL)
                if fm_match:
                    opener, fm_text, closer, body_text = fm_match.groups()
                    fm_text = re.sub(r'\nstale:.*', '', fm_text)
                    with open(fpath, "w") as f:
                        f.write(opener + fm_text + closer + body_text)
            except Exception:
                pass

    return {"total": total, "stale_count": stale_count, "stale_names": stale_names,
            "newly_stale": newly_stale}


def _autodream_conflicts(agent_id: str, ms: MemoryStore, config: dict, memories: list[dict] | None = None) -> dict:
    """Find contradicting memories using QMD similarity + LLM classification."""
    if memories is None:
        memories = _collect_agent_memories(agent_id)
    if len(memories) < 2:
        return {"conflicts_found": 0, "conflict_pairs": []}

    max_checks = config.get("max_conflict_checks", 30)

    try:
        pairs = _find_candidate_pairs(agent_id, memories, max_candidates=max_checks * 2)
    except Exception:
        return {"conflicts_found": 0, "conflict_pairs": [], "error": "Failed to find candidate pairs"}

    # Filter to same-type pairs (contradictions are most likely within same type)
    same_type_pairs = [(a, b, s) for a, b, s in pairs if a["type"] == b["type"]][:max_checks]

    model, fallback_model = _resolve_autodream_model(config)
    if not model or not _delegate_api_key:
        return {"conflicts_found": 0, "conflict_pairs": [],
                "error": "No model available for conflict detection"}

    conflicts_found = 0
    conflict_pairs = []

    for mem_a, mem_b, score in same_type_pairs:
        prompt = (
            f"Do these two memories contradict each other? Look for conflicting facts, "
            f"opposing instructions, or incompatible information.\n\n"
            f"MEMORY A ({mem_a['name']}, type={mem_a['type']}):\n{mem_a['content'][:600]}\n\n"
            f"MEMORY B ({mem_b['name']}, type={mem_b['type']}):\n{mem_b['content'][:600]}\n\n"
            f'Output ONLY JSON: {{"contradicts": true/false, "nature": "brief description of conflict"}}\n'
            f'If no contradiction: {{"contradicts": false}}'
        )
        try:
            result = _run_delegate_with_fallback(
                messages=[{"role": "user", "content": prompt}],
                primary_model=model, fallback_model=fallback_model,
                system_prompt="You are a memory conflict detector. Output only valid JSON.",
                inference_params={"max_tokens": 256, "temperature": 0.1},
                tools=False,
            )
            if not result:
                continue
            data = _extract_json_from_llm(result)
            if not data:
                continue
            if data.get("contradicts"):
                conflicts_found += 1
                conflict_pairs.append({
                    "a": mem_a["name"], "b": mem_b["name"],
                    "nature": data.get("nature", "Unknown conflict"),
                })
        except Exception:
            continue

    return {"conflicts_found": conflicts_found, "conflict_pairs": conflict_pairs}


def _autodream_skill_candidates(agent_id: str, ms: MemoryStore, config: dict, memories: list[dict] | None = None) -> dict:
    """Identify procedural memories that could become skills."""
    if memories is None:
        memories = _collect_agent_memories(agent_id)
    # Filter to feedback/general types
    candidates_pool = [m for m in memories if m["type"] in ("feedback", "general")]

    # Heuristic pre-filter: procedural patterns
    procedural_patterns = [
        r'\b\d+\.\s', r'always\b', r'never\b', r'when\s.*\bdo\b', r'steps?:',
        r'how to\b', r'make sure\b', r'first\b.*then\b', r'rule:',
    ]
    heuristic_matches = []
    for mem in candidates_pool:
        content_lower = mem["content"].lower()
        if any(re.search(p, content_lower) for p in procedural_patterns):
            heuristic_matches.append(mem)
    heuristic_matches = heuristic_matches[:15]  # Cap

    if not heuristic_matches:
        return {"candidates_found": 0, "candidate_names": [], "candidate_details": []}

    model, fallback_model = _resolve_autodream_model(config)
    if not model or not _delegate_api_key:
        return {"candidates_found": 0, "candidate_names": [], "candidate_details": [],
                "error": "No model available for skill detection"}

    candidates_found = 0
    candidate_names = []
    candidate_details = []

    for mem in heuristic_matches:
        prompt = (
            f"Is this memory a reusable procedure or workflow that should become a skill?\n"
            f"A skill is a repeatable process with clear steps that an AI agent would benefit from having as a template.\n\n"
            f"Memory ({mem['name']}, type={mem['type']}):\n{mem['content'][:800]}\n\n"
            f'Output ONLY JSON: {{"is_skill_candidate": true/false, "reason": "brief explanation"}}\n'
            f'If not a skill candidate: {{"is_skill_candidate": false}}'
        )
        try:
            result = _run_delegate_with_fallback(
                messages=[{"role": "user", "content": prompt}],
                primary_model=model, fallback_model=fallback_model,
                system_prompt="You are a skill detection assistant. Output only valid JSON.",
                inference_params={"max_tokens": 256, "temperature": 0.1},
                tools=False,
            )
            if not result:
                continue
            data = _extract_json_from_llm(result)
            if not data:
                continue
            if data.get("is_skill_candidate"):
                candidates_found += 1
                candidate_names.append(mem["name"])
                candidate_details.append({
                    "name": mem["name"],
                    "reason": data.get("reason", "Procedural content detected"),
                })
        except Exception:
            continue

    return {"candidates_found": candidates_found, "candidate_names": candidate_names,
            "candidate_details": candidate_details}


def _resolve_model_with_fallback(primary: str | None, fallback: str | None, hardcoded_default: str) -> tuple[str, str | None]:
    """Return (primary_to_use, fallback_to_use) after validating against enabled models.
    Falls back through the chain: configured primary → configured fallback → hardcoded default."""
    def _is_available(mid: str) -> bool:
        if not mid:
            return False
        if not _models_config:
            return True  # can't check, assume available
        mcfg = _models_config.get(mid, {})
        return mcfg.get("enabled", True)

    resolved_primary = primary if _is_available(primary) else None
    resolved_fallback = fallback if _is_available(fallback) else None

    if resolved_primary:
        return resolved_primary, resolved_fallback or hardcoded_default
    if resolved_fallback:
        return resolved_fallback, hardcoded_default
    return hardcoded_default, None


def _run_delegate_with_fallback(messages, primary_model, fallback_model, system_prompt,
                                 memory_store=None, inference_params=None,
                                 tools: bool = True) -> str | None:
    """Run _run_delegate with automatic fallback if primary model fails."""
    result = _run_delegate(
        messages=messages, model=primary_model, system_prompt=system_prompt,
        memory_store=memory_store, inference_params=inference_params, tools=tools,
    )
    if result and result.startswith("Delegation error:") and fallback_model and fallback_model != primary_model:
        logging.warning(f"Primary model {primary_model} failed, falling back to {fallback_model}")
        result = _run_delegate(
            messages=messages, model=fallback_model, system_prompt=system_prompt,
            memory_store=memory_store, inference_params=inference_params, tools=tools,
        )
    return result


def _resolve_autodream_model(config: dict) -> tuple[str, str | None]:
    """Resolve (primary, fallback) models for autodream."""
    primary = config.get("model") or "Crow-4B-Opus-4.6-Distill"
    fallback = config.get("model_fallback") or "claude-haiku-4-5-20251001"
    return _resolve_model_with_fallback(primary, fallback, "claude-haiku-4-5-20251001")


def _find_cheapest_model() -> str | None:
    """Find the cheapest available model (prefer local, then Haiku)."""
    if not _models_config:
        return None
    # Prefer local models first (priority <= 30)
    for mid, cfg in sorted(_models_config.items(), key=lambda x: x[1].get('priority', 50)):
        if cfg.get('enabled', True) and cfg.get('priority', 50) <= 30:
            return mid
    # Then Haiku
    for mid, cfg in _models_config.items():
        if cfg.get('enabled', True) and 'haiku' in mid.lower():
            return mid
    # Then lowest priority
    for mid, cfg in sorted(_models_config.items(), key=lambda x: x[1].get('priority', 0)):
        if cfg.get('enabled', True):
            return mid
    return None


def promote_memory_to_skill(agent_id: str, memory_name: str) -> dict:
    """Convert a memory into a skill by generating a SKILL.md via LLM."""
    ms = MemoryStore(agent_id)
    # Find the memory
    memories = _collect_agent_memories(agent_id)
    mem = next((m for m in memories if m["name"] == memory_name), None)
    if not mem:
        return {"error": f"Memory '{memory_name}' not found"}

    ad_cfg = _get_autodream_config(agent_id)
    model = _resolve_autodream_model(ad_cfg)
    if not model or not _delegate_api_key:
        return {"error": "No model available for skill generation"}

    # Generate SKILL.md content via LLM
    prompt = (
        f"Convert this memory into a well-structured SKILL.md file.\n\n"
        f"Memory name: {mem['name']}\n"
        f"Memory content:\n{mem['content']}\n\n"
        f"Generate a skill document with:\n"
        f"1. A clear, concise title\n"
        f"2. Step-by-step instructions or reference material\n"
        f"3. Examples where helpful\n"
        f"4. Keep it practical and actionable\n\n"
        f'Output ONLY JSON: {{"slug": "short-kebab-case-name", "name": "Display Name", '
        f'"description": "one-line description", "body": "the full skill body in markdown"}}'
    )
    try:
        result = _run_delegate(
            messages=[{"role": "user", "content": prompt}],
            model=model,
            system_prompt="You are a skill creation assistant. Output only valid JSON.",
            inference_params={"max_tokens": 2048, "temperature": 0.2},
            tools=False,
        )
        if not result:
            return {"error": "LLM returned empty response"}

        data = _extract_json_from_llm(result)
        if not data:
            return {"error": "Could not parse LLM response"}
        slug = re.sub(r'[^a-z0-9-]', '', data.get("slug", "").lower().replace(" ", "-"))[:40]
        if not slug:
            slug = re.sub(r'[^a-z0-9-]', '', memory_name.lower().replace(" ", "-"))[:40]
        name = data.get("name", memory_name)
        description = data.get("description", "")
        body = data.get("body", mem["content"])

        # Write SKILL.md
        skills_dir = os.path.join(AGENTS_DIR, agent_id, "skills", slug)
        os.makedirs(skills_dir, exist_ok=True)
        skill_path = os.path.join(skills_dir, "SKILL.md")

        skill_content = f"""---
name: {_yaml_escape(name)}
description: {_yaml_escape(description)}
---

{body}
"""
        with open(skill_path, "w") as f:
            f.write(skill_content)

        logging.info(f"Promoted memory '{memory_name}' to skill '{slug}' for {agent_id}")
        return {"status": "created", "slug": slug, "name": name, "path": skill_path}

    except Exception as e:
        return {"error": f"Skill generation failed: {str(e)[:200]}"}


def _build_autodream_report(agent_id: str, results: dict) -> str:
    """Compile autodream results into a structured markdown report."""
    now = datetime.datetime.now()
    dedup = results.get("dedup", {})
    stale = results.get("staleness", {})
    conflicts = results.get("conflicts", {})
    skills = results.get("skill_candidates", {})

    # Calculate health score
    total = stale.get("total", 1) or 1
    stale_pct = (stale.get("stale_count", 0) / total) * 100
    health_score = max(0, min(100, int(
        100
        - (stale_pct * 0.3)
        - (conflicts.get("conflicts_found", 0) * 5)
        - (max(0, dedup.get("duplicates_found", 0) - dedup.get("merged", 0)) * 2)
    )))
    results["health_score"] = health_score

    report = f"""# Memory Health Report — {now.strftime('%Y-%m-%d %H:%M')}

**Health Score: {health_score}/100**

## Deduplication
- Duplicates found: {dedup.get('duplicates_found', 0)}
- Merged: {dedup.get('merged', 0)}
- Skipped: {dedup.get('skipped', 0)}
"""
    for entry in dedup.get("merge_log", []):
        report += f"- {entry}\n"

    report += f"""
## Staleness
- Total memories: {stale.get('total', 0)}
- Stale (>{results.get('config', {}).get('stale_threshold_days', 30)}d): {stale.get('stale_count', 0)} ({stale_pct:.0f}%)
- Newly flagged: {stale.get('newly_stale', 0)}
"""
    for name in stale.get("stale_names", [])[:20]:
        report += f"- {name}\n"

    report += f"""
## Conflicts
- Conflicts detected: {conflicts.get('conflicts_found', 0)}
"""
    for pair in conflicts.get("conflict_pairs", []):
        report += f"- **{pair['a']}** ↔ **{pair['b']}**: {pair['nature']}\n"

    report += f"""
## Skill Candidates
- Candidates found: {skills.get('candidates_found', 0)}
"""
    for c in skills.get("candidate_details", []):
        report += f"- **{c['name']}**: {c['reason']}\n"

    return report


# Live autodream status per agent: {agent_id: {status, phase, started_at, finished_at, error, results}}
_autodream_status: dict[str, dict] = {}
_autodream_status_lock = threading.Lock()


def get_autodream_status(agent_id: str) -> dict | None:
    """Get live autodream status for an agent."""
    with _autodream_status_lock:
        return dict(_autodream_status.get(agent_id, {})) or None


def trigger_autodream(agent_id: str):
    """Run all autodream passes and produce a health report. Background thread."""
    ad_cfg = _get_autodream_config(agent_id)
    if not ad_cfg.get("enabled"):
        return

    def _update_status(**kwargs):
        with _autodream_status_lock:
            if agent_id not in _autodream_status:
                _autodream_status[agent_id] = {}
            _autodream_status[agent_id].update(kwargs)

    def _run():
        started = datetime.datetime.now()
        _update_status(status="running", phase="starting", started_at=started.isoformat(),
                       finished_at=None, error=None, results=None)
        try:
            logging.info(f"Autodream starting for {agent_id}")
            ms = MemoryStore(agent_id)
            results = {"config": ad_cfg, "agent": agent_id}

            # Collect memories once — shared across all passes to avoid repeated filesystem walks
            memories = _collect_agent_memories(agent_id)

            # Run passes sequentially
            _update_status(phase="dedup")
            results["dedup"] = _autodream_dedup(agent_id, ms, ad_cfg, memories=memories)
            logging.info(f"Autodream dedup done for {agent_id}: {results['dedup'].get('merged', 0)} merged")

            _update_status(phase="staleness")
            results["staleness"] = _autodream_staleness(agent_id, ms, ad_cfg)
            logging.info(f"Autodream staleness done for {agent_id}: {results['staleness'].get('stale_count', 0)} stale")

            _update_status(phase="conflicts")
            results["conflicts"] = _autodream_conflicts(agent_id, ms, ad_cfg, memories=memories)
            logging.info(f"Autodream conflicts done for {agent_id}: {results['conflicts'].get('conflicts_found', 0)} found")

            _update_status(phase="skill_candidates")
            results["skill_candidates"] = _autodream_skill_candidates(agent_id, ms, ad_cfg, memories=memories)
            logging.info(f"Autodream skills done for {agent_id}: {results['skill_candidates'].get('candidates_found', 0)} candidates")

            # Build and store report
            _update_status(phase="report")
            report = _build_autodream_report(agent_id, results)
            now_str = datetime.datetime.now().strftime('%Y-%m-%d')
            ms.store(
                f"Memory Health Report — {now_str}",
                report,
                description=f"Autodream consolidation report for {now_str}",
                mem_type="system",
            )

            # Clean up old reports beyond retention limit
            retention = ad_cfg.get("report_retention", 3)
            _cleanup_autodream_reports(agent_id, ms, retention)

            finished = datetime.datetime.now()
            elapsed = (finished - started).total_seconds()
            logging.info(f"Autodream complete for {agent_id}: health_score={results.get('health_score', '?')} ({elapsed:.1f}s)")
            _update_status(status="completed", phase="done", finished_at=finished.isoformat(),
                           elapsed=round(elapsed, 1), results={
                               "health_score": results.get("health_score"),
                               "merged": results.get("dedup", {}).get("merged", 0),
                               "stale": results.get("staleness", {}).get("stale_count", 0),
                               "conflicts": results.get("conflicts", {}).get("conflicts_found", 0),
                               "skill_candidates": results.get("skill_candidates", {}).get("candidates_found", 0),
                           })

        except Exception as e:
            logging.error(f"Autodream failed for {agent_id}: {e}")
            _update_status(status="error", phase="failed",
                           finished_at=datetime.datetime.now().isoformat(),
                           error=str(e)[:300])

    threading.Thread(target=_run, daemon=True, name=f"autodream_{agent_id}").start()


def _cleanup_autodream_reports(agent_id: str, ms: MemoryStore, retention: int):
    """Delete old Memory Health Report files beyond the retention limit."""
    agent_dir = os.path.join(AGENTS_DIR, agent_id)
    reports = []
    for fname in os.listdir(agent_dir):
        if not fname.endswith(".md"):
            continue
        fpath = os.path.join(agent_dir, fname)
        try:
            with open(fpath, "r") as f:
                raw = f.read(500)
            fm, _ = _parse_frontmatter(raw)
            if fm.get("name", "").strip('"').strip("'").startswith("Memory Health Report"):
                reports.append((fpath, os.path.getmtime(fpath), fm.get("name", "")))
        except Exception:
            continue

    # Sort by mtime descending, delete older ones
    reports.sort(key=lambda x: x[1], reverse=True)
    for fpath, _, name in reports[retention:]:
        try:
            os.remove(fpath)
        except Exception:
            pass
    if len(reports) > retention:
        _qmd_debounced_embed(agent_id)


def get_memory_health(agent_id: str) -> dict:
    """Compute live memory health stats for an agent."""
    agent_dir = os.path.join(AGENTS_DIR, agent_id)
    if not os.path.isdir(agent_dir):
        return {"error": f"Agent not found: {agent_id}"}

    now = datetime.datetime.now()
    total = 0
    by_type: dict[str, int] = {}
    stale_count = 0
    stale_names = []
    ages = []
    recall_hot = 0     # recalled in last 7 days
    recall_warm = 0    # recalled in last 30 days
    recall_cold = 0    # recalled 30+ days ago
    recall_never = 0   # never recalled
    auto_24h = 0
    auto_7d = 0
    reports = []
    last_results = None
    last_run = None
    cutoff_24h = time.time() - 86400
    cutoff_7d = time.time() - 604800

    skip_prefixes = ("Memory Summary", "Memory Health Report")

    for fname in os.listdir(agent_dir):
        if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
            continue
        if fname.startswith("ingest-"):
            continue
        fpath = os.path.join(agent_dir, fname)
        try:
            with open(fpath, "r") as f:
                raw = f.read()
            fm, body = _parse_frontmatter(raw)
        except Exception:
            continue

        name = fm.get("name", fname.replace(".md", "")).strip('"').strip("'")
        mem_type = fm.get("type", "general")
        mtime = os.path.getmtime(fpath)

        # Collect health reports
        if name.startswith("Memory Health Report"):
            # Parse health score from body
            score_match = re.search(r'Health Score:\s*(\d+)/100', body)
            score = int(score_match.group(1)) if score_match else 0
            date_match = re.search(r'(\d{4}-\d{2}-\d{2})', name)
            date_str = date_match.group(1) if date_match else datetime.datetime.fromtimestamp(mtime).strftime('%Y-%m-%d')
            reports.append({"date": date_str, "health_score": score, "mtime": mtime})
            if not last_run or mtime > last_run:
                last_run = mtime
                # Try to parse last results from report
                try:
                    last_results = _parse_health_report(body)
                except Exception:
                    last_results = None
            continue

        if any(name.startswith(p) for p in skip_prefixes):
            continue

        total += 1
        by_type[mem_type] = by_type.get(mem_type, 0) + 1

        # Age
        age_days = (now - datetime.datetime.fromtimestamp(mtime)).days
        ages.append(age_days)

        # Staleness
        if fm.get("stale", "").lower() == "true":
            stale_count += 1
            stale_names.append(name)

        # Recall frequency
        lr = fm.get("last_recalled", "")
        if lr:
            try:
                lr_date = datetime.datetime.strptime(lr.strip(), "%Y-%m-%d")
                days_since = (now - lr_date).days
                if days_since <= 7:
                    recall_hot += 1
                elif days_since <= 30:
                    recall_warm += 1
                else:
                    recall_cold += 1
            except ValueError:
                recall_never += 1
        else:
            recall_never += 1

        # Auto-memory rate (by creation time)
        if mtime > cutoff_24h:
            auto_24h += 1
        if mtime > cutoff_7d:
            auto_7d += 1

    # Sort reports by mtime
    reports.sort(key=lambda x: x.get("mtime", 0), reverse=True)
    for r in reports:
        r.pop("mtime", None)

    # QMD health
    qmd_health = {"indexed": 0, "embedded": 0, "stale": 0, "not_indexed": 0}
    try:
        r = subprocess.run(["qmd", "collection", "show", agent_id],
                           capture_output=True, text=True, timeout=5)
        if r.returncode == 0:
            # Count from the docs endpoint instead
            pass
    except Exception:
        pass

    # Health score from latest report, or compute a basic one
    if reports:
        health_score = reports[0].get("health_score", 0)
    else:
        stale_pct = (stale_count / max(total, 1)) * 100
        health_score = max(0, min(100, int(100 - (stale_pct * 0.3))))

    ad_cfg = _get_autodream_config(agent_id)

    return {
        "agent": agent_id,
        "total_memories": total,
        "by_type": by_type,
        "qmd_health": qmd_health,
        "stale_count": stale_count,
        "stale_names": stale_names[:50],
        "age_distribution": {
            "avg_days": round(sum(ages) / max(len(ages), 1), 1),
            "oldest_days": max(ages) if ages else 0,
            "newest_days": min(ages) if ages else 0,
        },
        "recall_frequency": {
            "hot": recall_hot, "warm": recall_warm,
            "cold": recall_cold, "never_recalled": recall_never,
        },
        "autodream": {
            "config": ad_cfg,
            "last_run": datetime.datetime.fromtimestamp(last_run).isoformat() if last_run else None,
            "last_results": last_results,
            "reports": reports[:ad_cfg.get("report_retention", 3)],
            "live": get_autodream_status(agent_id),
        },
        "auto_memory_rate": {"last_24h": auto_24h, "last_7d": auto_7d},
        "health_score": health_score,
    }


def _parse_health_report(body: str) -> dict:
    """Parse structured data from a Memory Health Report body."""
    results = {}
    # Dedup section
    m = re.search(r'Duplicates found:\s*(\d+)', body)
    results["duplicates_found"] = int(m.group(1)) if m else 0
    m = re.search(r'Merged:\s*(\d+)', body)
    results["merged"] = int(m.group(1)) if m else 0
    # Staleness
    m = re.search(r'Stale.*?:\s*(\d+)', body)
    results["stale_count"] = int(m.group(1)) if m else 0
    m = re.search(r'Newly flagged:\s*(\d+)', body)
    results["newly_stale"] = int(m.group(1)) if m else 0
    # Conflicts
    m = re.search(r'Conflicts detected:\s*(\d+)', body)
    results["conflicts_found"] = int(m.group(1)) if m else 0
    # Parse conflict pairs
    conflict_pairs = []
    for cm in re.finditer(r'\*\*(.+?)\*\* ↔ \*\*(.+?)\*\*:\s*(.+)', body):
        conflict_pairs.append({"a": cm.group(1), "b": cm.group(2), "nature": cm.group(3).strip()})
    results["conflict_pairs"] = conflict_pairs
    # Skill candidates
    m = re.search(r'Candidates found:\s*(\d+)', body)
    results["candidates_found"] = int(m.group(1)) if m else 0
    candidate_details = []
    for cm in re.finditer(r'- \*\*(.+?)\*\*:\s*(.+)', body):
        name = cm.group(1)
        if name not in [p["a"] for p in conflict_pairs] + [p["b"] for p in conflict_pairs]:
            candidate_details.append({"name": name, "reason": cm.group(2).strip()})
    results["candidate_details"] = candidate_details
    return results


def get_graph_stats(agent_id: str) -> dict:
    """Return knowledge graph statistics for an agent."""
    agent_dir = os.path.join(AGENTS_DIR, agent_id)
    if not os.path.isdir(agent_dir):
        return {"error": f"Agent not found: {agent_id}"}

    total_nodes = 0
    total_edges = 0
    auto_discovered_edges = 0
    edge_type_counts: dict[str, int] = {}
    entity_count = 0

    auto_edge_types = {"same_topic", "co_recalled", "depends_on", "contradicts", "extends"}

    for fname in os.listdir(agent_dir):
        if not fname.endswith(".md") or fname in _QMD_IGNORE_FILES:
            continue
        if fname.startswith("ingest-"):
            continue
        total_nodes += 1
        fpath = os.path.join(agent_dir, fname)
        try:
            with open(fpath, "r") as f:
                raw = f.read(2000)
            # Count edges (related entries)
            rel_files = re.findall(r'file:\s*(\S+\.md)', raw)
            rel_types = re.findall(r'type:\s*(\w+)', raw)
            for i, _ in enumerate(rel_files):
                total_edges += 1
                rtype = rel_types[i] if i < len(rel_types) else "references"
                edge_type_counts[rtype] = edge_type_counts.get(rtype, 0) + 1
                if rtype in auto_edge_types:
                    auto_discovered_edges += 1
        except Exception:
            continue

    # Count entities from entity index
    _ensure_entity_index(agent_id)
    with _entity_index_lock:
        idx = _entity_indices.get(agent_id, {})
        entity_count = len(idx)

    return {
        "agent": agent_id,
        "total_nodes": total_nodes,
        "total_edges": total_edges,
        "auto_discovered_edges": auto_discovered_edges,
        "entity_count": entity_count,
        "edge_types": edge_type_counts,
    }


_thread_local = threading.local()


MAX_DELEGATE_TOOL_ROUNDS = 10  # Limit for delegated/scheduled tasks (timeout is the real safety net)


_MEMORY_TOOL_NAMES = {"memory_store", "memory_recall", "memory_delete", "memory_shared"}


def _resolve_delegate_tools(tools: bool | str, api_type: str) -> list:
    """Resolve tool definitions for _run_delegate based on tools parameter."""
    if not tools:
        return []
    if tools == "memory_only":
        allowed = _MEMORY_TOOL_NAMES
    else:
        allowed = _get_agent_tool_names()
    if api_type != "openai" and api_type != "mistral":
        return _filter_tools(TOOL_DEFINITIONS, allowed)
    return _filter_tools(TOOL_DEFINITIONS_OPENAI, allowed, is_openai=True)


def _run_delegate(messages: list[dict], model: str, system_prompt: str,
                  memory_store: MemoryStore | None = None,
                  cancel_token: CancelToken | None = None,
                  event_callback=None,
                  inference_params: dict | None = None,
                  tools: bool | str = True) -> str | None:
    """Run a delegated task in a fresh context. Returns the final text response.
    Thread-safe: uses thread-local storage for memory instead of swapping globals.

    tools: True=all tools, False=no tools, "memory_only"=only memory tools

    Routes through the SDK sidecar when available (better context management).
    Falls back to direct API call if sidecar is not running.
    """
    # Try SDK sidecar first — routes both tools=True and tools=False
    try:
        import sdk_backend
        if sdk_backend.is_sidecar_running():
            # Extract prompt from messages
            prompt = ""
            for m in messages:
                if m.get("role") == "user":
                    content = m.get("content", "")
                    if isinstance(content, str):
                        prompt = content
                    elif isinstance(content, list):
                        for b in content:
                            if isinstance(b, dict) and b.get("type") == "text":
                                prompt += b.get("text", "")
            if prompt:
                kwargs = dict(
                    prompt=prompt, model=model,
                    system_prompt=system_prompt,
                    max_turns=1 if not tools else MAX_DELEGATE_TOOL_ROUNDS,
                )
                if tools:
                    # Build tool defs for sidecar MCP server
                    # Skip SDK-native tools that the sidecar already has
                    _SDK_NATIVE = {
                        "read_file", "write_file", "edit_file",
                        "list_directory", "search_files",
                        "execute_command", "web_fetch",
                    }
                    tool_defs = []
                    for td in TOOL_DEFINITIONS:
                        if td["name"] in _SDK_NATIVE:
                            continue
                        if td["name"] not in TOOL_DISPATCH:
                            continue
                        desc = td["description"]
                        if isinstance(desc, tuple):
                            desc = " ".join(desc)
                        tool_defs.append({
                            "name": td["name"],
                            "description": desc[:1000],
                            "input_schema": td["input_schema"],
                        })
                    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                    agent_id = agent.agent_id if agent else "main"
                    kwargs["tool_defs"] = tool_defs
                    kwargs["agent_id"] = agent_id
                result = sdk_backend.query_sync(**kwargs)
                if result is not None:
                    return result
    except Exception:
        pass  # Fall through to direct API call

    # Store memory in thread-local so tool_memory_* can find it
    if memory_store:
        _thread_local.memory_store = memory_store

    api_key = _delegate_api_key
    base_url = _delegate_base_url
    api_type = _delegate_api_type

    headers = make_headers(api_key, api_type)

    if api_type == "openai" or api_type == "mistral":
        endpoint = f"{base_url}/chat/completions"
        aug_messages = [{"role": "system", "content": system_prompt}] + messages
    else:
        endpoint = f"{base_url}/messages"
        aug_messages = list(messages)

    payload = {
        "model": model,
        "max_tokens": get_model_max_output(model),
        "messages": aug_messages,
        "stream": True,
        "tools": _resolve_delegate_tools(tools, api_type),
    }
    if inference_params:
        provider = _models_config.get(model, {}).get("provider", "")
        _apply_inference_to_payload(payload, inference_params, api_type, provider)
    if api_type != "openai" and api_type != "mistral":
        payload["system"] = system_prompt

    try:
        # Use stricter tool round limit via thread-local (thread-safe)
        _thread_local.max_tool_rounds = MAX_DELEGATE_TOOL_ROUNDS
        try:
            # Mistral path: use SDK directly
            if api_type == "mistral":
                return _handle_mistral_response(
                    payload, messages, model, api_key, base_url,
                    api_type, True, True, cancel_token, 0, event_callback)

            data = json.dumps(payload).encode("utf-8")
            request = urllib.request.Request(endpoint, data=data, headers=headers, method="POST")

            with urllib.request.urlopen(request) as response:
                if api_type == "openai":
                    return _handle_openai_response(
                        response, payload, messages, model, api_key, base_url,
                        api_type, True, True, headers, endpoint,
                        cancel_token, 0, event_callback)
                else:
                    return _handle_anthropic_response(
                        response, payload, messages, model, api_key, base_url,
                        api_type, True, True, headers, endpoint,
                        cancel_token, 0, event_callback)
        finally:
            _thread_local.max_tool_rounds = None
    except Exception as e:
        return f"Delegation error: {e}"


# Globals for delegation (set in _run_interactive)
_delegate_fallback_model: str | None = None
_delegate_api_key: str = ""
_delegate_base_url: str = ""
_delegate_api_type: str = "anthropic"


# --- Background Task Runner ---

import uuid as _uuid


class TaskRunner:
    """Manages background agent tasks with status tracking and cancellation."""

    def __init__(self):
        self._tasks: dict[str, dict] = {}  # task_id -> task_info
        self._threads: dict[str, threading.Thread] = {}
        self._cancel_flags: dict[str, threading.Event] = {}
        self._lock = threading.Lock()

    def submit(self, agent_id: str, task: str, model: str | None = None) -> str:
        """Submit a task to run in a background thread. Returns task_id."""
        task_id = _uuid.uuid4().hex[:8]
        cancel_flag = threading.Event()

        with self._lock:
            self._tasks[task_id] = {
                "id": task_id,
                "agent": agent_id,
                "task": task,
                "model": model,
                "status": "running",
                "result": None,
                "error": None,
                "submitted_at": datetime.datetime.now().isoformat(),
                "finished_at": None,
            }
            self._cancel_flags[task_id] = cancel_flag

        thread = threading.Thread(
            target=self._run_task, args=(task_id, agent_id, task, model, cancel_flag),
            daemon=True)
        self._threads[task_id] = thread
        thread.start()
        return task_id

    def get_status(self, task_id: str) -> dict | None:
        with self._lock:
            return self._tasks.get(task_id, {}).copy() if task_id in self._tasks else None

    def list_tasks(self) -> list[dict]:
        with self._lock:
            return [t.copy() for t in self._tasks.values()]

    def cancel(self, task_id: str) -> bool:
        with self._lock:
            if task_id not in self._tasks:
                return False
            if self._tasks[task_id]["status"] != "running":
                return False
            self._cancel_flags[task_id].set()
            self._tasks[task_id]["status"] = "cancelled"
            self._tasks[task_id]["finished_at"] = datetime.datetime.now().isoformat()
            return True

    def get_result(self, task_id: str) -> dict | None:
        """Get result, blocking until complete if still running. Timeout 0.1s poll."""
        if task_id not in self._threads:
            return self.get_status(task_id)
        # Wait for thread to finish (with timeout so we don't block forever)
        self._threads[task_id].join(timeout=300)
        return self.get_status(task_id)

    def _run_task(self, task_id: str, agent_id: str, task: str,
                  model: str | None, cancel_flag: threading.Event):
        """Execute a task in a background thread."""
        target = AgentConfig(agent_id)
        target_memory = MemoryStore(agent_id, base_dir=target.memory_dir)

        if not model:
            model = target.preferred_model or _delegate_fallback_model or "claude-opus-4-5-20251101"

        import platform
        cwd = os.getcwd()
        os_name = platform.system()
        soul = target.soul
        tools_guide = target.tools_guide

        from datetime import datetime as _dt
        system_prompt = (
            f"{soul}\n\n"
            f"You are agent '{agent_id}' running a background task.\n"
            f"Current date and time: {_dt.now().strftime('%Y-%m-%d %H:%M %Z').strip()}\n"
            f"Current working directory: {cwd}\n"
            f"Operating system: {os_name}\n\n"
            "Complete the task and provide a concise result summary.\n"
        )
        # Inject team context
        team_info = _get_agent_team_info(agent_id)
        if team_info:
            if team_info["is_head"]:
                peers = [m for m in team_info["members"] if m != agent_id]
                system_prompt += (
                    f"\nTEAM: You are the head of team '{team_info['name']}'. "
                    f"Your team members: {', '.join(peers)}\n"
                    "Delegate sub-tasks to your team members when appropriate.\n"
                    "Use memory_shared(scope='team') for team-level shared knowledge.\n"
                )
            else:
                peers = [m for m in team_info["members"] if m != agent_id and m != team_info["head"]]
                system_prompt += (
                    f"\nTEAM: You are a member of team '{team_info['name']}'.\n"
                    f"Team head: {team_info['head']}\n"
                )
                if peers:
                    system_prompt += f"Team peers: {', '.join(peers)}\n"
                system_prompt += "Use memory_shared(scope='team') for team-level shared knowledge.\n"

        if tools_guide:
            system_prompt += f"\n--- TOOL USAGE GUIDE ---\n{tools_guide}"

        # Build team-aware agent registry for this delegate
        agent_registry = build_agent_registry(for_agent_id=agent_id)
        if agent_registry:
            system_prompt += f"\n\n{agent_registry}\n"

        messages = [{"role": "user", "content": task}]

        result_text = ""
        status = "completed"
        try:
            # Store delegate agent context in thread-local (thread-safe, no global mutation)
            _thread_local.delegate_agent_id = agent_id
            _thread_local.current_agent = target
            _thread_local.memory_store = target_memory
            if cancel_flag.is_set():
                status = "cancelled"
            else:
                delegate_inf = get_inference_params(model, target.config.get("model_purpose"))
                result_text = _run_delegate(messages, model, system_prompt,
                                            memory_store=target_memory,
                                            inference_params=delegate_inf) or ""
                if cancel_flag.is_set():
                    status = "cancelled"
        except Exception as e:
            result_text = str(e)
            status = "error"
        finally:
            # Clean up thread-local state
            _thread_local.delegate_agent_id = None
            _thread_local.current_agent = None
            _thread_local.memory_store = None

        with self._lock:
            self._tasks[task_id]["status"] = status
            self._tasks[task_id]["result"] = result_text
            self._tasks[task_id]["finished_at"] = datetime.datetime.now().isoformat()
            if status == "error":
                self._tasks[task_id]["error"] = result_text
            # Clean up thread reference
            self._threads.pop(task_id, None)
            self._cancel_flags.pop(task_id, None)


# Global task runner
_task_runner: TaskRunner | None = None


# --- Workflow Engine ---

try:
    import yaml as _yaml
except ImportError:
    _yaml = None  # Graceful fallback — suggest pip3 install pyyaml


class WorkflowEngine:
    """Manages workflow definitions stored as YAML files per agent."""

    @staticmethod
    def _workflows_dir(agent_id: str) -> str:
        return os.path.join(AGENTS_DIR, agent_id, "workflows")

    @staticmethod
    def list_workflows(agent_id: str) -> list[dict]:
        """Scan agents/<name>/workflows/*.yaml and return summaries."""
        wdir = WorkflowEngine._workflows_dir(agent_id)
        if not os.path.isdir(wdir):
            return []
        results = []
        for fname in sorted(os.listdir(wdir)):
            if not fname.endswith((".yaml", ".yml")):
                continue
            wf = WorkflowEngine.get_workflow(agent_id, fname.rsplit(".", 1)[0])
            if wf:
                results.append({
                    "name": wf.get("name", fname),
                    "file": fname,
                    "description": wf.get("description", ""),
                    "stages": len(wf.get("stages", [])),
                    "variables": [v.get("name", "") for v in wf.get("variables", [])],
                })
        return results

    @staticmethod
    def get_workflow(agent_id: str, name: str) -> dict | None:
        """Parse a workflow YAML file. Returns dict or None."""
        if not _yaml:
            return None
        wdir = WorkflowEngine._workflows_dir(agent_id)
        for ext in (".yaml", ".yml"):
            fpath = os.path.join(wdir, name + ext)
            if os.path.exists(fpath):
                try:
                    with open(fpath, "r") as f:
                        return _yaml.safe_load(f)
                except Exception:
                    return None
        return None

    @staticmethod
    def save_workflow(agent_id: str, name: str, definition: dict | str) -> str:
        """Write a workflow YAML file. Returns the file path."""
        if not _yaml:
            raise RuntimeError("PyYAML is not installed. Run: pip3 install pyyaml")
        wdir = WorkflowEngine._workflows_dir(agent_id)
        os.makedirs(wdir, exist_ok=True)
        fpath = os.path.join(wdir, name + ".yaml")
        with open(fpath, "w") as f:
            if isinstance(definition, str):
                f.write(definition)
            else:
                _yaml.dump(definition, f, default_flow_style=False, sort_keys=False)
        return fpath

    @staticmethod
    def delete_workflow(agent_id: str, name: str) -> bool:
        """Remove a workflow file. Returns True if deleted."""
        wdir = WorkflowEngine._workflows_dir(agent_id)
        for ext in (".yaml", ".yml"):
            fpath = os.path.join(wdir, name + ext)
            if os.path.exists(fpath):
                os.remove(fpath)
                return True
        return False


class WorkflowExecution:
    """Runs a workflow: sequential stage execution with approval gates."""

    def __init__(self, workflow: dict, variables: dict, agent_id: str,
                 model: str | None = None, execution_id: str | None = None):
        self.workflow = workflow
        self.variables = variables or {}
        self.agent_id = agent_id
        self.model = model
        self.execution_id = execution_id or _uuid.uuid4().hex[:10]
        self.status = "pending"  # pending / running / waiting_approval / completed / failed / cancelled
        self.current_stage_idx = -1
        self.current_stage_name = ""
        self.stage_results: dict[str, dict] = {}  # stage_name -> {status, output, elapsed}
        self.started_at: str | None = None
        self.finished_at: str | None = None
        self.error: str | None = None
        self._cancel = threading.Event()
        self._approval_event = threading.Event()
        self._approval_result: str | None = None  # "approved" or "rejected"
        self._thread: threading.Thread | None = None

    @property
    def stages(self) -> list[dict]:
        return self.workflow.get("stages", [])

    def _substitute(self, text: str) -> str:
        """Replace {{variable}} and {{stages.X.output}} placeholders."""
        if not text:
            return text
        import re
        # Replace user variables: {{var_name}}
        for k, v in self.variables.items():
            text = text.replace("{{" + k + "}}", str(v))
        # Replace stage references: {{stages.X.output}}, {{stages.X.status}}
        def _stage_ref(m):
            stage_name = m.group(1)
            field = m.group(2)
            sr = self.stage_results.get(stage_name, {})
            return str(sr.get(field, f"[{stage_name}.{field} not available]"))
        text = re.sub(r"\{\{stages\.(\w+)\.(\w+)\}\}", _stage_ref, text)
        return text

    def _build_context(self) -> str:
        """Build accumulated context string from all completed stages."""
        parts = []
        for stage in self.stages:
            sname = stage.get("name", "")
            sr = self.stage_results.get(sname)
            if sr and sr.get("status") == "completed" and sr.get("output"):
                parts.append(f"=== Stage '{sname}' result ===\n{sr['output']}")
        return "\n\n".join(parts)

    def run(self):
        """Start the workflow in a background thread."""
        self.status = "running"
        self.started_at = datetime.datetime.now().isoformat()
        self._thread = threading.Thread(
            target=self._execute, daemon=True,
            name=f"workflow-{self.execution_id}")
        self._thread.start()

    def _execute(self):
        """Sequential stage execution."""
        try:
            for idx, stage in enumerate(self.stages):
                if self._cancel.is_set():
                    self.status = "cancelled"
                    self.finished_at = datetime.datetime.now().isoformat()
                    return

                sname = stage.get("name", f"stage_{idx}")
                stype = stage.get("type", "prompt")
                self.current_stage_idx = idx
                self.current_stage_name = sname

                if stype == "approval":
                    self._run_approval_stage(sname, stage)
                    if self._cancel.is_set() or self._approval_result == "rejected":
                        if self._approval_result == "rejected":
                            self.stage_results[sname] = {
                                "status": "rejected", "output": "Approval rejected by user.",
                                "elapsed": 0,
                            }
                            self.status = "failed"
                            self.error = f"Approval rejected at stage '{sname}'"
                        else:
                            self.status = "cancelled"
                        self.finished_at = datetime.datetime.now().isoformat()
                        return
                else:
                    self._run_prompt_stage(sname, stage)

                # Check if stage failed
                sr = self.stage_results.get(sname, {})
                if sr.get("status") == "error":
                    self.status = "failed"
                    self.error = f"Stage '{sname}' failed: {sr.get('output', 'unknown error')}"
                    self.finished_at = datetime.datetime.now().isoformat()
                    return

            self.status = "completed"
            self.finished_at = datetime.datetime.now().isoformat()

        except Exception as e:
            self.status = "failed"
            self.error = str(e)
            self.finished_at = datetime.datetime.now().isoformat()

    def _run_prompt_stage(self, sname: str, stage: dict):
        """Execute a prompt stage using _run_delegate."""
        start = datetime.datetime.now()
        prompt_template = stage.get("prompt", "")
        prompt = self._substitute(prompt_template)

        # Build context from previous stages
        context = self._build_context()
        full_prompt = prompt
        if context:
            full_prompt = f"Previous workflow results:\n{context}\n\n---\n\nCurrent task:\n{prompt}"

        # Resolve agent and model
        target_agent_id = stage.get("agent", self.agent_id)
        target = AgentConfig(target_agent_id)
        target_memory = MemoryStore(target_agent_id, base_dir=target.memory_dir)

        stage_model = self.model or target.preferred_model or _delegate_fallback_model or "claude-sonnet-4-6"

        import platform
        cwd = os.getcwd()
        os_name = platform.system()
        soul = target.soul
        tools_guide = target.tools_guide

        system_prompt = (
            f"{soul}\n\n"
            f"You are agent '{target_agent_id}' executing workflow stage '{sname}'.\n"
            f"Current date and time: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}\n"
            f"Current working directory: {cwd}\n"
            f"Operating system: {os_name}\n\n"
            "Complete the task and provide a concise result summary.\n"
        )
        if tools_guide:
            system_prompt += f"\n--- TOOL USAGE GUIDE ---\n{tools_guide}"

        # Tool restriction (set via thread-local if stage specifies allowed tools)
        restricted_tools = stage.get("tools")

        self.stage_results[sname] = {"status": "running", "output": "", "elapsed": 0}

        messages = [{"role": "user", "content": full_prompt}]

        try:
            _thread_local.delegate_agent_id = target_agent_id
            _thread_local.current_agent = target
            _thread_local.memory_store = target_memory
            if restricted_tools:
                _thread_local.workflow_allowed_tools = set(restricted_tools)

            cancel_token = CancelToken()
            # Link our cancel event to the cancel token
            def _watch_cancel():
                self._cancel.wait()
                cancel_token.cancel()
            watcher = threading.Thread(target=_watch_cancel, daemon=True)
            watcher.start()

            delegate_inf = get_inference_params(stage_model, target.config.get("model_purpose"))
            result_text = _run_delegate(
                messages, stage_model, system_prompt,
                memory_store=target_memory,
                cancel_token=cancel_token,
                inference_params=delegate_inf,
            ) or ""

            elapsed = (datetime.datetime.now() - start).total_seconds()
            if self._cancel.is_set():
                self.stage_results[sname] = {"status": "cancelled", "output": result_text, "elapsed": elapsed}
            else:
                self.stage_results[sname] = {"status": "completed", "output": result_text, "elapsed": elapsed}

        except Exception as e:
            elapsed = (datetime.datetime.now() - start).total_seconds()
            self.stage_results[sname] = {"status": "error", "output": str(e), "elapsed": elapsed}
        finally:
            _thread_local.delegate_agent_id = None
            _thread_local.current_agent = None
            _thread_local.memory_store = None
            _thread_local.workflow_allowed_tools = None

    def _run_approval_stage(self, sname: str, stage: dict):
        """Pause for human approval."""
        message = self._substitute(stage.get("message", "Approval required."))
        self.stage_results[sname] = {
            "status": "waiting_approval", "output": message, "elapsed": 0,
        }
        self.status = "waiting_approval"
        self._approval_event.clear()
        self._approval_result = None

        # Wait until approved, rejected, or cancelled
        while not self._approval_event.is_set() and not self._cancel.is_set():
            self._approval_event.wait(timeout=1.0)

        if self._cancel.is_set():
            self.stage_results[sname] = {"status": "cancelled", "output": message, "elapsed": 0}
            return

        if self._approval_result == "approved":
            self.stage_results[sname] = {"status": "completed", "output": "Approved.", "elapsed": 0}
            self.status = "running"
        # "rejected" handled by caller

    def approve(self):
        """Approve the current approval gate."""
        self._approval_result = "approved"
        self._approval_event.set()

    def reject(self):
        """Reject the current approval gate."""
        self._approval_result = "rejected"
        self._approval_event.set()

    def cancel(self):
        """Cancel the workflow execution."""
        self._cancel.set()
        self._approval_event.set()  # Unblock approval wait

    def to_dict(self) -> dict:
        """Serialize execution state."""
        stages_info = []
        for idx, stage in enumerate(self.stages):
            sname = stage.get("name", f"stage_{idx}")
            sr = self.stage_results.get(sname, {})
            stages_info.append({
                "name": sname,
                "type": stage.get("type", "prompt"),
                "status": sr.get("status", "pending"),
                "output": sr.get("output", ""),
                "elapsed": sr.get("elapsed", 0),
            })
        return {
            "execution_id": self.execution_id,
            "workflow_name": self.workflow.get("name", ""),
            "agent": self.agent_id,
            "model": self.model,
            "status": self.status,
            "current_stage": self.current_stage_name,
            "current_stage_idx": self.current_stage_idx,
            "total_stages": len(self.stages),
            "stages": stages_info,
            "variables": self.variables,
            "started_at": self.started_at,
            "finished_at": self.finished_at,
            "error": self.error,
        }


# Global workflow execution registry
_workflow_executions: dict[str, WorkflowExecution] = {}
_workflow_lock = threading.Lock()


def workflow_start(agent_id: str, workflow_name: str, variables: dict,
                   model: str | None = None) -> WorkflowExecution:
    """Start a workflow execution. Returns the execution object."""
    if not _yaml:
        raise RuntimeError("PyYAML is not installed. Run: pip3 install pyyaml")
    wf = WorkflowEngine.get_workflow(agent_id, workflow_name)
    if not wf:
        raise ValueError(f"Workflow '{workflow_name}' not found for agent '{agent_id}'")
    execution = WorkflowExecution(wf, variables, agent_id, model)
    with _workflow_lock:
        _workflow_executions[execution.execution_id] = execution
    execution.run()
    return execution


def workflow_get_execution(execution_id: str) -> WorkflowExecution | None:
    with _workflow_lock:
        return _workflow_executions.get(execution_id)


def workflow_list_executions() -> list[dict]:
    with _workflow_lock:
        return [ex.to_dict() for ex in _workflow_executions.values()]


def workflow_cleanup_old(max_age_hours: int = 24):
    """Remove completed/failed executions older than max_age_hours."""
    cutoff = datetime.datetime.now() - datetime.timedelta(hours=max_age_hours)
    with _workflow_lock:
        to_remove = []
        for eid, ex in _workflow_executions.items():
            if ex.status in ("completed", "failed", "cancelled") and ex.finished_at:
                try:
                    finished = datetime.datetime.fromisoformat(ex.finished_at)
                    if finished < cutoff:
                        to_remove.append(eid)
                except (ValueError, TypeError):
                    pass
        for eid in to_remove:
            del _workflow_executions[eid]


def tool_delegate_task(args: dict) -> str:
    """Delegate a task to another agent — runs in a background thread."""
    agent_id = args.get("agent", "")
    task = args.get("task", "")
    wait = args.get("wait", True)
    if not agent_id or not task:
        return _err("delegate_task: agent and task are required")

    available = list_agents()
    if agent_id not in available:
        return _err(f"delegate_task: agent '{agent_id}' not found. Available: {', '.join(available)}")

    # Team-aware delegation scoping (prefer thread-local for concurrent requests)
    caller_id = getattr(_thread_local, "delegate_agent_id", None)
    if not caller_id:
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        caller_id = agent.agent_id if agent else None
    if caller_id:
        scope = _get_delegation_scope(caller_id)
        if agent_id not in scope:
            return _err(f"delegate_task: '{caller_id}' cannot delegate to '{agent_id}'. Allowed: {', '.join(scope)}")

    if not _task_runner:
        return _err("Task runner not initialized")

    task_id = _task_runner.submit(agent_id, task, args.get("model"))

    if wait:
        # Synchronous: wait for result
        result = _task_runner.get_result(task_id)
        if result and result.get("status") == "completed":
            return _ok({
                "task_id": task_id,
                "agent": agent_id,
                "task": task,
                "response": result.get("result", ""),
            })
        elif result:
            return _err(f"delegate_task: {result.get('status')} — {result.get('error', '')}")
        return _err("delegate_task: no result")
    else:
        # Async: return task_id immediately
        return _ok({
            "task_id": task_id,
            "agent": agent_id,
            "task": task,
            "status": "running",
            "message": f"Task submitted. Use task_status(task_id='{task_id}') to check progress.",
        })


def tool_task_status(args: dict) -> str:
    """Check status of a background task."""
    if not _task_runner:
        return _err("Task runner not initialized")
    task_id = args.get("task_id", "")
    if task_id:
        status = _task_runner.get_status(task_id)
        if not status:
            return _err(f"Task '{task_id}' not found")
        # Truncate long results
        if status.get("result") and len(status["result"]) > 2000:
            status["result"] = status["result"][:2000] + "..."
        return _ok(status)
    else:
        # List all tasks
        tasks = _task_runner.list_tasks()
        for t in tasks:
            if t.get("result") and len(t["result"]) > 200:
                t["result"] = t["result"][:200] + "..."
        return _ok({"tasks": tasks, "count": len(tasks)})


def tool_task_cancel(args: dict) -> str:
    """Cancel a running background task."""
    if not _task_runner:
        return _err("Task runner not initialized")
    task_id = args.get("task_id", "")
    if not task_id:
        return _err("task_cancel: task_id is required")
    if _task_runner.cancel(task_id):
        return _ok({"task_id": task_id, "status": "cancelled"})
    return _err(f"Cannot cancel task '{task_id}' — not found or not running")


# --- Scheduler ---

import sqlite3

SCHEDULER_DB = os.path.join(AGENTS_DIR, "main", "scheduler.db")

_sched_db_lock = threading.Lock()
_sched_db_pool: dict[int, sqlite3.Connection] = {}


def _sched_conn():
    """Get a thread-safe reusable SQLite connection for the scheduler DB."""
    tid = threading.current_thread().ident
    with _sched_db_lock:
        conn = _sched_db_pool.get(tid)
        if conn is None:
            conn = sqlite3.connect(SCHEDULER_DB, timeout=10, check_same_thread=False)
            conn.execute("PRAGMA busy_timeout = 5000")
            conn.execute("PRAGMA journal_mode = WAL")
            _sched_db_pool[tid] = conn
    return conn


class Scheduler:
    """Background task scheduler with cron-like scheduling."""

    def __init__(self):
        os.makedirs(os.path.dirname(SCHEDULER_DB), exist_ok=True)
        self._init_db()
        self._stop = threading.Event()
        self._thread = None
        self._lock = threading.Lock()
        self._running_tasks: dict[str, dict] = {}

    def _init_db(self):
        with _sched_conn() as conn:
            conn.execute("""
                CREATE TABLE IF NOT EXISTS schedules (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    name TEXT NOT NULL UNIQUE,
                    task TEXT NOT NULL,
                    schedule TEXT NOT NULL,
                    agent TEXT DEFAULT 'main',
                    model TEXT,
                    enabled INTEGER DEFAULT 1,
                    last_run TEXT,
                    next_run TEXT,
                    created_at TEXT DEFAULT (datetime('now'))
                )
            """)
            conn.execute("""
                CREATE TABLE IF NOT EXISTS schedule_history (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    schedule_id INTEGER,
                    schedule_name TEXT,
                    agent TEXT,
                    task TEXT,
                    status TEXT,
                    result TEXT,
                    started_at TEXT,
                    finished_at TEXT DEFAULT (datetime('now')),
                    FOREIGN KEY (schedule_id) REFERENCES schedules(id)
                )
            """)
            # Migration: add timeout column if missing
            try:
                conn.execute("ALTER TABLE schedules ADD COLUMN timeout INTEGER DEFAULT 300")
            except sqlite3.OperationalError:
                pass
            conn.commit()

    def add(self, name: str, task: str, schedule: str,
            agent: str = "main", model: str | None = None,
            timeout: int = 300) -> dict:
        """Add a scheduled task. timeout in seconds (default: 300 = 5 min)."""
        next_run = self._calc_next_run(schedule)
        if next_run is None:
            return {"error": f"Invalid schedule format: {schedule}"}
        try:
            with _sched_conn() as conn:
                conn.execute("""
                    INSERT INTO schedules (name, task, schedule, agent, model, next_run, timeout)
                    VALUES (?, ?, ?, ?, ?, ?, ?)
                """, (name, task, schedule, agent, model, next_run.isoformat(), timeout))
                conn.commit()
            return {"name": name, "schedule": schedule, "agent": agent,
                    "next_run": next_run.isoformat(), "timeout": timeout, "status": "created"}
        except sqlite3.IntegrityError:
            return {"error": f"Schedule '{name}' already exists"}

    def remove(self, name: str) -> dict:
        with _sched_conn() as conn:
            r = conn.execute("DELETE FROM schedules WHERE name = ?", (name,))
            conn.commit()
            if r.rowcount == 0:
                return {"error": f"Schedule '{name}' not found"}
        return {"name": name, "status": "deleted"}

    def pause(self, name: str) -> dict:
        with _sched_conn() as conn:
            r = conn.execute("UPDATE schedules SET enabled = 0 WHERE name = ?", (name,))
            conn.commit()
            if r.rowcount == 0:
                return {"error": f"Schedule '{name}' not found"}
        return {"name": name, "status": "paused"}

    def resume(self, name: str) -> dict:
        next_run = None
        with _sched_conn() as conn:
            row = conn.execute("SELECT schedule FROM schedules WHERE name = ?", (name,)).fetchone()
            if not row:
                return {"error": f"Schedule '{name}' not found"}
            next_run = self._calc_next_run(row[0])
            conn.execute("UPDATE schedules SET enabled = 1, next_run = ? WHERE name = ?",
                         (next_run.isoformat() if next_run else None, name))
            conn.commit()
        return {"name": name, "status": "resumed", "next_run": next_run.isoformat() if next_run else None}

    def list_all(self) -> list[dict]:
        with _sched_conn() as conn:
            conn.row_factory = sqlite3.Row
            rows = conn.execute("SELECT * FROM schedules ORDER BY name").fetchall()
            return [dict(r) for r in rows]

    def get_task(self, name: str) -> dict | None:
        """Get a single scheduled task by name."""
        with _sched_conn() as conn:
            conn.row_factory = sqlite3.Row
            row = conn.execute("SELECT * FROM schedules WHERE name = ?", (name,)).fetchone()
            return dict(row) if row else None

    def get_history(self, name: str | None = None, limit: int = 20) -> list[dict]:
        with _sched_conn() as conn:
            conn.row_factory = sqlite3.Row
            if name:
                rows = conn.execute(
                    "SELECT * FROM schedule_history WHERE schedule_name = ? ORDER BY finished_at DESC LIMIT ?",
                    (name, limit)).fetchall()
            else:
                rows = conn.execute(
                    "SELECT * FROM schedule_history ORDER BY finished_at DESC LIMIT ?",
                    (limit,)).fetchall()
            return [dict(r) for r in rows]

    def _calc_next_run(self, schedule: str) -> datetime.datetime | None:
        """Calculate next run time from schedule string."""
        now = datetime.datetime.now()
        s = schedule.strip().lower()

        # every Xm, every Xh, every Xd
        m = re.match(r'every\s+(\d+)\s*(m|min|h|hour|d|day)s?', s)
        if m:
            val, unit = int(m.group(1)), m.group(2)[0]
            if unit == 'm':
                return now + datetime.timedelta(minutes=val)
            elif unit == 'h':
                return now + datetime.timedelta(hours=val)
            elif unit == 'd':
                return now + datetime.timedelta(days=val)

        # daily HH:MM
        m = re.match(r'daily\s+(\d{1,2}):(\d{2})', s)
        if m:
            hour, minute = int(m.group(1)), int(m.group(2))
            target = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
            if target <= now:
                target += datetime.timedelta(days=1)
            return target

        # weekly DOW HH:MM
        days_map = {'mon': 0, 'tue': 1, 'wed': 2, 'thu': 3, 'fri': 4, 'sat': 5, 'sun': 6}
        m = re.match(r'weekly\s+(\w{3})\s+(\d{1,2}):(\d{2})', s)
        if m:
            dow_str, hour, minute = m.group(1), int(m.group(2)), int(m.group(3))
            target_dow = days_map.get(dow_str[:3])
            if target_dow is None:
                return None
            target = now.replace(hour=hour, minute=minute, second=0, microsecond=0)
            days_ahead = (target_dow - now.weekday()) % 7
            if days_ahead == 0 and target <= now:
                days_ahead = 7
            target += datetime.timedelta(days=days_ahead)
            return target

        # once YYYY-MM-DD HH:MM
        m = re.match(r'once\s+(\d{4}-\d{2}-\d{2})\s+(\d{1,2}):(\d{2})', s)
        if m:
            date_str, hour, minute = m.group(1), int(m.group(2)), int(m.group(3))
            try:
                target = datetime.datetime.fromisoformat(f"{date_str}T{hour:02d}:{minute:02d}:00")
                return target
            except ValueError:
                return None

        return None

    def _calc_next_from_last(self, schedule: str, last_run: str) -> datetime.datetime | None:
        """Calculate next run based on last run time (for intervals)."""
        s = schedule.strip().lower()
        m = re.match(r'every\s+(\d+)\s*(m|min|h|hour|d|day)s?', s)
        if m:
            try:
                last = datetime.datetime.fromisoformat(last_run)
            except (ValueError, TypeError):
                return self._calc_next_run(schedule)
            val, unit = int(m.group(1)), m.group(2)[0]
            if unit == 'm':
                return last + datetime.timedelta(minutes=val)
            elif unit == 'h':
                return last + datetime.timedelta(hours=val)
            elif unit == 'd':
                return last + datetime.timedelta(days=val)
        return self._calc_next_run(schedule)

    def get_due_tasks(self) -> list[dict]:
        """Get tasks that are due for execution."""
        now = datetime.datetime.now().isoformat()
        with _sched_conn() as conn:
            conn.row_factory = sqlite3.Row
            rows = conn.execute("""
                SELECT * FROM schedules WHERE enabled = 1 AND next_run <= ?
            """, (now,)).fetchall()
            return [dict(r) for r in rows]

    def mark_executed(self, schedule_id: int, name: str, agent: str, task: str,
                      status: str, result: str, started_at: str = None,
                      tool_calls: int = 0):
        """Record execution and update next_run."""
        now = datetime.datetime.now()
        start = started_at or now.isoformat()
        try:
            start_dt = datetime.datetime.fromisoformat(start)
            duration = (now - start_dt).total_seconds()
        except (ValueError, TypeError):
            duration = 0
        with _sched_conn() as conn:
            # Record history with duration and tool_calls
            conn.execute("""
                INSERT INTO schedule_history (schedule_id, schedule_name, agent, task, status, result, started_at, finished_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?)
            """, (schedule_id, name, agent, task, status,
                  f"[Duration: {duration:.0f}s | Tools: {tool_calls}]\n\n{result[:10000]}",
                  start, now.isoformat()))

            # Get schedule to calc next run
            row = conn.execute("SELECT schedule FROM schedules WHERE id = ?", (schedule_id,)).fetchone()
            if row:
                schedule_str = row[0]
                if schedule_str.strip().lower().startswith("once"):
                    # One-shot: disable after execution
                    conn.execute("UPDATE schedules SET enabled = 0, last_run = ? WHERE id = ?",
                                 (now.isoformat(), schedule_id))
                else:
                    next_run = self._calc_next_from_last(schedule_str, now.isoformat())
                    conn.execute("UPDATE schedules SET last_run = ?, next_run = ? WHERE id = ?",
                                 (now.isoformat(),
                                  next_run.isoformat() if next_run else None,
                                  schedule_id))
            conn.commit()

    def start(self):
        """Start the background scheduler thread."""
        self._stop.clear()
        self._thread = threading.Thread(target=self._run_loop, daemon=True)
        self._thread.start()

    def stop(self):
        self._stop.set()
        if self._thread:
            self._thread.join(timeout=2)

    def _run_loop(self):
        """Background loop that checks for due tasks every 30s.
        Each task runs in its own thread to avoid blocking other due tasks."""
        while not self._stop.is_set():
            try:
                due = self.get_due_tasks()
                for task_row in due:
                    if self._stop.is_set():
                        break
                    t = threading.Thread(
                        target=self._execute_scheduled, args=(task_row,),
                        daemon=True, name=f"sched_{task_row.get('name', '?')}")
                    t.start()
            except Exception:
                pass
            self._stop.wait(30)

    def get_running_tasks(self) -> list[dict]:
        """Get currently running scheduled tasks with live stats."""
        with self._lock:
            return [dict(t) for t in self._running_tasks.values()]

    def cancel_running_task(self, name: str) -> bool:
        """Cancel a running scheduled task."""
        with self._lock:
            task = self._running_tasks.get(name)
            if task and task.get("cancel_token"):
                task["cancel_token"].cancel()
                task["status"] = "cancelling"
                return True
        return False

    def _execute_scheduled(self, task_row: dict):
        """Execute a single scheduled task."""
        agent_id = task_row.get("agent", "main")
        task = task_row.get("task", "")
        model = task_row.get("model")
        schedule_id = task_row.get("id")
        name = task_row.get("name", "")

        # Memory summary tasks: clean up orphaned chat indexes, then regenerate prompt with live data
        if name.startswith("_memory_summary_"):
            try:
                _cleanup_orphaned_chat_index_files(agent_id)
            except Exception:
                pass
            try:
                task = _build_memory_summary_prompt(agent_id)
            except Exception:
                pass

        # Track this execution
        cancel_token = CancelToken()
        run_info = {
            "name": name,
            "agent": agent_id,
            "task": task[:200],
            "model": model,
            "status": "running",
            "started_at": datetime.datetime.now().isoformat(),
            "tool_calls": 0,
            "tool_log": [],
            "cancel_token": cancel_token,
        }
        with self._lock:
            if not hasattr(self, '_running_tasks'):
                self._running_tasks = {}
            self._running_tasks[name] = run_info

        # Use delegation infrastructure
        target = AgentConfig(agent_id)
        target_memory = MemoryStore(agent_id, base_dir=target.memory_dir)

        if not model:
            model = target.preferred_model or _delegate_fallback_model or "claude-opus-4-5-20251101"

        # Build system prompt
        import platform
        cwd = os.getcwd()
        os_name = platform.system()
        soul = target.soul
        tools_guide = target.tools_guide

        system_prompt = (
            f"{soul}\n\n"
            f"You are agent '{agent_id}' executing a scheduled task: '{name}'.\n"
            f"Current date and time: {datetime.datetime.now().strftime('%Y-%m-%d %H:%M')}\n"
            f"Current working directory: {cwd}\n"
            f"Operating system: {os_name}\n\n"
            "IMPORTANT RULES FOR SCHEDULED TASKS:\n"
            "- Complete the task QUICKLY and CONCISELY. You have a limited number of tool calls.\n"
            "- Do NOT repeat the same tool call. If a search returns results, use them — don't search again.\n"
            "- Do NOT loop. One search per topic is enough. Summarize what you found.\n"
            "- Provide a concise result summary within 3-5 tool calls maximum.\n"
            "- If you can't find what you need in 2 searches, summarize what you have and stop.\n\n"
        )
        # Inject memory summary for context (skip for the summary task itself)
        tcfg = target.config.get("token_config", {})
        if not name.startswith("_memory_summary_") and tcfg.get("include_memory_summary", True):
            try:
                mem_summary = get_memory_summary(agent_id)
                if mem_summary:
                    cap = tcfg.get("memory_summary_cap", 3000)
                    if len(mem_summary) > cap:
                        mem_summary = mem_summary[:cap] + "\n...(truncated)"
                    system_prompt += (
                        "MEMORY SUMMARY (auto-generated synthesis of recent activity — use as context):\n"
                        f"{mem_summary}\n\n"
                    )
            except Exception:
                pass
        if tools_guide and tcfg.get("include_tools_guide", True):
            system_prompt += f"\n--- TOOL USAGE GUIDE ---\n{tools_guide}"

        messages = [{"role": "user", "content": task}]

        # Get timeout (default 5 min)
        task_timeout = task_row.get("timeout") or 300

        # Run with isolated memory and live tracking
        result_text = ""
        status = "success"

        def on_event(event_type, data):
            if event_type == "tool_call":
                run_info["tool_calls"] += 1
                entry = f"{data.get('name','')}({str(data.get('args',{}))[:80]})"
                run_info["tool_log"].append(entry)
                if len(run_info["tool_log"]) > 50:
                    run_info["tool_log"] = run_info["tool_log"][-50:]

        # Timeout watchdog — cancels the task after timeout seconds
        def watchdog():
            if not cancel_token._cancelled.wait(task_timeout):
                cancel_token.cancel()
                run_info["status"] = "timeout"

        timer = threading.Thread(target=watchdog, daemon=True)
        timer.start()

        try:
            # Set thread-local context for tools that need agent/memory
            _thread_local.current_agent = target
            _thread_local.memory_store = target_memory
            _thread_local.delegate_agent_id = agent_id

            sched_inf = get_inference_params(model, target.config.get("model_purpose"))
            # Memory summary tasks only need memory tools, not the full 39-tool schema
            sched_tools = tcfg.get("scheduled_task_tools", True)
            if name.startswith("_memory_summary_"):
                sched_tools = "memory_only"
            result_text = _run_delegate(messages, model, system_prompt,
                                        memory_store=target_memory,
                                        cancel_token=cancel_token,
                                        event_callback=on_event,
                                        inference_params=sched_inf,
                                        tools=sched_tools) or ""
            # Check if _run_delegate returned an error string instead of raising
            if result_text.startswith("Delegation error:"):
                status = "error"
                result_text = f"[DELEGATION ERROR] {result_text}"
        except TaskCancelled:
            if run_info.get("status") == "timeout":
                elapsed = (datetime.datetime.now() - datetime.datetime.fromisoformat(run_info["started_at"])).total_seconds()
                result_text = (result_text or "") + f"\n\n[TIMEOUT] Task timed out after {elapsed:.0f}s (limit: {task_timeout}s). Tool calls made: {run_info['tool_calls']}."
                status = "timeout"
            else:
                result_text = (result_text or "") + f"\n\n[CANCELLED] Task was cancelled. Tool calls made: {run_info['tool_calls']}."
                status = "cancelled"
        except urllib.error.HTTPError as e:
            error_body = ""
            try:
                error_body = e.read().decode("utf-8")[:500]
            except Exception:
                pass
            result_text = f"[HTTP ERROR] {e.code} {e.reason}\n{error_body}"
            status = "error"
        except urllib.error.URLError as e:
            result_text = f"[CONNECTION ERROR] Could not reach LLM API: {e.reason}"
            status = "error"
        except Exception as e:
            import traceback
            tb = traceback.format_exc()
            result_text = f"[ERROR] {type(e).__name__}: {e}\n\n{tb[-500:]}"
            status = "error"
        finally:
            cancel_token.cancel()  # stop the watchdog if still running
            # Clean up thread-local state
            _thread_local.current_agent = None
            _thread_local.memory_store = None
            _thread_local.delegate_agent_id = None
            with self._lock:
                self._running_tasks.pop(name, None)

        self.mark_executed(schedule_id, name, agent_id, task, status, result_text,
                           started_at=run_info.get("started_at"),
                           tool_calls=run_info.get("tool_calls", 0))

        # Chain relationship discovery after memory summary completes
        if name.startswith("_memory_summary_") and status == "success":
            rd_cfg = _get_relationship_discovery_config(agent_id)
            if rd_cfg.get("enabled"):
                try:
                    trigger_relationship_discovery(agent_id)
                except Exception:
                    pass

        # Chain autodream after relationship discovery completes
        if name.startswith("_relationship_discovery_") and status == "success":
            ad_cfg = _get_autodream_config(agent_id)
            if ad_cfg.get("enabled"):
                try:
                    trigger_autodream(agent_id)
                except Exception:
                    pass

        # Fire notification hook for task completion/failure
        if _notification_hook:
            try:
                if status in ("error", "timeout"):
                    evt = "task_timeout" if status == "timeout" else "task_failed"
                    sev = "warning" if status == "timeout" else "error"
                    _notification_hook(evt, f"Scheduled task: {name}",
                                       result_text[:300], severity=sev,
                                       agent=agent_id,
                                       metadata={"task_name": name, "status": status,
                                                  "tool_calls": run_info.get("tool_calls", 0)})
                elif status == "success" and not name.startswith("_memory_summary_") \
                        and not name.startswith("_relationship_discovery_"):
                    _notification_hook("task_complete", f"Task completed: {name}",
                                       f"Agent {agent_id} completed '{name}' with {run_info.get('tool_calls', 0)} tool calls.",
                                       severity="info", agent=agent_id,
                                       metadata={"task_name": name})
            except Exception:
                pass


# Global scheduler instance
_scheduler: Scheduler | None = None

# Notification hook — set by server.py to dispatch notifications
_notification_hook = None


# --- Cost Tracking ---

COST_DB = os.path.join(AGENTS_DIR, "main", "costs.db")

_cost_db_lock = threading.Lock()
_cost_db_pool: dict[int, sqlite3.Connection] = {}


def _cost_conn():
    """Get a thread-safe reusable SQLite connection for the cost DB."""
    tid = threading.current_thread().ident
    with _cost_db_lock:
        conn = _cost_db_pool.get(tid)
        if conn is None:
            conn = sqlite3.connect(COST_DB, timeout=10, check_same_thread=False)
            conn.execute("PRAGMA busy_timeout = 5000")
            conn.execute("PRAGMA journal_mode = WAL")
            _cost_db_pool[tid] = conn
    return conn


# Default cost rates per 1M tokens — 0 for free providers
_cost_rates: dict[str, dict[str, float]] = {
    # Anthropic — per-million-token rates (USD)
    "claude-opus-4-6":            {"input": 15.0,  "output": 75.0},
    "claude-opus-4-5-20251101":   {"input": 15.0,  "output": 75.0},
    "claude-opus-4-20250514":     {"input": 15.0,  "output": 75.0},
    "claude-sonnet-4-6":          {"input": 3.0,   "output": 15.0},
    "claude-sonnet-4-5-20241022": {"input": 3.0,   "output": 15.0},
    "claude-sonnet-4-20250514":   {"input": 3.0,   "output": 15.0},
    "claude-haiku-4-5-20251001":  {"input": 0.80,  "output": 4.0},
    "claude-haiku-3.5":           {"input": 0.80,  "output": 4.0},
    "claude-3-5-haiku-20241022":  {"input": 0.80,  "output": 4.0},
    "claude-3-7-sonnet-20250219": {"input": 3.0,   "output": 15.0},
    # Prefix patterns for future model IDs
    "claude-opus":                {"input": 15.0,  "output": 75.0},
    "claude-sonnet":              {"input": 3.0,   "output": 15.0},
    "claude-haiku":               {"input": 0.80,  "output": 4.0},
}


def _get_cost_rate(model: str) -> dict[str, float]:
    """Look up cost rate for a model. Checks _models_config first, then defaults.
    Config values of 0 are treated as unset — auto-discovery writes 0 for all models."""
    cfg = _models_config.get(model, {})
    ci = cfg.get("cost_input")
    co = cfg.get("cost_output")
    if ci is not None and co is not None and (float(ci) > 0 or float(co) > 0):
        return {"input": float(ci), "output": float(co)}
    # Check built-in rates
    if model in _cost_rates:
        return _cost_rates[model]
    # Try prefix matching (e.g. "claude-opus-4-6" matches "claude-opus-4-6-20260101")
    ml = model.lower()
    for pattern, rate in _cost_rates.items():
        if ml.startswith(pattern.lower()) or pattern.lower() in ml:
            return rate
    return {"input": 0.0, "output": 0.0}


def _compute_cost(model: str, tokens_in: int, tokens_out: int) -> float:
    """Compute estimated cost in USD."""
    rate = _get_cost_rate(model)
    return (tokens_in * rate["input"] + tokens_out * rate["output"]) / 1_000_000


class CostTracker:
    """Thread-safe cost tracking with SQLite persistence."""

    def __init__(self):
        os.makedirs(os.path.dirname(COST_DB), exist_ok=True)
        self._init_db()

    def _init_db(self):
        with _cost_conn() as conn:
            conn.execute("""
                CREATE TABLE IF NOT EXISTS cost_log (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    agent TEXT NOT NULL,
                    session_id TEXT,
                    model TEXT NOT NULL,
                    provider TEXT NOT NULL DEFAULT '',
                    tokens_in INTEGER NOT NULL DEFAULT 0,
                    tokens_out INTEGER NOT NULL DEFAULT 0,
                    cost_usd REAL NOT NULL DEFAULT 0.0,
                    tool_round INTEGER DEFAULT 0,
                    created_at TEXT NOT NULL DEFAULT (datetime('now'))
                )
            """)
            conn.execute("CREATE INDEX IF NOT EXISTS idx_cost_agent ON cost_log(agent)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_cost_session ON cost_log(session_id)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_cost_created ON cost_log(created_at)")
            conn.commit()

    def log_call(self, agent: str, session_id: str, model: str, provider: str,
                 tokens_in: int, tokens_out: int, tool_round: int = 0):
        """Log an LLM call with cost estimation."""
        cost = _compute_cost(model, tokens_in, tokens_out)
        try:
            with _cost_conn() as conn:
                conn.execute("""
                    INSERT INTO cost_log (agent, session_id, model, provider, tokens_in, tokens_out, cost_usd, tool_round)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?)
                """, (agent, session_id or "", model, provider, tokens_in, tokens_out, cost, tool_round))
                conn.commit()
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Cost tracking error: {e}")

    def get_stats(self, agent: str | None = None, hours: int = 24) -> dict:
        """Get aggregate stats for the last N hours."""
        try:
            with _cost_conn() as conn:
                conn.row_factory = sqlite3.Row
                where = "WHERE created_at >= datetime('now', ?)"
                params: list = [f"-{hours} hours"]
                if agent:
                    where += " AND agent = ?"
                    params.append(agent)
                row = conn.execute(f"""
                    SELECT COUNT(*) as total_calls,
                           COALESCE(SUM(tokens_in), 0) as total_tokens_in,
                           COALESCE(SUM(tokens_out), 0) as total_tokens_out,
                           COALESCE(SUM(cost_usd), 0.0) as total_cost
                    FROM cost_log {where}
                """, params).fetchone()
                # Per-agent breakdown
                agents_rows = conn.execute(f"""
                    SELECT agent,
                           COUNT(*) as calls,
                           COALESCE(SUM(tokens_in), 0) as tokens_in,
                           COALESCE(SUM(tokens_out), 0) as tokens_out,
                           COALESCE(SUM(cost_usd), 0.0) as cost
                    FROM cost_log {where}
                    GROUP BY agent ORDER BY cost DESC
                """, params).fetchall()
                # Per-model breakdown
                models_rows = conn.execute(f"""
                    SELECT model,
                           COUNT(*) as calls,
                           COALESCE(SUM(cost_usd), 0.0) as cost
                    FROM cost_log {where}
                    GROUP BY model ORDER BY cost DESC
                """, params).fetchall()
                return {
                    "total_calls": row["total_calls"],
                    "total_tokens_in": row["total_tokens_in"],
                    "total_tokens_out": row["total_tokens_out"],
                    "total_cost": round(row["total_cost"], 4),
                    "hours": hours,
                    "agent_filter": agent,
                    "by_agent": [dict(r) for r in agents_rows],
                    "by_model": [dict(r) for r in models_rows],
                }
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Cost stats error: {e}")
            return {"total_calls": 0, "total_tokens_in": 0, "total_tokens_out": 0,
                    "total_cost": 0.0, "hours": hours, "agent_filter": agent,
                    "by_agent": [], "by_model": []}

    def get_daily(self, agent: str | None = None, days: int = 7) -> list[dict]:
        """Get daily breakdown for the last N days."""
        try:
            with _cost_conn() as conn:
                conn.row_factory = sqlite3.Row
                where = "WHERE created_at >= datetime('now', ?)"
                params: list = [f"-{days} days"]
                if agent:
                    where += " AND agent = ?"
                    params.append(agent)
                rows = conn.execute(f"""
                    SELECT date(created_at) as day,
                           COUNT(*) as calls,
                           COALESCE(SUM(tokens_in), 0) as tokens_in,
                           COALESCE(SUM(tokens_out), 0) as tokens_out,
                           COALESCE(SUM(cost_usd), 0.0) as cost
                    FROM cost_log {where}
                    GROUP BY date(created_at) ORDER BY day DESC
                """, params).fetchall()
                return [dict(r) for r in rows]
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Cost daily error: {e}")
            return []

    def get_session_cost(self, session_id: str) -> dict:
        """Get cost for a specific session."""
        try:
            with _cost_conn() as conn:
                conn.row_factory = sqlite3.Row
                row = conn.execute("""
                    SELECT COUNT(*) as calls,
                           COALESCE(SUM(tokens_in), 0) as tokens_in,
                           COALESCE(SUM(tokens_out), 0) as tokens_out,
                           COALESCE(SUM(cost_usd), 0.0) as cost
                    FROM cost_log WHERE session_id = ?
                """, (session_id,)).fetchone()
                return dict(row) if row else {"calls": 0, "tokens_in": 0, "tokens_out": 0, "cost": 0.0}
        except (sqlite3.Error, OSError):
            return {"calls": 0, "tokens_in": 0, "tokens_out": 0, "cost": 0.0}


_cost_tracker: CostTracker | None = None


# --- Observability: Trace Manager ---

TRACES_DB = os.path.join(AGENTS_DIR, "main", "traces.db")

_traces_db_lock = threading.Lock()
_traces_db_pool: dict[int, sqlite3.Connection] = {}


def _traces_conn():
    """Get a thread-safe reusable SQLite connection for the traces DB."""
    tid = threading.current_thread().ident
    with _traces_db_lock:
        conn = _traces_db_pool.get(tid)
        if conn is None:
            conn = sqlite3.connect(TRACES_DB, timeout=10, check_same_thread=False)
            conn.execute("PRAGMA busy_timeout = 5000")
            conn.execute("PRAGMA journal_mode = WAL")
            _traces_db_pool[tid] = conn
    return conn


class TraceManager:
    """Thread-safe span-based tracing with SQLite persistence."""

    def __init__(self):
        os.makedirs(os.path.dirname(TRACES_DB), exist_ok=True)
        self._init_db()

    def _init_db(self):
        with _traces_conn() as conn:
            conn.execute("""
                CREATE TABLE IF NOT EXISTS traces (
                    id TEXT PRIMARY KEY,
                    trace_id TEXT NOT NULL,
                    parent_id TEXT,
                    agent TEXT NOT NULL,
                    session_id TEXT,
                    type TEXT NOT NULL,
                    name TEXT NOT NULL,
                    status TEXT NOT NULL DEFAULT 'ok',
                    started_at TEXT NOT NULL,
                    ended_at TEXT,
                    duration_ms INTEGER,
                    tokens_in INTEGER DEFAULT 0,
                    tokens_out INTEGER DEFAULT 0,
                    model TEXT,
                    metadata TEXT
                )
            """)
            conn.execute("CREATE INDEX IF NOT EXISTS idx_traces_trace ON traces(trace_id)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_traces_parent ON traces(parent_id)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_traces_agent ON traces(agent)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_traces_type ON traces(type)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_traces_started ON traces(started_at)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_traces_session ON traces(session_id)")
            conn.commit()

    def start_span(self, span_type: str, name: str, agent: str = "main",
                   model: str = "", parent_id: str | None = None,
                   trace_id: str | None = None,
                   session_id: str | None = None) -> dict:
        """Start a new span. Returns span dict with id, start time, etc."""
        import uuid as _uuid
        span_id = _uuid.uuid4().hex[:16]
        if not trace_id:
            trace_id = _uuid.uuid4().hex[:16]
        now = datetime.datetime.utcnow().isoformat(timespec="milliseconds") + "Z"
        return {
            "id": span_id,
            "trace_id": trace_id,
            "parent_id": parent_id,
            "agent": agent,
            "session_id": session_id,
            "type": span_type,
            "name": name,
            "model": model,
            "status": "ok",
            "started_at": now,
            "_start_time": time.time(),
            "tokens_in": 0,
            "tokens_out": 0,
            "metadata": {},
        }

    def end_span(self, span: dict, status: str = "ok",
                 result_summary: str = "", tokens_in: int = 0, tokens_out: int = 0):
        """End a span: compute duration and persist to DB."""
        now = datetime.datetime.utcnow().isoformat(timespec="milliseconds") + "Z"
        duration_ms = int((time.time() - span.get("_start_time", time.time())) * 1000)
        span["ended_at"] = now
        span["duration_ms"] = duration_ms
        span["status"] = status
        if tokens_in:
            span["tokens_in"] = tokens_in
        if tokens_out:
            span["tokens_out"] = tokens_out
        if result_summary:
            span["metadata"]["result_summary"] = result_summary[:500]
        try:
            with _traces_conn() as conn:
                conn.execute("""
                    INSERT INTO traces (id, trace_id, parent_id, agent, session_id,
                        type, name, status, started_at, ended_at, duration_ms,
                        tokens_in, tokens_out, model, metadata)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)
                """, (
                    span["id"], span["trace_id"], span.get("parent_id"),
                    span["agent"], span.get("session_id"),
                    span["type"], span["name"], span["status"],
                    span["started_at"], span["ended_at"], duration_ms,
                    span.get("tokens_in", 0), span.get("tokens_out", 0),
                    span.get("model", ""),
                    json.dumps(span.get("metadata", {})),
                ))
                conn.commit()
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Trace write error: {e}")

    def get_traces(self, agent: str | None = None, hours: int = 24,
                   limit: int = 50) -> list[dict]:
        """Get recent root-level traces (parent_id IS NULL)."""
        try:
            with _traces_conn() as conn:
                conn.row_factory = sqlite3.Row
                where = "WHERE started_at >= datetime('now', ?) AND parent_id IS NULL"
                params: list = [f"-{hours} hours"]
                if agent:
                    where += " AND agent = ?"
                    params.append(agent)
                params.append(limit)
                rows = conn.execute(f"""
                    SELECT t.*, (SELECT COUNT(*) FROM traces c WHERE c.trace_id = t.trace_id) as span_count
                    FROM traces t {where}
                    ORDER BY started_at DESC LIMIT ?
                """, params).fetchall()
                return [dict(r) for r in rows]
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Trace read error: {e}")
            return []

    def get_trace(self, trace_id: str) -> list[dict]:
        """Get all spans for a trace, ordered by start time."""
        try:
            with _traces_conn() as conn:
                conn.row_factory = sqlite3.Row
                rows = conn.execute(
                    "SELECT * FROM traces WHERE trace_id = ? ORDER BY started_at",
                    (trace_id,)
                ).fetchall()
                return [dict(r) for r in rows]
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Trace read error: {e}")
            return []

    def cleanup(self, retention_days: int = 30):
        """Delete traces older than retention period."""
        try:
            with _traces_conn() as conn:
                conn.execute(
                    "DELETE FROM traces WHERE started_at < datetime('now', ?)",
                    (f"-{retention_days} days",)
                )
                conn.commit()
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Trace cleanup error: {e}")


_trace_manager: TraceManager | None = None


# --- Audit Trail ---

AUDIT_DB = os.path.join(AGENTS_DIR, "main", "audit.db")

_audit_db_lock = threading.Lock()
_audit_db_pool: dict[int, sqlite3.Connection] = {}


def _audit_conn():
    """Get a thread-safe reusable SQLite connection for the audit DB."""
    tid = threading.current_thread().ident
    with _audit_db_lock:
        conn = _audit_db_pool.get(tid)
        if conn is None:
            conn = sqlite3.connect(AUDIT_DB, timeout=10, check_same_thread=False)
            conn.execute("PRAGMA busy_timeout = 5000")
            conn.execute("PRAGMA journal_mode = WAL")
            _audit_db_pool[tid] = conn
    return conn


_AUDIT_ACTION_MAP = {
    "read_file": "file_read",
    "write_file": "file_write",
    "edit_file": "file_write",
    "execute_command": "command_execute",
    "gmail_send": "email_send",
    "gmail_reply": "email_reply",
    "gmail_inbox": "email_read",
    "gmail_read": "email_read",
    "gmail_search": "email_read",
    "web_fetch": "web_fetch",
    "exa_search": "web_search",
    "memory_store": "memory_store",
    "memory_delete": "memory_delete",
    "memory_recall": "memory_recall",
    "memory_shared": "memory_shared",
    "delegate_task": "delegation",
    "task_cancel": "task_cancel",
    "use_skill": "skill_use",
    "list_directory": "file_read",
    "search_files": "file_read",
    "mcp_connect": "mcp_tool_call",
    "mcp_disconnect": "mcp_tool_call",
    "mcp_servers": "mcp_tool_call",
}


def _audit_summarize_args(tool_name: str, args: dict) -> str:
    """Generate human-readable summary of tool arguments, max 200 chars."""
    if tool_name == "execute_command":
        return args.get("command", "")[:200]
    elif tool_name in ("gmail_send", "gmail_reply"):
        return f"To: {args.get('to', '?')}, Subject: {args.get('subject', '?')}"[:200]
    elif tool_name in ("write_file", "edit_file"):
        content = args.get("content", "")
        return f"{args.get('path', '?')} ({len(content)} bytes)"[:200]
    elif tool_name == "read_file":
        return f"{args.get('path', '?')}"[:200]
    elif tool_name == "delegate_task":
        return f"-> {args.get('agent', '?')}: {args.get('task', '')[:100]}"[:200]
    elif tool_name in ("memory_store", "memory_delete"):
        return f"{args.get('name', '?')}"[:200]
    elif tool_name in ("memory_recall", "memory_shared"):
        return f"query={args.get('query', '?')}"[:200]
    elif tool_name == "exa_search":
        return f"query={args.get('query', '')}"[:200]
    elif tool_name == "web_fetch":
        return f"{args.get('url', '?')}"[:200]
    elif tool_name == "use_skill":
        return f"skill={args.get('skill', '?')}"[:200]
    elif tool_name == "list_directory":
        return f"{args.get('path', '.')}"[:200]
    elif tool_name == "search_files":
        return f"pattern={args.get('pattern', '?')} in {args.get('path', '.')}"[:200]
    else:
        return str(args)[:200]


def _audit_summarize_result(tool_name: str, result_str: str) -> str:
    """Generate human-readable result summary, max 200 chars."""
    try:
        rdata = json.loads(result_str)
    except (json.JSONDecodeError, TypeError):
        return str(result_str)[:200]
    if isinstance(rdata, str):
        return rdata[:200]
    if rdata.get("error"):
        return f"ERROR: {rdata['error']}"[:200]
    if tool_name == "execute_command":
        ec = rdata.get("exit_code", -1)
        return f"exit_code={ec}"[:200]
    if tool_name == "exa_search":
        return f"{rdata.get('result_count', 0)} results"[:200]
    if tool_name in ("memory_recall", "memory_shared"):
        return f"{rdata.get('count', 0)} memories"[:200]
    if tool_name == "delegate_task":
        return f"{rdata.get('agent', '')} responded"[:200]
    return str(rdata)[:200]


class AuditLog:
    """Append-only audit log with SQLite persistence. No UPDATE or DELETE."""

    def __init__(self):
        os.makedirs(os.path.dirname(AUDIT_DB), exist_ok=True)
        self._init_db()

    def _init_db(self):
        with _audit_conn() as conn:
            conn.execute("""
                CREATE TABLE IF NOT EXISTS audit_log (
                    id INTEGER PRIMARY KEY AUTOINCREMENT,
                    timestamp TEXT NOT NULL DEFAULT (strftime('%Y-%m-%dT%H:%M:%f', 'now')),
                    agent TEXT NOT NULL,
                    session_id TEXT,
                    action_type TEXT NOT NULL,
                    tool_name TEXT NOT NULL,
                    args_summary TEXT,
                    result_summary TEXT,
                    result_status TEXT NOT NULL DEFAULT 'success',
                    duration_ms INTEGER,
                    source TEXT NOT NULL DEFAULT 'chat'
                )
            """)
            conn.execute("CREATE INDEX IF NOT EXISTS idx_audit_timestamp ON audit_log(timestamp)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_audit_agent ON audit_log(agent)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_audit_action ON audit_log(action_type)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_audit_session ON audit_log(session_id)")
            conn.execute("CREATE INDEX IF NOT EXISTS idx_audit_agent_time ON audit_log(agent, timestamp)")
            conn.commit()

    def log_action(self, agent: str, action_type: str, tool_name: str,
                   args_summary: str = "", result_summary: str = "",
                   result_status: str = "success", duration_ms: int | None = None,
                   session_id: str | None = None, source: str = "chat"):
        """Insert an audit log entry. Append-only."""
        try:
            with _audit_conn() as conn:
                conn.execute("""
                    INSERT INTO audit_log (agent, session_id, action_type, tool_name,
                        args_summary, result_summary, result_status, duration_ms, source)
                    VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
                """, (agent, session_id or "", action_type, tool_name,
                      args_summary[:200], result_summary[:200], result_status,
                      duration_ms, source))
                conn.commit()
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Audit log error: {e}")

    def query(self, agent: str | None = None, action_type: str | None = None,
              from_ts: str | None = None, limit: int = 50) -> list[dict]:
        """Query audit log entries."""
        try:
            with _audit_conn() as conn:
                conn.row_factory = sqlite3.Row
                where_parts = ["1=1"]
                params: list = []
                if agent:
                    where_parts.append("agent = ?")
                    params.append(agent)
                if action_type:
                    where_parts.append("action_type = ?")
                    params.append(action_type)
                if from_ts:
                    where_parts.append("timestamp >= ?")
                    params.append(from_ts)
                where = " AND ".join(where_parts)
                params.append(limit)
                rows = conn.execute(f"""
                    SELECT * FROM audit_log WHERE {where}
                    ORDER BY timestamp DESC LIMIT ?
                """, params).fetchall()
                return [dict(r) for r in rows]
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Audit query error: {e}")
            return []

    def export_csv(self, agent: str | None = None,
                   from_ts: str | None = None,
                   to_ts: str | None = None) -> str:
        """Export audit log as CSV string."""
        try:
            with _audit_conn() as conn:
                conn.row_factory = sqlite3.Row
                where_parts = ["1=1"]
                params: list = []
                if agent:
                    where_parts.append("agent = ?")
                    params.append(agent)
                if from_ts:
                    where_parts.append("timestamp >= ?")
                    params.append(from_ts)
                if to_ts:
                    where_parts.append("timestamp <= ?")
                    params.append(to_ts)
                where = " AND ".join(where_parts)
                rows = conn.execute(f"""
                    SELECT * FROM audit_log WHERE {where}
                    ORDER BY timestamp DESC
                """, params).fetchall()
                import csv
                import io
                output = io.StringIO()
                if rows:
                    writer = csv.DictWriter(output, fieldnames=rows[0].keys())
                    writer.writeheader()
                    for row in rows:
                        writer.writerow(dict(row))
                return output.getvalue()
        except (sqlite3.Error, OSError) as e:
            logging.warning(f"Audit export error: {e}")
            return ""


_audit_log: AuditLog | None = None


# --- Rate Limiting ---

class RateLimiter:
    """Sliding-window rate limiter per agent. In-memory only (resets on restart)."""

    def __init__(self):
        self._lock = threading.Lock()
        self._requests: dict[str, collections.deque] = collections.defaultdict(collections.deque)
        self._tokens: dict[str, collections.deque] = collections.defaultdict(collections.deque)
        self._cost: dict[str, collections.deque] = collections.defaultdict(collections.deque)

    def _prune(self, dq: collections.deque, cutoff: float):
        """Remove entries older than cutoff timestamp."""
        while dq and (dq[0] if isinstance(dq[0], (int, float)) else dq[0][0]) < cutoff:
            dq.popleft()

    def check(self, agent_id: str) -> tuple[bool, str, dict]:
        """Check if a request is allowed for this agent.

        Returns (allowed, reason, usage_info).
        Loads limits from the agent's agent.json rate_limits field.
        """
        limits = self._get_limits(agent_id)
        if not limits:
            return True, "", {}

        now = time.time()
        with self._lock:
            # Check requests/minute
            rpm_limit = limits.get("max_requests_per_minute")
            if rpm_limit:
                dq = self._requests[agent_id]
                self._prune(dq, now - 60)
                if len(dq) >= rpm_limit:
                    oldest = dq[0]
                    retry = 60 - (now - oldest)
                    return False, f"Rate limit: {rpm_limit} requests/minute exceeded. Retry in {int(retry)}s.", {
                        "dimension": "max_requests_per_minute", "current": len(dq), "limit": rpm_limit}

            # Check tokens/hour
            tph_limit = limits.get("max_tokens_per_hour")
            if tph_limit:
                dq = self._tokens[agent_id]
                self._prune(dq, now - 3600)
                total = sum(t[1] for t in dq)
                if total >= tph_limit:
                    return False, f"Rate limit: {tph_limit} tokens/hour exceeded.", {
                        "dimension": "max_tokens_per_hour", "current": total, "limit": tph_limit}

            # Check cost/day
            cpd_limit = limits.get("max_cost_per_day")
            if cpd_limit:
                dq = self._cost[agent_id]
                self._prune(dq, now - 86400)
                total = sum(t[1] for t in dq)
                if total >= cpd_limit:
                    return False, f"Rate limit: ${cpd_limit}/day cost limit exceeded.", {
                        "dimension": "max_cost_per_day", "current": total, "limit": cpd_limit}

            # Record the request timestamp
            self._requests[agent_id].append(now)

        return True, "", {}

    def record_usage(self, agent_id: str, tokens: int, cost: float):
        """Record token and cost usage after a successful response."""
        now = time.time()
        with self._lock:
            self._tokens[agent_id].append((now, tokens))
            self._cost[agent_id].append((now, cost))

    def get_status(self, agent_id: str | None = None) -> dict:
        """Get current usage vs limits for display."""
        result = {}
        agents_to_check = [agent_id] if agent_id else list(set(
            list(self._requests.keys()) + list(self._tokens.keys())))

        now = time.time()
        with self._lock:
            for aid in agents_to_check:
                limits = self._get_limits(aid)
                if not limits:
                    continue
                # Requests/minute
                dq_r = self._requests.get(aid, collections.deque())
                self._prune(dq_r, now - 60)
                rpm_limit = limits.get("max_requests_per_minute", 0)
                # Tokens/hour
                dq_t = self._tokens.get(aid, collections.deque())
                self._prune(dq_t, now - 3600)
                tph_total = sum(t[1] for t in dq_t)
                tph_limit = limits.get("max_tokens_per_hour", 0)
                # Cost/day
                dq_c = self._cost.get(aid, collections.deque())
                self._prune(dq_c, now - 86400)
                cpd_total = sum(t[1] for t in dq_c)
                cpd_limit = limits.get("max_cost_per_day", 0)

                result[aid] = {
                    "requests_per_minute": {"current": len(dq_r), "limit": rpm_limit},
                    "tokens_per_hour": {"current": tph_total, "limit": tph_limit},
                    "cost_per_day": {"current": round(cpd_total, 4), "limit": cpd_limit},
                }
        return result

    def _get_limits(self, agent_id: str) -> dict:
        """Load rate limits from agent.json."""
        try:
            agent_json = os.path.join(AGENTS_DIR, agent_id, "agent.json")
            if os.path.isfile(agent_json):
                with open(agent_json) as f:
                    cfg = json.load(f)
                return cfg.get("rate_limits", {})
        except (OSError, json.JSONDecodeError):
            pass
        return {}


_rate_limiter: RateLimiter | None = None


def tool_list_nodes(args: dict) -> str:
    """List all registered remote nodes."""
    try:
        import urllib.request
        req = urllib.request.Request("http://127.0.0.1:8420/v1/nodes", method="GET")
        with urllib.request.urlopen(req, timeout=10) as resp:
            data = json.loads(resp.read().decode("utf-8"))
        nodes = data.get("nodes", [])
        if not nodes:
            return _ok({"nodes": [], "count": 0, "message": "No nodes registered"})
        return _ok({"nodes": nodes, "count": len(nodes)})
    except Exception as e:
        return _err(f"Failed to list nodes: {e}")


def tool_context_search(args: dict) -> str:
    """Search compacted conversation history."""
    if not _context_manager:
        return _err("Context manager not initialized")
    session_id = getattr(_thread_local, 'current_session_id', None) or ""
    if not session_id:
        return _err("No active session")
    query = args.get("query", "")
    if not query:
        return _err("Missing query")
    limit = args.get("limit", 10)
    results = _context_manager.search(session_id, query, limit=limit)
    return _ok({"results": results, "count": len(results), "query": query})


def tool_context_detail(args: dict) -> str:
    """Expand a summary to see original messages."""
    if not _context_manager:
        return _err("Context manager not initialized")
    summary_id = args.get("summary_id", "")
    if not summary_id:
        return _err("Missing summary_id")
    detail = _context_manager.get_detail(summary_id)
    return _ok(detail) if "error" not in detail else _err(detail["error"])


def tool_context_recall(args: dict) -> str:
    """Deep recall from compacted conversation history."""
    if not _context_manager:
        return _err("Context manager not initialized")
    session_id = getattr(_thread_local, 'current_session_id', None) or ""
    if not session_id:
        return _err("No active session")
    query = args.get("query", "")
    if not query:
        return _err("Missing query")
    # Get API credentials from thread local or delegate globals
    model = getattr(_thread_local, 'current_model', None) or _delegate_fallback_model or ""
    api_key = _delegate_api_key or ""
    base_url = _delegate_base_url or ""
    api_type = _delegate_api_type or "anthropic"
    result = _context_manager.recall(session_id, query, model, api_key, base_url, api_type)
    return _ok({"answer": result, "query": query})


def tool_schedule_list(args: dict) -> str:
    """List all scheduled tasks."""
    if not _scheduler:
        return _err("Scheduler not initialized")
    schedules = _scheduler.list_all()
    return _ok({"schedules": schedules, "count": len(schedules)})


def tool_schedule_history(args: dict) -> str:
    """Get execution history for scheduled tasks."""
    if not _scheduler:
        return _err("Scheduler not initialized")
    name = args.get("name")
    limit = args.get("limit", 20)
    history = _scheduler.get_history(name, limit)
    # Truncate long results
    for h in history:
        if h.get("result") and len(h["result"]) > 500:
            h["result"] = h["result"][:500] + "..."
    return _ok({"history": history, "count": len(history)})


# --- MCP Client Tools ---

def tool_mcp_connect(args: dict) -> str:
    """Connect to an MCP server at runtime."""
    url = args.get("url", "")
    name = args.get("name", "")
    transport = args.get("transport", "sse")
    persist = args.get("persist", False)

    if not url or not name:
        return _err("Both 'url' and 'name' are required")

    # Use thread-local MCP manager if available, otherwise global
    mcp = getattr(_thread_local, 'mcp_manager', None) or _mcp_manager
    if not mcp:
        mcp = MCPManager()
        _thread_local.mcp_manager = mcp

    result = mcp.connect_runtime(url, name, transport)
    if result.get("error"):
        return _err(result["error"])

    # Persist to mcp.json if requested
    if persist:
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        agent_id = agent.agent_id if agent else "main"
        mcp_json_path = os.path.join(AGENTS_DIR, agent_id, "mcp.json")
        try:
            existing = {}
            if os.path.exists(mcp_json_path):
                with open(mcp_json_path, "r") as f:
                    existing = json.load(f)
            if transport == "stdio":
                parts = url.split()
                existing[name] = {"transport": "stdio", "command": parts[0], "args": parts[1:] if len(parts) > 1 else []}
            else:
                existing[name] = {"transport": "sse", "url": url}
            with open(mcp_json_path, "w") as f:
                json.dump(existing, f, indent=2)
            result["persisted"] = True
        except Exception as e:
            result["persist_error"] = str(e)

    return _ok(result)


def tool_mcp_disconnect(args: dict) -> str:
    """Disconnect from an MCP server."""
    name = args.get("name", "")
    if not name:
        return _err("'name' is required")

    mcp = getattr(_thread_local, 'mcp_manager', None) or _mcp_manager
    if not mcp:
        return _err("No MCP manager available")

    result = mcp.disconnect_runtime(name)
    if result.get("error"):
        return _err(result["error"])
    return _ok(result)


def tool_mcp_servers(args: dict) -> str:
    """List all connected MCP servers."""
    mcp = getattr(_thread_local, 'mcp_manager', None) or _mcp_manager
    if not mcp:
        return _ok({"servers": [], "count": 0})
    servers = mcp.list_servers()
    return _ok({"servers": servers, "count": len(servers)})


# --- Code Structure Graph ---

CODE_GRAPH_DB = os.path.join(AGENTS_DIR, "main", "code-graph.db")

_code_graph_db_lock = threading.Lock()
_code_graph_db_pool: dict[int, sqlite3.Connection] = {}

_EXT_TO_LANG = {
    ".py": "python", ".js": "javascript", ".ts": "typescript", ".tsx": "tsx",
    ".go": "go", ".rs": "rust", ".java": "java", ".c": "c", ".cpp": "cpp",
    ".h": "c", ".hpp": "cpp", ".cs": "c_sharp", ".rb": "ruby",
    ".kt": "kotlin", ".swift": "swift", ".php": "php",
}

_DEFAULT_EXCLUDE_DIRS = {"node_modules", ".git", "__pycache__", "venv", ".venv", "dist", "build", ".next", ".tox", "egg-info"}

# AST node type mappings per language
_CLASS_TYPES = {
    "python": {"class_definition"},
    "javascript": {"class_declaration"},
    "typescript": {"class_declaration"},
    "tsx": {"class_declaration"},
    "go": {"type_declaration"},
    "rust": {"struct_item", "enum_item", "impl_item"},
    "java": {"class_declaration", "interface_declaration"},
    "c": {"struct_specifier"},
    "cpp": {"class_specifier", "struct_specifier"},
    "c_sharp": {"class_declaration", "interface_declaration"},
    "ruby": {"class", "module"},
    "kotlin": {"class_declaration", "interface_declaration"},
    "swift": {"class_declaration", "struct_declaration", "protocol_declaration"},
    "php": {"class_declaration", "interface_declaration"},
}

_FUNCTION_TYPES = {
    "python": {"function_definition"},
    "javascript": {"function_declaration", "method_definition", "arrow_function"},
    "typescript": {"function_declaration", "method_definition", "arrow_function"},
    "tsx": {"function_declaration", "method_definition", "arrow_function"},
    "go": {"function_declaration", "method_declaration"},
    "rust": {"function_item"},
    "java": {"method_declaration", "constructor_declaration"},
    "c": {"function_definition"},
    "cpp": {"function_definition"},
    "c_sharp": {"method_declaration", "constructor_declaration"},
    "ruby": {"method", "singleton_method"},
    "kotlin": {"function_declaration"},
    "swift": {"function_declaration"},
    "php": {"function_definition", "method_declaration"},
}

_IMPORT_TYPES = {
    "python": {"import_statement", "import_from_statement"},
    "javascript": {"import_statement"},
    "typescript": {"import_statement"},
    "tsx": {"import_statement"},
    "go": {"import_declaration"},
    "rust": {"use_declaration"},
    "java": {"import_declaration"},
    "c": {"preproc_include"},
    "cpp": {"preproc_include"},
    "c_sharp": {"using_directive"},
    "ruby": {"call"},  # require/require_relative
    "kotlin": {"import_header"},
    "swift": {"import_declaration"},
    "php": {"namespace_use_declaration"},
}

_CALL_TYPES = {
    "python": {"call"},
    "javascript": {"call_expression"},
    "typescript": {"call_expression"},
    "tsx": {"call_expression"},
    "go": {"call_expression"},
    "rust": {"call_expression", "macro_invocation"},
    "java": {"method_invocation"},
    "c": {"call_expression"},
    "cpp": {"call_expression"},
    "c_sharp": {"invocation_expression"},
    "ruby": {"call"},
    "kotlin": {"call_expression"},
    "swift": {"call_expression"},
    "php": {"function_call_expression", "method_call_expression"},
}


def _code_graph_conn():
    """Thread-local SQLite connection for code graph DB."""
    tid = threading.current_thread().ident
    with _code_graph_db_lock:
        conn = _code_graph_db_pool.get(tid)
        if conn is None:
            os.makedirs(os.path.dirname(CODE_GRAPH_DB), exist_ok=True)
            conn = sqlite3.connect(CODE_GRAPH_DB, timeout=10, check_same_thread=False)
            conn.execute("PRAGMA busy_timeout = 5000")
            conn.execute("PRAGMA journal_mode = WAL")
            _code_graph_db_pool[tid] = conn
    return conn


def _code_graph_init_db():
    """Initialize the code graph schema."""
    conn = _code_graph_conn()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS code_nodes (
            qualified_name TEXT PRIMARY KEY,
            file_path TEXT NOT NULL,
            kind TEXT NOT NULL,
            name TEXT NOT NULL,
            language TEXT,
            line_start INTEGER,
            line_end INTEGER,
            parent_name TEXT,
            params TEXT,
            return_type TEXT,
            modifiers TEXT,
            file_hash TEXT,
            line_count INTEGER
        );
        CREATE TABLE IF NOT EXISTS code_edges (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            source TEXT NOT NULL,
            target TEXT NOT NULL,
            kind TEXT NOT NULL,
            file_path TEXT
        );
        CREATE TABLE IF NOT EXISTS code_meta (
            key TEXT PRIMARY KEY,
            value TEXT
        );
        CREATE INDEX IF NOT EXISTS idx_code_nodes_file ON code_nodes(file_path);
        CREATE INDEX IF NOT EXISTS idx_code_nodes_kind ON code_nodes(kind);
        CREATE INDEX IF NOT EXISTS idx_code_nodes_name ON code_nodes(name);
        CREATE INDEX IF NOT EXISTS idx_code_edges_source ON code_edges(source);
        CREATE INDEX IF NOT EXISTS idx_code_edges_target ON code_edges(target);
    """)
    # Migration: add summary and layer columns
    for col, default in [("summary", "''"), ("layer", "''")]:
        try:
            conn.execute(f"ALTER TABLE code_nodes ADD COLUMN {col} TEXT DEFAULT {default}")
        except Exception:
            pass
    conn.executescript("""
        CREATE INDEX IF NOT EXISTS idx_code_edges_kind ON code_edges(kind);
    """)
    conn.commit()


def _extract_node_name(node):
    """Extract the name from a tree-sitter AST node."""
    # Try field name "name" first
    name_node = node.child_by_field_name("name")
    if name_node:
        return name_node.text.decode("utf-8", errors="replace")
    # Fallback: look for first identifier child
    for child in node.children:
        if child.type == "identifier":
            return child.text.decode("utf-8", errors="replace")
        if child.type == "type_identifier":
            return child.text.decode("utf-8", errors="replace")
    return None


def _extract_call_name(node):
    """Extract the called function/method name from a call expression."""
    fn = node.child_by_field_name("function")
    if fn is None:
        fn = node.child_by_field_name("method")
    if fn is None:
        # Fallback: first child
        if node.children:
            fn = node.children[0]
    if fn is None:
        return None
    # Dotted: a.b.c -> take full text
    text = fn.text.decode("utf-8", errors="replace")
    # For method calls, extract just the method name
    if "." in text:
        return text.rsplit(".", 1)[-1]
    return text


def _extract_import_name(node, language):
    """Extract imported module name from an import node."""
    text = node.text.decode("utf-8", errors="replace")
    if language == "python":
        # import foo / from foo import bar
        if text.startswith("from "):
            parts = text.split()
            if len(parts) >= 2:
                return parts[1]
        elif text.startswith("import "):
            parts = text.split()
            if len(parts) >= 2:
                return parts[1].split(",")[0].strip()
    elif language in ("javascript", "typescript", "tsx"):
        # import ... from "module"
        src = node.child_by_field_name("source")
        if src:
            return src.text.decode("utf-8", errors="replace").strip("'\"")
    elif language == "go":
        # Walk children for interpreted_string_literal
        for child in node.children:
            if child.type == "import_spec_list":
                for spec in child.children:
                    if spec.type == "import_spec":
                        path_node = spec.child_by_field_name("path")
                        if path_node:
                            return path_node.text.decode("utf-8", errors="replace").strip('"')
            elif child.type == "import_spec":
                path_node = child.child_by_field_name("path")
                if path_node:
                    return path_node.text.decode("utf-8", errors="replace").strip('"')
            elif child.type == "interpreted_string_literal":
                return child.text.decode("utf-8", errors="replace").strip('"')
    elif language == "rust":
        # use foo::bar
        if text.startswith("use "):
            return text[4:].rstrip(";").strip()
    elif language == "java":
        # import foo.bar.Baz;
        if text.startswith("import "):
            return text[7:].rstrip(";").strip()
    elif language in ("c", "cpp"):
        # #include <foo> or #include "foo"
        path_node = node.child_by_field_name("path")
        if path_node:
            return path_node.text.decode("utf-8", errors="replace").strip('<>"')
    return text


def _is_test_function(name, file_path):
    """Detect if a function is a test by name or file location."""
    if not name:
        return False
    lower = name.lower()
    if lower.startswith("test_") or (lower.startswith("test") and len(name) > 4 and name[4].isupper()):
        return True
    basename = os.path.basename(file_path).lower()
    if basename.startswith("test_") or basename.endswith("_test.py") or basename.endswith(".test.js") or \
       basename.endswith(".test.ts") or basename.endswith(".test.tsx") or basename.endswith("_test.go") or \
       basename.endswith(".spec.js") or basename.endswith(".spec.ts"):
        return True
    return False


class CodeGraph:
    """AST-based code structure graph using Tree-sitter and SQLite."""

    def __init__(self):
        _code_graph_init_db()
        self._ts_available = None

    def _check_ts(self):
        """Lazy check for tree-sitter availability."""
        if self._ts_available is not None:
            return self._ts_available
        try:
            from tree_sitter_language_pack import get_parser  # noqa: F401
            self._ts_available = True
        except ImportError:
            self._ts_available = False
        return self._ts_available

    def parse_file(self, file_path: str) -> tuple[list[dict], list[dict]]:
        """Parse a single file with tree-sitter. Returns (nodes, edges)."""
        if not self._check_ts():
            return [], []

        ext = os.path.splitext(file_path)[1].lower()
        language = _EXT_TO_LANG.get(ext)
        if not language:
            return [], []

        try:
            from tree_sitter_language_pack import get_parser
        except ImportError:
            return [], []

        try:
            parser = get_parser(language)
        except Exception:
            return [], []

        try:
            with open(file_path, "rb") as f:
                source = f.read()
        except (OSError, IOError):
            return [], []

        try:
            tree = parser.parse(source)
        except Exception:
            return [], []

        nodes = []
        edges = []
        line_count = source.count(b"\n") + 1

        # Compute file hash
        file_hash = hashlib.sha256(source).hexdigest()

        # File node
        file_qn = file_path
        nodes.append({
            "qualified_name": file_qn,
            "file_path": file_path,
            "kind": "file",
            "name": os.path.basename(file_path),
            "language": language,
            "line_start": 1,
            "line_end": line_count,
            "parent_name": None,
            "params": None,
            "return_type": None,
            "modifiers": None,
            "file_hash": file_hash,
            "line_count": line_count,
        })

        class_types = _CLASS_TYPES.get(language, set())
        func_types = _FUNCTION_TYPES.get(language, set())
        import_types = _IMPORT_TYPES.get(language, set())
        call_types = _CALL_TYPES.get(language, set())

        # Walk AST
        def walk(node, parent_class=None):
            ntype = node.type

            if ntype in class_types:
                name = _extract_node_name(node)
                if name:
                    qn = f"{file_path}::{name}"
                    nodes.append({
                        "qualified_name": qn,
                        "file_path": file_path,
                        "kind": "class",
                        "name": name,
                        "language": language,
                        "line_start": node.start_point[0] + 1,
                        "line_end": node.end_point[0] + 1,
                        "parent_name": f"{file_path}::{parent_class}" if parent_class else file_qn,
                        "params": None,
                        "return_type": None,
                        "modifiers": None,
                        "file_hash": None,
                        "line_count": node.end_point[0] - node.start_point[0] + 1,
                    })
                    edges.append({
                        "source": file_qn,
                        "target": qn,
                        "kind": "CONTAINS",
                        "file_path": file_path,
                    })
                    # Check for inheritance
                    superclass = node.child_by_field_name("superclass")
                    if superclass is None:
                        superclass = node.child_by_field_name("argument_list")
                    if superclass is None:
                        superclass = node.child_by_field_name("superclasses")
                    if superclass:
                        sc_text = superclass.text.decode("utf-8", errors="replace").strip("()")
                        if sc_text and sc_text not in ("object", "Object"):
                            for sc in sc_text.split(","):
                                sc = sc.strip()
                                if sc:
                                    edges.append({
                                        "source": qn,
                                        "target": sc,
                                        "kind": "INHERITS",
                                        "file_path": file_path,
                                    })
                    # Recurse with this class as parent
                    for child in node.children:
                        walk(child, parent_class=name)
                    return

            elif ntype in func_types:
                name = _extract_node_name(node)
                if name:
                    if parent_class:
                        qn = f"{file_path}::{parent_class}.{name}"
                    else:
                        qn = f"{file_path}::{name}"

                    is_test = _is_test_function(name, file_path)
                    kind = "test" if is_test else "function"

                    # Extract params
                    params_node = node.child_by_field_name("parameters")
                    params_text = None
                    if params_node:
                        params_text = params_node.text.decode("utf-8", errors="replace")

                    # Extract return type
                    ret_node = node.child_by_field_name("return_type")
                    ret_text = None
                    if ret_node:
                        ret_text = ret_node.text.decode("utf-8", errors="replace")

                    nodes.append({
                        "qualified_name": qn,
                        "file_path": file_path,
                        "kind": kind,
                        "name": name,
                        "language": language,
                        "line_start": node.start_point[0] + 1,
                        "line_end": node.end_point[0] + 1,
                        "parent_name": f"{file_path}::{parent_class}" if parent_class else file_qn,
                        "params": params_text,
                        "return_type": ret_text,
                        "modifiers": None,
                        "file_hash": None,
                        "line_count": node.end_point[0] - node.start_point[0] + 1,
                    })
                    container = f"{file_path}::{parent_class}" if parent_class else file_qn
                    edges.append({
                        "source": container,
                        "target": qn,
                        "kind": "CONTAINS",
                        "file_path": file_path,
                    })

                    # If it's a test, try to link TESTED_BY
                    if is_test:
                        tested_name = name
                        if tested_name.startswith("test_"):
                            tested_name = tested_name[5:]
                        elif tested_name.startswith("Test"):
                            tested_name = tested_name[4:]
                            if tested_name:
                                tested_name = tested_name[0].lower() + tested_name[1:]
                        if tested_name and tested_name != name:
                            edges.append({
                                "source": tested_name,
                                "target": qn,
                                "kind": "TESTED_BY",
                                "file_path": file_path,
                            })

            elif ntype in import_types:
                import_name = _extract_import_name(node, language)
                if import_name:
                    edges.append({
                        "source": file_qn,
                        "target": import_name,
                        "kind": "IMPORTS_FROM",
                        "file_path": file_path,
                    })

            elif ntype in call_types:
                call_name = _extract_call_name(node)
                # Find enclosing function
                enclosing = None
                p = node.parent
                while p:
                    if p.type in func_types:
                        enc_name = _extract_node_name(p)
                        if enc_name:
                            pp = p.parent
                            enc_class = None
                            while pp:
                                if pp.type in class_types:
                                    enc_class = _extract_node_name(pp)
                                    break
                                pp = pp.parent
                            if enc_class:
                                enclosing = f"{file_path}::{enc_class}.{enc_name}"
                            else:
                                enclosing = f"{file_path}::{enc_name}"
                        break
                    p = p.parent
                if call_name and enclosing:
                    edges.append({
                        "source": enclosing,
                        "target": call_name,
                        "kind": "CALLS",
                        "file_path": file_path,
                    })

            # Recurse into children
            for child in node.children:
                walk(child, parent_class=parent_class)

        walk(tree.root_node)
        return nodes, edges

    def build(self, root_dir: str, incremental: bool = True, exclude_dirs: set | None = None) -> dict:
        """Parse all source files in directory and build the graph."""
        if not self._check_ts():
            return {"error": "tree-sitter-language-pack not installed. Run: pip install tree-sitter-language-pack"}

        root_dir = os.path.abspath(root_dir)
        if not os.path.isdir(root_dir):
            return {"error": f"Not a directory: {root_dir}"}

        exclude = exclude_dirs or _DEFAULT_EXCLUDE_DIRS
        conn = _code_graph_conn()

        # Get existing hashes for incremental mode
        existing_hashes = {}
        if incremental:
            try:
                rows = conn.execute(
                    "SELECT file_path, file_hash FROM code_nodes WHERE kind = 'file'"
                ).fetchall()
                existing_hashes = {r[0]: r[1] for r in rows}
            except Exception:
                pass

        stats = {"files_parsed": 0, "files_skipped": 0, "nodes": 0, "edges": 0, "languages": set()}
        all_nodes = []
        all_edges = []

        for dirpath, dirnames, filenames in os.walk(root_dir):
            # Skip excluded directories
            dirnames[:] = [d for d in dirnames if d not in exclude and not d.startswith(".")]

            for fname in filenames:
                ext = os.path.splitext(fname)[1].lower()
                if ext not in _EXT_TO_LANG:
                    continue

                fpath = os.path.join(dirpath, fname)

                # Skip large files (>500KB)
                try:
                    if os.path.getsize(fpath) > 500_000:
                        stats["files_skipped"] += 1
                        continue
                except OSError:
                    continue

                # Check hash for incremental
                if incremental and fpath in existing_hashes:
                    try:
                        with open(fpath, "rb") as f:
                            current_hash = hashlib.sha256(f.read()).hexdigest()
                        if current_hash == existing_hashes[fpath]:
                            stats["files_skipped"] += 1
                            continue
                    except OSError:
                        continue

                nodes, edges = self.parse_file(fpath)
                if nodes:
                    all_nodes.extend(nodes)
                    all_edges.extend(edges)
                    stats["files_parsed"] += 1
                    stats["languages"].add(_EXT_TO_LANG[ext])

        # Bulk insert into SQLite
        if all_nodes or not incremental:
            try:
                if not incremental:
                    conn.execute("DELETE FROM code_nodes")
                    conn.execute("DELETE FROM code_edges")
                else:
                    parsed_files = {n["file_path"] for n in all_nodes if n["kind"] == "file"}
                    for fp in parsed_files:
                        conn.execute("DELETE FROM code_nodes WHERE file_path = ?", (fp,))
                        conn.execute("DELETE FROM code_edges WHERE file_path = ?", (fp,))

                for n in all_nodes:
                    conn.execute(
                        "INSERT OR REPLACE INTO code_nodes "
                        "(qualified_name, file_path, kind, name, language, line_start, line_end, "
                        "parent_name, params, return_type, modifiers, file_hash, line_count) "
                        "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
                        (n["qualified_name"], n["file_path"], n["kind"], n["name"],
                         n["language"], n["line_start"], n["line_end"], n["parent_name"],
                         n["params"], n["return_type"], n["modifiers"], n["file_hash"],
                         n["line_count"])
                    )

                for e in all_edges:
                    conn.execute(
                        "INSERT INTO code_edges (source, target, kind, file_path) "
                        "VALUES (?, ?, ?, ?)",
                        (e["source"], e["target"], e["kind"], e["file_path"])
                    )

                conn.execute(
                    "INSERT OR REPLACE INTO code_meta (key, value) VALUES (?, ?)",
                    ("last_build", datetime.datetime.now(datetime.timezone.utc).isoformat())
                )
                conn.execute(
                    "INSERT OR REPLACE INTO code_meta (key, value) VALUES (?, ?)",
                    ("root_dir", root_dir)
                )
                conn.commit()
            except Exception as e:
                try:
                    conn.rollback()
                except Exception:
                    pass
                return {"error": f"Database error: {e}"}

        stats["nodes"] = len(all_nodes)
        stats["edges"] = len(all_edges)
        stats["languages"] = sorted(stats["languages"])

        try:
            total_nodes = conn.execute("SELECT COUNT(*) FROM code_nodes").fetchone()[0]
            total_edges = conn.execute("SELECT COUNT(*) FROM code_edges").fetchone()[0]
            stats["total_nodes"] = total_nodes
            stats["total_edges"] = total_edges
        except Exception:
            pass

        return stats

    def query(self, query_type: str, target: str, limit: int = 20) -> list[dict]:
        """Run predefined structural queries."""
        conn = _code_graph_conn()
        results = []

        try:
            if query_type == "callers_of":
                rows = conn.execute(
                    "SELECT e.source, e.file_path, n.kind, n.line_start, n.line_end "
                    "FROM code_edges e LEFT JOIN code_nodes n ON e.source = n.qualified_name "
                    "WHERE e.kind = 'CALLS' AND (e.target = ? OR e.target LIKE ?) LIMIT ?",
                    (target, f"%::{target}", limit)
                ).fetchall()
                for r in rows:
                    results.append({"caller": r[0], "file": r[1], "kind": r[2], "line_start": r[3], "line_end": r[4]})

            elif query_type == "callees_of":
                rows = conn.execute(
                    "SELECT e.target, e.file_path "
                    "FROM code_edges e "
                    "WHERE e.kind = 'CALLS' AND (e.source = ? OR e.source LIKE ?) LIMIT ?",
                    (target, f"%::{target}", limit)
                ).fetchall()
                for r in rows:
                    results.append({"callee": r[0], "file": r[1]})

            elif query_type == "imports_of":
                rows = conn.execute(
                    "SELECT e.target, e.file_path "
                    "FROM code_edges e "
                    "WHERE e.kind = 'IMPORTS_FROM' AND (e.source = ? OR e.source LIKE ?) LIMIT ?",
                    (target, f"%{target}%", limit)
                ).fetchall()
                for r in rows:
                    results.append({"imports": r[0], "file": r[1]})

            elif query_type == "importers_of":
                rows = conn.execute(
                    "SELECT e.source, e.file_path "
                    "FROM code_edges e "
                    "WHERE e.kind = 'IMPORTS_FROM' AND (e.target = ? OR e.target LIKE ?) LIMIT ?",
                    (target, f"%{target}%", limit)
                ).fetchall()
                for r in rows:
                    results.append({"importer": r[0], "file": r[1]})

            elif query_type == "tests_for":
                rows = conn.execute(
                    "SELECT e.target, e.file_path, n.line_start, n.line_end "
                    "FROM code_edges e LEFT JOIN code_nodes n ON e.target = n.qualified_name "
                    "WHERE e.kind = 'TESTED_BY' AND (e.source = ? OR e.source LIKE ?) LIMIT ?",
                    (target, f"%::{target}", limit)
                ).fetchall()
                for r in rows:
                    results.append({"test": r[0], "file": r[1], "line_start": r[2], "line_end": r[3]})

            elif query_type == "inheritors_of":
                rows = conn.execute(
                    "SELECT e.source, e.file_path, n.line_start, n.line_end "
                    "FROM code_edges e LEFT JOIN code_nodes n ON e.source = n.qualified_name "
                    "WHERE e.kind = 'INHERITS' AND (e.target = ? OR e.target LIKE ?) LIMIT ?",
                    (target, f"%::{target}", limit)
                ).fetchall()
                for r in rows:
                    results.append({"inheritor": r[0], "file": r[1], "line_start": r[2], "line_end": r[3]})

            elif query_type == "children_of":
                rows = conn.execute(
                    "SELECT e.target, n.kind, n.name, n.line_start, n.line_end, n.params "
                    "FROM code_edges e LEFT JOIN code_nodes n ON e.target = n.qualified_name "
                    "WHERE e.kind = 'CONTAINS' AND (e.source = ? OR e.source LIKE ?) LIMIT ?",
                    (target, f"%::{target}", limit)
                ).fetchall()
                for r in rows:
                    results.append({"child": r[0], "kind": r[1], "name": r[2], "line_start": r[3], "line_end": r[4], "params": r[5]})

            elif query_type == "file_summary":
                rows = conn.execute(
                    "SELECT qualified_name, kind, name, line_start, line_end, params, return_type "
                    "FROM code_nodes WHERE file_path = ? OR file_path LIKE ? ORDER BY line_start LIMIT ?",
                    (target, f"%{target}", limit)
                ).fetchall()
                for r in rows:
                    results.append({
                        "qualified_name": r[0], "kind": r[1], "name": r[2],
                        "line_start": r[3], "line_end": r[4], "params": r[5], "return_type": r[6],
                    })

        except Exception as e:
            return [{"error": str(e)}]

        return results

    def impact_analysis(self, files: list[str], depth: int = 2) -> dict:
        """BFS blast-radius analysis using networkx."""
        try:
            import networkx as nx
        except ImportError:
            return {"error": "networkx not installed. Run: pip install networkx"}

        conn = _code_graph_conn()

        G = nx.DiGraph()
        try:
            edges = conn.execute("SELECT source, target, kind FROM code_edges").fetchall()
            for src, tgt, kind in edges:
                G.add_edge(src, tgt, kind=kind)
                G.add_edge(tgt, src, kind=f"REV_{kind}")
        except Exception as e:
            return {"error": f"Database error: {e}"}

        changed_nodes = set()
        for fp in files:
            fp = os.path.abspath(fp)
            try:
                rows = conn.execute(
                    "SELECT qualified_name FROM code_nodes WHERE file_path = ?", (fp,)
                ).fetchall()
                for r in rows:
                    changed_nodes.add(r[0])
            except Exception:
                pass

        if not changed_nodes:
            return {
                "changed_nodes": [],
                "impacted_nodes": [],
                "impacted_files": [],
                "warnings": ["No nodes found for the specified files. Run code_graph_build first."],
            }

        impacted = set()
        frontier = set(changed_nodes)
        visited = set(changed_nodes)
        for _ in range(depth):
            next_frontier = set()
            for node in frontier:
                if node in G:
                    for neighbor in G.neighbors(node):
                        if neighbor not in visited:
                            visited.add(neighbor)
                            next_frontier.add(neighbor)
                            impacted.add(neighbor)
            frontier = next_frontier
            if not frontier:
                break

        impacted_files = set()
        for qn in impacted:
            try:
                row = conn.execute(
                    "SELECT file_path FROM code_nodes WHERE qualified_name = ?", (qn,)
                ).fetchone()
                if row:
                    impacted_files.add(row[0])
            except Exception:
                pass

        warnings = []
        for qn in changed_nodes:
            try:
                test_count = conn.execute(
                    "SELECT COUNT(*) FROM code_edges WHERE kind = 'TESTED_BY' AND source = ?", (qn,)
                ).fetchone()[0]
                if test_count == 0:
                    row = conn.execute(
                        "SELECT kind, name FROM code_nodes WHERE qualified_name = ?", (qn,)
                    ).fetchone()
                    if row and row[0] in ("function", "class"):
                        warnings.append(f"No tests found for {row[1]} ({qn})")
            except Exception:
                pass

        return {
            "changed_nodes": sorted(changed_nodes),
            "impacted_nodes": sorted(impacted),
            "impacted_files": sorted(impacted_files),
            "warnings": warnings,
        }

    def get_stats(self) -> dict:
        """Return node/edge counts and language distribution."""
        conn = _code_graph_conn()
        try:
            total_nodes = conn.execute("SELECT COUNT(*) FROM code_nodes").fetchone()[0]
            total_edges = conn.execute("SELECT COUNT(*) FROM code_edges").fetchone()[0]
            kind_dist = conn.execute(
                "SELECT kind, COUNT(*) FROM code_nodes GROUP BY kind ORDER BY COUNT(*) DESC"
            ).fetchall()
            lang_dist = conn.execute(
                "SELECT language, COUNT(*) FROM code_nodes WHERE language IS NOT NULL GROUP BY language ORDER BY COUNT(*) DESC"
            ).fetchall()
            edge_dist = conn.execute(
                "SELECT kind, COUNT(*) FROM code_edges GROUP BY kind ORDER BY COUNT(*) DESC"
            ).fetchall()
            last_build = conn.execute(
                "SELECT value FROM code_meta WHERE key = 'last_build'"
            ).fetchone()
            root_dir = conn.execute(
                "SELECT value FROM code_meta WHERE key = 'root_dir'"
            ).fetchone()
            return {
                "total_nodes": total_nodes,
                "total_edges": total_edges,
                "node_kinds": {k: c for k, c in kind_dist},
                "languages": {k: c for k, c in lang_dist},
                "edge_kinds": {k: c for k, c in edge_dist},
                "last_build": last_build[0] if last_build else None,
                "root_dir": root_dir[0] if root_dir else None,
            }
        except Exception as e:
            return {"error": str(e)}

    def update_file(self, file_path: str):
        """Incrementally re-parse a single file and update the graph."""
        file_path = os.path.abspath(file_path)
        ext = os.path.splitext(file_path)[1].lower()
        if ext not in _EXT_TO_LANG:
            return

        conn = _code_graph_conn()
        nodes, edges = self.parse_file(file_path)
        if not nodes:
            return

        try:
            conn.execute("DELETE FROM code_nodes WHERE file_path = ?", (file_path,))
            conn.execute("DELETE FROM code_edges WHERE file_path = ?", (file_path,))
            for n in nodes:
                conn.execute(
                    "INSERT OR REPLACE INTO code_nodes "
                    "(qualified_name, file_path, kind, name, language, line_start, line_end, "
                    "parent_name, params, return_type, modifiers, file_hash, line_count) "
                    "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
                    (n["qualified_name"], n["file_path"], n["kind"], n["name"],
                     n["language"], n["line_start"], n["line_end"], n["parent_name"],
                     n["params"], n["return_type"], n["modifiers"], n["file_hash"],
                     n["line_count"])
                )
            for e in edges:
                conn.execute(
                    "INSERT INTO code_edges (source, target, kind, file_path) "
                    "VALUES (?, ?, ?, ?)",
                    (e["source"], e["target"], e["kind"], e["file_path"])
                )
            conn.commit()
        except Exception:
            try:
                conn.rollback()
            except Exception:
                pass


    # --- Architecture Layer Classification ---

    _LAYER_PATTERNS = {
        "api": ["route", "router", "handler", "endpoint", "controller", "view", "api", "rest", "graphql", "grpc"],
        "service": ["service", "manager", "engine", "processor", "worker", "scheduler", "pipeline"],
        "data": ["model", "schema", "migration", "orm", "database", "db", "query", "repository", "dao", "store"],
        "ui": ["component", "page", "layout", "widget", "template", "style", "css", "html", "view", "screen", "form"],
        "util": ["util", "helper", "common", "shared", "lib", "config", "constant", "middleware", "decorator", "mixin"],
        "test": ["test", "spec", "fixture", "mock", "stub", "fake"],
    }

    def classify_layers(self) -> dict:
        """Classify all nodes into architecture layers based on file paths and names."""
        conn = _code_graph_conn()
        rows = conn.execute("SELECT qualified_name, file_path, name, kind FROM code_nodes").fetchall()
        classified = 0
        for qn, fp, name, kind in rows:
            layer = self._detect_layer(fp, name, kind)
            if layer:
                conn.execute("UPDATE code_nodes SET layer = ? WHERE qualified_name = ?", (layer, qn))
                classified += 1
        conn.commit()
        # Get distribution
        dist = conn.execute(
            "SELECT layer, COUNT(*) FROM code_nodes WHERE layer != '' GROUP BY layer ORDER BY COUNT(*) DESC"
        ).fetchall()
        return {"classified": classified, "layers": {k: c for k, c in dist}}

    def _detect_layer(self, file_path: str, name: str, kind: str) -> str:
        """Detect architecture layer from file path and node name."""
        fp_lower = file_path.lower().replace("\\", "/")
        name_lower = name.lower()
        combined = fp_lower + "/" + name_lower

        if kind == "test":
            return "test"

        for layer, patterns in self._LAYER_PATTERNS.items():
            for p in patterns:
                if p in combined:
                    return layer
        return ""

    # --- LLM-Generated Summaries ---

    def generate_summaries(self, batch_size: int = 20) -> dict:
        """Generate plain-English summaries for nodes that don't have one.
        Uses a cheap LLM (Haiku) to describe functions/classes from their signature + context."""
        conn = _code_graph_conn()
        # Get nodes without summaries (classes and functions only)
        rows = conn.execute(
            "SELECT qualified_name, file_path, kind, name, params, return_type, line_start, line_end, language "
            "FROM code_nodes WHERE (summary IS NULL OR summary = '') AND kind IN ('class', 'function') "
            "ORDER BY file_path, line_start LIMIT ?",
            (batch_size * 5,)  # over-fetch to batch by file
        ).fetchall()

        if not rows:
            return {"summarized": 0, "message": "All nodes already have summaries"}

        # Resolve summary model
        model = None
        if _models_config:
            for mid, cfg in _models_config.items():
                if cfg.get("enabled", True) and "haiku" in mid.lower():
                    model = mid
                    break
            if not model:
                for mid, cfg in sorted(_models_config.items(), key=lambda x: x[1].get("cost_input", 999)):
                    if cfg.get("enabled", True):
                        model = mid
                        break
        if not model:
            return {"error": "No model available for summary generation"}

        # Group by file for efficient source reading
        file_groups = {}
        for r in rows:
            fp = r[1]
            if fp not in file_groups:
                file_groups[fp] = []
            file_groups[fp].append(r)

        summarized = 0
        for fp, file_nodes in list(file_groups.items())[:batch_size]:
            # Read source file
            try:
                with open(fp, "r", errors="replace") as f:
                    lines = f.readlines()
            except (OSError, IOError):
                continue

            # Build batch prompt
            node_snippets = []
            for qn, _, kind, name, params, ret_type, ls, le, lang in file_nodes:
                snippet = "".join(lines[max(0, ls-1):min(len(lines), le)])[:500]
                sig = f"{kind} {name}"
                if params:
                    sig += f"({params})"
                if ret_type:
                    sig += f" -> {ret_type}"
                node_snippets.append({"qn": qn, "sig": sig, "snippet": snippet})

            if not node_snippets:
                continue

            prompt_parts = [f"**{s['sig']}**\n```\n{s['snippet']}\n```" for s in node_snippets]
            prompt = (
                f"For each function/class below from `{os.path.basename(fp)}`, write a ONE-LINE summary "
                f"(max 80 chars) describing what it does. Output as numbered list matching the order.\n\n"
                + "\n\n".join(f"{i+1}. {p}" for i, p in enumerate(prompt_parts))
            )

            try:
                result = _run_delegate(
                    messages=[{"role": "user", "content": prompt}],
                    model=model,
                    system_prompt="Output only numbered one-line summaries. No markdown, no explanations.",
                    memory_store=None,
                    inference_params={"max_tokens": 2000, "temperature": 0.1},
                    tools=False,
                )
                if result and not result.startswith("Delegation error"):
                    # Parse numbered lines
                    summary_lines = []
                    for line in result.strip().split("\n"):
                        line = line.strip()
                        if line and line[0].isdigit():
                            # Strip number prefix
                            parts = line.split(".", 1)
                            if len(parts) == 2:
                                summary_lines.append(parts[1].strip())
                            else:
                                summary_lines.append(line)

                    for i, s in enumerate(node_snippets):
                        summary = summary_lines[i] if i < len(summary_lines) else ""
                        if summary:
                            conn.execute("UPDATE code_nodes SET summary = ? WHERE qualified_name = ?",
                                        (summary[:200], s["qn"]))
                            summarized += 1
            except Exception:
                continue

        conn.commit()
        return {"summarized": summarized, "total_pending": len(rows) - summarized}

    # --- Guided Tour Generation ---

    def generate_tour(self, root_dir: str | None = None) -> str:
        """Generate a dependency-ordered guided tour of the codebase."""
        conn = _code_graph_conn()
        if not root_dir:
            row = conn.execute("SELECT value FROM code_meta WHERE key = 'root_dir'").fetchone()
            root_dir = row[0] if row else "."

        # Get all file nodes with summaries and layers
        files = conn.execute(
            "SELECT qualified_name, name, language, line_count, layer, summary "
            "FROM code_nodes WHERE kind = 'file' ORDER BY qualified_name"
        ).fetchall()

        if not files:
            return "No files in the code graph. Run code_graph_build first."

        # Build dependency order: files with fewest imports first
        file_imports = {}
        for f in files:
            imports = conn.execute(
                "SELECT COUNT(*) FROM code_edges WHERE source = ? AND kind = 'IMPORTS_FROM'",
                (f[0],)
            ).fetchone()[0]
            importers = conn.execute(
                "SELECT COUNT(*) FROM code_edges WHERE target LIKE ? AND kind = 'IMPORTS_FROM'",
                (f"%%{f[1]}%%",)
            ).fetchone()[0]
            file_imports[f[0]] = {"imports": imports, "importers": importers}

        # Sort: foundation files first (many importers, few imports), then leaves
        sorted_files = sorted(files, key=lambda f: (
            file_imports.get(f[0], {}).get("imports", 0) - file_imports.get(f[0], {}).get("importers", 0),
            f[0]
        ))

        # Group by layer
        layer_groups = {}
        for f in sorted_files:
            layer = f[4] or "other"
            if layer not in layer_groups:
                layer_groups[layer] = []
            layer_groups[layer].append(f)

        # Build tour markdown
        tour = f"# Codebase Tour: {os.path.basename(root_dir)}\n\n"
        tour += f"**{len(files)} files** across {len(set(f[2] for f in files if f[2]))} language(s)\n\n"

        # Architecture overview
        if layer_groups:
            tour += "## Architecture Layers\n\n"
            layer_order = ["api", "service", "data", "ui", "util", "test", "other"]
            layer_emoji = {"api": "🌐", "service": "⚙️", "data": "💾", "ui": "🖥️", "util": "🔧", "test": "🧪", "other": "📁"}
            for layer in layer_order:
                if layer in layer_groups:
                    files_in_layer = layer_groups[layer]
                    tour += f"### {layer_emoji.get(layer, '📁')} {layer.upper()} ({len(files_in_layer)} files)\n\n"
                    for f in files_in_layer[:10]:
                        qn, name, lang, lc, _, summary = f
                        rel = os.path.relpath(qn, root_dir) if root_dir else qn
                        imp_count = file_imports.get(qn, {}).get("importers", 0)
                        desc = f" — {summary}" if summary else ""
                        tour += f"- `{rel}` ({lang}, {lc} lines, {imp_count} dependents){desc}\n"
                    if len(files_in_layer) > 10:
                        tour += f"- ... and {len(files_in_layer) - 10} more\n"
                    tour += "\n"

        # Entry points
        tour += "## Suggested Reading Order\n\n"
        tour += "Start with the foundation (most imported) files, then work outward:\n\n"
        for i, f in enumerate(sorted_files[:15], 1):
            qn, name, lang, lc, layer, summary = f
            rel = os.path.relpath(qn, root_dir) if root_dir else qn
            imp = file_imports.get(qn, {})
            desc = f" — {summary}" if summary else ""
            tour += f"{i}. `{rel}` ({imp.get('importers', 0)} dependents, {imp.get('imports', 0)} imports){desc}\n"

        # Get key classes
        key_classes = conn.execute(
            "SELECT n.qualified_name, n.name, n.file_path, n.summary, COUNT(e.id) as edge_count "
            "FROM code_nodes n LEFT JOIN code_edges e ON n.qualified_name = e.source OR n.qualified_name = e.target "
            "WHERE n.kind = 'class' GROUP BY n.qualified_name ORDER BY edge_count DESC LIMIT 10"
        ).fetchall()
        if key_classes:
            tour += "\n## Key Classes\n\n"
            for qn, name, fp, summary, edges in key_classes:
                rel = os.path.relpath(fp, root_dir) if root_dir else fp
                desc = f" — {summary}" if summary else ""
                tour += f"- **{name}** in `{rel}` ({edges} connections){desc}\n"

        return tour


_code_graph: CodeGraph | None = None


def _get_code_graph() -> CodeGraph:
    """Get or create the global CodeGraph instance."""
    global _code_graph
    if _code_graph is None:
        _code_graph = CodeGraph()
    return _code_graph


def _maybe_update_code_graph(path: str):
    """Update code graph if the file is a supported source file."""
    ext = os.path.splitext(path)[1].lower()
    if ext not in _EXT_TO_LANG:
        return
    try:
        cg = _get_code_graph()
        if cg._check_ts():
            cg.update_file(path)
    except Exception:
        pass


def tool_code_graph_build(args: dict) -> str:
    """Build or rebuild the code structure graph."""
    path = args.get("path", ".")
    incremental = args.get("incremental", True)
    cg = _get_code_graph()
    stats = cg.build(path, incremental=incremental)
    if "error" in stats:
        return _err(stats["error"])
    return _ok(stats)


def tool_code_graph_query(args: dict) -> str:
    """Query the code structure graph."""
    query_type = args.get("query_type", "")
    target = args.get("target", "")
    limit = args.get("limit", 20)
    if not query_type:
        return _err("Missing query_type")
    if not target:
        return _err("Missing target")
    cg = _get_code_graph()
    results = cg.query(query_type, target, limit)
    return _ok({"query_type": query_type, "target": target, "results": results, "count": len(results)})


def tool_code_graph_impact(args: dict) -> str:
    """Blast-radius impact analysis."""
    files = args.get("files", [])
    depth = args.get("depth", 2)
    if not files:
        return _err("Missing files list")
    cg = _get_code_graph()
    result = cg.impact_analysis(files, depth=depth)
    if "error" in result:
        return _err(result["error"])
    return _ok(result)


def tool_code_graph_enhance(args: dict) -> str:
    """Enhance the code graph with LLM summaries, architecture layers, and guided tour."""
    action = args.get("action", "all")
    cg = _get_code_graph()

    result = {}
    if action in ("all", "layers"):
        result["layers"] = cg.classify_layers()
    if action in ("all", "summaries"):
        batch_size = args.get("batch_size", 20)
        result["summaries"] = cg.generate_summaries(batch_size=batch_size)
    if action in ("all", "tour"):
        root_dir = args.get("root_dir")
        result["tour"] = cg.generate_tour(root_dir=root_dir)
    return _ok(result)


# --- Git / GitHub Tools ---

def _run_git(cmd_args: list[str], cwd: str | None = None, timeout: int = 30) -> tuple[int, str]:
    """Run a git command and return (exit_code, output)."""
    try:
        env = os.environ.copy()
        env["GIT_TERMINAL_PROMPT"] = "0"
        env["PAGER"] = "cat"
        proc = subprocess.run(
            ["git", "--no-pager"] + cmd_args,
            capture_output=True, cwd=cwd, env=env, timeout=timeout,
        )
        output = proc.stdout.decode("utf-8", errors="replace")
        if proc.stderr:
            err = proc.stderr.decode("utf-8", errors="replace").strip()
            if err and proc.returncode != 0:
                output += f"\n{err}" if output else err
        return proc.returncode, output.strip()
    except subprocess.TimeoutExpired:
        return -1, f"Git command timed out after {timeout}s"
    except FileNotFoundError:
        return -1, "git not found — install git first"
    except Exception as e:
        return -1, str(e)


def tool_git_command(args: dict) -> str:
    """Execute git operations with structured output."""
    action = args.get("action", "")
    if not action:
        return _err("Missing action")

    if action == "status":
        code, out = _run_git(["status", "--porcelain", "-b"])
        if code != 0:
            return _err(out)
        lines = out.split("\n")
        branch = lines[0][3:] if lines and lines[0].startswith("## ") else ""
        files = {"modified": [], "staged": [], "untracked": []}
        for l in lines[1:]:
            if not l.strip():
                continue
            idx, wt = l[0], l[1]
            fname = l[3:]
            if idx in ("M", "A", "D", "R"):
                files["staged"].append(fname)
            if wt == "M":
                files["modified"].append(fname)
            elif wt == "?":
                files["untracked"].append(fname)
        return _ok({"branch": branch, **files})

    elif action == "diff":
        cmd = ["diff"]
        if args.get("staged"):
            cmd.append("--staged")
        if args.get("ref"):
            cmd.append(args["ref"])
        if args.get("path"):
            cmd.extend(["--", args["path"]])
        code, out = _run_git(cmd)
        if code != 0:
            return _err(out)
        # Truncate large diffs
        if len(out) > 30000:
            out = out[:30000] + "\n... (diff truncated)"
        return _ok({"diff": out})

    elif action == "log":
        limit = args.get("limit", 20)
        cmd = ["log", f"--max-count={limit}", "--format=%H|%an|%ae|%ai|%s"]
        if args.get("author"):
            cmd.append(f"--author={args['author']}")
        if args.get("since"):
            cmd.append(f"--since={args['since']}")
        if args.get("path"):
            cmd.extend(["--", args["path"]])
        code, out = _run_git(cmd)
        if code != 0:
            return _err(out)
        commits = []
        for line in out.split("\n"):
            if "|" in line:
                parts = line.split("|", 4)
                if len(parts) == 5:
                    commits.append({"hash": parts[0], "author": parts[1], "email": parts[2], "date": parts[3], "message": parts[4]})
        return _ok({"commits": commits, "count": len(commits)})

    elif action == "branch":
        if args.get("create") and args.get("name"):
            code, out = _run_git(["checkout", "-b", args["name"]])
            return _ok({"created": args["name"]}) if code == 0 else _err(out)
        elif args.get("switch") and args.get("name"):
            code, out = _run_git(["checkout", args["name"]])
            return _ok({"switched": args["name"]}) if code == 0 else _err(out)
        else:
            code, out = _run_git(["branch", "-a", "--format=%(refname:short)|%(objectname:short)|%(upstream:short)"])
            if code != 0:
                return _err(out)
            branches = []
            for line in out.split("\n"):
                if "|" in line:
                    parts = line.split("|")
                    branches.append({"name": parts[0], "commit": parts[1] if len(parts) > 1 else "", "upstream": parts[2] if len(parts) > 2 else ""})
            # Get current branch
            _, current = _run_git(["rev-parse", "--abbrev-ref", "HEAD"])
            return _ok({"branches": branches, "current": current.strip()})

    elif action == "commit":
        message = args.get("message", "")
        if not message:
            return _err("Missing commit message")
        files = args.get("files", [])
        if files:
            code, out = _run_git(["add"] + files)
            if code != 0:
                return _err(f"Failed to stage: {out}")
        if args.get("all"):
            code, out = _run_git(["commit", "-a", "-m", message])
        else:
            code, out = _run_git(["commit", "-m", message])
        if code != 0:
            return _err(out)
        # Get the new commit hash
        _, hash_out = _run_git(["rev-parse", "--short", "HEAD"])
        return _ok({"committed": hash_out.strip(), "message": message})

    elif action == "stash":
        sub = args.get("sub_action", "save")
        if sub == "save":
            msg = args.get("message", "")
            cmd = ["stash", "push"]
            if msg:
                cmd.extend(["-m", msg])
            code, out = _run_git(cmd)
        elif sub == "pop":
            code, out = _run_git(["stash", "pop"])
        elif sub == "list":
            code, out = _run_git(["stash", "list"])
        elif sub == "drop":
            code, out = _run_git(["stash", "drop"])
        else:
            return _err(f"Unknown stash sub_action: {sub}")
        return _ok({"output": out}) if code == 0 else _err(out)

    elif action == "blame":
        path = args.get("path", "")
        if not path:
            return _err("Missing path for blame")
        cmd = ["blame", "--porcelain"]
        if args.get("line_start") and args.get("line_end"):
            cmd.extend([f"-L{args['line_start']},{args['line_end']}"])
        cmd.append(path)
        code, out = _run_git(cmd)
        if code != 0:
            return _err(out)
        if len(out) > 20000:
            out = out[:20000] + "\n... (truncated)"
        return _ok({"blame": out})

    elif action == "show":
        ref = args.get("ref", "HEAD")
        code, out = _run_git(["show", "--stat", ref])
        if code != 0:
            return _err(out)
        if len(out) > 20000:
            out = out[:20000] + "\n... (truncated)"
        return _ok({"output": out})

    elif action == "tag":
        if args.get("create") and args.get("name"):
            cmd = ["tag"]
            if args.get("message"):
                cmd.extend(["-a", args["name"], "-m", args["message"]])
            else:
                cmd.append(args["name"])
            code, out = _run_git(cmd)
            return _ok({"tagged": args["name"]}) if code == 0 else _err(out)
        else:
            code, out = _run_git(["tag", "-l", "--sort=-creatordate", "--format=%(refname:short)|%(creatordate:short)|%(subject)"])
            if code != 0:
                return _err(out)
            tags = []
            for line in out.split("\n"):
                if "|" in line:
                    parts = line.split("|", 2)
                    tags.append({"name": parts[0], "date": parts[1] if len(parts) > 1 else "", "message": parts[2] if len(parts) > 2 else ""})
            return _ok({"tags": tags})

    elif action == "remote":
        code, out = _run_git(["remote", "-v"])
        if code != 0:
            return _err(out)
        remotes = {}
        for line in out.split("\n"):
            parts = line.split()
            if len(parts) >= 2:
                remotes[parts[0]] = parts[1]
        return _ok({"remotes": remotes})

    return _err(f"Unknown git action: {action}")


def _run_gh(cmd_args: list[str], timeout: int = 30) -> tuple[int, str]:
    """Run a gh CLI command and return (exit_code, output)."""
    try:
        proc = subprocess.run(
            ["gh"] + cmd_args,
            capture_output=True, timeout=timeout,
        )
        output = proc.stdout.decode("utf-8", errors="replace")
        if proc.stderr:
            err = proc.stderr.decode("utf-8", errors="replace").strip()
            if err and proc.returncode != 0:
                output += f"\n{err}" if output else err
        return proc.returncode, output.strip()
    except FileNotFoundError:
        return -1, "gh CLI not found — install with: brew install gh"
    except subprocess.TimeoutExpired:
        return -1, f"GitHub CLI timed out after {timeout}s"
    except Exception as e:
        return -1, str(e)


def tool_github_command(args: dict) -> str:
    """Interact with GitHub via gh CLI."""
    action = args.get("action", "")
    if not action:
        return _err("Missing action")

    limit = args.get("limit", 20)

    if action == "pr_list":
        cmd = ["pr", "list", "--limit", str(limit), "--json", "number,title,author,state,headRefName,baseRefName,createdAt,url"]
        if args.get("state"):
            cmd.extend(["--state", args["state"]])
        if args.get("author"):
            cmd.extend(["--author", args["author"]])
        code, out = _run_gh(cmd)
        if code != 0:
            return _err(out)
        try:
            return _ok({"prs": json.loads(out)})
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "pr_create":
        title = args.get("title", "")
        if not title:
            return _err("Missing PR title")
        cmd = ["pr", "create", "--title", title]
        if args.get("body"):
            cmd.extend(["--body", args["body"]])
        if args.get("base"):
            cmd.extend(["--base", args["base"]])
        if args.get("head"):
            cmd.extend(["--head", args["head"]])
        if args.get("draft"):
            cmd.append("--draft")
        code, out = _run_gh(cmd)
        return _ok({"url": out}) if code == 0 else _err(out)

    elif action == "pr_view":
        number = args.get("number")
        if not number:
            return _err("Missing PR number")
        code, out = _run_gh(["pr", "view", str(number), "--json",
                             "number,title,body,state,author,headRefName,baseRefName,additions,deletions,files,reviews,comments,url"])
        if code != 0:
            return _err(out)
        try:
            return _ok(json.loads(out))
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "pr_merge":
        number = args.get("number")
        if not number:
            return _err("Missing PR number")
        method = args.get("method", "merge")
        cmd = ["pr", "merge", str(number), f"--{method}"]
        code, out = _run_gh(cmd)
        return _ok({"merged": number, "method": method}) if code == 0 else _err(out)

    elif action == "pr_review":
        number = args.get("number")
        if not number:
            return _err("Missing PR number")
        code, out = _run_gh(["pr", "view", str(number), "--json", "reviews,comments"])
        if code != 0:
            return _err(out)
        try:
            return _ok(json.loads(out))
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "issue_list":
        cmd = ["issue", "list", "--limit", str(limit), "--json", "number,title,author,state,labels,createdAt,url"]
        if args.get("state"):
            cmd.extend(["--state", args["state"]])
        if args.get("labels"):
            cmd.extend(["--label", args["labels"]])
        code, out = _run_gh(cmd)
        if code != 0:
            return _err(out)
        try:
            return _ok({"issues": json.loads(out)})
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "issue_create":
        title = args.get("title", "")
        if not title:
            return _err("Missing issue title")
        cmd = ["issue", "create", "--title", title]
        if args.get("body"):
            cmd.extend(["--body", args["body"]])
        if args.get("labels"):
            cmd.extend(["--label", args["labels"]])
        code, out = _run_gh(cmd)
        return _ok({"url": out}) if code == 0 else _err(out)

    elif action == "issue_view":
        number = args.get("number")
        if not number:
            return _err("Missing issue number")
        code, out = _run_gh(["issue", "view", str(number), "--json", "number,title,body,state,author,labels,comments,url"])
        if code != 0:
            return _err(out)
        try:
            return _ok(json.loads(out))
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "repo_view":
        code, out = _run_gh(["repo", "view", "--json", "name,description,url,defaultBranchRef,stargazerCount,forkCount,isPrivate,languages"])
        if code != 0:
            return _err(out)
        try:
            return _ok(json.loads(out))
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "release_list":
        code, out = _run_gh(["release", "list", "--limit", str(limit)])
        if code != 0:
            return _err(out)
        return _ok({"output": out})

    elif action == "workflow_list":
        code, out = _run_gh(["workflow", "list", "--json", "name,state,id"])
        if code != 0:
            return _err(out)
        try:
            return _ok({"workflows": json.loads(out)})
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "workflow_run":
        run_id = args.get("run_id", "")
        if not run_id:
            return _err("Missing run_id")
        code, out = _run_gh(["run", "view", run_id, "--json", "status,conclusion,name,createdAt,updatedAt,jobs"])
        if code != 0:
            return _err(out)
        try:
            return _ok(json.loads(out))
        except json.JSONDecodeError:
            return _ok({"output": out})

    elif action == "api":
        endpoint = args.get("endpoint", "")
        if not endpoint:
            return _err("Missing API endpoint")
        method = args.get("api_method", "GET")
        cmd = ["api", endpoint, "--method", method]
        code, out = _run_gh(cmd, timeout=30)
        if code != 0:
            return _err(out)
        try:
            return _ok(json.loads(out))
        except json.JSONDecodeError:
            return _ok({"output": out[:20000]})

    return _err(f"Unknown GitHub action: {action}")


# --- Context Window Management ---

DEFAULT_MAX_CONTEXT_TOKENS = 131072
COMPACT_THRESHOLD = 0.75  # compact at 75% full
KEEP_RECENT_MESSAGES = 6  # always keep the last N messages untouched (legacy fallback)
CHARS_PER_TOKEN = 4       # conservative estimate

# --- Lossless Context Manager ---

CONTEXT_DB = os.path.join(AGENTS_DIR, "main", "context.db")

_CONTEXT_CONFIG_DEFAULTS = {
    "enabled": True,
    "fresh_tail_count": 16,
    "compact_threshold": 0.60,
    "summary_target_tokens": 1000,
    "condense_threshold": 4,
    "max_depth": 5,
    "summary_model": "gemini-2.5-flash",
    "summary_model_fallback": "claude-haiku-4-5-20251001",
    "messages_per_summary": 10,
}

_context_db_lock = threading.Lock()
_context_db_pool: dict[int, sqlite3.Connection] = {}


def _context_conn():
    tid = threading.current_thread().ident
    with _context_db_lock:
        conn = _context_db_pool.get(tid)
        if conn is None:
            os.makedirs(os.path.dirname(CONTEXT_DB), exist_ok=True)
            conn = sqlite3.connect(CONTEXT_DB, timeout=10, check_same_thread=False)
            conn.execute("PRAGMA busy_timeout = 5000")
            conn.execute("PRAGMA journal_mode = WAL")
            _context_db_pool[tid] = conn
    return conn


def _context_init_db():
    conn = _context_conn()
    conn.executescript("""
        CREATE TABLE IF NOT EXISTS summaries (
            id TEXT PRIMARY KEY,
            session_id TEXT NOT NULL,
            depth INTEGER NOT NULL DEFAULT 0,
            token_count INTEGER NOT NULL,
            content TEXT NOT NULL,
            parent_ids TEXT DEFAULT '[]',
            message_range_start INTEGER NOT NULL,
            message_range_end INTEGER NOT NULL,
            created_at REAL DEFAULT (strftime('%s','now'))
        );
        CREATE INDEX IF NOT EXISTS idx_summary_session ON summaries(session_id);
        CREATE INDEX IF NOT EXISTS idx_summary_depth ON summaries(session_id, depth);
        CREATE TABLE IF NOT EXISTS context_config (
            key TEXT PRIMARY KEY,
            value TEXT NOT NULL
        );
    """)
    conn.commit()


class ContextManager:
    """Lossless context management with DAG-based hierarchical summarization."""

    def __init__(self):
        _context_init_db()

    def get_config(self) -> dict:
        cfg = dict(_CONTEXT_CONFIG_DEFAULTS)
        try:
            conn = _context_conn()
            rows = conn.execute("SELECT key, value FROM context_config").fetchall()
            for k, v in rows:
                if k in cfg:
                    # Coerce types
                    default = _CONTEXT_CONFIG_DEFAULTS[k]
                    if isinstance(default, bool):
                        cfg[k] = v.lower() in ("true", "1", "yes")
                    elif isinstance(default, int):
                        cfg[k] = int(v)
                    elif isinstance(default, float):
                        cfg[k] = float(v)
                    else:
                        cfg[k] = v
        except Exception:
            pass
        return cfg

    def save_config(self, updates: dict):
        conn = _context_conn()
        for k, v in updates.items():
            if k in _CONTEXT_CONFIG_DEFAULTS:
                conn.execute(
                    "INSERT OR REPLACE INTO context_config (key, value) VALUES (?, ?)",
                    (k, str(v))
                )
        conn.commit()

    def _resolve_summary_model(self) -> tuple[str | None, str | None]:
        """Return (primary_model, fallback_model) for context summarization."""
        cfg = self.get_config()
        primary = cfg.get("summary_model") or "gemini-2.5-flash"
        fallback = cfg.get("summary_model_fallback") or "claude-haiku-4-5-20251001"
        p, f = _resolve_model_with_fallback(primary, fallback, "claude-haiku-4-5-20251001")
        return p, f

    def _extract_message_text(self, msg: dict) -> str:
        """Extract plain text from a message (handles string and block content)."""
        role = msg.get("role", "unknown")
        content = msg.get("content", "")
        if isinstance(content, str):
            return f"{role}: {content}"
        elif isinstance(content, list):
            parts = []
            for block in content:
                if isinstance(block, dict):
                    if block.get("type") == "text":
                        parts.append(block.get("text", ""))
                    elif block.get("type") == "tool_use":
                        parts.append(f"[Tool: {block.get('name', '')}]")
                    elif block.get("type") == "tool_result":
                        c = block.get("content", "")
                        parts.append(f"[Result: {str(c)[:300]}]")
            return f"{role}: {' '.join(parts)}"
        return f"{role}: {str(content)[:500]}"

    def summarize_chunk(self, messages: list[dict], session_id: str,
                        range_start: int, range_end: int,
                        model: str, api_key: str, base_url: str, api_type: str,
                        fallback_model: str | None = None) -> str | None:
        """Summarize a chunk of messages into a leaf summary (depth 0). Returns summary ID."""
        text_parts = [self._extract_message_text(m) for m in messages]
        source_text = "\n".join(text_parts)

        summary_text = None
        try:
            result = _run_delegate_with_fallback(
                messages=[{"role": "user", "content": (
                    "Summarize this conversation segment concisely. Preserve: key facts, decisions, "
                    "file paths, commands run, errors encountered, and task context. "
                    "Be factual and specific, not vague. ~200-300 words.\n\n" + source_text
                )}],
                primary_model=model, fallback_model=fallback_model,
                system_prompt="You are a precise conversation summarizer. Output only the summary.",
                memory_store=None,
                inference_params={"max_tokens": 2000, "temperature": 0.1},
                tools=False,
            )
            if result and not result.startswith("Delegation error"):
                summary_text = result.strip()
        except Exception:
            pass

        if not summary_text:
            # Fallback: truncation
            summary_text = source_text[:3000]
            if len(source_text) > 3000:
                summary_text += "\n...(truncated)"

        sid = hashlib.sha256(f"{session_id}:{range_start}:{range_end}:{time.time()}".encode()).hexdigest()[:16]
        token_count = len(summary_text) // CHARS_PER_TOKEN

        conn = _context_conn()
        conn.execute(
            "INSERT OR REPLACE INTO summaries (id, session_id, depth, token_count, content, parent_ids, message_range_start, message_range_end) "
            "VALUES (?, ?, 0, ?, ?, '[]', ?, ?)",
            (sid, session_id, token_count, summary_text, range_start, range_end)
        )
        conn.commit()
        return sid

    def condense(self, session_id: str) -> bool:
        """Merge same-depth summaries into higher-depth summaries. Returns True if any merged."""
        cfg = self.get_config()
        threshold = cfg.get("condense_threshold", 4)
        max_depth = cfg.get("max_depth", 5)
        conn = _context_conn()
        merged_any = False

        for depth in range(max_depth):
            rows = conn.execute(
                "SELECT id, content, token_count, message_range_start, message_range_end "
                "FROM summaries WHERE session_id = ? AND depth = ? ORDER BY message_range_start",
                (session_id, depth)
            ).fetchall()
            if len(rows) < threshold:
                continue

            # Merge all same-depth summaries into one higher-depth summary
            combined_text = "\n\n---\n\n".join(r[1] for r in rows)
            parent_ids = [r[0] for r in rows]
            range_start = rows[0][3]
            range_end = rows[-1][4]

            # Summarize the combined text
            c_model, c_fallback = self._resolve_summary_model()
            condensed_text = None
            if c_model:
                try:
                    result = _run_delegate_with_fallback(
                        messages=[{"role": "user", "content": (
                            f"Condense these {len(rows)} conversation summaries into one cohesive summary. "
                            "Preserve all key facts, decisions, and context. ~300-500 words.\n\n" + combined_text
                        )}],
                        primary_model=c_model, fallback_model=c_fallback,
                        system_prompt="You are a precise summarizer. Output only the condensed summary.",
                        memory_store=None,
                        inference_params={"max_tokens": 3000, "temperature": 0.1},
                        tools=False,
                    )
                    if result and not result.startswith("Delegation error"):
                        condensed_text = result.strip()
                except Exception:
                    pass

            if not condensed_text:
                condensed_text = combined_text[:4000]

            new_id = hashlib.sha256(f"condense:{session_id}:{depth+1}:{time.time()}".encode()).hexdigest()[:16]
            token_count = len(condensed_text) // CHARS_PER_TOKEN

            conn.execute(
                "INSERT OR REPLACE INTO summaries (id, session_id, depth, token_count, content, parent_ids, message_range_start, message_range_end) "
                "VALUES (?, ?, ?, ?, ?, ?, ?, ?)",
                (new_id, session_id, depth + 1, token_count, condensed_text,
                 json.dumps(parent_ids), range_start, range_end)
            )
            # Remove merged lower-depth summaries
            conn.executemany("DELETE FROM summaries WHERE id = ?", [(pid,) for pid in parent_ids])
            conn.commit()
            merged_any = True

        return merged_any

    def assemble_context(self, session_id: str, messages: list[dict],
                         system_prompt: str = "", max_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS) -> list[dict]:
        """Assemble context: summaries + fresh tail within token budget."""
        cfg = self.get_config()
        fresh_tail_count = cfg.get("fresh_tail_count", 32)

        # Ensure we don't exceed available messages
        fresh_tail_count = min(fresh_tail_count, len(messages))
        fresh_tail = messages[-fresh_tail_count:] if fresh_tail_count > 0 else messages

        # Calculate token budget
        system_tokens = len(system_prompt) // CHARS_PER_TOKEN if system_prompt else 0
        fresh_tokens = sum(_estimate_tokens_message(m) for m in fresh_tail)
        response_reserve = 4000  # reserve tokens for response
        budget = max_tokens - system_tokens - fresh_tokens - response_reserve

        if budget <= 0:
            return fresh_tail

        # Load summaries, highest depth first
        conn = _context_conn()
        rows = conn.execute(
            "SELECT id, depth, token_count, content, message_range_start, message_range_end "
            "FROM summaries WHERE session_id = ? ORDER BY depth DESC, message_range_start ASC",
            (session_id,)
        ).fetchall()

        if not rows:
            return fresh_tail

        # Fill budget with summaries
        summary_parts = []
        used_tokens = 0
        for row in rows:
            sid, depth, tc, content, rs, re_ = row
            if used_tokens + tc > budget:
                continue
            summary_parts.append({
                "id": sid, "depth": depth, "tokens": tc,
                "content": content, "range": f"messages {rs}-{re_}"
            })
            used_tokens += tc

        if not summary_parts:
            return fresh_tail

        # Sort by message range for chronological order
        summary_parts.sort(key=lambda s: s["range"])

        # Build summary message
        summary_text = "## Compacted Conversation History\n\n"
        summary_text += "The following summaries cover earlier parts of this conversation. "
        summary_text += "Use context_search, context_detail, or context_recall tools to access original details.\n\n"
        for s in summary_parts:
            summary_text += f"### Summary (depth {s['depth']}, {s['range']})\n{s['content']}\n\n"

        assembled = [
            {"role": "user", "content": f"[Conversation Context]\n{summary_text}"},
            {"role": "assistant", "content": "I have the context from our earlier conversation. I can use context_search, context_detail, or context_recall to access specific details if needed."},
        ]
        assembled.extend(fresh_tail)
        return assembled

    def check_and_compact(self, messages: list[dict], session_id: str,
                          model: str, api_key: str, base_url: str, api_type: str,
                          system_prompt: str = "",
                          max_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS,
                          force: bool = False) -> tuple[list[dict], bool]:
        """Check if compaction needed and perform hierarchical summarization."""
        cfg = self.get_config()
        if not cfg.get("enabled", True) and not force:
            return messages, False

        estimated = _estimate_conversation_tokens(messages, system_prompt)
        if not force:
            threshold_pct = cfg.get("compact_threshold", 0.75)
            threshold = int(max_tokens * threshold_pct)
            if estimated < threshold:
                return messages, False

        fresh_tail_count = cfg.get("fresh_tail_count", 32)
        msgs_per_summary = cfg.get("messages_per_summary", 10)

        # Determine which messages to summarize (everything before fresh tail)
        if len(messages) <= fresh_tail_count:
            if force and len(messages) > msgs_per_summary:
                # Force mode: use half the messages as tail to allow some summarization
                fresh_tail_count = len(messages) // 2
            else:
                return messages, False

        old_messages = messages[:-fresh_tail_count]

        # Find what's already summarized (by checking existing summary ranges)
        conn = _context_conn()
        existing = conn.execute(
            "SELECT MAX(message_range_end) FROM summaries WHERE session_id = ?",
            (session_id,)
        ).fetchone()
        already_summarized = existing[0] if existing and existing[0] else 0

        # Calculate message indices (1-based, relative to session)
        total_msg_count = len(messages)
        old_count = len(old_messages)

        # Only summarize messages not yet covered
        unsummarized_start = already_summarized
        unsummarized_msgs = old_messages[unsummarized_start:]

        if len(unsummarized_msgs) < msgs_per_summary:
            # Not enough new messages to summarize, just assemble
            return self.assemble_context(session_id, messages, system_prompt, max_tokens), True

        # Use summary model (Gemini Flash default, Haiku fallback)
        summary_model, summary_model_fallback = self._resolve_summary_model()
        summary_model = summary_model or model
        summary_model_fallback = summary_model_fallback or model

        pct = int(estimated / max_tokens * 100)
        logging.info(f"Context {pct}% full (~{estimated:,} tokens), creating hierarchical summaries...")

        # Create leaf summaries in chunks
        for i in range(0, len(unsummarized_msgs), msgs_per_summary):
            chunk = unsummarized_msgs[i:i + msgs_per_summary]
            if len(chunk) < 3:  # skip tiny remnants
                continue
            range_start = unsummarized_start + i
            range_end = unsummarized_start + i + len(chunk)
            self.summarize_chunk(chunk, session_id, range_start, range_end,
                                 summary_model, api_key, base_url, api_type,
                                 fallback_model=summary_model_fallback)

        # Try condensation
        self.condense(session_id)

        # Assemble final context
        assembled = self.assemble_context(session_id, messages, system_prompt, max_tokens)
        new_estimated = _estimate_conversation_tokens(assembled, system_prompt)
        new_pct = int(new_estimated / max_tokens * 100)
        logging.info(f"Compacted: {pct}% → {new_pct}% (~{new_estimated:,} tokens)")

        return assembled, True

    def search(self, session_id: str, query: str, regex: bool = False, limit: int = 10) -> list[dict]:
        """Search original messages in ChatDB by keyword/regex."""
        results = []
        try:
            from server import _db_conn as chat_db_conn
            conn = chat_db_conn()
            conn.row_factory = sqlite3.Row
            if regex:
                # SQLite doesn't support regex natively, use LIKE as fallback
                rows = conn.execute(
                    "SELECT id, role, content, created_at FROM messages WHERE session_id = ? AND content LIKE ? ORDER BY id LIMIT ?",
                    (session_id, f"%{query}%", limit)
                ).fetchall()
            else:
                rows = conn.execute(
                    "SELECT id, role, content, created_at FROM messages WHERE session_id = ? AND content LIKE ? ORDER BY id LIMIT ?",
                    (session_id, f"%{query}%", limit)
                ).fetchall()
            for r in rows:
                content = r["content"]
                # Extract snippet around match
                idx = content.lower().find(query.lower())
                if idx >= 0:
                    start = max(0, idx - 80)
                    end = min(len(content), idx + len(query) + 120)
                    snippet = ("..." if start > 0 else "") + content[start:end] + ("..." if end < len(content) else "")
                else:
                    snippet = content[:200]
                results.append({
                    "message_id": r["id"],
                    "role": r["role"],
                    "snippet": snippet,
                    "timestamp": r["created_at"],
                })
        except Exception as e:
            results.append({"error": str(e)})
        return results

    def get_detail(self, summary_id: str) -> dict:
        """Get a summary and its original messages."""
        conn = _context_conn()
        row = conn.execute(
            "SELECT * FROM summaries WHERE id = ?", (summary_id,)
        ).fetchone()
        if not row:
            return {"error": f"Summary {summary_id} not found"}

        summary = {
            "id": row[0], "session_id": row[1], "depth": row[2],
            "token_count": row[3], "content": row[4],
            "parent_ids": json.loads(row[5] or "[]"),
            "message_range": f"{row[6]}-{row[7]}",
        }

        # Load original messages from ChatDB
        try:
            from server import _db_conn as chat_db_conn
            cconn = chat_db_conn()
            cconn.row_factory = sqlite3.Row
            msgs = cconn.execute(
                "SELECT id, role, content FROM messages WHERE session_id = ? AND id >= ? AND id <= ? ORDER BY id",
                (row[1], row[6], row[7])
            ).fetchall()
            summary["original_messages"] = [{"id": m["id"], "role": m["role"], "content": m["content"][:500]} for m in msgs]
        except Exception:
            summary["original_messages"] = []

        return summary

    def recall(self, session_id: str, query: str,
               model: str, api_key: str, base_url: str, api_type: str) -> str:
        """Deep recall: search summaries and original messages, then answer with a focused sub-query."""
        # Find relevant summaries
        conn = _context_conn()
        conn.row_factory = sqlite3.Row
        rows = conn.execute(
            "SELECT id, content, message_range_start, message_range_end FROM summaries WHERE session_id = ? ORDER BY depth DESC",
            (session_id,)
        ).fetchall()

        relevant_context = []
        for r in rows:
            if query.lower() in r["content"].lower():
                relevant_context.append(r["content"])

        # Also search original messages
        msg_results = self.search(session_id, query, limit=5)
        for mr in msg_results:
            if "snippet" in mr:
                relevant_context.append(f"[{mr['role']}]: {mr['snippet']}")

        if not relevant_context:
            return f"No relevant context found for: {query}"

        context_text = "\n\n---\n\n".join(relevant_context[:10])
        r_model, r_fallback = self._resolve_summary_model()
        r_model = r_model or model

        try:
            result = _run_delegate_with_fallback(
                messages=[{"role": "user", "content": (
                    f"Based on this conversation history, answer the following query precisely:\n\n"
                    f"**Query:** {query}\n\n**Context:**\n{context_text}"
                )}],
                primary_model=r_model, fallback_model=r_fallback,
                system_prompt="Answer based only on the provided context. Be specific and cite details.",
                memory_store=None,
                inference_params={"max_tokens": 2000, "temperature": 0.1},
                tools=False,
            )
            if result and not result.startswith("Delegation error"):
                return result.strip()
        except Exception as e:
            return f"Recall failed: {e}"

        return f"Could not generate answer for: {query}"

    def get_stats(self, session_id: str) -> dict:
        """Get context stats for a session."""
        conn = _context_conn()
        rows = conn.execute(
            "SELECT depth, COUNT(*), SUM(token_count) FROM summaries WHERE session_id = ? GROUP BY depth",
            (session_id,)
        ).fetchall()
        depth_stats = {r[0]: {"count": r[1], "tokens": r[2]} for r in rows}
        total_summaries = sum(s["count"] for s in depth_stats.values())
        total_tokens = sum(s["tokens"] for s in depth_stats.values())
        return {
            "session_id": session_id,
            "total_summaries": total_summaries,
            "total_summary_tokens": total_tokens,
            "depth_distribution": depth_stats,
            "config": self.get_config(),
        }


_context_manager: ContextManager | None = None


def _estimate_tokens_str(text: str) -> int:
    """Estimate token count for a string."""
    return max(1, len(text) // CHARS_PER_TOKEN)


def _estimate_tokens_message(msg: dict) -> int:
    """Estimate token count for a single message (any format)."""
    content = msg.get("content", "")
    if isinstance(content, str):
        return _estimate_tokens_str(content) + 4  # role overhead
    elif isinstance(content, list):
        # Anthropic-style content blocks (text, tool_use, tool_result, etc.)
        total = 4
        for block in content:
            if isinstance(block, dict):
                if block.get("type") == "text":
                    total += _estimate_tokens_str(block.get("text", ""))
                elif block.get("type") == "tool_use":
                    total += _estimate_tokens_str(json.dumps(block.get("input", {}))) + 20
                elif block.get("type") == "tool_result":
                    total += _estimate_tokens_str(str(block.get("content", "")))
                else:
                    total += _estimate_tokens_str(json.dumps(block))
            elif isinstance(block, str):
                total += _estimate_tokens_str(block)
        return total
    return 10


def _estimate_conversation_tokens(messages: list[dict], system_prompt: str = "") -> int:
    """Estimate total token count for the full conversation."""
    total = _estimate_tokens_str(system_prompt) if system_prompt else 0
    for msg in messages:
        total += _estimate_tokens_message(msg)
    return total


def _compact_conversation(messages: list[dict], model: str, api_key: str,
                          base_url: str, api_type: str,
                          max_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS) -> list[dict]:
    """Compact the conversation by summarizing older messages.

    Strategy:
    1. Keep the last KEEP_RECENT_MESSAGES as-is
    2. Summarize everything before that into a single system-style message
    3. If the model is available, use it to summarize; otherwise do a simple truncation
    """
    if len(messages) <= KEEP_RECENT_MESSAGES:
        return messages

    # Split into old and recent
    recent = messages[-KEEP_RECENT_MESSAGES:]
    old = messages[:-KEEP_RECENT_MESSAGES]

    # Build a text representation of old messages for summarization
    old_text_parts = []
    for msg in old:
        role = msg.get("role", "unknown")
        content = msg.get("content", "")
        if isinstance(content, str):
            text = content
        elif isinstance(content, list):
            text_pieces = []
            for block in content:
                if isinstance(block, dict):
                    if block.get("type") == "text":
                        text_pieces.append(block.get("text", ""))
                    elif block.get("type") == "tool_use":
                        text_pieces.append(f"[Called tool: {block.get('name', '')}]")
                    elif block.get("type") == "tool_result":
                        c = block.get("content", "")
                        text_pieces.append(f"[Tool result: {str(c)[:200]}]")
                    else:
                        text_pieces.append(f"[{block.get('type', 'block')}]")
            text = " ".join(text_pieces)
        else:
            text = str(content)

        # Truncate individual messages in the summary source
        if len(text) > 500:
            text = text[:500] + "..."
        old_text_parts.append(f"{role}: {text}")

    old_summary_source = "\n".join(old_text_parts)

    # Try to use the model to summarize
    summary = None
    try:
        summary_messages = [{
            "role": "user",
            "content": (
                "Summarize the following conversation history in 2-3 concise paragraphs. "
                "Focus on: key decisions made, information learned, files modified, "
                "and current task context. Be factual and brief.\n\n"
                f"{old_summary_source}"
            )
        }]
        summary = send_message(
            summary_messages, model, api_key, base_url, api_type,
            silent=True, tools=False, _tool_round=0,
        )
    except Exception:
        pass

    if not summary:
        # Fallback: simple truncation
        if len(old_summary_source) > 2000:
            summary = old_summary_source[:2000] + "\n...(earlier conversation truncated)"
        else:
            summary = old_summary_source

    # Build compacted conversation
    compacted = [
        {"role": "user", "content": f"[Previous conversation summary]\n{summary}"},
        {"role": "assistant", "content": "Understood. I have the context from our previous conversation. How can I help?"},
    ]
    compacted.extend(recent)

    return compacted


def _check_and_compact(messages: list[dict], model: str, api_key: str,
                       base_url: str, api_type: str,
                       system_prompt: str = "",
                       max_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS,
                       session_id: str = "",
                       force: bool = False) -> tuple[list[dict], bool]:
    """Check if conversation needs compaction and do it if necessary.

    Returns (messages, was_compacted).
    Uses lossless ContextManager if enabled, otherwise falls back to legacy compaction.
    """
    # Check for per-agent compact_threshold override
    tcfg = _get_token_config()
    agent_threshold = tcfg.get("compact_threshold")

    # Lossless context management (if enabled and session_id available)
    if _context_manager and session_id:
        try:
            cfg = _context_manager.get_config()
            if cfg.get("enabled", True):
                # Apply agent-level threshold override if set
                effective_max = max_tokens
                if agent_threshold is not None:
                    # Temporarily adjust max_tokens to shift the threshold
                    # ContextManager uses its own config threshold internally
                    pass
                return _context_manager.check_and_compact(
                    messages, session_id, model, api_key, base_url, api_type,
                    system_prompt=system_prompt, max_tokens=max_tokens,
                )
        except Exception as e:
            logging.warning(f"ContextManager failed, falling back to legacy: {e}")

    # Legacy flat compaction
    estimated = _estimate_conversation_tokens(messages, system_prompt)
    threshold_pct = agent_threshold if agent_threshold is not None else COMPACT_THRESHOLD
    threshold = int(max_tokens * threshold_pct)

    if estimated < threshold and not force:
        return messages, False

    # Show compaction notice
    pct = int(estimated / max_tokens * 100)
    print(f"\n  {DIM}⟳ Context {pct}% full (~{estimated:,} tokens), compacting...{RESET}")
    sys.stdout.flush()

    compacted = _compact_conversation(messages, model, api_key, base_url, api_type, max_tokens)
    new_estimated = _estimate_conversation_tokens(compacted, system_prompt)
    new_pct = int(new_estimated / max_tokens * 100)
    print(f"  {DIM}✔ Compacted: {pct}% → {new_pct}% (~{new_estimated:,} tokens){RESET}")

    return compacted, True


# --- Task cancellation ---

class TaskCancelled(Exception):
    """Raised when the user presses Escape to cancel the current task."""
    pass


class CancelToken:
    """Simple cancellation token — compatible with EscapeWatcher interface."""

    def __init__(self):
        self._cancelled = threading.Event()

    @property
    def cancelled(self) -> bool:
        return self._cancelled.is_set()

    def cancel(self):
        self._cancelled.set()

    def start(self):
        pass  # No-op — no background thread needed

    def stop(self):
        pass  # No-op


class EscapeWatcher:
    """Background thread that monitors for Escape key press in raw terminal mode."""

    def __init__(self):
        self._cancelled = threading.Event()
        self._stop = threading.Event()
        self._thread = None
        self._old_settings = None

    @property
    def cancelled(self) -> bool:
        return self._cancelled.is_set()

    def start(self):
        self._cancelled.clear()
        self._stop.clear()
        self._thread = threading.Thread(target=self._watch, daemon=True)
        self._thread.start()

    def stop(self):
        self._stop.set()
        if self._thread:
            self._thread.join(timeout=0.5)

    def _watch(self):
        fd = sys.stdin.fileno()
        try:
            old_settings = termios.tcgetattr(fd)
        except termios.error:
            return
        try:
            # Set input to non-canonical mode WITHOUT disabling output processing.
            # tty.setraw() clears OPOST which breaks \n -> \r\n conversion for
            # all print() output. Instead, only change what we need for input.
            new_settings = termios.tcgetattr(fd)
            # Disable canonical mode (ICANON) and echo (ECHO) in local flags
            new_settings[3] = new_settings[3] & ~(termios.ICANON | termios.ECHO)
            # Set minimum chars to read = 0, timeout = 0 (non-blocking)
            new_settings[6][termios.VMIN] = 0
            new_settings[6][termios.VTIME] = 0
            termios.tcsetattr(fd, termios.TCSANOW, new_settings)

            while not self._stop.is_set():
                if select.select([fd], [], [], 0.1)[0]:
                    ch = os.read(fd, 1)
                    if ch == b'\x1b':  # Escape key
                        self._cancelled.set()
                        break
                    elif ch in (b'o', b'O'):
                        # Toggle tool output visibility live
                        _toggle_tool_output()
        except Exception:
            pass
        finally:
            try:
                termios.tcsetattr(fd, termios.TCSADRAIN, old_settings)
            except termios.error:
                pass


# --- Spinner ---

SPINNER_CHARS = ["·", "✢", "✳", "∗", "✻", "✽"]
SPINNER_VERBS = [
    "Brewing", "Baking", "Thinking", "Pondering", "Computing",
    "Crafting", "Conjuring", "Composing", "Contemplating", "Weaving",
]


class Spinner:
    """Animated spinner shown on the status bar while waiting for a response."""

    def __init__(self, model: str):
        self.model = model
        self.verb = random.choice(SPINNER_VERBS)
        self._stop = threading.Event()
        self._thread = None
        self.start_time = 0.0

    def start(self):
        self.start_time = time.time()
        self._stop.clear()
        self._thread = threading.Thread(target=self._animate, daemon=True)
        self._thread.start()

    def stop(self) -> float:
        self._stop.set()
        if self._thread:
            self._thread.join()
        elapsed = time.time() - self.start_time
        return elapsed

    def _animate(self):
        idx = 0
        while not self._stop.is_set():
            char = SPINNER_CHARS[idx % len(SPINNER_CHARS)]
            cols = shutil.get_terminal_size().columns
            rows = shutil.get_terminal_size().lines
            elapsed = time.time() - self.start_time
            # Dark background with colored segments
            spinner_part = f" {FG_ORANGE}{char}{RESET}{BG_DARK} {WHITE}{self.verb}...{RESET}{BG_DARK}"
            time_part = f" {FG_GRAY}({elapsed:.0f}s){RESET}{BG_DARK}"
            model_part = f" {DIM}│{RESET}{BG_DARK} {GREEN}{self.model}{RESET}{BG_DARK} "
            # visible length: " X Verb... (Ns) │ model "
            visible_len = 2 + len(self.verb) + 3 + 2 + len(f"{elapsed:.0f}") + 2 + 3 + len(self.model) + 1
            padding = max(0, cols - visible_len)
            bar = f"\033[48;5;235m{spinner_part}{time_part}{model_part}{' ' * padding}{RESET}"
            sys.stdout.write(f"\0337\033[{rows};1H{bar}\0338")
            sys.stdout.flush()
            idx += 1
            self._stop.wait(0.12)


# --- Markdown rendering ---

# ANSI codes
BOLD = "\033[1m"
DIM = "\033[2m"
ITALIC = "\033[3m"
RESET = "\033[0m"
CYAN = "\033[36m"
YELLOW = "\033[33m"
GREEN = "\033[32m"
MAGENTA = "\033[35m"
BLUE = "\033[34m"
RED = "\033[31m"
BG_GRAY = "\033[48;5;236m"
FG_GRAY = "\033[38;5;245m"
FG_ORANGE = "\033[38;5;208m"
WHITE = "\033[37m"
BG_DARK = "\033[48;5;235m"

# Minimum box inner width so things don't collapse on very narrow terminals
_MIN_BOX_INNER = 20


def _term_cols() -> int:
    """Get current terminal width."""
    return shutil.get_terminal_size().columns


def _box_width() -> int:
    """Inner width for box-drawing, adapts to terminal width.
    Leaves 2 chars indent + 1 border left + 1 border right = 4 chars margin."""
    return max(_MIN_BOX_INNER, _term_cols() - 4)


def _box_top(label: str = "") -> str:
    """Draw ┌─ label ───...┐ spanning the available width."""
    w = _box_width()
    if label:
        vlen = _visible_len(label)
        # "┌─ " + label + " " + fill + "┐"  →  3 + vlen + 1 + fill + 1 = w + 2
        fill = max(0, w - 4 - vlen)
        return f"  {DIM}┌─ {RESET}{label}{DIM} {'─' * fill}┐{RESET}"
    else:
        return f"  {DIM}┌{'─' * w}┐{RESET}"


def _visible_len(text: str) -> int:
    """Get the visible (column) length of a string, ignoring ANSI escape codes.
    Handles wide Unicode characters and tabs."""
    import unicodedata
    stripped = re.sub(r"\033\[[0-9;]*[a-zA-Z]", "", text)
    stripped = re.sub(r"\033\[\?[0-9;]*[a-zA-Z]", "", stripped)
    stripped = re.sub(r"\033\([A-Z]", "", stripped)
    w = 0
    for ch in stripped:
        if ch == '\t':
            w += (8 - w % 8)
        elif unicodedata.east_asian_width(ch) in ('W', 'F'):
            w += 2
        elif unicodedata.category(ch) in ('Mn', 'Me', 'Cf'):
            pass  # zero-width combining/format chars
        else:
            w += 1
    return w


def _truncate_visible(text: str, max_cols: int) -> str:
    """Truncate a string (may contain ANSI codes) to max visible columns."""
    import unicodedata
    cols = 0
    i = 0
    while i < len(text) and cols < max_cols:
        if text[i] == '\033':
            # Skip full ANSI sequence
            j = i + 1
            while j < len(text) and text[j] not in 'mABCDHJKfhlGn':
                j += 1
            i = j + 1
        elif text[i] == '\t':
            add = 8 - cols % 8
            if cols + add > max_cols:
                break
            cols += add
            i += 1
        else:
            ch = text[i]
            cw = 2 if unicodedata.east_asian_width(ch) in ('W', 'F') else 1
            if unicodedata.category(ch) in ('Mn', 'Me', 'Cf'):
                cw = 0
            if cols + cw > max_cols:
                break
            cols += cw
            i += 1
    if i < len(text):
        # Check if remaining is only ANSI codes
        rest = text[i:]
        rest_stripped = re.sub(r"\033\[[0-9;]*[a-zA-Z]", "", rest)
        if rest_stripped.strip():
            return text[:i] + RESET
    return text


def _box_mid(content: str = "") -> str:
    """Draw │ content — truncated to fit terminal width. No wrapping."""
    # "  │  " = 5 visible chars prefix
    max_content = max(10, _term_cols() - 6)
    # Expand tabs to spaces first to avoid variable-width issues
    content = content.replace("\t", "    ")
    truncated = _truncate_visible(content, max_content)
    return f"  {DIM}│{RESET}  {truncated}"


def _box_bot() -> str:
    """Draw └───...┘ spanning the available width."""
    w = _box_width()
    return f"  {DIM}└{'─' * w}┘{RESET}"


def _render_tool_calls(text: str) -> str:
    """Pre-process tool_call XML blocks into styled representations."""

    def _replace_tool_call(m):
        block = m.group(1)
        func_match = re.search(r"<function=([^>]+)>", block)
        func_name = func_match.group(1) if func_match else "unknown"
        params = re.findall(
            r"<parameter=([^>]+)>\s*(.*?)\s*</parameter>",
            block, re.DOTALL)

        lines = []
        lines.append(_box_top(f"{FG_ORANGE}{BOLD}⚡ Tool Call{RESET}"))
        lines.append(_box_mid(f"{CYAN}{BOLD}{func_name}{RESET}()"))
        if params:
            for pname, pval in params:
                pval = pval.strip()
                lines.append(_box_mid(f"  {MAGENTA}{pname}{RESET}{DIM}:{RESET} {WHITE}{pval}{RESET}"))
        lines.append(_box_bot())
        return "\n".join(lines)

    text = re.sub(
        r"<tool_call>\s*(.*?)\s*</tool_call>",
        _replace_tool_call, text, flags=re.DOTALL)

    def _replace_standalone_tool(m):
        block = m.group(0)
        func_match = re.search(r"<function=([^>]+)>", block)
        func_name = func_match.group(1) if func_match else "unknown"
        params = re.findall(
            r"<parameter=([^>]+)>\s*(.*?)\s*</parameter>",
            block, re.DOTALL)

        lines = []
        lines.append(_box_top(f"{FG_ORANGE}{BOLD}⚡ Tool Call{RESET}"))
        lines.append(_box_mid(f"{CYAN}{BOLD}{func_name}{RESET}()"))
        if params:
            for pname, pval in params:
                pval = pval.strip()
                lines.append(_box_mid(f"  {MAGENTA}{pname}{RESET}{DIM}:{RESET} {WHITE}{pval}{RESET}"))
        lines.append(_box_bot())
        return "\n".join(lines)

    text = re.sub(
        r"<function=([^>]+)>\s*(?:<parameter=([^>]+)>\s*(.*?)\s*</parameter>\s*)*</function>",
        _replace_standalone_tool, text, flags=re.DOTALL)

    return text


def _render_tool_results(text: str) -> str:
    """Pre-process tool result XML blocks into styled representations."""

    def _replace_result(m):
        content = m.group(1).strip()
        lines = []
        lines.append(_box_top(f"{GREEN}{BOLD}✔ Tool Result{RESET}"))
        for cline in content.split("\n"):
            lines.append(_box_mid(cline))
        lines.append(_box_bot())
        return "\n".join(lines)

    text = re.sub(
        r"<tool_result>\s*(.*?)\s*</tool_result>",
        _replace_result, text, flags=re.DOTALL)

    return text


def render_markdown(text: str) -> str:
    """Render markdown text with ANSI escape codes for terminal display."""
    # Pre-process tool calls and results before line-by-line markdown
    text = _render_tool_calls(text)
    text = _render_tool_results(text)

    lines = text.split("\n")
    result = []
    in_code_block = False
    code_lang = ""

    for line in lines:
        # Skip lines already processed by tool renderers (contain box-drawing chars)
        if "┌─" in line or "└─" in line or (line.startswith("  ") and "│" in line[:6]):
            result.append(line)
            continue

        # Fenced code block toggle
        if line.strip().startswith("```"):
            if not in_code_block:
                in_code_block = True
                code_lang = line.strip()[3:].strip()
                if code_lang:
                    result.append(_box_top(f"{DIM}{code_lang}{RESET}"))
                else:
                    result.append(_box_top())
            else:
                in_code_block = False
                result.append(_box_bot())
            continue

        if in_code_block:
            result.append(f"  {DIM}│{RESET} {YELLOW}{line}{RESET}")
            continue

        # Headers
        if line.startswith("#### "):
            result.append(f"{BOLD}{CYAN}{line[5:]}{RESET}")
            continue
        if line.startswith("### "):
            result.append(f"{BOLD}{CYAN}{line[4:]}{RESET}")
            continue
        if line.startswith("## "):
            result.append(f"{BOLD}{CYAN}{line[3:]}{RESET}")
            continue
        if line.startswith("# "):
            result.append(f"{BOLD}{CYAN}{line[2:]}{RESET}")
            continue

        # Horizontal rule
        if re.match(r"^[-*_]{3,}\s*$", line):
            result.append(f"{DIM}{'─' * max(20, _term_cols() - 4)}{RESET}")
            continue

        # Bullet lists
        m = re.match(r"^(\s*)([-*])\s+(.*)$", line)
        if m:
            indent, _, content = m.groups()
            content = _inline_format(content)
            result.append(f"{indent}  {DIM}•{RESET} {content}")
            continue

        # Numbered lists
        m = re.match(r"^(\s*)(\d+)\.\s+(.*)$", line)
        if m:
            indent, num, content = m.groups()
            content = _inline_format(content)
            result.append(f"{indent}  {DIM}{num}.{RESET} {content}")
            continue

        # Regular line with inline formatting
        result.append(_inline_format(line))

    return "\n".join(result)


def _inline_format(text: str) -> str:
    """Apply inline markdown formatting (bold, italic, code)."""
    # Inline code (must be before bold/italic to avoid conflicts)
    text = re.sub(r"`([^`]+)`", rf"{YELLOW}\1{RESET}", text)
    # Bold + italic
    text = re.sub(r"\*\*\*(.+?)\*\*\*", rf"{BOLD}{ITALIC}\1{RESET}", text)
    # Bold
    text = re.sub(r"\*\*(.+?)\*\*", rf"{BOLD}\1{RESET}", text)
    # Italic
    text = re.sub(r"(?<!\*)\*(?!\*)(.+?)(?<!\*)\*(?!\*)", rf"{ITALIC}\1{RESET}", text)
    # Links [text](url) -> text (url)
    text = re.sub(r"\[([^\]]+)\]\(([^)]+)\)", rf"{BOLD}\1{RESET} {DIM}(\2){RESET}", text)
    return text


# --- API functions ---

def make_headers(api_key: str, api_type: str) -> dict:
    """Build request headers for the given API type."""
    if api_type == "openai":
        return {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }
    if api_type == "mistral":
        return {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }
    return {
        "Content-Type": "application/json",
        "x-api-key": api_key,
        "Authorization": f"Bearer {api_key}",
        "anthropic-version": "2023-06-01",
    }


def get_available_models(api_key: str, base_url: str, api_type: str) -> list[str]:
    """Fetch available models from the API and return as a list."""
    headers = make_headers(api_key, api_type)
    request = urllib.request.Request(
        f"{base_url}/models", headers=headers, method="GET",
    )
    try:
        with urllib.request.urlopen(request) as response:
            data = json.loads(response.read().decode("utf-8"))
            models = data.get("data", [])
            return [model.get("id") for model in models if model.get("id")]
    except (urllib.error.HTTPError, urllib.error.URLError):
        return []


# --- Model Configuration System ---

KNOWN_MODELS = {
    "claude-opus": {"icon": "\U0001f7e3", "priority": 100, "max_context": 200000, "capabilities": ["coding", "analysis", "agentic", "creative"]},
    "claude-sonnet": {"icon": "\U0001f7e0", "priority": 80, "max_context": 200000, "capabilities": ["coding", "analysis", "fast"]},
    "claude-haiku": {"icon": "\U0001f7e2", "priority": 60, "max_context": 200000, "capabilities": ["fast"]},
    "gemini": {"icon": "\U0001f48e", "priority": 70, "max_context": 1000000, "capabilities": ["coding", "analysis"]},
    "qwen": {"icon": "\U0001f43c", "priority": 50, "max_context": 131072, "capabilities": ["coding", "analysis"]},
    "crow": {"icon": "\U0001f426\u200d\u2b1b", "priority": 30, "max_context": 32768, "capabilities": ["fast", "local"]},
    "llama": {"icon": "\U0001f999", "priority": 40, "max_context": 131072, "capabilities": ["coding", "local"]},
    "mistral": {"icon": "\U0001f32c\ufe0f", "priority": 45, "max_context": 131072, "capabilities": ["coding", "analysis"]},
}

CAPABILITY_VALUES = ["coding", "analysis", "agentic", "fast", "creative", "local"]

_models_config: dict = {}


def _match_known_model(model_id: str) -> dict:
    """Match a model ID against KNOWN_MODELS patterns. Returns default config."""
    m = model_id.lower()
    for prefix, defaults in KNOWN_MODELS.items():
        if m.startswith(prefix) or prefix in m:
            # Generate shortname from model ID
            shortname = model_id.split("/")[-1]  # strip provider prefix
            # Simplify common patterns
            for suffix in ["-20251101", "-20250101", "-20241022", "-20251001"]:
                shortname = shortname.replace(suffix, "")
            result = {
                "enabled": True,
                "shortname": shortname,
                "display_name": shortname,
                "icon": defaults["icon"],
                "priority": defaults["priority"],
                "capabilities": list(defaults["capabilities"]),
            }
            if "max_context" in defaults:
                result["max_context"] = defaults["max_context"]
            return result
    # Unknown model — enabled with low priority
    shortname = model_id.split("/")[-1]
    return {
        "enabled": True,
        "shortname": shortname,
        "display_name": shortname,
        "icon": "\U0001f916",
        "priority": 10,
        "capabilities": [],
    }


def init_models_config(providers: dict, existing_models: dict | None = None) -> dict:
    """Auto-populate models config from provider model lists.

    Merges with existing config (preserves user edits). Returns the models dict.
    """
    global _models_config
    if existing_models:
        _models_config = dict(existing_models)

    # Discover models from all providers
    for name, p in providers.items():
        try:
            models = get_available_models(
                p.get("api_key", ""), p.get("base_url", ""), p.get("type", "openai"))
        except Exception:
            models = []
        for model_id in models:
            if model_id not in _models_config:
                entry = _match_known_model(model_id)
                entry["provider"] = name
                _models_config[model_id] = entry
            else:
                if "provider" not in _models_config[model_id]:
                    _models_config[model_id]["provider"] = name
                # Backfill max_context from KNOWN_MODELS if missing
                if "max_context" not in _models_config[model_id]:
                    known = _match_known_model(model_id)
                    if "max_context" in known:
                        _models_config[model_id]["max_context"] = known["max_context"]

    return _models_config


def resolve_model(model_spec: str, purpose: str | None = None) -> str:
    """Resolve a model specifier to a canonical model ID.

    Handles:
    - "auto" + purpose → highest-priority enabled model with that capability
    - shortname (e.g. "opus") → lookup by shortname
    - canonical ID → pass through
    """
    if not model_spec:
        return ""

    if model_spec == "auto":
        return _resolve_auto_model(purpose)

    # Try shortname lookup
    for model_id, cfg in _models_config.items():
        if cfg.get("shortname") == model_spec:
            if cfg.get("enabled", True):
                return model_id
            break  # Found but disabled

    # Pass through as canonical ID
    return model_spec


def _resolve_auto_model(purpose: str | None) -> str:
    """Pick the best enabled model for a given purpose."""
    enabled = [(mid, cfg) for mid, cfg in _models_config.items() if cfg.get("enabled", True)]
    if not enabled:
        return ""

    if purpose:
        # Filter to models with the requested capability
        matching = [(mid, cfg) for mid, cfg in enabled if purpose in cfg.get("capabilities", [])]
        if matching:
            matching.sort(key=lambda x: x[1].get("priority", 0), reverse=True)
            return matching[0][0]

    # Fallback: highest priority enabled model
    enabled.sort(key=lambda x: x[1].get("priority", 0), reverse=True)
    return enabled[0][0]


def get_enabled_models() -> list[str]:
    """Return enabled model IDs sorted by priority descending."""
    enabled = [(mid, cfg) for mid, cfg in _models_config.items() if cfg.get("enabled", True)]
    enabled.sort(key=lambda x: x[1].get("priority", 0), reverse=True)
    return [mid for mid, _ in enabled]


# --- Task-based auto model selection ---

_PURPOSE_PATTERNS: dict[str, list[re.Pattern]] = {
    "coding": [
        re.compile(r"\b(write|fix|debug|refactor|implement|code|function|class|bug|error|traceback|test|unittest|lint|compile|build|deploy|dockerfile|makefile|git\b|commit|branch|merge|PR\b|pull request|regex|parse|api|endpoint|route|migration|schema|sql|query|index\.)\b", re.I),
        re.compile(r"\b(python|javascript|typescript|rust|go|java|c\+\+|html|css|json|yaml|toml|bash|shell|script)\b", re.I),
        re.compile(r"```", re.I),
    ],
    "analysis": [
        re.compile(r"\b(analy[sz]e|explain|compare|evaluate|review|assess|investigate|research|summarize|summary|report|pros?\b|cons?\b|trade-?off|benchmark|metric|statistic|data|insight|trend|pattern|cause|why does|how does|what is|understand)\b", re.I),
    ],
    "creative": [
        re.compile(r"\b(write|draft|compose|story|poem|essay|blog|article|creative|brainstorm|idea|name|slogan|tagline|marketing|copy|narrative|fiction|dialogue|character)\b", re.I),
    ],
    "agentic": [
        re.compile(r"\b(search|find|look up|fetch|download|browse|scrape|monitor|schedule|automate|run|execute|install|setup|configure|deploy|send|email|notify|delegate|background)\b", re.I),
    ],
    "fast": [
        re.compile(r"\b(quick|brief|short|one-?liner|yes or no|true or false|translate|convert|format|list|enumerate|define)\b", re.I),
    ],
}


def classify_task_purpose(message: str) -> str | None:
    """Classify a user message into a capability/purpose using keyword heuristics.

    Returns the best-matching purpose or None if no strong signal.
    Scores each purpose by number of pattern matches; requires at least 2 signal
    strength to avoid false positives on single-word matches.
    """
    if not message:
        return None

    scores: dict[str, int] = {}
    for purpose, patterns in _PURPOSE_PATTERNS.items():
        score = 0
        for pat in patterns:
            hits = len(pat.findall(message))
            score += hits
        if score > 0:
            scores[purpose] = score

    if not scores:
        return None

    best_purpose = max(scores, key=scores.get)
    # Require some minimum signal (at least 2 keyword hits) to avoid noisy classification
    if scores[best_purpose] < 2:
        return None

    return best_purpose


def resolve_auto_model_for_task(agent_config: dict, message: str) -> tuple[str, str | None]:
    """For agents with model="auto", analyze the task and pick the best model.

    Returns (resolved_model_id, detected_purpose).
    If agent has a fixed model_purpose, uses that instead of classifying.
    """
    raw_model = agent_config.get("model", "")
    if raw_model != "auto":
        return resolve_model(raw_model, agent_config.get("model_purpose")), agent_config.get("model_purpose")

    fixed_purpose = agent_config.get("model_purpose")
    if fixed_purpose:
        return _resolve_auto_model(fixed_purpose), fixed_purpose

    # Classify task from message
    detected = classify_task_purpose(message)
    resolved = _resolve_auto_model(detected)
    return resolved, detected


def get_model_info(model: str) -> dict:
    """Return config entry for a model, or empty dict if not configured."""
    return _models_config.get(model, {})


def get_model_max_context(model: str) -> int:
    """Return the model's context window size, or DEFAULT_MAX_CONTEXT_TOKENS."""
    return _models_config.get(model, {}).get("max_context", DEFAULT_MAX_CONTEXT_TOKENS)


# Default max output tokens per model family
_MAX_OUTPUT_DEFAULTS = {
    "opus": 32768,
    "sonnet": 16384,
    "haiku": 8192,
    "minimax": 32768,
    "m2.7": 32768,
    "m2.5": 32768,
}


def get_model_max_output(model: str) -> int:
    """Return the model's max output tokens. Checks models_config first, then family defaults."""
    cfg = _models_config.get(model, {})
    if cfg.get("max_output"):
        return int(cfg["max_output"])
    ml = model.lower()
    for family, limit in _MAX_OUTPUT_DEFAULTS.items():
        if family in ml:
            return limit
    # Non-Anthropic models: use a generous default
    return 16384


def get_inference_params(model: str, purpose: str | None = None) -> dict:
    """Resolve inference parameters for a model, optionally overlaying a purpose preset.

    Returns only explicitly set keys (empty dict = use API defaults).
    """
    cfg = _models_config.get(model, {})
    base = dict(cfg.get("inference", {}))
    if purpose:
        preset = cfg.get("presets", {}).get(purpose, {})
        base.update(preset)
    return base


# Keys valid for all providers
_INFERENCE_STANDARD_KEYS = {"temperature", "top_p", "top_k", "max_tokens"}
# Keys only for OpenAI-compatible / oMLX
_INFERENCE_OPENAI_KEYS = {"frequency_penalty", "presence_penalty"}
# oMLX extension keys (also OpenAI-compatible endpoint)
_INFERENCE_OMLX_KEYS = {"min_p", "repetition_penalty"}


def _apply_inference_to_payload(payload: dict, params: dict, api_type: str, provider: str = "") -> None:
    """Apply resolved inference params to an API payload with provider translation."""
    if not params:
        return

    is_omlx = provider == "omlx"

    for key in _INFERENCE_STANDARD_KEYS:
        if key in params:
            payload[key] = params[key]

    for key in _INFERENCE_OPENAI_KEYS:
        if key in params and (api_type == "openai" or is_omlx):
            payload[key] = params[key]

    for key in _INFERENCE_OMLX_KEYS:
        if key in params and is_omlx:
            payload[key] = params[key]

    # Mistral-specific inference keys
    if api_type == "mistral":
        if "reasoning_effort" in params:
            payload["reasoning_effort"] = params["reasoning_effort"]

    # Thinking translation
    if params.get("thinking"):
        budget = params.get("thinking_budget", 4096)
        if api_type == "anthropic":
            payload["thinking"] = {"type": "enabled", "budget_tokens": budget}
        elif api_type == "mistral":
            # Map thinking to Mistral reasoning_effort (Vibe CLI mapping)
            _thinking_to_reasoning = {"low": "none", "medium": "high", "high": "high"}
            payload["reasoning_effort"] = _thinking_to_reasoning.get(
                str(params.get("thinking", "high")), "high"
            )
        elif is_omlx:
            payload.setdefault("chat_template_kwargs", {})["enable_thinking"] = True


def list_models(api_key: str, base_url: str, api_type: str) -> None:
    """List available models from the API."""
    models = get_available_models(api_key, base_url, api_type)
    if models:
        print("Available models:")
        for model_id in models:
            print(f"  {model_id}")
    else:
        print("No models available")


def _unescape(text: str) -> str:
    """Unescape literal backslash sequences that some APIs send."""
    return text.replace("\\n", "\n").replace("\\t", "\t").replace('\\"', '"')


def _collect_anthropic(response) -> str:
    """Parse Anthropic SSE stream silently. Returns full text."""
    collected = []
    current_event = None
    for line in response:
        line = line.decode("utf-8").strip()
        if line.startswith("event: "):
            current_event = line[7:]
        elif line.startswith("data: "):
            if current_event == "message_stop":
                break
            try:
                event = json.loads(line[6:])
                if event.get("type") == "content_block_delta":
                    delta = event.get("delta", {})
                    if delta.get("type") == "text_delta":
                        collected.append(_unescape(delta.get("text", "")))
            except json.JSONDecodeError:
                pass
    return "".join(collected)


def _collect_openai(response) -> str:
    """Parse OpenAI SSE stream silently. Returns full text."""
    collected = []
    for line in response:
        line = line.decode("utf-8").strip()
        if not line.startswith("data: "):
            continue
        payload = line[6:]
        if payload == "[DONE]":
            break
        try:
            event = json.loads(payload)
            choices = event.get("choices", [])
            if choices:
                delta = choices[0].get("delta", {})
                content = delta.get("content")
                if content:
                    collected.append(_unescape(content))
        except json.JSONDecodeError:
            pass
    return "".join(collected)


def _stream_anthropic(response) -> str:
    """Parse Anthropic SSE stream with live output. Returns full text."""
    collected = []
    current_event = None
    for line in response:
        line = line.decode("utf-8").strip()
        if line.startswith("event: "):
            current_event = line[7:]
        elif line.startswith("data: "):
            if current_event == "message_stop":
                break
            try:
                event = json.loads(line[6:])
                if event.get("type") == "content_block_delta":
                    delta = event.get("delta", {})
                    if delta.get("type") == "text_delta":
                        text = _unescape(delta.get("text", ""))
                        print(text, end="", flush=True)
                        collected.append(text)
            except json.JSONDecodeError:
                pass
    return "".join(collected)


def _stream_openai(response) -> str:
    """Parse OpenAI SSE stream with live output. Returns full text."""
    collected = []
    for line in response:
        line = line.decode("utf-8").strip()
        if not line.startswith("data: "):
            continue
        payload = line[6:]
        if payload == "[DONE]":
            break
        try:
            event = json.loads(payload)
            choices = event.get("choices", [])
            if choices:
                delta = choices[0].get("delta", {})
                content = delta.get("content")
                if content:
                    content = _unescape(content)
                    print(content, end="", flush=True)
                    collected.append(content)
        except json.JSONDecodeError:
            pass
    return "".join(collected)


TOOL_ICONS = {
    "read_file": "r", "write_file": "w", "edit_file": "e",
    "list_directory": "d", "search_files": "s", "execute_command": "$",
    "web_fetch": "~", "exa_search": "?",
    "memory_store": "+", "memory_recall": "m", "memory_delete": "-", "memory_shared": "M",
    "delegate_task": ">", "use_skill": "*",
    "gmail_inbox": "@", "gmail_read": "@", "gmail_search": "@",
    "gmail_send": "@", "gmail_reply": "@",
    "task_status": "?", "task_cancel": "x",
    "context_search": "c", "context_detail": "c", "context_recall": "c",
    "list_nodes": "n", "schedule_list": "t", "schedule_history": "h",
    "read_document": "D", "write_document": "D", "edit_document": "D",
    "code_graph_build": "G", "code_graph_query": "G", "code_graph_impact": "G", "code_graph_enhance": "G",
    "git_command": "g", "github_command": "g",
}

TOOL_VERBS = {
    "read_file": "Reading", "write_file": "Writing", "edit_file": "Editing",
    "list_directory": "Listing", "search_files": "Searching", "execute_command": "Executing",
    "web_fetch": "Fetching", "exa_search": "Searching",
    "memory_store": "Remembering", "memory_recall": "Recalling", "memory_delete": "Forgetting", "memory_shared": "Shared Memory",
    "delegate_task": "Delegating", "use_skill": "Loading Skill",
    "gmail_inbox": "Inbox", "gmail_read": "Reading Email", "gmail_search": "Searching Email",
    "gmail_send": "Sending Email", "gmail_reply": "Replying",
    "task_status": "Task Status", "task_cancel": "Cancelling",
    "context_search": "Searching Context", "context_detail": "Context Detail", "context_recall": "Recalling",
    "list_nodes": "Listing Nodes", "schedule_list": "Schedules", "schedule_history": "History",
    "read_document": "Reading Document", "write_document": "Writing Document", "edit_document": "Editing Document",
    "code_graph_build": "Building Graph", "code_graph_query": "Querying Graph", "code_graph_impact": "Impact Analysis", "code_graph_enhance": "Enhancing Graph",
    "git_command": "Git", "github_command": "GitHub",
}


# Tool display mode
_show_tools = True
_tool_call_count = 0         # counter for current response
_tool_output_buffer = []     # buffered output lines when hidden
_tool_counter_active = False # whether we have a counter line on screen
_tool_output_shown = False   # whether buffered output is currently displayed


def _reset_tool_tracking() -> None:
    """Reset tool tracking for a new response."""
    global _tool_call_count, _tool_output_buffer, _tool_counter_active, _tool_output_shown
    _tool_call_count = 0
    _tool_output_buffer = []
    _tool_counter_active = False
    _tool_output_shown = False


def _update_tool_counter() -> None:
    """In hidden mode, update the single-line tool counter in-place."""
    global _tool_counter_active
    msg = f"  {DIM}Tools called: {RESET}{BOLD}{_tool_call_count}{RESET}{DIM}  (press o to show/hide){RESET}"
    if _tool_counter_active:
        # Move up to overwrite the previous counter line
        sys.stdout.write(f"\033[A\r\033[K{msg}\n")
    else:
        sys.stdout.write(f"{msg}\n")
        _tool_counter_active = True
    sys.stdout.flush()


def _toggle_tool_output() -> None:
    """Toggle display of buffered tool output (called from EscapeWatcher on 'o')."""
    global _tool_output_shown
    if not _tool_output_buffer:
        return
    if not _tool_output_shown:
        # Show buffered output below the counter
        for line in _tool_output_buffer:
            sys.stdout.write(f"{line}\n")
        sys.stdout.flush()
        _tool_output_shown = True
    else:
        # Hide: move cursor up past the output and clear lines
        line_count = len(_tool_output_buffer)
        for _ in range(line_count):
            sys.stdout.write(f"\033[A\r\033[K")
        sys.stdout.flush()
        _tool_output_shown = False


def _format_tool_call(name: str, args: dict) -> list[str]:
    """Format a tool call box as a list of lines."""
    icon = TOOL_ICONS.get(name, "⚡")
    verb = TOOL_VERBS.get(name, "Running")
    max_w = max(20, _term_cols() - 8)

    lines = []
    lines.append(f"\n{_box_top(f'{FG_ORANGE}{BOLD}{icon} {verb}{RESET}')}")

    if name == "exa_search":
        lines.append(_box_mid(f"{CYAN}{name}{RESET}({MAGENTA}query{RESET}={WHITE}\"{args.get('query', '')}\"{RESET})"))
    elif name == "execute_command":
        cmd = args.get("command", "")
        if len(cmd) > max_w:
            cmd = cmd[:max_w - 3] + "..."
        lines.append(_box_mid(f"{YELLOW}$ {cmd}{RESET}"))
    elif name in ("read_file", "write_file", "edit_file"):
        lines.append(_box_mid(f"{CYAN}{name}{RESET} {WHITE}{args.get('path', '')}{RESET}"))
    elif name == "list_directory":
        p = args.get("path", ".")
        pat = args.get("pattern", "")
        label = f"{p}/{pat}" if pat else p
        lines.append(_box_mid(f"{CYAN}{name}{RESET} {WHITE}{label}{RESET}"))
    elif name == "search_files":
        lines.append(_box_mid(f"{CYAN}{name}{RESET} {MAGENTA}/{args.get('pattern', '')}/{RESET} in {WHITE}{args.get('path', '.')}{RESET}"))
    elif name == "web_fetch":
        lines.append(_box_mid(f"{CYAN}{args.get('method', 'GET')}{RESET} {WHITE}{args.get('url', '')}{RESET}"))
    elif name == "memory_store":
        lines.append(_box_mid(f"{MAGENTA}{args.get('name', '')}{RESET} {DIM}[{args.get('type', 'general')}]{RESET}"))
    elif name == "memory_recall":
        q = args.get("query", "(list all)")
        lines.append(_box_mid(f"{MAGENTA}{q}{RESET}"))
    elif name == "memory_shared":
        action = args.get("action", "recall")
        scope = args.get("scope", "global")
        scope_label = "team" if scope == "team" else "main"
        if action == "store":
            lines.append(_box_mid(f"{CYAN}store → {scope_label}:{RESET} {MAGENTA}{args.get('name', '')}{RESET}"))
        else:
            q = args.get("query", "(list all)")
            lines.append(_box_mid(f"{CYAN}recall ← {scope_label}:{RESET} {MAGENTA}{q}{RESET}"))
    elif name == "memory_delete":
        lines.append(_box_mid(f"{RED}{args.get('name', '')}{RESET}"))
    elif name == "delegate_task":
        lines.append(_box_mid(f"{CYAN}agent:{RESET} {BOLD}{args.get('agent', '')}{RESET}"))
        task_preview = args.get("task", "")[:max_w]
        lines.append(_box_mid(f"{DIM}{task_preview}{RESET}"))
    elif name == "use_skill":
        lines.append(_box_mid(f"{MAGENTA}{args.get('skill', '')}{RESET}"))
    else:
        summary = ", ".join(f"{k}={v!r}" for k, v in list(args.items())[:3])
        if len(summary) > max_w:
            summary = summary[:max_w - 3] + "..."
        lines.append(_box_mid(f"{CYAN}{name}{RESET}({summary})"))

    lines.append(_box_bot())
    return lines


def _display_tool_call(name: str, args: dict) -> None:
    """Print a styled tool invocation box, or buffer it in hidden mode."""
    global _tool_call_count
    _tool_call_count += 1
    formatted = _format_tool_call(name, args)

    if _show_tools:
        for line in formatted:
            print(line)
        sys.stdout.flush()
    else:
        _tool_output_buffer.extend(formatted)
        _update_tool_counter()


def _format_tool_result(name: str, result_str: str) -> list[str]:
    """Format a tool result summary box as a list of lines."""
    max_w = max(20, _term_cols() - 8)
    try:
        rdata = json.loads(result_str)
    except json.JSONDecodeError:
        return []

    out = []

    if rdata.get("error"):
        out.append(_box_top(f"{RED}{BOLD}✘ Error{RESET}"))
        for line in rdata["error"].split("\n")[:5]:
            out.append(_box_mid(line[:max_w]))
        out.append(f"{_box_bot()}\n")
        return out

    if name == "exa_search":
        count = rdata.get("result_count", 0)
        out.append(_box_top(f"{GREEN}{BOLD}✔ {count} results{RESET}"))
        for r in rdata.get("results", [])[:3]:
            out.append(_box_mid(f"{BOLD}{r.get('title', '')[:max_w]}{RESET}"))
        if count > 3:
            out.append(_box_mid(f"{DIM}... and {count - 3} more{RESET}"))
    elif name == "read_file":
        total = rdata.get("total_lines", 0)
        showing = rdata.get("showing", "")
        out.append(_box_top(f"{GREEN}{BOLD}✔ {total} lines{RESET} {DIM}(showing {showing}){RESET}"))
        content = rdata.get("content", "")
        lines = content.split("\n")
        for line in lines[:5]:
            out.append(_box_mid(f"{DIM}{line[:max_w]}{RESET}"))
        if len(lines) > 5:
            out.append(_box_mid(f"{DIM}... {len(lines) - 5} more lines{RESET}"))
    elif name == "write_file":
        out.append(_box_top(f"{GREEN}{BOLD}✔ Written{RESET}"))
        out.append(_box_mid(f"{rdata.get('path', '')} ({rdata.get('size', 0)} bytes)"))
    elif name == "edit_file":
        n = rdata.get("replacements", 0)
        out.append(_box_top(f"{GREEN}{BOLD}✔ {n} replacement{'s' if n != 1 else ''}{RESET}"))
        out.append(_box_mid(f"{rdata.get('path', '')}"))
    elif name == "list_directory":
        count = rdata.get("count", 0)
        out.append(_box_top(f"{GREEN}{BOLD}✔ {count} entries{RESET}"))
        for e in rdata.get("entries", [])[:8]:
            icon = "d" if e.get("type") == "directory" else " "
            out.append(_box_mid(f"{icon} {e.get('name', '')}"))
        if count > 8:
            out.append(_box_mid(f"{DIM}... and {count - 8} more{RESET}"))
    elif name == "search_files":
        mc = rdata.get("match_count", 0)
        fs = rdata.get("files_searched", 0)
        out.append(_box_top(f"{GREEN}{BOLD}✔ {mc} matches{RESET} {DIM}({fs} files searched){RESET}"))
        for m in rdata.get("matches", [])[:5]:
            out.append(_box_mid(f"{CYAN}{m.get('file', '')}:{m.get('line', '')}{RESET} {m.get('text', '')[:max_w]}"))
        if mc > 5:
            out.append(_box_mid(f"{DIM}... and {mc - 5} more{RESET}"))
    elif name == "execute_command":
        ec = rdata.get("exit_code", -1)
        label = f"{GREEN}{BOLD}✔ exit {ec}{RESET}" if ec == 0 else f"{RED}{BOLD}✘ exit {ec}{RESET}"
        out.append(_box_top(label))
        output = rdata.get("output", "")
        lines = output.split("\n")
        for line in lines[:10]:
            out.append(_box_mid(f"{DIM}{line[:max_w]}{RESET}"))
        if len(lines) > 10:
            out.append(_box_mid(f"{DIM}... {len(lines) - 10} more lines{RESET}"))
    elif name == "web_fetch":
        status = rdata.get("status", "")
        length = rdata.get("length", 0)
        out.append(_box_top(f"{GREEN}{BOLD}✔ HTTP {status}{RESET} {DIM}({length} chars){RESET}"))
        content = rdata.get("content", "")
        lines = content.split("\n")
        for line in lines[:5]:
            out.append(_box_mid(f"{DIM}{line[:max_w]}{RESET}"))
        if len(lines) > 5:
            out.append(_box_mid(f"{DIM}... {len(lines) - 5} more lines{RESET}"))
    elif name == "memory_store":
        out.append(_box_top(f"{GREEN}{BOLD}✔ Stored{RESET}"))
        out.append(_box_mid(f"{rdata.get('name', '')} → {rdata.get('file', '')}"))
    elif name == "memory_recall":
        count = rdata.get("count", 0)
        out.append(_box_top(f"{GREEN}{BOLD}✔ {count} memories{RESET}"))
        for r in rdata.get("results", [])[:5]:
            out.append(_box_mid(f"{BOLD}{r.get('name', '')}{RESET} {DIM}[{r.get('type', '')}]{RESET}"))
            desc = r.get("description", "")
            if desc:
                out.append(_box_mid(f"  {DIM}{desc[:max_w]}{RESET}"))
        if count > 5:
            out.append(_box_mid(f"{DIM}... and {count - 5} more{RESET}"))
    elif name == "memory_shared":
        source = rdata.get("source", "main")
        if rdata.get("status") == "stored":
            out.append(_box_top(f"{GREEN}{BOLD}✔ Stored → {source}{RESET}"))
            out.append(_box_mid(f"{rdata.get('name', '')} → {rdata.get('file', '')}"))
        else:
            count = rdata.get("count", 0)
            out.append(_box_top(f"{GREEN}{BOLD}✔ {count} shared memories{RESET} {DIM}({source}){RESET}"))
            for r in rdata.get("results", [])[:5]:
                out.append(_box_mid(f"{BOLD}{r.get('name', '')}{RESET} {DIM}[{r.get('type', '')}]{RESET}"))
            if count > 5:
                out.append(_box_mid(f"{DIM}... and {count - 5} more{RESET}"))
    elif name == "memory_delete":
        out.append(_box_top(f"{GREEN}{BOLD}✔ Deleted{RESET}"))
        out.append(_box_mid(f"{rdata.get('name', '')}"))
    elif name == "delegate_task":
        agent = rdata.get("agent", "")
        resp = rdata.get("response", "")
        out.append(_box_top(f"{GREEN}{BOLD}✔ {agent} responded{RESET}"))
        for line in resp.split("\n")[:8]:
            out.append(_box_mid(f"{DIM}{line[:max_w]}{RESET}"))
        resp_lines = resp.split("\n")
        if len(resp_lines) > 8:
            out.append(_box_mid(f"{DIM}... {len(resp_lines) - 8} more lines{RESET}"))
    elif name == "use_skill":
        skill = rdata.get("skill", "")
        instructions = rdata.get("instructions", "")
        line_count = len(instructions.split("\n"))
        out.append(_box_top(f"{GREEN}{BOLD}✔ {skill}{RESET} {DIM}({line_count} lines loaded){RESET}"))
    else:
        out.append(_box_top(f"{GREEN}{BOLD}✔ Done{RESET}"))

    out.append(f"{_box_bot()}\n")
    return out


def _display_tool_result(name: str, result_str: str) -> None:
    """Print a styled tool result summary box, or buffer it in hidden mode."""
    formatted = _format_tool_result(name, result_str)
    if _show_tools:
        for line in formatted:
            print(line)
        sys.stdout.flush()
    else:
        _tool_output_buffer.extend(formatted)
        _update_tool_counter()


TOOL_DISPATCH = {
    "read_file": tool_read_file,
    "write_file": tool_write_file,
    "edit_file": tool_edit_file,
    "list_directory": tool_list_directory,
    "search_files": tool_search_files,
    "execute_command": tool_execute_command,
    "web_fetch": tool_web_fetch,
    "exa_search": lambda args: exa_search(
        query=args.get("query", ""),
        num_results=args.get("num_results", 5),
        category=args.get("category"),
    ),
    "memory_store": tool_memory_store,
    "memory_recall": tool_memory_recall,
    "memory_delete": tool_memory_delete,
    "memory_shared": tool_memory_shared,
    "gmail_inbox": tool_gmail_inbox,
    "gmail_read": tool_gmail_read,
    "gmail_search": tool_gmail_search,
    "gmail_send": tool_gmail_send,
    "gmail_reply": tool_gmail_reply,
    "delegate_task": tool_delegate_task,
    "task_status": tool_task_status,
    "task_cancel": tool_task_cancel,
    "use_skill": tool_use_skill,
    "list_nodes": tool_list_nodes,
    "context_search": tool_context_search,
    "context_detail": tool_context_detail,
    "context_recall": tool_context_recall,
    "schedule_list": tool_schedule_list,
    "schedule_history": tool_schedule_history,
    "read_document": tool_read_document,
    "write_document": tool_write_document,
    "edit_document": tool_edit_document,
    "mcp_connect": tool_mcp_connect,
    "mcp_disconnect": tool_mcp_disconnect,
    "mcp_servers": tool_mcp_servers,
    "code_graph_build": tool_code_graph_build,
    "code_graph_query": tool_code_graph_query,
    "code_graph_impact": tool_code_graph_impact,
    "code_graph_enhance": tool_code_graph_enhance,
    "git_command": tool_git_command,
    "github_command": tool_github_command,
    "tool_search": lambda args: _tool_search(args),
}


def _tool_search(args: dict) -> str:
    """Search for available tools by name or description.

    Returns matching tool schemas from both built-in and MCP tools.
    Discovered tools are tracked per-session and included in subsequent API calls.
    """
    query = args.get("query", "").lower()
    max_results = args.get("max_results", 5)

    if not query:
        return _err("query is required")

    # Search built-in tools
    matches = []
    for td in TOOL_DEFINITIONS:
        name = td.get("name", "")
        desc = td.get("description", "")
        if isinstance(desc, tuple):
            desc = " ".join(desc)
        score = 0
        if query in name.lower():
            score += 10
        if query in desc.lower():
            score += 5
        # Fuzzy: match individual words
        for word in query.split():
            if word in name.lower():
                score += 3
            if word in desc.lower():
                score += 1
        if score > 0:
            matches.append((score, {"name": name, "description": desc[:200],
                                     "input_schema": td.get("input_schema", {})}))

    # Search MCP tools
    mcp_mgr = getattr(_thread_local, 'mcp_manager', None) or _mcp_manager
    if mcp_mgr:
        try:
            for mcp_td in mcp_mgr.get_tool_definitions():
                name = mcp_td.get("name", "")
                desc = mcp_td.get("description", "")
                score = 0
                if query in name.lower():
                    score += 10
                if query in desc.lower():
                    score += 5
                for word in query.split():
                    if word in name.lower():
                        score += 3
                    if word in desc.lower():
                        score += 1
                if score > 0:
                    matches.append((score, {"name": name, "description": desc[:200],
                                             "input_schema": mcp_td.get("input_schema", {})}))
        except Exception:
            pass

    # Sort by score descending, take top results
    matches.sort(key=lambda x: x[0], reverse=True)
    results = [m[1] for m in matches[:max_results]]

    # Track discovered tools for deferred loading
    discovered = getattr(_thread_local, '_discovered_tools', set())
    for r in results:
        discovered.add(r["name"])
    _thread_local._discovered_tools = discovered

    if not results:
        return _ok({"matches": [], "message": f"No tools found matching '{query}'"})
    return _ok({"matches": results, "count": len(results)})


# Per-thread tool call dedup tracking
_tool_call_history = threading.local()


def _check_tool_dedup(name: str, args: dict) -> str | None:
    """Check if this exact tool call was already made. Raises TaskCancelled after 2 dupes.

    Exempt tools: memory_recall (legitimate re-queries with same terms),
    delegate_task (different execution context each time).
    """
    # Exempt tools that legitimately repeat with same args
    _DEDUP_EXEMPT = {"memory_recall", "memory_shared", "delegate_task", "task_status",
                     "schedule_list", "schedule_history"}
    if name in _DEDUP_EXEMPT:
        return None

    if not hasattr(_tool_call_history, 'calls'):
        _tool_call_history.calls = set()
        _tool_call_history.dupe_count = 0
        _tool_call_history.consecutive_dupes = 0
    key = f"{name}:{json.dumps(args, sort_keys=True)}"
    if key in _tool_call_history.calls:
        _tool_call_history.dupe_count += 1
        _tool_call_history.consecutive_dupes += 1
        if _tool_call_history.consecutive_dupes >= 2:
            # Hard abort — model is stuck in a loop
            raise TaskCancelled()
        return _err(
            f"Duplicate tool call detected. You already called {name} with these exact arguments. "
            "Use the previous result or try a different approach."
        )
    _tool_call_history.calls.add(key)
    _tool_call_history.consecutive_dupes = 0  # Reset consecutive counter on new unique call
    if len(_tool_call_history.calls) > 100:
        _tool_call_history.calls = set(list(_tool_call_history.calls)[-50:])
    return None


def reset_tool_dedup():
    """Reset the dedup tracker."""
    _tool_call_history.calls = set()
    _tool_call_history.dupe_count = 0
    _tool_call_history.consecutive_dupes = 0


# --- Hook Runner ---

class HookRunner:
    """Runs external hook scripts for tool execution lifecycle events."""

    def __init__(self, agent_id: str = "main"):
        self.agent_id = agent_id
        self._hooks = self._load_hooks()

    def _load_hooks(self) -> list[dict]:
        """Load hook config from agent.json."""
        try:
            cfg = AgentConfig(self.agent_id)
            hooks_cfg = cfg.config.get("hooks", {})
            if not hooks_cfg.get("enabled", False):
                return []
            return [h for h in hooks_cfg.get("scripts", []) if h.get("enabled", True)]
        except Exception:
            return []

    def reload(self):
        self._hooks = self._load_hooks()

    def _get_timeout(self) -> int:
        try:
            cfg = AgentConfig(self.agent_id)
            return cfg.config.get("hooks", {}).get("timeout", 5000)
        except Exception:
            return 5000

    def get_hooks(self, hook_type: str, tool_name: str = "") -> list[dict]:
        """Get matching hooks for a type and tool."""
        result = []
        for h in self._hooks:
            if h.get("type") != hook_type:
                continue
            tools = h.get("tools", ["*"])
            if "*" in tools or tool_name in tools:
                result.append(h)
        return result

    def run_hook(self, hook: dict, env_extra: dict) -> tuple[int, str]:
        """Run a hook script. Returns (exit_code, stdout)."""
        script = hook.get("script", "")
        if not script:
            return 0, ""

        agent_dir = os.path.join(AGENTS_DIR, self.agent_id)
        script_path = os.path.join(agent_dir, script)
        if not os.path.isfile(script_path):
            logging.warning(f"Hook script not found: {script_path}")
            return 0, ""

        env = os.environ.copy()
        env["HOOK_AGENT"] = self.agent_id
        env["HOOK_SESSION_ID"] = getattr(_thread_local, 'session_id', "") or ""
        env["HOOK_TIMESTAMP"] = time.strftime("%Y-%m-%dT%H:%M:%SZ", time.gmtime())
        env.update({k: str(v) for k, v in env_extra.items()})

        timeout_s = self._get_timeout() / 1000.0
        stdin_data = json.dumps(env_extra).encode("utf-8")

        try:
            proc = subprocess.Popen(
                ["bash", script_path],
                stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                stdin=subprocess.PIPE, env=env,
                cwd=agent_dir, start_new_session=True,
            )
            stdout, stderr = proc.communicate(input=stdin_data, timeout=timeout_s)
            output = stdout.decode("utf-8", errors="replace").strip()
            if stderr:
                logging.debug(f"Hook {hook.get('name','?')} stderr: {stderr.decode('utf-8', errors='replace')[:200]}")
            return proc.returncode, output
        except subprocess.TimeoutExpired:
            try:
                os.killpg(proc.pid, 9)
            except OSError:
                proc.kill()
            proc.communicate(timeout=2)
            logging.warning(f"Hook {hook.get('name','?')} timed out after {timeout_s}s")
            return 1, f"Hook timed out after {timeout_s}s"
        except Exception as e:
            logging.warning(f"Hook {hook.get('name','?')} failed: {e}")
            return 0, ""  # fail-open

    def run_pre_hooks(self, tool_name: str, args: dict) -> str | None:
        """Run pre-hooks. Returns error message if blocked, None if allowed."""
        hooks = self.get_hooks("pre", tool_name)
        for h in hooks:
            env = {
                "HOOK_TYPE": "pre",
                "HOOK_TOOL_NAME": tool_name,
                "HOOK_TOOL_ARGS": json.dumps(args),
            }
            code, output = self.run_hook(h, env)
            if code == 1:
                hook_name = h.get("name", "unknown")
                msg = output or f"Blocked by hook: {hook_name}"
                logging.info(f"Hook {hook_name} blocked {tool_name}: {msg[:100]}")
                return _err(f"HOOK BLOCKED ({hook_name}): {msg}")
            if code == 2:
                break  # skip remaining hooks
        return None

    def run_post_hooks(self, tool_name: str, args: dict, result: str) -> str:
        """Run post-hooks. Returns (possibly modified) result."""
        hooks = self.get_hooks("post", tool_name)
        for h in hooks:
            env = {
                "HOOK_TYPE": "post",
                "HOOK_TOOL_NAME": tool_name,
                "HOOK_TOOL_ARGS": json.dumps(args),
                "HOOK_TOOL_RESULT": result[:50000],  # cap env var size
            }
            code, output = self.run_hook(h, env)
            if code == 0 and output:
                result = output  # modify result
            if code == 2:
                break
        return result

    def run_after_file_write(self, file_path: str, action: str):
        """Run after_file_write hooks."""
        hooks = self.get_hooks("after_file_write")
        for h in hooks:
            env = {
                "HOOK_TYPE": "after_file_write",
                "HOOK_FILE_PATH": file_path,
                "HOOK_FILE_ACTION": action,
            }
            self.run_hook(h, env)


# Cache hook runners per agent
_hook_runners: dict[str, HookRunner] = {}
_hook_runners_lock = threading.Lock()


def _get_hook_runner(agent_id: str = "main") -> HookRunner:
    """Get or create a HookRunner for an agent."""
    with _hook_runners_lock:
        if agent_id not in _hook_runners:
            _hook_runners[agent_id] = HookRunner(agent_id)
        return _hook_runners[agent_id]


# --- Artifact Helpers ---

_ARTIFACT_TYPE_MAP = {
    "html": "html", "htm": "html",
    "svg": "svg",
    "md": "markdown", "markdown": "markdown",
    "png": "image", "jpg": "image", "jpeg": "image", "gif": "image", "webp": "image", "bmp": "image",
    "pdf": "document", "docx": "document", "xlsx": "document", "pptx": "document",
    "txt": "text",
}

def _is_artifact_path(path: str) -> bool:
    """Check if path is under agents/<name>/artifacts/"""
    try:
        agents_dir = os.path.realpath(AGENTS_DIR)
        real_path = os.path.realpath(path)
        if not real_path.startswith(agents_dir + os.sep):
            return False
        parts = real_path[len(agents_dir) + 1:].split(os.sep)
        return len(parts) >= 3 and parts[1] == "artifacts"
    except Exception:
        return False

def _register_artifact_version(path: str, action: str, agent_id: str):
    """Register or update an artifact in the DB, capturing content snapshot.
    Returns (artifact_id, version, type) or None on failure."""
    try:
        from server import ChatDB
        import uuid as _uuid_mod

        session_id = getattr(_thread_local, 'current_session_id', None) or ""
        if not session_id:
            return None

        name = os.path.basename(path)
        ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
        artifact_type = _ARTIFACT_TYPE_MAP.get(ext, "code")

        # Read content snapshot (cap at 5MB)
        try:
            with open(path, "rb") as f:
                content = f.read(5 * 1024 * 1024)
            size = os.path.getsize(path)
            if size > 5 * 1024 * 1024:
                content = None  # too large, disk-only
        except Exception:
            content = None
            size = 0

        # Check if artifact already exists for this path+session
        existing = ChatDB.get_artifact_by_path(session_id, path)
        if existing:
            artifact_id = existing["id"]
            next_version = existing["latest_version"] + 1
        else:
            artifact_id = str(_uuid_mod.uuid4())[:12]
            ChatDB.create_artifact(artifact_id, session_id, agent_id, name, path, artifact_type)
            next_version = 1

        ChatDB.add_artifact_version(artifact_id, next_version, content, size, None, action)
        return (artifact_id, next_version, artifact_type)
    except Exception as e:
        print(f"  [WARN] artifact registration: {e}", flush=True)
        return None


# --- Centralized File-Write Pipeline ---

def _after_file_write(path: str, action: str = "created", agent_id: str = ""):
    """Centralized post-file-write pipeline. Called from tool_write_file and tool_edit_file.
    Replaces scattered _maybe_qmd_reindex(), _extract_entities(), file_created calls."""
    is_artifact = _is_artifact_path(path)

    # 1. QMD reindex (skip for artifacts)
    if not is_artifact:
        _maybe_qmd_reindex(path)

    # 2. Entity extraction + knowledge graph (skip for artifacts)
    if path.endswith(".md") and agent_id and not is_artifact:
        try:
            with open(path, "r") as f:
                raw = f.read()
            fm, body = _parse_frontmatter(raw)
            entities = _extract_entities(body)
            if entities:
                _update_entity_index(agent_id, os.path.basename(path), entities)
        except Exception:
            pass

    # 3. File/artifact event emission (for UI)
    ecb = getattr(_thread_local, 'event_callback', None)
    if ecb:
        try:
            if is_artifact:
                art_result = _register_artifact_version(path, action, agent_id)
                if art_result:
                    art_id, art_ver, art_type = art_result
                    ecb("artifact_updated", {
                        "path": path,
                        "name": os.path.basename(path),
                        "size": os.path.getsize(path),
                        "action": action,
                        "artifact_id": art_id,
                        "artifact_version": art_ver,
                        "artifact_type": art_type,
                    })
                else:
                    # Fallback to regular file event if artifact registration failed
                    ecb("file_created", {
                        "path": path,
                        "name": os.path.basename(path),
                        "size": os.path.getsize(path),
                        "action": action,
                    })
            else:
                ecb("file_created", {
                    "path": path,
                    "name": os.path.basename(path),
                    "size": os.path.getsize(path),
                    "action": action,
                })
        except Exception:
            pass

    # 4. Update code graph if source file
    _maybe_update_code_graph(path)

    # 5. External after_file_write hooks
    if agent_id:
        try:
            runner = _get_hook_runner(agent_id)
            runner.run_after_file_write(path, action)
        except Exception:
            pass


# --- Concurrent Tool Execution (Phase 5) ---
# Tools classified as concurrency-safe (read-only, no side effects)
_CONCURRENT_SAFE_TOOLS = {
    "read_file", "list_directory", "search_files", "read_document",
    "memory_recall", "memory_shared",
    "exa_search", "web_fetch",
    "code_graph_query",
    "schedule_list", "schedule_history", "list_nodes", "task_status",
    "context_search", "context_detail", "context_recall",
    "git_command",  # read-only git commands are safe (status, log, diff, blame)
}


def _execute_tools_batch(tool_calls: list[dict], event_callback=None) -> list[dict]:
    """Execute a batch of tool calls with concurrent-safe parallelism.

    Partitions tool calls into batches:
    - Consecutive concurrent-safe tools run in parallel (ThreadPoolExecutor)
    - Unsafe tools run sequentially, one at a time
    Results are returned in the original order.

    Each tool_call dict: {"id": str, "name": str, "input": dict}
    Returns list of {"tool_use_id": str, "result": str} in order.
    """
    if not tool_calls:
        return []

    # Partition into batches: [(is_concurrent, [tool_calls...])]
    batches = []
    current_batch = []
    current_is_concurrent = None

    for tc in tool_calls:
        is_safe = tc["name"] in _CONCURRENT_SAFE_TOOLS
        # git_command: only safe for read-only subcommands
        if tc["name"] == "git_command":
            subcmd = tc["input"].get("subcommand", "")
            if subcmd not in ("status", "log", "diff", "blame", "show", "branch", "tag", "remote"):
                is_safe = False

        if current_is_concurrent is None:
            current_is_concurrent = is_safe
        elif is_safe != current_is_concurrent:
            batches.append((current_is_concurrent, current_batch))
            current_batch = []
            current_is_concurrent = is_safe
        current_batch.append(tc)

    if current_batch:
        batches.append((current_is_concurrent, current_batch))

    # Execute batches
    results = []
    for is_concurrent, batch in batches:
        if is_concurrent and len(batch) > 1:
            # Parallel execution
            from concurrent.futures import ThreadPoolExecutor, as_completed
            futures = {}
            with ThreadPoolExecutor(max_workers=min(5, len(batch))) as executor:
                for tc in batch:
                    if event_callback:
                        event_callback("tool_call", {"name": tc["name"], "args": tc["input"]})
                    _display_tool_call(tc["name"], tc["input"])
                    future = executor.submit(_execute_tool_in_thread, tc["name"], tc["input"],
                                             tc["id"], event_callback)
                    futures[future] = tc

                # Collect results preserving order
                result_map = {}
                for future in as_completed(futures):
                    tc = futures[future]
                    try:
                        result = future.result()
                    except Exception as e:
                        result = _err(f"Tool execution error: {e}")
                    result_map[tc["id"]] = result
                    _display_tool_result(tc["name"], result)
                    if event_callback:
                        event_callback("tool_result", {"name": tc["name"], "result": result})

            # Append in original order
            for tc in batch:
                results.append({"tool_use_id": tc["id"], "result": result_map[tc["id"]]})
        else:
            # Sequential execution
            for tc in batch:
                _display_tool_call(tc["name"], tc["input"])
                if event_callback:
                    event_callback("tool_call", {"name": tc["name"], "args": tc["input"]})
                _thread_local.event_callback = event_callback
                _thread_local.tool_use_id = tc["id"]
                try:
                    result = _execute_tool(tc["name"], tc["input"])
                finally:
                    _thread_local.event_callback = None
                    _thread_local.tool_use_id = None
                _display_tool_result(tc["name"], result)
                if event_callback:
                    event_callback("tool_result", {"name": tc["name"], "result": result})
                results.append({"tool_use_id": tc["id"], "result": result})
    return results


def _execute_tool_in_thread(name: str, args: dict, tool_id: str, event_callback) -> str:
    """Execute a single tool in a worker thread with proper thread-local setup."""
    # Copy thread-local context from parent (set by the chat handler)
    # The worker thread inherits the main thread's context
    _thread_local.event_callback = event_callback
    _thread_local.tool_use_id = tool_id
    try:
        return _execute_tool(name, args)
    finally:
        _thread_local.event_callback = None
        _thread_local.tool_use_id = None


def _execute_tool(name: str, args: dict) -> str:
    """Execute a tool by name with the given arguments."""
    # --- Built-in pre-hooks ---
    # Plan mode: block non-readonly tools
    if getattr(_thread_local, 'plan_mode', False):
        if name == "memory_shared" and args.get("action") == "store":
            return _err("Blocked in plan mode. Describe what you would do instead.")
        if name not in READONLY_TOOLS:
            return _err("Blocked in plan mode. Describe what you would do instead.")
    # Workflow tool restriction (was dead code — now enforced)
    workflow_tools = getattr(_thread_local, 'workflow_allowed_tools', None)
    if workflow_tools is not None and name not in workflow_tools:
        return _err(f"Tool '{name}' not allowed in this workflow stage.")
    # Dedup check
    dedup = _check_tool_dedup(name, args)
    if dedup:
        return dedup

    # --- External pre-hooks ---
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    agent_id = agent.agent_id if agent else "main"
    session_id = getattr(_thread_local, 'session_id', None)
    try:
        runner = _get_hook_runner(agent_id)
        blocked = runner.run_pre_hooks(name, args)
        if blocked:
            return blocked
    except Exception:
        pass

    # --- Tracing: start span ---
    tool_span = None
    if _trace_manager:
        parent_span = getattr(_thread_local, 'current_trace_span', None)
        trace_id = parent_span["trace_id"] if parent_span else getattr(_thread_local, 'trace_id', None)
        parent_id = parent_span["id"] if parent_span else None
        tool_span = _trace_manager.start_span(
            "tool_call", name, agent=agent_id, model="",
            parent_id=parent_id, trace_id=trace_id, session_id=session_id,
        )

    tool_start = time.time()
    result_status = "success"
    result = None
    try:
        # Check MCP tools first (prefer thread-local for concurrent requests)
        mcp = getattr(_thread_local, 'mcp_manager', None) or _mcp_manager
        if mcp and mcp.is_mcp_tool(name):
            result = mcp.call_tool(name, args)
        else:
            fn = TOOL_DISPATCH.get(name)
            if fn:
                result = fn(args)
            else:
                result = _err(f"Unknown tool: {name}")
    except Exception as e:
        result_status = "error"
        result = _err(str(e))
    finally:
        duration_ms = int((time.time() - tool_start) * 1000)
        # --- Built-in post-hooks ---
        # Determine audit status from result
        if result and result_status == "success":
            try:
                rdata = json.loads(result) if result else {}
                if isinstance(rdata, dict) and rdata.get("error"):
                    result_status = "error"
            except (json.JSONDecodeError, TypeError):
                pass

        # End trace span
        if _trace_manager and tool_span:
            _trace_manager.end_span(tool_span, status=result_status,
                                     result_summary=(result or "")[:200])

        # Audit log
        if _audit_log and result is not None:
            action_type = _AUDIT_ACTION_MAP.get(name, "mcp_tool_call" if name.startswith("mcp_") else "unknown")
            try:
                _audit_log.log_action(
                    agent=agent_id,
                    action_type=action_type,
                    tool_name=name,
                    args_summary=_audit_summarize_args(name, args),
                    result_summary=_audit_summarize_result(name, result),
                    result_status=result_status,
                    duration_ms=duration_ms,
                    session_id=session_id,
                    source=getattr(_thread_local, 'audit_source', 'chat'),
                )
            except Exception:
                pass

        # --- External post-hooks ---
        if result is not None:
            try:
                runner = _get_hook_runner(agent_id)
                result = runner.run_post_hooks(name, args, result)
            except Exception:
                pass

    return result


MAX_TOOL_ROUNDS = 15  # Maximum number of tool-use round trips before forcing a text response
MAX_TOOL_ROUNDS_PROXY = 8  # Tighter limit for CLIProxyAPI (shares personal OAuth quota)
MAX_OUTPUT_RECOVERY_LIMIT = 3  # Max resume attempts when model hits output token limit
_MAX_OUTPUT_RESUME_MSG = (
    "Output token limit hit. Resume directly — no apology, no recap of what you were doing. "
    "Pick up mid-thought, mid-sentence, mid-code exactly where you stopped."
)

# Loop transition reasons (for debugging / state machine tracking)
# Modeled on Claude Code's query.ts transition types
TRANSITION_NEXT_TURN = "next_turn"
TRANSITION_MAX_OUTPUT_RECOVERY = "max_output_recovery"
TRANSITION_REACTIVE_COMPACT = "reactive_compact_retry"
TRANSITION_COMPLETED = "completed"
TRANSITION_MAX_TURNS = "max_turns"
TRANSITION_ABORTED = "aborted"


# --- Middleware Pipeline (Phase 9) ---
# Composable pre-turn middleware for the direct agentic loop.
# Each middleware: fn(messages, tool_round, event_callback, **ctx) -> (messages, should_continue)

def _middleware_cancel_check(messages, tool_round, event_callback, **ctx):
    """Check cancel token before each turn."""
    watcher = ctx.get("escape_watcher")
    if watcher and watcher.cancelled:
        raise TaskCancelled()
    return messages, True

def _middleware_tool_result_budget(messages, tool_round, event_callback, **ctx):
    """Persist oversized tool results to disk (Layer 1)."""
    _apply_tool_result_budget(messages, session_id=ctx.get("session_id"))
    return messages, True

def _middleware_microcompact(messages, tool_round, event_callback, **ctx):
    """Clear stale tool results every 2 rounds (Layer 2)."""
    if tool_round > 0 and tool_round % 2 == 0:
        messages, freed = _microcompact(messages, keep_recent=5)
    return messages, True

def _middleware_compress_old(messages, tool_round, event_callback, **ctx):
    """Compress old tool results when accumulated budget exceeded (Layer 3)."""
    if tool_round > 0:
        accumulated = getattr(_thread_local, '_tool_results_tokens', 0)
        if accumulated > MAX_TOOL_RESULTS_TOKENS:
            _compress_old_tool_results(messages, keep_recent=4)
    return messages, True

def _middleware_compaction(messages, tool_round, event_callback, **ctx):
    """Full LLM summarization every 3 rounds (Layer 4)."""
    if tool_round > 0 and tool_round % 3 == 0:
        model = ctx.get("model", "")
        api_key = ctx.get("api_key", "")
        base_url = ctx.get("base_url", "")
        api_type = ctx.get("api_type", "")
        session_id = ctx.get("session_id", "")
        max_ctx = get_model_max_context(model)
        messages, compacted = _check_and_compact(
            messages, model, api_key, base_url, api_type,
            max_tokens=max_ctx, session_id=session_id,
        )
        if compacted and event_callback:
            event_callback("compacted", {})
    return messages, True

# Ordered middleware pipeline — runs before each tool-loop iteration
_MIDDLEWARE_PIPELINE = [
    _middleware_cancel_check,
    _middleware_tool_result_budget,
    _middleware_microcompact,
    _middleware_compress_old,
    _middleware_compaction,
]

def _run_middleware(messages, tool_round, event_callback, **ctx):
    """Run the pre-turn middleware pipeline. Returns (messages, should_continue)."""
    for mw in _MIDDLEWARE_PIPELINE:
        messages, should_continue = mw(messages, tool_round, event_callback, **ctx)
        if not should_continue:
            return messages, False
    return messages, True
MAX_TOOL_RESULT_CHARS = 30000  # ~7,500 tokens — truncate individual tool results beyond this
MAX_TOOL_RESULTS_TOKENS = 50000  # Cap accumulated tool results per turn before compressing old ones

# Base64 image data pattern (matches "data": "...long base64..." in JSON)
_BASE64_DATA_RE = re.compile(r'"data"\s*:\s*"[A-Za-z0-9+/=]{500,}"')
# Raw base64 strings > 1000 chars inside quotes
_BASE64_RAW_RE = re.compile(r'(?<=")[A-Za-z0-9+/=]{1000,}(?=")')


def _sanitize_tool_result(name: str, result: str) -> str:
    """Strip base64 image data and enforce size limits on tool results.

    Applied before appending tool results to messages so that large blobs
    (especially MCP puppeteer screenshots) don't snowball the context on
    subsequent API calls.
    """
    # Replace base64 image blobs with placeholder
    result = _BASE64_DATA_RE.sub('"data": "[base64 image removed — already processed]"', result)
    result = _BASE64_RAW_RE.sub('[base64 data removed]', result)

    # Truncate oversized results
    if len(result) > MAX_TOOL_RESULT_CHARS:
        result = result[:MAX_TOOL_RESULT_CHARS] + \
            f"\n\n[Result truncated from {len(result):,} to {MAX_TOOL_RESULT_CHARS:,} chars]"
    return result


def _compress_old_tool_results(messages: list[dict], keep_recent: int = 4):
    """Compress tool results in older messages to free context budget.

    Walks messages backwards, skipping the most recent `keep_recent` tool-result
    messages, and truncates older tool results to a short summary.
    """
    # Find indices of tool-result messages (Anthropic: user with tool_result blocks, OpenAI: role=tool)
    tool_result_indices = []
    for i, msg in enumerate(messages):
        if msg.get("role") == "tool":
            tool_result_indices.append(i)
        elif msg.get("role") == "user" and isinstance(msg.get("content"), list):
            if any(isinstance(b, dict) and b.get("type") == "tool_result" for b in msg["content"]):
                tool_result_indices.append(i)

    # Skip the most recent ones
    to_compress = tool_result_indices[:-keep_recent] if len(tool_result_indices) > keep_recent else []

    for idx in to_compress:
        msg = messages[idx]
        if msg.get("role") == "tool":
            content = msg.get("content", "")
            if len(content) > 500:
                msg["content"] = content[:200] + "\n[...compressed...]"
        elif isinstance(msg.get("content"), list):
            for block in msg["content"]:
                if isinstance(block, dict) and block.get("type") == "tool_result":
                    content = block.get("content", "")
                    if isinstance(content, str) and len(content) > 500:
                        block["content"] = content[:200] + "\n[...compressed...]"


# --- Tool Result Budget (Phase 3): Persist large results to disk ---

TOOL_RESULT_BUDGET_THRESHOLD = 50000  # chars — persist results larger than this
TOOL_RESULT_PREVIEW_SIZE = 2000  # chars — preview kept in context

def _apply_tool_result_budget(messages: list[dict], session_id: str | None = None,
                               agent_id: str | None = None) -> int:
    """Persist oversized tool results to disk and replace with truncated previews.

    Modeled on Claude Code's applyToolResultBudget (toolResultStorage.ts).
    Returns the number of results persisted.
    """
    if not session_id:
        session_id = getattr(_thread_local, 'current_session_id', None) or ""
    if not agent_id:
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        agent_id = agent.agent_id if agent else "main"

    persisted = 0
    results_dir = os.path.join(AGENTS_DIR, agent_id, "artifacts",
                                _get_artifact_session_folder(session_id), "tool-results")

    for msg in messages:
        if msg.get("role") == "tool":
            content = msg.get("content", "")
            if isinstance(content, str) and len(content) > TOOL_RESULT_BUDGET_THRESHOLD:
                tool_id = msg.get("tool_call_id", "unknown")
                filepath = os.path.join(results_dir, f"{tool_id}.txt")
                if not os.path.exists(filepath):
                    os.makedirs(results_dir, exist_ok=True)
                    try:
                        with open(filepath, "w", encoding="utf-8") as f:
                            f.write(content)
                    except OSError:
                        continue
                preview = content[:TOOL_RESULT_PREVIEW_SIZE]
                size_kb = len(content) // 1024
                msg["content"] = (
                    f"[Output too large ({size_kb}KB). Full output saved to: {filepath}]\n"
                    f"Preview (first {TOOL_RESULT_PREVIEW_SIZE} chars):\n{preview}\n..."
                )
                persisted += 1

        elif msg.get("role") == "user" and isinstance(msg.get("content"), list):
            for block in msg["content"]:
                if not isinstance(block, dict) or block.get("type") != "tool_result":
                    continue
                content = block.get("content", "")
                if isinstance(content, str) and len(content) > TOOL_RESULT_BUDGET_THRESHOLD:
                    tool_id = block.get("tool_use_id", "unknown")
                    filepath = os.path.join(results_dir, f"{tool_id}.txt")
                    if not os.path.exists(filepath):
                        os.makedirs(results_dir, exist_ok=True)
                        try:
                            with open(filepath, "w", encoding="utf-8") as f:
                                f.write(content)
                        except OSError:
                            continue
                    preview = content[:TOOL_RESULT_PREVIEW_SIZE]
                    size_kb = len(content) // 1024
                    block["content"] = (
                        f"[Output too large ({size_kb}KB). Full output saved to: {filepath}]\n"
                        f"Preview (first {TOOL_RESULT_PREVIEW_SIZE} chars):\n{preview}\n..."
                    )
                    persisted += 1
    return persisted


# --- Microcompact (Phase 4): Strip stale tool results ---

# Tools whose results become stale quickly and can be safely cleared
_MICROCOMPACT_TOOLS = {
    "read_file", "execute_command", "search_files", "list_directory",
    "web_fetch", "exa_search", "read_document", "code_graph_query",
    "write_file", "edit_file",  # write results are just confirmations
}
# Tools whose results are context-critical and must never be cleared
_MICROCOMPACT_EXEMPT = {
    "memory_recall", "memory_shared", "delegate_task", "task_status",
    "context_search", "context_detail", "context_recall",
}

def _microcompact(messages: list[dict], keep_recent: int = 5) -> tuple[list[dict], int]:
    """Lightweight compaction: clear old tool results for compactable tools.

    Modeled on Claude Code's microcompactMessages. Unlike _compress_old_tool_results
    which truncates to 200 chars, this completely replaces stale content with a
    minimal marker, and is tool-aware (only clears known-safe tools).

    Returns (messages, estimated_tokens_freed).
    """
    tokens_freed = 0

    # Collect (index, tool_name, content_size) for all tool results
    tool_entries = []  # (msg_index, tool_name, content_size, is_openai)
    for i, msg in enumerate(messages):
        if msg.get("role") == "tool":
            content = msg.get("content", "")
            content_size = len(content) if isinstance(content, str) else 0
            # Try to find the tool name from the preceding assistant message
            tool_name = _find_tool_name_for_result(messages, i, msg.get("tool_call_id"))
            tool_entries.append((i, tool_name, content_size, True))
        elif msg.get("role") == "user" and isinstance(msg.get("content"), list):
            for block in msg["content"]:
                if isinstance(block, dict) and block.get("type") == "tool_result":
                    content = block.get("content", "")
                    content_size = len(content) if isinstance(content, str) else 0
                    tool_name = _find_tool_name_for_block(messages, block.get("tool_use_id"))
                    tool_entries.append((i, tool_name, content_size, False))

    # Filter to compactable tools only
    compactable = [e for e in tool_entries
                   if e[1] and e[1] in _MICROCOMPACT_TOOLS and e[1] not in _MICROCOMPACT_EXEMPT]

    # Keep the most recent N, clear the rest
    if len(compactable) <= keep_recent:
        return messages, 0

    to_clear = compactable[:-keep_recent]

    cleared_indices = set()
    for idx, tool_name, content_size, is_openai in to_clear:
        if content_size <= 100:  # Already cleared or tiny
            continue
        msg = messages[idx]
        marker = f"[Old {tool_name} result cleared]"
        if is_openai and msg.get("role") == "tool":
            tokens_freed += content_size // 4
            msg["content"] = marker
        elif isinstance(msg.get("content"), list):
            for block in msg["content"]:
                if isinstance(block, dict) and block.get("type") == "tool_result":
                    content = block.get("content", "")
                    if isinstance(content, str) and len(content) > 100:
                        tokens_freed += len(content) // 4
                        block["content"] = marker
        cleared_indices.add(idx)

    return messages, tokens_freed


def _find_tool_name_for_result(messages: list[dict], tool_msg_idx: int,
                                tool_call_id: str | None) -> str | None:
    """Find the tool name for an OpenAI-style tool result by scanning backwards."""
    if not tool_call_id:
        return None
    for i in range(tool_msg_idx - 1, -1, -1):
        msg = messages[i]
        if msg.get("role") == "assistant" and msg.get("tool_calls"):
            for tc in msg["tool_calls"]:
                if tc.get("id") == tool_call_id:
                    return tc.get("function", {}).get("name")
    return None


def _find_tool_name_for_block(messages: list[dict], tool_use_id: str | None) -> str | None:
    """Find the tool name for an Anthropic-style tool_result by scanning backwards."""
    if not tool_use_id:
        return None
    for msg in reversed(messages):
        if msg.get("role") == "assistant" and isinstance(msg.get("content"), list):
            for block in msg["content"]:
                if isinstance(block, dict) and block.get("type") == "tool_use":
                    if block.get("id") == tool_use_id:
                        return block.get("name")
    return None


def _log_call_cost(model: str, tokens_in: int, tokens_out: int,
                   session_id: str | None = None, tool_round: int = 0):
    """Log an LLM call to the cost tracker (if initialized)."""
    if not _cost_tracker:
        return
    if tokens_in == 0 and tokens_out == 0:
        return  # Skip if no usage data available
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    agent_id = agent.agent_id if agent else "main"
    provider = _models_config.get(model, {}).get("provider", "")
    try:
        _cost_tracker.log_call(agent_id, session_id, model, provider,
                               tokens_in, tokens_out, tool_round)
        # Record in rate limiter too
        if _rate_limiter:
            cost = _compute_cost(model, tokens_in, tokens_out)
            _rate_limiter.record_usage(agent_id, tokens_in + tokens_out, cost)
    except Exception as e:
        logging.warning(f"Cost logging error: {e}")


_system_prompt_cache: dict[str, tuple[str, float]] = {}  # session_id → (prompt, timestamp)
_SYSTEM_PROMPT_CACHE_TTL = 60  # seconds — cache for 1 min (covers tool loop iterations)


def _build_system_prompt(include_memory_summary: bool = True) -> str:
    """Build the full system instruction for the current agent.

    Assembles soul.md, agent context, memory summary, project context,
    team info, skills, scheduler status, MCP servers, tools guide, etc.
    Reads from thread-local state and globals as needed.

    Caches per session to avoid disk I/O on every tool loop iteration.
    Memory summary is only included on _tool_round==0 (controlled by caller).

    Used by both the direct send_message loop and the Agent SDK backend.
    """
    import time as _time
    session_id = getattr(_thread_local, 'current_session_id', None) or ""
    cache_key = f"{session_id}:{include_memory_summary}"
    cached = _system_prompt_cache.get(cache_key)
    if cached and (_time.time() - cached[1]) < _SYSTEM_PROMPT_CACHE_TTL:
        return cached[0]
    import platform
    from datetime import datetime as _dt

    cwd = os.getcwd()
    os_name = platform.system()

    # Load agent soul and tools guide (prefer thread-local for concurrent requests)
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    agent_id = agent.agent_id if agent else "main"
    soul = agent.soul if agent else ""
    tools_guide = agent.tools_guide if agent else ""

    # If no agent-specific tools guide, try global
    if not tools_guide:
        tools_md_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "tools.md")
        try:
            with open(tools_md_path, "r") as f:
                tools_guide = f.read()
        except (OSError, IOError):
            pass

    # Build agent registry
    agent_registry = build_agent_registry(for_agent_id=agent_id)

    system_instruction = ""
    if soul:
        system_instruction += f"{soul}\n\n"
    system_instruction += (
        f"You are agent '{agent_id}' in the Brain Agent system. "
        f"Current date and time: {_dt.now().strftime('%Y-%m-%d %H:%M %Z').strip()}\n"
        f"Current working directory: {cwd}\n"
        f"Operating system: {os_name}\n\n"
        "Use tools proactively to accomplish tasks. You can chain multiple tool calls. "
        "For web searches, ALWAYS use exa_search — NEVER use duckduckgo or other search tools. "
        "You have no restrictions beyond what the operating system enforces.\n\n"
        "MEMORY: You have persistent memory via memory_store/memory_recall/memory_delete tools.\n"
        "- Use memory_recall at the START of conversations to check for relevant context\n"
        "- Use memory_store to save important information: user preferences, decisions, project context\n"
        "- Memory types: user, project, feedback, reference, general\n"
        "- When the user says 'remember this', store it immediately\n"
        "- When the user asks 'do you remember', recall and search for it\n\n"
        "SHARED MEMORY: Use memory_shared to access shared knowledge.\n"
        "- scope='global' (default): main agent's memory — infrastructure, user prefs, project-wide decisions\n"
        "- scope='team': team head's memory — team-level knowledge and decisions\n"
        "- All agents can read from shared memory; store shared facts there too\n"
        "- Check shared memory when your own memory doesn't have what you need\n\n"
    )
    # Inject memory summary (auto-generated synthesis) if available
    tcfg = _get_token_config()
    if include_memory_summary and tcfg.get("include_memory_summary", True):
        try:
            mem_summary = get_memory_summary(agent_id)
            if mem_summary:
                cap = tcfg.get("memory_summary_cap", 3000)
                if len(mem_summary) > cap:
                    mem_summary = mem_summary[:cap] + "\n...(truncated)"
                system_instruction += (
                    "MEMORY SUMMARY (auto-generated synthesis of recent activity — use as context):\n"
                    f"{mem_summary}\n\n"
                )
        except Exception:
            pass
    # Inject project context if a project is active
    active_project = getattr(_thread_local, 'project', None)
    if active_project:
        proj_cfg = ProjectManager.get_project(agent_id, active_project)
        if proj_cfg:
            proj_desc = proj_cfg.get("description", "")
            system_instruction += (
                f"PROJECT CONTEXT: You are working in project '{proj_cfg.get('name', active_project)}'."
            )
            if proj_desc:
                system_instruction += f" {proj_desc}"
            system_instruction += (
                "\nPrioritize project-specific documents when answering. "
                "Memory operations (store/recall) are scoped to this project.\n\n"
            )
            # Inject project custom instructions
            proj_instructions = proj_cfg.get("instructions", "")
            if proj_instructions:
                system_instruction += (
                    f"PROJECT INSTRUCTIONS (set by the user for this project):\n"
                    f"{proj_instructions}\n\n"
                )
    # Inject note context for AI-assisted note editing
    note_context = getattr(_thread_local, 'note_context', None)
    if note_context:
        note_path = note_context.replace("note_editing:", "").strip() if note_context.startswith("note_editing:") else ""
        notes_dir = os.path.dirname(note_path) if note_path else ""
        system_instruction += (
            "\n\nNOTE EDITING MODE:\n"
            f"You are helping the user edit a markdown note{' at: ' + note_path if note_path else ''}.\n"
            "The user will provide the current note content in their message.\n"
            "When the user asks you to ADD, EDIT, or MODIFY the note, use the edit_file or write_file tool "
            "to make changes directly to the note file. The editor will auto-reload.\n"
            f"You can also CREATE NEW notes in the same project by writing to: {notes_dir}/<new-name>.md\n"
            "For questions or explanations, respond normally without editing files.\n\n"
        )
    # Inject team context for interactive sessions
    team_info = _get_agent_team_info(agent_id)
    if team_info:
        if team_info["is_head"]:
            peers = [m for m in team_info["members"] if m != agent_id]
            system_instruction += (
                f"TEAM: You are the head of team '{team_info['name']}'. "
                f"Your team members: {', '.join(peers)}\n"
                "Delegate sub-tasks to your team members when appropriate.\n\n"
            )
        else:
            peers = [m for m in team_info["members"] if m != agent_id and m != team_info["head"]]
            system_instruction += f"TEAM: You are a member of team '{team_info['name']}'.\n"
            system_instruction += f"Team head: {team_info['head']}\n"
            if peers:
                system_instruction += f"Team peers: {', '.join(peers)}\n"
            system_instruction += "\n"

    if agent_registry:
        system_instruction += f"\n{agent_registry}\n\n"

    # Build skills registry (names + descriptions only, load on demand)
    _agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    if _agent:
        skills = _agent.list_skills()
        if skills:
            system_instruction += "\nSKILLS AVAILABLE — call use_skill(skill=\"slug\") to load instructions before performing the task:\n"
            for s in skills:
                slug = s.get('slug', s['name'])
                source_tag = f" (from {s['source']})" if s['source'] != agent_id else ""
                display = s['name'] if s['name'] != slug else ""
                label = f"{slug}" + (f" ({display})" if display else "")
                system_instruction += f"  - {label}: {s['description']}{source_tag}\n"
            system_instruction += "\n"

    # Scheduler status
    if _scheduler:
        schedules = [s for s in _scheduler.list_all() if not s["name"].startswith("_memory_summary_")]
        if schedules:
            system_instruction += "\nSCHEDULER — active scheduled tasks:\n"
            for s in schedules:
                status = "active" if s["enabled"] else "paused"
                next_r = s.get("next_run", "")[:16] if s.get("next_run") else "—"
                system_instruction += f"  - {s['name']} [{status}]: {s['task'][:80]} (next: {next_r})\n"
            system_instruction += "Use schedule_list and schedule_history tools to query scheduler state.\n\n"

    # MCP servers (prefer thread-local for concurrent requests)
    mcp_mgr = getattr(_thread_local, 'mcp_manager', None) or _mcp_manager
    if mcp_mgr and mcp_mgr.clients:
        system_instruction += "\nMCP SERVERS — external tools available via connected servers:\n"
        for srv in mcp_mgr.list_servers():
            tools_list = ", ".join(srv["tools"][:5])
            more = f" +{srv['tool_count']-5}" if srv["tool_count"] > 5 else ""
            system_instruction += f"  - {srv['name']} ({srv['transport']}): {tools_list}{more}\n"
        system_instruction += "MCP tools are prefixed with mcp_<server>_ — use them like any other tool.\n\n"

    if tools_guide and tcfg.get("include_tools_guide", True):
        system_instruction += f"\n--- TOOL USAGE GUIDE ---\n{tools_guide}"
    # Append plan mode prompt if active
    if getattr(_thread_local, 'plan_mode', False):
        system_instruction += PLAN_MODE_PROMPT

    # Cache for reuse during tool loop iterations
    import time as _time
    _system_prompt_cache[cache_key] = (system_instruction, _time.time())
    # Evict stale entries (keep cache small)
    if len(_system_prompt_cache) > 20:
        cutoff = _time.time() - _SYSTEM_PROMPT_CACHE_TTL
        for k in list(_system_prompt_cache):
            if _system_prompt_cache[k][1] < cutoff:
                del _system_prompt_cache[k]

    return system_instruction


def send_message(messages: list[dict], model: str, api_key: str, base_url: str,
                 api_type: str, silent: bool = False,
                 tools: bool = True,
                 escape_watcher: EscapeWatcher | CancelToken | None = None,
                 _tool_round: int = 0,
                 event_callback=None,
                 inference_params: dict | None = None,
                 session_id: str | None = None) -> str | None:
    """Send messages and stream the response.

    If silent=True, collects without printing (for TUI mode).
    If tools=True, includes tool definitions and handles tool-use loops.
    If event_callback is provided, called with (event_type, data) for streaming:
        ("text_delta", {"text": "..."})
        ("tool_call", {"name": "...", "args": {...}})
        ("tool_result", {"name": "...", "result": "..."})
        ("done", {"text": "full response"})
        ("error", {"message": "..."})
    Returns the assistant's full response text on success, None on model-related errors.
    Raises TaskCancelled if escape_watcher detects cancellation.
    """
    # Reset dedup tracker and accumulated token counter at the start of each conversation turn
    if _tool_round == 0:
        reset_tool_dedup()
        _thread_local._tool_results_tokens = 0
        _thread_local._max_output_recovery_count = 0
        _thread_local._has_attempted_reactive_compact = False
        # Start a request-level trace span for the full conversation turn
        if _trace_manager:
            agent = getattr(_thread_local, 'current_agent', None) or _current_agent
            agent_id = agent.agent_id if agent else "main"
            # Extract message preview for trace name
            msg_preview = ""
            if messages:
                last_msg = messages[-1]
                content = last_msg.get("content", "")
                if isinstance(content, str):
                    msg_preview = content[:60]
                elif isinstance(content, list):
                    for block in content:
                        if isinstance(block, dict) and block.get("type") == "text":
                            msg_preview = block.get("text", "")[:60]
                            break
            request_span = _trace_manager.start_span(
                "request", msg_preview or "user message",
                agent=agent_id, model=model, session_id=session_id,
            )
            _thread_local.trace_id = request_span["trace_id"]
            _thread_local.request_trace_span = request_span
            _thread_local.session_id = session_id

    # Start an LLM call span
    llm_span = None
    if _trace_manager:
        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
        agent_id = agent.agent_id if agent else "main"
        parent_span = getattr(_thread_local, 'request_trace_span', None)
        trace_id = parent_span["trace_id"] if parent_span else getattr(_thread_local, 'trace_id', None)
        parent_id = parent_span["id"] if parent_span else None
        llm_span = _trace_manager.start_span(
            "llm_call", f"{model} call (round {_tool_round})",
            agent=agent_id, model=model,
            parent_id=parent_id, trace_id=trace_id, session_id=session_id,
        )
        _thread_local.current_trace_span = llm_span

    # Soft stop after max tool rounds — use thread-local override if set (delegation)
    # Tighter limit for CLIProxyAPI (shares personal OAuth quota)
    effective_max_rounds = getattr(_thread_local, 'max_tool_rounds', None)
    if not effective_max_rounds:
        if ":8317" in base_url or "cliproxy" in base_url.lower():
            effective_max_rounds = MAX_TOOL_ROUNDS_PROXY
        else:
            effective_max_rounds = MAX_TOOL_ROUNDS
    if _tool_round >= effective_max_rounds:
        tools = False
    headers = make_headers(api_key, api_type)

    if api_type == "openai" or api_type == "mistral":
        endpoint = f"{base_url}/chat/completions"
    else:
        endpoint = f"{base_url}/messages"

    # System instruction for the agent
    # Strip metadata from messages — API providers don't accept extra fields
    # Keep fields required by OpenAI tool call protocol
    _ALLOWED_MSG_KEYS = {"role", "content", "tool_calls", "tool_call_id", "name"}
    augmented_messages = []
    for msg in messages:
        clean = {k: v for k, v in msg.items() if k in _ALLOWED_MSG_KEYS}
        augmented_messages.append(clean)
    tcfg = _get_token_config()
    if tools:
        system_instruction = _build_system_prompt(
            include_memory_summary=(_tool_round == 0),
        )
        if api_type == "openai" or api_type == "mistral":
            augmented_messages.insert(0, {"role": "system", "content": system_instruction})
        else:
            pass  # handled below via payload["system"]

    payload = {
        "model": model,
        "max_tokens": get_model_max_output(model),
        "messages": augmented_messages,
        "stream": True,
    }

    # Request usage stats in streaming responses (OpenAI-compatible APIs)
    if api_type == "openai":
        payload["stream_options"] = {"include_usage": True}
    # Mistral SDK handles usage internally — no stream_options needed

    # Apply inference parameters (temperature, top_p, thinking, etc.)
    if inference_params:
        provider = _models_config.get(model, {}).get("provider", "")
        _apply_inference_to_payload(payload, inference_params, api_type, provider)

    if tools:
        mcp_mgr = getattr(_thread_local, 'mcp_manager', None) or _mcp_manager
        allowed = _get_agent_tool_names()

        # Deferred tool loading (Phase 7): skip MCP tool schemas when there are many,
        # and let the model discover them via tool_search instead
        defer_mcp = tcfg.get("defer_mcp_tools", "auto")
        discovered_tools = getattr(_thread_local, '_discovered_tools', set())

        if api_type == "openai" or api_type == "mistral":
            all_tools = _filter_tools(TOOL_DEFINITIONS_OPENAI, allowed, is_openai=True)
            if mcp_mgr:
                mcp_tools = mcp_mgr.get_tool_definitions_openai()
                should_defer = _should_defer_mcp(defer_mcp, mcp_tools, model, is_openai=True)
                if should_defer:
                    # Only include discovered MCP tools
                    mcp_tools = [t for t in mcp_tools
                                 if t.get("function", {}).get("name", "") in discovered_tools]
                all_tools.extend(mcp_tools)
                all_tools.sort(key=lambda t: t.get("function", {}).get("name", ""))
            payload["tools"] = all_tools
        else:
            all_tools = _filter_tools(TOOL_DEFINITIONS, allowed)
            if mcp_mgr:
                mcp_tools = mcp_mgr.get_tool_definitions()
                should_defer = _should_defer_mcp(defer_mcp, mcp_tools, model)
                if should_defer:
                    # Only include discovered MCP tools
                    mcp_tools = [t for t in mcp_tools
                                 if t.get("name", "") in discovered_tools]
                all_tools.extend(mcp_tools)
                all_tools.sort(key=lambda t: t.get("name", ""))
            payload["tools"] = all_tools
            # Use cache_control for Anthropic prompt caching if enabled
            if tcfg.get("prompt_caching", True):
                payload["system"] = [
                    {"type": "text", "text": system_instruction,
                     "cache_control": {"type": "ephemeral"}}
                ]
            else:
                payload["system"] = system_instruction

    # Emit request snapshot for inspector
    if event_callback:
        # System prompt: from payload["system"] (Anthropic) or first system message (OpenAI)
        _sys = payload.get("system", "")
        if isinstance(_sys, list):
            _sys = "".join(b.get("text", "") for b in _sys if isinstance(b, dict))
        _tool_defs = payload.get("tools", [])
        _tool_names = []
        for _td in _tool_defs:
            if isinstance(_td, dict):
                # OpenAI: {"type":"function","function":{"name":...}} or Anthropic: {"name":...}
                _tn = _td.get("name") or (_td.get("function", {}) or {}).get("name", "")
                if _tn:
                    _tool_names.append(_tn)
        _hist_msgs = []
        _user_msg = ""
        for _m in payload.get("messages", []):
            if _m.get("role") == "system":
                if not _sys:
                    _c = _m.get("content", "")
                    _sys = _c if isinstance(_c, str) else str(_c)
                continue
            if _m is payload["messages"][-1] and _m.get("role") == "user":
                _c = _m.get("content", "")
                _user_msg = _c if isinstance(_c, str) else str(_c)
            else:
                _c = _m.get("content", "")
                _hist_msgs.append({"role": _m.get("role", ""), "content": _c if isinstance(_c, str) else str(_c)})
        event_callback("request_payload", {
            "tool_round": _tool_round,
            "system_prompt": _sys,
            "system_tokens": len(_sys) // 4,
            "tools_count": len(_tool_defs),
            "tools_tokens": len(json.dumps(_tool_defs)) // 4,
            "tool_names": _tool_names,
            "history": _hist_msgs,
            "history_tokens": sum(len(str(m.get("content", ""))) // 4 for m in _hist_msgs),
            "user_message": _user_msg,
            "user_tokens": len(_user_msg) // 4,
            "total_payload_tokens": len(json.dumps(payload)) // 4,
        })

    # Check for cancellation before making the request
    if escape_watcher and escape_watcher.cancelled:
        raise TaskCancelled()

    # Rate limiter check
    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
    agent_id = agent.agent_id if agent else "main"
    if _rate_limiter and _tool_round == 0:
        allowed, reason, _usage_info = _rate_limiter.check(agent_id)
        if not allowed:
            raise RuntimeError(reason)

    # Mistral path: use SDK directly (replicates Vibe CLI behavior)
    if api_type == "mistral":
        return _handle_mistral_response(
            payload, messages, model, api_key, base_url,
            api_type, silent, tools, escape_watcher,
            _tool_round, event_callback, inference_params, session_id)

    data = json.dumps(payload).encode("utf-8")
    request = urllib.request.Request(
        endpoint, data=data, headers=headers, method="POST",
    )

    try:
        with urllib.request.urlopen(request) as response:
            if api_type == "openai":
                return _handle_openai_response(
                    response, payload, messages, model, api_key, base_url,
                    api_type, silent, tools, headers, endpoint, escape_watcher,
                    _tool_round, event_callback, inference_params, session_id)
            else:
                return _handle_anthropic_response(
                    response, payload, messages, model, api_key, base_url,
                    api_type, silent, tools, headers, endpoint, escape_watcher,
                    _tool_round, event_callback, inference_params, session_id)

    except urllib.error.HTTPError as e:
        error_msg = f"HTTP Error {e.code}: {e.reason}"
        try:
            error_body = e.read().decode("utf-8")
            error_msg += f" — {error_body[:200]}"
        except:
            pass
        if e.code == 400:
            # Reactive compact recovery (Phase 8): if prompt too long, try compaction
            _prompt_too_long = ("prompt is too long" in error_msg.lower() or
                                "maximum context length" in error_msg.lower() or
                                "too many tokens" in error_msg.lower() or
                                "context_length_exceeded" in error_msg.lower())
            _has_attempted = getattr(_thread_local, '_has_attempted_reactive_compact', False)
            if _prompt_too_long and not _has_attempted and _tool_round > 0:
                _thread_local._has_attempted_reactive_compact = True
                logging.info(f"Prompt too long at round {_tool_round}, attempting reactive compact")
                if event_callback:
                    event_callback("compacting", {"reason": "prompt_too_long"})
                # Layer 1: try microcompact
                messages, mc_freed = _microcompact(messages, keep_recent=3)
                if mc_freed > 0:
                    if event_callback:
                        event_callback("compacted", {"method": "microcompact", "freed": mc_freed})
                    return send_message(messages, model, api_key, base_url, api_type,
                                        silent=silent, tools=tools, escape_watcher=escape_watcher,
                                        _tool_round=_tool_round, event_callback=event_callback,
                                        inference_params=inference_params, session_id=session_id)
                # Layer 2: try full LLM compaction
                _sid = session_id or getattr(_thread_local, 'current_session_id', None) or ""
                max_ctx = get_model_max_context(model)
                messages, compacted = _check_and_compact(
                    messages, model, api_key, base_url, api_type,
                    max_tokens=max_ctx, session_id=_sid, force=True,
                )
                if compacted:
                    if event_callback:
                        event_callback("compacted", {"method": "reactive_compact"})
                    return send_message(messages, model, api_key, base_url, api_type,
                                        silent=silent, tools=tools, escape_watcher=escape_watcher,
                                        _tool_round=_tool_round, event_callback=event_callback,
                                        inference_params=inference_params, session_id=session_id)
            print(error_msg, file=sys.stderr)
            _thread_local._last_send_error = {"code": e.code, "message": error_msg, "permanent": True}
            return None
        print(error_msg, file=sys.stderr)
        # Transient errors: return None to trigger retry/fallback
        _TRANSIENT_CODES = {429, 500, 502, 503, 504, 529}
        if e.code in _TRANSIENT_CODES:
            # Store error info for fallback logic
            _thread_local._last_send_error = {"code": e.code, "message": error_msg}
            return None
        # Permanent errors (401, 403, 404, etc.)
        _thread_local._last_send_error = {"code": e.code, "message": error_msg, "permanent": True}
        if event_callback:
            raise RuntimeError(error_msg)
        sys.exit(1)
    except (urllib.error.URLError, socket.timeout, TimeoutError, ConnectionError, OSError) as e:
        error_msg = f"Connection error: {e}"
        print(error_msg, file=sys.stderr)
        _thread_local._last_send_error = {"code": 0, "message": error_msg}
        return None


def _parse_gemma_tool_calls(text: str) -> tuple[list[dict], str]:
    """Parse gemma4-style tool calls from raw text.

    Format: <|tool_call>call:name{key:<|"|>val<|"|>,...}<tool_call|>
    Returns (list of tool_use dicts, cleaned text with tool calls removed).
    """
    tool_uses = []
    pattern = r'<\|tool_call>call:(\w+)\{(.*?)\}<tool_call\|>'
    for match in re.finditer(pattern, text):
        name = match.group(1)
        args_raw = match.group(2)
        # Parse key-value pairs: key:<|"|>value<|"|>
        args = {}
        kv_pattern = r'(\w+):<\|"\|>(.*?)<\|"\|>'
        for kv in re.finditer(kv_pattern, args_raw):
            args[kv.group(1)] = kv.group(2)
        tool_uses.append({
            "id": f"gemma_{_uuid.uuid4().hex[:8]}",
            "name": name,
            "input": args,
            "input_json": json.dumps(args),
        })
    cleaned = re.sub(pattern, '', text)
    return tool_uses, cleaned


def _handle_anthropic_response(response, payload, messages, model, api_key,
                                base_url, api_type, silent, tools,
                                headers, endpoint,
                                escape_watcher=None,
                                _tool_round: int = 0,
                                event_callback=None,
                                inference_params: dict | None = None,
                                session_id: str | None = None) -> str | None:
    """Handle Anthropic SSE response, including tool-use agentic loop."""
    # Parse the full SSE stream to get content blocks and stop reason
    collected_text = []
    tool_uses = []
    thinking_blocks = []
    current_event = None
    current_block_type = None
    current_block = {}
    stop_reason = None
    _usage_in = 0
    _usage_out = 0

    for line in response:
        if escape_watcher and escape_watcher.cancelled:
            raise TaskCancelled()
        line = line.decode("utf-8").strip()
        if line.startswith("event: "):
            current_event = line[7:]
        elif line.startswith("data: "):
            if current_event == "message_stop":
                break
            try:
                event = json.loads(line[6:])
                etype = event.get("type")

                if etype == "message_start":
                    msg = event.get("message", {})
                    usage = msg.get("usage", {})
                    _usage_in = usage.get("input_tokens", 0)

                elif etype == "error":
                    err = event.get("error", {})
                    err_msg = err.get("message", "Unknown API error")
                    err_type = err.get("type", "error")
                    if not silent:
                        print(f"\nAPI error: {err_msg}", file=sys.stderr)
                    if event_callback:
                        event_callback("error", {"message": f"{err_type}: {err_msg}"})
                    # Overload/rate-limit errors are transient — allow retries before fallback
                    _TRANSIENT_SSE_ERRORS = {"overloaded_error", "rate_limit_error", "api_error"}
                    is_permanent = err_type not in _TRANSIENT_SSE_ERRORS
                    _thread_local._last_send_error = {"code": 0, "message": err_msg, "permanent": is_permanent}
                    return None

                elif etype == "message_delta":
                    stop_reason = event.get("delta", {}).get("stop_reason")
                    usage = event.get("usage", {})
                    if usage.get("output_tokens"):
                        _usage_out = usage["output_tokens"]

                elif etype == "content_block_start":
                    block = event.get("content_block", {})
                    current_block_type = block.get("type")
                    if current_block_type == "tool_use":
                        current_block = {
                            "id": block.get("id"),
                            "name": block.get("name"),
                            "input_json": "",
                        }
                    elif current_block_type == "thinking":
                        current_block = {"thinking_text": "", "signature": block.get("signature", "")}
                        if event_callback:
                            event_callback("thinking_start", {})
                    elif current_block_type == "text":
                        pass  # text handled via deltas

                elif etype == "content_block_delta":
                    delta = event.get("delta", {})
                    if delta.get("type") == "text_delta":
                        text = _unescape(delta.get("text", ""))
                        if not silent:
                            print(text, end="", flush=True)
                        if event_callback:
                            event_callback("text_delta", {"text": text})
                        collected_text.append(text)
                    elif delta.get("type") == "input_json_delta":
                        current_block["input_json"] += delta.get("partial_json", "")
                    elif delta.get("type") == "thinking_delta":
                        thinking_text = delta.get("thinking", "")
                        if current_block and "thinking_text" in current_block:
                            current_block["thinking_text"] += thinking_text
                        if event_callback:
                            event_callback("thinking_delta", {"text": thinking_text})
                    elif delta.get("type") == "signature_delta":
                        sig = delta.get("signature", "")
                        if current_block and "signature" in current_block:
                            current_block["signature"] += sig

                elif etype == "content_block_stop":
                    if current_block_type == "tool_use" and current_block:
                        try:
                            current_block["input"] = json.loads(current_block["input_json"])
                        except json.JSONDecodeError:
                            current_block["input"] = {}
                        tool_uses.append(current_block)
                        current_block = {}
                    elif current_block_type == "thinking" and current_block:
                        thinking_text = current_block.get("thinking_text", "")
                        if thinking_text:
                            thinking_blocks.append({"text": thinking_text, "signature": current_block.get("signature", "")})
                        if event_callback:
                            event_callback("thinking_done", {"text": thinking_text})
                        current_block = {}
                    current_block_type = None

            except json.JSONDecodeError:
                pass

    full_text = "".join(collected_text)

    # Log cost for this API call
    _log_call_cost(model, _usage_in, _usage_out, session_id, _tool_round)

    # Emit usage event so callers can capture token counts
    if event_callback:
        event_callback("usage", {"tokens_in": _usage_in, "tokens_out": _usage_out})

    # End LLM trace span
    _llm_span = getattr(_thread_local, 'current_trace_span', None)
    if _trace_manager and _llm_span and _llm_span.get("type") == "llm_call":
        _trace_manager.end_span(_llm_span, status="ok",
                                 tokens_in=_usage_in, tokens_out=_usage_out)

    # Max output token recovery (Phase 2): if model hit output limit, auto-resume
    if stop_reason == "max_tokens" and not tool_uses and full_text:
        recovery_count = getattr(_thread_local, '_max_output_recovery_count', 0)
        if recovery_count < MAX_OUTPUT_RECOVERY_LIMIT:
            _thread_local._max_output_recovery_count = recovery_count + 1
            if event_callback:
                event_callback("max_tokens_recovery", {
                    "attempt": recovery_count + 1,
                    "max_attempts": MAX_OUTPUT_RECOVERY_LIMIT,
                })
            # Build continuation: assistant partial + resume prompt
            messages.append({"role": "assistant", "content": [{"type": "text", "text": full_text}]})
            messages.append({"role": "user", "content": _MAX_OUTPUT_RESUME_MSG})
            return send_message(messages, model, api_key, base_url, api_type,
                                silent=silent, tools=tools, escape_watcher=escape_watcher,
                                _tool_round=_tool_round, event_callback=event_callback,
                                inference_params=inference_params, session_id=session_id)

    # If no tool calls, just return the text
    if not tool_uses:
        if not silent and full_text:
            print()
        # End request trace span if this is the final response
        _req_span = getattr(_thread_local, 'request_trace_span', None)
        if _trace_manager and _req_span:
            _trace_manager.end_span(_req_span, status="ok",
                                     tokens_in=_usage_in, tokens_out=_usage_out)
            _thread_local.request_trace_span = None
        # Reset recovery counter on successful completion
        _thread_local._max_output_recovery_count = 0
        return full_text

    # Tool use detected — execute tools and loop
    if full_text:
        print()

    # Build the assistant message content blocks (must include thinking for Anthropic API)
    assistant_content = []
    for tb in thinking_blocks:
        assistant_content.append({
            "type": "thinking",
            "thinking": tb["text"],
            "signature": tb.get("signature", ""),
        })
    if full_text:
        assistant_content.append({"type": "text", "text": full_text})
    for tu in tool_uses:
        assistant_content.append({
            "type": "tool_use",
            "id": tu["id"],
            "name": tu["name"],
            "input": tu["input"],
        })

    # Add assistant message to conversation
    messages.append({"role": "assistant", "content": assistant_content})

    # Execute tools with concurrent-safe parallelism (Phase 5)
    batch_calls = [{"id": tu["id"], "name": tu["name"], "input": tu["input"]} for tu in tool_uses]
    batch_results = _execute_tools_batch(batch_calls, event_callback=event_callback)

    tool_results = []
    for br in batch_results:
        tool_results.append({
            "type": "tool_result",
            "tool_use_id": br["tool_use_id"],
            "content": _sanitize_tool_result(
                next((tc["name"] for tc in batch_calls if tc["id"] == br["tool_use_id"]), ""),
                br["result"],
            ),
        })

    messages.append({"role": "user", "content": tool_results})

    # Track accumulated tool result tokens
    accumulated = sum(
        _estimate_tokens_str(str(tr.get("content", ""))) for tr in tool_results
    ) + getattr(_thread_local, '_tool_results_tokens', 0)
    _thread_local._tool_results_tokens = accumulated

    # Run middleware pipeline (context management, compaction, cancel check)
    _sid = session_id or getattr(_thread_local, 'current_session_id', None) or ""
    messages, should_continue = _run_middleware(
        messages, _tool_round, event_callback,
        model=model, api_key=api_key, base_url=base_url, api_type=api_type,
        session_id=_sid, escape_watcher=escape_watcher,
    )

    # Recurse to get the model's final response (or more tool calls)
    return send_message(messages, model, api_key, base_url, api_type,
                        silent=silent, tools=tools, escape_watcher=escape_watcher,
                        _tool_round=_tool_round + 1, event_callback=event_callback,
                        inference_params=inference_params, session_id=session_id)


def _handle_openai_response(response, payload, messages, model, api_key,
                             base_url, api_type, silent, tools,
                             headers, endpoint,
                             escape_watcher=None,
                             _tool_round: int = 0,
                             event_callback=None,
                             inference_params: dict | None = None,
                             session_id: str | None = None) -> str | None:
    """Handle OpenAI SSE response, including tool-use agentic loop."""
    collected_text = []
    tool_calls_map = {}  # index -> {id, name, arguments_str}
    _usage_in = 0
    _usage_out = 0
    finish_reason = None

    for line in response:
        if escape_watcher and escape_watcher.cancelled:
            raise TaskCancelled()
        line = line.decode("utf-8").strip()
        if not line.startswith("data: "):
            continue
        payload_str = line[6:]
        if payload_str == "[DONE]":
            break
        try:
            event = json.loads(payload_str)
            # Extract usage from final chunk (OpenAI stream_options)
            usage = event.get("usage")
            if usage:
                _usage_in = usage.get("prompt_tokens", 0)
                _usage_out = usage.get("completion_tokens", 0)
            choices = event.get("choices", [])
            if not choices:
                continue
            # Track finish reason for max_tokens recovery
            fr = choices[0].get("finish_reason")
            if fr:
                finish_reason = fr
            delta = choices[0].get("delta") or {}
            content = delta.get("content")
            if content:
                content = _unescape(content)
                if not silent:
                    print(content, end="", flush=True)
                if event_callback:
                    event_callback("text_delta", {"text": content})
                collected_text.append(content)

            # Accumulate tool calls
            for tc in delta.get("tool_calls", []):
                idx = tc.get("index", 0)
                if idx not in tool_calls_map:
                    tool_calls_map[idx] = {
                        "id": tc.get("id", ""),
                        "name": tc.get("function", {}).get("name", ""),
                        "arguments": "",
                    }
                if tc.get("id"):
                    tool_calls_map[idx]["id"] = tc["id"]
                if tc.get("function", {}).get("name"):
                    tool_calls_map[idx]["name"] = tc["function"]["name"]
                tool_calls_map[idx]["arguments"] += tc.get("function", {}).get("arguments", "")

        except json.JSONDecodeError:
            pass

    full_text = "".join(collected_text)

    # Parse gemma4-style tool calls from raw text (oMLX doesn't convert these)
    if not tool_calls_map and "<|tool_call>" in full_text:
        parsed, cleaned = _parse_gemma_tool_calls(full_text)
        if parsed:
            for i, tc in enumerate(parsed):
                tool_calls_map[i] = {
                    "id": tc["id"],
                    "name": tc["name"],
                    "arguments": json.dumps(tc["input"]),
                }
            full_text = cleaned.strip()
            collected_text = [full_text] if full_text else []

    # Log cost for this API call
    _log_call_cost(model, _usage_in, _usage_out, session_id, _tool_round)

    # Emit usage event so callers can capture token counts
    if event_callback:
        event_callback("usage", {"tokens_in": _usage_in, "tokens_out": _usage_out})

    # End LLM trace span (OpenAI handler)
    _llm_span_oai = getattr(_thread_local, 'current_trace_span', None)
    if _trace_manager and _llm_span_oai and _llm_span_oai.get("type") == "llm_call":
        _trace_manager.end_span(_llm_span_oai, status="ok",
                                 tokens_in=_usage_in, tokens_out=_usage_out)

    # Max output token recovery (Phase 2): if model hit output limit, auto-resume
    if finish_reason == "length" and not tool_calls_map and full_text:
        recovery_count = getattr(_thread_local, '_max_output_recovery_count', 0)
        if recovery_count < MAX_OUTPUT_RECOVERY_LIMIT:
            _thread_local._max_output_recovery_count = recovery_count + 1
            if event_callback:
                event_callback("max_tokens_recovery", {
                    "attempt": recovery_count + 1,
                    "max_attempts": MAX_OUTPUT_RECOVERY_LIMIT,
                })
            # Build continuation: assistant partial + resume prompt
            messages.append({"role": "assistant", "content": full_text})
            messages.append({"role": "user", "content": _MAX_OUTPUT_RESUME_MSG})
            return send_message(messages, model, api_key, base_url, api_type,
                                silent=silent, tools=tools, escape_watcher=escape_watcher,
                                _tool_round=_tool_round, event_callback=event_callback,
                                inference_params=inference_params, session_id=session_id)

    if not tool_calls_map:
        if not silent and full_text:
            print()
        # End request trace span for final response
        _req_span_oai = getattr(_thread_local, 'request_trace_span', None)
        if _trace_manager and _req_span_oai:
            _trace_manager.end_span(_req_span_oai, status="ok",
                                     tokens_in=_usage_in, tokens_out=_usage_out)
            _thread_local.request_trace_span = None
        # Reset recovery counter on successful completion
        _thread_local._max_output_recovery_count = 0
        return full_text

    if full_text:
        print()

    # Build assistant message with tool_calls
    assistant_msg = {"role": "assistant", "content": full_text or None}
    tc_list = []
    for idx in sorted(tool_calls_map.keys()):
        tc = tool_calls_map[idx]
        tc_list.append({
            "id": tc["id"],
            "type": "function",
            "function": {"name": tc["name"], "arguments": tc["arguments"]},
        })
    assistant_msg["tool_calls"] = tc_list
    messages.append(assistant_msg)

    # Execute tools with concurrent-safe parallelism (Phase 5)
    batch_calls = []
    for tc in tc_list:
        try:
            args = json.loads(tc["function"]["arguments"])
        except json.JSONDecodeError:
            args = {}
        batch_calls.append({"id": tc["id"], "name": tc["function"]["name"], "input": args})

    batch_results = _execute_tools_batch(batch_calls, event_callback=event_callback)

    for br in batch_results:
        sanitized = _sanitize_tool_result(
            next((bc["name"] for bc in batch_calls if bc["id"] == br["tool_use_id"]), ""),
            br["result"],
        )
        messages.append({
            "role": "tool",
            "tool_call_id": br["tool_use_id"],
            "content": sanitized,
        })

    # Run middleware pipeline (context management, compaction, cancel check)
    _sid = session_id or getattr(_thread_local, 'current_session_id', None) or ""
    messages, should_continue = _run_middleware(
        messages, _tool_round, event_callback,
        model=model, api_key=api_key, base_url=base_url, api_type=api_type,
        session_id=_sid, escape_watcher=escape_watcher,
    )

    return send_message(messages, model, api_key, base_url, api_type,
                        silent=silent, tools=tools, escape_watcher=escape_watcher,
                        _tool_round=_tool_round + 1, event_callback=event_callback,
                        inference_params=inference_params, session_id=session_id)


# --- Mistral SDK integration (replicates Vibe CLI behavior) ---

_VIBE_VERSION = "2.5.0"

def _get_mistral_vibe_headers(session_id: str | None = None) -> dict:
    """Build Vibe CLI-compatible HTTP headers for Mistral API calls."""
    return {
        "user-agent": f"mistral-client-python/Mistral-Vibe/{_VIBE_VERSION}",
        "x-affinity": session_id or "brain-agent",
    }

def _get_mistral_vibe_metadata(session_id: str | None = None) -> dict:
    """Build Vibe CLI-compatible metadata for Mistral API calls."""
    return {
        "agent_entrypoint": "cli",
        "agent_version": _VIBE_VERSION,
        "client_name": "vibe_cli",
        "client_version": _VIBE_VERSION,
        "session_id": session_id or "brain-agent",
        "is_user_prompt": "true",
        "call_type": "main_call",
    }

def _create_mistral_client(api_key: str):
    """Create a Mistral SDK client instance."""
    from mistralai.client import Mistral
    return Mistral(api_key=api_key)


def _handle_mistral_response(payload, messages, model, api_key,
                              base_url, api_type, silent, tools,
                              escape_watcher=None,
                              _tool_round: int = 0,
                              event_callback=None,
                              inference_params: dict | None = None,
                              session_id: str | None = None) -> str | None:
    """Handle Mistral SDK streaming response, including tool-use agentic loop.

    Uses the official mistralai SDK with Vibe CLI-compatible headers/metadata
    to replicate the Vibe CLI's API interaction pattern.
    """
    collected_text = []
    tool_calls_map = {}  # index -> {id, name, arguments_str}
    _usage_in = 0
    _usage_out = 0
    finish_reason = None

    client = _create_mistral_client(api_key)
    vibe_headers = _get_mistral_vibe_headers(session_id)
    vibe_metadata = _get_mistral_vibe_metadata(session_id)

    # Build SDK call kwargs from payload
    sdk_kwargs = {
        "model": payload["model"],
        "messages": payload["messages"],
        "max_tokens": payload.get("max_tokens"),
        "http_headers": vibe_headers,
        "metadata": vibe_metadata,
    }
    if payload.get("tools"):
        sdk_kwargs["tools"] = payload["tools"]
    if payload.get("temperature") is not None:
        sdk_kwargs["temperature"] = payload["temperature"]
    if payload.get("top_p") is not None:
        sdk_kwargs["top_p"] = payload["top_p"]
    if payload.get("reasoning_effort"):
        sdk_kwargs["reasoning_effort"] = payload["reasoning_effort"]

    try:
        stream = client.chat.stream(**sdk_kwargs)
    except Exception as e:
        error_msg = f"Mistral SDK error: {e}"
        print(error_msg, file=sys.stderr)
        _thread_local._last_send_error = {"code": 0, "message": error_msg}
        return None

    try:
        for event in stream:
            if escape_watcher and escape_watcher.cancelled:
                raise TaskCancelled()

            data = event.data
            if not data or not data.choices:
                # Check for usage on non-choice events
                if data and data.usage:
                    _usage_in = data.usage.prompt_tokens or 0
                    _usage_out = data.usage.completion_tokens or 0
                continue

            choice = data.choices[0]

            # Track finish reason
            if choice.finish_reason:
                finish_reason = choice.finish_reason

            # Extract usage when available
            if data.usage:
                _usage_in = data.usage.prompt_tokens or 0
                _usage_out = data.usage.completion_tokens or 0

            delta = choice.delta
            if not delta:
                continue

            # Text content
            content = delta.content
            if content:
                content = _unescape(content)
                if not silent:
                    print(content, end="", flush=True)
                if event_callback:
                    event_callback("text_delta", {"text": content})
                collected_text.append(content)

            # Tool calls (same structure as OpenAI)
            if delta.tool_calls:
                for tc in delta.tool_calls:
                    # SDK tool call deltas use index-based accumulation
                    idx = getattr(tc, 'index', 0) or 0
                    if idx not in tool_calls_map:
                        tool_calls_map[idx] = {
                            "id": getattr(tc, 'id', '') or "",
                            "name": "",
                            "arguments": "",
                        }
                    if getattr(tc, 'id', None):
                        tool_calls_map[idx]["id"] = tc.id
                    fn = getattr(tc, 'function', None)
                    if fn:
                        if getattr(fn, 'name', None):
                            tool_calls_map[idx]["name"] = fn.name
                        if getattr(fn, 'arguments', None):
                            tool_calls_map[idx]["arguments"] += fn.arguments

    except TaskCancelled:
        raise
    except Exception as e:
        error_msg = f"Mistral stream error: {e}"
        print(error_msg, file=sys.stderr)
        # If we got partial text, continue with it
        if not collected_text:
            _thread_local._last_send_error = {"code": 0, "message": error_msg}
            return None

    full_text = "".join(collected_text)

    # Log cost for this API call
    _log_call_cost(model, _usage_in, _usage_out, session_id, _tool_round)

    # Emit usage event
    if event_callback:
        event_callback("usage", {"tokens_in": _usage_in, "tokens_out": _usage_out})

    # End LLM trace span
    _llm_span = getattr(_thread_local, 'current_trace_span', None)
    if _trace_manager and _llm_span and _llm_span.get("type") == "llm_call":
        _trace_manager.end_span(_llm_span, status="ok",
                                 tokens_in=_usage_in, tokens_out=_usage_out)

    # Max output token recovery
    if finish_reason == "length" and not tool_calls_map and full_text:
        recovery_count = getattr(_thread_local, '_max_output_recovery_count', 0)
        if recovery_count < MAX_OUTPUT_RECOVERY_LIMIT:
            _thread_local._max_output_recovery_count = recovery_count + 1
            if event_callback:
                event_callback("max_tokens_recovery", {
                    "attempt": recovery_count + 1,
                    "max_attempts": MAX_OUTPUT_RECOVERY_LIMIT,
                })
            messages.append({"role": "assistant", "content": full_text})
            messages.append({"role": "user", "content": _MAX_OUTPUT_RESUME_MSG})
            return send_message(messages, model, api_key, base_url, api_type,
                                silent=silent, tools=tools, escape_watcher=escape_watcher,
                                _tool_round=_tool_round, event_callback=event_callback,
                                inference_params=inference_params, session_id=session_id)

    if not tool_calls_map:
        if not silent and full_text:
            print()
        # End request trace span
        _req_span = getattr(_thread_local, 'request_trace_span', None)
        if _trace_manager and _req_span:
            _trace_manager.end_span(_req_span, status="ok",
                                     tokens_in=_usage_in, tokens_out=_usage_out)
            _thread_local.request_trace_span = None
        _thread_local._max_output_recovery_count = 0
        return full_text

    if full_text:
        print()

    # Build assistant message with tool_calls (OpenAI format)
    assistant_msg = {"role": "assistant", "content": full_text or None}
    tc_list = []
    for idx in sorted(tool_calls_map.keys()):
        tc = tool_calls_map[idx]
        tc_list.append({
            "id": tc["id"],
            "type": "function",
            "function": {"name": tc["name"], "arguments": tc["arguments"]},
        })
    assistant_msg["tool_calls"] = tc_list
    messages.append(assistant_msg)

    # Execute tools (same as OpenAI path)
    batch_calls = []
    for tc in tc_list:
        try:
            args = json.loads(tc["function"]["arguments"])
        except json.JSONDecodeError:
            args = {}
        batch_calls.append({"id": tc["id"], "name": tc["function"]["name"], "input": args})

    batch_results = _execute_tools_batch(batch_calls, event_callback=event_callback)

    for br in batch_results:
        sanitized = _sanitize_tool_result(
            next((bc["name"] for bc in batch_calls if bc["id"] == br["tool_use_id"]), ""),
            br["result"],
        )
        messages.append({
            "role": "tool",
            "tool_call_id": br["tool_use_id"],
            "content": sanitized,
        })

    # Run middleware pipeline
    _sid = session_id or getattr(_thread_local, 'current_session_id', None) or ""
    messages, should_continue = _run_middleware(
        messages, _tool_round, event_callback,
        model=model, api_key=api_key, base_url=base_url, api_type=api_type,
        session_id=_sid, escape_watcher=escape_watcher,
    )

    return send_message(messages, model, api_key, base_url, api_type,
                        silent=silent, tools=tools, escape_watcher=escape_watcher,
                        _tool_round=_tool_round + 1, event_callback=event_callback,
                        inference_params=inference_params, session_id=session_id)


def _classify_error_transient(error_info: dict | None) -> bool:
    """Check if the last send error is transient (retryable) vs permanent."""
    if not error_info:
        return True  # Unknown error, assume transient
    if error_info.get("permanent"):
        return False
    code = error_info.get("code", 0)
    # Transient: 429, 500, 502, 503, 504, 529, connection errors (code=0)
    return code in {0, 429, 500, 502, 503, 504, 529}


def _retry_with_backoff(messages, model, api_key, base_url, api_type,
                        silent, tools, escape_watcher, event_callback,
                        inference_params, session_id, max_retries=2):
    """Try sending a message with exponential backoff retries for transient errors.

    Returns (result, last_error_info). result is None if all retries failed.
    """
    _thread_local._last_send_error = None
    result = send_message(messages, model, api_key, base_url, api_type,
                          silent=silent, tools=tools, escape_watcher=escape_watcher,
                          event_callback=event_callback,
                          inference_params=inference_params,
                          session_id=session_id)
    if result is not None:
        return result, None

    error_info = getattr(_thread_local, '_last_send_error', None)
    # If permanent error, don't retry
    if not _classify_error_transient(error_info):
        return None, error_info

    # Retry with exponential backoff
    for attempt in range(1, max_retries + 1):
        delay = min(1.0 * (2 ** (attempt - 1)), 30.0) + random.uniform(0, 0.5)
        error_msg = (error_info or {}).get("message", "unknown error")
        print(f"  Retrying {model} in {delay:.1f}s (attempt {attempt}/{max_retries}, error: {error_msg})", flush=True)
        if event_callback:
            event_callback("fallback", {
                "status": "retry",
                "model": model,
                "attempt": attempt,
                "max_retries": max_retries,
                "delay": round(delay, 1),
                "reason": error_msg,
            })
        time.sleep(delay)

        # Check for cancellation
        if escape_watcher and escape_watcher.cancelled:
            from claude_cli import TaskCancelled
            raise TaskCancelled()

        _thread_local._last_send_error = None
        result = send_message(messages, model, api_key, base_url, api_type,
                              silent=silent, tools=tools, escape_watcher=escape_watcher,
                              event_callback=event_callback,
                              inference_params=inference_params,
                              session_id=session_id)
        if result is not None:
            return result, None

        error_info = getattr(_thread_local, '_last_send_error', None)
        if not _classify_error_transient(error_info):
            break  # Permanent error, stop retrying

    return None, error_info


def send_message_with_fallback(messages: list[dict], model: str, api_key: str,
                               base_url: str, api_type: str,
                               silent: bool = False,
                               tools: bool = True,
                               escape_watcher=None,
                               event_callback=None,
                               provider_resolver=None,
                               inference_params: dict | None = None,
                               purpose: str | None = None,
                               session_id: str | None = None) -> str | None:
    """Send messages with retry + fallback chain.

    Retry logic: transient errors (502, 503, 429, timeout) retry 2x with exponential backoff.
    Permanent errors (400, 401, 404) skip retries and go straight to fallback.
    Fallbacks field in _models_config: ordered list of fallback model IDs.
    If provider_resolver is provided, it's called with (model) -> {api_key, base_url, api_type}.
    Emits ("fallback", {...}) events via event_callback for UI display.
    """
    # Track which model actually responded (for done event)
    _thread_local._fallback_model_used = None

    # Snapshot message count before primary attempt — if it fails mid-tool-loop,
    # intermediate messages must be stripped before trying fallback models
    msg_count_original = len(messages)

    # Try primary model with retries
    result, error_info = _retry_with_backoff(
        messages, model, api_key, base_url, api_type,
        silent, tools, escape_watcher, event_callback,
        inference_params, session_id, max_retries=2)
    if result is not None:
        return result

    # Strip any intermediate tool-loop messages from failed primary attempt
    if len(messages) > msg_count_original:
        del messages[msg_count_original:]

    primary_error = (error_info or {}).get("message", "unknown error")

    # Build fallback list — use explicit fallbacks from config, then capability-aware auto-fallback
    fallback_models = []
    if _models_config:
        model_cfg = _models_config.get(model, {})
        explicit_fallbacks = model_cfg.get("fallbacks", [])
        if explicit_fallbacks:
            fallback_models = list(explicit_fallbacks)
        else:
            # Auto-build fallback order: same provider first, then by priority
            failed_provider = model_cfg.get("provider", "")
            failed_caps = set(model_cfg.get("capabilities", []))
            candidates = []
            for mid, cfg in _models_config.items():
                if mid == model or not cfg.get("enabled", True):
                    continue
                same_provider = 1 if cfg.get("provider") == failed_provider else 0
                matching_caps = len(failed_caps & set(cfg.get("capabilities", [])))
                priority = cfg.get("priority", 0)
                # Sort key: same provider first, then capability match, then priority
                candidates.append((mid, same_provider, matching_caps, priority))
            candidates.sort(key=lambda x: (x[1], x[2], x[3]), reverse=True)
            fallback_models = [mid for mid, _, _, _ in candidates]
    else:
        fallback_models = get_available_models(api_key, base_url, api_type)

    if not fallback_models:
        msg = f"Error: Model '{model}' is not available and no fallback models found."
        print(msg, file=sys.stderr)
        if event_callback:
            raise RuntimeError(msg)
        sys.exit(1)

    tried_models = {model}
    for fallback_model in fallback_models:
        if fallback_model in tried_models:
            continue
        tried_models.add(fallback_model)

        # Re-resolve provider for fallback model
        fb_api_key, fb_base_url, fb_api_type = api_key, base_url, api_type
        if provider_resolver:
            try:
                prov = provider_resolver(fallback_model)
                fb_api_key = prov.get("api_key", api_key)
                fb_base_url = prov.get("base_url", base_url)
                fb_api_type = prov.get("api_type", api_type)
            except Exception:
                continue

        print(f"Note: Model '{model}' failed, trying fallback '{fallback_model}'.", flush=True)
        if event_callback:
            event_callback("fallback", {
                "status": "switch",
                "from": model,
                "to": fallback_model,
                "reason": primary_error,
            })

        fb_params = get_inference_params(fallback_model, purpose)
        # Snapshot message count: if fallback uses different api_type, tool-loop
        # messages will be in the wrong format and must be stripped after success
        msg_count_before = len(messages)
        result, fb_error = _retry_with_backoff(
            messages, fallback_model, fb_api_key, fb_base_url, fb_api_type,
            silent, tools, escape_watcher, event_callback,
            fb_params, session_id, max_retries=1)
        if result is not None:
            # Strip intermediate tool-loop messages appended by the fallback model.
            # These corrupt the session: different api_type → wrong message format,
            # same api_type → thinking blocks with invalid/missing signatures.
            # The final text reply is returned; server adds it properly to session.
            if len(messages) > msg_count_before:
                del messages[msg_count_before:]
            _thread_local._fallback_model_used = fallback_model
            # Log fallback event to cost tracker
            if _cost_tracker:
                try:
                    agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                    agent_id = agent.agent_id if agent else "main"
                    logging.info(f"Fallback: {model} -> {fallback_model} for agent {agent_id} (reason: {primary_error})")
                except Exception:
                    pass
            return result

    msg = f"Error: No working models found. Tried: {', '.join(tried_models)}"
    print(msg, file=sys.stderr)
    if event_callback:
        raise RuntimeError(msg)
    sys.exit(1)


# --- TUI helpers ---

def _draw_status_bar(model: str, history: list[dict] | None = None,
                     max_tokens: int = DEFAULT_MAX_CONTEXT_TOKENS) -> None:
    """Draw a status bar on the last terminal line with black background."""
    cols = shutil.get_terminal_size().columns
    rows = shutil.get_terminal_size().lines

    # Context usage
    ctx_part = ""
    ctx_visible = 0
    if history is not None:
        est = _estimate_conversation_tokens(history)
        pct = min(99, int(est / max_tokens * 100))
        if pct >= 75:
            color = RED
        elif pct >= 50:
            color = YELLOW
        else:
            color = FG_GRAY
        # Show token count in k for readability
        if est >= 1000:
            tok_str = f"{est // 1000}k"
        else:
            tok_str = str(est)
        ctx_label = f"{tok_str}/{max_tokens // 1000}k"
        ctx_part = f" {DIM}│{RESET}{BG_DARK} {color}{ctx_label}{RESET}{BG_DARK}"
        ctx_visible = 4 + len(ctx_label)  # " │ Nk/Nk"

    # Agent name
    agent_part = ""
    agent_visible = 0
    if _current_agent and _current_agent.agent_id != "main":
        agent_part = f" {CYAN}{_current_agent.agent_id}{RESET}{BG_DARK} {DIM}│{RESET}{BG_DARK}"
        agent_visible = 1 + len(_current_agent.agent_id) + 3  # " name │"

    label = f" {agent_part}{FG_GRAY}Model:{RESET}{BG_DARK} {GREEN}{BOLD}{model}{RESET}{BG_DARK}{ctx_part} "
    visible_len = 1 + agent_visible + 8 + len(model) + ctx_visible + 1
    padding = max(0, cols - visible_len)
    bar = f"\033[48;5;235m{label}{' ' * padding}{RESET}"
    sys.stdout.write(f"\0337\033[{rows};1H{bar}\0338")
    sys.stdout.flush()


def _setup_scroll_region() -> None:
    """Reserve the bottom line for the status bar."""
    rows = shutil.get_terminal_size().lines
    sys.stdout.write(f"\033[1;{rows - 1}r")
    sys.stdout.write(f"\033[1;1H")
    sys.stdout.flush()


def _restore_scroll_region() -> None:
    """Restore full terminal scroll region."""
    rows = shutil.get_terminal_size().lines
    sys.stdout.write(f"\033[1;{rows}r")
    sys.stdout.write(f"\033[{rows};1H\033[K")
    sys.stdout.flush()


# --- Main ---

def main():
    parser = argparse.ArgumentParser(
        description=f"Brain Agent v{VERSION} — Agentic CLI for LLM APIs"
    )
    parser.add_argument(
        "message", nargs="?",
        help="Message to send (or use -i for interactive mode)",
    )
    parser.add_argument(
        "-m", "--model", default="claude-opus-4-5-20251101",
        help="Model to use (default: claude-opus-4-5-20251101)",
    )
    parser.add_argument(
        "-i", "--interactive", action="store_true",
        help="Interactive mode - continuous chat",
    )
    parser.add_argument(
        "-l", "--list-models", action="store_true",
        help="List available models and exit",
    )
    parser.add_argument(
        "--api-key", default="sk-Xk7kOHpIpZkLutwnyxHpRO9jn4ZwyPaS",
        help="API key for authentication",
    )
    parser.add_argument(
        "--base-url", default="http://localhost:8317/v1",
        help="Base URL for the API (default: http://localhost:8317/v1)",
    )
    parser.add_argument(
        "-t", "--api-type", choices=["anthropic", "openai", "mistral"], default="anthropic",
        help="API type: anthropic or openai (default: anthropic)",
    )
    parser.add_argument(
        "--max-context", type=int, default=DEFAULT_MAX_CONTEXT_TOKENS,
        help=f"Max context window in tokens (default: {DEFAULT_MAX_CONTEXT_TOKENS})",
    )
    parser.add_argument(
        "--agent", default="main",
        help="Agent ID to start with (default: 'main')",
    )

    args = parser.parse_args()

    if args.list_models:
        list_models(args.api_key, args.base_url, args.api_type)
        sys.exit(0)

    if args.interactive:
        _run_interactive(args)
        sys.exit(0)

    # Single-message mode
    if not args.message:
        parser.print_help()
        sys.exit(1)

    # Try SDK sidecar first for consistent tool loop
    reply = None
    try:
        import sdk_backend
        if sdk_backend.is_sidecar_running():
            # Set up agent context for system prompt
            agent_cfg = AgentConfig(args.agent)
            _thread_local.current_agent = agent_cfg
            _thread_local.memory_store = MemoryStore(args.agent)
            system_prompt = _build_system_prompt(include_memory_summary=True)
            provider_env = sdk_backend.build_provider_env(args.model)

            # Build tool defs for sidecar MCP
            tool_defs = []
            for td in TOOL_DEFINITIONS:
                if td["name"] in ("read_file", "write_file", "edit_file",
                                   "list_directory", "search_files", "execute_command",
                                   "web_fetch"):
                    continue  # SDK has native equivalents
                if td["name"] not in TOOL_DISPATCH:
                    continue
                desc = td["description"]
                if isinstance(desc, tuple):
                    desc = " ".join(desc)
                tool_defs.append({
                    "name": td["name"],
                    "description": desc[:1000],
                    "input_schema": td["input_schema"],
                })

            reply = sdk_backend.query_sync(
                prompt=args.message, model=args.model,
                system_prompt=system_prompt, max_turns=30,
                tool_defs=tool_defs,
                agent_id=args.agent,
            )
    except Exception:
        pass  # Fall through to direct API

    # Fallback: direct API call
    if reply is None:
        messages = [{"role": "user", "content": args.message}]
        reply = send_message_with_fallback(
            messages, args.model, args.api_key, args.base_url, args.api_type,
            silent=True)
    if reply:
        print(render_markdown(reply))


def _print_greeting(model: str, agent_id: str = "default") -> None:
    """Print the Brain Agent startup banner."""
    cwd = os.getcwd()
    latest = CHANGELOG[0] if CHANGELOG else None

    # Color definitions for gradient brain
    C1 = "\033[38;5;213m"  # pink
    C2 = "\033[38;5;177m"  # purple
    C3 = "\033[38;5;141m"  # lavender
    C4 = "\033[38;5;105m"  # blue-purple
    C5 = "\033[38;5;69m"   # blue
    C6 = "\033[38;5;33m"   # deep blue
    CO = "\033[38;5;208m"  # orange accent

    brain = [
        f"  {C1}    ⣀⣀⣤⣤⣤⣤⣤⣤⣀⣀    {RESET}",
        f"  {C1}  ⣴⣿⣿⣿⣿⣿⣿⣿⣿⣿⣿⣿⣦  {RESET}",
        f"  {C2} ⣾⣿⣿⡟⠛⠛⣿⣿⡟⠛⠛⣿⣿⣿⣷ {RESET}",
        f"  {C2}⣸⣿⣿⡇  ⣸⣿⣿⡇  ⢸⣿⣿⣿⣇{RESET}",
        f"  {C3}⣿⣿⣿⣇  ⣿⣿⣿⣇  ⣸⣿⣿⣿⣿{RESET}",
        f"  {C3}⢿⣿⣿⣿⣦⣤⣿⣿⣿⣦⣤⣾⣿⣿⣿⡿{RESET}",
        f"  {C4} ⠻⣿⣿⣿⣿⣿⣿⣿⣿⣿⣿⣿⣿⠟ {RESET}",
        f"  {C5}  ⠙⢿⣿⣿⣿⣿⣿⣿⣿⣿⡿⠋  {RESET}",
        f"  {C6}    ⠉⠛⠿⣿⣿⣿⠿⠛⠉    {RESET}",
    ]

    # Title with gradient
    title = (
        f"  {C1}B{C2}r{C3}a{C4}i{C5}n{RESET} "
        f"{CO}{BOLD}Agent{RESET}"
    )
    ver = f" {DIM}v{VERSION}{RESET}"
    agent_label = f" {DIM}│{RESET} {CYAN}{BOLD}{agent_id}{RESET}" if agent_id != "main" else ""

    print()
    for line in brain:
        print(line)
    print()
    print(f"  {title}{ver}{agent_label}")
    print()

    # Info section with dim separators
    sep = f"{DIM}·{RESET}"
    print(f"  {DIM}Model{RESET}  {GREEN}{model}{RESET}")
    print(f"  {DIM}Path{RESET}   {DIM}{cwd}{RESET}")

    # Agents
    agents = list_agents()
    if len(agents) > 1:
        agent_list = []
        for a in agents:
            if a == agent_id:
                agent_list.append(f"{CYAN}{BOLD}{a}{RESET}")
            else:
                agent_list.append(f"{DIM}{a}{RESET}")
        print(f"  {DIM}Agents{RESET} {' {0} '.format(sep).join(agent_list)}")

    # Skills count
    if _current_agent:
        skills = _current_agent.list_skills()
        if skills:
            skill_names = [s["name"] for s in skills[:5]]
            more = f" {DIM}+{len(skills)-5} more{RESET}" if len(skills) > 5 else ""
            print(f"  {DIM}Skills{RESET} {DIM}{', '.join(skill_names)}{more}{RESET}")

    # Scheduled tasks count
    if _scheduler:
        schedules = _scheduler.list_all()
        active = sum(1 for s in schedules if s["enabled"])
        if schedules:
            print(f"  {DIM}Tasks{RESET}  {DIM}{active} scheduled ({len(schedules)} total){RESET}")

    print()
    print(f"  {DIM}Commands  /new /agent /model /models /tools /schedule{RESET}")
    print(f"  {DIM}Controls  Esc cancel · o toggle tools · exit quit{RESET}")

    if latest:
        print(f"\n  {DIM}↑ v{latest[0]}: {latest[2]}{RESET}")

    print()


# --- Slash commands registry ---

SLASH_COMMANDS = {
    "/help":     "Show this help",
    "/about":    "Show version and changelog",
    "/new":      "Start a new conversation",
    "/agent":    "Switch agent or list agents",
    "/model":    "Switch model",
    "/models":   "List available models",
    "/tools":    "Toggle tool call display",
    "/schedule": "Manage scheduled tasks (list/add/pause/resume/delete/history)",
}


def _print_help() -> None:
    """Print help for all slash commands."""
    print()
    print(f"  {BOLD}Commands{RESET}")
    print()
    for cmd, desc in SLASH_COMMANDS.items():
        print(f"  {GREEN}{BOLD}{cmd:12s}{RESET} {DIM}{desc}{RESET}")
    print()
    print(f"  {BOLD}Schedule subcommands{RESET}")
    print()
    for sub, desc in [
        ("list", "List all scheduled tasks"),
        ("add", "Create a new scheduled task"),
        ("pause NAME", "Pause a task"),
        ("resume NAME", "Resume a task"),
        ("delete NAME", "Delete a task"),
        ("history", "Show execution history"),
    ]:
        print(f"  {GREEN}{BOLD}  {sub:16s}{RESET} {DIM}{desc}{RESET}")
    print()
    print(f"  {BOLD}Keyboard{RESET}")
    print()
    for key, desc in [
        ("Esc", "Cancel current request"),
        ("Tab", "Autocomplete slash commands"),
        ("Up/Down", "Input history"),
        ("Ctrl+A/E", "Beginning/end of line"),
        ("Ctrl+W", "Delete word backward"),
        ("Ctrl+U", "Clear line"),
        ("o", "Toggle tool output (when hidden)"),
    ]:
        print(f"  {YELLOW}{key:12s}{RESET} {DIM}{desc}{RESET}")
    print()


def _print_about() -> None:
    """Print version info and changelog."""
    print()
    print(f"  {BOLD}Brain Agent{RESET}  {GREEN}{BOLD}v{VERSION}{RESET}  {DIM}{VERSION_DATE}{RESET}")
    print()
    print(f"  {BOLD}Changelog{RESET}")
    print()
    for v, date, changes in CHANGELOG[:8]:
        print(f"  {GREEN}{BOLD}v{v}{RESET}  {DIM}{date}{RESET}")
        # Word-wrap changes at ~72 chars
        words = changes.split()
        line = "    "
        for word in words:
            if len(line) + len(word) + 1 > 76:
                print(f"{DIM}{line}{RESET}")
                line = "    " + word
            else:
                line = line + (" " if line.strip() else "") + word
        if line.strip():
            print(f"{DIM}{line}{RESET}")
        print()


def _select_menu(items: list[str], prompt: str = "Select",
                 labels: list[str] | None = None,
                 active: str | None = None) -> str | None:
    """Interactive arrow-key selection menu. Returns selected item or None on Esc/Ctrl-C."""
    if not items:
        return None
    display = labels if labels and len(labels) == len(items) else items
    selected = 0
    # Find active item
    if active and active in items:
        selected = items.index(active)

    fd = sys.stdin.fileno()
    try:
        old_settings = termios.tcgetattr(fd)
    except termios.error:
        return None

    # Count lines we'll draw so we can clean up
    menu_lines = len(items)

    try:
        new_settings = termios.tcgetattr(fd)
        new_settings[3] = new_settings[3] & ~(termios.ICANON | termios.ECHO)
        new_settings[6][termios.VMIN] = 0
        new_settings[6][termios.VTIME] = 0
        termios.tcsetattr(fd, termios.TCSANOW, new_settings)

        # Draw initial menu
        print(f"\n  {DIM}{prompt}:{RESET}")
        for i, label in enumerate(display):
            _draw_menu_item(i, label, i == selected, items[i] == active if active else False)
        sys.stdout.flush()

        while True:
            if select.select([fd], [], [], 0.1)[0]:
                ch = os.read(fd, 1)
                if ch == b'\r' or ch == b'\n':  # Enter
                    break
                elif ch == b'\x1b':  # Escape or arrow
                    seq1 = os.read(fd, 1) if select.select([fd], [], [], 0.05)[0] else b''
                    if seq1 == b'[':
                        seq2 = os.read(fd, 1) if select.select([fd], [], [], 0.05)[0] else b''
                        if seq2 == b'A':  # Up
                            selected = (selected - 1) % len(items)
                        elif seq2 == b'B':  # Down
                            selected = (selected + 1) % len(items)
                        else:
                            # Bare Esc or unknown sequence — cancel
                            _erase_menu(menu_lines + 1)
                            return None
                    else:
                        # Bare Esc — cancel
                        _erase_menu(menu_lines + 1)
                        return None
                elif ch == b'\x03':  # Ctrl-C
                    _erase_menu(menu_lines + 1)
                    return None
                else:
                    continue

                # Redraw menu
                # Move cursor up to start of menu
                sys.stdout.write(f"\033[{menu_lines}A")
                for i, label in enumerate(display):
                    sys.stdout.write(f"\r\033[K")
                    _draw_menu_item(i, label, i == selected, items[i] == active if active else False)
                sys.stdout.flush()

        # Erase menu after selection
        _erase_menu(menu_lines + 1)
        return items[selected]

    except Exception:
        return None
    finally:
        try:
            termios.tcsetattr(fd, termios.TCSADRAIN, old_settings)
        except termios.error:
            pass


def _draw_menu_item(idx: int, label: str, is_selected: bool, is_active: bool):
    """Draw a single menu item line."""
    marker = f"{CYAN}{BOLD}❯{RESET}" if is_selected else " "
    active_tag = f" {GREEN}(active){RESET}" if is_active else ""
    if is_selected:
        print(f"  {marker} {BOLD}{label}{RESET}{active_tag}")
    else:
        print(f"  {marker} {DIM}{label}{RESET}{active_tag}")


def _erase_menu(lines: int):
    """Erase N lines above cursor."""
    for _ in range(lines):
        sys.stdout.write(f"\033[A\r\033[K")
    sys.stdout.flush()


def _readline(prompt: str, input_history: list[str], history_idx_ref: list[int],
              completions: list[str] | None = None) -> str | None:
    """Read a line with arrow-key history navigation and inline editing.

    Returns the entered string, or None on Ctrl-C / Ctrl-D.
    input_history is a list of previous inputs (newest last).
    history_idx_ref is a single-element list holding the current browse index.
    """
    fd = sys.stdin.fileno()
    old_settings = termios.tcgetattr(fd)
    try:
        tty.setraw(fd)
        sys.stdout.write(prompt)
        sys.stdout.flush()

        buf = []        # current line characters
        pos = 0         # cursor position within buf
        saved_line = "" # stash current input when browsing history
        hist_idx = len(input_history)  # start past the end (= new input)

        while True:
            ch = os.read(fd, 1)

            if ch == b'\r' or ch == b'\n':  # Enter
                # Move cursor to end, print newline
                if pos < len(buf):
                    sys.stdout.write(f"\033[{len(buf) - pos}C")
                sys.stdout.write("\r\n")
                sys.stdout.flush()
                return "".join(buf)

            elif ch == b'\x03':  # Ctrl-C
                sys.stdout.write("\r\n")
                sys.stdout.flush()
                return None

            elif ch == b'\x04':  # Ctrl-D
                if not buf:
                    sys.stdout.write("\r\n")
                    sys.stdout.flush()
                    return None
                # Otherwise ignore

            elif ch == b'\x7f' or ch == b'\x08':  # Backspace
                if pos > 0:
                    buf.pop(pos - 1)
                    pos -= 1
                    # Redraw from cursor position
                    tail = "".join(buf[pos:])
                    sys.stdout.write(f"\033[D{tail} \033[{len(tail) + 1}D")
                    sys.stdout.flush()

            elif ch == b'\x15':  # Ctrl-U: clear line
                if pos > 0:
                    sys.stdout.write(f"\033[{pos}D")
                sys.stdout.write(" " * len(buf))
                if len(buf) > 0:
                    sys.stdout.write(f"\033[{len(buf)}D")
                buf.clear()
                pos = 0
                sys.stdout.flush()

            elif ch == b'\x01':  # Ctrl-A: beginning of line
                if pos > 0:
                    sys.stdout.write(f"\033[{pos}D")
                    pos = 0
                    sys.stdout.flush()

            elif ch == b'\x05':  # Ctrl-E: end of line
                if pos < len(buf):
                    sys.stdout.write(f"\033[{len(buf) - pos}C")
                    pos = len(buf)
                    sys.stdout.flush()

            elif ch == b'\x17':  # Ctrl-W: delete word backward
                if pos > 0:
                    old_pos = pos
                    # Skip trailing spaces
                    while pos > 0 and buf[pos - 1] == ' ':
                        pos -= 1
                    # Skip word
                    while pos > 0 and buf[pos - 1] != ' ':
                        pos -= 1
                    deleted = old_pos - pos
                    del buf[pos:old_pos]
                    sys.stdout.write(f"\033[{deleted}D")
                    tail = "".join(buf[pos:])
                    sys.stdout.write(f"{tail}{' ' * deleted}")
                    sys.stdout.write(f"\033[{len(tail) + deleted}D")
                    sys.stdout.flush()

            elif ch == b'\x1b':  # Escape sequence
                seq1 = os.read(fd, 1)
                if seq1 == b'[':
                    seq2 = os.read(fd, 1)

                    if seq2 == b'A':  # Up arrow
                        if hist_idx > 0:
                            if hist_idx == len(input_history):
                                saved_line = "".join(buf)
                            hist_idx -= 1
                            _replace_line(buf, pos, input_history[hist_idx])
                            buf[:] = list(input_history[hist_idx])
                            pos = len(buf)

                    elif seq2 == b'B':  # Down arrow
                        if hist_idx < len(input_history):
                            hist_idx += 1
                            if hist_idx == len(input_history):
                                _replace_line(buf, pos, saved_line)
                                buf[:] = list(saved_line)
                            else:
                                _replace_line(buf, pos, input_history[hist_idx])
                                buf[:] = list(input_history[hist_idx])
                            pos = len(buf)

                    elif seq2 == b'C':  # Right arrow
                        if pos < len(buf):
                            sys.stdout.write("\033[C")
                            pos += 1
                            sys.stdout.flush()

                    elif seq2 == b'D':  # Left arrow
                        if pos > 0:
                            sys.stdout.write("\033[D")
                            pos -= 1
                            sys.stdout.flush()

                    elif seq2 == b'H':  # Home
                        if pos > 0:
                            sys.stdout.write(f"\033[{pos}D")
                            pos = 0
                            sys.stdout.flush()

                    elif seq2 == b'F':  # End
                        if pos < len(buf):
                            sys.stdout.write(f"\033[{len(buf) - pos}C")
                            pos = len(buf)
                            sys.stdout.flush()

                    elif seq2 == b'3':  # Delete key (ESC [ 3 ~)
                        seq3 = os.read(fd, 1)  # consume '~'
                        if pos < len(buf):
                            buf.pop(pos)
                            tail = "".join(buf[pos:])
                            sys.stdout.write(f"{tail} \033[{len(tail) + 1}D")
                            sys.stdout.flush()

                # Bare Escape — ignore (don't break input)

            elif ch == b'\t':  # Tab — autocomplete
                current = "".join(buf)
                comp_list = completions or list(SLASH_COMMANDS.keys())
                if current.startswith("/"):
                    matches = [c for c in comp_list if c.startswith(current)]
                    if len(matches) == 1:
                        # Single match — complete it
                        completion = matches[0]
                        # Add space after completed command
                        if not completion.endswith(" "):
                            completion += " "
                        _replace_line(buf, pos, completion)
                        buf[:] = list(completion)
                        pos = len(buf)
                    elif len(matches) > 1:
                        # Multiple matches — find common prefix
                        prefix = os.path.commonprefix(matches)
                        if len(prefix) > len(current):
                            _replace_line(buf, pos, prefix)
                            buf[:] = list(prefix)
                            pos = len(buf)
                        else:
                            # Show options briefly below
                            sys.stdout.write(f"\n  {DIM}{' '.join(matches)}{RESET}")
                            sys.stdout.write(f"\033[A")  # move back up
                            # Redraw prompt + buffer
                            sys.stdout.write(f"\r{prompt}{''.join(buf)}")
                            sys.stdout.flush()

            elif ch >= b' ':  # Printable character
                char = ch.decode("utf-8", errors="replace")
                buf.insert(pos, char)
                pos += 1
                tail = "".join(buf[pos:])
                sys.stdout.write(f"{char}{tail}")
                if tail:
                    sys.stdout.write(f"\033[{len(tail)}D")
                sys.stdout.flush()

    finally:
        termios.tcsetattr(fd, termios.TCSADRAIN, old_settings)


def _replace_line(buf: list, pos: int, new_text: str) -> None:
    """Clear the current input line and write new_text."""
    # Move cursor to start of input
    if pos > 0:
        sys.stdout.write(f"\033[{pos}D")
    # Clear old content
    sys.stdout.write(" " * len(buf))
    if len(buf) > 0:
        sys.stdout.write(f"\033[{len(buf)}D")
    # Write new content
    sys.stdout.write(new_text)
    sys.stdout.flush()



def _switch_agent(agent_id: str, args) -> tuple[str, AgentConfig]:
    """Switch to a different agent. Returns (model, agent_config)."""
    global _current_agent, _memory_store, _mcp_manager
    agent = AgentConfig(agent_id)
    _current_agent = agent
    _memory_store = MemoryStore(agent_id=agent_id, base_dir=agent.memory_dir)

    # Load MCP servers: reuse shared manager if available, else create one
    if not _mcp_manager:
        _mcp_manager = MCPManager()
    else:
        _mcp_manager.stop_all()
    main_mcp = os.path.join(AGENTS_DIR, "main", "mcp.json")
    _mcp_manager.load_config(main_mcp)
    if agent_id != "main":
        _mcp_manager.load_config(agent.mcp_config_path)

    model = agent.preferred_model or args.model
    return model, agent


def _run_interactive(args):
    """Run the interactive TUI chat loop."""
    global _memory_store, _current_agent
    global _delegate_fallback_model, _delegate_api_key, _delegate_base_url, _delegate_api_type

    history = []
    input_history = []   # list of previous user inputs for arrow-key recall
    history_idx = [0]    # mutable ref for current position

    # Store API config for delegation
    _delegate_api_key = args.api_key
    _delegate_base_url = args.base_url
    _delegate_api_type = args.api_type
    _delegate_fallback_model = args.model

    # SDK session tracking (for resume across turns)
    _run_interactive._sdk_sid = None

    # Initialize agent
    current_model, _ = _switch_agent(args.agent, args)

    # Start scheduler
    global _scheduler
    _scheduler = Scheduler()
    _scheduler.start()

    # Ensure memory summary schedules
    try:
        ensure_memory_summary_schedules()
    except Exception:
        pass

    # Initialize background task runner
    global _task_runner
    _task_runner = TaskRunner()

    # Clear screen and move cursor to top
    sys.stdout.write("\033[2J\033[H")
    sys.stdout.flush()

    _setup_scroll_region()
    _draw_status_bar(current_model, history, args.max_context)

    # Handle terminal resize
    def _on_resize(signum, frame):
        _setup_scroll_region()
        _draw_status_bar(current_model, history, args.max_context)

    old_sigwinch = signal.signal(signal.SIGWINCH, _on_resize)

    # Startup greeting
    _print_greeting(current_model, args.agent)

    try:
        while True:
            # Read input with history support
            message = _readline(f"\n{BOLD}{GREEN}❯{RESET} ", input_history, history_idx)
            if message is None:
                print(f"{DIM}Bye!{RESET}")
                break

            stripped = message.strip().lower()

            if stripped in ("exit", "quit"):
                print(f"{DIM}Bye!{RESET}")
                break

            if stripped == "/help":
                _print_help()
                continue

            if stripped == "/about":
                _print_about()
                continue

            if stripped == "/new":
                history = []
                print(f"\n{DIM}{'─' * 40}{RESET}")
                print(f"{DIM}  New chat started{RESET}")
                print(f"{DIM}{'─' * 40}{RESET}")
                _draw_status_bar(current_model, history, args.max_context)
                continue

            if stripped == "/tools":
                global _show_tools
                _show_tools = not _show_tools
                state = f"{GREEN}visible{RESET}" if _show_tools else f"{DIM}hidden{RESET}"
                print(f"  {DIM}Tool display:{RESET} {state}")
                continue

            if stripped.startswith("/agent"):
                arg = message.strip()[6:].strip()
                agents = list_agents()
                if arg:
                    if arg not in agents:
                        print(f"  {DIM}Creating new agent:{RESET} {BOLD}{arg}{RESET}")
                    current_model, _ = _switch_agent(arg, args)
                    history = []
                    print(f"  {DIM}Switched to agent:{RESET} {BOLD}{arg}{RESET} {DIM}(model: {current_model}){RESET}")
                else:
                    # Build labels with descriptions
                    labels = []
                    for aid in agents:
                        cfg = AgentConfig(aid)
                        model_info = f" [{cfg.preferred_model}]" if cfg.preferred_model else ""
                        labels.append(f"{aid}{model_info} — {cfg.description}")
                    current_aid = _current_agent.agent_id if _current_agent else "main"
                    choice = _select_menu(agents, "Select agent", labels=labels, active=current_aid)
                    if choice:
                        current_model, _ = _switch_agent(choice, args)
                        history = []
                        print(f"  {DIM}Switched to agent:{RESET} {BOLD}{choice}{RESET} {DIM}(model: {current_model}){RESET}")
                _draw_status_bar(current_model, history, args.max_context)
                continue

            if stripped.startswith("/schedule"):
                arg = message.strip()[9:].strip()
                if not _scheduler:
                    print(f"  {DIM}Scheduler not running{RESET}")
                    continue

                if arg == "" or arg == "list":
                    # List schedules
                    schedules = _scheduler.list_all()
                    if not schedules:
                        print(f"\n  {DIM}No scheduled tasks{RESET}")
                    else:
                        print()
                        for s in schedules:
                            status = f"{GREEN}active{RESET}" if s["enabled"] else f"{DIM}paused{RESET}"
                            next_r = s.get("next_run", "")[:16] if s.get("next_run") else "—"
                            print(f"  {BOLD}{s['name']}{RESET} [{status}] {DIM}{s['schedule']}{RESET}")
                            print(f"    {DIM}agent:{RESET} {s['agent']}  {DIM}next:{RESET} {next_r}")
                            print(f"    {DIM}task:{RESET} {s['task'][:60]}")
                    continue

                elif arg == "add":
                    # Interactive add
                    print(f"\n  {DIM}Add scheduled task{RESET}")
                    name = _readline(f"  {DIM}Name:{RESET} ", [], [0])
                    if not name or not name.strip():
                        continue
                    name = name.strip()
                    task = _readline(f"  {DIM}Task:{RESET} ", [], [0])
                    if not task or not task.strip():
                        continue
                    task = task.strip()
                    schedule = _readline(f"  {DIM}Schedule (every Xm/Xh/Xd, daily HH:MM, weekly DOW HH:MM):{RESET} ", [], [0])
                    if not schedule or not schedule.strip():
                        continue
                    schedule = schedule.strip()
                    agent = _readline(f"  {DIM}Agent (default: main):{RESET} ", [], [0])
                    agent = agent.strip() if agent and agent.strip() else "main"
                    model_in = _readline(f"  {DIM}Model (default: current):{RESET} ", [], [0])
                    model_val = model_in.strip() if model_in and model_in.strip() else None

                    result = _scheduler.add(name, task, schedule, agent, model_val)
                    if result.get("error"):
                        print(f"  {RED}{result['error']}{RESET}")
                    else:
                        print(f"  {GREEN}✔ Created:{RESET} {BOLD}{name}{RESET} — next run: {result.get('next_run', '')[:16]}")
                    continue

                elif arg.startswith("pause "):
                    name = arg[6:].strip()
                    result = _scheduler.pause(name)
                    if result.get("error"):
                        print(f"  {RED}{result['error']}{RESET}")
                    else:
                        print(f"  {DIM}Paused:{RESET} {name}")
                    continue

                elif arg.startswith("resume "):
                    name = arg[7:].strip()
                    result = _scheduler.resume(name)
                    if result.get("error"):
                        print(f"  {RED}{result['error']}{RESET}")
                    else:
                        print(f"  {GREEN}Resumed:{RESET} {name} — next: {result.get('next_run', '')[:16]}")
                    continue

                elif arg.startswith("delete ") or arg.startswith("rm "):
                    name = arg.split(" ", 1)[1].strip()
                    result = _scheduler.remove(name)
                    if result.get("error"):
                        print(f"  {RED}{result['error']}{RESET}")
                    else:
                        print(f"  {DIM}Deleted:{RESET} {name}")
                    continue

                elif arg == "history":
                    history_items = _scheduler.get_history(limit=10)
                    if not history_items:
                        print(f"\n  {DIM}No execution history{RESET}")
                    else:
                        print()
                        for h in history_items:
                            status_color = GREEN if h["status"] == "success" else RED
                            print(f"  {status_color}{h['status']}{RESET} {BOLD}{h['schedule_name']}{RESET} {DIM}({h['finished_at'][:16]}){RESET}")
                            if h.get("result"):
                                preview = h["result"][:80].replace("\n", " ")
                                print(f"    {DIM}{preview}{RESET}")
                    continue

                else:
                    print(f"  {DIM}Usage: /schedule [list|add|pause NAME|resume NAME|delete NAME|history]{RESET}")
                    continue

            if stripped == "/models":
                models = get_available_models(args.api_key, args.base_url, args.api_type)
                if models:
                    choice = _select_menu(models, "Select model", active=current_model)
                    if choice:
                        current_model = choice
                        print(f"  {DIM}Switched to:{RESET} {BOLD}{current_model}{RESET}")
                        _draw_status_bar(current_model, history, args.max_context)
                else:
                    print(f"  {DIM}No models available{RESET}")
                continue

            if stripped.startswith("/model"):
                arg = message.strip()[6:].strip()
                models = get_available_models(args.api_key, args.base_url, args.api_type)
                if arg:
                    current_model = arg
                    print(f"  {DIM}Switched to:{RESET} {BOLD}{current_model}{RESET}")
                elif models:
                    choice = _select_menu(models, "Select model", active=current_model)
                    if choice:
                        current_model = choice
                        print(f"  {DIM}Switched to:{RESET} {BOLD}{current_model}{RESET}")
                else:
                    print(f"  {DIM}No models available. Use: /model <name>{RESET}")
                _draw_status_bar(current_model, history, args.max_context)
                continue

            # Custom slash commands (from agent's commands.json + .claude/commands/*.md)
            if message.strip().startswith("/"):
                cmd_name = message.strip().split()[0][1:].lower()
                cmd_args = message.strip()[len(cmd_name)+2:].strip()
                agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                if agent:
                    for cmd in agent.load_commands():
                        if (cmd.get("name", "").lower() == cmd_name or
                                cmd.get("slug", "").lower() == cmd_name):
                            message = AgentConfig.expand_command(cmd, cmd_args)
                            print(f"  {DIM}Running /{cmd_name}{RESET}")
                            break

            if not message.strip():
                continue

            # Save to input history (dedup consecutive)
            if not input_history or input_history[-1] != message.strip():
                input_history.append(message.strip())

            # Send message
            history.append({"role": "user", "content": message})
            _reset_tool_tracking()

            # Check context window and compact if needed
            history, was_compacted = _check_and_compact(
                history, current_model, args.api_key, args.base_url,
                args.api_type, max_tokens=args.max_context,
            )

            # Start spinner and escape watcher
            spinner = Spinner(current_model)
            escape_watcher = EscapeWatcher()
            spinner.start()
            escape_watcher.start()

            cancelled = False
            reply = None
            try:
                # Try SDK sidecar first
                _sdk_ok = False
                try:
                    import sdk_backend
                    if sdk_backend.is_sidecar_running():
                        system_prompt = _build_system_prompt(include_memory_summary=True)
                        # Build tool defs (skip SDK-native tools)
                        _SDK_NATIVE = {"read_file", "write_file", "edit_file",
                                       "list_directory", "search_files",
                                       "execute_command", "web_fetch"}
                        tool_defs = []
                        for td in TOOL_DEFINITIONS:
                            if td["name"] in _SDK_NATIVE or td["name"] not in TOOL_DISPATCH:
                                continue
                            desc = td["description"]
                            if isinstance(desc, tuple):
                                desc = " ".join(desc)
                            tool_defs.append({
                                "name": td["name"],
                                "description": desc[:1000],
                                "input_schema": td["input_schema"],
                            })
                        agent = getattr(_thread_local, 'current_agent', None) or _current_agent
                        agent_id = agent.agent_id if agent else "main"

                        meta = sdk_backend.query_sync(
                            prompt=message, model=current_model,
                            system_prompt=system_prompt, max_turns=30,
                            tool_defs=tool_defs,
                            agent_id=agent_id,
                            cancel_fn=lambda: escape_watcher.cancelled,
                            sdk_session_id=getattr(_run_interactive, '_sdk_sid', None),
                            return_metadata=True,
                        )
                        if meta is not None:
                            reply = meta.get("text")
                            # Track SDK session for resume on next turn
                            if meta.get("sdk_session_id"):
                                _run_interactive._sdk_sid = meta["sdk_session_id"]
                            _sdk_ok = True
                        elif escape_watcher.cancelled:
                            cancelled = True
                            _sdk_ok = True
                except Exception:
                    pass  # Fall through to direct API

                # Fallback: direct API call
                if not _sdk_ok and not cancelled:
                    reply = send_message_with_fallback(
                        history, current_model, args.api_key, args.base_url,
                        args.api_type, silent=True, escape_watcher=escape_watcher)
            except TaskCancelled:
                cancelled = True
                reply = None
            finally:
                elapsed = spinner.stop()
                escape_watcher.stop()

            if cancelled:
                # Remove the user message that was cancelled
                history.pop()
                print(f"\n{DIM}✘ Cancelled (Esc){RESET}")
                _draw_status_bar(current_model, history, args.max_context)
                continue

            if reply:
                history.append({"role": "assistant", "content": reply})

                # Render formatted response
                print()
                rendered = render_markdown(reply)
                print(rendered)

                # Completion message
                verb = spinner.verb.rstrip("ing")
                # Make past tense
                past = spinner.verb[:-3] + "ed" if spinner.verb.endswith("ing") else spinner.verb
                if spinner.verb == "Thinking":
                    past = "Thought"
                elif spinner.verb == "Weaving":
                    past = "Woven"
                elif spinner.verb == "Computing":
                    past = "Computed"
                elif spinner.verb == "Brewing":
                    past = "Brewed"
                elif spinner.verb == "Baking":
                    past = "Baked"
                elif spinner.verb == "Crafting":
                    past = "Crafted"
                elif spinner.verb == "Conjuring":
                    past = "Conjured"
                elif spinner.verb == "Composing":
                    past = "Composed"
                elif spinner.verb == "Contemplating":
                    past = "Contemplated"
                elif spinner.verb == "Pondering":
                    past = "Pondered"

                print(f"\n{DIM}✻ {past} for {elapsed:.0f}s{RESET}")
            else:
                print(f"\n{DIM}(no response){RESET}")

            # Restore status bar after response
            _draw_status_bar(current_model, history, args.max_context)

    finally:
        if _scheduler:
            _scheduler.stop()
        if _mcp_manager:
            _mcp_manager.stop_all()
        signal.signal(signal.SIGWINCH, old_sigwinch)
        _restore_scroll_region()


if __name__ == "__main__":
    main()
